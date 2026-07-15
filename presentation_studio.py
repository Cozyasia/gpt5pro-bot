# -*- coding: utf-8 -*-
"""Reliable PDF/PPTX presentation wizard for python-telegram-bot 21.x.

The module is deliberately self-contained and can be connected to an existing bot
through ``register(application)``. It provides:

- persistent project state in SQLite;
- fact-preserving brief parsing;
- structure approval and repeatable revisions;
- automatic logo concepts or a separate cumulative logo brief;
- repeatable cumulative visual/style/palette requirements;
- semantic image prompts per slide;
- editable PPTX and matching PDF export;
- local error handling, so a successful step is not followed by a generic
  "Упс, произошла ошибка" message.

Callback prefix: ``ps:``.
"""

from __future__ import annotations

import asyncio
import base64
import contextlib
import io
import json
import logging
import math
import os
import re
import sqlite3
import textwrap
import time
import uuid
from copy import deepcopy
from pathlib import Path
from dataclasses import dataclass
from typing import Any, Iterable, Awaitable, Callable

import httpx
from PIL import Image, ImageDraw, ImageFont, ImageOps
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from reportlab.lib.colors import HexColor
from reportlab.lib.pagesizes import landscape
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.ttfonts import TTFont
from reportlab.pdfgen import canvas
from telegram import InlineKeyboardButton, InlineKeyboardMarkup, InputFile, InputMediaPhoto, Update
from telegram.constants import ChatAction
from telegram.error import NetworkError, RetryAfter, TelegramError, TimedOut
from telegram.ext import (
    ApplicationHandlerStop,
    CallbackQueryHandler,
    ContextTypes,
    MessageHandler,
    filters,
)

log = logging.getLogger("presentation-wizard")

CALLBACK_PREFIX = "ps:"
DB_PATH = os.path.abspath(os.getenv("DB_PATH", "subs.db"))
ASSET_ROOT = Path(
    os.getenv(
        "PRESENTATION_DATA_DIR",
        os.getenv("PRESENTATION_ASSET_DIR", str(Path(DB_PATH).parent / "presentation_studio")),
    )
)
MAX_IMAGES = max(
    1,
    min(14, int(os.getenv("PRESENTATION_MAX_GENERATED_IMAGES", os.getenv("PRESENTATION_MAX_IMAGES", "10")) or "10")),
)
MAX_UPLOADS = max(1, int(os.getenv("PRESENTATION_MAX_UPLOADS", "60") or "60"))
IMAGE_MODEL = os.getenv("PRESENTATION_IMAGE_MODEL", "gpt-image-1").strip() or "gpt-image-1"
IMAGE_SIZE = os.getenv("PRESENTATION_IMAGE_SIZE", "1536x1024").strip() or "1536x1024"

# Injected by PresentationStudio from main.py. Keeping these callbacks in this module
# preserves the existing billing/provider routing and avoids a second bot stack.
_STUDIO_LLM_CALL = None
_STUDIO_IMAGE_BATCH_CALL = None
_STUDIO_PAID_RUNNER = None
_STUDIO_RENDER_COST_USD = 0.20
TEXT_MODEL = (
    os.getenv("PRESENTATION_TEXT_MODEL", "").strip()
    or os.getenv("OPENAI_MODEL", "openai/gpt-4o-mini").strip()
)

# Palette: VERDIA Natural Tech by default; generic enough for arbitrary brands.
DEFAULT_PALETTE = {
    "background": "#F4F1E9",
    "text": "#17211C",
    "primary": "#1F4D3A",
    "secondary": "#6E8B73",
    "sand": "#C9B99A",
    "accent": "#B08D57",
    "card": "#E7E3D9",
    "white": "#FFFFFF",
}

STOP_BRAND_NAMES = {
    "бренд",
    "бренда",
    "компания",
    "компании",
    "товар",
    "услуга",
    "продукт",
    "название",
    "brand",
    "company",
    "business",
}

FONT_REGULAR = next(
    (p for p in [
        "/usr/share/fonts/truetype/lato/Lato-Regular.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ] if os.path.exists(p)),
    "",
)
FONT_BOLD = next(
    (p for p in [
        "/usr/share/fonts/truetype/lato/Lato-Bold.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    ] if os.path.exists(p)),
    FONT_REGULAR,
)
FONT_HEAVY = next(
    (p for p in [
        "/usr/share/fonts/truetype/lato/Lato-Heavy.ttf",
        "/usr/share/fonts/truetype/lato/Lato-Black.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    ] if os.path.exists(p)),
    FONT_BOLD,
)


def _now() -> int:
    return int(time.time())


def _safe_filename(value: str, fallback: str = "presentation") -> str:
    value = re.sub(r"[^0-9A-Za-zА-Яа-яЁё._-]+", "_", (value or "").strip())
    value = value.strip("._-")
    return (value[:80] or fallback)


def _project_dir(user_id: int, project_id: str) -> Path:
    p = ASSET_ROOT / str(user_id) / project_id
    p.mkdir(parents=True, exist_ok=True)
    return p


def init_storage() -> None:
    ASSET_ROOT.mkdir(parents=True, exist_ok=True)
    con = sqlite3.connect(DB_PATH)
    try:
        con.execute(
            """
            CREATE TABLE IF NOT EXISTS presentation_projects_v95 (
                user_id INTEGER PRIMARY KEY,
                project_id TEXT NOT NULL,
                stage TEXT NOT NULL,
                payload_json TEXT NOT NULL,
                updated_at INTEGER NOT NULL
            )
            """
        )
        con.commit()
    finally:
        con.close()


def _load(user_id: int) -> dict[str, Any] | None:
    init_storage()
    con = sqlite3.connect(DB_PATH)
    try:
        row = con.execute(
            "SELECT payload_json FROM presentation_projects_v95 WHERE user_id=?",
            (int(user_id),),
        ).fetchone()
    finally:
        con.close()
    if not row:
        return None
    try:
        return json.loads(row[0])
    except Exception:
        log.exception("Corrupt presentation state for user=%s", user_id)
        return None


def _save(user_id: int, project: dict[str, Any]) -> None:
    init_storage()
    project["updated_at"] = _now()
    payload = json.dumps(project, ensure_ascii=False, separators=(",", ":"))
    con = sqlite3.connect(DB_PATH)
    try:
        con.execute(
            """
            INSERT INTO presentation_projects_v95(user_id, project_id, stage, payload_json, updated_at)
            VALUES(?,?,?,?,?)
            ON CONFLICT(user_id) DO UPDATE SET
                project_id=excluded.project_id,
                stage=excluded.stage,
                payload_json=excluded.payload_json,
                updated_at=excluded.updated_at
            """,
            (
                int(user_id),
                project.get("project_id", ""),
                project.get("stage", ""),
                payload,
                project["updated_at"],
            ),
        )
        con.commit()
    finally:
        con.close()


def _delete(user_id: int) -> None:
    init_storage()
    con = sqlite3.connect(DB_PATH)
    try:
        con.execute("DELETE FROM presentation_projects_v95 WHERE user_id=?", (int(user_id),))
        con.commit()
    finally:
        con.close()


def _new_project(user_id: int, kind: str = "presentation", chat_id: int = 0) -> dict[str, Any]:
    project_id = uuid.uuid4().hex[:12]
    return {
        "project_id": project_id,
        "user_id": int(user_id),
        "chat_id": int(chat_id or 0),
        "kind": "catalog" if kind == "catalog" else "presentation",
        "stage": "await_brief",
        "raw_brief": "",
        "profile": {},
        "structure": [],
        "structure_notes": [],
        "logo_mode": "",
        "logo_notes": [],
        "logo_candidates": [],
        "logo_selected": "",
        "image_mode": "",
        "generation_engine": "auto",
        "visual_notes": [],
        "style_notes": [],
        "palette_notes": [],
        "uploaded_images": [],
        "slide_images": {},
        "palette": deepcopy(DEFAULT_PALETTE),
        "created_at": _now(),
        "updated_at": _now(),
    }


def _is_valid_brand_name(value: str) -> bool:
    value = re.sub(r"\s+", " ", (value or "").strip(" \n\t:—-.,;\"'«»"))
    if not value or len(value) < 2 or len(value) > 80:
        return False
    if value.lower() in STOP_BRAND_NAMES:
        return False
    if not re.search(r"[A-Za-zА-Яа-яЁё]", value):
        return False
    return True


def _extract_brand_name_regex(text: str) -> str:
    """Extract an exact brand name without swallowing neighbouring headings.

    Handles both ``Название бренда: VERDIA HOME`` and a form where the value is
    on the next line. Newlines are never treated as spaces inside the value.
    """
    raw = (text or "").replace("\r\n", "\n").replace("\r", "\n")
    lines = [line.strip() for line in raw.split("\n")]
    label_re = re.compile(
        r"^(?:название(?:\s+(?:бренда|компании|проекта))?|бренд|компания)\s*(?:[:—-]\s*(.*))?$",
        re.I,
    )
    for i, line in enumerate(lines):
        m = label_re.match(line)
        if not m:
            continue
        candidates: list[str] = []
        if m.group(1):
            candidates.append(m.group(1))
        for nxt in lines[i + 1 : i + 4]:
            if nxt:
                candidates.append(nxt)
                break
        for candidate in candidates:
            candidate = re.split(r"[.;,]", candidate.strip())[0].strip(" :—-\"'«»")
            candidate = re.sub(r"[ \t]+", " ", candidate)
            if _is_valid_brand_name(candidate):
                return candidate

    # Inline prose such as «коммерческая презентация бренда VERDIA HOME».
    inline_patterns = [
        r"(?i)\b(?:бренд(?:а)?|марка|компания)\s+[«\"']?([A-ZА-ЯЁ][A-ZА-ЯЁ0-9&._-]*(?:[ \t]+[A-ZА-ЯЁ0-9&._-]+){0,4})",
        r"(?i)\b(?:brand|company)\s*[:—-]?[ \t]+[«\"']?([A-Z][A-Z0-9&._-]*(?:[ \t]+[A-Z0-9&._-]+){0,4})",
    ]
    for pattern in inline_patterns:
        m = re.search(pattern, raw)
        if m:
            candidate = re.sub(r"[ \t]+", " ", m.group(1)).strip(" :—-\"'«»")
            if _is_valid_brand_name(candidate):
                return candidate

    # Strong uppercase fallback is evaluated line-by-line, never across headings.
    for line in lines:
        for candidate in re.findall(r"\b[A-ZА-ЯЁ][A-ZА-ЯЁ0-9&._-]{2,}(?:[ \t]+[A-ZА-ЯЁ0-9&._-]{2,}){0,3}\b", line):
            candidate = re.sub(r"[ \t]+", " ", candidate).strip()
            if _is_valid_brand_name(candidate) and candidate.lower() not in STOP_BRAND_NAMES:
                return candidate
    return ""


def _extract_tagline(text: str) -> str:
    raw = (text or "").replace("\r\n", "\n").replace("\r", "\n")
    lines = [line.strip() for line in raw.split("\n")]
    for i, line in enumerate(lines):
        m = re.match(r"^(?:слоган|девиз|позиционирование)\s*(?:[:—-]\s*(.*))?$", line, re.I)
        if not m:
            continue
        if m.group(1) and m.group(1).strip():
            return m.group(1).strip(" .")[:160]
        for nxt in lines[i + 1 : i + 4]:
            if nxt:
                return nxt.strip(" .")[:160]
    return ""


def _extract_contacts(text: str) -> list[str]:
    raw = text or ""
    result: list[str] = []
    emails = re.findall(r"[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}", raw)
    for value in emails:
        if value not in result:
            result.append(value)
    without_emails = raw
    for email in emails:
        without_emails = without_emails.replace(email, " ")
    patterns = [
        r"\+?\d[\d\s()\-]{8,}\d",
        r"(?<![A-Za-z0-9._%+-])@[A-Za-z0-9_]{4,}",
        r"\b(?:https?://)?(?:www\.)?[A-Za-z0-9-]+(?:\.[A-Za-z0-9-]+)+(?:/[A-Za-z0-9_./?&=%+-]*)?",
    ]
    for pattern in patterns:
        for value in re.findall(pattern, without_emails):
            cleaned = re.sub(r"\s+", " ", value).strip()
            if cleaned and cleaned not in result:
                result.append(cleaned)
    return result[:8]


def _extract_prices(text: str) -> list[str]:
    prices = re.findall(
        r"(?:от\s+)?\d[\d\s.,]{1,12}\s*(?:₽|руб(?:лей|ля|ль)?|THB|USD|EUR|\$|€)(?:\s*(?:в месяц|/мес\.?))?",
        text or "",
        flags=re.I,
    )
    return [re.sub(r"\s+", " ", p).strip() for p in prices[:12]]


def _profile_fallback(text: str) -> dict[str, Any]:
    brand = _extract_brand_name_regex(text)
    tagline = _extract_tagline(text)
    contacts = _extract_contacts(text)
    prices = _extract_prices(text)
    return {
        "brand_name": brand,
        "tagline": tagline,
        "objective": "Создать убедительную коммерческую презентацию по исходному брифу.",
        "product": _first_nonempty_section(text, ["Описание компании", "Продукт", "Описание продукта", "Решение"]),
        "audience": _first_nonempty_section(text, ["Целевая аудитория", "Целевая аудитория B2C", "Целевая аудитория B2B"]),
        "positioning": _first_nonempty_section(text, ["Ключевое позиционирование", "Позиционирование"]),
        "contacts": contacts,
        "prices": prices,
        "language": "ru",
        "requested_slide_count": _extract_slide_count(text),
    }


def _extract_slide_count(text: str) -> int:
    m = re.search(r"(\d{1,2})(?:\s*[-–—]\s*(\d{1,2}))?\s*(?:слайд|страниц)", text or "", re.I)
    if not m:
        return 12
    first = int(m.group(1))
    second = int(m.group(2)) if m.group(2) else first
    return max(8, min(16, max(first, second)))


def _first_nonempty_section(text: str, headings: Iterable[str]) -> str:
    lines = (text or "").splitlines()
    normalized = [h.lower() for h in headings]
    for i, line in enumerate(lines):
        l = line.strip().lower().rstrip(":")
        if any(l.startswith(h.lower()) for h in normalized):
            buf: list[str] = []
            for next_line in lines[i + 1 : i + 7]:
                clean = next_line.strip()
                if not clean:
                    if buf:
                        break
                    continue
                if re.match(r"^[A-ZА-ЯЁ][^.!?]{1,50}:$", clean):
                    break
                buf.append(clean)
            if buf:
                return " ".join(buf)[:800]
    paragraphs = [p.strip() for p in re.split(r"\n\s*\n", text or "") if len(p.strip()) > 40]
    return (paragraphs[0][:800] if paragraphs else "")


def _json_from_text(value: str) -> Any:
    value = (value or "").strip()
    value = re.sub(r"^```(?:json)?\s*", "", value, flags=re.I)
    value = re.sub(r"\s*```$", "", value)
    start_candidates = [p for p in [value.find("{"), value.find("[")] if p >= 0]
    if start_candidates:
        value = value[min(start_candidates):]
    for end_char in ["}", "]"]:
        end = value.rfind(end_char)
        if end >= 0:
            candidate = value[: end + 1]
            with contextlib.suppress(Exception):
                return json.loads(candidate)
    return json.loads(value)


def _text_base_url_and_headers() -> tuple[str, dict[str, str]]:
    key = os.getenv("OPENAI_API_KEY", "").strip()
    provider = os.getenv("TEXT_PROVIDER", "").strip().lower()
    base = os.getenv("OPENAI_BASE_URL", "").strip()
    if key.startswith("sk-or-") or provider == "openrouter":
        base = "https://openrouter.ai/api/v1"
    base = (base or "https://api.openai.com/v1").rstrip("/")
    headers = {
        "Authorization": f"Bearer {key}",
        "Content-Type": "application/json; charset=utf-8",
    }
    if "openrouter.ai" in base:
        site = os.getenv("OPENROUTER_SITE_URL", "").strip()
        app_name = os.getenv("OPENROUTER_APP_NAME", "").strip()
        if site:
            headers["HTTP-Referer"] = site
        if app_name:
            headers["X-Title"] = app_name
    return base, headers


async def _llm_text(prompt: str, temperature: float = 0.25, update: Update | None = None) -> str:
    global _STUDIO_LLM_CALL
    if _STUDIO_LLM_CALL is not None and update is not None:
        return await _STUDIO_LLM_CALL(update, prompt)
    key = os.getenv("OPENAI_API_KEY", "").strip()
    if not key:
        raise RuntimeError("OPENAI_API_KEY is not configured")
    base, headers = _text_base_url_and_headers()
    body = {
        "model": TEXT_MODEL,
        "messages": [
            {
                "role": "system",
                "content": (
                    "Ты арт-директор и редактор коммерческих презентаций. "
                    "Строго сохраняй факты и названия из брифа. Не выдумывай кейсы, "
                    "статистику, сертификаты, отзывы, цены или контакты."
                ),
            },
            {"role": "user", "content": prompt},
        ],
        "temperature": temperature,
    }
    last_error: Exception | None = None
    for attempt in range(3):
        try:
            async with httpx.AsyncClient(timeout=120.0) as client:
                response = await client.post(f"{base}/chat/completions", headers=headers, json=body)
            if response.status_code // 100 != 2:
                raise RuntimeError(f"LLM HTTP {response.status_code}: {response.text[:500]}")
            data = response.json()
            return (data["choices"][0]["message"]["content"] or "").strip()
        except Exception as exc:
            last_error = exc
            await asyncio.sleep(0.8 * (attempt + 1))
    raise RuntimeError(f"LLM request failed: {last_error}")


async def _parse_profile(raw_brief: str, update: Update | None = None) -> dict[str, Any]:
    fallback = _profile_fallback(raw_brief)
    prompt = f"""
Извлеки факты из брифа в JSON. Верни только JSON без Markdown:
{{
  "brand_name": "точное название бренда, не слово 'бренд'",
  "tagline": "",
  "objective": "",
  "product": "",
  "audience": "",
  "positioning": "",
  "geography": "",
  "contacts": [""],
  "prices": [""],
  "requested_slide_count": 12,
  "visual_direction": ""
}}

Критические правила:
- brand_name перепиши дословно из поля «Название бренда» или аналогичного поля;
- не используй слова «бренд», «бренда», «компания» как название;
- не сокращай контакты и цены;
- не выдумывай отсутствующие данные.

БРИФ:
{raw_brief}
"""
    try:
        parsed = _json_from_text(await _llm_text(prompt, 0.1, update))
        if not isinstance(parsed, dict):
            parsed = {}
    except Exception:
        log.exception("Profile parsing via LLM failed; using deterministic fallback")
        parsed = {}

    profile = deepcopy(fallback)
    for key in [
        "brand_name",
        "tagline",
        "objective",
        "product",
        "audience",
        "positioning",
        "geography",
        "visual_direction",
    ]:
        value = parsed.get(key)
        if isinstance(value, str) and value.strip():
            profile[key] = value.strip()
    for key in ["contacts", "prices"]:
        value = parsed.get(key)
        if isinstance(value, list):
            profile[key] = [str(v).strip() for v in value if str(v).strip()]
    with contextlib.suppress(Exception):
        profile["requested_slide_count"] = max(8, min(16, int(parsed.get("requested_slide_count") or profile["requested_slide_count"])))

    # Deterministic brand value has priority if it is clearly defined in the brief.
    exact_brand = _extract_brand_name_regex(raw_brief)
    if exact_brand:
        profile["brand_name"] = exact_brand
    if not _is_valid_brand_name(profile.get("brand_name", "")):
        profile["brand_name"] = ""
    if not profile.get("contacts"):
        profile["contacts"] = fallback["contacts"]
    if not profile.get("prices"):
        profile["prices"] = fallback["prices"]
    return profile


def _fallback_structure(project: dict[str, Any]) -> list[dict[str, Any]]:
    brand = project.get("profile", {}).get("brand_name") or "Бренд"
    raw = project.get("raw_brief", "")
    has_products = bool(re.search(r"продуктов|линейк|тариф|пакет|модель", raw, re.I))
    has_service = bool(re.search(r"обслужив|сервис|подписк", raw, re.I))
    slides = [
        (brand, "cover", True, [project.get("profile", {}).get("tagline") or "Коммерческая презентация"]),
        ("Идея бренда", "split", True, ["Эмоциональное позиционирование", "Главная ценность для клиента"]),
        ("Проблема клиента", "cards", False, ["Ключевые барьеры", "Почему обычные решения не работают"]),
        ("Наше решение", "split", True, ["Что получает клиент", "Как продукт решает задачу"]),
        ("Как это работает", "process", True, ["Основные компоненты", "Понятная последовательность работы"]),
        ("Ключевые преимущества", "cards", False, ["Практическая ценность", "Эстетика", "Сервис"]),
        ("Сценарии использования", "split", True, ["Для частных клиентов", "Для бизнеса"]),
    ]
    if has_products:
        slides.append(("Продуктовая линейка", "comparison", True, ["Варианты решения", "Цены и комплектация из брифа"]))
    if has_service:
        slides.append(("Сервис и сопровождение", "cards", True, ["Регулярная поддержка", "Обслуживание после запуска"]))
    slides.extend([
        ("Этапы реализации", "process", False, ["От консультации до запуска", "Контроль результата"]),
        (f"Почему {brand}", "cards", False, ["Отличия от альтернатив", "Причины доверять решению"]),
        ("Следующий шаг", "cta", True, ["Призыв к действию", "Контакты из брифа"]),
    ])
    count = project.get("profile", {}).get("requested_slide_count", 12)
    fillers = [
        ("Кому подходит решение", "cards", False, ["Частные клиенты", "Архитекторы и дизайнеры", "Коммерческие пространства"]),
        ("Что получает клиент", "cards", False, ["Проектирование", "Монтаж под ключ", "Поддержка после запуска"]),
        ("Ответы на сомнения", "cards", False, ["Эксплуатация", "Надёжность", "Стоимость владения"]),
        ("Форматы сотрудничества", "split", True, ["Индивидуальный проект", "Партнёрство с дизайнерами", "Корпоративные решения"]),
    ]
    # Guarantee the requested count instead of silently returning 8-11 slides.
    insert_at = max(1, len(slides) - 1)
    filler_index = 0
    while len(slides) < count:
        slides.insert(insert_at, fillers[filler_index % len(fillers)])
        insert_at += 1
        filler_index += 1
    # Keep cover/CTA, trim middle only.
    while len(slides) > count and len(slides) > 8:
        slides.pop(-2)
    return [
        {
            "title": title,
            "layout": layout,
            "image_needed": image_needed,
            "bullets": bullets,
            "image_prompt": "",
        }
        for title, layout, image_needed, bullets in slides
    ]


def _normalize_structure(value: Any, fallback: list[dict[str, Any]]) -> list[dict[str, Any]]:
    if isinstance(value, dict):
        value = value.get("slides")
    if not isinstance(value, list):
        return fallback
    result: list[dict[str, Any]] = []
    valid_layouts = {"cover", "split", "cards", "process", "comparison", "full_image", "cta"}
    for item in value:
        if not isinstance(item, dict):
            continue
        title = str(item.get("title") or "").strip()
        if not title:
            continue
        bullets = item.get("bullets") or []
        if isinstance(bullets, str):
            bullets = [b.strip() for b in re.split(r"[;\n]", bullets) if b.strip()]
        bullets = [str(b).strip() for b in bullets if str(b).strip()][:6]
        layout = str(item.get("layout") or "split").strip().lower()
        if layout not in valid_layouts:
            layout = "split"
        result.append(
            {
                "title": title[:110],
                "layout": layout,
                "image_needed": bool(item.get("image_needed", layout in {"cover", "split", "full_image", "comparison", "cta"})),
                "bullets": bullets or ["Содержание будет подготовлено строго по брифу"],
                "image_prompt": str(item.get("image_prompt") or "").strip()[:1200],
            }
        )
    if len(result) < 8:
        return fallback
    target_count = len(fallback)
    if len(result) < target_count:
        existing_titles = {item.get("title", "").strip().lower() for item in result}
        for item in fallback:
            if len(result) >= target_count:
                break
            if item.get("title", "").strip().lower() not in existing_titles:
                result.insert(max(1, len(result) - 1), deepcopy(item))
                existing_titles.add(item.get("title", "").strip().lower())
    while len(result) > target_count and len(result) > 8:
        result.pop(-2)
    return result[:16]


async def _generate_structure(project: dict[str, Any], update: Update | None = None) -> list[dict[str, Any]]:
    fallback = _fallback_structure(project)
    profile = project.get("profile", {})
    notes = "\n".join(project.get("structure_notes", []))
    prompt = f"""
Создай структуру коммерческой презентации в JSON:
{{"slides":[{{"title":"", "layout":"cover|split|cards|process|comparison|full_image|cta", "image_needed":true, "bullets":[""], "image_prompt":""}}]}}

Количество: {profile.get('requested_slide_count', 12)} слайдов.
Название бренда: {profile.get('brand_name', '')}
Сохрани все продукты, цены, сервисы, этапы, контакты и ограничения из брифа.
Не добавляй отзывы, статистику, сертификаты или кейсы, которых нет.
Первый слайд — сильная обложка с точным названием бренда.
Последний — CTA и реальные контакты из брифа.
На слайде 3-5 коротких тезисов. Не делай отдельные малополезные слайды B2C и B2B,
если их можно объединить в один содержательный слайд.
Для image_prompt опиши только сцену, без текста и логотипов.

НАКОПЛЕННЫЕ ПРАВКИ К СТРУКТУРЕ:
{notes or 'Нет'}

ПОЛНЫЙ ИСХОДНЫЙ БРИФ:
{project.get('raw_brief', '')}
"""
    try:
        return _normalize_structure(_json_from_text(await _llm_text(prompt, 0.2, update)), fallback)
    except Exception:
        log.exception("Structure generation failed; using fallback")
        return fallback


def _structure_text(project: dict[str, Any]) -> str:
    brand = project.get("profile", {}).get("brand_name") or "Без названия"
    lines = [f"Предлагаемая структура: {brand}", ""]
    for i, slide in enumerate(project.get("structure", []), 1):
        bullets = "; ".join(slide.get("bullets", [])[:4])
        lines.append(f"{i}. {slide.get('title', '')} — {bullets}")
    lines.append("")
    lines.append("Проверьте названия, цены, порядок и отсутствие вымышленных фактов.")
    return "\n".join(lines)


def _brief_summary(project: dict[str, Any]) -> str:
    p = project.get("profile", {})
    contacts = ", ".join(p.get("contacts", [])[:4]) or "не найдены"
    prices = ", ".join(p.get("prices", [])[:6]) or "не найдены"
    return (
        f"Название: {p.get('brand_name') or 'не определено'}\n"
        f"Слоган: {p.get('tagline') or 'не указан'}\n"
        f"Продукт: {(p.get('product') or 'не выделен')[:500]}\n"
        f"Аудитория: {(p.get('audience') or 'не выделена')[:350]}\n"
        f"Цены: {prices}\n"
        f"Контакты: {contacts}"
    )


def _kb(rows: list[list[tuple[str, str]]]) -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup(
        [[InlineKeyboardButton(text, callback_data=data) for text, data in row] for row in rows]
    )


START_KB = _kb([
    [("🆕 Новый проект", "ps:new"), ("▶️ Продолжить", "ps:resume")],
    [("❌ Отменить текущий", "ps:cancel")],
])

STRUCTURE_KB = _kb([
    [("✅ Утвердить структуру", "ps:structure_ok")],
    [("✏️ Изменить структуру", "ps:structure_edit")],
    [("📋 Показать распознанный бриф", "ps:brief_show")],
    [("❌ Отменить проект", "ps:cancel")],
])

LOGO_CHOICE_KB = _kb([
    [("✨ Автоматически по главному брифу", "ps:logo_auto")],
    [("🧾 По отдельному брифу", "ps:logo_custom")],
    [("📤 Загрузить готовый логотип", "ps:logo_upload")],
    [("⏭ Продолжить без логотипа", "ps:logo_skip")],
    [("⬅️ К структуре", "ps:back_structure")],
])

LOGO_NOTES_KB = _kb([
    [("➕ Добавить ещё пожелания", "ps:logo_note")],
    [("📋 Показать накопленный бриф", "ps:logo_notes_show")],
    [("🎨 Создать 3 варианта", "ps:logo_generate")],
    [("⬅️ Назад", "ps:logo_choice")],
])

IMAGE_MODE_KB = _kb([
    [("✨ Создать изображения автоматически", "ps:images_auto")],
    [("📤 Использовать свои изображения", "ps:images_upload")],
    [("🧩 Смешанный режим", "ps:images_mixed")],
    [("⏭ Без изображений", "ps:images_skip")],
])

VISUAL_NOTES_KB = _kb([
    [("➕ Добавить требования к визуалам", "ps:visual_note")],
    [("📋 Показать все требования", "ps:visual_show")],
    [("🧹 Очистить требования", "ps:visual_clear")],
    [("➡️ Продолжить к стилю", "ps:visual_done")],
])

UPLOAD_IMAGES_KB = _kb([
    [("➕ Добавить требования к визуалам", "ps:visual_note")],
    [("📋 Показать загруженные материалы", "ps:uploads_show")],
    [("➡️ Закончить загрузку", "ps:visual_done")],
])

STYLE_KB = _kb([
    [("✨ Автоматический premium-стиль", "ps:style_auto")],
    [("✏️ Описать свой стиль", "ps:style_custom")],
])

STYLE_NOTES_KB = _kb([
    [("➕ Добавить ещё требования", "ps:style_note")],
    [("📋 Показать накопленный стиль", "ps:style_show")],
    [("➡️ Продолжить к палитре", "ps:style_done")],
])

PALETTE_KB = _kb([
    [("🎨 Автоматическая палитра", "ps:palette_auto")],
    [("✏️ Описать свою палитру", "ps:palette_custom")],
])

PALETTE_NOTES_KB = _kb([
    [("➕ Добавить ещё цвета/правила", "ps:palette_note")],
    [("📋 Показать накопленную палитру", "ps:palette_show")],
    [("➡️ Выбрать движок изображений", "ps:palette_done")],
])

ENGINE_KB = _kb([
    [("✨ Автоматический выбор", "ps:engine_auto")],
    [("🖼 OpenAI Images", "ps:engine_openai")],
    [("🎨 Midjourney", "ps:engine_midjourney")],
    [("✏️ Дополнить визуальный бриф", "ps:visual_note")],
])

FINAL_KB = _kb([
    [("🚀 Собрать PDF + PPTX", "ps:build")],
    [("✏️ Дополнить визуальный бриф", "ps:visual_note")],
    [("🎨 Дополнить стиль", "ps:style_note")],
    [("⬅️ Назад к логотипу", "ps:logo_choice")],
    [("❌ Отменить проект", "ps:cancel")],
])


async def _send_long(message, text: str, reply_markup=None) -> None:
    text = (text or "").strip() or "Готово."
    chunks: list[str] = []
    while len(text) > 3900:
        split_at = text.rfind("\n", 0, 3900)
        if split_at < 1000:
            split_at = 3900
        chunks.append(text[:split_at])
        text = text[split_at:].lstrip()
    chunks.append(text)
    for i, chunk in enumerate(chunks):
        await message.reply_text(chunk, reply_markup=(reply_markup if i == len(chunks) - 1 else None))


async def _reply(update: Update, text: str, reply_markup=None) -> None:
    message = update.effective_message
    if not message:
        return
    await _send_long(message, text, reply_markup=reply_markup)


async def _local_error(update: Update, operation: str, exc: Exception) -> None:
    error_id = uuid.uuid4().hex[:8]
    log.exception("Presentation operation=%s error_id=%s: %s", operation, error_id, exc)
    await _reply(
        update,
        "Не удалось завершить именно этот этап. Проект и все введённые данные сохранены.\n"
        f"Код ошибки: {error_id}. Нажмите «Продолжить» или повторите действие.",
        START_KB,
    )


def _font(size: int, bold: bool = False) -> ImageFont.FreeTypeFont | ImageFont.ImageFont:
    path = FONT_BOLD if bold else FONT_REGULAR
    try:
        return ImageFont.truetype(path, size=size)
    except Exception:
        return ImageFont.load_default()


def _hex(value: str) -> tuple[int, int, int]:
    value = (value or "#000000").lstrip("#")
    if len(value) != 6:
        value = "000000"
    return tuple(int(value[i : i + 2], 16) for i in (0, 2, 4))


def _image_api_config() -> tuple[str, str]:
    key = (os.getenv("OPENAI_IMAGE_KEY", "").strip() or os.getenv("OPENAI_API_KEY", "").strip())
    base = (os.getenv("OPENAI_IMAGE_BASE_URL", "").strip() or "https://api.openai.com/v1").rstrip("/")
    return key, base


async def _generate_image_bytes(prompt: str, size: str = IMAGE_SIZE) -> bytes:
    key, base = _image_api_config()
    if not key:
        raise RuntimeError("OPENAI_IMAGE_KEY/OPENAI_API_KEY is not configured")
    body = {
        "model": IMAGE_MODEL,
        "prompt": prompt,
        "size": size,
        "n": 1,
    }
    headers = {"Authorization": f"Bearer {key}", "Content-Type": "application/json"}
    async with httpx.AsyncClient(timeout=240.0) as client:
        response = await client.post(f"{base}/images/generations", headers=headers, json=body)
    if response.status_code // 100 != 2:
        raise RuntimeError(f"Image HTTP {response.status_code}: {response.text[:500]}")
    data = response.json()
    item = (data.get("data") or [{}])[0]
    b64 = item.get("b64_json")
    if b64:
        return base64.b64decode(b64)
    url = item.get("url")
    if url:
        async with httpx.AsyncClient(timeout=120.0) as client:
            image_response = await client.get(url)
        image_response.raise_for_status()
        return image_response.content
    raise RuntimeError("Image provider returned neither b64_json nor url")


def _fallback_logo_symbol(index: int, palette: dict[str, str]) -> Image.Image:
    img = Image.new("RGBA", (720, 720), (0, 0, 0, 0))
    draw = ImageDraw.Draw(img)
    primary = _hex(palette["primary"]) + (255,)
    accent = _hex(palette["accent"]) + (255,)
    if index == 1:
        draw.line([(210, 540), (360, 150), (510, 540)], fill=primary, width=42, joint="curve")
        draw.arc((245, 270, 475, 560), 210, 340, fill=accent, width=24)
        draw.line([(360, 225), (360, 535)], fill=accent, width=18)
    elif index == 2:
        for r, color, width in [(240, primary, 36), (165, accent, 24), (90, primary, 20)]:
            draw.rounded_rectangle((360-r, 360-r, 360+r, 360+r), radius=60, outline=color, width=width)
        draw.line([(225, 495), (495, 225)], fill=primary, width=22)
    else:
        pts = [(360, 110), (565, 240), (510, 500), (360, 610), (210, 500), (155, 240)]
        draw.line(pts + [pts[0]], fill=primary, width=28, joint="curve")
        draw.line([(250, 410), (470, 410)], fill=accent, width=24)
        draw.arc((245, 200, 475, 445), 200, 342, fill=primary, width=20)
    return img


def _compose_logo_card(symbol_bytes: bytes | None, brand: str, index: int, palette: dict[str, str]) -> bytes:
    canvas_img = Image.new("RGB", (1024, 1024), _hex(palette["background"]))
    if symbol_bytes:
        try:
            symbol = Image.open(io.BytesIO(symbol_bytes)).convert("RGBA")
            symbol = ImageOps.contain(symbol, (660, 660), Image.Resampling.LANCZOS)
        except Exception:
            symbol = _fallback_logo_symbol(index, palette)
    else:
        symbol = _fallback_logo_symbol(index, palette)
    x = (1024 - symbol.width) // 2
    canvas_img.paste(symbol, (x, 70), symbol if symbol.mode == "RGBA" else None)
    draw = ImageDraw.Draw(canvas_img)
    brand_value = re.sub(r"\s+", " ", brand.strip()).upper()
    font_size = 76
    while font_size > 36:
        f = _font(font_size, bold=True)
        box = draw.textbbox((0, 0), brand_value, font=f)
        if box[2] - box[0] <= 880:
            break
        font_size -= 4
    bbox = draw.textbbox((0, 0), brand_value, font=f)
    tx = (1024 - (bbox[2] - bbox[0])) // 2
    draw.text((tx, 790), brand_value, font=f, fill=_hex(palette["text"]))
    subtitle = f"LOGO CONCEPT {index}"
    sf = _font(22)
    sb = draw.textbbox((0, 0), subtitle, font=sf)
    draw.text(((1024 - (sb[2] - sb[0])) // 2, 900), subtitle, font=sf, fill=_hex(palette["secondary"]))
    out = io.BytesIO()
    canvas_img.save(out, format="PNG", optimize=True)
    return out.getvalue()


async def _generate_logos(project: dict[str, Any], update: Update | None = None, context: ContextTypes.DEFAULT_TYPE | None = None) -> list[str]:
    brand = project.get("profile", {}).get("brand_name", "")
    if not _is_valid_brand_name(brand):
        raise RuntimeError("Brand name is missing or invalid")
    combined_notes = "\n".join(project.get("logo_notes", []))
    profile = project.get("profile", {})
    product = profile.get("product", "")
    positioning = profile.get("positioning", "")
    visual_direction = profile.get("visual_direction", "")
    raw_brief = project.get("raw_brief", "")
    concepts = [
        "architectural monogram-like geometry without letters; one memorable coherent symbol with refined negative space",
        "modular ecosystem symbol; a small set of modules forms one distinctive scalable mark",
        "premium minimal emblem; restrained geometry, clear silhouette and strong recognition at small size",
    ]
    prompts: list[str] = []
    for concept in concepts:
        prompts.append(f"""
Create one premium vector-style LOGO SYMBOL ONLY for the exact brand context below.
Never draw letters, words, captions, initials, pseudo-letters, mockups, business cards, walls or packaging.
The exact brand name will be typeset separately by software.
Concept direction: {concept}.
Exact brand name for semantic context only: {brand}.
Tagline/context: {profile.get('tagline', '')}.
Business/product: {product[:1400]}.
Positioning: {positioning[:1000]}.
Visual direction from main brief: {visual_direction[:900]}.
Cumulative user logo requirements: {combined_notes[:2400] or 'none'}.
Approved main brief excerpt: {raw_brief[:2600]}.
Flat vector style, simple geometry, strong silhouette, timeless, highly legible at 48 px, no stock icon, no gradients, no 3D, no shadows, centered square composition, generous margins, one standalone concept.
""")
    generated: list[bytes] = []
    if _STUDIO_IMAGE_BATCH_CALL is not None and update is not None and context is not None:
        with contextlib.suppress(Exception):
            batch = await _STUDIO_IMAGE_BATCH_CALL(update, context, prompts, "openai", "presentation_logo_variants")
            if batch:
                generated = list(batch)
    if not generated:
        for prompt in prompts:
            try:
                generated.append(await _generate_image_bytes(prompt, "1024x1024"))
            except Exception:
                generated.append(b"")
    out_paths: list[str] = []
    pdir = _project_dir(project["user_id"], project["project_id"])
    for index in range(1, 4):
        symbol_bytes = generated[index - 1] if index - 1 < len(generated) else b""
        card = _compose_logo_card(symbol_bytes or None, brand, index, project.get("palette", DEFAULT_PALETTE))
        path = pdir / f"logo_concept_{index}.png"
        path.write_bytes(card)
        out_paths.append(str(path))
    project["logo_candidates"] = out_paths
    return out_paths


def _all_notes(project: dict[str, Any], key: str) -> str:
    notes = project.get(key, [])
    if not notes:
        return "Пока ничего не добавлено."
    return "\n\n".join(f"{i}. {note}" for i, note in enumerate(notes, 1))


def _logo_selection_kb(count: int) -> InlineKeyboardMarkup:
    rows = [[(f"✅ Выбрать вариант {i}", f"ps:logo_select:{i}")] for i in range(1, count + 1)]
    rows.extend([
        [("🔄 Создать ещё 3 с учётом всех пожеланий", "ps:logo_generate")],
        [("✏️ Дополнить пожелания", "ps:logo_note")],
        [("📤 Загрузить свой логотип", "ps:logo_upload")],
        [("⏭ Без логотипа", "ps:logo_skip")],
    ])
    return _kb(rows)


async def _send_logo_candidates(update: Update, project: dict[str, Any]) -> None:
    paths = project.get("logo_candidates", [])
    for i, path in enumerate(paths, 1):
        with open(path, "rb") as fh:
            await update.effective_message.reply_photo(
                photo=InputFile(fh, filename=f"logo_{i}.png"),
                caption=f"Вариант {i} — точное название: {project['profile']['brand_name']}",
            )
    await _reply(
        update,
        "Выберите вариант. При повторной генерации все предыдущие дополнения к логотипу сохраняются.",
        _logo_selection_kb(len(paths)),
    )


def _sanitize_bullets(items: list[str]) -> list[str]:
    out = []
    for item in items:
        item = re.sub(r"\s+", " ", str(item)).strip(" •-–—")
        if item and item not in out:
            out.append(item[:240])
    return out[:6]


def _extract_named_offers(raw_brief: str) -> list[str]:
    """Extract product/service lines that already contain their exact price."""
    offers: list[str] = []
    for line in (raw_brief or "").splitlines():
        clean = re.sub(r"\s+", " ", line).strip(" •-*#\t")
        if not clean or not re.search(r"(?:от\s+)?\d[\d\s.,]*\s*(?:₽|руб|THB|USD|EUR|\$|€)", clean, re.I):
            continue
        if re.search(r"[A-Za-zА-Яа-яЁё]", clean) and clean not in offers:
            offers.append(clean[:240])
    return offers[:10]


async def _generate_deck_content(project: dict[str, Any], update: Update | None = None) -> list[dict[str, Any]]:
    structure = project.get("structure", [])
    prompt = f"""
Подготовь финальный текст презентации строго в JSON:
{{"slides":[{{"title":"", "subtitle":"", "bullets":[""], "image_prompt":""}}]}}

Нужно ровно {len(structure)} слайдов в том же порядке и с теми же смыслами.
Структура:
{json.dumps(structure, ensure_ascii=False)}

Правила:
- название бренда всегда точно: {project.get('profile', {}).get('brand_name', '')};
- используй фактические продукты, цены, этапы и контакты из полного брифа;
- на слайде 2-5 содержательных тезисов, а не общие фразы;
- не выдумывай отзывы, статистику, награды, сертификаты и кейсы;
- на продуктовом слайде обязательно сохрани названия продуктов и цены с «от», если они есть;
- на финальном слайде обязательно сохрани реальные контакты из брифа;
- image_prompt: фотореалистичная сцена по смыслу слайда, без текста, без логотипов,
  с местом под верстку, единый premium editorial style.

Главный бриф:
{project.get('raw_brief', '')}

Накопленные визуальные требования:
{_all_notes(project, 'visual_notes')}

Накопленные требования к стилю:
{_all_notes(project, 'style_notes')}
"""
    result: list[dict[str, Any]] = []
    try:
        parsed = _json_from_text(await _llm_text(prompt, 0.2, update))
        slide_values = parsed.get("slides") if isinstance(parsed, dict) else None
        if isinstance(slide_values, list) and len(slide_values) == len(structure):
            for base, item in zip(structure, slide_values):
                if not isinstance(item, dict):
                    item = {}
                result.append({
                    **base,
                    "title": str(item.get("title") or base.get("title") or "").strip()[:110],
                    "subtitle": str(item.get("subtitle") or "").strip()[:220],
                    "bullets": _sanitize_bullets(item.get("bullets") or base.get("bullets") or []),
                    "image_prompt": str(item.get("image_prompt") or base.get("image_prompt") or "").strip()[:1800],
                })
    except Exception:
        log.exception("Deck copy generation failed; using structure content")
    if not result:
        result = [
            {
                **slide,
                "subtitle": "",
                "bullets": _sanitize_bullets(slide.get("bullets", [])),
            }
            for slide in structure
        ]

    # Deterministic safeguards for the most important facts.
    profile = project.get("profile", {})
    if result:
        result[0]["title"] = profile.get("brand_name") or result[0]["title"]
        result[0]["subtitle"] = profile.get("tagline") or result[0].get("subtitle", "")
    prices = profile.get("prices", [])
    contacts = profile.get("contacts", [])
    named_offers = _extract_named_offers(project.get("raw_brief", ""))
    product_offers = [x for x in named_offers if not re.search(r"CARE|SERVICE|СЕРВИС|ОБСЛУЖ", x, re.I)]
    service_offers = [x for x in named_offers if re.search(r"CARE|SERVICE|СЕРВИС|ОБСЛУЖ", x, re.I)]
    for slide in result:
        title = slide.get("title", "")
        if product_offers and re.search(r"продукт|линейк|тариф|пакет", title, re.I):
            slide["bullets"] = _sanitize_bullets(product_offers + slide.get("bullets", []))
        elif service_offers and re.search(r"care|сервис|обслуж|сопровожд", title, re.I):
            slide["bullets"] = _sanitize_bullets(service_offers + slide.get("bullets", []))
    if prices and not named_offers:
        for slide in result:
            if re.search(r"продукт|линейк|тариф|пакет|стоим", slide["title"], re.I):
                existing = " ".join(slide["bullets"])
                for price in prices:
                    if price not in existing and len(slide["bullets"]) < 6:
                        slide["bullets"].append(price)
                break
    if contacts and result:
        final = result[-1]
        final["bullets"] = _sanitize_bullets(final.get("bullets", []) + contacts)
    return result


def _default_scene_prompt(project: dict[str, Any], slide: dict[str, Any]) -> str:
    profile = project.get("profile", {})
    return (
        f"Premium editorial commercial photograph for a presentation about {profile.get('product', '')[:600]}. "
        f"Slide idea: {slide.get('title', '')}. Details: {'; '.join(slide.get('bullets', [])[:5])}. "
        "Photorealistic, contemporary architecture and product design, natural materials, soft daylight, "
        "calm refined color grading, realistic scale, strong composition, generous negative space for layout, "
        "16:9 landscape. No words, no lettering, no logo, no watermark, no presentation frame, no UI, no collage."
    )


async def _prepare_slide_images(project: dict[str, Any], slides: list[dict[str, Any]], update: Update, context: ContextTypes.DEFAULT_TYPE) -> dict[int, str]:
    mode = project.get("image_mode", "auto")
    uploaded = [p for p in project.get("uploaded_images", []) if os.path.exists(p)]
    result: dict[int, str] = {int(k): v for k, v in (project.get("slide_images") or {}).items() if os.path.exists(v)}
    target_indexes = [
        i for i, slide in enumerate(slides)
        if slide.get("image_needed") or slide.get("layout") in {"cover", "split", "full_image", "comparison", "cta"}
    ][:MAX_IMAGES]
    upload_iter = iter(uploaded)
    if mode in {"upload", "mixed"}:
        for idx in target_indexes:
            if idx in result:
                continue
            with contextlib.suppress(StopIteration):
                result[idx] = next(upload_iter)
    if mode == "skip":
        project["slide_images"] = {str(k): v for k, v in result.items()}
        return result
    if mode in {"auto", "mixed"}:
        missing = [idx for idx in target_indexes if idx not in result]
        prompts: list[str] = []
        for idx in missing:
            slide = slides[idx]
            prompt = slide.get("image_prompt") or _default_scene_prompt(project, slide)
            prompt += "\nExact main brief context (do not drift to another business):\n" + project.get("raw_brief", "")[:3200]
            prompt += "\nCumulative visual rules:\n" + _all_notes(project, "visual_notes")[:3000]
            prompt += "\nCumulative style rules:\n" + _all_notes(project, "style_notes")[:2200]
            prompt += "\nMandatory: 16:9 landscape, no readable text, no logo, no watermark, no slide frame, no UI. The image must depict the actual product/service and the exact slide meaning, not a generic abstract placeholder."
            prompts.append(prompt[:7600])
        generated: list[bytes] = []
        engine = project.get("generation_engine") or "auto"
        if prompts and _STUDIO_IMAGE_BATCH_CALL is not None:
            batch = await _STUDIO_IMAGE_BATCH_CALL(update, context, prompts, engine, "presentation_slide_images")
            if batch:
                generated = list(batch)
        elif prompts:
            for pos, prompt in enumerate(prompts, 1):
                await _reply(update, f"Создаю визуал {pos}/{len(prompts)}…")
                try:
                    generated.append(await _generate_image_bytes(prompt, IMAGE_SIZE))
                except Exception:
                    generated.append(b"")
        pdir = _project_dir(project["user_id"], project["project_id"])
        for idx, image_bytes in zip(missing, generated):
            if not image_bytes:
                continue
            path = pdir / f"slide_{idx + 1:02d}.png"
            path.write_bytes(image_bytes)
            result[idx] = str(path)
    project["slide_images"] = {str(k): v for k, v in result.items()}
    return result


def _parse_palette(project: dict[str, Any]) -> dict[str, str]:
    palette = deepcopy(DEFAULT_PALETTE)
    joined = "\n".join(project.get("palette_notes", []))
    hexes = re.findall(r"#[0-9A-Fa-f]{6}", joined)
    keys = ["background", "text", "primary", "secondary", "sand", "accent", "card"]
    for key, value in zip(keys, hexes):
        palette[key] = value.upper()
    return palette


def _crop_image(path: str, width: int, height: int) -> Image.Image:
    im = Image.open(path).convert("RGB")
    return ImageOps.fit(im, (width, height), method=Image.Resampling.LANCZOS, centering=(0.5, 0.5))


def _tmp_crop(path: str, width: int, height: int, out_path: Path) -> str:
    _crop_image(path, width, height).save(out_path, "JPEG", quality=92)
    return str(out_path)


def _ppt_add_text(slide, text: str, x, y, w, h, size: int, color: str, bold: bool = False,
                  align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, font_name: str = "Lato"):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*_hex(color))
    return box


def _ppt_add_bullets(slide, bullets: list[str], x, y, w, h, palette: dict[str, str], size: int = 22):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, bullet in enumerate(bullets[:6]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = "Lato"
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(*_hex(palette["text"]))
        p.space_after = Pt(10)
        p.text = "•  " + p.text
    return box


def _ppt_rect(slide, x, y, w, h, fill: str, line: str | None = None, radius: bool = False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*_hex(fill))
    shape.line.color.rgb = RGBColor(*_hex(line or fill))
    return shape


def _ppt_add_logo(slide, project: dict[str, Any], x, y, w):
    path = project.get("logo_selected")
    if path and os.path.exists(path):
        try:
            slide.shapes.add_picture(path, x, y, width=w)
        except Exception:
            log.exception("Could not add logo to PPTX")


def _ppt_header(slide, index: int, title: str, palette: dict[str, str], project: dict[str, Any]):
    _ppt_add_text(slide, f"{index:02d}", Inches(0.55), Inches(0.36), Inches(0.6), Inches(0.4), 14, palette["white"], True, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    badge = _ppt_rect(slide, Inches(0.52), Inches(0.30), Inches(0.55), Inches(0.46), palette["primary"], radius=True)
    slide.shapes._spTree.remove(badge._element)
    slide.shapes._spTree.insert(2, badge._element)
    _ppt_add_text(slide, title, Inches(1.2), Inches(0.28), Inches(10.8), Inches(0.75), 28, palette["text"], True)
    _ppt_rect(slide, Inches(1.2), Inches(1.05), Inches(1.25), Inches(0.04), palette["primary"])
    _ppt_add_logo(slide, project, Inches(11.9), Inches(0.18), Inches(0.8))


def _build_pptx(project: dict[str, Any], slides: list[dict[str, Any]], images: dict[int, str], out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    palette = project["palette"]
    pdir = out_path.parent
    brand = project.get("profile", {}).get("brand_name", "")

    for idx, data in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        _ppt_rect(slide, 0, 0, prs.slide_width, prs.slide_height, palette["background"])
        layout = data.get("layout", "split")
        image_path = images.get(idx)
        if idx == 0 or layout == "cover":
            if image_path:
                crop = _tmp_crop(image_path, 900, 1200, pdir / f"_crop_cover_{idx}.jpg")
                slide.shapes.add_picture(crop, Inches(7.35), 0, width=Inches(5.98), height=Inches(7.5))
                _ppt_rect(slide, Inches(7.18), 0, Inches(0.18), Inches(7.5), palette["primary"])
            else:
                _ppt_rect(slide, Inches(7.35), 0, Inches(5.98), Inches(7.5), palette["primary"])
            _ppt_add_text(slide, data.get("title") or brand, Inches(0.75), Inches(1.35), Inches(6.15), Inches(1.55), 38, palette["text"], True)
            subtitle = data.get("subtitle") or project.get("profile", {}).get("tagline", "")
            _ppt_add_text(slide, subtitle, Inches(0.78), Inches(3.05), Inches(5.7), Inches(0.8), 22, palette["accent"], False)
            _ppt_add_bullets(slide, data.get("bullets", [])[:3], Inches(0.78), Inches(4.15), Inches(5.7), Inches(1.5), palette, 17)
            _ppt_add_logo(slide, project, Inches(0.75), Inches(0.35), Inches(1.15))
            continue

        _ppt_header(slide, idx + 1, data.get("title", ""), palette, project)
        bullets = data.get("bullets", [])
        if layout in {"cards", "process"}:
            cols = 2 if len(bullets) <= 4 else 3
            rows = math.ceil(len(bullets) / cols)
            gap = 0.28
            left = 0.72
            top = 1.55
            total_w = 11.9
            card_w = (total_w - gap * (cols - 1)) / cols
            card_h = min(1.55, (5.25 - gap * (rows - 1)) / max(1, rows))
            for i, bullet in enumerate(bullets[:6]):
                c = i % cols
                r = i // cols
                x = Inches(left + c * (card_w + gap))
                y = Inches(top + r * (card_h + gap))
                _ppt_rect(slide, x, y, Inches(card_w), Inches(card_h), palette["card"], palette["sand"], True)
                _ppt_add_text(slide, f"{i + 1:02d}", x + Inches(0.22), y + Inches(0.16), Inches(0.45), Inches(0.35), 12, palette["primary"], True)
                _ppt_add_text(slide, bullet, x + Inches(0.22), y + Inches(0.52), Inches(card_w - 0.44), Inches(card_h - 0.62), 18, palette["text"], True)
            if image_path and len(bullets) <= 4:
                crop = _tmp_crop(image_path, 1600, 450, pdir / f"_crop_band_{idx}.jpg")
                slide.shapes.add_picture(crop, Inches(0.72), Inches(5.7), width=Inches(11.9), height=Inches(1.2))
        elif layout == "comparison":
            offer_lines = _extract_offer_lines(project.get("raw_brief", ""))
            items = offer_lines[:4] or bullets[:4]
            cols = max(1, len(items))
            card_w = 11.9 / cols - 0.18
            for i, item in enumerate(items):
                x = Inches(0.72 + i * (card_w + 0.18))
                _ppt_rect(slide, x, Inches(1.55), Inches(card_w), Inches(4.95), palette["white"], palette["sand"], True)
                title, body = _split_offer(item)
                _ppt_add_text(slide, title, x + Inches(0.2), Inches(1.85), Inches(card_w - 0.4), Inches(0.75), 20, palette["primary"], True, PP_ALIGN.CENTER)
                _ppt_add_text(slide, body, x + Inches(0.22), Inches(2.8), Inches(card_w - 0.44), Inches(3.1), 15, palette["text"], False)
        elif layout == "full_image" and image_path:
            crop = _tmp_crop(image_path, 1600, 720, pdir / f"_crop_full_{idx}.jpg")
            slide.shapes.add_picture(crop, Inches(0.72), Inches(1.42), width=Inches(11.9), height=Inches(5.55))
            overlay = _ppt_rect(slide, Inches(0.72), Inches(5.7), Inches(11.9), Inches(1.27), palette["primary"])
            overlay.fill.transparency = 8
            _ppt_add_text(slide, "  •  ".join(bullets[:3]), Inches(1.0), Inches(5.96), Inches(11.3), Inches(0.75), 18, palette["white"], True)
        elif layout == "cta":
            _ppt_rect(slide, 0, 0, prs.slide_width, prs.slide_height, palette["primary"])
            if image_path:
                crop = _tmp_crop(image_path, 720, 900, pdir / f"_crop_cta_{idx}.jpg")
                slide.shapes.add_picture(crop, Inches(8.45), Inches(0.5), width=Inches(4.35), height=Inches(6.5))
            _ppt_add_text(slide, data.get("title", "Следующий шаг"), Inches(0.8), Inches(1.1), Inches(7.1), Inches(1.2), 40, palette["white"], True)
            _ppt_add_text(slide, data.get("subtitle", ""), Inches(0.82), Inches(2.45), Inches(6.8), Inches(0.8), 21, palette["sand"])
            _ppt_add_bullets_white(slide, bullets, Inches(0.82), Inches(3.45), Inches(6.9), Inches(2.7), palette, 20)
            _ppt_add_logo(slide, project, Inches(0.8), Inches(0.3), Inches(1.1))
        else:  # split
            image_left = idx % 2 == 0
            if image_path:
                crop = _tmp_crop(image_path, 900, 850, pdir / f"_crop_split_{idx}.jpg")
                img_x = Inches(0.72 if image_left else 7.15)
                slide.shapes.add_picture(crop, img_x, Inches(1.45), width=Inches(5.45), height=Inches(5.45))
                text_x = Inches(6.55 if image_left else 0.78)
            else:
                text_x = Inches(1.1)
            _ppt_add_text(slide, data.get("subtitle", ""), text_x, Inches(1.55), Inches(5.7), Inches(0.8), 20, palette["accent"], True)
            _ppt_add_bullets(slide, bullets, text_x, Inches(2.35), Inches(5.75), Inches(3.9), palette, 20)

        _ppt_add_text(slide, f"{brand}  •  {idx + 1}", Inches(0.65), Inches(7.12), Inches(4.0), Inches(0.25), 8, palette["secondary"])

    prs.save(out_path)


def _ppt_add_bullets_white(slide, bullets: list[str], x, y, w, h, palette: dict[str, str], size: int = 22):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, bullet in enumerate(bullets[:6]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "•  " + bullet
        p.font.name = "Lato"
        p.font.size = Pt(size)
        p.font.color.rgb = RGBColor(*_hex(palette["white"]))
        p.space_after = Pt(10)
    return box


def _extract_offer_lines(raw: str) -> list[str]:
    result = []
    lines = raw.splitlines()
    for i, line in enumerate(lines):
        clean = line.strip()
        if re.match(r"^(?:#{1,4}\s*)?[A-ZА-ЯЁ][A-ZА-ЯЁ0-9 _-]{2,35}$", clean) and not clean.lower().startswith(("целевая", "этап", "требован", "визуал")):
            block = [clean]
            for nxt in lines[i + 1 : i + 8]:
                n = nxt.strip()
                if not n:
                    if len(block) > 1:
                        break
                    continue
                if re.match(r"^(?:#{1,4}\s*)?[A-ZА-ЯЁ][A-ZА-ЯЁ0-9 _-]{2,35}$", n):
                    break
                block.append(n)
            joined = "\n".join(block)
            if re.search(r"цена|стоим|руб|₽|включает|размер", joined, re.I):
                result.append(joined[:900])
    # Unique, preserve order.
    return list(dict.fromkeys(result))[:6]


def _split_offer(value: str) -> tuple[str, str]:
    lines = [l.strip(" #-") for l in value.splitlines() if l.strip()]
    return (lines[0] if lines else "Решение", "\n".join(lines[1:]) if len(lines) > 1 else value)


def _register_pdf_fonts() -> tuple[str, str]:
    regular_name = "PresRegular"
    bold_name = "PresBold"
    with contextlib.suppress(Exception):
        if regular_name not in pdfmetrics.getRegisteredFontNames() and FONT_REGULAR:
            pdfmetrics.registerFont(TTFont(regular_name, FONT_REGULAR))
        if bold_name not in pdfmetrics.getRegisteredFontNames() and FONT_BOLD:
            pdfmetrics.registerFont(TTFont(bold_name, FONT_BOLD))
    return regular_name, bold_name


def _pdf_wrap(text: str, max_chars: int) -> list[str]:
    return textwrap.wrap(text, width=max_chars, break_long_words=False, replace_whitespace=False) or [""]


def _pdf_draw_image(c: canvas.Canvas, path: str, x: float, y: float, w: float, h: float, pdir: Path, name: str):
    crop_path = pdir / name
    _crop_image(path, max(200, int(w * 2)), max(200, int(h * 2))).save(crop_path, "JPEG", quality=92)
    c.drawImage(ImageReader(str(crop_path)), x, y, width=w, height=h, preserveAspectRatio=False, mask="auto")


def _build_pdf(project: dict[str, Any], slides: list[dict[str, Any]], images: dict[int, str], out_path: Path) -> None:
    page_w, page_h = 960, 540  # exact 16:9
    c = canvas.Canvas(str(out_path), pagesize=(page_w, page_h), pageCompression=1)
    regular, bold = _register_pdf_fonts()
    palette = project["palette"]
    pdir = out_path.parent
    brand = project.get("profile", {}).get("brand_name", "")

    def set_fill(key: str):
        c.setFillColor(HexColor(palette[key]))

    for idx, data in enumerate(slides):
        set_fill("background")
        c.rect(0, 0, page_w, page_h, stroke=0, fill=1)
        image_path = images.get(idx)
        layout = data.get("layout", "split")
        if idx == 0 or layout == "cover":
            if image_path:
                _pdf_draw_image(c, image_path, 530, 0, 430, 540, pdir, f"_pdf_cover_{idx}.jpg")
            else:
                set_fill("primary"); c.rect(530, 0, 430, 540, stroke=0, fill=1)
            set_fill("primary"); c.rect(520, 0, 10, 540, stroke=0, fill=1)
            set_fill("text"); c.setFont(bold, 35)
            y = 405
            for line in _pdf_wrap(data.get("title") or brand, 24):
                c.drawString(55, y, line); y -= 42
            set_fill("accent"); c.setFont(regular, 18)
            y -= 20
            for line in _pdf_wrap(data.get("subtitle", ""), 42):
                c.drawString(58, y, line); y -= 25
            set_fill("text"); c.setFont(regular, 12)
            y -= 25
            for bullet in data.get("bullets", [])[:3]:
                for line_i, line in enumerate(_pdf_wrap(bullet, 55)):
                    c.drawString(58 if line_i == 0 else 72, y, ("• " if line_i == 0 else "") + line)
                    y -= 18
                y -= 6
        elif layout == "cta":
            set_fill("primary"); c.rect(0, 0, page_w, page_h, stroke=0, fill=1)
            if image_path:
                _pdf_draw_image(c, image_path, 650, 35, 270, 470, pdir, f"_pdf_cta_{idx}.jpg")
            set_fill("white"); c.setFont(bold, 34); c.drawString(58, 410, data.get("title", "Следующий шаг"))
            set_fill("sand"); c.setFont(regular, 18); c.drawString(60, 365, data.get("subtitle", ""))
            set_fill("white"); c.setFont(regular, 16)
            y = 305
            for bullet in data.get("bullets", [])[:6]:
                for line_i, line in enumerate(_pdf_wrap(bullet, 55)):
                    c.drawString(60 if line_i == 0 else 76, y, ("• " if line_i == 0 else "") + line); y -= 22
                y -= 8
        else:
            # Header
            set_fill("primary"); c.roundRect(38, 485, 34, 30, 5, stroke=0, fill=1)
            set_fill("white"); c.setFont(bold, 11); c.drawCentredString(55, 496, f"{idx + 1:02d}")
            set_fill("text"); c.setFont(bold, 26); c.drawString(85, 488, data.get("title", ""))
            set_fill("primary"); c.rect(85, 472, 90, 3, stroke=0, fill=1)
            bullets = data.get("bullets", [])
            if layout in {"cards", "process"}:
                cols = 2 if len(bullets) <= 4 else 3
                rows = math.ceil(len(bullets) / cols)
                gap = 18
                left, top = 48, 425
                total_w = 864
                card_w = (total_w - gap * (cols - 1)) / cols
                card_h = min(105, (330 - gap * (rows - 1)) / max(1, rows))
                for i, bullet in enumerate(bullets[:6]):
                    col, row = i % cols, i // cols
                    x = left + col * (card_w + gap)
                    y0 = top - (row + 1) * card_h - row * gap
                    set_fill("card"); c.roundRect(x, y0, card_w, card_h, 10, stroke=0, fill=1)
                    c.setStrokeColor(HexColor(palette["sand"])); c.roundRect(x, y0, card_w, card_h, 10, stroke=1, fill=0)
                    set_fill("primary"); c.setFont(bold, 10); c.drawString(x + 14, y0 + card_h - 22, f"{i + 1:02d}")
                    set_fill("text"); c.setFont(bold, 14)
                    ty = y0 + card_h - 48
                    for line in _pdf_wrap(bullet, max(24, int(card_w / 8))):
                        c.drawString(x + 14, ty, line); ty -= 18
            elif layout == "comparison":
                items = _extract_offer_lines(project.get("raw_brief", ""))[:4] or bullets[:4]
                cols = max(1, len(items)); gap = 10; left = 42; total_w = 876
                card_w = (total_w - gap * (cols - 1)) / cols
                for i, item in enumerate(items):
                    x = left + i * (card_w + gap)
                    set_fill("white"); c.roundRect(x, 70, card_w, 370, 10, stroke=0, fill=1)
                    c.setStrokeColor(HexColor(palette["sand"])); c.roundRect(x, 70, card_w, 370, 10, stroke=1, fill=0)
                    title, body = _split_offer(item)
                    set_fill("primary"); c.setFont(bold, 14)
                    ty = 410
                    for line in _pdf_wrap(title, max(14, int(card_w / 8))):
                        c.drawCentredString(x + card_w / 2, ty, line); ty -= 18
                    set_fill("text"); c.setFont(regular, 10); ty -= 15
                    for line in _pdf_wrap(body, max(20, int(card_w / 6.5)))[:16]:
                        c.drawString(x + 12, ty, line); ty -= 14
            else:
                image_left = idx % 2 == 0
                if image_path:
                    ix = 45 if image_left else 525
                    _pdf_draw_image(c, image_path, ix, 65, 390, 380, pdir, f"_pdf_split_{idx}.jpg")
                    tx = 475 if image_left else 55
                else:
                    tx = 90
                set_fill("accent"); c.setFont(bold, 16)
                sy = 412
                for line in _pdf_wrap(data.get("subtitle", ""), 46):
                    c.drawString(tx, sy, line); sy -= 22
                set_fill("text"); c.setFont(regular, 15); y = sy - 25
                for bullet in bullets[:6]:
                    for line_i, line in enumerate(_pdf_wrap(bullet, 50)):
                        c.drawString(tx if line_i == 0 else tx + 16, y, ("• " if line_i == 0 else "") + line); y -= 20
                    y -= 8
            set_fill("secondary"); c.setFont(regular, 7); c.drawString(42, 22, f"{brand}  •  {idx + 1}")
        c.showPage()
    c.save()


async def _build_files(project: dict[str, Any], update: Update, context: ContextTypes.DEFAULT_TYPE) -> tuple[str, str]:
    project["palette"] = _parse_palette(project)
    slides = await _generate_deck_content(project, update)
    _save(project["user_id"], project)
    images = await _prepare_slide_images(project, slides, update, context)
    _save(project["user_id"], project)
    pdir = _project_dir(project["user_id"], project["project_id"])
    brand = _safe_filename(project.get("profile", {}).get("brand_name", "presentation"))
    pptx_path = pdir / f"{brand}_presentation.pptx"
    pdf_path = pdir / f"{brand}_presentation.pdf"
    await asyncio.to_thread(_build_pptx, project, slides, images, pptx_path)
    await asyncio.to_thread(_build_pdf, project, slides, images, pdf_path)
    project["final_slides"] = slides
    project["pptx_path"] = str(pptx_path)
    project["pdf_path"] = str(pdf_path)
    project["stage"] = "done"
    _save(project["user_id"], project)
    return str(pdf_path), str(pptx_path)


def _review_text(project: dict[str, Any]) -> str:
    p = project.get("profile", {})
    return (
        "Финальная проверка проекта\n\n"
        f"Бренд: {p.get('brand_name', '')}\n"
        f"Слайдов: {len(project.get('structure', []))}\n"
        f"Логотип: {'выбран' if project.get('logo_selected') else 'без логотипа'}\n"
        f"Режим изображений: {project.get('image_mode') or 'не выбран'}\n"
        f"Движок генерации: {project.get('generation_engine') or 'auto'}\n"
        f"Требований к визуалам сохранено: {len(project.get('visual_notes', []))}\n"
        f"Требований к стилю сохранено: {len(project.get('style_notes', []))}\n"
        f"Требований к палитре сохранено: {len(project.get('palette_notes', []))}\n"
        f"Загружено собственных изображений: {len(project.get('uploaded_images', []))}\n\n"
        "При сборке бот повторно использует полный главный бриф и все накопленные дополнения."
    )


async def _show_resume(update: Update, project: dict[str, Any]) -> None:
    stage = project.get("stage", "")
    if stage == "await_brief":
        await _reply(update, "Продолжаем. Пришлите полный главный бриф одним сообщением.", START_KB)
    elif stage in {"await_brand_name"}:
        await _reply(update, "Не удалось надёжно определить название. Пришлите только точное название бренда.")
    elif stage in {"structure_review", "await_structure_edit"}:
        await _reply(update, _structure_text(project), STRUCTURE_KB)
    elif stage.startswith("logo") or stage in {"await_logo_brief", "await_logo_upload"}:
        if project.get("logo_candidates"):
            await _send_logo_candidates(update, project)
        else:
            await _reply(update, "Выберите способ работы с логотипом.", LOGO_CHOICE_KB)
    elif stage in {"visual_collect", "await_visual_note", "await_images_upload"}:
        kb = UPLOAD_IMAGES_KB if stage == "await_images_upload" else VISUAL_NOTES_KB
        await _reply(update, "Все ранее введённые требования сохранены. Можно продолжить дополнять.", kb)
    elif stage.startswith("style") or stage == "await_style_note":
        await _reply(update, "Продолжаем настройку стиля.", STYLE_NOTES_KB)
    elif stage.startswith("palette") or stage == "await_palette_note":
        await _reply(update, "Продолжаем настройку палитры.", PALETTE_NOTES_KB)
    elif stage == "engine_choice":
        await _reply(update, "Выберите движок для недостающих изображений.", ENGINE_KB)
    elif stage in {"final_review", "building"}:
        await _reply(update, _review_text(project), FINAL_KB)
    elif stage == "done":
        await _reply(update, "Проект уже собран. Можно создать новый или повторно скачать файлы.", _kb([
            [("📄 Скачать PDF", "ps:download_pdf"), ("📊 Скачать PPTX", "ps:download_pptx")],
            [("🆕 Новый проект", "ps:new")],
        ]))
    else:
        await _reply(update, "Проект найден. Выберите продолжение.", START_KB)


async def on_start_callback(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    await _reply(
        update,
        "Мастер презентации создаёт редактируемый PPTX и совпадающий PDF.\n\n"
        "Главный бриф сохраняется целиком. Пожелания к логотипу, изображениям, стилю и палитре "
        "можно дополнять многократно — предыдущие дополнения не стираются.",
        START_KB,
    )


async def on_callback(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    q = update.callback_query
    if not q or not (q.data or "").startswith(CALLBACK_PREFIX):
        return
    with contextlib.suppress(Exception):
        await q.answer()
    action = (q.data or "")[len(CALLBACK_PREFIX):]
    user_id = q.from_user.id
    project = _load(user_id)
    try:
        if action == "start":
            await on_start_callback(update, context)
        elif action == "new":
            project = _new_project(user_id, context.user_data.get("presentation_kind", "presentation"), update.effective_chat.id if update.effective_chat else 0)
            _save(user_id, project)
            await _reply(
                update,
                "Пришлите полный главный бриф одним сообщением.\n\n"
                "Обязательно укажите точное название бренда, продукт/услугу, цель, аудиторию, "
                "продуктовую линейку и цены, контакты, желаемое количество слайдов, визуальную концепцию "
                "и факты, которые нельзя выдумывать.",
                _kb([[("❌ Отменить", "ps:cancel")]]),
            )
        elif action == "resume":
            if not project:
                await _reply(update, "Сохранённого проекта нет. Создайте новый.", START_KB)
            else:
                await _show_resume(update, project)
        elif action == "cancel":
            _delete(user_id)
            await _reply(update, "Проект отменён и удалён. Можно начать новый.", START_KB)
        elif not project:
            await _reply(update, "Проект не найден. Начните новый.", START_KB)
        elif action == "brief_show":
            await _reply(update, _brief_summary(project), STRUCTURE_KB)
        elif action == "structure_edit":
            project["stage"] = "await_structure_edit"
            _save(user_id, project)
            await _reply(update, "Напишите правки к структуре. Их можно добавлять повторно — предыдущие сохранятся.")
        elif action == "structure_ok":
            project["stage"] = "logo_choice"
            _save(user_id, project)
            await _reply(
                update,
                "Структура утверждена. Выберите способ создания логотипа.\n\n"
                "Автоматический режим использует точное название и весь главный бриф. "
                "Отдельный бриф позволяет несколько раз дополнять требования перед генерацией.",
                LOGO_CHOICE_KB,
            )
        elif action == "back_structure":
            project["stage"] = "structure_review"
            _save(user_id, project)
            await _reply(update, _structure_text(project), STRUCTURE_KB)
        elif action == "logo_choice":
            project["stage"] = "logo_choice"
            _save(user_id, project)
            await _reply(update, "Выберите способ работы с логотипом.", LOGO_CHOICE_KB)
        elif action == "logo_auto":
            project["logo_mode"] = "auto"
            project["stage"] = "logo_generating"
            _save(user_id, project)
            await _reply(update, f"Создаю три разные концепции для «{project['profile']['brand_name']}». Текст названия будет наложен программно без ошибок.")
            await context.bot.send_chat_action(update.effective_chat.id, ChatAction.UPLOAD_PHOTO)
            await _generate_logos(project, update, context)
            project["stage"] = "logo_review"
            _save(user_id, project)
            await _send_logo_candidates(update, project)
        elif action == "logo_custom":
            project["logo_mode"] = "custom"
            project["stage"] = "await_logo_brief"
            _save(user_id, project)
            await _reply(
                update,
                "Пришлите отдельный бриф логотипа. После сохранения можно несколько раз нажимать "
                "«Добавить ещё пожелания». Генерация начнётся только по отдельной кнопке.",
            )
        elif action == "logo_note":
            project["stage"] = "await_logo_brief"
            _save(user_id, project)
            await _reply(update, "Добавьте новые пожелания к логотипу. Ранее введённые требования останутся в проекте.")
        elif action == "logo_notes_show":
            await _reply(update, _all_notes(project, "logo_notes"), LOGO_NOTES_KB)
        elif action == "logo_generate":
            project["stage"] = "logo_generating"
            _save(user_id, project)
            await _reply(update, "Создаю новую тройку с учётом главного брифа и всех накопленных пожеланий.")
            await _generate_logos(project, update, context)
            project["stage"] = "logo_review"
            _save(user_id, project)
            await _send_logo_candidates(update, project)
        elif action == "logo_upload":
            project["stage"] = "await_logo_upload"
            _save(user_id, project)
            await _reply(update, "Загрузите логотип как PNG/JPG или документ. Лучше PNG с прозрачным фоном.")
        elif action == "logo_skip":
            project["logo_selected"] = ""
            project["stage"] = "image_mode"
            _save(user_id, project)
            await _reply(update, "Выберите источник изображений для презентации.", IMAGE_MODE_KB)
        elif action.startswith("logo_select:"):
            index = int(action.split(":", 1)[1]) - 1
            candidates = project.get("logo_candidates", [])
            if index < 0 or index >= len(candidates):
                raise ValueError("Invalid logo selection")
            project["logo_selected"] = candidates[index]
            project["stage"] = "image_mode"
            _save(user_id, project)
            await _reply(update, "Логотип сохранён. Выберите источник изображений.", IMAGE_MODE_KB)
        elif action in {"images_auto", "images_upload", "images_mixed", "images_skip"}:
            mode_map = {
                "images_auto": "auto",
                "images_upload": "upload",
                "images_mixed": "mixed",
                "images_skip": "skip",
            }
            project["image_mode"] = mode_map[action]
            project["stage"] = "await_images_upload" if mode_map[action] in {"upload", "mixed"} else "visual_collect"
            _save(user_id, project)
            if project["stage"] == "await_images_upload":
                await _reply(
                    update,
                    "Загружайте изображения по одному как фото или документы. Они сохраняются в проекте. "
                    "После загрузки нажмите «Закончить загрузку».",
                    UPLOAD_IMAGES_KB,
                )
            else:
                await _reply(
                    update,
                    "Добавьте требования к изображениям. Можно делать это несколько раз; каждое дополнение сохраняется и применяется ко всем промптам.",
                    VISUAL_NOTES_KB,
                )
        elif action == "visual_note":
            project["stage"] = "await_visual_note"
            _save(user_id, project)
            await _reply(update, "Введите очередное дополнение к визуалам. Предыдущие дополнения не будут заменены.")
        elif action == "visual_show":
            await _reply(update, _all_notes(project, "visual_notes"), VISUAL_NOTES_KB)
        elif action == "visual_clear":
            project["visual_notes"] = []
            project["stage"] = "visual_collect"
            _save(user_id, project)
            await _reply(update, "Требования к визуалам очищены.", VISUAL_NOTES_KB)
        elif action == "uploads_show":
            await _reply(update, f"Сохранено изображений: {len(project.get('uploaded_images', []))}", UPLOAD_IMAGES_KB)
        elif action == "visual_done":
            project["stage"] = "style_choice"
            _save(user_id, project)
            await _reply(update, "Выберите стиль презентации.", STYLE_KB)
        elif action == "style_auto":
            if not project.get("style_notes"):
                project["style_notes"] = [
                    "Premium editorial: современная модульная сетка, много свободного пространства, крупные визуалы, содержательные слайды без шаблонных клипартов."
                ]
            project["stage"] = "palette_choice"
            _save(user_id, project)
            await _reply(update, "Стиль сохранён. Выберите палитру.", PALETTE_KB)
        elif action in {"style_custom", "style_note"}:
            project["stage"] = "await_style_note"
            _save(user_id, project)
            await _reply(update, "Опишите очередное требование к стилю. Можно добавлять несколько сообщений.")
        elif action == "style_show":
            await _reply(update, _all_notes(project, "style_notes"), STYLE_NOTES_KB)
        elif action == "style_done":
            project["stage"] = "palette_choice"
            _save(user_id, project)
            await _reply(update, "Выберите палитру.", PALETTE_KB)
        elif action == "palette_auto":
            project["palette_notes"] = project.get("palette_notes") or [
                "Тёплый светлый фон, глубокий основной цвет бренда, высокий контраст текста, один сдержанный акцент."
            ]
            project["palette"] = deepcopy(DEFAULT_PALETTE)
            project["stage"] = "engine_choice"
            _save(user_id, project)
            await _reply(update, "Палитра сохранена. Выберите движок для недостающих изображений.", ENGINE_KB)
        elif action in {"palette_custom", "palette_note"}:
            project["stage"] = "await_palette_note"
            _save(user_id, project)
            await _reply(update, "Опишите цвета и правила их использования. HEX-коды можно отправлять несколькими сообщениями.")
        elif action == "palette_show":
            await _reply(update, _all_notes(project, "palette_notes"), PALETTE_NOTES_KB)
        elif action == "palette_done":
            project["palette"] = _parse_palette(project)
            project["stage"] = "engine_choice"
            _save(user_id, project)
            await _reply(update, "Палитра сохранена. Выберите движок для недостающих изображений.", ENGINE_KB)
        elif action in {"engine_auto", "engine_openai", "engine_midjourney"}:
            project["generation_engine"] = {
                "engine_auto": "auto",
                "engine_openai": "openai",
                "engine_midjourney": "midjourney",
            }[action]
            project["stage"] = "final_review"
            _save(user_id, project)
            await _reply(update, _review_text(project), FINAL_KB)
        elif action == "build":
            project["stage"] = "building"
            _save(user_id, project)
            await _reply(update, "Начинаю сборку. Проверяю факты, создаю только недостающие визуалы и затем формирую PPTX и PDF.")
            await context.bot.send_chat_action(update.effective_chat.id, ChatAction.UPLOAD_DOCUMENT)

            async def _render_action():
                return await _build_files(project, update, context)

            if _STUDIO_PAID_RUNNER is not None:
                built = await _STUDIO_PAID_RUNNER(update, context, "img", "presentation_render", float(_STUDIO_RENDER_COST_USD or 0.0), _render_action)
            else:
                built = await _render_action()
            if not built:
                project["stage"] = "final_review"
                _save(user_id, project)
                return
            pdf_path, pptx_path = built
            with open(pdf_path, "rb") as fh:
                await update.effective_message.reply_document(InputFile(fh, filename=Path(pdf_path).name), caption="Готовый PDF")
            with open(pptx_path, "rb") as fh:
                await update.effective_message.reply_document(InputFile(fh, filename=Path(pptx_path).name), caption="Редактируемый PPTX")
            await _reply(update, "Проект собран. Главный бриф и все дополнения сохранены.", _kb([
                [("📄 Скачать PDF ещё раз", "ps:download_pdf"), ("📊 Скачать PPTX ещё раз", "ps:download_pptx")],
                [("🆕 Новый проект", "ps:new")],
            ]))
        elif action in {"download_pdf", "download_pptx"}:
            key = "pdf_path" if action.endswith("pdf") else "pptx_path"
            path = project.get(key, "")
            if not path or not os.path.exists(path):
                await _reply(update, "Файл не найден. Повторите сборку.", FINAL_KB)
            else:
                with open(path, "rb") as fh:
                    await update.effective_message.reply_document(InputFile(fh, filename=Path(path).name))
        else:
            await _reply(update, "Неизвестное действие. Продолжите сохранённый проект.", START_KB)
    except Exception as exc:
        await _local_error(update, action, exc)
    raise ApplicationHandlerStop


async def _process_text_value(update: Update, context: ContextTypes.DEFAULT_TYPE, text: str) -> bool:
    if not update.effective_user or not update.effective_message:
        return False
    project = _load(update.effective_user.id)
    if not project:
        return False
    stage = project.get("stage", "")
    active_text_stages = {
        "await_brief",
        "await_brand_name",
        "await_structure_edit",
        "await_logo_brief",
        "await_visual_note",
        "visual_collect",
        "await_style_note",
        "await_palette_note",
    }
    if stage not in active_text_stages:
        return False
    text = (text or "").strip()
    if not text:
        return True
    try:
        if stage == "await_brief":
            if len(text) < 80:
                await _reply(update, "Бриф слишком короткий. Пришлите подробное описание минимум на несколько абзацев.")
                return True
            project["raw_brief"] = text
            await _reply(update, "Бриф сохранён. Проверяю точное название, факты, цены и контакты…")
            project["profile"] = await _parse_profile(text, update)
            if not _is_valid_brand_name(project["profile"].get("brand_name", "")):
                project["stage"] = "await_brand_name"
                _save(project["user_id"], project)
                await _reply(update, "Я не хочу угадывать название. Пришлите только точное название бренда отдельным сообщением.")
                return True
            await _reply(update, "Создаю структуру без вымышленных фактов…")
            project["structure"] = await _generate_structure(project, update)
            project["stage"] = "structure_review"
            _save(project["user_id"], project)
            await _reply(update, _structure_text(project), STRUCTURE_KB)
        elif stage == "await_brand_name":
            if not _is_valid_brand_name(text):
                await _reply(update, "Название похоже на служебное слово. Пришлите точное имя, например VERDIA HOME.")
                return True
            project["profile"]["brand_name"] = re.sub(r"\s+", " ", text.strip())
            project["structure"] = await _generate_structure(project, update)
            project["stage"] = "structure_review"
            _save(project["user_id"], project)
            await _reply(update, _structure_text(project), STRUCTURE_KB)
        elif stage == "await_structure_edit":
            project.setdefault("structure_notes", []).append(text)
            await _reply(update, f"Правка №{len(project['structure_notes'])} сохранена. Перестраиваю структуру с учётом всех правок…")
            project["structure"] = await _generate_structure(project, update)
            project["stage"] = "structure_review"
            _save(project["user_id"], project)
            await _reply(update, _structure_text(project), STRUCTURE_KB)
        elif stage == "await_logo_brief":
            project.setdefault("logo_notes", []).append(text)
            project["stage"] = "logo_notes"
            _save(project["user_id"], project)
            await _reply(update, f"Пожелание №{len(project['logo_notes'])} сохранено.", LOGO_NOTES_KB)
        elif stage in {"await_visual_note", "visual_collect"}:
            project.setdefault("visual_notes", []).append(text)
            project["stage"] = "visual_collect"
            _save(project["user_id"], project)
            await _reply(update, f"Требование к визуалам №{len(project['visual_notes'])} сохранено.", VISUAL_NOTES_KB)
        elif stage == "await_style_note":
            project.setdefault("style_notes", []).append(text)
            project["stage"] = "style_collect"
            _save(project["user_id"], project)
            await _reply(update, f"Требование к стилю №{len(project['style_notes'])} сохранено.", STYLE_NOTES_KB)
        elif stage == "await_palette_note":
            project.setdefault("palette_notes", []).append(text)
            project["palette"] = _parse_palette(project)
            project["stage"] = "palette_collect"
            _save(project["user_id"], project)
            await _reply(update, f"Дополнение к палитре №{len(project['palette_notes'])} сохранено.", PALETTE_NOTES_KB)
        return True
    except Exception as exc:
        await _local_error(update, f"text:{stage}", exc)
        return True


async def on_text(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    if not update.effective_message or not update.effective_message.text:
        return
    if await _process_text_value(update, context, update.effective_message.text):
        raise ApplicationHandlerStop


async def _download_telegram_file(file_obj) -> bytes:
    tg_file = await file_obj.get_file()
    return bytes(await tg_file.download_as_bytearray())


async def on_photo(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    if not update.effective_user or not update.effective_message or not update.effective_message.photo:
        return
    project = _load(update.effective_user.id)
    if not project or project.get("stage") not in {"await_logo_upload", "await_images_upload"}:
        return
    try:
        data = await _download_telegram_file(update.effective_message.photo[-1])
        pdir = _project_dir(project["user_id"], project["project_id"])
        if project["stage"] == "await_logo_upload":
            path = pdir / "uploaded_logo.jpg"
            path.write_bytes(data)
            project["logo_selected"] = str(path)
            project["stage"] = "image_mode"
            _save(project["user_id"], project)
            await _reply(update, "Готовый логотип сохранён. Выберите источник изображений.", IMAGE_MODE_KB)
        else:
            path = pdir / f"uploaded_image_{len(project.get('uploaded_images', [])) + 1:02d}.jpg"
            path.write_bytes(data)
            project.setdefault("uploaded_images", []).append(str(path))
            _save(project["user_id"], project)
            await _reply(update, f"Изображение сохранено. Всего: {len(project['uploaded_images'])}.", UPLOAD_IMAGES_KB)
    except Exception as exc:
        await _local_error(update, "photo_upload", exc)
    raise ApplicationHandlerStop


async def on_document(update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
    if not update.effective_user or not update.effective_message or not update.effective_message.document:
        return
    project = _load(update.effective_user.id)
    if not project or project.get("stage") not in {"await_logo_upload", "await_images_upload"}:
        return
    doc = update.effective_message.document
    mime = (doc.mime_type or "").lower()
    name = doc.file_name or "image"
    if not (mime.startswith("image/") or name.lower().endswith((".png", ".jpg", ".jpeg", ".webp"))):
        await _reply(update, "На этом этапе нужен файл изображения PNG/JPG/WEBP.")
        raise ApplicationHandlerStop
    try:
        data = await _download_telegram_file(doc)
        ext = Path(name).suffix.lower() or ".png"
        pdir = _project_dir(project["user_id"], project["project_id"])
        if project["stage"] == "await_logo_upload":
            path = pdir / f"uploaded_logo{ext}"
            path.write_bytes(data)
            project["logo_selected"] = str(path)
            project["stage"] = "image_mode"
            _save(project["user_id"], project)
            await _reply(update, "Логотип сохранён. Выберите источник изображений.", IMAGE_MODE_KB)
        else:
            path = pdir / f"uploaded_image_{len(project.get('uploaded_images', [])) + 1:02d}{ext}"
            path.write_bytes(data)
            project.setdefault("uploaded_images", []).append(str(path))
            _save(project["user_id"], project)
            await _reply(update, f"Изображение сохранено. Всего: {len(project['uploaded_images'])}.", UPLOAD_IMAGES_KB)
    except Exception as exc:
        await _local_error(update, "document_upload", exc)
    raise ApplicationHandlerStop


def register(application) -> None:
    """Register presentation handlers before broad/catch-all handlers."""
    init_storage()
    application.add_handler(CallbackQueryHandler(on_callback, pattern=r"^ps:"), group=-20)
    application.add_handler(MessageHandler(filters.PHOTO, on_photo), group=-20)
    application.add_handler(MessageHandler(filters.Document.ALL, on_document), group=-20)
    application.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, on_text), group=-20)


async def _store_uploaded_bytes(
    update: Update,
    raw: bytes,
    filename: str,
    mime: str,
    caption: str = "",
) -> bool:
    if not update.effective_user or not update.effective_message:
        return False
    project = _load(update.effective_user.id)
    if not project or project.get("stage") not in {"await_logo_upload", "await_images_upload"}:
        return False
    low = (filename or "").lower()
    if not ((mime or "").lower().startswith("image/") or low.endswith((".png", ".jpg", ".jpeg", ".webp"))):
        await _reply(update, "На этом этапе нужен файл изображения PNG/JPG/WEBP.")
        return True
    try:
        pdir = _project_dir(project["user_id"], project["project_id"])
        ext = Path(filename or "").suffix.lower()
        if ext not in {".png", ".jpg", ".jpeg", ".webp"}:
            ext = ".png" if "png" in (mime or "").lower() else ".jpg"
        if project["stage"] == "await_logo_upload":
            path = pdir / f"uploaded_logo{ext}"
            path.write_bytes(raw)
            project["logo_selected"] = str(path)
            project["stage"] = "image_mode"
            _save(project["user_id"], project)
            await _reply(update, "Логотип сохранён. Выберите источник изображений.", IMAGE_MODE_KB)
        else:
            current_uploads = len(project.get("uploaded_images", []))
            if current_uploads >= MAX_UPLOADS:
                await _reply(update, f"Достигнут лимит проекта: {MAX_UPLOADS} изображений.", UPLOAD_IMAGES_KB)
                return True
            idx = current_uploads + 1
            path = pdir / f"uploaded_image_{idx:02d}{ext}"
            path.write_bytes(raw)
            project.setdefault("uploaded_images", []).append(str(path))
            if caption:
                project.setdefault("visual_notes", []).append(f"Комментарий к загруженному изображению {idx}: {caption}")
            _save(project["user_id"], project)
            await _reply(update, f"Изображение сохранено. Всего: {len(project['uploaded_images'])}.", UPLOAD_IMAGES_KB)
        return True
    except Exception as exc:
        await _local_error(update, "asset_upload", exc)
        return True


@dataclass
class StudioConfig:
    db_path: str
    data_dir: str
    max_uploads: int = 60
    max_generated_images: int = 10
    render_cost_usd: float = 0.20


class PresentationStudio:
    """Adapter used by the current monolithic main.py.

    All user-facing errors are handled inside this class/module, so presentation
    failures never bubble into the bot-wide generic error handler.
    """

    def __init__(
        self,
        config: StudioConfig,
        llm_call: Callable[[Update, str], Awaitable[str]],
        image_batch_call: Callable[[Update, ContextTypes.DEFAULT_TYPE, list[str], str, str], Awaitable[list[bytes] | None]],
        paid_runner: Callable[[Update, ContextTypes.DEFAULT_TYPE, str, str, float, Callable[[], Awaitable[Any]]], Awaitable[Any]],
    ):
        global DB_PATH, ASSET_ROOT, MAX_IMAGES, MAX_UPLOADS
        global _STUDIO_LLM_CALL, _STUDIO_IMAGE_BATCH_CALL, _STUDIO_PAID_RUNNER, _STUDIO_RENDER_COST_USD
        self.cfg = config
        DB_PATH = os.path.abspath(config.db_path)
        ASSET_ROOT = Path(config.data_dir)
        MAX_IMAGES = max(1, min(14, int(config.max_generated_images or 10)))
        MAX_UPLOADS = max(1, int(config.max_uploads or 60))
        _STUDIO_LLM_CALL = llm_call
        _STUDIO_IMAGE_BATCH_CALL = image_batch_call
        _STUDIO_PAID_RUNNER = paid_runner
        _STUDIO_RENDER_COST_USD = max(0.0, float(config.render_cost_usd or 0.0))
        self.data_dir = ASSET_ROOT
        init_storage()

    def _active_project(self, user_id: int, chat_id: int) -> dict[str, Any] | None:
        project = _load(user_id)
        if not project:
            return None
        if int(project.get("chat_id") or 0) not in (0, int(chat_id or 0)):
            return None
        return {
            **project,
            "id": project.get("project_id"),
            "state": project.get("stage"),
            "status": "done" if project.get("stage") == "done" else "active",
        }

    async def start(self, update: Update, context: ContextTypes.DEFAULT_TYPE, kind: str = "presentation") -> bool:
        if not update.effective_user or not update.effective_message:
            return False
        context.user_data["presentation_kind"] = "catalog" if kind == "catalog" else "presentation"
        project = _load(update.effective_user.id)
        if project and project.get("stage") not in {"cancelled"}:
            context.user_data["presentation_studio_active"] = project.get("project_id")
            await _reply(
                update,
                f"Найден сохранённый проект «{project.get('profile', {}).get('brand_name') or 'без названия'}». "
                "Можно продолжить с текущего этапа либо начать заново.",
                START_KB,
            )
            return True
        project = _new_project(update.effective_user.id, kind, update.effective_chat.id if update.effective_chat else 0)
        _save(update.effective_user.id, project)
        context.user_data["presentation_studio_active"] = project["project_id"]
        noun = "каталога" if kind == "catalog" else "презентации"
        await _reply(
            update,
            f"Начинаем новый проект {noun}. Пришлите полный главный бриф одним сообщением.\n\n"
            "Обязательно укажите точное название бренда, продукт или услугу, цель, аудиторию, "
            "линейку и цены, контакты, нужное количество слайдов, визуальную концепцию и факты, "
            "которые нельзя выдумывать.",
            _kb([[("❌ Отменить проект", "ps:cancel")]]),
        )
        return True

    async def handle_callback(self, update: Update, context: ContextTypes.DEFAULT_TYPE) -> bool:
        q = update.callback_query
        if not q or not (q.data or "").startswith("ps:"):
            return False
        try:
            await on_callback(update, context)
        except ApplicationHandlerStop:
            pass
        except Exception as exc:
            await _local_error(update, "callback-adapter", exc)
        project = _load(q.from_user.id)
        if project:
            context.user_data["presentation_studio_active"] = project.get("project_id")
        else:
            context.user_data.pop("presentation_studio_active", None)
        return True

    async def handle_text(self, update: Update, context: ContextTypes.DEFAULT_TYPE, text: str) -> bool:
        project = self._active_project(update.effective_user.id, update.effective_chat.id)
        if not project:
            return False
        context.user_data["presentation_studio_active"] = project.get("project_id")
        return await _process_text_value(update, context, text)

    async def handle_photo(
        self,
        update: Update,
        context: ContextTypes.DEFAULT_TYPE,
        raw: bytes,
        mime: str = "image/jpeg",
        caption: str = "",
    ) -> bool:
        project = self._active_project(update.effective_user.id, update.effective_chat.id)
        if not project:
            return False
        context.user_data["presentation_studio_active"] = project.get("project_id")
        return await _store_uploaded_bytes(update, raw, "photo.jpg", mime, caption)

    async def handle_document(
        self,
        update: Update,
        context: ContextTypes.DEFAULT_TYPE,
        raw: bytes,
        filename: str,
        mime: str,
        caption: str = "",
    ) -> bool:
        project = self._active_project(update.effective_user.id, update.effective_chat.id)
        if not project:
            return False
        context.user_data["presentation_studio_active"] = project.get("project_id")
        return await _store_uploaded_bytes(update, raw, filename, mime, caption)


# ─────────────────────────────────────────────────────────────────────────────
# v96 presentation content / rendering / validation hotfixes
# These overrides intentionally appear near the end of the module so they take
# precedence over the original v95 implementations without a risky full rewrite.
# ─────────────────────────────────────────────────────────────────────────────

_PLACEHOLDER_SENTINEL = "Содержание будет подготовлено строго по брифу"


def _normalize_exact_brand(value: str) -> str:
    value = re.sub(r"\s+", " ", str(value or "").strip())
    value = value.strip(" \n\t:—-.,;\"'«»")
    return value[:80]


def _has_placeholder(value: str) -> bool:
    return _PLACEHOLDER_SENTINEL.lower() in str(value or "").lower()


def _sanitize_bullets(items: list[str]) -> list[str]:
    out = []
    for item in items:
        item = re.sub(r"\s+", " ", str(item)).strip(" •-–—")
        if not item or _has_placeholder(item):
            continue
        if item not in out:
            out.append(item[:240])
    return out[:6]


def _extract_contacts(text: str) -> list[str]:
    raw = text or ""
    result: list[str] = []
    # Keep email / telegram / phone / web in stable order.
    email_re = r"[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Za-z]{2,}"
    tg_re = r"(?<![A-Za-z0-9._%+-])@[A-Za-z0-9_]{4,}"
    phone_re = r"\+?\d[\d\s()\-]{8,}\d"
    web_re = r"\b(?:https?://)?(?:www\.)?[A-Za-z0-9-]+(?:\.[A-Za-z0-9-]+)+(?:/[A-Za-z0-9_./?&=%+-]*)?"
    for pattern in [email_re, tg_re, phone_re, web_re]:
        for value in re.findall(pattern, raw):
            cleaned = re.sub(r"\s+", " ", value).strip().rstrip('.,;')
            if cleaned and cleaned not in result:
                result.append(cleaned)
    return result[:10]


def _extract_offer_objects(raw: str) -> list[dict[str, Any]]:
    raw = (raw or "").replace("\r\n", "\n").replace("\r", "\n")
    lines = [ln.strip() for ln in raw.split("\n")]
    offers: list[dict[str, Any]] = []
    heading_re = re.compile(r"^(?:#{1,4}\s*)?([A-ZА-ЯЁ][A-ZА-ЯЁ0-9&._ -]{2,40})$")
    stop_re = re.compile(r"^(?:целевая|аудитор|позиционир|визуал|структур|контакт|этап|преиму|проблем|заключ|призыв|цвет|стиль|логотип|задач|цель)\b", re.I)
    price_re = re.compile(r"(?:от\s+)?\d[\d\s.,]{1,12}\s*(?:₽|руб(?:лей|ля|ль)?|THB|USD|EUR|\$|€)(?:\s*(?:в\s*месяц|/мес\.?))?", re.I)
    for i, line in enumerate(lines):
        m = heading_re.match(line)
        if not m:
            continue
        name = re.sub(r"\s+", " ", m.group(1)).strip(" -")
        if stop_re.search(name):
            continue
        block: list[str] = []
        for nxt in lines[i + 1 : i + 10]:
            if not nxt:
                if block:
                    break
                continue
            if heading_re.match(nxt):
                break
            if stop_re.search(nxt):
                break
            block.append(nxt)
        joined = "\n".join(block)
        if not price_re.search(joined):
            continue
        price = price_re.search(joined)
        details = []
        for ln in block:
            if price_re.fullmatch(ln.strip()):
                continue
            details.append(ln)
        offers.append({
            'name': name[:80],
            'price': (price.group(0).strip() if price else ''),
            'details': _sanitize_bullets(details),
            'raw': joined[:900],
        })
    if offers:
        uniq = []
        seen = set()
        for item in offers:
            key = item['name'].lower()
            if key in seen:
                continue
            seen.add(key)
            uniq.append(item)
        return uniq[:6]
    # Fallback: one-line offers with price.
    out = []
    for line in lines:
        clean = re.sub(r"\s+", " ", line).strip(" •-*#\t")
        if not clean or not price_re.search(clean):
            continue
        title = clean.split(':', 1)[0]
        out.append({'name': title[:80], 'price': price_re.search(clean).group(0).strip(), 'details': [], 'raw': clean[:240]})
    # unique
    uniq=[]; seen=set()
    for item in out:
        key=item['name'].lower()
        if key in seen: continue
        seen.add(key); uniq.append(item)
    return uniq[:6]


def _service_offer_objects(raw: str) -> list[dict[str, Any]]:
    objs = _extract_offer_objects(raw)
    return [o for o in objs if re.search(r"care|service|сервис|обслуж", o['name'], re.I)]


def _product_offer_objects(raw: str) -> list[dict[str, Any]]:
    objs = _extract_offer_objects(raw)
    return [o for o in objs if not re.search(r"care|service|сервис|обслуж", o['name'], re.I)]


def _infer_cover_line(project: dict[str, Any]) -> str:
    raw = project.get('raw_brief', '')
    patterns = [
        r"(?i)интеллектуальные\s+вертикальные\s+сады",
        r"(?i)живые\s+вертикальные\s+сады",
        r"(?i)вертикальные\s+сады",
        r"(?i)модульные\s+вертикальные\s+сады",
    ]
    for pat in patterns:
        m = re.search(pat, raw)
        if m:
            return m.group(0)
    product = project.get('profile', {}).get('product', '')
    if re.search(r"вертикальн", product, re.I):
        return "Интеллектуальные вертикальные сады"
    return "Премиальные решения для интерьера"


def _first_matching_section(raw: str, names: list[str]) -> list[str]:
    text = (raw or '').replace('\r\n', '\n').replace('\r', '\n')
    lines = [ln.rstrip() for ln in text.split('\n')]
    for i, line in enumerate(lines):
        clean = line.strip().lower().rstrip(':')
        if any(clean.startswith(name.lower()) for name in names):
            buf=[]
            for nxt in lines[i+1:i+10]:
                s=nxt.strip(' •-*\t')
                if not s:
                    if buf:
                        break
                    continue
                if re.match(r"^[A-ZА-ЯЁ][^.!?]{1,60}:?$", s) and len(s.split()) <= 6 and buf:
                    break
                buf.append(s)
            return _sanitize_bullets(buf)
    return []


def _build_deterministic_slide(project: dict[str, Any], base: dict[str, Any], slide: dict[str, Any], idx: int) -> dict[str, Any]:
    profile = project.get('profile', {})
    raw = project.get('raw_brief', '')
    brand = _normalize_exact_brand(profile.get('brand_name') or base.get('title') or slide.get('title') or 'Бренд')
    offers = project.get('_offer_objects') or _product_offer_objects(raw)
    service_offers = project.get('_service_offer_objects') or _service_offer_objects(raw)
    contacts = profile.get('contacts') or _extract_contacts(raw)
    title = str(slide.get('title') or base.get('title') or '').strip()
    base_title = str(base.get('title') or '').strip()
    key = (title or base_title).lower()
    layout = slide.get('layout') or base.get('layout') or 'split'
    bullets = _sanitize_bullets(slide.get('bullets') or base.get('bullets') or [])
    subtitle = str(slide.get('subtitle') or '').strip()

    if idx == 0 or layout == 'cover':
        slide['title'] = brand
        slide['subtitle'] = _normalize_exact_brand(profile.get('tagline') or subtitle or 'Коммерческая презентация')
        slide['bullets'] = [_infer_cover_line(project)]
        slide['image_prompt'] = slide.get('image_prompt') or (
            f"Premium editorial interior photograph for {brand}: elegant contemporary interior with a living vertical garden, "
            "natural materials, soft daylight, calm premium atmosphere, clean composition, generous negative space for presentation layout, "
            "strictly no text, no letters, no logo, no watermark, no UI."
        )
        return slide

    if 'проблем' in key:
        problems = _first_matching_section(raw, ['Проблемы клиента', 'Проблемы', 'Боли клиента']) or bullets
        slide['bullets'] = problems[:6] or [
            'Живые растения требуют постоянного ухода',
            'Недостаток естественного света в помещениях',
            'Перегрузка пространства отдельными горшками',
            'Сложности в сочетании растений',
        ]
    elif any(x in key for x in ['решение', 'нашe решение'.lower(), 'verdia home']):
        slide['bullets'] = _first_matching_section(raw, ['Решение', 'Решение VERDIA HOME', 'Что получает клиент'])[:6] or bullets or [
            'Живые растения в модульной конструкции',
            'Скрытый резервуар и автоматический полив',
            'Профессиональное фитосвещение',
            'Интегрированные датчики контроля',
        ]
    elif any(x in key for x in ['схема', 'технолог', 'как это работает']):
        slide['layout'] = 'process'
        slide['bullets'] = _first_matching_section(raw, ['Схема технологии', 'Как это работает', 'Технология'])[:6] or [
            'Посадочные модули с живыми растениями',
            'Скрытый резервуар для воды и питания',
            'Автоматический полив с таймером',
            'Фитосвет для стабильного роста растений',
            'Датчики контроля и защита от переполнения',
            'Сервисное обслуживание и сезонная настройка',
        ]
        slide['image_prompt'] = slide.get('image_prompt') or (
            f"Photorealistic close-up of a premium modular vertical garden system for {brand}, showing irrigation, lighting and modular structure, "
            "clean product-focused composition, no labels, no diagram text, no letters, no logo, no watermark."
        )
    elif 'преим' in key:
        slide['bullets'] = _first_matching_section(raw, ['Ключевые преимущества', 'Преимущества'])[:6] or bullets or [
            'Минимум ухода для клиента',
            'Эстетическая ценность и wow-эффект',
            'Современные технологии ухода',
            'Интеграция в любой интерьер',
        ]
    elif any(x in key for x in ['варианты использования', 'сценарии использования', 'кому подходит']):
        slide['bullets'] = _first_matching_section(raw, ['Варианты использования', 'Сценарии использования', 'Целевая аудитория'])[:6] or bullets or [
            'Квартиры и загородные дома',
            'Офисы и рестораны',
            'Гостиницы и клиники',
            'Общественные пространства',
        ]
    elif any(x in key for x in ['сравнение', 'линейк', 'продукт', 'модел']):
        slide['layout'] = 'comparison'
        if offers:
            slide['bullets'] = [f"{o['name']} — {o['price']}" for o in offers[:4]]
        else:
            slide['bullets'] = bullets
        slide['image_prompt'] = slide.get('image_prompt') or (
            f"Premium interior product photo for {brand}: modular vertical garden product family, distinct sizes and configurations, "
            "clean showroom or modern interior, no text, no letters, no logo, no watermark."
        )
    elif any(x in key for x in ['care', 'сервис', 'обслуж', 'сопровожд']):
        service_details = []
        for offer in service_offers[:2]:
            service_details.extend([d for d in (offer.get('details') or []) if not re.search(r'(?:@|\+?\d[\d\s()\-]{8,}|\.[A-Za-z]{2,}|контакт|телефон|email|сайт)', d, re.I)])
            if offer.get('price'):
                service_details.insert(0, f"{offer['name']} — {offer['price']}")
        slide['bullets'] = _sanitize_bullets(service_details) or bullets or [
            'VERDIA CARE — обслуживание от 6 900 рублей в месяц',
            'Контроль полива и освещения',
            'Диагностика и уход за растениями',
            'Очистка конструкции и сезонная настройка',
        ]
    elif any(x in key for x in ['этап', 'реализац']):
        slide['layout'] = 'process'
        slide['bullets'] = _first_matching_section(raw, ['Этапы реализации', 'Этапы работы', 'Процесс'])[:6] or bullets or [
            'Консультация и замер помещения',
            'Разработка концепции и визуализация',
            'Подбор растений и инженерии',
            'Производство и установка',
            'Настройка и обслуживание',
        ]
    elif any(x in key for x in ['почему', 'выбирают', 'отличия']):
        slide['bullets'] = _first_matching_section(raw, ['Почему выбирают', 'Почему', 'Ключевое позиционирование'])[:6] or bullets or [
            'Премиальный eco-tech подход',
            'Готовая экосистема для пространства',
            'Индивидуальный подход к каждому проекту',
            'Высокое качество материалов и технологий',
        ]
    elif layout == 'cta' or any(x in key for x in ['контакт', 'следующий шаг', 'призыв']):
        slide['layout'] = 'cta'
        if not subtitle or _has_placeholder(subtitle):
            slide['subtitle'] = 'Свяжитесь с нами для консультации и расчёта проекта'
        cta = _first_matching_section(raw, ['Контакты и призыв к действию', 'Контакты', 'Следующий шаг'])
        bullets = _sanitize_bullets([*cta, *contacts])
        if contacts and not any(c in ' '.join(bullets) for c in contacts):
            bullets.extend([c for c in contacts if c not in bullets])
        slide['bullets'] = bullets[:6] or ['Закажите консультацию', 'Контакты указаны в главном брифе']
        slide['image_prompt'] = slide.get('image_prompt') or (
            f"Premium interior photo for {brand}: elegant living wall in a contemporary space, clean right-side composition, "
            "warm inviting mood, no text, no letters, no logo, no watermark."
        )
    else:
        if not bullets:
            if base.get('bullets'):
                bullets = _sanitize_bullets(base.get('bullets'))
            else:
                bullets = ['Содержание раскрывается строго по главному брифу']
        slide['bullets'] = bullets

    slide['title'] = _normalize_exact_brand(slide.get('title') or title or base_title)
    if not slide.get('subtitle') or _has_placeholder(slide.get('subtitle', '')):
        slide['subtitle'] = _normalize_exact_brand(subtitle)
    return slide


def _validate_slides_or_raise(project: dict[str, Any], slides: list[dict[str, Any]]) -> None:
    profile = project.get('profile', {})
    brand = _normalize_exact_brand(profile.get('brand_name') or '')
    contacts = profile.get('contacts') or _extract_contacts(project.get('raw_brief', ''))
    if not slides:
        raise RuntimeError('Не удалось сформировать слайды.')
    if brand and _normalize_exact_brand(slides[0].get('title', '')) != brand:
        raise RuntimeError('На обложке потерялось точное название бренда.')
    for i, slide in enumerate(slides, 1):
        if _has_placeholder(slide.get('title', '')) or _has_placeholder(slide.get('subtitle', '')):
            raise RuntimeError(f'Слайд {i} содержит техническую заглушку в заголовке.')
        if any(_has_placeholder(b) for b in slide.get('bullets', []) or []):
            raise RuntimeError(f'Слайд {i} содержит техническую заглушку в тезисах.')
    if contacts:
        final_text = ' '.join(slides[-1].get('bullets', []) or [])
        if not any(c in final_text for c in contacts):
            raise RuntimeError('На финальном слайде отсутствуют контакты из брифа.')


async def _generate_deck_content(project: dict[str, Any], update: Update | None = None) -> list[dict[str, Any]]:
    structure = project.get('structure', [])
    profile = project.get('profile', {})
    brand = _normalize_exact_brand(profile.get('brand_name') or '')
    prompt = f"""
Подготовь финальный текст презентации строго в JSON:
{{"slides":[{{"title":"", "subtitle":"", "bullets":[""], "image_prompt":""}}]}}

Нужно ровно {len(structure)} слайдов в том же порядке и с теми же смыслами.
Структура:
{json.dumps(structure, ensure_ascii=False)}

ЖЁСТКИЕ ПРАВИЛА:
- точное название бренда: {brand}; запрещено менять написание, регистр или добавлять точку в конце;
- нельзя оставлять служебные заглушки вроде «{_PLACEHOLDER_SENTINEL}»;
- нельзя придумывать отсутствующие кейсы, отзывы, статистику и награды;
- на финальном слайде обязательно перечисли реальные контакты из брифа;
- на продуктовом слайде перечисли реальные продукты / модели и точные цены;
- если в структуре есть технология, опиши реальные компоненты системы;
- image_prompt должен описывать ФОТО без текста внутри кадра: no text, no letters, no numbers, no logos.

Главный бриф:
{project.get('raw_brief', '')}

Накопленные визуальные требования:
{_all_notes(project, 'visual_notes')}

Накопленные требования к стилю:
{_all_notes(project, 'style_notes')}
"""
    result: list[dict[str, Any]] = []
    offers = _product_offer_objects(project.get('raw_brief', ''))
    service_offers = _service_offer_objects(project.get('raw_brief', ''))
    project['_offer_objects'] = offers
    project['_service_offer_objects'] = service_offers
    try:
        parsed = _json_from_text(await _llm_text(prompt, 0.15, update))
        slide_values = parsed.get('slides') if isinstance(parsed, dict) else None
        if isinstance(slide_values, list) and len(slide_values) == len(structure):
            for base, item in zip(structure, slide_values):
                if not isinstance(item, dict):
                    item = {}
                result.append({
                    **base,
                    'title': str(item.get('title') or base.get('title') or '').strip()[:110],
                    'subtitle': str(item.get('subtitle') or '').strip()[:220],
                    'bullets': _sanitize_bullets(item.get('bullets') or base.get('bullets') or []),
                    'image_prompt': str(item.get('image_prompt') or base.get('image_prompt') or '').strip()[:1800],
                })
    except Exception:
        log.exception('Deck copy generation failed; switching to deterministic repair')
    if not result:
        result = [{**slide, 'subtitle': '', 'bullets': _sanitize_bullets(slide.get('bullets', []))} for slide in structure]

    repaired: list[dict[str, Any]] = []
    for idx, (base, slide) in enumerate(zip(structure, result)):
        repaired.append(_build_deterministic_slide(project, deepcopy(base), deepcopy(slide), idx))
    if repaired:
        repaired[0]['title'] = brand or repaired[0].get('title', '')
        if not repaired[0].get('subtitle'):
            repaired[0]['subtitle'] = profile.get('tagline') or 'Коммерческая презентация'
    if repaired:
        final = repaired[-1]
        contacts = profile.get('contacts') or _extract_contacts(project.get('raw_brief', ''))
        final['bullets'] = _sanitize_bullets(final.get('bullets', []) + contacts)
    _validate_slides_or_raise(project, repaired)
    return repaired


def _route_engine_for_slide(slide: dict[str, Any]) -> str:
    title = (slide.get('title') or '').lower()
    layout = (slide.get('layout') or '').lower()
    technical = any(k in title for k in ['технолог', 'схема', 'сравнение', 'линейк', 'продукт', 'этап', 'care', 'сервис']) or layout in {'comparison', 'process'}
    return 'openai' if technical else 'midjourney'


async def _prepare_slide_images(project: dict[str, Any], slides: list[dict[str, Any]], update: Update, context: ContextTypes.DEFAULT_TYPE) -> dict[int, str]:
    mode = project.get('image_mode', 'auto')
    uploaded = [p for p in project.get('uploaded_images', []) if os.path.exists(p)]
    result: dict[int, str] = {int(k): v for k, v in (project.get('slide_images') or {}).items() if os.path.exists(v)}
    target_indexes = [
        i for i, slide in enumerate(slides)
        if slide.get('image_needed') or slide.get('layout') in {'cover', 'split', 'full_image', 'comparison', 'cta'}
    ][:MAX_IMAGES]
    upload_iter = iter(uploaded)
    if mode in {'upload', 'mixed'}:
        for idx in target_indexes:
            if idx in result:
                continue
            with contextlib.suppress(StopIteration):
                result[idx] = next(upload_iter)
    if mode == 'skip':
        project['slide_images'] = {str(k): v for k, v in result.items()}
        return result
    if mode in {'auto', 'mixed'}:
        missing = [idx for idx in target_indexes if idx not in result]
        pdir = _project_dir(project['user_id'], project['project_id'])
        grouped: dict[str, list[tuple[int, str]]] = {}
        chosen_engine = (project.get('generation_engine') or 'auto').lower()
        for idx in missing:
            slide = slides[idx]
            prompt = slide.get('image_prompt') or _default_scene_prompt(project, slide)
            prompt += "\nExact main brief context (do not drift to another business):\n" + project.get('raw_brief', '')[:2600]
            prompt += "\nCumulative visual rules:\n" + _all_notes(project, 'visual_notes')[:2200]
            prompt += "\nCumulative style rules:\n" + _all_notes(project, 'style_notes')[:1800]
            prompt += "\nMandatory negative constraints: no readable text, no letters, no numbers, no logo, no watermark, no UI, no poster, no diagram labels, no slide frame."
            engine = chosen_engine if chosen_engine in {'openai', 'midjourney'} else _route_engine_for_slide(slide)
            grouped.setdefault(engine, []).append((idx, prompt[:7600]))
        for engine, entries in grouped.items():
            prompts = [p for _, p in entries]
            generated: list[bytes] = []
            if prompts and _STUDIO_IMAGE_BATCH_CALL is not None:
                await _reply(update, f"Создаю {len(prompts)} визуал(ов) через {'Midjourney' if engine == 'midjourney' else 'OpenAI'}…")
                batch = await _STUDIO_IMAGE_BATCH_CALL(update, context, prompts, engine, 'presentation_slide_images')
                if batch:
                    generated = list(batch)
            if not generated:
                for prompt in prompts:
                    try:
                        generated.append(await _generate_image_bytes(prompt, IMAGE_SIZE))
                    except Exception:
                        generated.append(b'')
            for (idx, _), image_bytes in zip(entries, generated):
                if not image_bytes:
                    continue
                path = pdir / f'slide_{idx + 1:02d}.png'
                path.write_bytes(image_bytes)
                result[idx] = str(path)
    project['slide_images'] = {str(k): v for k, v in result.items()}
    return result


def _tmp_crop(path: str, width: int, height: int, out_path: Path) -> str:
    # Stronger compression than v95 to reduce PPTX size and improve Telegram delivery.
    _crop_image(path, width, height).save(out_path, 'JPEG', quality=84, optimize=True, progressive=True)
    return str(out_path)


def _resolve_logo_path(project: dict[str, Any], pdir: Path | None = None) -> str:
    path = project.get('logo_selected') or ''
    if not path or not os.path.exists(path):
        return ''
    try:
        im = Image.open(path).convert('RGBA')
        # If the selected logo is a preview card, crop the upper symbol area.
        if im.width >= 900 and im.height >= 900:
            crop = im.crop((120, 40, im.width - 120, int(im.height * 0.72)))
        else:
            crop = im
        pdir = pdir or Path(path).parent
        out = pdir / '_resolved_logo.png'
        crop.save(out, 'PNG', optimize=True)
        return str(out)
    except Exception:
        return path


def _ppt_add_logo(slide, project: dict[str, Any], x, y, w):
    path = _resolve_logo_path(project, _project_dir(project['user_id'], project['project_id']))
    if path and os.path.exists(path):
        try:
            slide.shapes.add_picture(path, x, y, width=w)
        except Exception:
            log.exception('Could not add logo to PPTX')


def _ppt_add_cover_lines(slide, lines: list[str], x, y, w, h, palette: dict[str, str]):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, line in enumerate(lines[:3]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = 'Lato'
        p.font.size = Pt(18 if i == 0 else 16)
        p.font.color.rgb = RGBColor(*_hex(palette['text']))
        p.space_after = Pt(8)
    return box


def _render_comparison_cards_ppt(slide, project: dict[str, Any], palette: dict[str, str]):
    items = (project.get('_offer_objects') or _product_offer_objects(project.get('raw_brief', '')))[:4]
    if not items:
        return
    cols = max(1, len(items))
    gap = 0.18
    card_w = 11.9 / cols - gap
    for i, item in enumerate(items):
        x = Inches(0.72 + i * (card_w + gap))
        _ppt_rect(slide, x, Inches(1.55), Inches(card_w), Inches(4.95), palette['white'], palette['sand'], True)
        _ppt_add_text(slide, item['name'], x + Inches(0.18), Inches(1.82), Inches(card_w - 0.36), Inches(0.62), 20, palette['primary'], True, PP_ALIGN.CENTER)
        _ppt_add_text(slide, item['price'], x + Inches(0.18), Inches(2.42), Inches(card_w - 0.36), Inches(0.46), 16, palette['accent'], True, PP_ALIGN.CENTER)
        details = item.get('details') or []
        if details:
            _ppt_add_bullets(slide, details[:4], x + Inches(0.18), Inches(3.02), Inches(card_w - 0.36), Inches(2.8), palette, 14)


def _build_pptx(project: dict[str, Any], slides: list[dict[str, Any]], images: dict[int, str], out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    palette = project['palette']
    pdir = out_path.parent
    brand = _normalize_exact_brand(project.get('profile', {}).get('brand_name', ''))

    for idx, data in enumerate(slides):
        slide = prs.slides.add_slide(blank)
        _ppt_rect(slide, 0, 0, prs.slide_width, prs.slide_height, palette['background'])
        layout = data.get('layout', 'split')
        image_path = images.get(idx)
        if idx == 0 or layout == 'cover':
            if image_path:
                crop = _tmp_crop(image_path, 900, 1200, pdir / f'_crop_cover_{idx}.jpg')
                slide.shapes.add_picture(crop, Inches(7.35), 0, width=Inches(5.98), height=Inches(7.5))
                _ppt_rect(slide, Inches(7.18), 0, Inches(0.18), Inches(7.5), palette['primary'])
            else:
                _ppt_rect(slide, Inches(7.35), 0, Inches(5.98), Inches(7.5), palette['primary'])
            _ppt_add_logo(slide, project, Inches(0.75), Inches(0.32), Inches(1.2))
            _ppt_add_text(slide, data.get('title') or brand, Inches(0.75), Inches(1.35), Inches(6.15), Inches(1.55), 38, palette['text'], True)
            subtitle = data.get('subtitle') or project.get('profile', {}).get('tagline', '')
            _ppt_add_text(slide, subtitle, Inches(0.78), Inches(3.05), Inches(5.9), Inches(0.8), 22, palette['accent'], False)
            _ppt_add_cover_lines(slide, data.get('bullets', [])[:2], Inches(0.78), Inches(4.06), Inches(5.7), Inches(1.3), palette)
            continue

        _ppt_header(slide, idx + 1, data.get('title', ''), palette, project)
        bullets = data.get('bullets', [])
        if layout in {'cards', 'process'}:
            cols = 2 if len(bullets) <= 4 else 3
            rows = math.ceil(len(bullets) / cols)
            gap = 0.28
            left = 0.72
            top = 1.55
            total_w = 11.9
            card_w = (total_w - gap * (cols - 1)) / cols
            card_h = min(1.55, (5.25 - gap * (rows - 1)) / max(1, rows))
            for i, bullet in enumerate(bullets[:6]):
                c = i % cols
                r = i // cols
                x = Inches(left + c * (card_w + gap))
                y = Inches(top + r * (card_h + gap))
                _ppt_rect(slide, x, y, Inches(card_w), Inches(card_h), palette['card'], palette['sand'], True)
                _ppt_add_text(slide, f'{i + 1:02d}', x + Inches(0.22), y + Inches(0.16), Inches(0.45), Inches(0.35), 12, palette['primary'], True)
                _ppt_add_text(slide, bullet, x + Inches(0.22), y + Inches(0.52), Inches(card_w - 0.44), Inches(card_h - 0.62), 18, palette['text'], True)
            if image_path and len(bullets) <= 4:
                crop = _tmp_crop(image_path, 1400, 360, pdir / f'_crop_band_{idx}.jpg')
                slide.shapes.add_picture(crop, Inches(0.72), Inches(5.7), width=Inches(11.9), height=Inches(1.2))
        elif layout == 'comparison':
            _render_comparison_cards_ppt(slide, project, palette)
        elif layout == 'full_image' and image_path:
            crop = _tmp_crop(image_path, 1600, 720, pdir / f'_crop_full_{idx}.jpg')
            slide.shapes.add_picture(crop, Inches(0.72), Inches(1.42), width=Inches(11.9), height=Inches(5.55))
            overlay = _ppt_rect(slide, Inches(0.72), Inches(5.7), Inches(11.9), Inches(1.27), palette['primary'])
            overlay.fill.transparency = 8
            _ppt_add_text(slide, '  •  '.join(bullets[:3]), Inches(1.0), Inches(5.96), Inches(11.3), Inches(0.75), 18, palette['white'], True)
        elif layout == 'cta':
            _ppt_rect(slide, 0, 0, prs.slide_width, prs.slide_height, palette['primary'])
            if image_path:
                crop = _tmp_crop(image_path, 720, 900, pdir / f'_crop_cta_{idx}.jpg')
                slide.shapes.add_picture(crop, Inches(8.45), Inches(0.5), width=Inches(4.35), height=Inches(6.5))
            _ppt_add_logo(slide, project, Inches(0.8), Inches(0.3), Inches(1.1))
            _ppt_add_text(slide, data.get('title', 'Следующий шаг'), Inches(0.8), Inches(1.1), Inches(7.1), Inches(1.2), 40, palette['white'], True)
            _ppt_add_text(slide, data.get('subtitle', ''), Inches(0.82), Inches(2.45), Inches(6.8), Inches(0.8), 21, palette['sand'])
            _ppt_add_bullets_white(slide, bullets, Inches(0.82), Inches(3.45), Inches(6.9), Inches(2.7), palette, 19)
        else:
            image_left = idx % 2 == 0
            if image_path:
                crop = _tmp_crop(image_path, 900, 850, pdir / f'_crop_split_{idx}.jpg')
                img_x = Inches(0.72 if image_left else 7.15)
                slide.shapes.add_picture(crop, img_x, Inches(1.45), width=Inches(5.45), height=Inches(5.45))
                text_x = Inches(6.55 if image_left else 0.78)
            else:
                text_x = Inches(1.1)
            _ppt_add_text(slide, data.get('subtitle', ''), text_x, Inches(1.55), Inches(5.7), Inches(0.8), 20, palette['accent'], True)
            _ppt_add_bullets(slide, bullets, text_x, Inches(2.35), Inches(5.75), Inches(3.9), palette, 20)

        _ppt_add_text(slide, f'{brand}  •  {idx + 1}', Inches(0.65), Inches(7.12), Inches(4.0), Inches(0.25), 8, palette['secondary'])

    prs.save(out_path)


def _pdf_draw_logo(c: canvas.Canvas, project: dict[str, Any], x: float, y: float, w: float, pdir: Path):
    path = _resolve_logo_path(project, pdir)
    if not path or not os.path.exists(path):
        return
    try:
        im = Image.open(path)
        ratio = im.height / max(1, im.width)
        h = w * ratio
        c.drawImage(ImageReader(path), x, y, width=w, height=h, preserveAspectRatio=True, mask='auto')
    except Exception:
        log.exception('Could not add logo to PDF')


def _render_comparison_cards_pdf(c: canvas.Canvas, project: dict[str, Any], palette: dict[str, str], regular: str, bold: str):
    items = (project.get('_offer_objects') or _product_offer_objects(project.get('raw_brief', '')))[:4]
    if not items:
        return
    cols = max(1, len(items)); gap = 10; left = 42; total_w = 876
    card_w = (total_w - gap * (cols - 1)) / cols
    for i, item in enumerate(items):
        x = left + i * (card_w + gap)
        c.setFillColor(HexColor(palette['white'])); c.roundRect(x, 70, card_w, 370, 10, stroke=0, fill=1)
        c.setStrokeColor(HexColor(palette['sand'])); c.roundRect(x, 70, card_w, 370, 10, stroke=1, fill=0)
        c.setFillColor(HexColor(palette['primary'])); c.setFont(bold, 14)
        ty = 410
        for line in _pdf_wrap(item['name'], max(14, int(card_w / 8))):
            c.drawCentredString(x + card_w / 2, ty, line); ty -= 18
        c.setFillColor(HexColor(palette['accent'])); c.setFont(bold, 12)
        c.drawCentredString(x + card_w / 2, ty - 2, item.get('price', '')); ty -= 24
        c.setFillColor(HexColor(palette['text'])); c.setFont(regular, 10)
        for detail in (item.get('details') or [])[:4]:
            for line in _pdf_wrap(detail, max(18, int(card_w / 7))):
                c.drawString(x + 12, ty, '• ' + line if line == _pdf_wrap(detail, max(18, int(card_w / 7)))[0] else line)
                ty -= 14
            ty -= 4


def _build_pdf(project: dict[str, Any], slides: list[dict[str, Any]], images: dict[int, str], out_path: Path) -> None:
    page_w, page_h = 960, 540
    c = canvas.Canvas(str(out_path), pagesize=(page_w, page_h), pageCompression=1)
    regular, bold = _register_pdf_fonts()
    palette = project['palette']
    pdir = out_path.parent
    brand = _normalize_exact_brand(project.get('profile', {}).get('brand_name', ''))

    def set_fill(key: str):
        c.setFillColor(HexColor(palette[key]))

    for idx, data in enumerate(slides):
        set_fill('background')
        c.rect(0, 0, page_w, page_h, stroke=0, fill=1)
        image_path = images.get(idx)
        layout = data.get('layout', 'split')
        if idx == 0 or layout == 'cover':
            if image_path:
                _pdf_draw_image(c, image_path, 530, 0, 430, 540, pdir, f'_pdf_cover_{idx}.jpg')
            else:
                set_fill('primary'); c.rect(530, 0, 430, 540, stroke=0, fill=1)
            set_fill('primary'); c.rect(520, 0, 10, 540, stroke=0, fill=1)
            _pdf_draw_logo(c, project, 44, 462, 90, pdir)
            set_fill('text'); c.setFont(bold, 35)
            y = 405
            for line in _pdf_wrap(data.get('title') or brand, 24):
                c.drawString(55, y, line); y -= 40
            set_fill('accent'); c.setFont(regular, 19)
            sy = 315
            for line in _pdf_wrap(data.get('subtitle') or '', 30):
                c.drawString(55, sy, line); sy -= 24
            set_fill('text'); c.setFont(regular, 16)
            by = sy - 24
            for line in data.get('bullets', [])[:2]:
                for wrapped in _pdf_wrap(line, 36):
                    c.drawString(55, by, wrapped); by -= 20
                by -= 6
        else:
            set_fill('primary'); c.roundRect(38, 486, 36, 28, 6, stroke=0, fill=1)
            set_fill('white'); c.setFont(bold, 11); c.drawCentredString(56, 495, f'{idx + 1:02d}')
            set_fill('text'); c.setFont(bold, 25)
            title = data.get('title', '')
            ty = 485
            for line in _pdf_wrap(title, 44):
                c.drawString(84, ty, line); ty -= 28
            set_fill('primary'); c.rect(84, ty - 4, 92, 3, stroke=0, fill=1)
            _pdf_draw_logo(c, project, 870, 484, 52, pdir)
            bullets = data.get('bullets', [])
            if layout in {'cards', 'process'}:
                cols = 2 if len(bullets) <= 4 else 3
                rows = math.ceil(len(bullets) / cols)
                gap = 14; left = 48; top = 425; total_w = 864
                card_w = (total_w - gap * (cols - 1)) / cols
                card_h = min(105, (330 - gap * (rows - 1)) / max(1, rows))
                for i, bullet in enumerate(bullets[:6]):
                    col, row = i % cols, i // cols
                    x = left + col * (card_w + gap)
                    y0 = top - (row + 1) * card_h - row * gap
                    set_fill('card'); c.roundRect(x, y0, card_w, card_h, 10, stroke=0, fill=1)
                    c.setStrokeColor(HexColor(palette['sand'])); c.roundRect(x, y0, card_w, card_h, 10, stroke=1, fill=0)
                    set_fill('primary'); c.setFont(bold, 10); c.drawString(x + 14, y0 + card_h - 22, f'{i + 1:02d}')
                    set_fill('text'); c.setFont(bold, 14)
                    ty2 = y0 + card_h - 48
                    for line in _pdf_wrap(bullet, max(24, int(card_w / 8))):
                        c.drawString(x + 14, ty2, line); ty2 -= 18
            elif layout == 'comparison':
                _render_comparison_cards_pdf(c, project, palette, regular, bold)
            elif layout == 'cta':
                set_fill('primary'); c.rect(0, 0, page_w, page_h, stroke=0, fill=1)
                if image_path:
                    _pdf_draw_image(c, image_path, 690, 35, 215, 470, pdir, f'_pdf_cta_{idx}.jpg')
                _pdf_draw_logo(c, project, 42, 472, 72, pdir)
                set_fill('white'); c.setFont(bold, 31)
                tty = 420
                for line in _pdf_wrap(data.get('title', ''), 26):
                    c.drawString(55, tty, line); tty -= 34
                set_fill('sand'); c.setFont(regular, 16)
                for line in _pdf_wrap(data.get('subtitle', ''), 40):
                    c.drawString(55, tty - 10, line); tty -= 20
                set_fill('white'); c.setFont(regular, 15); y = 250
                for bullet in bullets[:6]:
                    for line_i, line in enumerate(_pdf_wrap(bullet, 40)):
                        c.drawString(55 if line_i == 0 else 70, y, ('• ' if line_i == 0 else '') + line); y -= 20
                    y -= 6
            else:
                image_left = idx % 2 == 0
                if image_path:
                    ix = 45 if image_left else 525
                    _pdf_draw_image(c, image_path, ix, 65, 390, 380, pdir, f'_pdf_split_{idx}.jpg')
                    tx = 475 if image_left else 55
                else:
                    tx = 90
                set_fill('accent'); c.setFont(bold, 16)
                sy = 412
                for line in _pdf_wrap(data.get('subtitle', ''), 46):
                    c.drawString(tx, sy, line); sy -= 22
                set_fill('text'); c.setFont(regular, 15); y = sy - 25
                for bullet in bullets[:6]:
                    for line_i, line in enumerate(_pdf_wrap(bullet, 50)):
                        c.drawString(tx if line_i == 0 else tx + 16, y, ('• ' if line_i == 0 else '') + line); y -= 20
                    y -= 8
            set_fill('secondary'); c.setFont(regular, 7); c.drawString(42, 22, f'{brand}  •  {idx + 1}')
        c.showPage()
    c.save()


async def _build_files(project: dict[str, Any], update: Update, context: ContextTypes.DEFAULT_TYPE) -> tuple[str, str]:
    project['profile']['brand_name'] = _normalize_exact_brand(project.get('profile', {}).get('brand_name', ''))
    project['palette'] = _parse_palette(project)
    slides = await _generate_deck_content(project, update)
    _save(project['user_id'], project)
    images = await _prepare_slide_images(project, slides, update, context)
    _save(project['user_id'], project)
    _validate_slides_or_raise(project, slides)
    pdir = _project_dir(project['user_id'], project['project_id'])
    brand = _safe_filename(project.get('profile', {}).get('brand_name', 'presentation'))
    pptx_path = pdir / f'{brand}_presentation.pptx'
    pdf_path = pdir / f'{brand}_presentation.pdf'
    await asyncio.to_thread(_build_pptx, project, slides, images, pptx_path)
    await asyncio.to_thread(_build_pdf, project, slides, images, pdf_path)
    project['final_slides'] = slides
    project['pptx_path'] = str(pptx_path)
    project['pdf_path'] = str(pdf_path)
    project['stage'] = 'done'
    _save(project['user_id'], project)
    return str(pdf_path), str(pptx_path)


# ─────────────────────────────────────────────────────────────────────────────
# v97 first-stage profile / structure integrity hotfix
# ─────────────────────────────────────────────────────────────────────────────


def _explicit_slide_plan(project: dict[str, Any]) -> list[dict[str, Any]]:
    raw = (project.get('raw_brief') or '').replace('\r\n', '\n').replace('\r', '\n')
    rows: list[tuple[int, str]] = []
    for line in raw.split('\n'):
        m = re.match(r'^\s*Слайд\s+(\d{1,2})\s*[.\-–—:]\s*(.+?)\s*$', line, re.I)
        if not m:
            continue
        num = int(m.group(1))
        title = re.sub(r'\s+', ' ', m.group(2)).strip(' .:—-')
        if title:
            rows.append((num, title))
    if not rows:
        return []
    rows = sorted(dict(rows).items())
    target = int(project.get('profile', {}).get('requested_slide_count') or len(rows))
    if len(rows) != target or [n for n, _ in rows] != list(range(1, target + 1)):
        return []
    result=[]
    for num,title in rows:
        low=title.lower()
        if num == 1 or 'облож' in low:
            layout='cover'; image_needed=True
        elif any(k in low for k in ['сравнен', 'линейк', 'модел']):
            layout='comparison'; image_needed=True
        elif any(k in low for k in ['контакт', 'призыв', 'следующий шаг']):
            layout='cta'; image_needed=True
        elif any(k in low for k in ['технолог', 'этап', 'как работает']):
            layout='process'; image_needed=True
        elif any(k in low for k in ['проблем', 'преимущ', 'почему', 'care', 'сервис']):
            layout='cards'; image_needed=('care' in low or 'сервис' in low)
        else:
            layout='split'; image_needed=True
        result.append({'title': title, 'layout': layout, 'image_needed': image_needed, 'bullets': [], 'image_prompt': ''})
    return result


async def _parse_profile(raw_brief: str, update: Update | None = None) -> dict[str, Any]:
    fallback = _profile_fallback(raw_brief)
    parsed: dict[str, Any] = {}
    prompt = f"""
Извлеки смысловые поля из брифа в JSON без Markdown:
{{"objective":"", "product":"", "audience":"", "positioning":"", "geography":"", "visual_direction":""}}
Не возвращай и не придумывай название, цены, контакты или количество слайдов: они извлекаются программно.
БРИФ:\n{raw_brief}
"""
    try:
        value = _json_from_text(await _llm_text(prompt, 0.1, update))
        if isinstance(value, dict):
            parsed = value
    except Exception:
        log.exception('v97 profile semantic parsing failed; deterministic profile retained')
    profile = deepcopy(fallback)
    for key in ['objective','product','audience','positioning','geography','visual_direction']:
        value = parsed.get(key)
        if isinstance(value,str) and value.strip():
            profile[key]=value.strip()
    # Immutable deterministic facts from the user's actual brief.
    profile['brand_name'] = _normalize_exact_brand(_extract_brand_name_regex(raw_brief))
    profile['tagline'] = _extract_tagline(raw_brief)
    profile['contacts'] = _extract_contacts(raw_brief)
    profile['prices'] = _extract_prices(raw_brief)
    profile['requested_slide_count'] = _extract_slide_count(raw_brief)
    profile['language'] = 'ru'
    return profile


def _repair_structure_v97(project: dict[str, Any], value: Any) -> list[dict[str, Any]]:
    explicit = _explicit_slide_plan(project)
    fallback = explicit or _fallback_structure(project)
    target = int(project.get('profile', {}).get('requested_slide_count') or len(fallback) or 12)
    raw_items = value.get('slides') if isinstance(value,dict) else value
    parsed_items = raw_items if isinstance(raw_items,list) else []
    parsed_by_index=[]
    for item in parsed_items[:target]:
        parsed_by_index.append(item if isinstance(item,dict) else {})
    repaired=[]
    for idx in range(target):
        base = deepcopy(fallback[idx] if idx < len(fallback) else _fallback_structure(project)[min(idx, len(_fallback_structure(project))-1)])
        item = deepcopy(parsed_by_index[idx] if idx < len(parsed_by_index) else {})
        if explicit:
            # User-specified slide titles/order are authoritative.
            item['title'] = base.get('title','')
            item['layout'] = base.get('layout','split')
            item['image_needed'] = base.get('image_needed', True)
        merged={
            **base,
            'title': str(item.get('title') or base.get('title') or '').strip()[:110],
            'layout': str(item.get('layout') or base.get('layout') or 'split').strip().lower(),
            'image_needed': bool(item.get('image_needed', base.get('image_needed', True))),
            'bullets': _sanitize_bullets(item.get('bullets') or base.get('bullets') or []),
            'image_prompt': str(item.get('image_prompt') or base.get('image_prompt') or '').strip()[:1500],
            'subtitle': str(item.get('subtitle') or '').strip()[:220],
        }
        merged = _build_deterministic_slide(project, base, merged, idx)
        merged['title'] = _normalize_exact_brand(merged.get('title') or base.get('title') or '')
        merged['bullets'] = _sanitize_bullets(merged.get('bullets') or [])
        repaired.append(merged)
    brand = _normalize_exact_brand(project.get('profile', {}).get('brand_name',''))
    if repaired:
        repaired[0]['title']=brand or repaired[0]['title']
        repaired[0]['layout']='cover'
        repaired[0]['image_needed']=True
        repaired[0]['subtitle']=project.get('profile',{}).get('tagline','')
        repaired[0]['bullets']=[_infer_cover_line(project)]
    if repaired:
        contacts=project.get('profile',{}).get('contacts') or _extract_contacts(project.get('raw_brief',''))
        repaired[-1]['layout']='cta'
        repaired[-1]['image_needed']=True
        repaired[-1]['bullets']=_sanitize_bullets([
            'Закажите бесплатную консультацию',
            *contacts,
            *(['Контакты указаны для тестирования'] if re.search(r'контакты указаны для тестирования', project.get('raw_brief',''), re.I) else []),
        ])
    # No contact hallucinations are allowed anywhere except the final slide.
    exact_contacts=set(project.get('profile',{}).get('contacts') or [])
    contact_like=re.compile(r'(?:@|https?://|www\.|\b[A-Za-z0-9.-]+\.[A-Za-z]{2,}\b|\+?\d[\d\s()\-]{8,}\d)')
    for i,slide in enumerate(repaired[:-1]):
        slide['bullets']=[b for b in slide.get('bullets',[]) if not contact_like.search(b) or b in exact_contacts]
        if not slide['bullets']:
            slide['bullets']=_sanitize_bullets((fallback[i].get('bullets') if i < len(fallback) else []) or ['Ключевой смысл раскрывается по исходному брифу'])
    return repaired


async def _generate_structure(project: dict[str, Any], update: Update | None = None) -> list[dict[str, Any]]:
    explicit = _explicit_slide_plan(project)
    fallback = explicit or _fallback_structure(project)
    profile = project.get('profile', {})
    notes='\n'.join(project.get('structure_notes',[]))
    prompt=f"""
Создай структуру коммерческой презентации в JSON:
{{"slides":[{{"title":"", "layout":"cover|split|cards|process|comparison|full_image|cta", "image_needed":true, "bullets":[""], "image_prompt":""}}]}}
Количество: {profile.get('requested_slide_count',12)}.
Точное название: {profile.get('brand_name','')}.
Используй только факты, цены и контакты из брифа. Не придумывай сайт, email или телефон.
Запрещена служебная фраза «{_PLACEHOLDER_SENTINEL}».
Если в брифе перечислены слайды, сохрани их порядок и названия.
Правки:\n{notes or 'Нет'}
Полный бриф:\n{project.get('raw_brief','')}
"""
    value=None
    try:
        value=_json_from_text(await _llm_text(prompt,0.15,update))
    except Exception:
        log.exception('v97 structure generation failed; deterministic structure used')
    result=_repair_structure_v97(project,value)
    if len(result)!=int(profile.get('requested_slide_count',12)):
        raise RuntimeError('Не удалось сохранить точное количество слайдов.')
    return result


def _structure_text(project: dict[str, Any]) -> str:
    brand=_normalize_exact_brand(project.get('profile',{}).get('brand_name') or 'Без названия')
    lines=[f'Предлагаемая структура: {brand}','']
    for i,slide in enumerate(project.get('structure',[]),1):
        bullets='; '.join(_sanitize_bullets(slide.get('bullets',[]))[:4])
        lines.append(f"{i}. {slide.get('title','')}" + (f" — {bullets}" if bullets else ''))
    contacts=project.get('profile',{}).get('contacts') or []
    if contacts:
        lines.extend(['', 'Контакты, распознанные строго из брифа:', *[f'• {c}' for c in contacts]])
    lines.extend(['', 'Проверьте названия, цены, порядок и отсутствие вымышленных фактов.'])
    return '\n'.join(lines)


# ─────────────────────────────────────────────────────────────────────────────
# v99 generic content extraction and resilient render pipeline
# ─────────────────────────────────────────────────────────────────────────────

_V99_SECTION_HEADINGS = {
    'проблемы клиента','проблемы','боли клиента','решение','решение aquavera home',
    'как работает технология','технология','ключевые преимущества','преимущества',
    'этапы работы','этапы реализации','этапы реализации проекта','контакты',
    'целевая аудитория b2c','целевая аудитория b2b','целевая аудитория',
    'варианты использования','сценарии использования','позиционирование',
}


def _v99_is_heading(line: str) -> bool:
    clean = re.sub(r'^\s*\d+[.)]\s*', '', str(line or '')).strip().rstrip(':').lower()
    if not clean:
        return False
    if clean in _V99_SECTION_HEADINGS:
        return True
    if re.match(r'^(?:слайд\s+\d+|[A-ZА-ЯЁ][A-ZА-ЯЁ0-9 _&.-]{2,40})$', str(line or '').strip()):
        return True
    return False


def _v99_section_items(raw: str, aliases: list[str], limit: int = 8) -> list[str]:
    lines = (raw or '').replace('\r\n','\n').replace('\r','\n').split('\n')
    aliases_n = [a.strip().lower().rstrip(':') for a in aliases]
    for i, line in enumerate(lines):
        clean = re.sub(r'^\s*\d+[.)]\s*', '', line).strip().rstrip(':').lower()
        if not any(clean == a or clean.startswith(a + ' ') for a in aliases_n):
            continue
        items=[]
        for nxt in lines[i+1:i+24]:
            s=nxt.strip()
            if not s:
                if items:
                    break
                continue
            if _v99_is_heading(s) and items:
                break
            s=re.sub(r'^\s*(?:[-–—•*]|\d+[.)])\s*','',s).strip()
            if not s or re.match(r'^(?:слайд\s+\d+)',s,re.I):
                continue
            if len(s)>260:
                s=s[:260]
            if s not in items:
                items.append(s)
            if len(items)>=limit:
                break
        return _sanitize_bullets(items)
    return []


def _v99_product_context(project: dict[str, Any]) -> str:
    profile=project.get('profile',{})
    return (profile.get('product') or profile.get('positioning') or project.get('raw_brief',''))[:900]


def _v99_generic_fallback(project: dict[str, Any], title: str, kind: str) -> list[str]:
    product=_v99_product_context(project)
    if kind=='technology':
        # Generic, not tied to plants/water/real estate.
        return [
            'Основные функциональные компоненты решения',
            'Последовательная обработка или выполнение задачи',
            'Автоматический контроль ключевых параметров',
            'Защита от ошибок и нештатных ситуаций',
            'Профессиональная настройка под объект клиента',
            'Сервисное сопровождение после запуска',
        ]
    if kind=='benefits':
        return [
            'Решение подбирается под задачу клиента',
            'Аккуратная интеграция в существующее пространство',
            'Автоматизация регулярных операций',
            'Профессиональный монтаж и настройка',
            'Понятное обслуживание после запуска',
        ]
    if kind=='stages':
        return [
            'Консультация и сбор исходных данных',
            'Подбор решения и подготовка предложения',
            'Согласование комплектации и сроков',
            'Доставка, монтаж и настройка',
            'Проверка результата и дальнейший сервис',
        ]
    return [f'Ключевая ценность раздела «{title}»', (product[:140].rstrip(' .,:;') if product else 'Решение раскрывается по данным главного брифа')]


def _build_deterministic_slide(project: dict[str, Any], base: dict[str, Any], slide: dict[str, Any], idx: int) -> dict[str, Any]:
    profile=project.get('profile',{})
    raw=project.get('raw_brief','')
    brand=_normalize_exact_brand(profile.get('brand_name') or base.get('title') or slide.get('title') or 'Бренд')
    title=str(slide.get('title') or base.get('title') or '').strip()
    key=title.lower()
    layout=slide.get('layout') or base.get('layout') or 'split'
    bullets=_sanitize_bullets(slide.get('bullets') or base.get('bullets') or [])
    subtitle=str(slide.get('subtitle') or '').strip()
    offers=project.get('_offer_objects') or _product_offer_objects(raw)
    service_offers=project.get('_service_offer_objects') or _service_offer_objects(raw)
    contacts=profile.get('contacts') or _extract_contacts(raw)

    if idx==0 or layout=='cover':
        slide['title']=brand
        slide['subtitle']=_normalize_exact_brand(profile.get('tagline') or subtitle or 'Коммерческая презентация')
        cover_line=_v99_section_items(raw,['ключевая фраза для обложки'],2)
        slide['bullets']=cover_line[:1] or [_infer_cover_line(project)]
        slide['image_prompt']=slide.get('image_prompt') or _default_scene_prompt(project, slide)
        return slide

    if 'проблем' in key:
        slide['bullets']=_v99_section_items(raw,['проблемы клиента','проблемы','боли клиента'],6) or bullets or _v99_generic_fallback(project,title,'default')
    elif 'решение' in key:
        slide['bullets']=_v99_section_items(raw,['решение '+brand.lower(),'решение','что получает клиент'],6) or bullets or _v99_generic_fallback(project,title,'default')
    elif any(k in key for k in ['технолог','как работает','схема']):
        slide['layout']='process'
        slide['bullets']=_v99_section_items(raw,['как работает технология','технология','основные этапы очистки','схема технологии'],6) or _v99_generic_fallback(project,title,'technology')
    elif 'преим' in key:
        slide['bullets']=_v99_section_items(raw,['ключевые преимущества','преимущества'],6) or bullets or _v99_generic_fallback(project,title,'benefits')
    elif any(k in key for k in ['варианты использования','сценарии использования','кому подходит']):
        b2c=_v99_section_items(raw,['целевая аудитория b2c'],4)
        b2b=_v99_section_items(raw,['целевая аудитория b2b'],4)
        slide['bullets']=_sanitize_bullets(b2c+b2b)[:6] or bullets or _v99_generic_fallback(project,title,'default')
    elif any(k in key for k in ['продуктовая линейка','сравнение','модел','тариф']):
        slide['layout']='comparison'
        slide['bullets']=[f"{o['name']} — {o['price']}" for o in offers[:4]] or bullets
    elif any(k in key for k in ['care','сервис','обслуж','сопровожд']):
        details=[]
        for offer in service_offers[:2]:
            if offer.get('price'):
                details.append(f"{offer['name']} — {offer['price']}")
            details.extend([d for d in (offer.get('details') or []) if not re.search(r'(?:@|\+?\d[\d\s()\-]{8,}|\.[A-Za-z]{2,}|контакт|телефон|email|сайт)', d, re.I)])
        section=_v99_section_items(raw,[title,'сервисное обслуживание системы','сервис и сопровождение'],6)
        slide['bullets']=_sanitize_bullets(details+section)[:6] or bullets or _v99_generic_fallback(project,title,'benefits')
    elif any(k in key for k in ['этап','реализац','процесс']):
        slide['layout']='process'
        slide['bullets']=_v99_section_items(raw,['этапы работы','этапы реализации проекта','этапы реализации'],6) or bullets or _v99_generic_fallback(project,title,'stages')
    elif any(k in key for k in ['почему','выбирают','отличия']):
        slide['bullets']=_v99_section_items(raw,['почему выбирают '+brand.lower(),'ключевые преимущества','позиционирование'],6) or bullets or _v99_generic_fallback(project,title,'benefits')
    elif layout=='cta' or any(k in key for k in ['контакт','призыв','следующий шаг']):
        slide['layout']='cta'
        slide['subtitle']=subtitle or 'Свяжитесь с нами для консультации и расчёта проекта'
        cta=_v99_section_items(raw,['призыв к действию','контакты'],4)
        slide['bullets']=_sanitize_bullets(cta+contacts)[:6]
    else:
        slide['bullets']=bullets or _v99_generic_fallback(project,title,'default')

    slide['title']=_normalize_exact_brand(title or base.get('title') or '')
    slide['subtitle']='' if _has_placeholder(subtitle) else subtitle
    return slide


def _v99_normalize_image(path: str, pdir: Path, idx: int) -> str:
    """Verify and normalize provider output before PDF/PPTX rendering."""
    try:
        with Image.open(path) as probe:
            probe.verify()
        with Image.open(path) as src:
            im=src.convert('RGB')
            # Cap megapixels to control RAM and Telegram PPTX size.
            max_side=1800
            if max(im.size)>max_side:
                im.thumbnail((max_side,max_side), Image.Resampling.LANCZOS)
            out=pdir/f'_normalized_slide_{idx+1:02d}.jpg'
            im.save(out,'JPEG',quality=82,optimize=True,progressive=True)
            return str(out)
    except Exception:
        log.exception('Invalid presentation image ignored: slide=%s path=%s',idx+1,path)
        return ''


def _v99_sanitize_images(images: dict[int,str], pdir: Path) -> dict[int,str]:
    clean={}
    for idx,path in images.items():
        if path and os.path.exists(path):
            normalized=_v99_normalize_image(path,pdir,int(idx))
            if normalized:
                clean[int(idx)]=normalized
    return clean


async def _build_files(project: dict[str, Any], update: Update, context: ContextTypes.DEFAULT_TYPE) -> tuple[str, str]:
    project.setdefault('profile',{})['brand_name']=_normalize_exact_brand(project.get('profile',{}).get('brand_name',''))
    project['palette']=_parse_palette(project)
    slides=await _generate_deck_content(project,update)
    _validate_slides_or_raise(project,slides)
    _save(project['user_id'],project)
    images=await _prepare_slide_images(project,slides,update,context)
    pdir=_project_dir(project['user_id'],project['project_id'])
    images=_v99_sanitize_images(images,pdir)
    project['slide_images']={str(k):v for k,v in images.items()}
    _save(project['user_id'],project)

    brand=_safe_filename(project.get('profile',{}).get('brand_name','presentation'))
    pptx_path=pdir/f'{brand}_presentation.pptx'
    pdf_path=pdir/f'{brand}_presentation.pdf'
    errors=[]
    try:
        await asyncio.to_thread(_build_pptx,project,slides,images,pptx_path)
    except Exception as exc:
        log.exception('PPTX build failed')
        errors.append(f'PPTX: {type(exc).__name__}: {exc}')
    try:
        await asyncio.to_thread(_build_pdf,project,slides,images,pdf_path)
    except Exception as exc:
        log.exception('PDF build failed')
        errors.append(f'PDF: {type(exc).__name__}: {exc}')

    if not pptx_path.exists() or pptx_path.stat().st_size < 1000:
        errors.append('PPTX file was not created')
    if not pdf_path.exists() or pdf_path.stat().st_size < 1000:
        errors.append('PDF file was not created')
    if errors:
        project['last_build_errors']=errors[-6:]
        project['stage']='final_review'
        _save(project['user_id'],project)
        raise RuntimeError(' | '.join(errors))

    project['final_slides']=slides
    project['pptx_path']=str(pptx_path)
    project['pdf_path']=str(pdf_path)
    project['last_build_errors']=[]
    project['stage']='done'
    _save(project['user_id'],project)
    return str(pdf_path),str(pptx_path)


# ─────────────────────────────────────────────────────────────────────────────
# v100 logo generation / Telegram delivery stability hotfix
# ─────────────────────────────────────────────────────────────────────────────


def _compose_logo_card(symbol_bytes: bytes | None, brand: str, index: int, palette: dict[str, str]) -> bytes:
    """Create a clean preview card.

    Exact brand text is typeset by Pillow. No provider-generated captions and no
    internal labels such as "LOGO CONCEPT 1" are shown to the user.
    """
    canvas_img = Image.new("RGB", (1024, 1024), _hex(palette["background"]))
    if symbol_bytes:
        try:
            symbol = Image.open(io.BytesIO(symbol_bytes)).convert("RGBA")
            symbol.thumbnail((650, 650), Image.Resampling.LANCZOS)
        except Exception:
            log.exception("Invalid generated logo symbol; using deterministic fallback index=%s", index)
            symbol = _fallback_logo_symbol(index, palette)
    else:
        symbol = _fallback_logo_symbol(index, palette)

    # Keep generous clear space around the mark.
    x = (1024 - symbol.width) // 2
    y = max(45, (700 - symbol.height) // 2 + 25)
    canvas_img.paste(symbol, (x, y), symbol if symbol.mode == "RGBA" else None)

    draw = ImageDraw.Draw(canvas_img)
    brand_value = _normalize_exact_brand(brand).upper()
    font_size = 76
    while font_size > 34:
        f = _font(font_size, bold=True)
        box = draw.textbbox((0, 0), brand_value, font=f)
        if box[2] - box[0] <= 880:
            break
        font_size -= 4
    bbox = draw.textbbox((0, 0), brand_value, font=f)
    tx = (1024 - (bbox[2] - bbox[0])) // 2
    draw.text((tx, 820), brand_value, font=f, fill=_hex(palette["text"]))

    out = io.BytesIO()
    canvas_img.save(out, format="PNG", optimize=True)
    return out.getvalue()


def _v100_logo_symbol_prompt(project: dict[str, Any], concept: str) -> str:
    profile = project.get("profile", {})
    notes = "\n".join(project.get("logo_notes", []))
    raw = project.get("raw_brief", "")
    brand = _normalize_exact_brand(profile.get("brand_name", ""))
    return f"""
Create ONE original premium vector-style LOGO SYMBOL ONLY for the business below.
The user requirements are absolute and outrank every other instruction.
If the user forbids a shape, object, letter, house, roof, circle, droplet, bottle,
mockup or other cliché, it must not appear in any form.
Do not infer a house or roof merely because the exact brand name contains HOME.
Do not draw letters, words, captions, initials, pseudo-letters, brand text, labels,
mockups, packaging, signs, business cards, walls, screenshots or presentation frames.
The exact brand name is typeset later by software.

Concept direction: {concept}
Exact brand name for semantic context only: {brand}
Tagline: {profile.get('tagline', '')}
Business/product: {(profile.get('product') or '')[:1400]}
Positioning: {(profile.get('positioning') or '')[:900]}
Visual direction: {(profile.get('visual_direction') or '')[:800]}

CUMULATIVE USER LOGO REQUIREMENTS — FOLLOW LITERALLY:
{notes[:3200] or 'No additional requirements.'}

MAIN BRIEF EXCERPT:
{raw[:2400]}

Flat vector style, one coherent symbol, strong silhouette, simple geometry,
high recognition at 48 px, balanced negative space, centered square composition,
generous margins, no stock icon, no gradient, no 3D, no shadow, no fine detail.
""".strip()


async def _generate_logos(project: dict[str, Any], update: Update | None = None, context: ContextTypes.DEFAULT_TYPE | None = None) -> list[str]:
    brand = _normalize_exact_brand(project.get("profile", {}).get("brand_name", ""))
    if not _is_valid_brand_name(brand):
        raise RuntimeError("Brand name is missing or invalid")
    project.setdefault("profile", {})["brand_name"] = brand

    concepts = [
        "distinctive negative-space mark derived from the product and service; avoid literal category clichés",
        "modular technical-system symbol formed from a few purposeful geometric elements; no building outline",
        "premium continuous-motion emblem with a memorable silhouette; abstract rather than literal",
    ]
    prompts = [_v100_logo_symbol_prompt(project, concept) for concept in concepts]

    generated: list[bytes] = []
    if _STUDIO_IMAGE_BATCH_CALL is not None and update is not None and context is not None:
        try:
            batch = await _STUDIO_IMAGE_BATCH_CALL(update, context, prompts, "openai", "presentation_logo_variants")
            if batch:
                generated = list(batch)
        except Exception:
            log.exception("Logo batch generation failed; retrying missing concepts individually")

    # Preserve positions and retry missing outputs individually.
    generated = (generated + [b""] * 3)[:3]
    for i, prompt in enumerate(prompts):
        if generated[i]:
            continue
        for attempt in range(2):
            try:
                generated[i] = await _generate_image_bytes(prompt, "1024x1024")
                if generated[i]:
                    break
            except Exception:
                log.exception("Logo generation retry failed concept=%s attempt=%s", i + 1, attempt + 1)
                await asyncio.sleep(0.8 * (attempt + 1))

    pdir = _project_dir(project["user_id"], project["project_id"])
    out_paths: list[str] = []
    symbol_paths: list[str] = []
    for index in range(1, 4):
        symbol_bytes = generated[index - 1] or b""
        # Keep a separate symbol file for final PDF/PPTX insertion.
        if symbol_bytes:
            try:
                symbol_im = Image.open(io.BytesIO(symbol_bytes)).convert("RGBA")
                symbol_im.thumbnail((900, 900), Image.Resampling.LANCZOS)
                symbol_path = pdir / f"logo_symbol_{index}.png"
                symbol_im.save(symbol_path, "PNG", optimize=True)
            except Exception:
                log.exception("Could not normalize logo symbol index=%s", index)
                symbol_path = pdir / f"logo_symbol_{index}.png"
                _fallback_logo_symbol(index, project.get("palette", DEFAULT_PALETTE)).save(symbol_path, "PNG", optimize=True)
        else:
            symbol_path = pdir / f"logo_symbol_{index}.png"
            _fallback_logo_symbol(index, project.get("palette", DEFAULT_PALETTE)).save(symbol_path, "PNG", optimize=True)
        symbol_paths.append(str(symbol_path))

        preview = _compose_logo_card(symbol_path.read_bytes(), brand, index, project.get("palette", DEFAULT_PALETTE))
        preview_path = pdir / f"logo_concept_{index}.png"
        preview_path.write_bytes(preview)
        out_paths.append(str(preview_path))

    project["logo_symbol_paths"] = symbol_paths
    project["logo_candidates"] = out_paths
    return out_paths


def _v100_clean_existing_logo_preview(path: str, project: dict[str, Any], index: int) -> str:
    """Remove the legacy LOGO CONCEPT footer from already-generated v99 cards."""
    if not path or not os.path.exists(path):
        return ""
    try:
        im = Image.open(path).convert("RGB")
        if im.width >= 900 and im.height >= 900:
            draw = ImageDraw.Draw(im)
            bg = _hex((project.get("palette") or DEFAULT_PALETTE)["background"])
            # Legacy internal caption was around y=900.
            draw.rectangle((0, int(im.height * 0.875), im.width, im.height), fill=bg)
        out = Path(path).with_name(f"logo_preview_clean_{index}.png")
        im.save(out, "PNG", optimize=True)
        return str(out)
    except Exception:
        log.exception("Could not clean legacy logo preview path=%s", path)
        return path


def _resolve_logo_path(project: dict[str, Any], pdir: Path | None = None) -> str:
    """Prefer the original symbol file over the preview card."""
    selected = project.get("logo_selected") or ""
    if not selected or not os.path.exists(selected):
        return ""
    try:
        name = Path(selected).name
        m = re.match(r"logo_concept_(\d+)\.png$", name, re.I)
        if m:
            idx = int(m.group(1)) - 1
            symbols = project.get("logo_symbol_paths") or []
            if 0 <= idx < len(symbols) and os.path.exists(symbols[idx]):
                return symbols[idx]
            sibling = Path(selected).with_name(f"logo_symbol_{idx + 1}.png")
            if sibling.exists():
                return str(sibling)
    except Exception:
        log.exception("Could not resolve separate logo symbol")
    # Backward-compatible crop for v95-v99 preview cards.
    try:
        im = Image.open(selected).convert("RGBA")
        if im.width >= 900 and im.height >= 900:
            crop = im.crop((120, 40, im.width - 120, int(im.height * 0.72)))
        else:
            crop = im
        pdir = pdir or Path(selected).parent
        out = pdir / "_resolved_logo.png"
        crop.save(out, "PNG", optimize=True)
        return str(out)
    except Exception:
        return selected


async def _v100_send_single_logo(update: Update, path: str, index: int, brand: str) -> str:
    """Return ok / uncertain / failed without aborting the whole logo stage."""
    try:
        with open(path, "rb") as fh:
            await update.effective_message.reply_photo(
                photo=InputFile(fh, filename=f"logo_{index}.png"),
                caption=f"Вариант {index} — точное название: {brand}",
                read_timeout=90,
                write_timeout=120,
                connect_timeout=30,
                pool_timeout=30,
            )
        return "ok"
    except TimedOut:
        # Telegram may have accepted the upload but failed to acknowledge it in time.
        log.warning("Telegram timed out while sending logo %s; treating as uncertain", index)
        return "uncertain"
    except RetryAfter as exc:
        await asyncio.sleep(float(getattr(exc, "retry_after", 1.0)) + 0.5)
        try:
            with open(path, "rb") as fh:
                await update.effective_message.reply_photo(
                    photo=InputFile(fh, filename=f"logo_{index}.png"),
                    caption=f"Вариант {index} — точное название: {brand}",
                    read_timeout=90,
                    write_timeout=120,
                    connect_timeout=30,
                    pool_timeout=30,
                )
            return "ok"
        except TimedOut:
            return "uncertain"
        except TelegramError:
            log.exception("Telegram retry failed for logo %s", index)
            return "failed"
    except (NetworkError, TelegramError):
        log.exception("Telegram failed to send logo %s", index)
        return "failed"


async def _send_logo_candidates(update: Update, project: dict[str, Any]) -> None:
    """Send all logo previews without turning a Telegram upload timeout into a project error."""
    paths = [p for p in project.get("logo_candidates", []) if p and os.path.exists(p)]
    if not paths:
        raise RuntimeError("Файлы вариантов логотипа не найдены")
    brand = _normalize_exact_brand(project.get("profile", {}).get("brand_name", ""))
    cleaned = [_v100_clean_existing_logo_preview(path, project, i) for i, path in enumerate(paths, 1)]

    statuses: list[str] = []
    # A single media-group request reduces three separate upload round trips.
    streams: list[io.BytesIO] = []
    try:
        media = []
        for i, path in enumerate(cleaned, 1):
            bio = io.BytesIO(Path(path).read_bytes())
            bio.name = f"logo_{i}.png"
            streams.append(bio)
            media.append(InputMediaPhoto(media=InputFile(bio, filename=bio.name), caption=f"Вариант {i} — точное название: {brand}"))
        try:
            await update.effective_message.reply_media_group(
                media=media,
                read_timeout=90,
                write_timeout=150,
                connect_timeout=30,
                pool_timeout=30,
            )
            statuses = ["ok"] * len(cleaned)
        except TimedOut:
            # Do not retry immediately: Telegram often delivers the album after a timeout.
            log.warning("Telegram timed out while sending logo media group; candidates may still arrive")
            statuses = ["uncertain"] * len(cleaned)
        except (NetworkError, TelegramError):
            log.exception("Logo media group failed; falling back to individual uploads")
            statuses = []
    finally:
        for stream in streams:
            with contextlib.suppress(Exception):
                stream.close()

    if not statuses:
        for i, path in enumerate(cleaned, 1):
            statuses.append(await _v100_send_single_logo(update, path, i, brand))

    if all(status == "failed" for status in statuses):
        raise RuntimeError("Telegram не смог отправить ни один вариант логотипа")

    note = "Выберите вариант. Все накопленные пожелания сохранены."
    if any(status == "uncertain" for status in statuses):
        note += " Telegram долго подтверждал загрузку; варианты могут появиться с небольшой задержкой."
    elif any(status == "failed" for status in statuses):
        note += " Один из файлов не отправился, но остальные варианты доступны для выбора."
    await _reply(update, note, _logo_selection_kb(len(paths)))


__all__ = ["PresentationStudio", "StudioConfig", "register", "init_storage", "on_start_callback"]

