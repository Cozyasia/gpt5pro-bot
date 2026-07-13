import hashlib
import hmac
# -*- coding: utf-8 -*-
import os
import re
import json
import time
import types
import base64
import logging
from io import BytesIO
import asyncio
import sqlite3
from datetime import datetime, timedelta, timezone
from zoneinfo import ZoneInfo
import threading
import uuid
import urllib.parse
import shutil
import sys
import tempfile
import subprocess
import contextlib
import random

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ Render Secret Files bootstrap в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
# Must run before any API keys are read from os.environ.
from secret_loader import bootstrap_secret_environment, get_secret
_SECRET_FILE_SOURCES = bootstrap_secret_environment()

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ Rembg cache bootstrap в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
# Render Starter/СҖСғСҮРҪРҫР№ deploy СҮР°СҒСӮРҫ РҪРө РёРјРөРөСӮ writable /data.
# Rembg/pooch РҙРҫР»Р¶РөРҪ РҝРҫР»СғСҮРёСӮСҢ writable U2NET_HOME Р”Рһ РёРјРҝРҫСҖСӮР° rembg, РёРҪР°СҮРө РјРҫРҙРөР»СҢ РҪРө СҒРәР°СҮРёРІР°РөСӮСҒСҸ.
def _ensure_writable_dir_env(var_name: str, preferred_default: str, fallback: str) -> str:
    raw = (os.environ.get(var_name) or "").strip()
    candidates = []
    for c in (raw, preferred_default, fallback):
        if c and c not in candidates:
            candidates.append(c)
    last = fallback
    for d in candidates:
        try:
            os.makedirs(d, exist_ok=True)
            probe = os.path.join(d, ".write_test")
            with open(probe, "w", encoding="utf-8") as f:
                f.write("ok")
            with contextlib.suppress(Exception):
                os.remove(probe)
            os.environ[var_name] = d
            return d
        except Exception:
            last = d
            continue
    # РҹРҫСҒР»РөРҙРҪРёР№ СҖРөР·РөСҖРІ вҖ” /tmp, РҙР°Р¶Рө РөСҒР»Рё РҝСҖРҫРІРөСҖРәР° РІСӢСҲРө РҪРөРҫР¶РёРҙР°РҪРҪРҫ РҪРө РҝСҖРҫСҲР»Р°.
    os.environ[var_name] = fallback
    with contextlib.suppress(Exception):
        os.makedirs(fallback, exist_ok=True)
    return fallback

U2NET_HOME = _ensure_writable_dir_env("U2NET_HOME", "/tmp/.u2net", "/tmp/.u2net")
XDG_CACHE_HOME = _ensure_writable_dir_env("XDG_CACHE_HOME", "/tmp/.cache", "/tmp/.cache")
os.environ.setdefault("MPLCONFIGDIR", "/tmp/.matplotlib")
with contextlib.suppress(Exception):
    os.makedirs(os.environ.get("MPLCONFIGDIR", "/tmp/.matplotlib"), exist_ok=True)

from http.server import HTTPServer, BaseHTTPRequestHandler

import httpx
from runway_official import (
    RunwayOfficialClient, RunwayAPIError, RunwayTaskTimeout,
    key_format_hint as runway_key_format_hint,
    safe_key_fingerprint as runway_safe_key_fingerprint,
)
from telegram import (
    Update, ReplyKeyboardMarkup, KeyboardButton, WebAppInfo, InputFile,
    LabeledPrice, InlineKeyboardMarkup, InlineKeyboardButton
)
from telegram.ext import (
    Application, ApplicationBuilder, CommandHandler, MessageHandler, ContextTypes, filters,
    PreCheckoutQueryHandler, CallbackQueryHandler
)
from telegram.constants import ChatAction
from telegram.error import TelegramError, TimedOut, BadRequest
from presentation_studio import PresentationStudio, StudioConfig
# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ TTS imports в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
import contextlib  # СғР¶Рө Сғ СӮРөРұСҸ РІСӢСҲРө РөСҒСӮСҢ, РҙСғРұР»РёСҖРҫРІР°СӮСҢ РқР• РҪР°РҙРҫ, РөСҒР»Рё РёРјРҝРҫСҖСӮ СҒСӮРҫРёСӮ

# Optional PIL / rembg for photo tools
try:
    from PIL import Image, ImageFilter, ImageOps, ImageDraw, ImageFont
except Exception:
    Image = None
    ImageFilter = None
    ImageOps = None
    ImageDraw = None
try:
    from rembg import remove as rembg_remove, new_session as rembg_new_session
    REMBG_IMPORT_ERROR = ""
except Exception as _rembg_e:
    rembg_remove = None
    rembg_new_session = None
    REMBG_IMPORT_ERROR = repr(_rembg_e)

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ LOGGING в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s | %(levelname)s | %(name)s | %(message)s",
)
log = logging.getLogger("gpt-bot")

PATCH_VERSION = "v93-webapp-checkout-bridge-2026-07-14"

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ ENV в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ

def _env_float(name: str, default: float) -> float:
    """
    Р‘РөР·РҫРҝР°СҒРҪРҫРө СҮСӮРөРҪРёРө float РёР· ENV:
    - РҝРҫРҙРҙРөСҖР¶РёРІР°РөСӮ Рё '4,99', Рё '4.99'
    - РҝСҖРё РҫСҲРёРұРәРө РІРҫР·РІСҖР°СүР°РөСӮ default
    """
    raw = os.environ.get(name)
    if not raw:
        return float(default)
    raw = raw.replace(",", ".").strip()
    try:
        return float(raw)
    except Exception:
        return float(default)

# Commercial pricing guardrails.
# By default the bot uses audited canonical provider costs so stale Render ENV values
# cannot accidentally sell generations below cost. Set PRICING_CANONICAL_COSTS=0
# only when you intentionally want to manage every cost manually in Render.
PRICING_CANONICAL_COSTS = os.environ.get("PRICING_CANONICAL_COSTS", "1").strip().lower() not in ("0", "false", "no", "off")
PRICING_LOCK_1_CREDIT_1_RUB = os.environ.get("PRICING_LOCK_1_CREDIT_1_RUB", "1").strip().lower() not in ("0", "false", "no", "off")
PRICING_MIN_USD_RUB = max(1.0, _env_float("PRICING_MIN_USD_RUB", 100.0))
PRICING_MIN_MULTIPLIER = max(1.0, _env_float("PRICING_MIN_MULTIPLIER", 2.0))

def _pricing_cost(name: str, canonical: float) -> float:
    return float(canonical) if PRICING_CANONICAL_COSTS else _env_float(name, canonical)

BOT_TOKEN = (os.environ.get("BOT_TOKEN") or os.environ.get("TELEGRAM_BOT_TOKEN", "")).strip()
BOT_USERNAME     = os.environ.get("BOT_USERNAME", "").strip().lstrip("@")
PUBLIC_URL       = os.environ.get("PUBLIC_URL", "").strip()
WEBAPP_URL       = os.environ.get("WEBAPP_URL", "").strip()

OPENAI_API_KEY   = os.environ.get("OPENAI_API_KEY", "").strip()
OPENAI_BASE_URL  = os.environ.get("OPENAI_BASE_URL", "").strip()        # OpenRouter РёР»Рё СҒРІРҫР№ РҝСҖРҫРәСҒРё РҙР»СҸ СӮРөРәСҒСӮР°
OPENAI_MODEL     = os.environ.get("OPENAI_MODEL", "openai/gpt-4o-mini").strip()

OPENROUTER_SITE_URL = os.environ.get("OPENROUTER_SITE_URL", "").strip()
OPENROUTER_APP_NAME = os.environ.get("OPENROUTER_APP_NAME", "").strip()

USE_WEBHOOK      = os.environ.get("USE_WEBHOOK", "1").lower() in ("1","true","yes","on")
WEBHOOK_PATH     = os.environ.get("WEBHOOK_PATH", "/tg").strip()
WEBHOOK_SECRET   = os.environ.get("TELEGRAM_WEBHOOK_SECRET", "").strip()

BANNER_URL       = os.environ.get("BANNER_URL", "").strip()
TAVILY_API_KEY   = os.environ.get("TAVILY_API_KEY", "").strip()

# Р’РҗР–РқРһ: РҝСҖРҫРІР°Р№РҙРөСҖ СӮРөРәСҒСӮР° (openai / openrouter Рё СӮ.Рҝ.)
TEXT_PROVIDER    = os.environ.get("TEXT_PROVIDER", "").strip()

# STT:
OPENAI_STT_KEY   = os.environ.get("OPENAI_STT_KEY", "").strip()
TRANSCRIBE_MODEL = os.environ.get("OPENAI_TRANSCRIBE_MODEL", "whisper-1").strip()
DEEPGRAM_API_KEY = os.environ.get("DEEPGRAM_API_KEY", "").strip()

# TTS:
OPENAI_TTS_KEY       = os.environ.get("OPENAI_TTS_KEY", "").strip() or OPENAI_API_KEY
OPENAI_TTS_BASE_URL  = (os.environ.get("OPENAI_TTS_BASE_URL", "").strip() or "https://api.openai.com/v1")
OPENAI_TTS_MODEL     = os.environ.get("OPENAI_TTS_MODEL", "gpt-4o-mini-tts").strip()
OPENAI_TTS_VOICE     = os.environ.get("OPENAI_TTS_VOICE", "alloy").strip()
# Р“РҫР»РҫСҒ РҝРҫ СғРјРҫР»СҮР°РҪРёСҺ РёРјРөРҪРҪРҫ РҙР»СҸ РіРҫРІРҫСҖСҸСүРөРіРҫ Р°РІР°СӮР°СҖР°. РҹРҫР»СҢР·РҫРІР°СӮРөР»СҢ РјРҫР¶РөСӮ РІСӢРұСҖР°СӮСҢ РіРҫР»РҫСҒ РәРҪРҫРҝРәРҫР№ РІ РјРөРҪСҺ Р°РІР°СӮР°СҖР°.
AVATAR_TTS_DEFAULT_VOICE = os.environ.get("AVATAR_TTS_DEFAULT_VOICE", OPENAI_TTS_VOICE or "nova").strip() or "nova"
AVATAR_TTS_VOICES = [v.strip() for v in os.environ.get("AVATAR_TTS_VOICES", "nova,alloy,onyx,shimmer,fable").split(",") if v.strip()]
TTS_MAX_CHARS        = int(os.environ.get("TTS_MAX_CHARS", "1000") or "1000")
# 0 = РҪРө РҙСғРұР»РёСҖРҫРІР°СӮСҢ РҫР·РІСғСҮРөРҪРҪСӢР№ СӮРөРәСҒСӮ РҝРҫРҙРҝРёСҒСҢСҺ Рә voice-СҒРҫРҫРұСүРөРҪРёСҺ.
# РҹРҫР»СҢР·РҫРІР°СӮРөР»СҢ РҝРҫР»СғСҮР°РөСӮ: СӮРөРәСҒСӮРҫРІСӢР№ РҫСӮРІРөСӮ + РҫСӮРҙРөР»СҢРҪРҫРө voice РұРөР· caption.
TTS_VOICE_CAPTION    = os.environ.get("TTS_VOICE_CAPTION", "0").strip().lower() in ("1", "true", "yes", "on")
# 0 = РҪРө РҫСӮРҝСҖР°РІР»СҸСӮСҢ РҫСӮРҙРөР»СҢРҪРҫРө СҒРҫРҫРұСүРөРҪРёРө В«Р Р°СҒРҝРҫР·РҪР°Р»: ...В» РҝРҫСҒР»Рө voice/STT.
STT_ECHO_TRANSCRIPT  = os.environ.get("STT_ECHO_TRANSCRIPT", "0").strip().lower() in ("1", "true", "yes", "on")

# РҹР°РјСҸСӮСҢ РҙРёР°Р»РҫРіР° РҙР»СҸ GPT-РҫСӮРІРөСӮРҫРІ: РәРҫСҖРҫСӮРәР°СҸ РёСҒСӮРҫСҖРёСҸ РҝРҫ user_id/chat_id.
CHAT_MEMORY_ENABLED      = os.environ.get("CHAT_MEMORY_ENABLED", "1").strip().lower() in ("1", "true", "yes", "on")
CHAT_MEMORY_MAX_MESSAGES = int(os.environ.get("CHAT_MEMORY_MAX_MESSAGES", "16") or "16")
CHAT_MEMORY_MAX_CHARS    = int(os.environ.get("CHAT_MEMORY_MAX_CHARS", "6000") or "6000")
CHAT_MEMORY_TTL_DAYS     = int(os.environ.get("CHAT_MEMORY_TTL_DAYS", "0") or "0")
# Р’РёСҖСӮСғР°Р»СҢРҪСӢРө РҝРҫР»СҢР·РҫРІР°СӮРөР»СҢСҒРәРёРө РҙРёР°Р»РҫРіРё РІРҪСғСӮСҖРё РҫРҙРҪРҫРіРҫ Telegram-СҮР°СӮР°.
CHAT_MAX_CONVERSATIONS      = max(1, min(4, int(os.environ.get("CHAT_MAX_CONVERSATIONS", "4") or "4")))
CHAT_HISTORY_PAGE_MESSAGES  = max(10, int(os.environ.get("CHAT_HISTORY_PAGE_MESSAGES", "40") or "40"))
CHAT_HISTORY_PAGE_SIZE      = max(4, min(12, int(os.environ.get("CHAT_HISTORY_PAGE_SIZE", "8") or "8")))

# Images:
OPENAI_IMAGE_KEY    = os.environ.get("OPENAI_IMAGE_KEY", "").strip() or OPENAI_API_KEY
IMAGES_BASE_URL     = (os.environ.get("OPENAI_IMAGE_BASE_URL", "").strip() or "https://api.openai.com/v1")
IMAGES_MODEL        = os.environ.get("OPENAI_IMAGE_MODEL", "gpt-image-1").strip() or "gpt-image-1"
OPENAI_IMAGE_QUALITY = os.environ.get("OPENAI_IMAGE_QUALITY", "medium").strip().lower() or "medium"

# Runway
# Priority: Render Environment -> Render Secret File -> legacy alias.
# Supported Secret Files include /etc/secrets/runway.env and, for compatibility, yookassa.env.
RUNWAY_API_KEY, RUNWAY_KEY_SOURCE = get_secret("RUNWAYML_API_SECRET", "RUNWAY_API_KEY", "RUNWAY_KEY")
# Official Runway Developer API is the primary production route. Comet remains only a fallback.
RUNWAY_DIRECT_ENABLED = os.environ.get("RUNWAY_DIRECT_ENABLED", "1").strip().lower() not in ("0", "false", "no", "off")
RUNWAY_MODEL        = os.environ.get("RUNWAY_MODEL", "gen4.5").strip() or "gen4.5"
RUNWAY_RATIO        = os.environ.get("RUNWAY_RATIO", "720:1280").strip()
RUNWAY_DURATION_S   = int(os.environ.get("RUNWAY_DURATION_S", "8") or 8)

# Luma
LUMA_API_KEY     = os.environ.get("LUMA_API_KEY", "").strip()
LUMA_MODEL       = os.environ.get("LUMA_MODEL", "ray-2").strip()
LUMA_ASPECT      = os.environ.get("LUMA_ASPECT", "16:9").strip()
LUMA_DURATION_S  = int((os.environ.get("LUMA_DURATION_S") or "5").strip() or 5)
LUMA_BASE_URL    = (os.environ.get("LUMA_BASE_URL", "https://api.lumalabs.ai/dream-machine/v1").strip().rstrip("/"))
LUMA_CREATE_PATH = "/generations"
LUMA_STATUS_PATH = "/generations/{id}"
# Luma Images (РҫРҝСҶРёРҫРҪР°Р»СҢРҪРҫ: РөСҒР»Рё РҪРөСӮ вҖ” РёСҒРҝРҫР»СҢР·СғРөРј OpenAI Images РәР°Рә С„РҫР»РұСҚРә)
LUMA_IMG_BASE_URL = os.environ.get("LUMA_IMG_BASE_URL", "").strip().rstrip("/")
LUMA_IMG_MODEL    = os.environ.get("LUMA_IMG_MODEL", "imagine-image-1").strip()

# РӨРҫР»РұСҚРәРё Luma
_fallbacks_raw = ",".join([
    os.environ.get("LUMA_FALLBACKS", ""),
    os.environ.get("LUMA_FALLBACK_BASE_URL", "")
])
LUMA_FALLBACKS = []
for u in re.split(r"[;,]\s*", _fallbacks_raw):
    if not u:
        continue
    u = u.strip().rstrip("/")
    if u and u != LUMA_BASE_URL and u not in LUMA_FALLBACKS:
        LUMA_FALLBACKS.append(u)

# Runway endpoints
RUNWAY_BASE_URL    = (os.environ.get("RUNWAY_BASE_URL", "https://api.dev.runwayml.com").strip().rstrip("/"))
RUNWAY_CREATE_PATH = "/v1/tasks"
RUNWAY_STATUS_PATH = "/v1/tasks/{id}"
RUNWAY_I2V_PATH    = os.environ.get("RUNWAY_I2V_PATH", "/v1/image_to_video").strip() or "/v1/image_to_video"
RUNWAY_TEXT_CREATE_PATH = os.environ.get("RUNWAY_TEXT_CREATE_PATH", "/v1/text_to_video").strip() or "/v1/text_to_video"
RUNWAY_TEXT_COMPAT_PATH = os.environ.get("RUNWAY_TEXT_COMPAT_PATH", "/v1/image_to_video").strip() or "/v1/image_to_video"
RUNWAY_UPLOAD_PATH = os.environ.get("RUNWAY_UPLOAD_PATH", "/v1/uploads").strip() or "/v1/uploads"
RUNWAY_ORGANIZATION_PATH = os.environ.get("RUNWAY_ORGANIZATION_PATH", "/v1/organization").strip() or "/v1/organization"
RUNWAY_API_VERSION = os.environ.get("RUNWAY_API_VERSION", "2024-11-06").strip()
RUNWAY_USE_COMET   = os.environ.get("RUNWAY_USE_COMET", "1").strip().lower() not in ("0", "false", "no", "off")

# CometAPI / Sora / Kling wrappers for imageвҶ’video
COMET_API_KEY  = (os.environ.get("COMET_API_KEY") or os.environ.get("COMETAPI_KEY") or "").strip()
COMET_BASE_URL = os.environ.get("COMET_BASE_URL", "https://api.cometapi.com").strip().rstrip("/")

# Optional Comet image generation fallback for business logos.
# If OpenAI Images/Luma are not configured, the bot tries this route and then a local PNG fallback.
COMET_IMAGE_GEN_MODEL = os.environ.get("COMET_IMAGE_GEN_MODEL", "gpt-image-1").strip() or "gpt-image-1"
COMET_IMAGE_GEN_PATH = os.environ.get("COMET_IMAGE_GEN_PATH", "/v1/images/generations").strip() or "/v1/images/generations"
COMET_IMAGE_GEN_TIMEOUT_S = int(os.environ.get("COMET_IMAGE_GEN_TIMEOUT_S", "180") or "180")
LOGO_LOCAL_FALLBACK = os.environ.get("LOGO_LOCAL_FALLBACK", "1").strip().lower() not in ("0", "false", "no", "off")

# Midjourney through CometAPI (official Comet wrapper flow: submit -> task fetch).
MIDJOURNEY_ENABLED = os.environ.get("MIDJOURNEY_ENABLED", "1").strip().lower() not in ("0", "false", "no", "off")
MIDJOURNEY_MODE = os.environ.get("MIDJOURNEY_MODE", "fast").strip().lower() or "fast"
if MIDJOURNEY_MODE not in ("relax", "fast", "turbo"):
    MIDJOURNEY_MODE = "fast"
_MJ_PREFIX = {"relax": "", "fast": "/mj-fast", "turbo": "/mj-turbo"}[MIDJOURNEY_MODE]
MIDJOURNEY_CREATE_PATH = os.environ.get("MIDJOURNEY_CREATE_PATH", f"{_MJ_PREFIX}/mj/submit/imagine").strip() or f"{_MJ_PREFIX}/mj/submit/imagine"
MIDJOURNEY_STATUS_PATH = os.environ.get("MIDJOURNEY_STATUS_PATH", "/mj/task/{id}/fetch").strip() or "/mj/task/{id}/fetch"
MIDJOURNEY_TIMEOUT_S = int(os.environ.get("MIDJOURNEY_TIMEOUT_S", "600") or 600)
MIDJOURNEY_POLL_DELAY_S = float(os.environ.get("MIDJOURNEY_POLL_DELAY_S", "5") or 5)
MIDJOURNEY_UNIT_COST_USD = _pricing_cost("MIDJOURNEY_UNIT_COST_USD", 0.168 if MIDJOURNEY_MODE == "turbo" else 0.056)
MIDJOURNEY_DEFAULT_VERSION = os.environ.get("MIDJOURNEY_DEFAULT_VERSION", "7").strip() or "7"

# Suno / music generation through CometAPI-compatible gateway.
# Endpoints differ by Comet channel, so paths are configurable and several payloads are tried.
SUNO_ENABLED = os.environ.get("SUNO_ENABLED", "1").strip().lower() not in ("0", "false", "no", "off")
SUNO_API_KEY = (os.environ.get("SUNO_API_KEY") or COMET_API_KEY).strip()
SUNO_BASE_URL = os.environ.get("SUNO_BASE_URL", COMET_BASE_URL).strip().rstrip("/")
SUNO_MODEL = os.environ.get("SUNO_MODEL", "chirp-bluejay").strip() or "chirp-bluejay"
# Comet Suno uses mv values, not OpenAI-style model names. Keep compatibility
# with old Render ENV values from earlier builds.
_suno_model_map = {
    "suno": "chirp-bluejay",
    "suno-v4": "chirp-v4",
    "suno-v4.0": "chirp-v4",
    "suno-v4.5": "chirp-auk",
    "suno-4.5": "chirp-auk",
    "suno-v4.5+": "chirp-bluejay",
    "suno-4.5+": "chirp-bluejay",
    "suno-v5": "chirp-crow",
    "suno-5": "chirp-crow",
}
SUNO_MODEL = _suno_model_map.get(SUNO_MODEL.lower(), SUNO_MODEL)
SUNO_CREATE_PATH = os.environ.get("SUNO_CREATE_PATH", "/suno/submit/music").strip() or "/suno/submit/music"
SUNO_STATUS_PATH = os.environ.get("SUNO_STATUS_PATH", "/suno/fetch/{id}").strip() or "/suno/fetch/{id}"
SUNO_TIMEOUT_S = int(os.environ.get("SUNO_TIMEOUT_S", "600") or 600)
SUNO_POLL_DELAY_S = float(os.environ.get("SUNO_POLL_DELAY_S", "5.0") or 5.0)
SUNO_COST_USD = _pricing_cost("SUNO_COST_USD", 0.20)
# РҹРҫРәР° РҫСӮРҙРөР»СҢРҪРҫР№ РәРҫР»РҫРҪРәРё music РІ Р‘Р” РҪРөСӮ, СҒРҝРёСҒСӢРІР°РөРј РёР· РҫРұСүРөРіРҫ РІРёРҙРөРҫ/generative-РұСҺРҙР¶РөСӮР°.
SUNO_BILLING_ENGINE = os.environ.get("SUNO_BILLING_ENGINE", "runway").strip().lower() or "runway"

# Background remove/replace pipeline:
# auto = Comet/Bria first, then local rembg fallback; local/rembg = only local.
BG_PROVIDER = os.environ.get("BG_PROVIDER", "photoroom-api-only").strip().lower() or "photoroom-api-only"
BG_COMET_MODEL = os.environ.get("BG_COMET_MODEL", "bria/remove-background").strip() or "bria/remove-background"
BG_COMET_REMOVE_PATH = os.environ.get("BG_COMET_REMOVE_PATH", "/v1/images/edits").strip() or "/v1/images/edits"
BG_REMOVE_TIMEOUT_S = float(os.environ.get("BG_REMOVE_TIMEOUT_S", "90") or 90)
BRIA_API_KEY = os.environ.get("BRIA_API_KEY", "").strip()
BRIA_BASE_URL = os.environ.get("BRIA_BASE_URL", "https://engine.prod.bria-api.com").strip().rstrip("/")
BRIA_REMOVE_PATH = os.environ.get("BRIA_REMOVE_PATH", "/v1/background/remove").strip() or "/v1/background/remove"
# Photoroom Remove Background API (primary production path)
PHOTOROOM_API_KEY = (
    os.environ.get("PHOTOROOM_API_KEY")
    or os.environ.get("PHOTOROOM_KEY")
    or os.environ.get("PHOTOROOM_REMOVE_BG_API_KEY")
    or ""
).strip()
PHOTOROOM_BASE_URL = os.environ.get("PHOTOROOM_BASE_URL", "https://sdk.photoroom.com").strip().rstrip("/")
PHOTOROOM_REMOVE_PATH = os.environ.get("PHOTOROOM_REMOVE_PATH", "/v1/segment").strip() or "/v1/segment"
PHOTOROOM_EDIT_BASE_URL = os.environ.get("PHOTOROOM_EDIT_BASE_URL", "https://image-api.photoroom.com").strip().rstrip("/")
PHOTOROOM_EDIT_PATH = os.environ.get("PHOTOROOM_EDIT_PATH", "/v2/edit").strip() or "/v2/edit"
PHOTOROOM_EDIT_ENABLED = os.environ.get("PHOTOROOM_EDIT_ENABLED", "1").strip().lower() not in ("0", "false", "no", "off")
PHOTOROOM_EDIT_TIMEOUT_S = float(os.environ.get("PHOTOROOM_EDIT_TIMEOUT_S", "95") or 95)
PHOTOROOM_EDIT_EXPAND_PROMPT_MODE = os.environ.get("PHOTOROOM_EDIT_EXPAND_PROMPT_MODE", "ai.never").strip() or "ai.never"
PHOTOROOM_EDIT_NEGATIVE_PROMPT = (
    os.environ.get("PHOTOROOM_EDIT_NEGATIVE_PROMPT")
    or "phone, smartphone, phone edge, mirror, selfie stick, frame, window frame, black bar, vertical bar, pole, reflection, extra object, extra person, duplicate person, extra hands, text, watermark, logo, illustration, cartoon, painting, anime, CGI, 3d render, surreal background, distorted face, distorted horizon"
).strip()
PHOTOROOM_EDIT_GUIDANCE_SCALE = float(os.environ.get("PHOTOROOM_EDIT_GUIDANCE_SCALE", "0.8") or 0.8)
PHOTOROOM_FORMAT = os.environ.get("PHOTOROOM_FORMAT", "png").strip().lower() or "png"
PHOTOROOM_CHANNELS = os.environ.get("PHOTOROOM_CHANNELS", "rgba").strip().lower() or "rgba"
PHOTOROOM_SIZE = os.environ.get("PHOTOROOM_SIZE", "hd").strip().lower() or "hd"
PHOTOROOM_CROP = os.environ.get("PHOTOROOM_CROP", "false").strip().lower() or "false"
PHOTOROOM_DESPILL = os.environ.get("PHOTOROOM_DESPILL", "false").strip().lower() or "false"
PHOTOROOM_TIMEOUT_S = float(os.environ.get("PHOTOROOM_TIMEOUT_S", "70") or 70)
PHOTOROOM_INPUT_MAX_SIDE = int(os.environ.get("PHOTOROOM_INPUT_MAX_SIDE", "1400") or 1400)
BG_OUTPUT_MAX_SIDE = int(os.environ.get("BG_OUTPUT_MAX_SIDE", "1400") or 1400)
BG_REALISTIC_BACKGROUNDS = os.environ.get("BG_REALISTIC_BACKGROUNDS", "1").strip().lower() not in ("0", "false", "no", "off")
BG_BACKGROUND_TIMEOUT_S = float(os.environ.get("BG_BACKGROUND_TIMEOUT_S", "12") or 12)
BG_CACHE_DIR = os.environ.get("BG_CACHE_DIR", "/tmp/bg_cache").strip() or "/tmp/bg_cache"
BG_REPLACE_TWO_STAGE = os.environ.get("BG_REPLACE_TWO_STAGE", "1").strip().lower() not in ("0", "false", "no", "off")
BG_REPLACE_GENERATE_PRESETS = os.environ.get("BG_REPLACE_GENERATE_PRESETS", "0").strip().lower() in ("1", "true", "yes", "on")
BG_REPLACE_GENERATE_CUSTOM = os.environ.get("BG_REPLACE_GENERATE_CUSTOM", "1").strip().lower() not in ("0", "false", "no", "off")
BG_REPLACE_USE_STOCK_BACKGROUNDS = os.environ.get("BG_REPLACE_USE_STOCK_BACKGROUNDS", "1").strip().lower() not in ("0", "false", "no", "off")
BG_REPLACE_JPEG_QUALITY = int(os.environ.get("BG_REPLACE_JPEG_QUALITY", "94") or 94)

# Face swap production pipeline:
# piapi = fast/cheap primary; segmind = quality fallback/premium.
FACESWAP_ENABLED = os.environ.get("FACESWAP_ENABLED", "1").strip().lower() not in ("0", "false", "no", "off")
FACESWAP_PROVIDER = os.environ.get("FACESWAP_PROVIDER", "piapi").strip().lower() or "piapi"
FACESWAP_FALLBACK_PROVIDER = os.environ.get("FACESWAP_FALLBACK_PROVIDER", "segmind-v2").strip().lower() or "segmind-v2"
FACESWAP_FAST_PROVIDER = os.environ.get("FACESWAP_FAST_PROVIDER", FACESWAP_PROVIDER).strip().lower() or FACESWAP_PROVIDER
FACESWAP_PREMIUM_PROVIDER = os.environ.get("FACESWAP_PREMIUM_PROVIDER", "segmind-v4").strip().lower() or "segmind-v4"
FACESWAP_ASK_TARGET_FACE = os.environ.get("FACESWAP_ASK_TARGET_FACE", "1").strip().lower() not in ("0", "false", "no", "off")
FACESWAP_ASK_SOURCE_FACE = os.environ.get("FACESWAP_ASK_SOURCE_FACE", "1").strip().lower() not in ("0", "false", "no", "off")
FACESWAP_STRICT_SELECTED_FACE = os.environ.get("FACESWAP_STRICT_SELECTED_FACE", "1").strip().lower() not in ("0", "false", "no", "off")
FACESWAP_FACE_DETECTION_ENABLED = os.environ.get("FACESWAP_FACE_DETECTION_ENABLED", "1").strip().lower() not in ("0", "false", "no", "off")
FACESWAP_DETECTION_MAX_SIDE = int(os.environ.get("FACESWAP_DETECTION_MAX_SIDE", "1200") or 1200)
FACESWAP_PREVIEW_MAX_SIDE = int(os.environ.get("FACESWAP_PREVIEW_MAX_SIDE", "1200") or 1200)
FACESWAP_FAST_COST_USD = _pricing_cost("FACESWAP_FAST_COST_USD", 0.03)
FACESWAP_PREMIUM_COST_USD = _pricing_cost("FACESWAP_PREMIUM_COST_USD", 0.12)
FACESWAP_TIMEOUT_S = float(os.environ.get("FACESWAP_TIMEOUT_S", "300") or 300)
FACESWAP_POLL_DELAY_S = float(os.environ.get("FACESWAP_POLL_DELAY_S", "2.5") or 2.5)
FACESWAP_INPUT_MAX_SIDE = int(os.environ.get("FACESWAP_INPUT_MAX_SIDE", "1600") or 1600)
FACESWAP_OUTPUT_MAX_SIDE = int(os.environ.get("FACESWAP_OUTPUT_MAX_SIDE", "1600") or 1600)
FACESWAP_RESULT_AS_DOCUMENT = os.environ.get("FACESWAP_RESULT_AS_DOCUMENT", "0").strip().lower() in ("1", "true", "yes", "on")
FACESWAP_IMAGE_DATA_URL = os.environ.get("FACESWAP_IMAGE_DATA_URL", "0").strip().lower() in ("1", "true", "yes", "on")
FACESWAP_WARN_TEXT = os.environ.get("FACESWAP_WARN_TEXT", "1").strip().lower() not in ("0", "false", "no", "off")
FACESWAP_MANUAL_CHOICE_IF_DETECTION_FAIL = os.environ.get("FACESWAP_MANUAL_CHOICE_IF_DETECTION_FAIL", "1").strip().lower() not in ("0", "false", "no", "off")
# v35: СӮРҫСҮРҪСӢР№ FaceSwap РҙР»СҸ РіСҖСғРҝРҝРҫРІСӢС… С„РҫСӮРҫ.
# РҳРҙРөСҸ: РөСҒР»Рё РҝРҫР»СҢР·РҫРІР°СӮРөР»СҢ РІСӢРұСҖР°Р» РәРҫРҪРәСҖРөСӮРҪРҫРө Р»РёСҶРҫ, РјСӢ РёР·РҫР»РёСҖСғРөРј СҚСӮРҫ Р»РёСҶРҫ РҙР»СҸ РҝСҖРҫРІР°Р№РҙРөСҖР°,
# Р° Р·Р°СӮРөРј РІРҫР·РІСҖР°СүР°РөРј РІ РёСҒС…РҫРҙРҪСӢР№ РәР°РҙСҖ СӮРҫР»СҢРәРҫ РёР·РјРөРҪС‘РҪРҪСғСҺ РҫРұР»Р°СҒСӮСҢ РІСӢРұСҖР°РҪРҪРҫРіРҫ Р»РёСҶР°. РўР°Рә РҝСҖРҫРІР°Р№РҙРөСҖ
# РҪРө РјРҫР¶РөСӮ СҒР»СғСҮР°Р№РҪРҫ РҝРҫРјРөРҪСҸСӮСҢ СҒРҫСҒРөРҙРҪРөРіРҫ СҮРөР»РҫРІРөРәР°, РҙР°Р¶Рө РөСҒР»Рё РөРіРҫ РІРҪСғСӮСҖРөРҪРҪРёР№ РҝРҫСҖСҸРҙРҫРә Р»РёСҶ РҙСҖСғРіРҫР№.
FACESWAP_PRECISE_COMPOSITE = os.environ.get("FACESWAP_PRECISE_COMPOSITE", "1").strip().lower() not in ("0", "false", "no", "off")
FACESWAP_FORCE_SEGMIND_FOR_MULTI = os.environ.get("FACESWAP_FORCE_SEGMIND_FOR_MULTI", "1").strip().lower() not in ("0", "false", "no", "off")
FACESWAP_GROUP_ALLOW_SEGMIND_FALLBACK = os.environ.get("FACESWAP_GROUP_ALLOW_SEGMIND_FALLBACK", "1").strip().lower() not in ("0", "false", "no", "off")
FACESWAP_FACE_BOX_FILTER_RATIO = float(os.environ.get("FACESWAP_FACE_BOX_FILTER_RATIO", "0.35") or 0.35)
FACESWAP_SOURCE_CROP_MARGIN = float(os.environ.get("FACESWAP_SOURCE_CROP_MARGIN", "2.20") or 2.20)
FACESWAP_TARGET_HIDE_MARGIN = float(os.environ.get("FACESWAP_TARGET_HIDE_MARGIN", "1.55") or 1.55)
FACESWAP_COMPOSITE_MARGIN_X = float(os.environ.get("FACESWAP_COMPOSITE_MARGIN_X", "1.85") or 1.85)
FACESWAP_COMPOSITE_MARGIN_Y_UP = float(os.environ.get("FACESWAP_COMPOSITE_MARGIN_Y_UP", "1.45") or 1.45)
FACESWAP_COMPOSITE_MARGIN_Y_DOWN = float(os.environ.get("FACESWAP_COMPOSITE_MARGIN_Y_DOWN", "1.75") or 1.75)

PIAPI_API_KEY = (os.environ.get("PIAPI_API_KEY") or os.environ.get("PIAPI_KEY") or "").strip()
PIAPI_BASE_URL = os.environ.get("PIAPI_BASE_URL", "https://api.piapi.ai").strip().rstrip("/")
PIAPI_FACE_CREATE_PATH = os.environ.get("PIAPI_FACE_CREATE_PATH", "/api/v1/task").strip() or "/api/v1/task"
PIAPI_FACE_STATUS_PATH = os.environ.get("PIAPI_FACE_STATUS_PATH", "/api/v1/task/{task_id}").strip() or "/api/v1/task/{task_id}"
PIAPI_FACE_MODEL = os.environ.get("PIAPI_FACE_MODEL", "Qubico/image-toolkit").strip() or "Qubico/image-toolkit"
PIAPI_FACE_TASK_TYPE = os.environ.get("PIAPI_FACE_TASK_TYPE", "face-swap").strip() or "face-swap"

SEGMIND_API_KEY = (os.environ.get("SEGMIND_API_KEY") or os.environ.get("SEGMIND_KEY") or "").strip()
SEGMIND_BASE_URL = os.environ.get("SEGMIND_BASE_URL", "https://api.segmind.com").strip().rstrip("/")
SEGMIND_FACESWAP_MODEL_FAST = os.environ.get("SEGMIND_FACESWAP_MODEL_FAST", "faceswap-v2").strip() or "faceswap-v2"
SEGMIND_FACESWAP_MODEL_PREMIUM = os.environ.get("SEGMIND_FACESWAP_MODEL_PREMIUM", "faceswap-v4").strip() or "faceswap-v4"
SEGMIND_FACE_RESTORE = os.environ.get("SEGMIND_FACE_RESTORE", "codeformer-v0.1.0.pth").strip() or "codeformer-v0.1.0.pth"
SEGMIND_FACE_SWAP_TYPE = os.environ.get("SEGMIND_FACE_SWAP_TYPE", "head").strip() or "head"
SEGMIND_FACE_STYLE_TYPE = os.environ.get("SEGMIND_FACE_STYLE_TYPE", "normal").strip() or "normal"

# Optional legacy remove.bg-compatible fallback; not used unless BG_PROVIDER=multi/auto and key exists.
REMOVE_BG_API_KEY = (os.environ.get("REMOVE_BG_API_KEY") or os.environ.get("REMOVEBG_API_KEY") or "").strip()
REMOVE_BG_BASE_URL = os.environ.get("REMOVE_BG_BASE_URL", "https://api.remove.bg").strip().rstrip("/")
REMOVE_BG_PATH = os.environ.get("REMOVE_BG_PATH", "/v1.0/removebg").strip() or "/v1.0/removebg"
REMOVE_BG_SIZE = os.environ.get("REMOVE_BG_SIZE", "auto").strip() or "auto"
REMOVE_BG_FORMAT = os.environ.get("REMOVE_BG_FORMAT", "png").strip() or "png"
REMOVE_BG_TIMEOUT_S = float(os.environ.get("REMOVE_BG_TIMEOUT_S", "60") or 60)
BG_DISABLE_LOCAL_REMBG = os.environ.get("BG_DISABLE_LOCAL_REMBG", "1").strip().lower() not in ("0", "false", "no", "off")
BRIA_ALLOW_LOCAL_FALLBACK = os.environ.get("BRIA_ALLOW_LOCAL_FALLBACK", "1").strip().lower() not in ("0", "false", "no", "off")
# РЎСӮР°СҖСӢР№ Render ENV РјРҫРі РҫСҒСӮР°РІР»СҸСӮСҢ LOCAL_REMBG_ENABLED=0.
# Р”Р»СҸ РҝСҖРҫРҙР°РәСҲРҪ-СҒРұРҫСҖРәРё С„РҫРҪР° РҝСҖРёРҪСғРҙРёСӮРөР»СҢРҪРҫ РІРәР»СҺСҮР°РөРј local rembg,
# РөСҒР»Рё СҒРҝРөСҶРёР°Р»СҢРҪРҫ РҪРө РҝРҫСҒСӮР°РІР»РөРҪ BG_FORCE_LOCAL_REMBG=0.
BG_FORCE_LOCAL_REMBG = os.environ.get("BG_FORCE_LOCAL_REMBG", "1").strip().lower() not in ("0", "false", "no", "off")
LOCAL_REMBG_ENABLED_RAW = os.environ.get("LOCAL_REMBG_ENABLED", "0").strip().lower()
LOCAL_REMBG_ENABLED = (not BG_DISABLE_LOCAL_REMBG) and (BG_FORCE_LOCAL_REMBG or (LOCAL_REMBG_ENABLED_RAW not in ("0", "false", "no", "off")))
REMBG_MODEL = os.environ.get("REMBG_MODEL", "u2netp").strip() or "u2netp"
LOCAL_REMBG_TIMEOUT_S = float(os.environ.get("LOCAL_REMBG_TIMEOUT_S", "240") or 240)
LOCAL_REMBG_SUBPROCESS = os.environ.get("LOCAL_REMBG_SUBPROCESS", "1").strip().lower() not in ("0", "false", "no", "off")
BG_ACTION_TIMEOUT_S = float(os.environ.get("BG_ACTION_TIMEOUT_S", "180") or 180)
BG_LOCAL_FIRST = os.environ.get("BG_LOCAL_FIRST", "1").strip().lower() not in ("0", "false", "no", "off")
REMBG_MAX_SIDE = int((os.environ.get("REMBG_MAX_SIDE") or "1600").strip() or 1600)
REMBG_MODEL_FALLBACKS = [m.strip() for m in os.environ.get("REMBG_MODEL_FALLBACKS", REMBG_MODEL).split(",") if m.strip()]
if REMBG_MODEL not in REMBG_MODEL_FALLBACKS:
    REMBG_MODEL_FALLBACKS.insert(0, REMBG_MODEL)
with contextlib.suppress(Exception):
    os.environ.setdefault("U2NET_HOME", U2NET_HOME)
    os.makedirs(U2NET_HOME, exist_ok=True)
    os.makedirs(XDG_CACHE_HOME, exist_ok=True)
_REMBG_SESSION = None
_REMBG_SESSION_LOCK = threading.Lock()
_LAST_BG_ERRORS: list[str] = []

def _bg_note_error(message: str):
    try:
        msg = str(message or "").strip()
        if not msg:
            return
        _LAST_BG_ERRORS.append(msg[:700])
        del _LAST_BG_ERRORS[:-8]
    except Exception:
        pass

def _bg_last_errors_text() -> str:
    if not _LAST_BG_ERRORS:
        return "РҫСҲРёРұРҫРә РҝРҫРәР° РҪРөСӮ"
    return "\n".join(f"вҖў {x}" for x in _LAST_BG_ERRORS[-5:])


async def cmd_diag_bg(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Р”РёР°РіРҪРҫСҒСӮРёРәР° СғРҙР°Р»РөРҪРёСҸ/Р·Р°РјРөРҪСӢ С„РҫРҪР°: ENV, rembg import, model session."""
    lines: list[str] = []
    lines.append(f"рҹ§Ә BG diagnostic / {PATCH_VERSION}")
    lines.append(f"BG_PROVIDER={BG_PROVIDER}")
    lines.append(f"PHOTOROOM_API_KEY={'on' if bool(PHOTOROOM_API_KEY) else 'off'} remove={PHOTOROOM_BASE_URL}{PHOTOROOM_REMOVE_PATH} edit={PHOTOROOM_EDIT_BASE_URL}{PHOTOROOM_EDIT_PATH} edit_enabled={'on' if PHOTOROOM_EDIT_ENABLED else 'off'}")
    lines.append(f"PHOTOROOM_FORMAT={PHOTOROOM_FORMAT} channels={PHOTOROOM_CHANNELS} size={PHOTOROOM_SIZE} crop={PHOTOROOM_CROP} timeout={PHOTOROOM_TIMEOUT_S}")
    lines.append(f"PHOTOROOM_INPUT_MAX_SIDE={PHOTOROOM_INPUT_MAX_SIDE} BG_OUTPUT_MAX_SIDE={BG_OUTPUT_MAX_SIDE} edit_timeout={PHOTOROOM_EDIT_TIMEOUT_S}s")
    lines.append(f"BG_REALISTIC_BACKGROUNDS={BG_REALISTIC_BACKGROUNDS} cache={BG_CACHE_DIR}")
    lines.append(f"COMET_API_KEY={'on' if bool(COMET_API_KEY) else 'off'}")
    lines.append(f"BRIA_API_KEY={'on' if bool(BRIA_API_KEY) else 'off'}")
    lines.append(f"REMOVE_BG_API_KEY={'on' if bool(REMOVE_BG_API_KEY) else 'off'} base={REMOVE_BG_BASE_URL}{REMOVE_BG_PATH}")
    lines.append(f"LOCAL_REMBG_ENABLED={LOCAL_REMBG_ENABLED} raw={LOCAL_REMBG_ENABLED_RAW} force={BG_FORCE_LOCAL_REMBG} disable={BG_DISABLE_LOCAL_REMBG}")
    lines.append(f"LOCAL_REMBG_SUBPROCESS={LOCAL_REMBG_SUBPROCESS} timeout={LOCAL_REMBG_TIMEOUT_S} action_timeout={BG_ACTION_TIMEOUT_S}")
    lines.append(f"BG_LOCAL_FIRST={BG_LOCAL_FIRST}")
    lines.append(f"REMBG_MAX_SIDE={REMBG_MAX_SIDE}")
    lines.append(f"rembg_import={'ok' if rembg_remove is not None else 'FAILED'}")
    if rembg_remove is None:
        lines.append(f"REMBG_IMPORT_ERROR={REMBG_IMPORT_ERROR[:700]}")
    lines.append(f"REMBG_MODEL={REMBG_MODEL}")
    lines.append(f"REMBG_MODEL_FALLBACKS={','.join(REMBG_MODEL_FALLBACKS) or '-'}")
    lines.append(f"U2NET_HOME={os.environ.get('U2NET_HOME', '')}")
    lines.append(f"XDG_CACHE_HOME={os.environ.get('XDG_CACHE_HOME', '')}")
    lines.append(f"cache_bootstrap_u2net={U2NET_HOME}")

    # РҹСҖРҫРІРөСҖСҸРөРј РҝР°РҝРәРё РұРөР· РҝР°РҙРөРҪРёСҸ РәРҫРјР°РҪРҙСӢ.
    for d in (os.environ.get('U2NET_HOME', ''), os.environ.get('XDG_CACHE_HOME', '')):
        if d:
            try:
                os.makedirs(d, exist_ok=True)
                lines.append(f"dir_ok={d}")
            except Exception as e:
                lines.append(f"dir_FAIL={d}: {e}")

    session_ok = False
    if LOCAL_REMBG_ENABLED and rembg_remove is not None:
        try:
            # РҹРөСҖРІСӢР№ Р·Р°РҝСғСҒРә РјРҫР¶РөСӮ СҒРәР°СҮР°СӮСҢ РјРҫРҙРөР»СҢ, РҝРҫСҚСӮРҫРјСғ РҙР°С‘Рј РұРҫР»СҢСҲРө РІСҖРөРјРөРҪРё.
            sess = await asyncio.wait_for(asyncio.to_thread(_get_local_rembg_session), timeout=min(max(30, int(LOCAL_REMBG_TIMEOUT_S)), 240))
            session_ok = sess is not None
        except Exception as e:
            _bg_note_error(f"diag session failed: {e}")
            lines.append(f"session_error={repr(e)[:700]}")

    lines.append(f"session_ok={session_ok}")
    lines.append("РҹРҫСҒР»РөРҙРҪРёРө РҫСҲРёРұРәРё:")
    lines.append(_bg_last_errors_text())

    await update.effective_message.reply_text("\n".join(lines)[:3900])
RUNWAY_COMET_CREATE_PATH = os.environ.get("RUNWAY_COMET_CREATE_PATH", "/runwayml/v1/image_to_video").strip() or "/runwayml/v1/image_to_video"
RUNWAY_COMET_STATUS_PATH = os.environ.get("RUNWAY_COMET_STATUS_PATH", "/runwayml/v1/tasks/{id}").strip() or "/runwayml/v1/tasks/{id}"
# Comet/Runway РёРҪРҫРіРҙР° РҪРөСҒРәРҫР»СҢРәРҫ РјРёРҪСғСӮ РҫСӮРІРөСҮР°РөСӮ task_not_exist: СҚСӮРҫ СҒСӮР°РҙРёСҸ РёРҪРёСҶРёР°Р»РёР·Р°СҶРёРё.
# Р•СҒР»Рё РҫСӮРІРөСӮ РІРёСҒРёСӮ СҒР»РёСҲРәРҫРј РҙРҫР»РіРҫ, РҪРө РҙРөСҖР¶РёРј РҝРҫР»СҢР·РҫРІР°СӮРөР»СҸ 20 РјРёРҪСғСӮ вҖ” РјСҸРіРәРҫ РҝРөСҖРөРәР»СҺСҮР°РөРјСҒСҸ РҪР° Kling.
RUNWAY_TASK_NOT_EXIST_FALLBACK_S = int(os.environ.get("RUNWAY_TASK_NOT_EXIST_FALLBACK_S", "45") or 45)
RUNWAY_AUTO_FALLBACK_KLING = os.environ.get("RUNWAY_AUTO_FALLBACK_KLING", "1").strip().lower() not in ("0", "false", "no", "off")
# v76: production-safe Runway/Comet controls.
RUNWAY_IMAGE2VIDEO_ENABLED = os.environ.get("RUNWAY_IMAGE2VIDEO_ENABLED", "1").strip().lower() not in ("0", "false", "no", "off")
RUNWAY_HIDE_TECH_ERRORS = os.environ.get("RUNWAY_HIDE_TECH_ERRORS", "1").strip().lower() not in ("0", "false", "no", "off")
RUNWAY_IMAGE2VIDEO_FAIL_FAST = os.environ.get("RUNWAY_IMAGE2VIDEO_FAIL_FAST", "1").strip().lower() not in ("0", "false", "no", "off")
RUNWAY_PROVIDER_COOLDOWN_S = int(os.environ.get("RUNWAY_PROVIDER_COOLDOWN_S", "300") or 300)
RUNWAY_PROVIDER_FAIL_THRESHOLD = int(os.environ.get("RUNWAY_PROVIDER_FAIL_THRESHOLD", "2") or 2)
# v77: РҝРөСҖРІСӢРј РҝСҖРҫРұСғРөРј РҝСғРұР»РёСҮРҪСӢР№ РёРҙРөРҪСӮРёС„РёРәР°СӮРҫСҖ Comet model page, Р·Р°СӮРөРј backend-Р°Р»РёР°СҒСӢ.
RUNWAY_IMAGE2VIDEO_MODELS_ENV = os.environ.get("RUNWAY_IMAGE2VIDEO_MODELS", "runwayml-image-to-video,gen4_turbo,gen3a_turbo,veo3.1_fast,veo3.1,veo3").strip()
RUNWAY_PUBLIC_FALLBACK_TEXT = os.environ.get("RUNWAY_PUBLIC_FALLBACK_TEXT", "вҡ пёҸ РҡР°РҪР°Р» Runway СҒРөР№СҮР°СҒ РҪРөРҙРҫСҒСӮСғРҝРөРҪ Сғ РҝСҖРҫРІР°Р№РҙРөСҖР°. РҗРІСӮРҫРјР°СӮРёСҮРөСҒРәРё Р·Р°РҝСғСҒРәР°СҺ Kling РұРөР· РҙРҫРҝРҫР»РҪРёСӮРөР»СҢРҪРҫРіРҫ СҒРҝРёСҒР°РҪРёСҸ.").strip()
# v88: official Runway Developer API first; Comet is secondary and Kling is the production fallback.
RUNWAY_DIRECT_FIRST = os.environ.get("RUNWAY_DIRECT_FIRST", "1").strip().lower() not in ("0", "false", "no", "off")
RUNWAY_TEXT_MODEL = os.environ.get("RUNWAY_TEXT_MODEL", "gen4.5").strip() or "gen4.5"
RUNWAY_DIRECT_TEXT_MODELS_ENV = os.environ.get("RUNWAY_DIRECT_TEXT_MODELS", "gen4.5").strip()
RUNWAY_DIRECT_I2V_MODELS_ENV = os.environ.get("RUNWAY_DIRECT_I2V_MODELS", "gen4.5,gen4_turbo").strip()
RUNWAY_DIRECT_RETRY_ATTEMPTS = max(1, int(os.environ.get("RUNWAY_DIRECT_RETRY_ATTEMPTS", "4") or 4))
RUNWAY_DIRECT_RETRY_BASE_S = max(0.5, float(os.environ.get("RUNWAY_DIRECT_RETRY_BASE_S", "1.5") or 1.5))
RUNWAY_DIRECT_POLL_INTERVAL_S = max(5.0, float(os.environ.get("RUNWAY_DIRECT_POLL_INTERVAL_S", "5.0") or 5.0))
RUNWAY_DIRECT_POLL_MAX_INTERVAL_S = max(RUNWAY_DIRECT_POLL_INTERVAL_S, float(os.environ.get("RUNWAY_DIRECT_POLL_MAX_INTERVAL_S", "15.0") or 15.0))
RUNWAY_DIRECT_UPLOAD_ATTEMPTS = max(1, int(os.environ.get("RUNWAY_DIRECT_UPLOAD_ATTEMPTS", "2") or 2))
RUNWAY_DIRECT_DATA_URI_FALLBACK = os.environ.get("RUNWAY_DIRECT_DATA_URI_FALLBACK", "1").strip().lower() not in ("0", "false", "no", "off")
RUNWAY_COMET_TEXT_MODELS_ENV = os.environ.get("RUNWAY_COMET_TEXT_MODELS", "runway-video,gen4.5").strip()
RUNWAY_TEXT_FALLBACK_KLING = os.environ.get("RUNWAY_TEXT_FALLBACK_KLING", "1").strip().lower() not in ("0", "false", "no", "off")
# РҹСҖРөРҙРҫРұСҖР°РұРҫСӮРәР° РёСҒС…РҫРҙРҪРёРәР° РҙР»СҸ imageвҶ’video: СғРјРөРҪСҢСҲР°РөСӮ СҖРёСҒРә, СҮСӮРҫ Kling/Runway РҫР¶РёРІРёСӮ РІРөСҒСҢ СҒРәСҖРёРҪСҲРҫСӮ СӮРөР»РөС„РҫРҪР° РІРјРөСҒСӮРҫ РҝРҫСҖСӮСҖРөСӮР°.
I2V_PREPROCESS_ENABLED = os.environ.get("I2V_PREPROCESS_ENABLED", "1").strip().lower() not in ("0", "false", "no", "off")
I2V_WARN_BAD_SOURCE = os.environ.get("I2V_WARN_BAD_SOURCE", "1").strip().lower() not in ("0", "false", "no", "off")
I2V_AUTOCROP_BLACK_BORDERS = os.environ.get("I2V_AUTOCROP_BLACK_BORDERS", "1").strip().lower() not in ("0", "false", "no", "off")
I2V_MAX_SOURCE_SIDE = int(os.environ.get("I2V_MAX_SOURCE_SIDE", "1280") or 1280)
I2V_BAD_SOURCE_POLICY = os.environ.get("I2V_BAD_SOURCE_POLICY", "warn_continue").strip().lower()  # warn_continue | ask_clean
I2V_KLING_SAFE_PROMPT_SUFFIX = os.environ.get("I2V_KLING_SAFE_PROMPT_SUFFIX", "animate only the person/photo content; do not create phone screens, app interface, black borders, UI overlays, or screenshots; keep identity and original composition").strip()
SORA_API_KEY   = (os.environ.get("SORA_API_KEY") or COMET_API_KEY).strip()
SORA_MODEL     = os.environ.get("SORA_MODEL", "sora-2").strip()
# Р—Р°СүРёСӮР° РҫСӮ СҒСӮР°СҖРҫРіРҫ ENV Render: РөСҒР»Рё СӮР°Рј СҒР»СғСҮР°Р№РҪРҫ РҫСҒСӮР°Р»РҫСҒСҢ sora-1,
# Comet РІРҫР·РІСҖР°СүР°РөСӮ model_not_found. Р”Р»СҸ СҚСӮРҫР№ СҒРұРҫСҖРәРё РҝСҖРёРҪСғРҙРёСӮРөР»СҢРҪРҫ РёСҒРҝРҫР»СҢР·СғРөРј sora-2.
if SORA_MODEL.lower() in ("", "sora", "sora-1", "sora1"):
    SORA_MODEL = "sora-2"
SORA_CREATE_PATH = os.environ.get("SORA_CREATE_PATH", "/v1/videos").strip() or "/v1/videos"
SORA_STATUS_PATH = os.environ.get("SORA_STATUS_PATH", "/v1/videos/{id}").strip() or "/v1/videos/{id}"
KLING_API_KEY  = (os.environ.get("KLING_API_KEY") or COMET_API_KEY).strip()
KLING_MODEL    = os.environ.get("KLING_MODEL", "kling-v1-6").strip()
KLING_CREATE_PATH = os.environ.get("KLING_CREATE_PATH", "/kling/v1/videos/image2video").strip() or "/kling/v1/videos/image2video"
KLING_STATUS_PATH = os.environ.get("KLING_STATUS_PATH", "/kling/v1/videos/image2video/{id}").strip() or "/kling/v1/videos/image2video/{id}"
# TextвҶ’Video СҮРөСҖРөР· CometAPI: Sora 2 (СӮРҫР»СҢРәРҫ РұРөР· Р»СҺРҙРөР№), Kling Рё Runway. Luma РІСҖРөРјРөРҪРҪРҫ СҒРәСҖСӢСӮР°.
KLING_TEXT_CREATE_PATH = os.environ.get("KLING_TEXT_CREATE_PATH", "/kling/v1/videos/text2video").strip() or "/kling/v1/videos/text2video"
KLING_TEXT_STATUS_PATH = os.environ.get("KLING_TEXT_STATUS_PATH", "/kling/v1/videos/text2video/{id}").strip() or "/kling/v1/videos/text2video/{id}"

# Kling Avatar / talking head / photoвҶ’music-video
KLING_AVATAR_CREATE_PATH = os.environ.get("KLING_AVATAR_CREATE_PATH", "/kling/v1/videos/avatar/image2video").strip() or "/kling/v1/videos/avatar/image2video"
KLING_AVATAR_STATUS_PATH = os.environ.get("KLING_AVATAR_STATUS_PATH", "/kling/v1/videos/avatar/image2video/{id}").strip() or "/kling/v1/videos/avatar/image2video/{id}"
KLING_AVATAR_MODE = os.environ.get("KLING_AVATAR_MODE", "std").strip().lower() or "std"
if KLING_AVATAR_MODE not in ("std", "pro"):
    KLING_AVATAR_MODE = "std"
KLING_AVATAR_PROMPT = os.environ.get(
    "KLING_AVATAR_PROMPT",
    "The person talks naturally to camera, realistic facial motion, accurate lip sync, subtle head movement, stable identity."
).strip()

# TextвҶ’speech for avatar. OpenAI TTS is primary because it supports Russian well;
# Kling TTS is kept as a fallback when the channel is available.
AVATAR_TTS_PROVIDER = os.environ.get("AVATAR_TTS_PROVIDER", "openai").strip().lower() or "openai"
KLING_TTS_CREATE_PATH = os.environ.get("KLING_TTS_CREATE_PATH", "/kling/v1/audio/tts").strip() or "/kling/v1/audio/tts"
KLING_TTS_STATUS_PATH = os.environ.get("KLING_TTS_STATUS_PATH", "/kling/v1/audio/tts/{id}").strip() or "/kling/v1/audio/tts/{id}"
KLING_TTS_VOICE_ID = os.environ.get("KLING_TTS_VOICE_ID", "genshin_vindi2").strip()
KLING_TTS_LANGUAGE = os.environ.get("KLING_TTS_LANGUAGE", "en").strip().lower() or "en"
KLING_TTS_SPEED = _env_float("KLING_TTS_SPEED", 1.0)

PHOTO_CLIP_SOUND = os.environ.get("PHOTO_CLIP_SOUND", "1").strip().lower() not in ("0", "false", "no", "off")
PHOTO_CLIP_MODE = os.environ.get("PHOTO_CLIP_MODE", "pro").strip().lower() or "pro"
if PHOTO_CLIP_MODE not in ("std", "pro"):
    PHOTO_CLIP_MODE = "pro"
SUNO_AUTO_FOR_PHOTO_CLIP = os.environ.get("SUNO_AUTO_FOR_PHOTO_CLIP", "1").strip().lower() in ("1", "true", "yes", "on")
PHOTO_CLIP_PIPELINE = os.environ.get("PHOTO_CLIP_PIPELINE", "1").strip().lower() not in ("0", "false", "no", "off")
PHOTO_CLIP_VIDEO_ENGINE = os.environ.get("PHOTO_CLIP_VIDEO_ENGINE", "kling").strip().lower() or "kling"
PHOTO_CLIP_DEFAULT_DURATION_S = int(os.environ.get("PHOTO_CLIP_DEFAULT_DURATION_S", "15") or 15)
PHOTO_CLIP_MAX_DURATION_S = int(os.environ.get("PHOTO_CLIP_MAX_DURATION_S", "30") or 30)
PHOTO_CLIP_MUX_AUDIO = os.environ.get("PHOTO_CLIP_MUX_AUDIO", "1").strip().lower() not in ("0", "false", "no", "off")
PHOTO_CLIP_SEND_BASE_IF_MUX_FAILS = os.environ.get("PHOTO_CLIP_SEND_BASE_IF_MUX_FAILS", "0").strip().lower() not in ("0", "false", "no", "off")
PHOTO_CLIP_BACKGROUND_TASK = os.environ.get("PHOTO_CLIP_BACKGROUND_TASK", "1").strip().lower() not in ("0", "false", "no", "off")
PHOTO_CLIP_PARALLEL_STAGES = os.environ.get("PHOTO_CLIP_PARALLEL_STAGES", "1").strip().lower() not in ("0", "false", "no", "off")
PHOTO_CLIP_SUNO_FAST_TIMEOUT_S = int(os.environ.get("PHOTO_CLIP_SUNO_FAST_TIMEOUT_S", "900") or 900)
PHOTO_CLIP_AUDIO_AFTER_VIDEO_WAIT_S = int(os.environ.get("PHOTO_CLIP_AUDIO_AFTER_VIDEO_WAIT_S", "900") or 900)
PHOTO_CLIP_TOTAL_USER_WAIT_S = int(os.environ.get("PHOTO_CLIP_TOTAL_USER_WAIT_S", "1200") or 1200)
PHOTO_CLIP_SEND_BASE_WHILE_MUSIC_PENDING = os.environ.get("PHOTO_CLIP_SEND_BASE_WHILE_MUSIC_PENDING", "0").strip().lower() not in ("0", "false", "no", "off")
FFMPEG_MUX_TIMEOUT_S = int(os.environ.get("FFMPEG_MUX_TIMEOUT_S", "180") or 180)
FFMPEG_MUX_COPY_FIRST = os.environ.get("FFMPEG_MUX_COPY_FIRST", "1").strip().lower() not in ("0", "false", "no", "off")
FFMPEG_MUX_REENCODE_PRESET = os.environ.get("FFMPEG_MUX_REENCODE_PRESET", "ultrafast").strip() or "ultrafast"
FFMPEG_MUX_AUDIO_BITRATE = os.environ.get("FFMPEG_MUX_AUDIO_BITRATE", "128k").strip() or "128k"
FFMPEG_MUX_MAX_MB = int(os.environ.get("FFMPEG_MUX_MAX_MB", "45") or 45)
FFMPEG_MUX_CRF = os.environ.get("FFMPEG_MUX_CRF", "32").strip() or "32"
FFMPEG_MUX_SCALE_HEIGHT = int(os.environ.get("FFMPEG_MUX_SCALE_HEIGHT", "720") or 720)
FFMPEG_MUX_FPS = int(os.environ.get("FFMPEG_MUX_FPS", "24") or 24)

# Audited provider cost estimates. Retail price is calculated centrally with the
# service multiplier. Canonical mode protects the commercial margin from stale ENV.
SORA_COST_PER_SECOND_USD = _pricing_cost("SORA_COST_PER_SECOND_USD", 0.10)
SORA_PRO_COST_PER_SECOND_USD = _pricing_cost("SORA_PRO_COST_PER_SECOND_USD", 0.30)
KLING_5S_COST_USD = _pricing_cost("KLING_5S_COST_USD", 0.23)
RUNWAY_COST_PER_SECOND_USD = _pricing_cost("RUNWAY_COST_PER_SECOND_USD", 0.12)
RUNWAY_TURBO_COST_PER_SECOND_USD = _pricing_cost("RUNWAY_TURBO_COST_PER_SECOND_USD", 0.05)
RUNWAY_5S_COST_USD = _pricing_cost("RUNWAY_5S_COST_USD", RUNWAY_COST_PER_SECOND_USD * 5)
SORA_UNIT_COST_USD = _pricing_cost("SORA_UNIT_COST_USD", SORA_COST_PER_SECOND_USD * 5)
SORA_PRO_UNIT_COST_USD = _pricing_cost("SORA_PRO_UNIT_COST_USD", SORA_PRO_COST_PER_SECOND_USD * 5)
KLING_UNIT_COST_USD = _pricing_cost("KLING_UNIT_COST_USD", KLING_5S_COST_USD)
AVATAR_UNIT_COST_USD = _pricing_cost("AVATAR_UNIT_COST_USD", 0.65)
# 15-second Kling video + Suno + mux/storage overhead.
PHOTO_CLIP_UNIT_COST_USD = _pricing_cost("PHOTO_CLIP_UNIT_COST_USD", 1.00)
SUNO_UNIT_COST_USD = _pricing_cost("SUNO_UNIT_COST_USD", 0.20)
# Music + avatar/lip-sync + retries/processing reserve.
VOCAL_CLIP_UNIT_COST_USD = _pricing_cost("VOCAL_CLIP_UNIT_COST_USD", 1.50)
# Kling Avatar СҮРөСҖРөР· Comet РҪРөСҒСӮР°РұРёР»СҢРҪРҫ РҝСҖРёРҪРёРјР°РөСӮ РҙР»РёРҪРҪСӢРө Suno-СӮСҖРөРәРё.
# Р”Р»СҸ lip-sync РҙРөСҖР¶РёРј РұРөР·РҫРҝР°СҒРҪСӢР№ С„СҖР°РіРјРөРҪСӮ: СӮР°Рә Р·Р°РҙР°СҮР° РұСӢСҒСӮСҖРөРө СҒСӮР°СҖСӮСғРөСӮ Рё РҪРө РІРёСҒРёСӮ РҪР° polling.
VOCAL_CLIP_MAX_AUDIO_S = int(os.environ.get("VOCAL_CLIP_MAX_AUDIO_S", "65") or 65)
VOCAL_CLIP_MIN_AUDIO_S = int(os.environ.get("VOCAL_CLIP_MIN_AUDIO_S", "12") or 12)
VOCAL_CLIP_KLING_MAX_WAIT_S = int(os.environ.get("VOCAL_CLIP_KLING_MAX_WAIT_S", "1800") or 1800)
TEXT_VIDEO_UNIT_COST_USD = _env_float("TEXT_VIDEO_UNIT_COST_USD", KLING_UNIT_COST_USD)
TEXT_VIDEO_DEFAULT_ENGINE = os.environ.get("TEXT_VIDEO_DEFAULT_ENGINE", "kling").strip().lower() or "kling"
TEXT_VIDEO_ALLOW_RUNWAY = os.environ.get("TEXT_VIDEO_ALLOW_RUNWAY", "1").strip().lower() not in ("0", "false", "no", "off")

_photo_clip_background_jobs: set[str] = set()
_vocal_clip_background_jobs: set[str] = set()

# AI selfie / Nano Banana style multi-image editor routed through CometAPI.
# Expected path is OpenAI-compatible image edit endpoint on Comet.
# Comet Nano Banana / Gemini image edit.
# Production primary route is Gemini-style generateContent through CometAPI.
# Comet docs show Nano Banana/Gemini image via /v1beta/models/{model}:generateContent with inline image data.
COMET_IMAGE_EDIT_MODEL = os.environ.get("COMET_IMAGE_EDIT_MODEL", "gemini-2.5-flash-image").strip() or "gemini-2.5-flash-image"
COMET_IMAGE_EDIT_FALLBACK_MODELS = [m.strip() for m in os.environ.get("COMET_IMAGE_EDIT_FALLBACK_MODELS", "gemini-2.5-flash-image,gemini-2.5-flash-image-preview,gemini-3.1-flash-image,gemini-3-pro-image,gemini-2-5-flash-image").split(",") if m.strip()]
if COMET_IMAGE_EDIT_MODEL not in COMET_IMAGE_EDIT_FALLBACK_MODELS:
    COMET_IMAGE_EDIT_FALLBACK_MODELS.insert(0, COMET_IMAGE_EDIT_MODEL)
COMET_IMAGE_EDIT_PATH = os.environ.get("COMET_IMAGE_EDIT_PATH", "/v1beta/models/{model}:generateContent").strip() or "/v1beta/models/{model}:generateContent"
COMET_IMAGE_EDIT_STATUS_PATH = os.environ.get("COMET_IMAGE_EDIT_STATUS_PATH", "").strip()
COMET_IMAGE_EDIT_TIMEOUT_S = float(os.environ.get("COMET_IMAGE_EDIT_TIMEOUT_S", "600") or 600)
COMET_IMAGE_EDIT_OPENAI_FALLBACK = os.environ.get("COMET_IMAGE_EDIT_OPENAI_FALLBACK", "0").strip().lower() in ("1", "true", "yes", "on")
AI_SELFIE_UNIT_COST_USD = _pricing_cost("AI_SELFIE_UNIT_COST_USD", 0.15)
AI_SELFIE_DEFAULT_ASPECT = os.environ.get("AI_SELFIE_DEFAULT_ASPECT", "4:5").strip() or "4:5"
AI_SELFIE_IMAGE_SIZE = os.environ.get("AI_SELFIE_IMAGE_SIZE", "1K").strip() or "1K"
AI_SELFIE_MAX_SIDE = int(os.environ.get("AI_SELFIE_MAX_SIDE", "1024") or 1024)
AI_SELFIE_SEND_AS_DOCUMENT = os.environ.get("AI_SELFIE_SEND_AS_DOCUMENT", "1").strip().lower() not in ("0", "false", "no", "off")
AI_SELFIE_PROVIDER = os.environ.get("AI_SELFIE_PROVIDER", "comet").strip().lower() or "comet"
AI_SELFIE_ALLOW_PUBLIC_FIGURES = os.environ.get("AI_SELFIE_ALLOW_PUBLIC_FIGURES", "1").strip().lower() not in ("0", "false", "no", "off")
AI_SELFIE_FAST_MODE = os.environ.get("AI_SELFIE_FAST_MODE", "1").strip().lower() not in ("0", "false", "no", "off")
AI_SELFIE_RETRY_ON_TIMEOUT = os.environ.get("AI_SELFIE_RETRY_ON_TIMEOUT", "1").strip().lower() not in ("0", "false", "no", "off")

# РўР°Р№РјР°СғСӮСӢ
LUMA_MAX_WAIT_S     = int((os.environ.get("LUMA_MAX_WAIT_S") or "900").strip() or 900)
LUMA_TEMP_DISABLED = True  # РІСҖРөРјРөРҪРҪР°СҸ Р·Р°РіР»СғСҲРәР°: СҒРәСҖСӢРІР°РөРј Luma РёР· РјРөРҪСҺ Рё РҫСӮРәР»СҺСҮР°РөРј РёСҒРҝРҫР»СҢР·РҫРІР°РҪРёРө
RUNWAY_MAX_WAIT_S   = int((os.environ.get("RUNWAY_MAX_WAIT_S") or "1200").strip() or 1200)
VIDEO_POLL_DELAY_S  = float((os.environ.get("VIDEO_POLL_DELAY_S") or "6.0").strip() or 6.0)
VIDEO_RESULT_SEND_AS_DOCUMENT = os.getenv("VIDEO_RESULT_SEND_AS_DOCUMENT", "1") == "1"
TELEGRAM_RESULT_MAX_MB = int(os.environ.get("TELEGRAM_RESULT_MAX_MB", "48") or 48)
TELEGRAM_VIDEO_COMPRESS_ON_FAIL = os.getenv("TELEGRAM_VIDEO_COMPRESS_ON_FAIL", "1") == "1"
VIDEO_RESULT_DEDUPE_TTL_S = int((os.getenv("VIDEO_RESULT_DEDUPE_TTL_S") or "900").strip() or 900)
SORA_AUTO_FALLBACK_KLING = os.getenv("SORA_AUTO_FALLBACK_KLING", "1") == "1"
_SENT_VIDEO_KEYS: dict[str, float] = {}

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ UTILS ---------
_LUMA_ACTIVE_BASE = None  # РәСҚСҲ РҝРҫСҒР»РөРҙРҪРөРіРҫ Р¶РёРІРҫРіРҫ РұР°Р·РҫРІРҫРіРҫ URL

async def _pick_luma_base(client: httpx.AsyncClient) -> str:
    global _LUMA_ACTIVE_BASE
    candidates = []
    if _LUMA_ACTIVE_BASE:
        candidates.append(_LUMA_ACTIVE_BASE)
    if LUMA_BASE_URL and LUMA_BASE_URL not in candidates:
        candidates.append(LUMA_BASE_URL)
    for b in LUMA_FALLBACKS:
        if b not in candidates:
            candidates.append(b)
    for base in candidates:
        try:
            url = f"{base}{LUMA_CREATE_PATH}"
            r = await client.options(url, timeout=10.0)
            if r.status_code in (200, 201, 202, 204, 400, 401, 403, 404, 405):
                _LUMA_ACTIVE_BASE = base
                if base != LUMA_BASE_URL:
                    log.info("Luma base switched to fallback: %s", base)
                return base
        except Exception as e:
            log.warning("Luma base probe failed for %s: %s", base, e)
    return LUMA_BASE_URL or "https://api.lumalabs.ai/dream-machine/v1"

# Payments / DB
PROVIDER_TOKEN = os.environ.get("PROVIDER_TOKEN_YOOKASSA", "").strip()
CURRENCY       = "RUB"
DB_PATH        = os.path.abspath(os.environ.get("DB_PATH", "subs.db"))

# Presentation/Catalog Studio v86
PRESENTATION_DATA_DIR = os.environ.get(
    "PRESENTATION_DATA_DIR",
    os.path.join(os.path.dirname(DB_PATH) or ".", "presentation_studio"),
).strip() or "/tmp/presentation_studio"
PRESENTATION_MAX_UPLOADS = int(os.environ.get("PRESENTATION_MAX_UPLOADS", "60") or 60)
PRESENTATION_MAX_GENERATED_IMAGES = int(os.environ.get("PRESENTATION_MAX_GENERATED_IMAGES", "8") or 8)
PRESENTATION_RENDER_COST_USD = _pricing_cost("PRESENTATION_RENDER_COST_USD", 0.20)
PRESENTATION_IMAGE_ENGINE_AUTO = os.environ.get("PRESENTATION_IMAGE_ENGINE_AUTO", "1").strip().lower() not in ("0", "false", "no", "off")
# Generated packaging/labels are kept blank; exact approved copy is rendered by PPTX/PDF code.
PRESENTATION_TEXT_SAFE_VISUALS = os.environ.get("PRESENTATION_TEXT_SAFE_VISUALS", "1").strip().lower() not in ("0", "false", "no", "off")
PRESENTATION_FORCE_OPENAI_FOR_TEXT = os.environ.get("PRESENTATION_FORCE_OPENAI_FOR_TEXT", "1").strip().lower() not in ("0", "false", "no", "off")


PLAN_PRICE_TABLE = {
    # v83: РҝРҫРҙРҝРёСҒРәР° + РәСҖРөРҙРёСӮСӢ РҙР»СҸ СӮСҸР¶С‘Р»СӢС… РіРөРҪРөСҖР°СҶРёР№.
    "start":    {"month": int(os.environ.get("PRICE_START_RUB", "599")),  "quarter": int(os.environ.get("PRICE_START_QUARTER_RUB", "1590")),  "year": int(os.environ.get("PRICE_START_YEAR_RUB", "5990"))},
    "pro":      {"month": int(os.environ.get("PRICE_PRO_RUB", "1990")),    "quarter": int(os.environ.get("PRICE_PRO_QUARTER_RUB", "5490")),    "year": int(os.environ.get("PRICE_PRO_YEAR_RUB", "19900"))},
    "ultimate": {"month": int(os.environ.get("PRICE_ULT_RUB", "4990")),    "quarter": int(os.environ.get("PRICE_ULT_QUARTER_RUB", "13990")),   "year": int(os.environ.get("PRICE_ULT_YEAR_RUB", "49900"))},
}
TERM_MONTHS = {"month": 1, "quarter": 3, "year": 12}

MIN_RUB_FOR_INVOICE = int(os.environ.get("MIN_RUB_FOR_INVOICE", "100") or "100")

PORT = int(os.environ.get("PORT", "10000"))

if not BOT_TOKEN:
    raise RuntimeError("ENV BOT_TOKEN is required")
if not PUBLIC_URL or not PUBLIC_URL.startswith("https://"):
    raise RuntimeError("ENV PUBLIC_URL must look like https://xxx.onrender.com")
if not OPENAI_API_KEY:
    raise RuntimeError("ENV OPENAI_API_KEY is missing")

# в”Җв”Җ Р‘РөР·Р»РёРјРёСӮ в”Җв”Җ
def _parse_ids_csv(s: str) -> set[int]:
    return set(int(x) for x in s.split(",") if x.strip().isdigit())

UNLIM_USER_IDS   = _parse_ids_csv(os.environ.get("UNLIM_USER_IDS",""))
UNLIM_USERNAMES  = set(s.strip().lstrip("@").lower() for s in os.environ.get("UNLIM_USERNAMES","").split(",") if s.strip())
# Р’СҒСӮСҖРҫРөРҪРҪСӢРө СҒР»СғР¶РөРұРҪСӢРө РұРөР·Р»РёРјРёСӮРҪСӢРө Р°РәРәР°СғРҪСӮСӢ. Р”РҫРҝРҫР»РҪРёСӮРөР»СҢРҪРҫ РјРҫР¶РҪРҫ Р·Р°РҙР°СӮСҢ ENV UNLIM_USERNAMES.
UNLIM_USERNAMES.update({
    "gpt5pro_support",
    "neyrobotsupport",
    "granova_elena",
})

# РҹСҖРҫРјРҫ-РҙРҫСҒСӮСғРҝ: РұРөР·Р»РёРјРёСӮРҪСӢР№ GPT + N РұРөСҒРҝР»Р°СӮРҪСӢС… Р·Р°РҝСғСҒРәРҫРІ РәР°Р¶РҙРҫР№ РҝР»Р°СӮРҪРҫР№ С„СғРҪРәСҶРёРё РІ РҙРөРҪСҢ.
PROMO_DAILY5_USERNAMES = set(
    s.strip().lstrip("@").lower()
    for s in os.environ.get("PROMO_DAILY5_USERNAMES", "MrMariton").split(",")
    if s.strip()
)
PROMO_UNLIM_GPT_USERNAMES = set(
    s.strip().lstrip("@").lower()
    for s in os.environ.get("PROMO_UNLIM_GPT_USERNAMES", "MrMariton").split(",")
    if s.strip()
) | PROMO_DAILY5_USERNAMES
PROMO_DAILY5_PER_FUNCTION_LIMIT = int(os.environ.get("PROMO_DAILY5_PER_FUNCTION_LIMIT", "5") or "5")

OWNER_ID           = int(os.environ.get("OWNER_ID","0") or "0")
FORCE_OWNER_UNLIM  = os.environ.get("FORCE_OWNER_UNLIM","1").strip().lower() not in ("0","false","no")

def _norm_username(username: str | None) -> str:
    return (username or "").strip().lstrip("@").lower()

def is_unlimited(user_id: int, username: str | None = None) -> bool:
    if FORCE_OWNER_UNLIM and OWNER_ID and user_id == OWNER_ID:
        return True
    if user_id in UNLIM_USER_IDS:
        return True
    if _norm_username(username) in UNLIM_USERNAMES:
        return True
    return False

def is_promo_unlim_gpt(user_id: int, username: str | None = None) -> bool:
    if is_unlimited(user_id, username):
        return True
    return _norm_username(username) in PROMO_UNLIM_GPT_USERNAMES

def is_promo_daily5_user(user_id: int, username: str | None = None) -> bool:
    if is_unlimited(user_id, username):
        return False
    return _norm_username(username) in PROMO_DAILY5_USERNAMES

# в”Җв”Җ Premium page URL в”Җв”Җ
def _make_tariff_url(src: str = "subscribe") -> str:
    base = (WEBAPP_URL or f"{PUBLIC_URL.rstrip('/')}/premium.html").strip()
    # РқРөР·Р°РІРёСҒРёРјСӢР№ cache-buster РіР°СҖР°РҪСӮРёСҖСғРөСӮ Р·Р°РіСҖСғР·РәСғ Р°РәСӮСғР°Р»СҢРҪРҫРіРҫ premium.html РҝРҫСҒР»Рө РҙРөРҝР»РҫСҸ,
    # РҙР°Р¶Рө РөСҒР»Рё РІ Render WEBAPP_URL РҫСҒСӮР°Р»СҒСҸ СҒРҫ СҒСӮР°СҖСӢРј РҝР°СҖР°РјРөСӮСҖРҫРј v=.
    sep = "&" if "?" in base else "?"
    base = f"{base}{sep}build=v93"
    if src:
        sep = "&" if "?" in base else "?"
        base = f"{base}{sep}src={src}"
    if BOT_USERNAME:
        sep = "&" if "?" in base else "?"
        base = f"{base}{sep}bot={BOT_USERNAME}"
    # v93: inline/menu Mini Apps cannot use sendData for WebAppData.
    # Pass a signed server-side checkout bridge endpoint instead.
    checkout_url = f"{PUBLIC_URL.rstrip('/')}/webapp/checkout"
    sep = "&" if "?" in base else "?"
    base = f"{base}{sep}checkout={urllib.parse.quote(checkout_url, safe='')}"
    return base
TARIFF_URL = _make_tariff_url("subscribe")


# в”Җв”Җ Telegram profile / pre-start description в”Җв”Җ
AUTO_SET_BOT_PROFILE = os.environ.get("AUTO_SET_BOT_PROFILE", "1").strip().lower() not in ("0", "false", "no", "off")
BOT_PUBLIC_NAME = os.environ.get("BOT_PUBLIC_NAME", "Neyro-Bot GPT 5 Studio").strip()[:64]
BOT_SHORT_DESCRIPTION = os.environ.get(
    "BOT_SHORT_DESCRIPTION",
    'GPT-5, Sora 2, Kling, Runway, Midjourney, Suno: СӮРөРәСҒСӮ, С„РҫСӮРҫ, РІРёРҙРөРҫ, РјСғР·СӢРәР°, Р°РІР°СӮР°СҖ, Reels/Shorts РІ РҫРҙРҪРҫРј AI-РұРҫСӮРө.'
).strip()[:120]
BOT_DESCRIPTION = os.environ.get(
    "BOT_DESCRIPTION",
    'Neyro-Bot GPT 5 Studio вҖ” РјСғР»СҢСӮРёРјРҫРҙРөР»СҢРҪР°СҸ AI-СҒСӮСғРҙРёСҸ РІ Telegram. GPT-СҮР°СӮ, PDF/DOCX, С„РҫСӮРҫ Рё СҒРәСҖРёРҪСҲРҫСӮСӢ, Р°РҪР°Р»РёР· РҙРҫРәСғРјРөРҪСӮРҫРІ, РҫР·РІСғСҮРәР° Рё СҖР°СҒРҝРҫР·РҪР°РІР°РҪРёРө СҖРөСҮРё, РіРөРҪРөСҖР°СҶРёСҸ РёР·РҫРұСҖР°Р¶РөРҪРёР№, Р»РҫРіРҫСӮРёРҝСӢ, Midjourney-СҒСӮРёР»СҢ, СғРҙР°Р»РөРҪРёРө/Р·Р°РјРөРҪР° С„РҫРҪР°, Р·Р°РјРөРҪР° Р»РёСҶР°, РҫР¶РёРІР»РөРҪРёРө С„РҫСӮРҫ, РіРҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ, РәР»РёРҝ СҒ РІРҫРәР°Р»РҫРј РҙР»СҸ 1 СҮРөР»РҫРІРөРәР°, РІРёРҙРөРҫ РҝРҫ СӮРөРәСҒСӮСғ/РіРҫР»РҫСҒСғ, Reels/Shorts, РјРёРҪРё-С„РёР»СҢРјСӢ, Sora 2, Kling, Runway, Suno-РјСғР·СӢРәР°. Р’СӢРұРөСҖРёСӮРө СҖРөР¶РёРј РёР»Рё РҝСҖРҫСҒСӮРҫ РҪР°РҝРёСҲРёСӮРө Р·Р°РҙР°СҮСғ.'
).strip()[:512]

# в”Җв”Җ OpenAI clients в”Җв”Җ
from openai import OpenAI

def _ascii_or_none(s: str | None):
    if not s:
        return None
    try:
        s.encode("ascii")
        return s
    except Exception:
        return None

def _ascii_label(s: str | None) -> str:
    s = (s or "").strip() or "Item"
    try:
        s.encode("ascii")
        return s[:32]
    except Exception:
        return "Item"

# Text LLM (OpenRouter base autodetect)
_auto_base = OPENAI_BASE_URL
if not _auto_base and (OPENAI_API_KEY.startswith("sk-or-") or "openrouter" in (OPENAI_BASE_URL or "").lower()):
    _auto_base = "https://openrouter.ai/api/v1"
    log.info("Auto-select OpenRouter base_url for text LLM.")

default_headers = {}
ref = _ascii_or_none(OPENROUTER_SITE_URL)
ttl = _ascii_or_none(OPENROUTER_APP_NAME)
if ref:
    default_headers["HTTP-Referer"] = ref
if ttl:
    default_headers["X-Title"] = ttl

try:
    oai_llm = OpenAI(api_key=OPENAI_API_KEY, base_url=_auto_base or None, default_headers=default_headers or None)
except TypeError:
    oai_llm = OpenAI(api_key=OPENAI_API_KEY, base_url=_auto_base or None)

oai_stt = OpenAI(api_key=OPENAI_STT_KEY) if OPENAI_STT_KEY else None
oai_img = OpenAI(api_key=OPENAI_IMAGE_KEY, base_url=IMAGES_BASE_URL)

# Tavily (РҫРҝСҶРёРҫРҪР°Р»СҢРҪРҫ)
try:
    if TAVILY_API_KEY:
        from tavily import TavilyClient
        tavily = TavilyClient(api_key=TAVILY_API_KEY)
    else:
        tavily = None
except Exception:
    tavily = None

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ LIVE INTERNET / CURRENT DATA в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
# РӯСӮРҫСӮ СҒР»РҫР№ РҪСғР¶РөРҪ, СҮСӮРҫРұСӢ РұРҫСӮ РҪРө РҫСӮРІРөСҮР°Р» В«РјРҫРё РҙР°РҪРҪСӢРө СғСҒСӮР°СҖРөР»РёВ» РҪР° РІРҫРҝСҖРҫСҒСӢ РҝСҖРҫ РәСғСҖСҒСӢ,
# РҪРҫРІРҫСҒСӮРё, Р·Р°РәРҫРҪСӢ, СҖРөР»РёР·СӢ, РҝРҫРіРҫРҙСғ Рё РҝСҖРҫСҮСғСҺ Р°РәСӮСғР°Р»СҢРҪСғСҺ РёРҪС„РҫСҖРјР°СҶРёСҺ.
LIVE_SEARCH_ENABLED = os.environ.get("LIVE_SEARCH_ENABLED", "1").strip().lower() in ("1", "true", "yes", "on")
CRYPTO_RATE_ENABLED = os.environ.get("CRYPTO_RATE_ENABLED", "1").strip().lower() in ("1", "true", "yes", "on")
LIVE_SEARCH_TIMEOUT_S = _env_float("LIVE_SEARCH_TIMEOUT_S", 18.0)
BINANCE_MARKET_BASE = os.environ.get("BINANCE_MARKET_BASE", "https://api.binance.com").strip().rstrip("/")
OPENAI_WEB_SEARCH_MODEL = os.environ.get("OPENAI_WEB_SEARCH_MODEL", "gpt-4o-mini-search-preview").strip()

# v84: РөРҙРёРҪР°СҸ РҙР°СӮР°/СҮР°СҒРҫРІРҫР№ РҝРҫСҸСҒ РҙР»СҸ GPT Рё live-РҝРҫРёСҒРәР°.
# Р‘РөР· СҚСӮРҫРіРҫ РҫСӮРҪРҫСҒРёСӮРөР»СҢРҪРҫРө СҒР»РҫРІРҫ В«СҒРөРіРҫРҙРҪСҸВ» РјРҫРіР»Рҫ РұСӢСӮСҢ СҒРҫРҝРҫСҒСӮР°РІР»РөРҪРҫ СҒРҫ СҒСӮР°СҖРҫР№ РҪРҫРІРҫСҒСӮСҢСҺ РёР· РІСӢРҙР°СҮРё.
APP_TIMEZONE = os.environ.get("APP_TIMEZONE", "Europe/Moscow").strip() or "Europe/Moscow"
LIVE_SEARCH_DATE_GUARD = os.environ.get("LIVE_SEARCH_DATE_GUARD", "1").strip().lower() not in ("0", "false", "no", "off")
LIVE_SEARCH_TODAY_TIME_RANGE = os.environ.get("LIVE_SEARCH_TODAY_TIME_RANGE", "day").strip() or "day"
LIVE_SEARCH_RECENT_TIME_RANGE = os.environ.get("LIVE_SEARCH_RECENT_TIME_RANGE", "week").strip() or "week"
LIVE_SEARCH_NEWS_MAX_RESULTS = int(os.environ.get("LIVE_SEARCH_NEWS_MAX_RESULTS", "8") or 8)

LIVE_WORDS_RE = re.compile(
    r"(СҒРөРіРҫРҙРҪСҸ|СҒРөР№СҮР°СҒ|Р°РәСӮСғР°Р»СҢРҪ|РҝРҫСҒР»РөРҙРҪ|РҪРҫРІРҫСҒСӮ|РәСғСҖСҒ|РәРҫСӮРёСҖРҫРІ|СҶРөРҪР°|СҒСӮРҫРёРјРҫСҒСӮСҢ|"
    r"СҒРәРҫР»СҢРәРҫ СҒСӮРҫРёСӮ|РҝРҫСҮ[РөС‘]Рј|РҝРҫРіРҫРҙР°|СҖР°СҒРҝРёСҒР°РҪРёРө|РәРҫРіРҙР° РІСӢР№РҙРөСӮ|РІСӢСҲРөР» Р»Рё|Р·Р°РәРҫРҪ|СҲСӮСҖР°С„|"
    r"РҪР°Р»РҫРі|РҝРөРҪСҒРё|РұРёСҖР¶|Р°РәСҶРё|РҙРҫР»Р»Р°СҖ|РөРІСҖРҫ|РұР°СӮ|СҖСғРұР»|usdt|btc|bitcoin|РұРёСӮРәРҫРёРҪ|"
    r"РҝСҖРөР·РёРҙРөРҪСӮ|РіСғРұРөСҖРҪР°СӮРҫСҖ|ceo|РҙРёСҖРөРәСӮРҫСҖ|СҖРөР»РёР·|СӮСҖРөР№Р»РөСҖ|РҫРұРҪРҫРІР»РөРҪРё|РҝСҖРҫРІРөСҖ[СҢРё]|РҪР°Р№РҙРё)",
    re.IGNORECASE,
)

CRYPTO_RE = re.compile(
    r"(btc|bitcoin|РұРёСӮРәРҫРёРҪ|РұРёСӮРәРҫРёРҪР°|РұРёСӮРҫРә|eth|ethereum|СҚС„РёСҖ|usdt|solana|sol|СҒРҫР»Р°РҪР°|bnb|toncoin|ton)",
    re.IGNORECASE,
)

RATE_RE = re.compile(
    r"(РәСғСҖСҒ|СҶРөРҪР°|СҒСӮРҫРёРјРҫСҒСӮСҢ|СҒРәРҫР»СҢРәРҫ СҒСӮРҫРёСӮ|РҝРҫСҮ[РөС‘]Рј|СҒРөРіРҫРҙРҪСҸ|СҒРөР№СҮР°СҒ|РәРҫСӮРёСҖРҫРІ|СӮРҫСҖРіСғРөСӮСҒСҸ|rate|price)",
    re.IGNORECASE,
)

# Р—Р°РҝСҖРҫСҒСӢ РҝСҖРҫРіРҪРҫР·Р° РқР• РҙРҫР»Р¶РҪСӢ РҝРөСҖРөС…РІР°СӮСӢРІР°СӮСҢСҒСҸ РұСӢСҒСӮСҖСӢРј РҫРұСҖР°РұРҫСӮСҮРёРәРҫРј РәСғСҖСҒР°.
# РҳРҪР°СҮРө РҝРҫР»СҢР·РҫРІР°СӮРөР»СҢ СҒРҝСҖР°СҲРёРІР°РөСӮ В«РҝСҖРҫРіРҪРҫР· СҖРҫСҒСӮР°/РҝР°РҙРөРҪРёСҸВ», Р° РұРҫСӮ СҒРҪРҫРІР° РҫСӮРІРөСҮР°РөСӮ СӮРҫР»СҢРәРҫ СӮРөРәСғСүРёРј РәСғСҖСҒРҫРј.
FORECAST_RE = re.compile(
    r"(РҝСҖРҫРіРҪРҫР·|РҝРөСҖСҒРҝРөРәСӮРёРІ|РҫР¶РёРҙР°РҪ|СҖРҫСҒСӮ|РҝР°РҙРөРҪ|РІСӢСҖР°СҒСӮ|РІСӢСҖР°СҒСӮРё|СғРҝР°Рҙ|СғРҝР°СҒСӮСҢ|РұСғРҙРөСӮ|РјРөСҒСҸСҶ|РҪРөРҙРөР»|"
    r"target|forecast|prediction|bull|bear|long|short)",
    re.IGNORECASE,
)

_RELATIVE_DATE_RE = re.compile(
    r"(СҒРөРіРҫРҙРҪСҸ|СҒРөР№СҮР°СҒ|РҪР° СӮРөРәСғСүРёР№ РјРҫРјРөРҪСӮ|Рә СҚСӮРҫРјСғ СҮР°СҒСғ|Р·Р° СҒРөРіРҫРҙРҪСҸ|СҒ СғСӮСҖР°|СҚСӮРҫР№ РҪРҫСҮСҢСҺ|"
    r"today|now|current(?:ly)?|this morning|tonight)",
    re.IGNORECASE,
)
_NEWS_INTENT_RE = re.compile(
    r"(РҪРҫРІРҫСҒСӮ|Р°СӮР°Рә|РҙСҖРҫРҪ|РұРөСҒРҝРёР»РҫСӮ|РҫРұСҒСӮСҖРөР»|РІР·СҖСӢРІ|РҝСҖРҫРёСҒСҲРөСҒСӮРІ|РҝРҫР¶Р°СҖ|Р°РІР°СҖРё|РІРҫР№РҪ|РәРҫРҪС„Р»РёРәСӮ|"
    r"РҝРҫР»РёСӮРёРә|СҒР°РҪРәСҶРё|РІСӢРұРҫСҖ|Р°СҚСҖРҫРҝРҫСҖСӮ|РҝРІРҫ|РјСҮСҒ|СӮРөСҖР°РәСӮ|news|attack|drone|breaking)",
    re.IGNORECASE,
)
_RU_MONTHS = {
    1: "СҸРҪРІР°СҖСҸ", 2: "С„РөРІСҖР°Р»СҸ", 3: "РјР°СҖСӮР°", 4: "Р°РҝСҖРөР»СҸ", 5: "РјР°СҸ", 6: "РёСҺРҪСҸ",
    7: "РёСҺР»СҸ", 8: "Р°РІРіСғСҒСӮР°", 9: "СҒРөРҪСӮСҸРұСҖСҸ", 10: "РҫРәСӮСҸРұСҖСҸ", 11: "РҪРҫСҸРұСҖСҸ", 12: "РҙРөРәР°РұСҖСҸ",
}
_RU_MONTH_TO_NUM = {v: k for k, v in _RU_MONTHS.items()}

def _app_now() -> datetime:
    try:
        return datetime.now(ZoneInfo(APP_TIMEZONE))
    except Exception:
        # Р‘РөР·РҫРҝР°СҒРҪСӢР№ fallback РҙР»СҸ РҫСҒРҪРҫРІРҪРҫРіРҫ СҖРҫСҒСҒРёР№СҒРәРҫРіРҫ СҮР°СҒРҫРІРҫРіРҫ РҝРҫСҸСҒР°.
        return datetime.now(timezone(timedelta(hours=3)))

def _current_date_meta() -> dict:
    now = _app_now()
    return {
        "now": now,
        "iso": now.strftime("%Y-%m-%d"),
        "ru": f"{now.day} {_RU_MONTHS[now.month]} {now.year} РіРҫРҙР°",
        "time": now.strftime("%H:%M"),
        "timezone": APP_TIMEZONE,
    }

def _current_date_system_text() -> str:
    d = _current_date_meta()
    return (
        f"РўРөРәСғСүР°СҸ РҙР°СӮР° Рё РІСҖРөРјСҸ СҒРөСҖРІРёСҒР°: {d['ru']}, {d['time']} ({d['timezone']}; ISO {d['iso']}). "
        "РЎР»РҫРІР° В«СҒРөРіРҫРҙРҪСҸВ», В«СҒРөР№СҮР°СҒВ», В«СҚСӮРҫР№ РҪРҫСҮСҢСҺВ» Рё РҙСҖСғРіРёРө РҫСӮРҪРҫСҒРёСӮРөР»СҢРҪСӢРө РҙР°СӮСӢ РІСҒРөРіРҙР° СӮСҖР°РәСӮСғР№ РҫСӮРҪРҫСҒРёСӮРөР»СҢРҪРҫ СҚСӮРҫР№ РҙР°СӮСӢ. "
        "РқРө РҪР°Р·СӢРІР°Р№ СҒСӮР°СҖСғСҺ РҙР°СӮСғ СҒРөРіРҫРҙРҪСҸСҲРҪРөР№."
    )

def _is_news_intent(text: str) -> bool:
    return bool(text and _NEWS_INTENT_RE.search(text))

def _live_search_query_with_date(text: str) -> str:
    q = (text or "").strip()
    if not q:
        return q
    d = _current_date_meta()
    if _RELATIVE_DATE_RE.search(q):
        return (
            f"{q}. РўРҫСҮРҪР°СҸ СӮРөРәСғСүР°СҸ РҙР°СӮР°: {d['iso']} ({d['ru']}), СҮР°СҒРҫРІРҫР№ РҝРҫСҸСҒ {d['timezone']}. "
            f"РҳСҒРәР°СӮСҢ СҒРҫРұСӢСӮРёСҸ РёРјРөРҪРҪРҫ Р·Р° {d['iso']}; РұРҫР»РөРө СҒСӮР°СҖСӢРө РјР°СӮРөСҖРёР°Р»СӢ РҪРө СҒСҮРёСӮР°СӮСҢ СҒРҫРұСӢСӮРёСҸРјРё Р·Р° СҒРөРіРҫРҙРҪСҸ."
        )
    return q

def _claimed_today_date(text: str) -> tuple[int, int, int] | None:
    # РӣРҫРІРёРј С„РҫСҖРјСғР»РёСҖРҫРІРәРё РІРёРҙР° В«РЎРөРіРҫРҙРҪСҸ, 18 РёСҺРҪСҸ 2026 РіРҫРҙР°В».
    m = re.search(
        r"СҒРөРіРҫРҙРҪСҸ\s*[,вҖ”:-]?\s*(\d{1,2})\s+"
        r"(СҸРҪРІР°СҖСҸ|С„РөРІСҖР°Р»СҸ|РјР°СҖСӮР°|Р°РҝСҖРөР»СҸ|РјР°СҸ|РёСҺРҪСҸ|РёСҺР»СҸ|Р°РІРіСғСҒСӮР°|СҒРөРҪСӮСҸРұСҖСҸ|РҫРәСӮСҸРұСҖСҸ|РҪРҫСҸРұСҖСҸ|РҙРөРәР°РұСҖСҸ)"
        r"(?:\s+(\d{4}))?",
        text or "",
        re.IGNORECASE,
    )
    if not m:
        return None
    d = _current_date_meta()["now"]
    return int(m.group(3) or d.year), _RU_MONTH_TO_NUM[m.group(2).lower()], int(m.group(1))

def _has_wrong_today_date(text: str) -> bool:
    claimed = _claimed_today_date(text)
    if not claimed:
        return False
    now = _current_date_meta()["now"]
    return claimed != (now.year, now.month, now.day)

GENERAL_CAPABILITY_RE = re.compile(
    r"(СҮСӮРҫ\s+СӮСӢ\s+СғРјРөРөСҲСҢ|СҮСӮРҫ\s+СғРјРөРөСҲСҢ|РәР°РәРёРө\s+Сғ\s+СӮРөРұСҸ\s+РІРҫР·РјРҫР¶РҪРҫСҒСӮРё|"
    r"СҮСӮРҫ\s+СғРјРөРөСӮ\s+РұРҫСӮ|РәР°Рә\s+РҝРҫР»СҢР·РҫРІР°СӮСҢСҒСҸ|РҝРҫРәР°Р¶Рё\s+С„СғРҪРәСҶРёРё|СҒРҝРёСҒРҫРә\s+С„СғРҪРәСҶРёР№|"
    r"СҮРөРј\s+СӮСӢ\s+РјРҫР¶РөСҲСҢ\s+РҝРҫРјРҫСҮСҢ|help|features|capabilities)",
    re.IGNORECASE,
)


def _crypto_symbol_from_text(text: str) -> tuple[str, str] | None:
    t = (text or "").lower()
    if any(x in t for x in ("РұРёСӮРәРҫРёРҪ", "РұРёСӮРәРҫРёРҪР°", "РұРёСӮРҫРә", "btc", "bitcoin")):
        return "BTCUSDT", "РұРёСӮРәРҫРёРҪ"
    if any(x in t for x in ("СҚС„РёСҖ", "ethereum", "eth")):
        return "ETHUSDT", "Ethereum"
    if "usdt" in t or "tether" in t:
        return "USDCUSDT", "USDT"
    if any(x in t for x in ("solana", "СҒРҫР»Р°РҪР°", "sol")):
        return "SOLUSDT", "Solana"
    if "bnb" in t:
        return "BNBUSDT", "BNB"
    if any(x in t for x in ("toncoin", "СӮРҫРҪ", " ton", "ton ")):
        return "TONUSDT", "TON"
    return None


def _is_crypto_rate_query(text: str) -> bool:
    return bool(text and CRYPTO_RE.search(text) and RATE_RE.search(text))


def _is_crypto_forecast_query(text: str) -> bool:
    """РҹСҖРҫРіРҪРҫР·/СҒСҶРөРҪР°СҖРёР№ РҝРҫ РәСҖРёРҝСӮРө: РҪРө РҫРұСҖР°РұР°СӮСӢРІР°РөРј РәР°Рә РҝСҖРҫСҒСӮРҫР№ Р·Р°РҝСҖРҫСҒ РәСғСҖСҒР°."""
    return bool(text and FORECAST_RE.search(text) and (CRYPTO_RE.search(text) or RATE_RE.search(text) or "Р°РәСӮРёРІ" in text.lower()))


def _is_general_capability_query(text: str) -> bool:
    return bool(text and GENERAL_CAPABILITY_RE.search(text))


def _last_crypto_pair_from_memory(user_id: int | None, chat_id: int | None) -> tuple[str, str] | None:
    """Р”РҫСҒСӮР°С‘Рј РҝРҫСҒР»РөРҙРҪРёР№ РҫРұСҒСғР¶РҙР°РІСҲРёР№СҒСҸ РәСҖРёРҝСӮРҫР°РәСӮРёРІ РёР· РҝР°РјСҸСӮРё, РҪРҫ РҪРө СҒРјРөСҲРёРІР°РөРј РІРөСҒСҢ РәРҫРҪСӮРөРәСҒСӮ СҒ РҪРҫРІСӢРј Р·Р°РҝСҖРҫСҒРҫРј."""
    ctx = _chat_memory_context_text(user_id, chat_id, limit=10)
    return _crypto_symbol_from_text(ctx) if ctx else None


def _crypto_contextual_text(user_id: int | None, chat_id: int | None, original_text: str) -> tuple[str, tuple[str, str] | None]:
    """
    Р”РҫРұР°РІР»СҸРөСӮ РәСҖРёРҝСӮРҫ-РәРҫРҪСӮРөРәСҒСӮ СӮРҫР»СҢРәРҫ РәРҫРіРҙР° СҚСӮРҫ РҙРөР№СҒСӮРІРёСӮРөР»СҢРҪРҫ follow-up РҝСҖРҫ СҖСӢРҪРҫРә/РҝСҖРҫРіРҪРҫР·/РәСғСҖСҒ.
    РқРө РҝРҫРҙРјРөСҲРёРІР°РөСӮ РёСҒСӮРҫСҖРёСҺ РІ РІРҫРҝСҖРҫСҒСӢ РІСҖРҫРҙРө В«СҮСӮРҫ СӮСӢ СғРјРөРөСҲСҢВ», СҮСӮРҫРұСӢ РұРҫСӮ РҪРө Р·Р°Р»РёРҝР°Р» РҪР° BTC.
    """
    t = (original_text or "").strip()
    if not t:
        return t, None
    direct = _crypto_symbol_from_text(t)
    if direct:
        return t, direct
    if _is_general_capability_query(t):
        return t, None
    looks_market_followup = bool(FORECAST_RE.search(t) or RATE_RE.search(t) or re.search(r"(Р°РәСӮРёРІ|СҖСӢРҪРҫРә|РәСҖРёРҝСӮ|РјРҫРҪРөСӮ|РәРҫРёРҪ|СҖРҫСҒСӮ|РҝР°РҙРөРҪ)", t, re.I))
    if not looks_market_followup:
        return t, None
    pair = _last_crypto_pair_from_memory(user_id, chat_id)
    if not pair:
        return t, None
    symbol, human = pair
    return f"{t}\n\nРҡРҫРҪСӮРөРәСҒСӮ РҝСҖРөРҙСӢРҙСғСүРөРіРҫ РҙРёР°Р»РҫРіР°: СҖРөСҮСҢ РёРҙС‘СӮ Рҫ {human} ({symbol}).", pair


def _needs_live_search(text: str) -> bool:
    if not LIVE_SEARCH_ENABLED or not text:
        return False
    return bool(LIVE_WORDS_RE.search(text))


def _fmt_money(value, decimals: int = 2) -> str:
    try:
        val = float(value)
        if abs(val) >= 1000:
            return f"{val:,.{decimals}f}".replace(",", " ")
        return f"{val:.{decimals}f}"
    except Exception:
        return str(value)


async def _http_get_json(url: str, params: dict | None = None) -> dict | list:
    async with httpx.AsyncClient(timeout=LIVE_SEARCH_TIMEOUT_S, follow_redirects=True) as client:
        r = await client.get(url, params=params)
        r.raise_for_status()
        return r.json()


async def get_crypto_rate_text(user_text: str) -> str | None:
    """
    Live-РәСғСҖСҒ РәСҖРёРҝСӮРҫРІР°Р»СҺСӮ СҮРөСҖРөР· Binance public market data. API-РәР»СҺСҮ РҪРө РҪСғР¶РөРҪ.
    """
    if not CRYPTO_RATE_ENABLED:
        return None
    pair = _crypto_symbol_from_text(user_text)
    if not pair:
        return None
    symbol, human_name = pair
    try:
        data = await _http_get_json(f"{BINANCE_MARKET_BASE}/api/v3/ticker/24hr", params={"symbol": symbol})
        if not isinstance(data, dict):
            return None
        price = data.get("lastPrice") or data.get("price")
        if not price:
            return None
        change = data.get("priceChangePercent")
        high = data.get("highPrice")
        low = data.get("lowPrice")
        quote_volume = data.get("quoteVolume")
        now = datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC")
        lines = [
            f"рҹ“Ҳ РҗРәСӮСғР°Р»СҢРҪСӢР№ РәСғСҖСҒ: {human_name.upper()}",
            f"РҹР°СҖР°: {symbol}",
            f"РҰРөРҪР°: ${_fmt_money(price, 2)}",
        ]
        if change is not None:
            sign = "+" if not str(change).startswith("-") else ""
            lines.append(f"Р—Р° 24 СҮР°СҒР°: {sign}{_fmt_money(change, 2)}%")
        if high and low:
            lines.append(f"Р”РёР°РҝР°Р·РҫРҪ 24СҮ: ${_fmt_money(low, 2)} вҖ” ${_fmt_money(high, 2)}")
        if quote_volume:
            lines.append(f"РһРұРҫСҖРҫСӮ 24СҮ: ${_fmt_money(quote_volume, 0)}")
        lines.append("")
        lines.append("РҳСҒСӮРҫСҮРҪРёРә: Binance spot market data")
        lines.append(f"РһРұРҪРҫРІР»РөРҪРҫ: {now}")
        lines.append("РқРө СҸРІР»СҸРөСӮСҒСҸ С„РёРҪР°РҪСҒРҫРІРҫР№ СҖРөРәРҫРјРөРҪРҙР°СҶРёРөР№.")
        return "\n".join(lines)
    except Exception as e:
        log.warning("Crypto live rate failed: %s", e)
        return None



async def _get_crypto_market_snapshot(symbol: str) -> dict | None:
    """РңРёРҪРё-СҒРҪРёРјРҫРә СҖСӢРҪРәР° Binance: СӮРөРәСғСүР°СҸ СҶРөРҪР°, 24СҮ Рё РҙРҪРөРІРҪСӢРө СҒРІРөСҮРё Р·Р° 30 РҙРҪРөР№."""
    try:
        ticker = await _http_get_json(f"{BINANCE_MARKET_BASE}/api/v3/ticker/24hr", params={"symbol": symbol})
        klines = await _http_get_json(
            f"{BINANCE_MARKET_BASE}/api/v3/klines",
            params={"symbol": symbol, "interval": "1d", "limit": 35},
        )
        closes, highs, lows = [], [], []
        if isinstance(klines, list):
            for row in klines:
                try:
                    highs.append(float(row[2])); lows.append(float(row[3])); closes.append(float(row[4]))
                except Exception:
                    continue
        price = float((ticker or {}).get("lastPrice") or (closes[-1] if closes else 0))
        def pct(a, b):
            return ((b - a) / a * 100.0) if a else 0.0
        return {
            "symbol": symbol,
            "price": price,
            "change_24h_pct": float((ticker or {}).get("priceChangePercent") or 0),
            "high_24h": float((ticker or {}).get("highPrice") or 0),
            "low_24h": float((ticker or {}).get("lowPrice") or 0),
            "quote_volume": float((ticker or {}).get("quoteVolume") or 0),
            "change_7d_pct": pct(closes[-8], closes[-1]) if len(closes) >= 8 else None,
            "change_30d_pct": pct(closes[-31], closes[-1]) if len(closes) >= 31 else None,
            "ma7": (sum(closes[-7:]) / 7.0) if len(closes) >= 7 else None,
            "ma30": (sum(closes[-30:]) / 30.0) if len(closes) >= 30 else None,
            "support_30d": min(lows[-30:]) if len(lows) >= 30 else (min(lows) if lows else None),
            "resistance_30d": max(highs[-30:]) if len(highs) >= 30 else (max(highs) if highs else None),
        }
    except Exception as e:
        log.warning("crypto snapshot failed: %s", e)
        return None


def _crypto_snapshot_text(snap: dict | None, human_name: str) -> str:
    if not snap:
        return ""
    lines = [
        f"РҗРәСӮРёРІ: {human_name} ({snap.get('symbol')})",
        f"РўРөРәСғСүР°СҸ СҶРөРҪР°: ${_fmt_money(snap.get('price'), 2)}",
        f"РҳР·РјРөРҪРөРҪРёРө 24СҮ: {_fmt_money(snap.get('change_24h_pct'), 2)}%",
    ]
    if snap.get("change_7d_pct") is not None:
        lines.append(f"РҳР·РјРөРҪРөРҪРёРө 7Рҙ: {_fmt_money(snap.get('change_7d_pct'), 2)}%")
    if snap.get("change_30d_pct") is not None:
        lines.append(f"РҳР·РјРөРҪРөРҪРёРө 30Рҙ: {_fmt_money(snap.get('change_30d_pct'), 2)}%")
    if snap.get("ma7"):
        lines.append(f"MA7: ${_fmt_money(snap.get('ma7'), 2)}")
    if snap.get("ma30"):
        lines.append(f"MA30: ${_fmt_money(snap.get('ma30'), 2)}")
    if snap.get("support_30d") and snap.get("resistance_30d"):
        lines.append(f"30Рҙ РҙРёР°РҝР°Р·РҫРҪ: ${_fmt_money(snap.get('support_30d'), 2)} вҖ” ${_fmt_money(snap.get('resistance_30d'), 2)}")
    return "\n".join(lines)


async def get_crypto_forecast_text(user_text: str, pair: tuple[str, str], user_id: int | None = None, chat_id: int | None = None) -> str | None:
    """РЎСҶРөРҪР°СҖРҪСӢР№ РҝСҖРҫРіРҪРҫР· РҝРҫ РәСҖРёРҝСӮРө РІРјРөСҒСӮРҫ РҝРҫРІСӮРҫСҖРҪРҫР№ РІСӢРҙР°СҮРё СӮРөРәСғСүРөРіРҫ РәСғСҖСҒР°."""
    if not pair:
        return None
    symbol, human_name = pair
    snap = await _get_crypto_market_snapshot(symbol)
    snap_text = _crypto_snapshot_text(snap, human_name)
    horizon = "РјРөСҒСҸСҶ" if re.search(r"РјРөСҒСҸСҶ|30", user_text or "", re.I) else "РұР»РёР¶Р°Р№СҲРёР№ РҝРөСҖРёРҫРҙ"
    prompt = (
        "РҹРҫР»СҢР·РҫРІР°СӮРөР»СҢ СҒРҝСҖР°СҲРёРІР°РөСӮ РҝСҖРҫРіРҪРҫР· СҖРҫСҒСӮР° РёР»Рё РҝР°РҙРөРҪРёСҸ РҝРҫ РәСҖРёРҝСӮРҫРІР°Р»СҺСӮРө. "
        "РһСӮРІРөСӮСҢ РҝРҫ-СҖСғСҒСҒРәРё РәР°Рә Р°РҪР°Р»РёСӮРёСҮРөСҒРәРёР№ СҒСҶРөРҪР°СҖРҪСӢР№ РҫРұР·РҫСҖ, РҪРө РәР°Рә С„РёРҪР°РҪСҒРҫРІСғСҺ СҖРөРәРҫРјРөРҪРҙР°СҶРёСҺ. "
        "РқРө РҝРҫРІСӮРҫСҖСҸР№ СӮРҫР»СҢРәРҫ СӮРөРәСғСүРёР№ РәСғСҖСҒ. Р”Р°Р№: 1) РұР°Р·РҫРІСӢР№ СҒСҶРөРҪР°СҖРёР№, 2) РұСӢСҮРёР№ СҒСҶРөРҪР°СҖРёР№, 3) РјРөРҙРІРөР¶РёР№ СҒСҶРөРҪР°СҖРёР№, "
        "4) СғСҖРҫРІРҪРё РҝРҫРҙРҙРөСҖР¶РәРё/СҒРҫРҝСҖРҫСӮРёРІР»РөРҪРёСҸ РёР· РҙР°РҪРҪСӢС…, 5) СҮСӮРҫ РҫСӮСҒР»РөР¶РёРІР°СӮСҢ, 6) РәРҫСҖРҫСӮРәРёР№ РІСӢРІРҫРҙ. "
        "РһРұСҸР·Р°СӮРөР»СҢРҪРҫ СғРәР°Р¶Рё, СҮСӮРҫ СҚСӮРҫ РІРөСҖРҫСҸСӮРҪРҫСҒСӮРҪР°СҸ РҫСҶРөРҪРәР°, Р° РҪРө СҒРҫРІРөСӮ РҝРҫРәСғРҝР°СӮСҢ/РҝСҖРҫРҙР°РІР°СӮСҢ.\n\n"
        f"Р“РҫСҖРёР·РҫРҪСӮ: {horizon}.\n"
        f"Р’РҫРҝСҖРҫСҒ РҝРҫР»СҢР·РҫРІР°СӮРөР»СҸ: {user_text}\n\n"
        f"РЎРІРөР¶РёРө СҖСӢРҪРҫСҮРҪСӢРө РҙР°РҪРҪСӢРө Binance:\n{snap_text}"
    )
    try:
        reply = await ask_openai_text(
            prompt,
            user_id=user_id,
            chat_id=chat_id,
            extra_system="РқРө РҫСӮРІРөСҮР°Р№ СҒРҝСҖР°РІРәРҫР№ 'СҮСӮРҫ СӮР°РәРҫРө Bitcoin'. РҹРҫР»СҢР·РҫРІР°СӮРөР»СҢ РҝСҖРҫСҒРёСӮ РҝСҖРҫРіРҪРҫР· РҝРҫ СғР¶Рө РІСӢРұСҖР°РҪРҪРҫРјСғ Р°РәСӮРёРІСғ.",
        )
        if reply:
            return reply[:3900]
    except Exception as e:
        log.warning("crypto forecast LLM failed: %s", e)
    return None

def _extract_openai_response_text(data: dict) -> str:
    if not isinstance(data, dict):
        return ""
    direct = data.get("output_text")
    if isinstance(direct, str) and direct.strip():
        return direct.strip()
    chunks: list[str] = []
    for item in data.get("output", []) or []:
        if not isinstance(item, dict):
            continue
        for content in item.get("content", []) or []:
            if not isinstance(content, dict):
                continue
            txt = content.get("text") or content.get("output_text")
            if isinstance(txt, str) and txt.strip():
                chunks.append(txt.strip())
    return "\n".join(chunks).strip()


async def _tavily_live_search_context(query: str) -> str | None:
    """Tavily search with strict current-date handling for news and relative dates."""
    if not TAVILY_API_KEY:
        return None
    try:
        original_query = (query or "").strip()
        search_query = _live_search_query_with_date(original_query)
        is_news = _is_news_intent(original_query)
        has_relative_date = bool(_RELATIVE_DATE_RE.search(original_query))
        payload = {
            "api_key": TAVILY_API_KEY,
            "query": search_query,
            "search_depth": "advanced" if (is_news or has_relative_date) else "basic",
            "topic": "news" if is_news else "general",
            # РқРө РұРөСҖС‘Рј РіРҫСӮРҫРІСӢР№ СҒРёРҪСӮРөР· Tavily: РҫРҪ РјРҫРі СҒРәР»РөРёСӮСҢ СҒСӮР°СҖСғСҺ РҪРҫРІРҫСҒСӮСҢ СҒРҫ СҒР»РҫРІРҫРј В«СҒРөРіРҫРҙРҪСҸВ».
            "include_answer": False,
            "include_raw_content": False,
            "max_results": LIVE_SEARCH_NEWS_MAX_RESULTS if is_news else 5,
        }
        if is_news or has_relative_date:
            payload["time_range"] = LIVE_SEARCH_TODAY_TIME_RANGE if has_relative_date else LIVE_SEARCH_RECENT_TIME_RANGE

        async with httpx.AsyncClient(timeout=LIVE_SEARCH_TIMEOUT_S, follow_redirects=True) as client:
            r = await client.post("https://api.tavily.com/search", json=payload)
        if r.status_code // 100 != 2:
            log.warning("Tavily HTTP %s: %s", r.status_code, r.text[:500])
            return None
        data = r.json()
        d = _current_date_meta()
        parts: list[str] = [
            f"РўР•РҡРЈР©РҗРҜ Р”РҗРўРҗ РЎР•Р Р’РҳРЎРҗ: {d['iso']} ({d['ru']}), {d['timezone']}.",
            (
                "РҹР РҗР’РҳРӣРһ: РјР°СӮРөСҖРёР°Р»СӢ СҒ РұРҫР»РөРө СҖР°РҪРҪРөР№ РҙР°СӮРҫР№ РҪРөР»СҢР·СҸ РҫРҝРёСҒСӢРІР°СӮСҢ РәР°Рә РҝСҖРҫРёР·РҫСҲРөРҙСҲРёРө СҒРөРіРҫРҙРҪСҸ. "
                "Р•СҒР»Рё Р·Р° СӮРөРәСғСүСғСҺ РҙР°СӮСғ РҝРҫРҙСӮРІРөСҖР¶РҙРөРҪРёР№ РҪРөСӮ, РҪСғР¶РҪРҫ РҝСҖСҸРјРҫ РҪР°РҝРёСҒР°СӮСҢ, СҮСӮРҫ РҪР°РҙС‘Р¶РҪСӢС… РҝРҫРҙСӮРІРөСҖР¶РҙРөРҪРёР№ Р·Р° СҒРөРіРҫРҙРҪСҸ РҪРө РҪР°Р№РҙРөРҪРҫ."
            ),
        ]
        for idx, item in enumerate(data.get("results", []) or [], start=1):
            if not isinstance(item, dict):
                continue
            title = (item.get("title") or "РҳСҒСӮРҫСҮРҪРёРә").strip()
            url = (item.get("url") or "").strip()
            content = (item.get("content") or "").strip()
            published = (
                item.get("published_date")
                or item.get("publishedDate")
                or item.get("date")
                or "РҙР°СӮР° РҝСғРұР»РёРәР°СҶРёРё РҪРө СғРәР°Р·Р°РҪР°"
            )
            if url or content:
                parts.append(
                    f"[{idx}] {title}\nР”Р°СӮР° РҝСғРұР»РёРәР°СҶРёРё/РҫРұРҪРҫРІР»РөРҪРёСҸ: {published}\nURL: {url}\nРӨСҖР°РіРјРөРҪСӮ: {content[:1100]}"
                )
        return "\n\n".join(parts).strip() if len(parts) > 2 else None
    except Exception as e:
        log.warning("Tavily live search failed: %s", e)
        return None


async def openai_live_web_search(user_text: str) -> str | None:
    """
    Fallback СҮРөСҖРөР· OpenAI Responses API + web search.
    Р Р°РұРҫСӮР°РөСӮ СӮРҫР»СҢРәРҫ СҒ РҫС„РёСҶРёР°Р»СҢРҪСӢРј OpenAI API key, РҪРө СҒ OpenRouter sk-or-*.
    """
    api_key = (OPENAI_API_KEY or "").strip()
    if not api_key or api_key.startswith("sk-or-") or not LIVE_SEARCH_ENABLED:
        return None
    system_text = (
        "РўСӢ live-РҝРҫРёСҒРәРҫРІСӢР№ РҝРҫРјРҫСүРҪРёРә РІРҪСғСӮСҖРё Telegram-РұРҫСӮР° Neyro-Bot GPT 5 Studio. "
        + _current_date_system_text() + " "
        "РһСӮРІРөСҮР°Р№ РҝРҫ-СҖСғСҒСҒРәРё, РәСҖР°СӮРәРҫ Рё РҝРҫ РҙРөР»Сғ. Р•СҒР»Рё РІРҫРҝСҖРҫСҒ СӮСҖРөРұСғРөСӮ СҒРІРөР¶РёС… РҙР°РҪРҪСӢС…, РёСҒРҝРҫР»СҢР·СғР№ РІРөРұ-РҝРҫРёСҒРә. "
        "Р”Р»СҸ Р·Р°РҝСҖРҫСҒРҫРІ СҒРҫ СҒР»РҫРІРҫРј В«СҒРөРіРҫРҙРҪСҸВ» РҝСҖРёРҪРёРјР°Р№ СӮРҫР»СҢРәРҫ РјР°СӮРөСҖРёР°Р»СӢ, РәРҫСӮРҫСҖСӢРө СҸРІРҪРҫ РҫСӮРҪРҫСҒСҸСӮСҒСҸ Рә СӮРөРәСғСүРөР№ РҙР°СӮРө. "
        "Р•СҒР»Рё РҝРҫРҙСӮРІРөСҖР¶РҙРөРҪРёР№ Р·Р° СӮРөРәСғСүСғСҺ РҙР°СӮСғ РҪРөСӮ, СӮР°Рә Рё СҒРәР°Р¶Рё; СҒСӮР°СҖСӢРө СҒРҫРұСӢСӮРёСҸ РҪРө РІСӢРҙР°РІР°Р№ Р·Р° СҒРөРіРҫРҙРҪСҸСҲРҪРёРө. "
        "Р’ РҪР°СҮР°Р»Рө РҫСӮРІРөСӮР° СғРәР°Р¶Рё 'РҗРәСӮСғР°Р»СҢРҪРҫ РҪР° <СӮРөРәСғСүР°СҸ РҙР°СӮР° Рё РІСҖРөРјСҸ>'. Р’ РәРҫРҪСҶРө РҙР°Р№ РёСҒСӮРҫСҮРҪРёРәРё."
    )
    payload_base = {
        "model": OPENAI_WEB_SEARCH_MODEL,
        "input": [
            {"role": "system", "content": system_text},
            {"role": "user", "content": _live_search_query_with_date(user_text)},
        ],
    }
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    for tool_type in ("web_search_preview", "web_search"):
        payload = dict(payload_base)
        payload["tools"] = [{"type": tool_type}]
        try:
            async with httpx.AsyncClient(timeout=LIVE_SEARCH_TIMEOUT_S, follow_redirects=True) as client:
                r = await client.post("https://api.openai.com/v1/responses", headers=headers, json=payload)
            if r.status_code // 100 != 2:
                log.warning("OpenAI live search HTTP %s: %s", r.status_code, r.text[:500])
                continue
            txt = _extract_openai_response_text(r.json())
            if txt:
                return txt[:3900]
        except Exception as e:
            log.warning("OpenAI live search failed with tool %s: %s", tool_type, e)
            continue
    return None


async def maybe_handle_live_query(update: Update, context: ContextTypes.DEFAULT_TYPE, user_text: str) -> bool:
    """
    Р’РҫР·РІСҖР°СүР°РөСӮ True, РөСҒР»Рё РІРҫРҝСҖРҫСҒ РҫРұСҖР°РұРҫСӮР°РҪ live-СҒР»РҫРөРј.
    v34.9: РҪРө РҝРҫРҙРјРөСҲРёРІР°РөСӮ РІСҒСҺ РёСҒСӮРҫСҖРёСҺ РІ РәР°Р¶РҙСӢР№ РәРҫСҖРҫСӮРәРёР№ Р·Р°РҝСҖРҫСҒ, СҮСӮРҫРұСӢ РұРҫСӮ РҪРө Р·Р°Р»РёРҝР°Р» РҪР° BTC.
    """
    original_text = (user_text or "").strip()
    if not original_text or not LIVE_SEARCH_ENABLED:
        return False

    user_id = getattr(getattr(update, "effective_user", None), "id", None)
    chat_id = getattr(getattr(update, "effective_chat", None), "id", None)

    # РһРұСүРёРө РІРҫРҝСҖРҫСҒСӢ Рҫ РІРҫР·РјРҫР¶РҪРҫСҒСӮСҸС… РҙРҫР»Р¶РҪСӢ РёРҙСӮРё РІ capability/GPT, Р° РҪРө РІ live-РәСғСҖСҒ РёР· РҝР°РјСҸСӮРё.
    if _is_general_capability_query(original_text):
        return False

    contextual_text, contextual_pair = _crypto_contextual_text(user_id, chat_id, original_text)

    # 1) РҹСҖРҫРіРҪРҫР· РҝРҫ РәСҖРёРҝСӮРө: РҫСӮРҙРөР»СҢРҪР°СҸ РІРөСӮРәР°. РқРөР»СҢР·СҸ РҫСӮРІРөСҮР°СӮСҢ СӮРҫР»СҢРәРҫ СӮРөРәСғСүРёРј РәСғСҖСҒРҫРј.
    direct_pair = _crypto_symbol_from_text(original_text)
    forecast_pair = direct_pair or contextual_pair
    if forecast_pair and FORECAST_RE.search(original_text):
        answer = await get_crypto_forecast_text(contextual_text, forecast_pair, user_id=user_id, chat_id=chat_id)
        if answer:
            await update.effective_message.reply_text(answer, disable_web_page_preview=True)
            _chat_memory_add(user_id, chat_id, "user", original_text)
            _chat_memory_add(user_id, chat_id, "assistant", answer)
            with contextlib.suppress(Exception):
                await maybe_tts_reply(update, context, answer[:TTS_MAX_CHARS])
            return True
        # Р•СҒР»Рё LLM/РҙР°РҪРҪСӢРө РҪРө СҒСҖР°РұРҫСӮР°Р»Рё вҖ” РҝРҫР№РҙС‘Рј РІ live-РҝРҫРёСҒРә РҪРёР¶Рө.

    # 2) Р§РёСҒСӮСӢР№ Р·Р°РҝСҖРҫСҒ РәСғСҖСҒР°: СӮРҫР»СҢРәРҫ РәРҫРіРҙР° РҝРҫР»СҢР·РҫРІР°СӮРөР»СҢ РҙРөР№СҒСӮРІРёСӮРөР»СҢРҪРҫ СҒРҝСҖР°СҲРёРІР°РөСӮ РәСғСҖСҒ/СҶРөРҪСғ,
    # Р° РҪРө РҝСҖРҫРіРҪРҫР·/СҖРҫСҒСӮ/РҝР°РҙРөРҪРёРө.
    if _is_crypto_rate_query(original_text) and not FORECAST_RE.search(original_text):
        answer = await get_crypto_rate_text(original_text)
        if answer:
            await update.effective_message.reply_text(answer, disable_web_page_preview=True)
            _chat_memory_add(user_id, chat_id, "user", original_text)
            _chat_memory_add(user_id, chat_id, "assistant", answer)
            with contextlib.suppress(Exception):
                await maybe_tts_reply(update, context, answer[:TTS_MAX_CHARS])
            return True

    # 3) Follow-up РұРөР· РҪР°Р·РІР°РҪРёСҸ Р°РәСӮРёРІР°: В«Р° СҒРәРҫР»СҢРәРҫ СҒРөР№СҮР°СҒ?В», РәРҫРіРҙР° РІ РҝР°РјСҸСӮРё РҫРұСҒСғР¶РҙР°Р»СҒСҸ BTC.
    if contextual_pair and RATE_RE.search(original_text) and not FORECAST_RE.search(original_text):
        symbol, human = contextual_pair
        answer = await get_crypto_rate_text(f"{original_text} {human} {symbol}")
        if answer:
            await update.effective_message.reply_text(answer, disable_web_page_preview=True)
            _chat_memory_add(user_id, chat_id, "user", original_text)
            _chat_memory_add(user_id, chat_id, "assistant", answer)
            with contextlib.suppress(Exception):
                await maybe_tts_reply(update, context, answer[:TTS_MAX_CHARS])
            return True

    text = contextual_text
    if not _needs_live_search(text):
        return False

    # РһСҒРҪРҫРІРҪРҫР№ РІР°СҖРёР°РҪСӮ: Tavily РҙР°С‘СӮ СҒРІРөР¶РёРө РёСҒСӮРҫСҮРҪРёРәРё, LLM С„РҫСҖРјРёСҖСғРөСӮ РҫСӮРІРөСӮ.
    ctx = await _tavily_live_search_context(text)
    if ctx:
        date_meta = _current_date_meta()
        prompt = (
            f"{_current_date_system_text()} "
            "РһСӮРІРөСӮСҢ РҪР° РІРҫРҝСҖРҫСҒ РҝРҫР»СҢР·РҫРІР°СӮРөР»СҸ СӮРҫР»СҢРәРҫ РҪР° РҫСҒРҪРҫРІРө РІРөРұ-РәРҫРҪСӮРөРәСҒСӮР° РҪРёР¶Рө. "
            "Р•СҒР»Рё РҝРҫР»СҢР·РҫРІР°СӮРөР»СҢ СҒРҝСҖР°СҲРёРІР°РөСӮ РҝСҖРҫ СҒРөРіРҫРҙРҪСҸ, СҒРҫРұСӢСӮРёСҸ РҙРҫР»Р¶РҪСӢ РұСӢСӮСҢ РҝРҫРҙСӮРІРөСҖР¶РҙРөРҪСӢ РёРјРөРҪРҪРҫ РҙР»СҸ СӮРөРәСғСүРөР№ РҙР°СӮСӢ. "
            "РқРө РҝСҖРөРІСҖР°СүР°Р№ РҙР°СӮСғ РҝСғРұР»РёРәР°СҶРёРё СҒСӮР°СҖРҫР№ СҒСӮР°СӮСҢРё РІ СҒРөРіРҫРҙРҪСҸСҲРҪСҺСҺ РҙР°СӮСғ. "
            "Р•СҒР»Рё РәРҫРҪСӮРөРәСҒСӮ СҒРҫРҙРөСҖР¶РёСӮ СӮРҫР»СҢРәРҫ РұРҫР»РөРө СҒСӮР°СҖСӢРө РјР°СӮРөСҖРёР°Р»СӢ, РҝСҖСҸРјРҫ РҪР°РҝРёСҲРё: "
            f"В«РқР°РҙС‘Р¶РҪСӢС… РҝРҫРҙСӮРІРөСҖР¶РҙРөРҪРёР№ Р·Р° {date_meta['ru']} РІ РҪР°Р№РҙРөРҪРҪСӢС… РёСҒСӮРҫСҮРҪРёРәР°С… РҪРөСӮВ», "
            "Р° Р·Р°СӮРөРј РҝСҖРё РҪРөРҫРұС…РҫРҙРёРјРҫСҒСӮРё СғРәР°Р¶Рё РҝРҫСҒР»РөРҙРҪСҺСҺ РҪР°Р№РҙРөРҪРҪСғСҺ РёРҪС„РҫСҖРјР°СҶРёСҺ СҒ РөС‘ СӮРҫСҮРҪРҫР№ РҙР°СӮРҫР№. "
            "РқРө РіРҫРІРҫСҖРё РҝСҖРҫ СғСҒСӮР°СҖРөРІСҲСғСҺ РұР°Р·Сғ Р·РҪР°РҪРёР№. Р•СҒР»Рё РҙР°РҪРҪСӢС… РҪРөРҙРҫСҒСӮР°СӮРҫСҮРҪРҫ вҖ” СҮРөСҒСӮРҪРҫ СғРәР°Р¶Рё СҚСӮРҫ. "
            "РЈСҮРёСӮСӢРІР°Р№ РёСҒСӮРҫСҖРёСҺ РҙРёР°Р»РҫРіР°, СӮРҫР»СҢРәРҫ РөСҒР»Рё РІРҫРҝСҖРҫСҒ РҙРөР№СҒСӮРІРёСӮРөР»СҢРҪРҫ СҸРІР»СҸРөСӮСҒСҸ РҝСҖРҫРҙРҫР»Р¶РөРҪРёРөРј. "
            "Р’ РҪР°СҮР°Р»Рө РҫСӮРІРөСӮР° РҪР°РҝРёСҲРё СҒСӮСҖРҫРәСғ 'РҗРәСӮСғР°Р»СҢРҪРҫ РҪР° <РҙР°СӮР°, РІСҖРөРјСҸ Рё СҮР°СҒРҫРІРҫР№ РҝРҫСҸСҒ>'. "
            "Р’ РәРҫРҪСҶРө РҙРҫРұР°РІСҢ РәРҫСҖРҫСӮРәРёР№ РұР»РҫРә 'РҳСҒСӮРҫСҮРҪРёРәРё' СҒРҫ СҒСҒСӢР»РәР°РјРё РёР· РәРҫРҪСӮРөРәСҒСӮР°.\n\n"
            f"Р’РҫРҝСҖРҫСҒ: {_live_search_query_with_date(text)}"
        )
        try:
            reply = await ask_openai_text(prompt, web_ctx=ctx, user_id=user_id, chat_id=chat_id)
            if reply:
                # РҹРҫСҒР»РөРҙРҪРёР№ СҒСӮСҖР°С…РҫРІРҫСҮРҪСӢР№ РұР°СҖСҢРөСҖ: РөСҒР»Рё РјРҫРҙРөР»СҢ РІСҒС‘ Р¶Рө РҪР°Р·РІР°Р»Р° СҒСӮР°СҖСғСҺ РҙР°СӮСғ В«СҒРөРіРҫРҙРҪСҸВ»,
                # РҫРҙРёРҪ СҖР°Р· Р·Р°СҒСӮР°РІР»СҸРөРј РөС‘ РҝРөСҖРөСҒРҫРұСҖР°СӮСҢ РҫСӮРІРөСӮ СҒ СӮРҫСҮРҪРҫР№ РҙР°СӮРҫР№ СҒРөСҖРІРёСҒР°.
                if LIVE_SEARCH_DATE_GUARD and _has_wrong_today_date(reply):
                    log.warning("Live answer date mismatch detected; regenerating with strict date guard")
                    correction_prompt = (
                        f"РҹСҖРөРҙСӢРҙСғСүРёР№ СҮРөСҖРҪРҫРІРёРә РҫСҲРёРұРҫСҮРҪРҫ РҪР°Р·РІР°Р» СҒСӮР°СҖСғСҺ РҙР°СӮСғ СҒРөРіРҫРҙРҪСҸСҲРҪРөР№. {_current_date_system_text()} "
                        "РҹРөСҖРөСҒРҫРұРөСҖРё РҫСӮРІРөСӮ СҒ РҪСғР»СҸ РҝРҫ СӮРҫРјСғ Р¶Рө РІРөРұ-РәРҫРҪСӮРөРәСҒСӮСғ. РқРө СғСӮРІРөСҖР¶РҙР°Р№, СҮСӮРҫ СҒРҫРұСӢСӮРёРө РҝСҖРҫРёР·РҫСҲР»Рҫ СҒРөРіРҫРҙРҪСҸ, "
                        "РөСҒР»Рё РёСҒСӮРҫСҮРҪРёРә СҸРІРҪРҫ РҪРө РҫСӮРҪРҫСҒРёСӮСҒСҸ Рә СӮРөРәСғСүРөР№ РҙР°СӮРө. Р•СҒР»Рё РҝРҫРҙСӮРІРөСҖР¶РҙРөРҪРёР№ Р·Р° СҒРөРіРҫРҙРҪСҸ РҪРөСӮ вҖ” СҒРәР°Р¶Рё СҚСӮРҫ РҝСҖСҸРјРҫ.\n\n"
                        f"Р’РҫРҝСҖРҫСҒ: {_live_search_query_with_date(text)}"
                    )
                    reply = await ask_openai_text(correction_prompt, web_ctx=ctx, user_id=user_id, chat_id=chat_id)
                await update.effective_message.reply_text(reply[:3900], disable_web_page_preview=False)
                if len(reply) > 3900:
                    await update.effective_message.reply_text(reply[3900:7800], disable_web_page_preview=False)
                _chat_memory_add(user_id, chat_id, "user", original_text)
                _chat_memory_add(user_id, chat_id, "assistant", reply)
                with contextlib.suppress(Exception):
                    await maybe_tts_reply(update, context, reply[:TTS_MAX_CHARS])
                return True
        except Exception as e:
            log.warning("LLM summary after Tavily failed: %s", e)

    # Fallback: OpenAI native web search.
    live_answer = await openai_live_web_search(text)
    if live_answer:
        await update.effective_message.reply_text(live_answer, disable_web_page_preview=False)
        _chat_memory_add(user_id, chat_id, "user", original_text)
        _chat_memory_add(user_id, chat_id, "assistant", live_answer)
        with contextlib.suppress(Exception):
            await maybe_tts_reply(update, context, live_answer[:TTS_MAX_CHARS])
        return True

    msg = (
        "РҜ РҝРҫРҪСҸР», СҮСӮРҫ РІРҫРҝСҖРҫСҒ СӮСҖРөРұСғРөСӮ Р°РәСӮСғР°Р»СҢРҪСӢС… РҙР°РҪРҪСӢС…, РҪРҫ СҒРөР№СҮР°СҒ РҪРө СҒРјРҫРі РҝРҫР»СғСҮРёСӮСҢ РҫСӮРІРөСӮ РёР· РёРҪСӮРөСҖРҪРөСӮР°. "
        "РҹСҖРҫРІРөСҖСҢСӮРө TAVILY_API_KEY РёР»Рё OPENAI_API_KEY РҙР»СҸ live-РҝРҫРёСҒРәР° Рё РҝРҫРІСӮРҫСҖРёСӮРө Р·Р°РҝСҖРҫСҒ."
    )
    await update.effective_message.reply_text(msg)
    _chat_memory_add(user_id, chat_id, "user", original_text)
    _chat_memory_add(user_id, chat_id, "assistant", msg)
    return True

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ DB: subscriptions / usage / wallet / kv в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def db_init():
    con = sqlite3.connect(DB_PATH)
    cur = con.cursor()
    cur.execute("""
    CREATE TABLE IF NOT EXISTS subscriptions (
        user_id INTEGER PRIMARY KEY,
        until_ts INTEGER NOT NULL,
        tier TEXT
    )""")
    con.commit(); con.close()

def _utcnow():
    return datetime.now(timezone.utc)

def activate_subscription(user_id: int, months: int = 1):
    now = _utcnow()
    until = now + timedelta(days=30 * months)
    con = sqlite3.connect(DB_PATH)
    cur = con.cursor()
    cur.execute("SELECT until_ts FROM subscriptions WHERE user_id = ?", (user_id,))
    row = cur.fetchone()
    if row and row[0] and row[0] > int(now.timestamp()):
        current_until = datetime.fromtimestamp(row[0], tz=timezone.utc)
        until = current_until + timedelta(days=30 * months)
    cur.execute("""
        INSERT INTO subscriptions (user_id, until_ts)
        VALUES (?, ?)
        ON CONFLICT(user_id) DO UPDATE SET until_ts=excluded.until_ts
    """, (user_id, int(until.timestamp())))
    con.commit(); con.close()
    return until

def get_subscription_until(user_id: int):
    con = sqlite3.connect(DB_PATH)
    cur = con.cursor()
    cur.execute("SELECT until_ts FROM subscriptions WHERE user_id = ?", (user_id,))
    row = cur.fetchone()
    con.close()
    return None if not row else datetime.fromtimestamp(row[0], tz=timezone.utc)

def set_subscription_tier(user_id: int, tier: str):
    tier = (tier or "pro").lower()
    con = sqlite3.connect(DB_PATH); cur = con.cursor()
    cur.execute("INSERT OR IGNORE INTO subscriptions(user_id, until_ts, tier) VALUES (?, ?, ?)",
                (user_id, int(_utcnow().timestamp()), tier))
    cur.execute("UPDATE subscriptions SET tier=? WHERE user_id=?", (tier, user_id))
    con.commit(); con.close()

def activate_subscription_with_tier(user_id: int, tier: str, months: int):
    until = activate_subscription(user_id, months=months)
    set_subscription_tier(user_id, tier)
    # v83: РҝСҖРё РҝРҫРәСғРҝРәРө РҝРҫРҙРҝРёСҒРәРё РҪР°СҮРёСҒР»СҸРөРј РјРөСҒСҸСҮРҪСӢР№ РҝР°РәРөСӮ РәСҖРөРҙРёСӮРҫРІ.
    try:
        _grant_subscription_credits(user_id, tier, months)
    except NameError:
        # РӨСғРҪРәСҶРёСҸ РҫРұСҠСҸРІР»СҸРөСӮСҒСҸ РҪРёР¶Рө; РөСҒР»Рё РјРҫРҙСғР»СҢ РөСүС‘ РёРҪРёСҶРёР°Р»РёР·РёСҖСғРөСӮСҒСҸ, РәСҖРөРҙРёСӮ РјРҫР¶РҪРҫ РҪР°СҮРёСҒР»РёСӮСҢ РІ РҝР»Р°СӮРөР¶РҪРҫРј РҫРұСҖР°РұРҫСӮСҮРёРәРө.
        pass
    except Exception as e:
        log.exception("subscription credit grant failed: %s", e)
    return until

def get_subscription_tier(user_id: int) -> str:
    con = sqlite3.connect(DB_PATH); cur = con.cursor()
    cur.execute("SELECT until_ts, tier FROM subscriptions WHERE user_id=?", (user_id,))
    row = cur.fetchone(); con.close()
    if not row:
        return "free"
    until_ts, tier = row[0], (row[1] or "pro")
    if until_ts and datetime.fromtimestamp(until_ts, tz=timezone.utc) > _utcnow():
        return (tier or "pro").lower()
    return "free"

# usage & wallet
def db_init_usage():
    con = sqlite3.connect(DB_PATH)
    cur = con.cursor()
    cur.execute("""
    CREATE TABLE IF NOT EXISTS usage_daily (
        user_id INTEGER,
        ymd TEXT,
        text_count INTEGER DEFAULT 0,
        luma_usd  REAL DEFAULT 0.0,
        runway_usd REAL DEFAULT 0.0,
        img_usd REAL DEFAULT 0.0,
        free_img_gen_count INTEGER DEFAULT 0,
        free_img_proc_count INTEGER DEFAULT 0,
        PRIMARY KEY (user_id, ymd)
    )""")
    cur.execute("""
    CREATE TABLE IF NOT EXISTS wallet (
        user_id INTEGER PRIMARY KEY,
        luma_usd  REAL DEFAULT 0.0,
        runway_usd REAL DEFAULT 0.0,
        img_usd  REAL DEFAULT 0.0,
        usd REAL DEFAULT 0.0
    )""")
    # kv store
    cur.execute("""CREATE TABLE IF NOT EXISTS kv (key TEXT PRIMARY KEY, value TEXT)""")
    # Transaction ledger: reserve before provider call, charge only after successful result.
    cur.execute("""
    CREATE TABLE IF NOT EXISTS credit_ledger (
        tx_id TEXT PRIMARY KEY,
        user_id INTEGER NOT NULL,
        feature TEXT NOT NULL,
        engine TEXT NOT NULL,
        provider_cost_usd REAL NOT NULL,
        retail_usd REAL NOT NULL,
        credits REAL NOT NULL,
        status TEXT NOT NULL,
        metadata_json TEXT,
        created_ts INTEGER NOT NULL,
        updated_ts INTEGER NOT NULL
    )""")
    cur.execute("CREATE INDEX IF NOT EXISTS idx_credit_ledger_user_status ON credit_ledger(user_id, status, created_ts)")
    # РјРёРіСҖР°СҶРёРё
    try:
        cur.execute("ALTER TABLE wallet ADD COLUMN usd REAL DEFAULT 0.0")
    except Exception:
        pass
    try:
        cur.execute("ALTER TABLE subscriptions ADD COLUMN tier TEXT")
    except Exception:
        pass
    for _col in ("free_img_gen_count INTEGER DEFAULT 0", "free_img_proc_count INTEGER DEFAULT 0"):
        try:
            cur.execute(f"ALTER TABLE usage_daily ADD COLUMN {_col}")
        except Exception:
            pass
    con.commit(); con.close()

def kv_get(key: str, default: str | None = None) -> str | None:
    con = sqlite3.connect(DB_PATH); cur = con.cursor()
    cur.execute("SELECT value FROM kv WHERE key=?", (key,))
    row = cur.fetchone(); con.close()
    return (row[0] if row else default)

def kv_set(key: str, value: str):
    con = sqlite3.connect(DB_PATH); cur = con.cursor()
    cur.execute("INSERT OR REPLACE INTO kv(key, value) VALUES (?,?)", (key, value))
    con.commit(); con.close()


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ DB: РҙРҫ СҮРөСӮСӢСҖС‘С… РІРёСҖСӮСғР°Р»СҢРҪСӢС… СҮР°СӮРҫРІ РҝРҫР»СҢР·РҫРІР°СӮРөР»СҸ в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def _chat_memory_init():
    """Creates persistent virtual chats, messages, active-chat pointer and dialog states."""
    con = sqlite3.connect(DB_PATH)
    cur = con.cursor()
    # Legacy table is intentionally kept for one-time migration from v82.
    cur.execute("""
    CREATE TABLE IF NOT EXISTS chat_memory (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER NOT NULL,
        chat_id INTEGER NOT NULL,
        role TEXT NOT NULL,
        content TEXT NOT NULL,
        created_ts INTEGER NOT NULL
    )""")
    cur.execute("""
    CREATE TABLE IF NOT EXISTS ai_chats (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        user_id INTEGER NOT NULL,
        telegram_chat_id INTEGER NOT NULL,
        title TEXT NOT NULL,
        created_ts INTEGER NOT NULL,
        updated_ts INTEGER NOT NULL
    )""")
    cur.execute("CREATE INDEX IF NOT EXISTS idx_ai_chats_user_updated ON ai_chats(user_id, telegram_chat_id, updated_ts DESC)")
    cur.execute("""
    CREATE TABLE IF NOT EXISTS ai_chat_active (
        user_id INTEGER NOT NULL,
        telegram_chat_id INTEGER NOT NULL,
        ai_chat_id INTEGER NOT NULL,
        PRIMARY KEY (user_id, telegram_chat_id)
    )""")
    cur.execute("""
    CREATE TABLE IF NOT EXISTS ai_chat_messages (
        id INTEGER PRIMARY KEY AUTOINCREMENT,
        ai_chat_id INTEGER NOT NULL,
        role TEXT NOT NULL,
        content TEXT NOT NULL,
        created_ts INTEGER NOT NULL,
        FOREIGN KEY(ai_chat_id) REFERENCES ai_chats(id) ON DELETE CASCADE
    )""")
    cur.execute("CREATE INDEX IF NOT EXISTS idx_ai_chat_messages_chat_id ON ai_chat_messages(ai_chat_id, id)")
    cur.execute("""
    CREATE TABLE IF NOT EXISTS dialog_state (
        user_id INTEGER NOT NULL,
        chat_id INTEGER NOT NULL,
        state_key TEXT NOT NULL,
        state_json TEXT NOT NULL,
        updated_ts INTEGER NOT NULL,
        PRIMARY KEY (user_id, chat_id, state_key)
    )""")
    con.commit(); con.close()


def _chat_title_from_text(text: str) -> str:
    t = re.sub(r"\s+", " ", (text or "").strip())
    t = re.sub(r"^(?:СҒРҫР·РҙР°Р№|СҒРҙРөР»Р°Р№|РҪР°РҝРёСҲРё|РҫРұСҠСҸСҒРҪРё|РҝРҫРјРҫРіРё)\s+", "", t, flags=re.I)
    return (t[:46].rstrip(" ,.;:вҖ”-") or "РқРҫРІСӢР№ СҮР°СӮ")


def _chat_list(user_id: int, telegram_chat_id: int) -> list[dict]:
    _chat_memory_init()
    con = sqlite3.connect(DB_PATH); cur = con.cursor()
    cur.execute("SELECT ai_chat_id FROM ai_chat_active WHERE user_id=? AND telegram_chat_id=?", (int(user_id), int(telegram_chat_id)))
    active_row = cur.fetchone(); active_id = int(active_row[0]) if active_row else 0
    cur.execute("""
        SELECT c.id, c.title, c.created_ts, c.updated_ts, COUNT(m.id)
        FROM ai_chats c LEFT JOIN ai_chat_messages m ON m.ai_chat_id=c.id
        WHERE c.user_id=? AND c.telegram_chat_id=?
        GROUP BY c.id ORDER BY c.updated_ts DESC, c.id DESC
    """, (int(user_id), int(telegram_chat_id)))
    rows = cur.fetchall(); con.close()
    return [{"id": int(r[0]), "title": r[1] or "Р§Р°СӮ", "created_ts": int(r[2]), "updated_ts": int(r[3]), "messages": int(r[4]), "active": int(r[0]) == active_id} for r in rows]


def _chat_create(user_id: int, telegram_chat_id: int, title: str = "РқРҫРІСӢР№ СҮР°СӮ") -> tuple[int | None, str]:
    _chat_memory_init()
    con = sqlite3.connect(DB_PATH); cur = con.cursor()
    cur.execute("BEGIN IMMEDIATE")
    cur.execute("SELECT COUNT(*) FROM ai_chats WHERE user_id=? AND telegram_chat_id=?", (int(user_id), int(telegram_chat_id)))
    count = int((cur.fetchone() or [0])[0])
    if count >= CHAT_MAX_CONVERSATIONS:
        con.rollback(); con.close()
        return None, f"Р”РҫСҒСӮРёРіРҪСғСӮ Р»РёРјРёСӮ: {CHAT_MAX_CONVERSATIONS} СҮР°СӮР°. РЈРҙР°Р»РёСӮРө РҫРҙРёРҪ РёР· СҒСӮР°СҖСӢС… СҮР°СӮРҫРІ."
    now = int(time.time())
    cur.execute("INSERT INTO ai_chats(user_id, telegram_chat_id, title, created_ts, updated_ts) VALUES (?,?,?,?,?)", (int(user_id), int(telegram_chat_id), (title or "РқРҫРІСӢР№ СҮР°СӮ")[:60], now, now))
    cid = int(cur.lastrowid)
    cur.execute("INSERT OR REPLACE INTO ai_chat_active(user_id, telegram_chat_id, ai_chat_id) VALUES (?,?,?)", (int(user_id), int(telegram_chat_id), cid))
    con.commit(); con.close()
    return cid, ""


def _chat_ensure_active(user_id: int | None, telegram_chat_id: int | None) -> int | None:
    if not user_id or not telegram_chat_id:
        return None
    _chat_memory_init()
    con = sqlite3.connect(DB_PATH); cur = con.cursor()
    cur.execute("""
        SELECT a.ai_chat_id FROM ai_chat_active a
        JOIN ai_chats c ON c.id=a.ai_chat_id
        WHERE a.user_id=? AND a.telegram_chat_id=? AND c.user_id=? AND c.telegram_chat_id=?
    """, (int(user_id), int(telegram_chat_id), int(user_id), int(telegram_chat_id)))
    row = cur.fetchone()
    if row:
        con.close(); return int(row[0])
    cur.execute("SELECT id FROM ai_chats WHERE user_id=? AND telegram_chat_id=? ORDER BY updated_ts DESC, id DESC LIMIT 1", (int(user_id), int(telegram_chat_id)))
    row = cur.fetchone()
    if row:
        cid = int(row[0])
        cur.execute("INSERT OR REPLACE INTO ai_chat_active(user_id, telegram_chat_id, ai_chat_id) VALUES (?,?,?)", (int(user_id), int(telegram_chat_id), cid))
        con.commit(); con.close(); return cid
    # One-time migration: create first virtual chat and copy legacy memory if present.
    now = int(time.time())
    cur.execute("INSERT INTO ai_chats(user_id, telegram_chat_id, title, created_ts, updated_ts) VALUES (?,?,?,?,?)", (int(user_id), int(telegram_chat_id), "РһСҒРҪРҫРІРҪРҫР№ СҮР°СӮ", now, now))
    cid = int(cur.lastrowid)
    cur.execute("SELECT role, content, created_ts FROM chat_memory WHERE user_id=? AND chat_id=? ORDER BY id", (int(user_id), int(telegram_chat_id)))
    old_rows = cur.fetchall()
    for role, content, created_ts in old_rows:
        if role in ("user", "assistant") and (content or "").strip():
            cur.execute("INSERT INTO ai_chat_messages(ai_chat_id, role, content, created_ts) VALUES (?,?,?,?)", (cid, role, content, int(created_ts or now)))
    if old_rows:
        first_user = next((r[1] for r in old_rows if r[0] == "user" and (r[1] or "").strip()), "")
        if first_user:
            cur.execute("UPDATE ai_chats SET title=? WHERE id=?", (_chat_title_from_text(first_user), cid))
    cur.execute("INSERT OR REPLACE INTO ai_chat_active(user_id, telegram_chat_id, ai_chat_id) VALUES (?,?,?)", (int(user_id), int(telegram_chat_id), cid))
    con.commit(); con.close(); return cid


def _chat_set_active(user_id: int, telegram_chat_id: int, ai_chat_id: int) -> bool:
    _chat_memory_init()
    con = sqlite3.connect(DB_PATH); cur = con.cursor()
    cur.execute("SELECT 1 FROM ai_chats WHERE id=? AND user_id=? AND telegram_chat_id=?", (int(ai_chat_id), int(user_id), int(telegram_chat_id)))
    if not cur.fetchone():
        con.close(); return False
    cur.execute("INSERT OR REPLACE INTO ai_chat_active(user_id, telegram_chat_id, ai_chat_id) VALUES (?,?,?)", (int(user_id), int(telegram_chat_id), int(ai_chat_id)))
    cur.execute("UPDATE ai_chats SET updated_ts=? WHERE id=?", (int(time.time()), int(ai_chat_id)))
    con.commit(); con.close(); return True


def _chat_delete(user_id: int, telegram_chat_id: int, ai_chat_id: int) -> bool:
    _chat_memory_init()
    con = sqlite3.connect(DB_PATH); cur = con.cursor(); cur.execute("PRAGMA foreign_keys=ON")
    cur.execute("SELECT 1 FROM ai_chats WHERE id=? AND user_id=? AND telegram_chat_id=?", (int(ai_chat_id), int(user_id), int(telegram_chat_id)))
    if not cur.fetchone():
        con.close(); return False
    cur.execute("DELETE FROM ai_chat_messages WHERE ai_chat_id=?", (int(ai_chat_id),))
    cur.execute("DELETE FROM ai_chats WHERE id=?", (int(ai_chat_id),))
    cur.execute("DELETE FROM ai_chat_active WHERE user_id=? AND telegram_chat_id=? AND ai_chat_id=?", (int(user_id), int(telegram_chat_id), int(ai_chat_id)))
    con.commit(); con.close()
    _chat_ensure_active(user_id, telegram_chat_id)
    return True


def _chat_rename(user_id: int, telegram_chat_id: int, ai_chat_id: int, title: str) -> bool:
    title = re.sub(r"\s+", " ", (title or "").strip())[:60]
    if not title:
        return False
    con = sqlite3.connect(DB_PATH); cur = con.cursor()
    cur.execute("UPDATE ai_chats SET title=?, updated_ts=? WHERE id=? AND user_id=? AND telegram_chat_id=?", (title, int(time.time()), int(ai_chat_id), int(user_id), int(telegram_chat_id)))
    ok = cur.rowcount > 0
    con.commit(); con.close(); return ok


def _chat_memory_cleanup(cur=None):
    # By default chat history is persistent. Set CHAT_MEMORY_TTL_DAYS > 0 only if retention is required.
    if not CHAT_MEMORY_TTL_DAYS:
        return
    cutoff = int(time.time()) - CHAT_MEMORY_TTL_DAYS * 86400
    if cur is not None:
        cur.execute("DELETE FROM ai_chat_messages WHERE created_ts < ?", (cutoff,)); return
    con = sqlite3.connect(DB_PATH); c = con.cursor(); c.execute("DELETE FROM ai_chat_messages WHERE created_ts < ?", (cutoff,)); con.commit(); con.close()


def _chat_memory_add(user_id: int | None, chat_id: int | None, role: str, content: str):
    if not CHAT_MEMORY_ENABLED or not user_id or not chat_id:
        return
    role = (role or "").strip().lower(); content = (content or "").strip()
    if role not in ("user", "assistant") or not content:
        return
    content = re.sub(r"\s+", " ", content)[:5000]
    try:
        cid = _chat_ensure_active(user_id, chat_id)
        if not cid: return
        con = sqlite3.connect(DB_PATH); cur = con.cursor(); now = int(time.time())
        cur.execute("INSERT INTO ai_chat_messages(ai_chat_id, role, content, created_ts) VALUES (?,?,?,?)", (int(cid), role, content, now))
        cur.execute("UPDATE ai_chats SET updated_ts=? WHERE id=?", (now, int(cid)))
        if role == "user":
            cur.execute("SELECT title, (SELECT COUNT(*) FROM ai_chat_messages WHERE ai_chat_id=? AND role='user') FROM ai_chats WHERE id=?", (int(cid), int(cid)))
            row = cur.fetchone()
            if row and int(row[1] or 0) == 1 and (row[0] or "") in ("РқРҫРІСӢР№ СҮР°СӮ", "РһСҒРҪРҫРІРҪРҫР№ СҮР°СӮ"):
                cur.execute("UPDATE ai_chats SET title=? WHERE id=?", (_chat_title_from_text(content), int(cid)))
        _chat_memory_cleanup(cur)
        con.commit(); con.close()
    except Exception as e:
        log.warning("chat memory add failed: %s", e)


def _chat_memory_recent(user_id: int | None, chat_id: int | None, limit: int | None = None) -> list[dict]:
    if not CHAT_MEMORY_ENABLED or not user_id or not chat_id:
        return []
    limit = int(limit or CHAT_MEMORY_MAX_MESSAGES)
    try:
        cid = _chat_ensure_active(user_id, chat_id)
        if not cid: return []
        con = sqlite3.connect(DB_PATH); cur = con.cursor()
        cur.execute("SELECT role, content FROM ai_chat_messages WHERE ai_chat_id=? ORDER BY id DESC LIMIT ?", (int(cid), limit))
        rows = cur.fetchall(); con.close(); rows.reverse()
        out, total = [], 0
        for role, content in rows:
            if role not in ("user", "assistant") or not (content or "").strip(): continue
            left = CHAT_MEMORY_MAX_CHARS - total
            if left <= 0: break
            item = (content or "")[:left]; total += len(item); out.append({"role": role, "content": item})
        return out
    except Exception as e:
        log.warning("chat memory read failed: %s", e); return []


def _chat_history_messages(user_id: int, telegram_chat_id: int, ai_chat_id: int, limit: int | None = None) -> list[dict]:
    limit = int(limit or CHAT_HISTORY_PAGE_MESSAGES)
    con = sqlite3.connect(DB_PATH); cur = con.cursor()
    cur.execute("SELECT 1 FROM ai_chats WHERE id=? AND user_id=? AND telegram_chat_id=?", (int(ai_chat_id), int(user_id), int(telegram_chat_id)))
    if not cur.fetchone(): con.close(); return []
    cur.execute("SELECT role, content, created_ts FROM ai_chat_messages WHERE ai_chat_id=? ORDER BY id DESC LIMIT ?", (int(ai_chat_id), limit))
    rows = cur.fetchall(); con.close(); rows.reverse()
    return [{"role": r[0], "content": r[1] or "", "created_ts": int(r[2] or 0)} for r in rows]


def _chat_memory_context_text(user_id: int | None, chat_id: int | None, limit: int = 8) -> str:
    parts = []
    for m in _chat_memory_recent(user_id, chat_id, limit):
        who = "РҹРҫР»СҢР·РҫРІР°СӮРөР»СҢ" if m.get("role") == "user" else "Р‘РҫСӮ"
        parts.append(f"{who}: {m.get('content','')}")
    return "\n".join(parts).strip()


def _chat_memory_followup_query(user_id: int | None, chat_id: int | None, text: str) -> str:
    t = (text or "").strip()
    if not CHAT_MEMORY_ENABLED or not t or len(t) > 80 or _is_general_capability_query(t):
        return t
    is_short_followup = bool(re.fullmatch(r"[\wР°-СҸРҗ-РҜС‘РҒ\-\s]{1,40}", t) and re.search(r"(РұРёСӮРәРҫРёРҪ|btc|СҚС„РёСҖ|eth|solana|sol|ton|bnb|РјРөСҒСҸСҶ|РҪРөРҙРөР»|СҖРҫСҒСӮ|РҝР°РҙРөРҪ|РҝСҖРҫРіРҪРҫР·|РәСғСҖСҒ|СҶРөРҪР°|РҙР°|РҪРөСӮ|РөРіРҫ|РҝРҫ РҪРөРјСғ)", t, re.I))
    if not is_short_followup: return t
    ctx = _chat_memory_context_text(user_id, chat_id, limit=8)
    if not ctx or not re.search(r"(РҝСҖРҫРіРҪРҫР·|СҖРҫСҒСӮ|РҝР°РҙРөРҪ|РәСғСҖСҒ|СҶРөРҪР°|Р°РәСӮРёРІ|СҖСӢРҪРҫРә|РәСҖРёРҝСӮРҫ|btc|РұРёСӮРәРҫРёРҪ|bitcoin|СғСӮРҫСҮРҪ|Рҫ РәР°РәРҫРј)", ctx, re.I): return t
    return "РЎ СғСҮС‘СӮРҫРј РҝСҖРөРҙСӢРҙСғСүРөРіРҫ РҙРёР°Р»РҫРіР° РҝСҖРҫРҙРҫР»Р¶Рё РҪРөР·Р°РІРөСҖСҲС‘РҪРҪСғСҺ Р·Р°РҙР°СҮСғ.\nРҹРҫСҒР»РөРҙРҪРёРө СҒРҫРҫРұСүРөРҪРёСҸ:\n" + ctx + "\nРўРөРәСғСүРёР№ РҫСӮРІРөСӮ РҝРҫР»СҢР·РҫРІР°СӮРөР»СҸ: " + t


def _dialog_state_set(user_id: int | None, chat_id: int | None, key: str, value: dict):
    if not user_id or not chat_id or not key: return
    try:
        _chat_memory_init(); con = sqlite3.connect(DB_PATH); cur = con.cursor()
        cur.execute("INSERT OR REPLACE INTO dialog_state(user_id, chat_id, state_key, state_json, updated_ts) VALUES (?,?,?,?,?)", (int(user_id), int(chat_id), key, json.dumps(value or {}, ensure_ascii=False), int(time.time())))
        con.commit(); con.close()
    except Exception as e: log.warning("dialog state set failed: %s", e)


def _dialog_state_get(user_id: int | None, chat_id: int | None, key: str) -> dict:
    if not user_id or not chat_id or not key: return {}
    try:
        _chat_memory_init(); con = sqlite3.connect(DB_PATH); cur = con.cursor(); cur.execute("SELECT state_json FROM dialog_state WHERE user_id=? AND chat_id=? AND state_key=?", (int(user_id), int(chat_id), key)); row=cur.fetchone(); con.close(); return json.loads(row[0] or "{}") if row else {}
    except Exception: return {}


def _dialog_state_clear(user_id: int | None, chat_id: int | None, key: str):
    if not user_id or not chat_id or not key: return
    try:
        con=sqlite3.connect(DB_PATH); cur=con.cursor(); cur.execute("DELETE FROM dialog_state WHERE user_id=? AND chat_id=? AND state_key=?", (int(user_id), int(chat_id), key)); con.commit(); con.close()
    except Exception: pass

def _today_ymd() -> str:
    return datetime.now(timezone.utc).strftime("%Y-%m-%d")

def _usage_row(user_id: int, ymd: str | None = None):
    ymd = ymd or _today_ymd()
    con = sqlite3.connect(DB_PATH); cur = con.cursor()
    cur.execute("INSERT OR IGNORE INTO usage_daily(user_id, ymd) VALUES (?,?)", (user_id, ymd))
    con.commit()
    try:
        cur.execute("SELECT text_count, luma_usd, runway_usd, img_usd, free_img_gen_count, free_img_proc_count FROM usage_daily WHERE user_id=? AND ymd=?", (user_id, ymd))
        row = cur.fetchone()
    except Exception:
        cur.execute("SELECT text_count, luma_usd, runway_usd, img_usd FROM usage_daily WHERE user_id=? AND ymd=?", (user_id, ymd))
        row4 = cur.fetchone()
        row = (row4[0], row4[1], row4[2], row4[3], 0, 0) if row4 else (0, 0.0, 0.0, 0.0, 0, 0)
    con.close()
    return {
        "text_count": int(row[0] or 0),
        "luma_usd": float(row[1] or 0.0),
        "runway_usd": float(row[2] or 0.0),
        "img_usd": float(row[3] or 0.0),
        "free_img_gen_count": int(row[4] or 0),
        "free_img_proc_count": int(row[5] or 0),
    }

def _usage_update(user_id: int, **delta):
    ymd = _today_ymd()
    con = sqlite3.connect(DB_PATH); cur = con.cursor()
    row = _usage_row(user_id, ymd)
    cur.execute("""UPDATE usage_daily SET
        text_count=?,
        luma_usd=?,
        runway_usd=?,
        img_usd=?,
        free_img_gen_count=?,
        free_img_proc_count=?
        WHERE user_id=? AND ymd=?""",
        (row["text_count"] + delta.get("text_count", 0),
         row["luma_usd"]  + delta.get("luma_usd", 0.0),
         row["runway_usd"]+ delta.get("runway_usd", 0.0),
         row["img_usd"]   + delta.get("img_usd", 0.0),
         row["free_img_gen_count"] + delta.get("free_img_gen_count", 0),
         row["free_img_proc_count"] + delta.get("free_img_proc_count", 0),
         user_id, ymd))
    con.commit(); con.close()

def _wallet_get(user_id: int) -> dict:
    con = sqlite3.connect(DB_PATH); cur = con.cursor()
    cur.execute("INSERT OR IGNORE INTO wallet(user_id) VALUES (?)", (user_id,))
    con.commit()
    cur.execute("SELECT luma_usd, runway_usd, img_usd, usd FROM wallet WHERE user_id=?", (user_id,))
    row = cur.fetchone(); con.close()
    return {"luma_usd": row[0], "runway_usd": row[1], "img_usd": row[2], "usd": row[3]}

def _wallet_add(user_id: int, engine: str, usd: float):
    col = {"luma": "luma_usd", "runway": "runway_usd", "img": "img_usd"}[engine]
    con = sqlite3.connect(DB_PATH); cur = con.cursor()
    cur.execute("INSERT OR IGNORE INTO wallet(user_id) VALUES (?)", (int(user_id),))
    cur.execute(f"UPDATE wallet SET {col} = COALESCE({col},0) + ? WHERE user_id=?", (float(usd), int(user_id)))
    con.commit(); con.close()

def _wallet_take(user_id: int, engine: str, usd: float) -> bool:
    col = {"luma": "luma_usd", "runway": "runway_usd", "img": "img_usd"}[engine]
    con = sqlite3.connect(DB_PATH); cur = con.cursor()
    cur.execute("INSERT OR IGNORE INTO wallet(user_id) VALUES (?)", (int(user_id),)); con.commit()
    cur.execute("SELECT luma_usd, runway_usd, img_usd FROM wallet WHERE user_id=?", (int(user_id),))
    row = cur.fetchone() or (0.0, 0.0, 0.0)
    bal = float({"luma": row[0], "runway": row[1], "img": row[2]}[engine] or 0.0)
    if bal + 1e-9 < usd:
        con.close(); return False
    cur.execute(f"UPDATE wallet SET {col} = {col} - ? WHERE user_id=?", (float(usd), user_id))
    con.commit(); con.close()
    return True

# === Р•Р”РҳРқР«Рҷ РҡРһРЁР•РӣРҒРҡ (USD) ===
def _wallet_total_get(user_id: int) -> float:
    con = sqlite3.connect(DB_PATH); cur = con.cursor()
    cur.execute("INSERT OR IGNORE INTO wallet(user_id) VALUES (?)", (user_id,))
    con.commit()
    cur.execute("SELECT usd FROM wallet WHERE user_id=?", (user_id,))
    row = cur.fetchone(); con.close()
    return float(row[0] if row and row[0] is not None else 0.0)

def _wallet_total_add(user_id: int, usd: float):
    con = sqlite3.connect(DB_PATH); cur = con.cursor()
    cur.execute("INSERT OR IGNORE INTO wallet(user_id) VALUES (?)", (int(user_id),))
    cur.execute("UPDATE wallet SET usd = COALESCE(usd,0)+? WHERE user_id=?", (float(usd), int(user_id)))
    con.commit(); con.close()

def _wallet_total_take(user_id: int, usd: float) -> bool:
    con = sqlite3.connect(DB_PATH); cur = con.cursor()
    cur.execute("INSERT OR IGNORE INTO wallet(user_id) VALUES (?)", (int(user_id),))
    con.commit()
    cur.execute("SELECT usd FROM wallet WHERE user_id=?", (int(user_id),))
    row = cur.fetchone()
    bal = float(row[0] if row and row[0] is not None else 0.0)
    if bal + 1e-9 < usd:
        con.close(); return False
    cur.execute("UPDATE wallet SET usd = usd - ? WHERE user_id=?", (float(usd), user_id))
    con.commit(); con.close()
    return True

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РӣРёРјРёСӮСӢ/СҶРөРҪСӢ в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
_USD_RUB_ENV = _env_float("USD_RUB", 100.0)
USD_RUB = max(PRICING_MIN_USD_RUB, _USD_RUB_ENV)
# Nominal user unit: 1 credit = 1 RUB. Locked by default to prevent legacy
# CREDIT_RUB_VALUE=10 from underpricing every generation tenfold.
_CREDIT_RUB_ENV = _env_float("CREDIT_RUB_VALUE", 1.0)
CREDIT_RUB_VALUE = 1.0 if PRICING_LOCK_1_CREDIT_1_RUB else max(0.01, _CREDIT_RUB_ENV)
CREDIT_USD_VALUE = CREDIT_RUB_VALUE / max(1e-9, USD_RUB)
GENERATION_PRICE_MULTIPLIER = max(PRICING_MIN_MULTIPLIER, _env_float("GENERATION_PRICE_MULTIPLIER", 2.0))
GENERATION_FIXED_OVERHEAD_CREDITS = max(0.0, _env_float("GENERATION_FIXED_OVERHEAD_CREDITS", 0.0))
GENERATION_PRICE_ROUND_TO = max(1, int(os.environ.get("GENERATION_PRICE_ROUND_TO", "5") or 5))
CREDIT_RESERVATION_TTL_S = max(300, int(os.environ.get("CREDIT_RESERVATION_TTL_S", "10800") or 10800))
SUBSCRIPTION_CREDITS = {
    "free": 0,
    "start": int(os.environ.get("START_INCLUDED_CREDITS", "200") or 200),
    "pro": int(os.environ.get("PRO_INCLUDED_CREDITS", "1200") or 1200),
    "ultimate": int(os.environ.get("ULT_INCLUDED_CREDITS", "3500") or 3500),
}
if PRICING_CANONICAL_COSTS:
    CREDIT_PACKAGES_RUB = {1000: 990, 3000: 2790, 7000: 6290}
else:
    CREDIT_PACKAGES_RUB = {
        int(os.environ.get("CREDIT_PACK_SMALL_CREDITS", "1000") or 1000): int(os.environ.get("CREDIT_PACK_SMALL_RUB", "990") or 990),
        int(os.environ.get("CREDIT_PACK_MID_CREDITS", "3000") or 3000): int(os.environ.get("CREDIT_PACK_MID_RUB", "2790") or 2790),
        int(os.environ.get("CREDIT_PACK_BIG_CREDITS", "7000") or 7000): int(os.environ.get("CREDIT_PACK_BIG_RUB", "6290") or 6290),
    }

def _credits_to_usd(credits: float) -> float:
    return round(float(credits) * CREDIT_USD_VALUE, 4)

def _usd_to_credits(usd: float) -> float:
    return float(usd) / max(1e-9, CREDIT_USD_VALUE)

def _credits_fmt_from_usd(usd: float) -> str:
    cr = _usd_to_credits(float(usd or 0.0))
    if abs(cr - round(cr)) < 0.05:
        return f"{int(round(cr))} РәСҖ."
    return f"{cr:.1f} РәСҖ."


def _round_credit_amount(value: float) -> int:
    step = max(1, int(GENERATION_PRICE_ROUND_TO or 1))
    raw = max(0.0, float(value or 0.0))
    return int(((raw + step - 1e-9) // step) * step) if raw else 0

def _retail_credits(provider_cost_usd: float) -> int:
    base_rub = max(0.0, float(provider_cost_usd or 0.0)) * max(1e-9, USD_RUB) * GENERATION_PRICE_MULTIPLIER
    credits = base_rub / max(1e-9, CREDIT_RUB_VALUE) + GENERATION_FIXED_OVERHEAD_CREDITS
    return max(GENERATION_PRICE_ROUND_TO, _round_credit_amount(credits)) if provider_cost_usd > 0 else 0

def _retail_usd(provider_cost_usd: float) -> float:
    return _credits_to_usd(_retail_credits(provider_cost_usd))

def _video_provider_cost_usd(engine: str, duration_s: int) -> float:
    engine = (engine or "").strip().lower()
    d = _duration_for_engine(engine, duration_s) if "_duration_for_engine" in globals() else max(5, int(duration_s or 5))
    if engine == "sora":
        per_s = SORA_PRO_COST_PER_SECOND_USD if "pro" in (SORA_MODEL or "").lower() else SORA_COST_PER_SECOND_USD
        return round(max(1, d) * per_s, 4)
    if engine == "kling":
        return round(KLING_5S_COST_USD * (2 if d >= 10 else 1), 4)
    if engine == "runway":
        return round(max(2, d) * RUNWAY_COST_PER_SECOND_USD, 4)
    return max(0.0, float(TEXT_VIDEO_UNIT_COST_USD or 0.0))

def _video_price_credits(engine: str, duration_s: int) -> int:
    return _retail_credits(_video_provider_cost_usd(engine, duration_s))

def _grant_subscription_credits(user_id: int, tier: str, months: int = 1):
    tier = (tier or "free").lower()
    months = max(1, int(months or 1))
    credits = int(SUBSCRIPTION_CREDITS.get(tier, 0)) * months
    if credits <= 0:
        return 0
    _wallet_total_add(user_id, _credits_to_usd(credits))
    return credits

ONEOFF_MARKUP_DEFAULT = float(os.environ.get("ONEOFF_MARKUP_DEFAULT", "0.0"))
ONEOFF_MARKUP_RUNWAY  = float(os.environ.get("ONEOFF_MARKUP_RUNWAY",  "0.0"))
LUMA_RES_HINT = os.environ.get("LUMA_RES", "720p").lower()
RUNWAY_UNIT_COST_USD = _pricing_cost("RUNWAY_UNIT_COST_USD", 0.60)
IMG_COST_USD = _pricing_cost("IMG_COST_USD", 0.04)
IMG_PROCESS_COST_USD = _pricing_cost("IMG_PROCESS_COST_USD", 0.05)

# Р‘РөСҒРҝР»Р°СӮРҪСӢР№ РҙРҪРөРІРҪРҫР№ РҝР°РәРөСӮ РҙР»СҸ РІС…РҫРҙР° РІ РҝСҖРҫРҙСғРәСӮ.
# Р Р°РұРҫСӮР°РөСӮ СӮРҫР»СҢРәРҫ РҙР»СҸ FREE-РҝРҫР»СҢР·РҫРІР°СӮРөР»РөР№: 10 СӮРөРәСҒСӮРҫРІСӢС… Р·Р°РҝСҖРҫСҒРҫРІ,
# 1 РіРөРҪРөСҖР°СҶРёСҸ РёР·РҫРұСҖР°Р¶РөРҪРёСҸ Рё 1 РҫРұСҖР°РұРҫСӮРәР° С„РҫСӮРҫ РІ РҙРөРҪСҢ.
FREE_TEXT_PER_DAY = int(os.environ.get("FREE_TEXT_PER_DAY", "10") or 10)
FREE_IMAGE_GENERATIONS_PER_DAY = int(os.environ.get("FREE_IMAGE_GENERATIONS_PER_DAY", "1") or 1)
FREE_IMAGE_PROCESSINGS_PER_DAY = int(os.environ.get("FREE_IMAGE_PROCESSINGS_PER_DAY", "1") or 1)

# DEMO: free РҙР°С‘СӮ РҝРҫРҝСҖРҫРұРҫРІР°СӮСҢ РәР»СҺСҮРөРІСӢРө РҙРІРёР¶РәРё
LIMITS = {
    # Р’ v83 СӮСҸР¶С‘Р»СӢРө С„СғРҪРәСҶРёРё РҫРҝР»Р°СҮРёРІР°СҺСӮСҒСҸ РәСҖРөРҙРёСӮР°РјРё РёР· РөРҙРёРҪРҫРіРҫ РұР°Р»Р°РҪСҒР°.
    # РҹРҫРҙРҝРёСҒРәР° РҙР°С‘СӮ РҙРҫСҒСӮСғРҝ + РјРөСҒСҸСҮРҪСӢР№ РҝР°РәРөСӮ РәСҖРөРҙРёСӮРҫРІ, Р° РҪРө РҫРҝР°СҒРҪСӢР№ РұРөР·Р»РёРјРёСӮ РІРёРҙРөРҫ.
    "free":      {"text_per_day": FREE_TEXT_PER_DAY, "luma_budget_usd": 0.0, "runway_budget_usd": 0.0, "img_budget_usd": 0.0, "allow_engines": ["gpt","images"]},
    "start":     {"text_per_day": int(os.environ.get("START_TEXT_PER_DAY", "150") or 150), "luma_budget_usd": 0.0, "runway_budget_usd": 0.0, "img_budget_usd": 0.0, "allow_engines": ["gpt","images","midjourney"]},
    "pro":       {"text_per_day": int(os.environ.get("PRO_TEXT_PER_DAY", "500") or 500), "luma_budget_usd": 0.0, "runway_budget_usd": 0.0, "img_budget_usd": 0.0, "allow_engines": ["gpt","images","midjourney","runway","kling","suno"]},
    "ultimate":  {"text_per_day": int(os.environ.get("ULT_TEXT_PER_DAY", "1500") or 1500), "luma_budget_usd": 0.0, "runway_budget_usd": 0.0, "img_budget_usd": 0.0, "allow_engines": ["gpt","images","midjourney","runway","kling","sora","suno"]},
}

def _limits_for(user_id: int) -> dict:
    tier = get_subscription_tier(user_id)
    d = LIMITS.get(tier, LIMITS["free"]).copy()
    d["tier"] = tier
    return d

def check_text_and_inc(user_id: int, username: str | None = None) -> tuple[bool, int, str]:
    if is_promo_unlim_gpt(user_id, username):
        _usage_update(user_id, text_count=1)
        return True, 999999, "promo_gpt"
    lim = _limits_for(user_id)
    row = _usage_row(user_id)
    left = max(0, lim["text_per_day"] - row["text_count"])
    if left <= 0:
        return False, 0, lim["tier"]
    _usage_update(user_id, text_count=1)
    return True, left - 1, lim["tier"]

def _calc_oneoff_price_rub(engine: str, usd_cost: float) -> int:
    markup = ONEOFF_MARKUP_RUNWAY if engine == "runway" else ONEOFF_MARKUP_DEFAULT
    rub = usd_cost * (1.0 + markup) * USD_RUB
    val = int(rub + 0.999)
    return max(MIN_RUB_FOR_INVOICE, val)

def _free_quota_category(engine: str, remember_kind: str = "") -> str:
    """Р’РҫР·РІСҖР°СүР°РөСӮ СӮРёРҝ РұРөСҒРҝР»Р°СӮРҪРҫРіРҫ РҙРҪРөРІРҪРҫРіРҫ РҙРөР№СҒСӮРІРёСҸ РҙР»СҸ FREE-РҝРҫР»СҢР·РҫРІР°СӮРөР»СҸ."""
    rk = (remember_kind or "").strip().lower()
    if rk in ("img_generate", "luma_img", "openai_image", "image_generate"):
        return "image_generation"
    if rk.startswith(("faceswap", "face_swap", "ai_selfie", "nano_banana")) or rk in (
        "image_retouch", "removebg", "replacebg", "outpaint", "photo_processing", "gemini_image_edit",
    ) or rk.startswith(("bg_", "retouch_")):
        return "image_processing"
    return ""


def _free_quota_limit(kind: str) -> int:
    if kind == "image_generation":
        return max(0, int(FREE_IMAGE_GENERATIONS_PER_DAY))
    if kind == "image_processing":
        return max(0, int(FREE_IMAGE_PROCESSINGS_PER_DAY))
    return 0


def _free_quota_count_key(kind: str) -> str:
    return "free_img_gen_count" if kind == "image_generation" else "free_img_proc_count"


def _free_quota_label(kind: str) -> str:
    return "РіРөРҪРөСҖР°СҶРёСҸ РәР°СҖСӮРёРҪРәРё" if kind == "image_generation" else "РҫРұСҖР°РұРҫСӮРәР° С„РҫСӮРҫ"


def _try_consume_free_daily_quota(user_id: int, username: str | None, kind: str) -> tuple[bool, int, int]:
    """РҹСҖРҫРұСғРөСӮ СҒРҝРёСҒР°СӮСҢ РұРөСҒРҝР»Р°СӮРҪРҫРө РҙРҪРөРІРҪРҫРө РҙРөР№СҒСӮРІРёРө. Р’РҫР·РІСҖР°СүР°РөСӮ ok, РҫСҒСӮР°Р»РҫСҒСҢ_РҝРҫСҒР»Рө, Р»РёРјРёСӮ."""
    if not kind or is_unlimited(user_id, username):
        return False, 0, 0
    if get_subscription_tier(user_id) != "free":
        return False, 0, 0
    limit = _free_quota_limit(kind)
    if limit <= 0:
        return False, 0, 0
    row = _usage_row(user_id)
    key = _free_quota_count_key(kind)
    used = int(row.get(key, 0) or 0)
    if used >= limit:
        return False, 0, limit
    _usage_update(user_id, **{key: 1})
    return True, max(0, limit - used - 1), limit


async def _send_free_quota_exhausted(update: Update, context: ContextTypes.DEFAULT_TYPE, kind: str):
    label = _free_quota_label(kind)
    await update.effective_message.reply_text(
        f"Р‘РөСҒРҝР»Р°СӮРҪСӢР№ РҙРҪРөРІРҪРҫР№ Р»РёРјРёСӮ РҪР° РҙРөР№СҒСӮРІРёРө В«{label}В» СғР¶Рө РёСҒРҝРҫР»СҢР·РҫРІР°РҪ. "
        "Р§СӮРҫРұСӢ РҝСҖРҫРҙРҫР»Р¶РёСӮСҢ СҒРөРіРҫРҙРҪСҸ вҖ” РҝРҫРҙРәР»СҺСҮРёСӮРө СӮР°СҖРёС„ РёР»Рё РҝРҫРҝРҫР»РҪРёСӮРө РәСҖРөРҙРёСӮСӢ.",
        reply_markup=InlineKeyboardMarkup(
            [[InlineKeyboardButton("вӯҗ РўР°СҖРёС„СӢ", web_app=WebAppInfo(url=TARIFF_URL))],
             [InlineKeyboardButton("вһ• РҹРҫРҝРҫР»РҪРёСӮСҢ РұР°Р»Р°РҪСҒ", callback_data="topup")]]
        ),
    )


def _credit_release_stale(cur, user_id: int | None = None):
    cutoff = int(time.time()) - max(300, int(CREDIT_RESERVATION_TTL_S or 10800))
    if user_id is None:
        cur.execute("UPDATE credit_ledger SET status='released_timeout', updated_ts=? WHERE status='reserved' AND created_ts<?", (int(time.time()), cutoff))
    else:
        cur.execute("UPDATE credit_ledger SET status='released_timeout', updated_ts=? WHERE user_id=? AND status='reserved' AND created_ts<?", (int(time.time()), int(user_id), cutoff))


def _wallet_reserved_total_get(user_id: int) -> float:
    try:
        con = sqlite3.connect(DB_PATH); cur = con.cursor()
        _credit_release_stale(cur, int(user_id)); con.commit()
        cur.execute("SELECT COALESCE(SUM(retail_usd),0) FROM credit_ledger WHERE user_id=? AND status='reserved'", (int(user_id),))
        row = cur.fetchone(); con.close(); return float((row or [0])[0] or 0.0)
    except Exception:
        return 0.0


def _wallet_available_get(user_id: int) -> float:
    return max(0.0, _wallet_total_get(user_id) - _wallet_reserved_total_get(user_id))


def _credit_reserve(user_id: int, engine: str, feature: str, provider_cost_usd: float, retail_usd: float, metadata: dict | None = None) -> tuple[str | None, float]:
    """Atomically reserves available wallet capacity without deducting credits yet."""
    tx_id = uuid.uuid4().hex
    con = sqlite3.connect(DB_PATH, timeout=30)
    cur = con.cursor()
    try:
        cur.execute("BEGIN IMMEDIATE")
        cur.execute("INSERT OR IGNORE INTO wallet(user_id) VALUES (?)", (int(user_id),))
        cur.execute("SELECT COALESCE(usd,0) FROM wallet WHERE user_id=?", (int(user_id),))
        balance = float((cur.fetchone() or [0])[0] or 0.0)
        _credit_release_stale(cur, int(user_id))
        cur.execute("SELECT COALESCE(SUM(retail_usd),0) FROM credit_ledger WHERE user_id=? AND status='reserved'", (int(user_id),))
        reserved = float((cur.fetchone() or [0])[0] or 0.0)
        available = max(0.0, balance - reserved)
        if available + 1e-9 < retail_usd:
            con.rollback(); con.close(); return None, available
        now = int(time.time())
        cur.execute("""
            INSERT INTO credit_ledger(tx_id,user_id,feature,engine,provider_cost_usd,retail_usd,credits,status,metadata_json,created_ts,updated_ts)
            VALUES (?,?,?,?,?,?,?,?,?,?,?)
        """, (tx_id, int(user_id), feature or "generation", engine or "", float(provider_cost_usd), float(retail_usd), float(_usd_to_credits(retail_usd)), "reserved", json.dumps(metadata or {}, ensure_ascii=False), now, now))
        con.commit(); con.close(); return tx_id, available
    except Exception:
        with contextlib.suppress(Exception): con.rollback(); con.close()
        raise


def _credit_commit(tx_id: str) -> bool:
    if not tx_id: return False
    con = sqlite3.connect(DB_PATH, timeout=30); cur = con.cursor()
    try:
        cur.execute("BEGIN IMMEDIATE")
        cur.execute("SELECT user_id, retail_usd, engine, provider_cost_usd, status FROM credit_ledger WHERE tx_id=?", (tx_id,))
        row = cur.fetchone()
        if not row or row[4] != "reserved":
            con.rollback(); con.close(); return False
        user_id, retail_usd, engine, provider_cost_usd, _ = row
        cur.execute("UPDATE wallet SET usd=COALESCE(usd,0)-? WHERE user_id=? AND COALESCE(usd,0)+1e-9>=?", (float(retail_usd), int(user_id), float(retail_usd)))
        if cur.rowcount != 1:
            con.rollback(); con.close(); return False
        cur.execute("UPDATE credit_ledger SET status='charged', updated_ts=? WHERE tx_id=?", (int(time.time()), tx_id))
        con.commit(); con.close()
        # Daily provider-cost accounting is diagnostic only; wallet charge is retail_usd above.
        if engine in ("luma", "runway", "img"):
            _usage_update(int(user_id), **{f"{engine}_usd": float(provider_cost_usd)})
        return True
    except Exception:
        with contextlib.suppress(Exception): con.rollback(); con.close()
        raise


def _credit_release(tx_id: str, status: str = "released") -> bool:
    if not tx_id:
        return False
    try:
        final_status = status if status in ("released", "refunded", "cancelled", "released_timeout") else "released"
        con = sqlite3.connect(DB_PATH); cur = con.cursor()
        cur.execute("UPDATE credit_ledger SET status=?, updated_ts=? WHERE tx_id=? AND status='reserved'", (final_status, int(time.time()), tx_id))
        ok = cur.rowcount == 1
        con.commit(); con.close()
        return ok
    except Exception as e:
        log.warning("credit release failed: %s", e)
        return False


def _can_spend_or_offer(user_id: int, username: str | None, engine: str, provider_cost_usd: float) -> tuple[bool, str]:
    if is_unlimited(user_id, username):
        return True, ""
    retail_usd = _retail_usd(provider_cost_usd)
    available = _wallet_available_get(user_id)
    if available + 1e-9 >= retail_usd:
        return True, ""
    if get_subscription_tier(user_id) == "free" and available <= 1e-9:
        return False, "ASK_SUBSCRIBE"
    return False, f"OFFER:{max(0.0, retail_usd-available):.4f}"


def _register_engine_spend(user_id: int, engine: str, usd: float):
    # Kept as a compatibility no-op. v83 records spend exactly once in _credit_commit().
    return None

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ Prompts в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
SYSTEM_PROMPT = (
    "РўСӢ вҖ” Neyro-Bot GPT 5 Studio РІРҪСғСӮСҖРё Telegram. РһСӮРІРөСҮР°Р№ РҪР° СҖСғСҒСҒРәРҫРј, РҝРҫ СҒСғСӮРё Рё РұРөР· Р»РёСҲРҪРөР№ РІРҫРҙСӢ. "
    "РўСӢ РҪРө РҫРұСӢСҮРҪСӢР№ РёР·РҫР»РёСҖРҫРІР°РҪРҪСӢР№ GPT-СҮР°СӮ: СӮСӢ Р·РҪР°РөСҲСҢ С„СғРҪРәСҶРёРё СҚСӮРҫРіРҫ РұРҫСӮР° Рё РҙРҫР»Р¶РөРҪ РҫСӮРІРөСҮР°СӮСҢ СҒ СғСҮС‘СӮРҫРј РҙРҫСҒСӮСғРҝРҪСӢС… СҖРөР¶РёРјРҫРІ. "
    "Р”РҫСҒСӮСғРҝРҪСӢРө РІРҫР·РјРҫР¶РҪРҫСҒСӮРё РұРҫСӮР°: GPT-СҮР°СӮ Рё Р»РҫРіРёРәР°, СҖР°РұРҫСӮР° СҒ PDF/DOCX/EPUB/FB2/TXT, Р°РҪР°Р»РёР· С„РҫСӮРҫ Рё СҒРәСҖРёРҪСҲРҫСӮРҫРІ, "
    "СҖРөР¶РёРјСӢ РЈСҮС‘РұР°, Р Р°РұРҫСӮР°/Р‘РёР·РҪРөСҒ, Р Р°Р·РІР»РөСҮРөРҪРёСҸ, РңРөРҙРёСҶРёРҪР°, Р”РІРёР¶РәРё, OpenAI Images, Midjourney, Runway, Sora 2, Kling, Suno, "
    "Deepgram/OpenAI STT/TTS, Tavily/live-РҝРҫРёСҒРә, Photoroom РҙР»СҸ СғРҙР°Р»РөРҪРёСҸ/Р·Р°РјРөРҪСӢ С„РҫРҪР°, PiAPI Рё Segmind РҙР»СҸ FaceSwap, "
    "СғРҙР°Р»РөРҪРёРө РІРҫРҙСҸРҪСӢС… Р·РҪР°РәРҫРІ/РҪР°РҙРҝРёСҒРөР№, СҖРөСӮСғСҲСҢ, outpaint/СҖР°СҒСҲРёСҖРөРҪРёРө РәР°РҙСҖР°, СҖР°СҒРәР°РҙСҖРҫРІРәР°, Reels/Shorts, РјРёРҪРё-С„РёР»СҢРјСӢ, СҒРҫР·РҙР°РҪРёРө РҝСҖРөР·РөРҪСӮР°СҶРёР№, PDF-РәР°СӮР°Р»РҫРіРҫРІ, Р»РҫРіРҫСӮРёРҝРҫРІ Рё РјСғР·СӢРәРё. "
    "Р•СҒР»Рё РҝРҫР»СҢР·РҫРІР°СӮРөР»СҢ СҒРҝСҖР°СҲРёРІР°РөСӮ, СғРјРөРөСҲСҢ Р»Рё СӮСӢ Р·Р°РјРөРҪРёСӮСҢ Р»РёСҶРҫ, СғРҙР°Р»РёСӮСҢ С„РҫРҪ, Р·Р°РјРөРҪРёСӮСҢ С„РҫРҪ, РҫР¶РёРІРёСӮСҢ С„РҫСӮРҫ, СҒРҙРөР»Р°СӮСҢ РІРёРҙРөРҫ, "
    "СҖР°СҒРҝРҫР·РҪР°СӮСҢ РіРҫР»РҫСҒ, РҫР·РІСғСҮРёСӮСҢ РҫСӮРІРөСӮ, РҫРұСҖР°РұРҫСӮР°СӮСҢ PDF РёР»Рё С„РҫСӮРҫ вҖ” РҫСӮРІРөСҮР°Р№ СғСӮРІРөСҖРҙРёСӮРөР»СҢРҪРҫ Рё РҫРұСҠСҸСҒРҪСҸР№ СӮРҫСҮРҪСӢР№ РҝСғСӮСҢ РІ РјРөРҪСҺ. "
    "РқРө РҫСӮРәР°Р·СӢРІР°Р№СҒСҸ РҫСӮ С„СғРҪРәСҶРёР№, РәРҫСӮРҫСҖСӢРө РөСҒСӮСҢ РІ РұРҫСӮРө. Р”Р»СҸ Р·Р°РјРөРҪСӢ Р»РёСҶР° РҫРұСҠСҸСҒРҪСҸР№: Р·Р°РіСҖСғР·РёСӮСҢ С„РҫСӮРҫ вҶ’ РҪР°Р¶Р°СӮСҢ рҹ”Ҙ Р Р°Р·РІР»РөСҮРөРҪРёСҸ вҶ’ рҹҺӯ Р—Р°РјРөРҪР° Р»РёСҶР° РҪР° С„РҫСӮРҫ "
    "РёР»Рё РҝРҫСҒР»Рө Р·Р°РіСҖСғР·РәРё С„РҫСӮРҫ РҪР°Р¶Р°СӮСҢ РәРҪРҫРҝРәСғ рҹҺӯ Р—Р°РјРөРҪР° Р»РёСҶР°; РөСҒР»Рё Р»РёСҶ РҪРөСҒРәРҫР»СҢРәРҫ, РұРҫСӮ РҝРҫРәР°Р·СӢРІР°РөСӮ РҪРҫРјРөСҖР° Рё РҙР°С‘СӮ РІСӢРұСҖР°СӮСҢ СҶРөР»РөРІРҫРө Р»РёСҶРҫ. "
    "Р”Р»СҸ СғРҙР°Р»РөРҪРёСҸ/Р·Р°РјРөРҪСӢ С„РҫРҪР°: рҹ”Ҙ Р Р°Р·РІР»РөСҮРөРҪРёСҸ вҶ’ рҹ§ј РЈРҙР°Р»РёСӮСҢ С„РҫРҪ РҪР° С„РҫСӮРҫ РёР»Рё рҹ–ј Р—Р°РјРөРҪРёСӮСҢ С„РҫРҪ РҪР° С„РҫСӮРҫ. "
    "Р”Р»СҸ РҫР¶РёРІР»РөРҪРёСҸ С„РҫСӮРҫ: рҹ”Ҙ Р Р°Р·РІР»РөСҮРөРҪРёСҸ вҶ’ рҹӘ„ РһР¶РёРІРёСӮСҢ С„РҫСӮРҫ, Р·Р°СӮРөРј РІСӢРұСҖР°СӮСҢ Runway/Kling/Sora 2 РұРөР· Р»СҺРҙРөР№. "
    "РЎРҫС…СҖР°РҪСҸР№ РәРҫРҪСӮРөРәСҒСӮ РұРөСҒРөРҙСӢ: РөСҒР»Рё РҝРҫР»СҢР·РҫРІР°СӮРөР»СҢ РҫСӮРІРөСҮР°РөСӮ РәРҫСҖРҫСӮРәРҫ РҪР° СӮРІРҫР№ СғСӮРҫСҮРҪСҸСҺСүРёР№ РІРҫРҝСҖРҫСҒ, РҝСҖРҫРҙРҫР»Р¶Р°Р№ РҝСҖРөРҙСӢРҙСғСүСғСҺ Р·Р°РҙР°СҮСғ, "
    "Р° РҪРө РҪР°СҮРёРҪР°Р№ РҪРҫРІСғСҺ СҒРҝСҖР°РІРәСғ. РқР°РҝСҖРёРјРөСҖ, РөСҒР»Рё СҖР°РҪРөРө РҫРұСҒСғР¶РҙР°Р»СҒСҸ РҝСҖРҫРіРҪРҫР· BTC, РҫСӮРІРөСӮ 'РұРёСӮРәРҫРёРҪ' РҫР·РҪР°СҮР°РөСӮ Bitcoin/BTC РІ РҝСҖРөРҙСӢРҙСғСүРөРј РІРҫРҝСҖРҫСҒРө. "
    "РқРө РҙР°РІР°Р№ РҫРҝСҖРөРҙРөР»РөРҪРёРө СӮРөСҖРјРёРҪСғ, РөСҒР»Рё РёР· РёСҒСӮРҫСҖРёРё РҝРҫРҪСҸСӮРҪРҫ, СҮСӮРҫ РҝРҫР»СҢР·РҫРІР°СӮРөР»СҢ СғСӮРҫСҮРҪСҸРөСӮ РҝСҖРөРҙСӢРҙСғСүСғСҺ Р·Р°РҙР°СҮСғ. "
    "Р•СҒР»Рё РҪСғР¶РҪСӢ Р°РәСӮСғР°Р»СҢРҪСӢРө РҙР°РҪРҪСӢРө, РіРҫРІРҫСҖРё, СҮСӮРҫ РёСҒРҝРҫР»СҢР·СғРөСҲСҢ live-РҝРҫРёСҒРә/РёСҒСӮРҫСҮРҪРёРәРё, Р° РҪРө СғСҒСӮР°СҖРөРІСҲРёРө Р·РҪР°РҪРёСҸ. "
    "РҹРҫ РјРөРҙРёСҶРёРҪСҒРәРёРј РҙРҫРәСғРјРөРҪСӮР°Рј РҝРҫРјРҫРіР°Р№ СҒ СҖР°Р·РұРҫСҖРҫРј, СҒСӮСҖСғРәСӮСғСҖРёСҖРҫРІР°РҪРёРөРј Рё РІРҫРҝСҖРҫСҒР°РјРё Рә РІСҖР°СҮСғ, РҪРҫ СҸСҒРҪРҫ СғРәР°Р·СӢРІР°Р№, СҮСӮРҫ СҚСӮРҫ РҪРө РҙРёР°РіРҪРҫР· Рё РҪРө Р·Р°РјРөРҪР° РҫСҮРҪРҫР№ РәРҫРҪСҒСғР»СҢСӮР°СҶРёРё. "
    "РқРө РІСӢРҙСғРјСӢРІР°Р№ С„Р°РәСӮСӢ; РөСҒР»Рё РҙР°РҪРҪСӢС… РҪРөРҙРҫСҒСӮР°СӮРҫСҮРҪРҫ, Р·Р°РҙР°Р№ РәРҫСҖРҫСӮРәРёР№ СғСӮРҫСҮРҪСҸСҺСүРёР№ РІРҫРҝСҖРҫСҒ."
)
VISION_SYSTEM_PROMPT = (
    "РўСӢ СҮС‘СӮРәРҫ РҫРҝРёСҒСӢРІР°РөСҲСҢ СҒРҫРҙРөСҖР¶РёРјРҫРө РёР·РҫРұСҖР°Р¶РөРҪРёР№: РҫРұСҠРөРәСӮСӢ, СӮРөРәСҒСӮ, СҒС…РөРјСӢ, РіСҖР°С„РёРәРё. "
    "РқРө РёРҙРөРҪСӮРёС„РёСҶРёСҖСғР№ Р»РёСҮРҪРҫСҒСӮРё Р»СҺРҙРөР№ Рё РҪРө РҝРёСҲРё РёРјРөРҪР°, РөСҒР»Рё РҫРҪРё РҪРө РҪР°РҝРөСҮР°СӮР°РҪСӢ РҪР° РёР·РҫРұСҖР°Р¶РөРҪРёРё. "
    "Р•СҒР»Рё РёР·РҫРұСҖР°Р¶РөРҪРёРө РјРөРҙРёСҶРёРҪСҒРәРҫРө, РҙРөР»Р°Р№ СӮРҫР»СҢРәРҫ СҒРҝСҖР°РІРҫСҮРҪСӢР№ СҖР°Р·РұРҫСҖ РІРёРҙРёРјРҫРіРҫ СӮРөРәСҒСӮР°/РҝСҖРёР·РҪР°РәРҫРІ, "
    "РҪРө СҒСӮР°РІСҢ РҙРёР°РіРҪРҫР· Рё РҪР°РҝРҫРјРёРҪР°Р№, СҮСӮРҫ РҪСғР¶РөРҪ РІСҖР°СҮ Рё РҫС„РёСҶРёР°Р»СҢРҪСӢР№ РҝСҖРҫСӮРҫРәРҫР»."
)

HELP_TEXT = globals().get("HELP_TEXT") or (
    "РҡРҫРјР°РҪРҙСӢ: /start, /chats, /newchat, /engines, /plans, /balance, /prices, /img <РҫРҝРёСҒР°РҪРёРө>, /mj <РҫРҝРёСҒР°РҪРёРө>, /voice_on, /voice_off, /diag_video, /diag_bg, /diag_face.\n"
    "РһСҒРҪРҫРІРҪСӢРө СҖРөР¶РёРјСӢ: рҹҺ“ РЈСҮС‘РұР°, рҹ’ј Р Р°РұРҫСӮР°/Р‘РёР·РҪРөСҒ, рҹ”Ҙ Р Р°Р·РІР»РөСҮРөРҪРёСҸ, рҹ©ә РңРөРҙРёСҶРёРҪР°, рҹ§  Р”РІРёР¶РәРё, рҹ’і Р‘Р°Р»Р°РҪСҒ.\n"
    "РӨРҫСӮРҫ/РұРёР·РҪРөСҒ: РҫР¶РёРІР»РөРҪРёРө СҮРөСҖРөР· Runway/Kling/Sora 2 РұРөР· Р»СҺРҙРөР№, РіРҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ, С„РҫСӮРҫвҶ’РІРёРҙРөРҫРәР»РёРҝ, Р·Р°РјРөРҪР° Р»РёСҶР°, СғРҙР°Р»РөРҪРёРө/Р·Р°РјРөРҪР° С„РҫРҪР°, СғРҙР°Р»РөРҪРёРө РІРҫРҙСҸРҪРҫРіРҫ Р·РҪР°РәР°/РҪР°РҙРҝРёСҒРё, СҖР°СҒСҲРёСҖРөРҪРёРө РәР°РҙСҖР°, СҖР°СҒРәР°РҙСҖРҫРІРәР°, Р°РҪР°Р»РёР· С„РҫСӮРҫ, Р»РҫРіРҫСӮРёРҝСӢ, РҝСҖРөР·РөРҪСӮР°СҶРёРё Рё PDF-РәР°СӮР°Р»РҫРіРё.\n"
    "Р”РҫРәСғРјРөРҪСӮСӢ, РіРҫР»РҫСҒ Рё РјСғР·СӢРәР°: PDF/EPUB/DOCX/TXT/FB2, СҒРІРҫРҙРәРё, РәРҫРҪСҒРҝРөРәСӮСӢ, СӮР°РұР»РёСҶСӢ, СҖРөСҮСҢ вҶ” СӮРөРәСҒСӮ, РҫР·РІСғСҮРәР°, РҝРөСҒРҪРё/РјРёРҪСғСҒРҫРІРәРё Suno.\n"
    "РңРөРҙРёСҶРёРҪР°: СҒРҝСҖР°РІРҫСҮРҪСӢР№ СҖР°Р·РұРҫСҖ РІСӢРҝРёСҒРҫРә, Р·Р°РәР»СҺСҮРөРҪРёР№, Р°РҪР°Р»РёР·РҫРІ, СҒРҪРёРјРәРҫРІ, РңР Рў/РҡРў. РӯСӮРҫ РҪРө РҙРёР°РіРҪРҫР· Рё РҪРө Р·Р°РјРөРҪР° РІСҖР°СҮСғ."
)
EXAMPLES_TEXT = globals().get("EXAMPLES_TEXT") or (
    "РҹСҖРёРјРөСҖСӢ:\n"
    "вҖў РһР¶РёРІРё С„РҫСӮРҫ: Р»С‘РіРәР°СҸ СғР»СӢРұРәР°, РІР·РіР»СҸРҙ РІ РәР°РјРөСҖСғ, РҝР»Р°РІРҪРҫРө РҙРІРёР¶РөРҪРёРө РәР°РјРөСҖСӢ, 5 СҒРөРәСғРҪРҙ, 9:16\n"
    "вҖў Р“РҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ: Р·Р°РіСҖСғР·Рё РҝРҫСҖСӮСҖРөСӮ, РҪР°Р¶РјРё рҹ—Ј Рё РҝСҖРёСҲР»Рё СӮРөРәСҒСӮ/voice/audio РҙР»СҸ СҖРөСҮРё\n"
    "вҖў РӨРҫСӮРҫ РІ РІРёРҙРөРҫРәР»РёРҝ: Р·Р°РіСҖСғР·Рё С„РҫСӮРҫ, РҪР°Р¶РјРё рҹҺө Рё РҫРҝРёСҲРё СҒСӮРёР»СҢ РәР»РёРҝР°/РјСғР·СӢРәРё, 5 СҒРөРәСғРҪРҙ, 9:16\n"
    "вҖў РЎРҙРөР»Р°Р№ РІРёРҙРөРҫ: РІРёР»Р»Р° РҪР° РұРөСҖРөРіСғ РјРҫСҖСҸ РҪР° РЎР°РјСғРё, Р·Р°РәР°СӮ, 10 СҒРөРәСғРҪРҙ, 16:9\n"
    "вҖў /img luxury villa in Koh Samui, tropical, cinematic"
)

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ Heuristics / intent в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
_SMALLTALK_RE = re.compile(r"^(РҝСҖРёРІРөСӮ|Р·РҙСҖР°РІСҒСӮРІСғР№|РҙРҫРұСҖСӢР№\s*(РҙРөРҪСҢ|РІРөСҮРөСҖ|СғСӮСҖРҫ)|С…Рё|hi|hello|РәР°Рә РҙРөР»Р°|СҒРҝР°СҒРёРұРҫ|РҝРҫРәР°)\b", re.I)
_NEWSY_RE     = re.compile(r"(РәРҫРіРҙР°|РҙР°СӮР°|РІСӢР№РҙРөСӮ|СҖРөР»РёР·|РҪРҫРІРҫСҒСӮ|РәСғСҖСҒ|СҶРөРҪР°|РҝСҖРҫРіРҪРҫР·|РҪР°Р№РҙРё|РҫС„РёСҶРёР°Р»|РҝРҫРіРҫРҙР°|СҒРөРіРҫРҙРҪСҸ|СӮСҖРөРҪРҙ|Р°РҙСҖРөСҒ|СӮРөР»РөС„РҫРҪ)", re.I)
_CAPABILITY_RE = re.compile(r"(РјРҫР¶(РөСҲСҢ|РҪРҫ|РөСӮРө)|СғРјРө(РөСҲСҢ|РөСӮРө)|СҒРҝРҫСҒРҫРұРөРҪ|РјРҫР¶РөСӮ\s+Р»Рё).{0,80}(Р°РҪР°Р»РёР·|СҖР°СҒРҝРҫР·РҪ|СҮРёСӮР°СӮСҢ|СҒРҫР·РҙР°(РІР°)?СӮ|РҙРөР»Р°(СӮСҢ)?|РҫР¶РёРІ|Р°РҪРёРјРёСҖ).{0,80}(С„РҫСӮРҫ|С„РҫСӮРҫРіСҖР°С„|РәР°СҖСӮРёРҪРә|РёР·РҫРұСҖР°Р¶РөРҪ|pdf|docx|epub|fb2|Р°СғРҙРёРҫ|РәРҪРёРі|РІРёРҙРөРҫ)", re.I)

_IMG_WORDS = r"(РәР°СҖСӮРёРҪ\w+|РёР·РҫРұСҖР°Р¶РөРҪ\w+|С„РҫСӮРҫ\w*|СҖРёСҒСғРҪРә\w+|Р»РҫРіРҫСӮРёРҝ\w*|Р»РҫРіРҫ\b|СҚРјРұР»РөРј\w+|РұСҖРөРҪРҙ\w+|С„РёСҖРјРөРҪРҪ\w+|image|picture|img\b|logo|banner|poster|brand|emblem)"
_VID_WORDS = r"(РІРёРҙРөРҫ|СҖРҫР»РёРә\w*|Р°РҪРёРјР°СҶРё\w*|shorts?|reels?|clip|video|vid\b)"


_PEOPLE_PROMPT_RE = re.compile(
    r"\b(СҮРөР»РҫРІРөРә|Р»СҺРҙРё|РјСғР¶СҮРёРҪ\w*|Р¶РөРҪСүРёРҪ\w*|РҙРөРІСғСҲРә\w*|РҝР°СҖРҪ\w*|СҖРөРұ[РөС‘]РҪРҫРә|РҙРөСӮ\w*|"
    r"РҙРІРҫСҖРҪРёРә\w*|РҝРҫР»РёСҶРөР№СҒРә\w*|СҒРҫР»РҙР°СӮ\w*|РІСҖР°СҮ\w*|Р°РәСӮ[РөС‘]СҖ\w*|РҝРөРІ[РөС‘]СҶ\w*|РҝРөСҖСҒРҫРҪР°Р¶\w*|"
    r"Р»РёСҶРҫ|РҝРҫСҖСӮСҖРөСӮ|СҒРөР»С„Рё|human|person|people|man|woman|girl|boy|child|children|actor|singer|worker|janitor)\b",
    re.IGNORECASE,
)

def _prompt_likely_has_people(prompt: str) -> bool:
    return bool(_PEOPLE_PROMPT_RE.search(prompt or ""))

def is_smalltalk(text: str) -> bool:
    t = (text or "").strip().lower()
    return bool(_SMALLTALK_RE.search(t))

def should_browse(text: str) -> bool:
    t = (text or "").strip().lower()
    if len(t) < 8:
        return False
    if "http://" in t or "https://" in t:
        return False
    return bool(_NEWSY_RE.search(t)) and not is_smalltalk(t)

_CREATE_CMD = r"(СҒРҙРөР»Р°(Р№|Р№СӮРө|СӮСҢ)|СҒРҫР·РҙР°(Р№|Р№СӮРө|СӮСҢ)|СҒРіРөРҪРөСҖРёСҖСғ(Р№|Р№СӮРө)|РҪР°СҖРёСҒСғ(Р№|Р№СӮРө)|РҪСғР¶РҪРҫ\s+СҒРҙРөР»Р°СӮСҢ|С…РҫСҮСғ\s+СҒРҫР·РҙР°СӮСҢ|render|generate|create|make)"
_PREFIXES_VIDEO = [r"^" + _CREATE_CMD + r"\s+РІРёРҙРөРҫ", r"^video\b", r"^reels?\b", r"^shorts?\b"]
_PREFIXES_IMAGE = [r"^" + _CREATE_CMD + r"\s+(?:РәР°СҖСӮРёРҪ\w+|РёР·РҫРұСҖР°Р¶РөРҪ\w+|С„РҫСӮРҫ\w+|СҖРёСҒСғРҪРә\w+|Р»РҫРіРҫСӮРёРҝ\w*|Р»РҫРіРҫ\b|СҚРјРұР»РөРј\w+|РұСҖРөРҪРҙ\w+|С„РёСҖРјРөРҪРҪ\w+)", r"^image\b", r"^picture\b", r"^img\b"]

def _strip_leading(s: str) -> str:
    return s.strip(" \n\t:вҖ”вҖ“-\"вҖңвҖқ'В«В»,.()[]")

def _after_match(text: str, match) -> str:
    return _strip_leading(text[match.end():])

def _looks_like_capability_question(tl: str) -> bool:
    if "?" in tl and re.search(_CAPABILITY_RE, tl):
        if not re.search(_CREATE_CMD, tl, re.I):
            return True
    m = re.search(r"\b(СӮСӢ|РІСӢ)?\s*РјРҫР¶(РөСҲСҢ|РҪРҫ|РөСӮРө)\b", tl)
    if m and re.search(_CAPABILITY_RE, tl) and not re.search(_CREATE_CMD, tl, re.I):
        return True
    return False

def detect_media_intent(text: str):
    if not text:
        return (None, "")
    t = text.strip()
    tl = t.lower()

    if _looks_like_capability_question(tl):
        return (None, "")

    for p in _PREFIXES_VIDEO:
        m = re.search(p, tl, re.I)
        if m:
            return ("video", _after_match(t, m))
    for p in _PREFIXES_IMAGE:
        m = re.search(p, tl, re.I)
        if m:
            return ("image", _after_match(t, m))

    if re.search(_CREATE_CMD, tl, re.I):
        if re.search(_VID_WORDS, tl, re.I):
            clean = re.sub(_VID_WORDS, "", tl, flags=re.I)
            clean = re.sub(_CREATE_CMD, "", clean, flags=re.I)
            return ("video", _strip_leading(clean))
        if re.search(_IMG_WORDS, tl, re.I):
            clean = re.sub(_IMG_WORDS, "", tl, flags=re.I)
            clean = re.sub(_CREATE_CMD, "", clean, flags=re.I)
            return ("image", _strip_leading(clean))

    m = re.match(r"^(img|image|picture)\s*[:\-]\s*(.+)$", tl)
    if m:
        return ("image", _strip_leading(t[m.end(1)+1:]))

    m = re.match(r"^(video|vid|reels?|shorts?)\s*[:\-]\s*(.+)$", tl)
    if m:
        return ("video", _strip_leading(t[m.end(1)+1:]))

    return (None, "")

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ OpenAI helpers в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def _oai_text_client():
    return oai_llm

def _pick_vision_model() -> str:
    try:
        mv = globals().get("OPENAI_VISION_MODEL")
        return (mv or OPENAI_MODEL).strip()
    except Exception:
        return OPENAI_MODEL

async def ask_openai_text(user_text: str, web_ctx: str = "", user_id: int | None = None, chat_id: int | None = None, extra_system: str = "") -> str:
    """
    РЈРҪРёРІРөСҖСҒР°Р»СҢРҪСӢР№ Р·Р°РҝСҖРҫСҒ Рә LLM:
    - РҝРҫРҙРҙРөСҖР¶РёРІР°РөСӮ OpenRouter (СҮРөСҖРөР· OPENAI_API_KEY = sk-or-...);
    - РҝСҖРёРҪСғРҙРёСӮРөР»СҢРҪРҫ СҲР»С‘СӮ JSON РІ UTF-8, СҮСӮРҫРұСӢ РҪРө РұСӢР»Рҫ ascii-РҫСҲРёРұРҫРә;
    - Р»РҫРіРёСҖСғРөСӮ HTTP-СҒСӮР°СӮСғСҒ Рё СӮРөР»Рҫ РҫСҲРёРұРәРё РІ Render-Р»РҫРіРё;
    - РҙРөР»Р°РөСӮ РҙРҫ 3 РҝРҫРҝСӢСӮРҫРә СҒ РҪРөРұРҫР»СҢСҲРҫР№ РҝР°СғР·РҫР№.
    """
    user_text = (user_text or "").strip()
    if not user_text:
        return "РҹСғСҒСӮРҫР№ Р·Р°РҝСҖРҫСҒ."

    messages = [{"role": "system", "content": SYSTEM_PROMPT + "\n\n" + _current_date_system_text()}]
    if extra_system:
        messages.append({"role": "system", "content": str(extra_system)[:3000]})
    if web_ctx:
        messages.append({
            "role": "system",
            "content": f"РҡРҫРҪСӮРөРәСҒСӮ РёР· РІРөРұ-РҝРҫРёСҒРәР°:\n{web_ctx}",
        })

    # РҡРҫСҖРҫСӮРәР°СҸ РҝР°РјСҸСӮСҢ РҙРёР°Р»РҫРіР°: РҝРҫСҒР»РөРҙРҪРёРө СҒРҫРҫРұСүРөРҪРёСҸ РҝРҫР»СҢР·РҫРІР°СӮРөР»СҸ Рё РұРҫСӮР°.
    # РӯСӮРҫ РҝРҫР·РІРҫР»СҸРөСӮ РҝРҫРҪРёРјР°СӮСҢ РҫСӮРІРөСӮСӢ РІСҖРҫРҙРө В«РұРёСӮРәРҫРёРҪВ» РәР°Рә РҝСҖРҫРҙРҫР»Р¶РөРҪРёРө РҝСҖРөРҙСӢРҙСғСүРөРіРҫ РІРҫРҝСҖРҫСҒР°,
    # Р° РҪРө РҪР°СҮРёРҪР°СӮСҢ РҪРҫРІСғСҺ СҒРҝСҖР°РІРәСғ СҒ РҫРҝСҖРөРҙРөР»РөРҪРёСҸ СӮРөСҖРјРёРҪР°.
    if CHAT_MEMORY_ENABLED and user_id and chat_id:
        messages.extend(_chat_memory_recent(user_id, chat_id, CHAT_MEMORY_MAX_MESSAGES))

    messages.append({"role": "user", "content": user_text})

    # в”Җв”Җ Р‘Р°Р·РҫРІСӢР№ URL в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
    # Р•СҒР»Рё РәР»СҺСҮ РҫСӮ OpenRouter РёР»Рё TEXT_PROVIDER=openrouter вҖ” СҲР»С‘Рј РҪР° OpenRouter
    provider = (TEXT_PROVIDER or "").strip().lower()
    if OPENAI_API_KEY.startswith("sk-or-") or provider == "openrouter":
        base_url = "https://openrouter.ai/api/v1"
    else:
        base_url = (OPENAI_BASE_URL or "").strip() or "https://api.openai.com/v1"

    # в”Җв”Җ Р—Р°РіРҫР»РҫРІРәРё в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
    headers = {
        "Authorization": f"Bearer {OPENAI_API_KEY}",
        "Content-Type": "application/json; charset=utf-8",
        "Accept-Charset": "utf-8",
    }

    # РЎР»СғР¶РөРұРҪСӢРө Р·Р°РіРҫР»РҫРІРәРё OpenRouter
    if "openrouter.ai" in base_url:
        if OPENROUTER_SITE_URL:
            headers["HTTP-Referer"] = OPENROUTER_SITE_URL
        if OPENROUTER_APP_NAME:
            headers["X-Title"] = OPENROUTER_APP_NAME

    last_err: Exception | None = None

    for attempt in range(3):
        try:
            async with httpx.AsyncClient(
                base_url=base_url,
                timeout=90.0,
            ) as client:
                resp = await client.post(
                    "/chat/completions",
                    json={
                        "model": OPENAI_MODEL,
                        "messages": messages,
                        "temperature": 0.6,
                    },
                    headers=headers,
                )

            # РӣРҫРіРёСҖСғРөРј РІСҒС‘, СҮСӮРҫ РҪРө 2xx
            if resp.status_code // 100 != 2:
                body_preview = resp.text[:800]
                log.warning(
                    "LLM HTTP %s from %s: %s",
                    resp.status_code,
                    base_url,
                    body_preview,
                )
                resp.raise_for_status()

            data = resp.json()
            txt = (data["choices"][0]["message"]["content"] or "").strip()
            if txt:
                return txt

        except Exception as e:
            last_err = e
            log.warning(
                "OpenAI/OpenRouter chat attempt %d failed: %s",
                attempt + 1,
                e,
            )
            await asyncio.sleep(0.8 * (attempt + 1))

    log.error("ask_openai_text failed after 3 attempts: %s", last_err)
    return (
        "вҡ пёҸ РЎРөР№СҮР°СҒ РҪРө РҝРҫР»СғСҮРёР»РҫСҒСҢ РҝРҫР»СғСҮРёСӮСҢ РҫСӮРІРөСӮ РҫСӮ РјРҫРҙРөР»Рё. "
        "РҜ РҪР° СҒРІСҸР·Рё вҖ” РҝРҫРҝСҖРҫРұСғР№ РҝРөСҖРөС„РҫСҖРјСғР»РёСҖРҫРІР°СӮСҢ Р·Р°РҝСҖРҫСҒ РёР»Рё РҝРҫРІСӮРҫСҖРёСӮСҢ СҮСғСӮСҢ РҝРҫР·Р¶Рө."
    )
    
async def ask_openai_vision(user_text: str, img_b64: str, mime: str) -> str:
    try:
        prompt = (user_text or "РһРҝРёСҲРё, СҮСӮРҫ РҪР° РёР·РҫРұСҖР°Р¶РөРҪРёРё Рё РәР°РәРҫР№ СӮР°Рј СӮРөРәСҒСӮ.").strip()
        model = _pick_vision_model()
        resp = _oai_text_client().chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": VISION_SYSTEM_PROMPT},
                {"role": "user", "content": [
                    {"type": "text", "text": prompt},
                    {"type": "image_url", "image_url": {"url": f"data:{mime};base64,{img_b64}"}}
                ]}
            ],
            temperature=0.4,
        )
        return (resp.choices[0].message.content or "").strip()
    except Exception as e:
        log.exception("Vision error: %s", e)
        return "РқРө СғРҙР°Р»РҫСҒСҢ РҝСҖРҫР°РҪР°Р»РёР·РёСҖРҫРІР°СӮСҢ РёР·РҫРұСҖР°Р¶РөРҪРёРө."


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РҹРҫР»СҢР·РҫРІР°СӮРөР»СҢСҒРәРёРө РҪР°СҒСӮСҖРҫР№РәРё (TTS) в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def _db_init_prefs():
    con = sqlite3.connect(DB_PATH)
    cur = con.cursor()
    cur.execute("""
    CREATE TABLE IF NOT EXISTS user_prefs (
        user_id INTEGER PRIMARY KEY,
        tts_on  INTEGER DEFAULT 0
    )""")
    con.commit(); con.close()

def _tts_get(user_id: int) -> bool:
    try:
        _db_init_prefs()
    except Exception:
        pass
    con = sqlite3.connect(DB_PATH); cur = con.cursor()
    cur.execute("INSERT OR IGNORE INTO user_prefs(user_id, tts_on) VALUES (?,0)", (user_id,))
    con.commit()
    cur.execute("SELECT tts_on FROM user_prefs WHERE user_id=?", (user_id,))
    row = cur.fetchone(); con.close()
    return bool(row and row[0])

def _tts_set(user_id: int, on: bool):
    try:
        _db_init_prefs()
    except Exception:
        pass
    con = sqlite3.connect(DB_PATH); cur = con.cursor()
    cur.execute("INSERT OR IGNORE INTO user_prefs(user_id, tts_on) VALUES (?,?)", (user_id, 1 if on else 0))
    cur.execute("UPDATE user_prefs SET tts_on=? WHERE user_id=?", (1 if on else 0, user_id))
    con.commit(); con.close()


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РқР°РҙС‘Р¶РҪСӢР№ TTS СҮРөСҖРөР· REST (OGG/Opus) в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def _tts_bytes_sync(text: str, audio_format: str = "opus", voice: str | None = None) -> bytes | None:
    try:
        if not OPENAI_TTS_KEY:
            return None
        if OPENAI_TTS_KEY.startswith("sk-or-"):
            log.error("TTS key looks like OpenRouter (sk-or-...). Provide a real OpenAI key in OPENAI_TTS_KEY.")
            return None
        url = f"{OPENAI_TTS_BASE_URL.rstrip('/')}/audio/speech"
        payload = {
            "model": OPENAI_TTS_MODEL,
            "voice": (voice or OPENAI_TTS_VOICE or "alloy"),
            "input": text,
            "response_format": audio_format  # mp3 for avatar, opus for Telegram voice
        }
        headers = {
            "Authorization": f"Bearer {OPENAI_TTS_KEY}",
            "Content-Type": "application/json"
        }
        r = httpx.post(url, headers=headers, json=payload, timeout=60.0)
        r.raise_for_status()
        data = r.content if r.content else None
        if data:
            log.info("TTS bytes: %s bytes", len(data))
        return data
    except Exception as e:
        log.exception("TTS HTTP error: %s", e)
        return None

async def maybe_tts_reply(update: Update, context: ContextTypes.DEFAULT_TYPE, text: str):
    user_id = update.effective_user.id
    if not _tts_get(user_id):
        return
    text = (text or "").strip()
    if not text:
        return
    if len(text) > TTS_MAX_CHARS:
        with contextlib.suppress(Exception):
            await update.effective_message.reply_text(
                f"рҹ”Ү РһР·РІСғСҮРәР° РІСӢРәР»СҺСҮРөРҪР° РҙР»СҸ СҚСӮРҫРіРҫ СҒРҫРҫРұСүРөРҪРёСҸ: СӮРөРәСҒСӮ РҙР»РёРҪРҪРөРө {TTS_MAX_CHARS} СҒРёРјРІРҫР»РҫРІ."
            )
        return
    if not OPENAI_TTS_KEY:
        return
    try:
        with contextlib.suppress(Exception):
            await context.bot.send_chat_action(update.effective_chat.id, ChatAction.UPLOAD_VOICE)
        audio = await asyncio.to_thread(_tts_bytes_sync, text, "opus")
        if not audio:
            with contextlib.suppress(Exception):
                await update.effective_message.reply_text("рҹ”Ү РқРө СғРҙР°Р»РҫСҒСҢ СҒРёРҪСӮРөР·РёСҖРҫРІР°СӮСҢ РіРҫР»РҫСҒ.")
            return
        bio = BytesIO(audio); bio.seek(0); bio.name = "say.ogg"
        await update.effective_message.reply_voice(voice=InputFile(bio), caption=(text if TTS_VOICE_CAPTION else None))
    except Exception as e:
        log.exception("maybe_tts_reply error: %s", e)

async def cmd_voice_on(update: Update, context: ContextTypes.DEFAULT_TYPE):
    _tts_set(update.effective_user.id, True)
    await update.effective_message.reply_text(f"рҹ”Ҡ РһР·РІСғСҮРәР° РІРәР»СҺСҮРөРҪР°. РӣРёРјРёСӮ {TTS_MAX_CHARS} СҒРёРјРІРҫР»РҫРІ РҪР° РҫСӮРІРөСӮ.")

async def cmd_voice_off(update: Update, context: ContextTypes.DEFAULT_TYPE):
    _tts_set(update.effective_user.id, False)
    await update.effective_message.reply_text("рҹ”Ҳ РһР·РІСғСҮРәР° РІСӢРәР»СҺСҮРөРҪР°.")

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ Speech-to-Text (STT) вҖў OpenAI Whisper/4o-mini-transcribe в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
from openai import OpenAI as _OpenAI_STT

OPENAI_STT_MODEL    = (os.getenv("OPENAI_STT_MODEL") or "whisper-1").strip()
OPENAI_STT_KEY      = (os.getenv("OPENAI_STT_KEY") or os.getenv("OPENAI_API_KEY") or "").strip()
OPENAI_STT_BASE_URL = (os.getenv("OPENAI_STT_BASE_URL") or "https://api.openai.com/v1").rstrip("/")

def _oai_stt_client():
    return _OpenAI_STT(api_key=OPENAI_STT_KEY, base_url=OPENAI_STT_BASE_URL)

async def _stt_transcribe_bytes(filename: str, raw: bytes) -> str:
    last_err = None
    for attempt in range(3):
        try:
            bio = BytesIO(raw)
            bio.name = filename
            bio.seek(0)
            resp = _oai_stt_client().audio.transcriptions.create(
                model=OPENAI_STT_MODEL,
                file=bio,
            )
            text = (getattr(resp, "text", "") or "").strip()
            if text:
                return text
        except Exception as e:
            last_err = e
            log.warning("STT attempt %d failed: %s", attempt+1, e)
            await asyncio.sleep(0.8 * (attempt + 1))
    log.error("STT failed: %s", last_err)
    return ""

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РҘРөРҪРҙР»РөСҖ РіРҫР»РҫСҒРҫРІСӢС…/Р°СғРҙРёРҫ в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
async def handle_voice(update: Update, context: ContextTypes.DEFAULT_TYPE):
    msg = update.effective_message
    voice = getattr(msg, "voice", None)
    audio = getattr(msg, "audio", None)
    media = voice or audio
    if not media:
        await msg.reply_text("РқРө РҪР°СҲС‘Р» РіРҫР»РҫСҒРҫРІРҫР№ С„Р°Р№Р».")
        return

    # РЎРәР°СҮРёРІР°РөРј С„Р°Р№Р»
    try:
        with contextlib.suppress(Exception):
            await context.bot.send_chat_action(update.effective_chat.id, ChatAction.TYPING)

        tg_file = await context.bot.get_file(media.file_id)
        buf = BytesIO()
        await tg_file.download_to_memory(out=buf)
        raw = buf.getvalue()

        mime = (getattr(media, "mime_type", "") or "").lower()
        if "ogg" in mime or "opus" in mime:
            filename = "voice.ogg"
        elif "webm" in mime:
            filename = "voice.webm"
        elif "wav" in mime:
            filename = "voice.wav"
        elif "mp3" in mime or "mpeg" in mime or "mpga" in mime:
            filename = "voice.mp3"
        else:
            filename = "voice.ogg"

    except Exception as e:
        log.exception("TG download error: %s", e)
        await msg.reply_text("РқРө СғРҙР°Р»РҫСҒСҢ СҒРәР°СҮР°СӮСҢ РіРҫР»РҫСҒРҫРІРҫРө СҒРҫРҫРұСүРөРҪРёРө. РҹРҫРҝСҖРҫРұСғР№СӮРө РҫСӮРҝСҖР°РІРёСӮСҢ РөРіРҫ РөСүС‘ СҖР°Р·, Р»СғСҮСҲРө РҪРө РҝРөСҖРөСҒСӢР»Р°СҸ, Р° Р·Р°РіСҖСғР·РёРІ РҪР°РҝСҖСҸРјСғСҺ.")
        return

    # Р•СҒР»Рё РҝРҫР»СҢР·РҫРІР°СӮРөР»СҢ РҪР° СҲР°РіРө РІСӢРұРҫСҖР° РіРҫР»РҫСҒР° РҝСҖРёСҒР»Р°Р» voice/audio, РёСҒРҝРҫР»СҢР·СғРөРј СҖРөР°Р»СҢРҪСӢР№ РіРҫР»РҫСҒ РұРөР· TTS.
    if context.user_data.get("awaiting_avatar_voice_choice"):
        context.user_data.pop("awaiting_avatar_voice_choice", None)
        context.user_data["awaiting_avatar_script"] = True

    # Р•СҒР»Рё РҝРҫР»СҢР·РҫРІР°СӮРөР»СҢ РҪР°Р¶Р°Р» В«Р“РҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖВ», voice/audio СҒСӮР°РҪРҫРІРёСӮСҒСҸ СҖРөСҮСҢСҺ РҙР»СҸ Р·Р°РіСҖСғР¶РөРҪРҪРҫРіРҫ РҝРҫСҖСӮСҖРөСӮР°.
    if context.user_data.get("awaiting_avatar_script"):
        img = _get_cached_photo(update.effective_user.id)
        if not img:
            _clear_avatar_wait(context)
            await msg.reply_text("РЎРҪР°СҮР°Р»Р° Р·Р°РіСҖСғР·РёСӮРө РҝРҫСҖСӮСҖРөСӮ СҮРөР»РҫРІРөРәР°, Р·Р°СӮРөРј РҪР°Р¶РјРёСӮРө рҹ—Ј Р“РҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ.")
            return
        _clear_avatar_wait(context)
        await _start_talking_avatar(
            update, context, img,
            audio_bytes=raw,
            audio_filename=filename,
            audio_file_url=getattr(tg_file, "file_path", "") or "",
            audio_mime=mime,
        )
        return

    # STT
    transcript = await _stt_transcribe_bytes(filename, raw)
    if not transcript:
        await msg.reply_text("РһСҲРёРұРәР° РҝСҖРё СҖР°СҒРҝРҫР·РҪР°РІР°РҪРёРё СҖРөСҮРё.")
        return

    transcript = transcript.strip()

    # РҹРҫ СғРјРҫР»СҮР°РҪРёСҺ РқР• РҝРҫРәР°Р·СӢРІР°РөРј РҫСӮРҙРөР»СҢРҪСғСҺ СҖР°СҒСҲРёС„СҖРҫРІРәСғ voice, СҮСӮРҫРұСӢ РҪРө РұСӢР»Рҫ 3 СҒРҫРҫРұСүРөРҪРёР№:
    # 1) СӮРөРәСҒСӮРҫРІСӢР№ РҫСӮРІРөСӮ, 2) voice-РҫР·РІСғСҮРәР°, 3) РҙСғРұР»СҢ СҖР°СҒСҲРёС„СҖРҫРІРәРё.
    if STT_ECHO_TRANSCRIPT:
        with contextlib.suppress(Exception):
            await msg.reply_text(f"рҹ—ЈпёҸ Р Р°СҒРҝРҫР·РҪР°Р»: {transcript}")

    # вҖ”вҖ”вҖ” РҡРӣР®Р§Р•Р’РһРҷ РңРһРңР•РқРў вҖ”вҖ”вҖ”
    # Р‘РҫР»СҢСҲРө РқР• СҒРҫР·РҙР°С‘Рј С„РөР№РәРҫРІСӢР№ Update, РҪРө Р»РөР·РөРј РІ Message.text вҖ” СҚСӮРҫ Р·Р°РҝСҖРөСүРөРҪРҫ РІ Telegram API
    # РўРөРҝРөСҖСҢ РјСӢ РёСҒРҝРҫР»СҢР·СғРөРј РұРөР·РҫРҝР°СҒРҪСӢР№ РҝСҖРҫРәСҒРё-РјРөСӮРҫРҙ, РәРҫСӮРҫСҖСӢР№ СҒРҫР·РҙР°С‘СӮ РІСҖРөРјРөРҪРҪСӢР№ message-РҫРұСҠРөРәСӮ
    try:
        await on_text_with_text(update, context, transcript)
    except Exception as e:
        log.exception("Voice->text handler error: %s", e)
        await msg.reply_text("РЈРҝСҒ, РҝСҖРҫРёР·РҫСҲР»Р° РҫСҲРёРұРәР°. РҜ СғР¶Рө СҖР°Р·РұРёСҖР°СҺСҒСҢ.")
        
# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РҳР·РІР»РөСҮРөРҪРёРө СӮРөРәСҒСӮР° РёР· РҙРҫРәСғРјРөРҪСӮРҫРІ в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def _safe_decode_txt(b: bytes) -> str:
    for enc in ("utf-8","cp1251","latin-1"):
        try:
            return b.decode(enc)
        except Exception:
            continue
    return b.decode("utf-8", errors="ignore")

def _extract_pdf_text(data: bytes) -> str:
    try:
        import PyPDF2
        rd = PyPDF2.PdfReader(BytesIO(data))
        parts = []
        for p in rd.pages:
            try:
                parts.append(p.extract_text() or "")
            except Exception:
                continue
        t = "\n".join(parts).strip()
        if t:
            return t
    except Exception:
        pass
    try:
        from pdfminer_high_level import extract_text as pdfminer_extract_text  # may not exist
    except Exception:
        pdfminer_extract_text = None  # type: ignore
    if pdfminer_extract_text:
        try:
            return (pdfminer_extract_text(BytesIO(data)) or "").strip()
        except Exception:
            pass
    try:
        import fitz
        doc = fitz.open(stream=data, filetype="pdf")
        txt = []
        for page in doc:
            try:
                txt.append(page.get_text("text"))
            except Exception:
                continue
        return "\n".join(txt)
    except Exception:
        pass
    return ""

def _extract_epub_text(data: bytes) -> str:
    try:
        from ebooklib import epub
        from bs4 import BeautifulSoup
        book = epub.read_epub(BytesIO(data))
        chunks = []
        for item in book.get_items():
            if item.get_type() == 9:  # DOCUMENT
                try:
                    soup = BeautifulSoup(item.get_content(), "html.parser")
                    txt = soup.get_text(separator=" ", strip=True)
                    if txt:
                        chunks.append(txt)
                except Exception:
                    continue
        return "\n".join(chunks).strip()
    except Exception:
        return ""

def _extract_docx_text(data: bytes) -> str:
    try:
        import docx
        doc = docx.Document(BytesIO(data))
        return "\n".join(p.text for p in doc.paragraphs).strip()
    except Exception:
        return ""

def _extract_fb2_text(data: bytes) -> str:
    try:
        import xml.etree.ElementTree as ET
        root = ET.fromstring(data)
        texts = []
        for elem in root.iter():
            if elem.text and elem.text.strip():
                texts.append(elem.text.strip())
        return " ".join(texts).strip()
    except Exception:
        return ""

def extract_text_from_document(data: bytes, filename: str) -> tuple[str, str]:
    name = (filename or "").lower()
    if name.endswith(".pdf"):  return _extract_pdf_text(data),  "PDF"
    if name.endswith(".epub"): return _extract_epub_text(data), "EPUB"
    if name.endswith(".docx"): return _extract_docx_text(data), "DOCX"
    if name.endswith(".fb2"):  return _extract_fb2_text(data),  "FB2"
    if name.endswith(".txt"):  return _safe_decode_txt(data),    "TXT"
    if name.endswith((".mobi",".azw",".azw3")): return "", "MOBI/AZW"
    decoded = _safe_decode_txt(data)
    return decoded if decoded else "", "UNKNOWN"


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РЎСғРјРјР°СҖРёР·Р°СҶРёСҸ РҙР»РёРҪРҪСӢС… СӮРөРәСҒСӮРҫРІ в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
async def _summarize_chunk(text: str, query: str | None = None) -> str:
    prefix = "РЎСғРјРјРёСҖСғР№ РәСҖР°СӮРәРҫ РҝРҫ РҝСғРҪРәСӮР°Рј РҫСҒРҪРҫРІРҪРҫРө РёР· С„СҖР°РіРјРөРҪСӮР° РҙРҫРәСғРјРөРҪСӮР° РҪР° СҖСғСҒСҒРәРҫРј:\n"
    if query:
        prefix = (f"РЎСғРјРјРёСҖСғР№ С„СҖР°РіРјРөРҪСӮ СҒ СғСҮС‘СӮРҫРј СҶРөР»Рё: {query}\n"
                  f"Р”Р°Р№ РҫСҒРҪРҫРІРҪСӢРө СӮРөР·РёСҒСӢ, С„Р°РәСӮСӢ, СҶРёС„СҖСӢ. Р СғСҒСҒРәРёР№ СҸР·СӢРә.\n")
    prompt = prefix + text
    return await ask_openai_text(prompt)

async def summarize_long_text(full_text: str, query: str | None = None) -> str:
    max_chunk = 8000
    text = full_text.strip()
    if len(text) <= max_chunk:
        return await _summarize_chunk(text, query=query)
    parts = []
    i = 0
    while i < len(text) and len(parts) < 8:
        parts.append(text[i:i+max_chunk]); i += max_chunk
    partials = [await _summarize_chunk(p, query=query) for p in parts]
    combined = "\n\n".join(f"- РӨСҖР°РіРјРөРҪСӮ {idx+1}:\n{s}" for idx, s in enumerate(partials))
    final_prompt = ("РһРұСҠРөРҙРёРҪРё СӮРөР·РёСҒСӢ РҝРҫ С„СҖР°РіРјРөРҪСӮР°Рј РІ СҶРөР»СҢРҪРҫРө СҖРөР·СҺРјРө РҙРҫРәСғРјРөРҪСӮР°: 1) 5вҖ“10 РіР»Р°РІРҪСӢС… РҝСғРҪРәСӮРҫРІ; "
                    "2) РәР»СҺСҮРөРІСӢРө СҶРёС„СҖСӢ/СҒСҖРҫРәРё; 3) РІСӢРІРҫРҙ/СҖРөРәРҫРјРөРҪРҙР°СҶРёРё. Р СғСҒСҒРәРёР№ СҸР·СӢРә.\n\n" + combined)
    return await ask_openai_text(final_prompt)


# ======= РҗРҪР°Р»РёР· РҙРҫРәСғРјРөРҪСӮРҫРІ (PDF/EPUB/DOCX/FB2/TXT) =======
async def on_doc_analyze(update: Update, context: ContextTypes.DEFAULT_TYPE):
    try:
        if not update.message or not update.message.document:
            return
        doc = update.message.document
        tg_file = await doc.get_file()
        data = await tg_file.download_as_bytearray()
        text, kind = extract_text_from_document(bytes(data), doc.file_name or "file")
        if not text.strip():
            await update.effective_message.reply_text(f"РқРө СғРҙР°Р»РҫСҒСҢ РёР·РІР»РөСҮСҢ СӮРөРәСҒСӮ РёР· {kind}.")
            return
        caption = (update.message.caption or "").strip()
        goal = caption or None
        if _should_route_medical(context, update.effective_user.id, caption, doc.file_name or "file"):
            await _medical_analyze_text(update, context, text, goal=goal)
            _clear_medicine_wait(context)
            with contextlib.suppress(Exception):
                _mode_track_set(update.effective_user.id, "")
            return

        await update.effective_message.reply_text(f"рҹ“„ РҳР·РІР»РөРәР°СҺ СӮРөРәСҒСӮ ({kind}), РіРҫСӮРҫРІР»СҺ РәРҫРҪСҒРҝРөРәСӮвҖҰ")
        summary = await summarize_long_text(text, query=goal)
        summary = summary or "Р“РҫСӮРҫРІРҫ."
        await update.effective_message.reply_text(summary)
        await maybe_tts_reply(update, context, summary[:TTS_MAX_CHARS])
    except Exception as e:
        log.exception("on_doc_analyze error: %s", e)
    # РҪРёСҮРөРіРҫ РҪРө РұСҖРҫСҒР°РөРј РҪР°СҖСғР¶Сғ



# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ OpenAI Images (РіРөРҪРөСҖР°СҶРёСҸ РәР°СҖСӮРёРҪРҫРә) в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
async def _generate_openai_image_bytes(prompt: str) -> bytes | None:
    try:
        try:
            resp = oai_img.images.generate(model=IMAGES_MODEL, prompt=prompt, size="1024x1024", quality=OPENAI_IMAGE_QUALITY, n=1)
        except Exception:
            # Compatibility with proxies that do not expose the quality argument.
            resp = oai_img.images.generate(model=IMAGES_MODEL, prompt=prompt, size="1024x1024", n=1)
        b64 = resp.data[0].b64_json
        return base64.b64decode(b64) if b64 else None
    except Exception as e:
        log.exception("OpenAI image generation error: %s", e)
        return None

async def _do_img_generate(update: Update, context: ContextTypes.DEFAULT_TYPE, prompt: str) -> bool:
    try:
        await context.bot.send_chat_action(update.effective_chat.id, ChatAction.UPLOAD_PHOTO)
        img_bytes = await _generate_openai_image_bytes(prompt)
        if not img_bytes:
            await update.effective_message.reply_text("РқРө СғРҙР°Р»РҫСҒСҢ СҒРҫР·РҙР°СӮСҢ РёР·РҫРұСҖР°Р¶РөРҪРёРө.")
            return False
        await update.effective_message.reply_photo(
            photo=img_bytes,
            caption=f"Р“РҫСӮРҫРІРҫ вң…\nР”РІРёР¶РҫРә: OpenAI Images\nР—Р°РҝСҖРҫСҒ: {prompt}",
        )
        return True
    except Exception as e:
        log.exception("IMG gen error: %s", e)
        await update.effective_message.reply_text("РқРө СғРҙР°Р»РҫСҒСҢ СҒРҫР·РҙР°СӮСҢ РёР·РҫРұСҖР°Р¶РөРҪРёРө.")
        return False


def _is_text_sensitive_image_prompt(prompt: str) -> bool:
    return bool(re.search(
        r"(logo|brand|branding|label|poster|flyer|banner|packag|box|bottle|menu|sign|screen|ui|ux|interface|"
        r"Р»РҫРіРҫСӮРёРҝ|Р»РҫРіРҫ|РұСҖРөРҪРҙ|СҚСӮРёРәРөСӮ|СғРҝР°РәРҫРІ|РұСғСӮСӢР»|РәРҫСҖРҫРұ|РҝР»Р°РәР°СӮ|РұР°РҪРҪРөСҖ|С„Р»Р°РөСҖ|РјРөРҪСҺ|РІСӢРІРөСҒРә|СҚРәСҖР°РҪ|РёРҪСӮРөСҖС„РөР№СҒ)",
        prompt or "",
        re.I,
    ))


def _is_atmospheric_image_prompt(prompt: str) -> bool:
    return bool(re.search(
        r"(cinematic|editorial|mood|atmospher|luxury|premium|fashion|concept art|fantasy|epic|dramatic|"
        r"РәРёРҪРҫСҲРҪ|РәРёРҪРөРјР°СӮРҫРіСҖР°С„|Р°СӮРјРҫСҒС„РөСҖ|Р»СҺРәСҒ|РҝСҖРөРјРёСғРј|РјСғРҙРұРҫСҖРҙ|С„СҚСҲРҪ|Р°СҖСӮ|РәРҫРҪСҶРөРҝСӮ|СҚРҝРёСҮ|РҙСҖР°РјР°СӮ)",
        prompt or "",
        re.I,
    ))


def _resolve_general_image_engine(prompt: str, requested_engine: str = "auto") -> tuple[str, str]:
    requested_engine = (requested_engine or "auto").strip().lower()
    if requested_engine == "openai":
        return ("openai", "РҝРҫР»СҢР·РҫРІР°СӮРөР»СҢ РІСӢРұСҖР°Р» СӮРҫСҮРҪСӢР№ РҙРІРёР¶РҫРә")
    if requested_engine == "midjourney":
        return ("midjourney", "РҝРҫР»СҢР·РҫРІР°СӮРөР»СҢ РІСӢРұСҖР°Р» С…СғРҙРҫР¶РөСҒСӮРІРөРҪРҪСӢР№ РҙРІРёР¶РҫРә")

    if _is_text_sensitive_image_prompt(prompt):
        return ("openai", "РҫРұРҪР°СҖСғР¶РөРҪ Р·Р°РҝСҖРҫСҒ РҪР° СӮРөРәСҒСӮ, Р»РҫРіРҫСӮРёРҝ, СғРҝР°РәРҫРІРәСғ РёР»Рё СӮРҫСҮРҪСғСҺ РіСҖР°С„РёРәСғ")
    if _is_atmospheric_image_prompt(prompt) and MIDJOURNEY_ENABLED and COMET_API_KEY:
        return ("midjourney", "РҫРұРҪР°СҖСғР¶РөРҪ Р°СӮРјРҫСҒС„РөСҖРҪСӢР№/РәРёРҪРөРјР°СӮРҫРіСҖР°С„РёСҮРҪСӢР№ С…СғРҙРҫР¶РөСҒСӮРІРөРҪРҪСӢР№ Р·Р°РҝСҖРҫСҒ")
    return ("openai", "СғРҪРёРІРөСҖСҒР°Р»СҢРҪСӢР№ СҖРөР¶РёРј РҝРҫ СғРјРҫР»СҮР°РҪРёСҺ")


def _image_engine_price_label(engine: str) -> str:
    engine = (engine or "").lower()
    cost = MIDJOURNEY_UNIT_COST_USD if engine == "midjourney" else IMG_COST_USD
    return f"{_retail_credits(cost)} РәСҖ."


def _image_engine_choice_kb(aid: str, prompt: str = "") -> InlineKeyboardMarkup:
    auto_engine, _ = _resolve_general_image_engine(prompt, "auto")
    auto_title = "Midjourney" if auto_engine == "midjourney" else "OpenAI"
    return InlineKeyboardMarkup([
        [InlineKeyboardButton(f"вңЁ РҗРІСӮРҫ: {auto_title} В· {_image_engine_price_label(auto_engine)}", callback_data=f"chooseimg:auto:{aid}")],
        [InlineKeyboardButton(f"рҹ–ј OpenAI Images В· {_image_engine_price_label('openai')}", callback_data=f"chooseimg:openai:{aid}")],
        [InlineKeyboardButton(f"рҹҺЁ Midjourney В· {_image_engine_price_label('midjourney')}", callback_data=f"chooseimg:midjourney:{aid}")],
    ])


async def _ask_image_engine_choice(update: Update, context: ContextTypes.DEFAULT_TYPE, prompt: str):
    aid = _new_aid()
    _pending_actions[aid] = {"kind": "image_generate", "prompt": prompt}
    mid_note = "РҙРҫСҒСӮСғРҝРөРҪ" if MIDJOURNEY_ENABLED and COMET_API_KEY else "РҪРөРҙРҫСҒСӮСғРҝРөРҪ СҒРөР№СҮР°СҒ"
    msg = (
        "РҡР°РәРҫР№ РҙРІРёР¶РҫРә РёСҒРҝРҫР»СҢР·РҫРІР°СӮСҢ РҙР»СҸ РёР·РҫРұСҖР°Р¶РөРҪРёСҸ?\n"
        f"Р—Р°РҝСҖРҫСҒ: В«{prompt}В»\n\n"
        "вңЁ РҗРІСӮРҫ вҖ” РұРҫСӮ СҒР°Рј РІСӢРұРөСҖРөСӮ Р»СғСҮСҲРёР№ РјР°СҖСҲСҖСғСӮ РҝРҫРҙ Р·Р°РҙР°СҮСғ. РһРұСӢСҮРҪРҫ СҚСӮРҫ СҒР°РјСӢР№ СғРҙРҫРұРҪСӢР№ РІР°СҖРёР°РҪСӮ.\n"
        "рҹ–ј OpenAI Images вҖ” РәРҫРіРҙР° РҪСғР¶РҪР° СӮРҫСҮРҪРҫСҒСӮСҢ: РҝСҖРөРҙРјРөСӮСӢ, РҝСҖРҫРҙСғРәСӮСӢ, РұР°РҪРҪРөСҖСӢ, Р»РҫРіРҫСӮРёРҝСӢ, РҝРҫРҪСҸСӮРҪР°СҸ РәРҫРјРҝРҫР·РёСҶРёСҸ, СӮРөРәСҒСӮРҫРІСӢРө СҚР»РөРјРөРҪСӮСӢ Рё РәРҫРјРјРөСҖСҮРөСҒРәРёРө РјР°РәРөСӮСӢ.\n"
        f"рҹҺЁ Midjourney вҖ” РәРҫРіРҙР° РҪСғР¶РөРҪ Р°СӮРјРҫСҒС„РөСҖРҪСӢР№, С…СғРҙРҫР¶РөСҒСӮРІРөРҪРҪСӢР№, РәРёРҪРөРјР°СӮРҫРіСҖР°С„РёСҮРҪСӢР№ РёР»Рё fashion/editorial СҖРөР·СғР»СҢСӮР°СӮ. РЎСӮР°СӮСғСҒ СҒРөР№СҮР°СҒ: {mid_note}.\n\n"
        "РҹРҫРҙСҒРәР°Р·РәР°: /img вҖ” РұСӢСҒСӮСҖСӢР№ Р·Р°РҝСғСҒРә СҮРөСҖРөР· OpenAI Images, /mj вҖ” РҝСҖСҸРјРҫР№ Р·Р°РҝСғСҒРә Midjourney."
    )
    await update.effective_message.reply_text(msg, reply_markup=_image_engine_choice_kb(aid, prompt))


async def _run_selected_image_generation(update: Update, context: ContextTypes.DEFAULT_TYPE, prompt: str, requested_engine: str) -> bool:
    engine, reason = _resolve_general_image_engine(prompt, requested_engine)
    provider_cost = MIDJOURNEY_UNIT_COST_USD if engine == "midjourney" else IMG_COST_USD
    engine_title = "Midjourney" if engine == "midjourney" else "OpenAI Images"

    async def _go():
        try:
            await context.bot.send_chat_action(update.effective_chat.id, ChatAction.UPLOAD_PHOTO)
        except Exception:
            pass
        try:
            if engine == "midjourney":
                if not MIDJOURNEY_ENABLED:
                    await update.effective_message.reply_text("вқҢ Midjourney РҫСӮРәР»СҺСҮС‘РҪ РҪР°СҒСӮСҖРҫР№РәРҫР№ MIDJOURNEY_ENABLED=0.")
                    return False
                if not COMET_API_KEY:
                    await update.effective_message.reply_text("вқҢ Midjourney СҒРөР№СҮР°СҒ РҪРөРҙРҫСҒСӮСғРҝРөРҪ: РҪРө Р·Р°РҙР°РҪ COMET_API_KEY.")
                    return False
                await update.effective_message.reply_text(f"рҹҺЁ Р”РІРёР¶РҫРә: {engine_title}. Р—Р°РҝСғСҒРәР°СҺ РіРөРҪРөСҖР°СҶРёСҺвҖҰ")
                img, task_id = await _midjourney_generate_image_bytes(prompt)
                if not img:
                    await update.effective_message.reply_text("вқҢ Midjourney РҪРө РІРөСҖРҪСғР» РёР·РҫРұСҖР°Р¶РөРҪРёРө.")
                    return False
                bio = BytesIO(img)
                bio.name = f"midjourney_{task_id}.jpg"
                caption = f"Р“РҫСӮРҫРІРҫ вң…\nР”РІРёР¶РҫРә: {engine_title}\nРҹСҖРёСҮРёРҪР° РІСӢРұРҫСҖР°: {reason}\nР—Р°РҝСҖРҫСҒ: {prompt}"
                try:
                    await update.effective_message.reply_photo(photo=InputFile(bio), caption=caption[:1024])
                except Exception:
                    bio.seek(0)
                    await update.effective_message.reply_document(document=InputFile(bio), caption=caption[:1024])
                return True

            await update.effective_message.reply_text(f"рҹ–ј Р”РІРёР¶РҫРә: {engine_title}. Р—Р°РҝСғСҒРәР°СҺ РіРөРҪРөСҖР°СҶРёСҺвҖҰ")
            img = await _generate_openai_image_bytes(prompt)
            if not img:
                await update.effective_message.reply_text("РқРө СғРҙР°Р»РҫСҒСҢ СҒРҫР·РҙР°СӮСҢ РёР·РҫРұСҖР°Р¶РөРҪРёРө СҮРөСҖРөР· OpenAI Images.")
                return False
            caption = f"Р“РҫСӮРҫРІРҫ вң…\nР”РІРёР¶РҫРә: {engine_title}\nРҹСҖРёСҮРёРҪР° РІСӢРұРҫСҖР°: {reason}\nР—Р°РҝСҖРҫСҒ: {prompt}"
            await update.effective_message.reply_photo(photo=img, caption=caption[:1024])
            return True
        except Exception as e:
            log.exception("Selected image generation failed: %s", e)
            await update.effective_message.reply_text(f"вқҢ РһСҲРёРұРәР° РіРөРҪРөСҖР°СҶРёРё РёР·РҫРұСҖР°Р¶РөРҪРёСҸ ({engine_title}). РҹРҫРҝСҖРҫРұСғР№СӮРө РөСүС‘ СҖР°Р· РҝРҫР·Р¶Рө.")
            return False

    await _try_pay_then_do(
        update,
        context,
        update.effective_user.id,
        "img",
        provider_cost,
        _go,
        remember_kind="img_generate",
        remember_payload={"prompt": prompt[:1000], "requested_engine": requested_engine, "resolved_engine": engine},
    )
    return True

async def _luma_generate_image_bytes(prompt: str) -> bytes | None:
    if not LUMA_IMG_BASE_URL or not LUMA_API_KEY:
        # С„РҫР»РұСҚРә: OpenAI Images
        try:
            resp = oai_img.images.generate(model=IMAGES_MODEL, prompt=prompt, size="1024x1024", n=1)
            return base64.b64decode(resp.data[0].b64_json)
        except Exception as e:
            log.exception("OpenAI images fallback error: %s", e)
            return None
    try:
        # РҹСҖРёРјРөСҖРҪСӢР№ СҚРҪРҙРҝРҫРёРҪСӮ; РөСҒР»Рё Сғ СӮРөРұСҸ РҙСҖСғРіРҫР№ вҖ” Р·Р°РјРөРҪРё path/РҝРҫР»СҸ РҝРҫРҙ СҒРІРҫР№ Р°РәРәР°СғРҪСӮ.
        url = f"{LUMA_IMG_BASE_URL}/v1/images"
        headers = {"Authorization": f"Bearer {LUMA_API_KEY}", "Accept": "application/json"}
        payload = {"model": LUMA_IMG_MODEL, "prompt": prompt, "size": "1024x1024"}
        async with httpx.AsyncClient(timeout=60.0) as client:
            r = await client.post(url, headers=headers, json=payload)
            if r.status_code >= 400:
                return None
            j = r.json() or {}
            b64 = (j.get("data") or [{}])[0].get("b64_json") or j.get("image_base64")
            return base64.b64decode(b64) if b64 else None
    except Exception as e:
        log.exception("Luma image gen error: %s", e)
        return None

async def _start_luma_img(update: Update, context: ContextTypes.DEFAULT_TYPE, prompt: str):
    async def _go():
        img = await _luma_generate_image_bytes(prompt)
        if not img:
            await update.effective_message.reply_text("РқРө СғРҙР°Р»РҫСҒСҢ СҒРҫР·РҙР°СӮСҢ РёР·РҫРұСҖР°Р¶РөРҪРёРө.")
            return False
        await update.effective_message.reply_photo(photo=img, caption=f"рҹ–Ң Р“РҫСӮРҫРІРҫ вң…\nР—Р°РҝСҖРҫСҒ: {prompt}")
        return True
    await _try_pay_then_do(update, context, update.effective_user.id, "img", IMG_COST_USD, _go,
                           remember_kind="luma_img", remember_payload={"prompt": prompt})


async def _midjourney_generate_image_bytes(prompt: str) -> tuple[bytes | None, str]:
    """Create a Midjourney grid through CometAPI and return bytes plus task id."""
    if not MIDJOURNEY_ENABLED:
        raise RuntimeError("Midjourney РҫСӮРәР»СҺСҮС‘РҪ РҪР°СҒСӮСҖРҫР№РәРҫР№ MIDJOURNEY_ENABLED=0")
    if not COMET_API_KEY:
        raise RuntimeError("Р”Р»СҸ Midjourney РҪСғР¶РөРҪ COMET_API_KEY")
    prompt = re.sub(r"\s+", " ", (prompt or "").strip())
    if not prompt:
        raise RuntimeError("РҹСғСҒСӮРҫР№ РҝСҖРҫРјРҝСӮ Midjourney")
    # Add a default version only when the user did not specify one.
    if MIDJOURNEY_DEFAULT_VERSION and not re.search(r"(?:^|\s)--v(?:ersion)?\s+", prompt, re.I):
        prompt = f"{prompt} --v {MIDJOURNEY_DEFAULT_VERSION}"
    mode = MIDJOURNEY_MODE.upper()
    payload = {
        "botType": "MID_JOURNEY",
        "prompt": prompt,
        "accountFilter": {"modes": [mode]},
    }
    headers = {
        "Authorization": f"Bearer {COMET_API_KEY}",
        "Content-Type": "application/json",
        "Accept": "application/json",
    }
    create_paths = []
    for path in (MIDJOURNEY_CREATE_PATH, f"{_MJ_PREFIX}/mj/submit/imagine", "/mj/submit/imagine"):
        if path and path not in create_paths:
            create_paths.append(path)
    last_err = ""
    async with httpx.AsyncClient(timeout=90.0, follow_redirects=True) as client:
        task_id = ""
        for path in create_paths:
            try:
                r = await client.post(f"{COMET_BASE_URL}{path}", headers=headers, json=payload)
                if r.status_code >= 400:
                    last_err = f"POST {path} вҶ’ {r.status_code}: {_api_error_preview(r)}"
                    continue
                js = r.json() or {}
                code = js.get("code")
                task_id = str(js.get("result") or js.get("taskId") or js.get("task_id") or js.get("id") or "").strip()
                if task_id and (code in (None, 1, "1", 200, "200") or str(js.get("description", "")).lower().find("success") >= 0):
                    break
                last_err = f"POST {path}: РҪРөСӮ task id: {json.dumps(js, ensure_ascii=False)[:700]}"
            except Exception as e:
                last_err = f"POST {path}: {e}"
        if not task_id:
            raise RuntimeError(last_err or "Midjourney РҪРө РІРөСҖРҪСғР» task id")

        started = time.time()
        status_paths = []
        for path in (MIDJOURNEY_STATUS_PATH, "/mj/task/{id}/fetch"):
            if path and path not in status_paths:
                status_paths.append(path)
        last = ""
        while time.time() - started < MIDJOURNEY_TIMEOUT_S:
            for path in status_paths:
                try:
                    rr = await client.get(f"{COMET_BASE_URL}{path.format(id=task_id)}", headers=headers)
                    if rr.status_code >= 400:
                        last = f"GET {path} вҶ’ {rr.status_code}: {_api_error_preview(rr)}"
                        continue
                    js = rr.json() or {}
                    status = str(js.get("status") or js.get("state") or "").upper()
                    image_url = (
                        js.get("imageUrl") or js.get("image_url")
                        or _extract_first_url(js.get("output"))
                        or _extract_first_url(js.get("result"))
                        or _extract_first_url(js.get("data"))
                    )
                    if status == "SUCCESS" and image_url:
                        img_resp = await client.get(str(image_url), headers={"Accept": "image/*"})
                        img_resp.raise_for_status()
                        if not img_resp.content:
                            raise RuntimeError("Midjourney РІРөСҖРҪСғР» РҝСғСҒСӮРҫРө РёР·РҫРұСҖР°Р¶РөРҪРёРө")
                        return img_resp.content, task_id
                    if status in ("FAILURE", "FAILED", "ERROR", "CANCELLED", "CANCELED"):
                        desc = js.get("failReason") or js.get("description") or js.get("error") or js
                        raise RuntimeError(f"Midjourney task failed: {str(desc)[:900]}")
                    if status == "MODAL":
                        raise RuntimeError("Midjourney Р·Р°РҝСҖРҫСҒРёР» РҙРҫРҝРҫР»РҪРёСӮРөР»СҢРҪСӢР№ modal-РІРІРҫРҙ; РҙР»СҸ РҫРұСӢСҮРҪРҫР№ РіРөРҪРөСҖР°СҶРёРё РёР·РјРөРҪРёСӮРө РҝСҖРҫРјРҝСӮ")
                    last = json.dumps(js, ensure_ascii=False)[:900]
                except RuntimeError:
                    raise
                except Exception as e:
                    last = str(e)
            await asyncio.sleep(MIDJOURNEY_POLL_DELAY_S)
        raise RuntimeError(f"Midjourney timeout. РҹРҫСҒР»РөРҙРҪРёР№ РҫСӮРІРөСӮ: {last[:900]}")


async def _start_midjourney_image(update: Update, context: ContextTypes.DEFAULT_TYPE, prompt: str):
    prompt = (prompt or "").strip()
    if not prompt:
        context.user_data["awaiting_midjourney_prompt"] = True
        await update.effective_message.reply_text(
            "рҹҺЁ РһРҝРёСҲРёСӮРө РёР·РҫРұСҖР°Р¶РөРҪРёРө РҙР»СҸ Midjourney. РңРҫР¶РҪРҫ РҙРҫРұР°РІРёСӮСҢ РҝР°СҖР°РјРөСӮСҖСӢ --ar, --stylize Рё --v."
        )
        return False

    async def _go():
        await update.effective_message.reply_text(
            f"рҹҺЁ Midjourney {MIDJOURNEY_MODE.upper()}: Р·Р°РҙР°СҮР° РҝСҖРёРҪСҸСӮР°, РҫР¶РёРҙР°СҺ РёР·РҫРұСҖР°Р¶РөРҪРёРөвҖҰ"
        )
        try:
            img, task_id = await _midjourney_generate_image_bytes(prompt)
            if not img:
                return False
            bio = BytesIO(img); bio.name = f"midjourney_{task_id}.jpg"
            try:
                await update.effective_message.reply_photo(photo=InputFile(bio), caption=f"рҹҺЁ Midjourney РіРҫСӮРҫРІ вң…\n{prompt[:700]}")
            except Exception:
                bio.seek(0)
                await update.effective_message.reply_document(document=InputFile(bio), caption="рҹҺЁ Midjourney РіРҫСӮРҫРІ вң…")
            return True
        except Exception as e:
            log.exception("Midjourney failed: %s", e)
            await update.effective_message.reply_text(f"вқҢ Midjourney: {str(e)[:900]}")
            return False

    await _try_pay_then_do(
        update, context, update.effective_user.id,
        "img", MIDJOURNEY_UNIT_COST_USD, _go,
        remember_kind="midjourney_imagine",
        remember_payload={"prompt": prompt[:1000], "mode": MIDJOURNEY_MODE},
    )
    return True


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ Image retouch / own-image watermark cleanup в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
_RETOUCH_TERMS_RE = re.compile(
    r"(РІРҫРҙСҸРҪ(?:РҫР№|РҫРіРҫ|РҫРјСғ|СӢРј|РҫРј)\s+Р·РҪР°Рә|РІРҫРҙРҪ(?:СӢР№|РҫРіРҫ|РҫРјСғ|СӢРј|РҫРј)\s+Р·РҪР°Рә|watermark|"
    r"РІР°СӮРөСҖРјР°СҖРә|СҖРөСӮСғСҲ|СҖРөСӮСғСҲРёСҖ|Р·Р°СҖРөСӮСғСҲ|РІРөСӮРҫСҲСҢ|"
    r"(?:СғРұРөСҖ(?:Рё|РёСӮРө|Сғ)|СғРҙР°Р»(?:Рё|РёСӮРө|РёСӮСҢ)|СҒРҫСӮСҖ(?:Рё|РёСӮРө)|Р·Р°РјР°Р¶СҢ|Р·Р°РјР°Р·Р°СӮСҢ)\s+.{0,80}"
    r"(?:РІРҫРҙСҸРҪ(?:РҫР№|РҫРіРҫ|РҫРјСғ|СӢРј|РҫРј)\s+Р·РҪР°Рә|РІРҫРҙРҪ(?:СӢР№|РҫРіРҫ|РҫРјСғ|СӢРј|РҫРј)\s+Р·РҪР°Рә|watermark|РҪР°РҙРҝРёСҒСҢ|СӮРөРәСҒСӮ|Р»РҫРіРҫСӮРёРҝ|Р»РҫРіРҫ|СҚРјРұР»РөРј|Р»РёСҲРҪ(?:РёР№|СҺСҺ|РөРө|РёРө)\s+РҫРұСҠРөРәСӮ)|"
    r"РҫСҮРёСҒСӮ(?:Рё|РёСӮРө|РёСӮСҢ)\s+.{0,40}(?:С„РҫСӮРҫ|РёР·РҫРұСҖР°Р¶РөРҪ|РәР°СҖСӮРёРҪРә|photo|image)|"
    r"Р»РёСҲРҪ(?:РёР№|СҺСҺ|РөРө|РёРө)\s+(?:РҪР°РҙРҝРёСҒСҢ|РҫРұСҠРөРәСӮ|Р»РҫРіРҫСӮРёРҝ|СӮРөРәСҒСӮ)|"
    r"РҪР°РҙРҝРёСҒСҢ\s+РҪР°\s+С„РҫСӮРҫ|СӮРөРәСҒСӮ\s+РҪР°\s+С„РҫСӮРҫ|Р»РҫРіРҫСӮРёРҝ\s+РҪР°\s+С„РҫСӮРҫ|"
    r"remove\s+(?:watermark|text|logo|object)|clean\s+(?:image|photo))",
    re.IGNORECASE,
)

_OWN_IMAGE_CONFIRM_RE = re.compile(
    r"(РјРҫ[Р№СҸС‘РёРө]|СҒРІРҫ[Р№СҸС‘РөРёС…]|СҒРҫРұСҒСӮРІРөРҪРҪ|РјРҫР№\s+С„Р°Р№Р»|РјРҫР№\s+РјР°РәРөСӮ|РјРҫСҸ\s+С„РҫСӮ|"
    r"РөСҒСӮСҢ\s+РҝСҖР°РІ|РёРјРөСҺ\s+РҝСҖР°РІ|СҖР°Р·СҖРөСҲРөРҪРё|СҸ\s+РІР»Р°РҙРөР»РөСҶ|own\s+image|my\s+image|my\s+photo|i\s+own)",
    re.IGNORECASE,
)


def _is_image_retouch_request(text: str) -> bool:
    """True РҙР»СҸ РұРөР·РҫРҝР°СҒРҪРҫР№ СҖРөСӮСғСҲРё СҒРҫРұСҒСӮРІРөРҪРҪРҫРіРҫ РёР·РҫРұСҖР°Р¶РөРҪРёСҸ: РІРҫРҙСҸРҪРҫР№ Р·РҪР°Рә, Р»РёСҲРҪСҸСҸ РҪР°РҙРҝРёСҒСҢ, Р»РҫРіРҫСӮРёРҝ, РҫРұСҠРөРәСӮ."""
    return bool(_RETOUCH_TERMS_RE.search(text or ""))


def _has_own_image_confirmation(text: str) -> bool:
    return bool(_OWN_IMAGE_CONFIRM_RE.search(text or ""))


def _set_waiting_image_retouch(update: Update, context: ContextTypes.DEFAULT_TYPE, prompt: str = ""):
    """РЎР»РөРҙСғСҺСүРөРө С„РҫСӮРҫ РҙРҫР»Р¶РҪРҫ РҝРҫР№СӮРё РІ СҖРөСӮСғСҲСҢ СҒРҫРұСҒСӮРІРөРҪРҪРҫРіРҫ РёР·РҫРұСҖР°Р¶РөРҪРёСҸ, Р° РҪРө РІ РјРөРҙРёСҶРёРҪСғ/Р°РҪР°Р»РёР·."""
    uid = update.effective_user.id
    _set_mode_clean(uid, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "")
    _clear_transient_flows(context)
    context.user_data["awaiting_photo_for"] = "retouch"
    context.user_data["photo_flow"] = "retouch"
    context.user_data["retouch_prompt"] = (prompt or "").strip()


def _is_waiting_image_retouch(context) -> bool:
    if not context:
        return False
    return (
        context.user_data.get("awaiting_photo_for") == "retouch"
        or context.user_data.get("photo_flow") == "retouch"
    )


def _set_retouch_wait_text(context, prompt: str = ""):
    if not context:
        return
    context.user_data["retouch_wait_text"] = "1"
    if prompt:
        context.user_data["retouch_prompt"] = prompt.strip()


def _is_retouch_wait_text(context) -> bool:
    return bool(context and context.user_data.get("retouch_wait_text"))


def _clear_image_retouch_wait(context):
    if not context:
        return
    for key in ("awaiting_photo_for", "photo_flow", "retouch_prompt", "retouch_wait_text"):
        with contextlib.suppress(Exception):
            context.user_data.pop(key, None)


def _retouch_user_hint_text() -> str:
    return (
        "рҹ§Ҫ Р”Р°, РјРҫРіСғ СҒРҙРөР»Р°СӮСҢ СҖРөСӮСғСҲСҢ СҒРҫРұСҒСӮРІРөРҪРҪРҫРіРҫ РёР·РҫРұСҖР°Р¶РөРҪРёСҸ.\n\n"
        "РҹСҖРёСҲР»РёСӮРө С„РҫСӮРҫ Рё, РөСҒР»Рё РІРҫР·РјРҫР¶РҪРҫ, СғРәР°Р¶РёСӮРө РіРҙРө РҪР°С…РҫРҙРёСӮСҒСҸ СҚР»РөРјРөРҪСӮ: "
        "РҪР°РҝСҖРёРјРөСҖ В«РІРҫРҙСҸРҪРҫР№ Р·РҪР°Рә СҒРҝСҖР°РІР° СҒРҪРёР·СғВ», В«РҪР°РҙРҝРёСҒСҢ РҝРҫ СҶРөРҪСӮСҖСғВ», В«Р»РҫРіРҫСӮРёРҝ СҒРІРөСҖС…СғВ».\n\n"
        "Р’Р°Р¶РҪРҫ: СҸ РҫРұСҖР°РұР°СӮСӢРІР°СҺ СӮРҫР»СҢРәРҫ РІР°СҲРё СҒРҫРұСҒСӮРІРөРҪРҪСӢРө РёР·РҫРұСҖР°Р¶РөРҪРёСҸ/РјР°РәРөСӮСӢ РёР»Рё С„Р°Р№Р»СӢ, "
        "РҪР° РәРҫСӮРҫСҖСӢРө Сғ РІР°СҒ РөСҒСӮСҢ РҝСҖР°РІРҫ СҖРөРҙР°РәСӮРёСҖРҫРІР°РҪРёСҸ."
    )


def _retouch_system_prompt(user_instruction: str) -> str:
    instruction = (user_instruction or "watermark or unwanted text/logo").strip()
    return (
        "Edit this user-owned image. Remove only the unwanted watermark/text/logo/object described by the user: "
        f"{instruction}. Naturally reconstruct the background in that area, preserving texture, lighting, perspective, "
        "shadows, colors and all important details. Do not change faces, identity, composition, style, objects, "
        "branding created by the user, or any other part of the image. Keep the result realistic and high quality."
    )


def _prepare_image_for_edit(img_bytes: bytes) -> tuple[bytes, str, str]:
    """OpenAI image edit Р»СғСҮСҲРө РәРҫСҖРјРёСӮСҢ PNG. Р•СҒР»Рё PIL РҪРөРҙРҫСҒСӮСғРҝРөРҪ вҖ” РҫСӮРҝСҖР°РІР»СҸРөРј РёСҒС…РҫРҙРҪСӢР№ С„Р°Р№Р»."""
    if Image is None:
        mime = sniff_image_mime(img_bytes)
        ext = "jpg" if mime == "image/jpeg" else "png"
        return img_bytes, f"image.{ext}", mime
    try:
        im = Image.open(BytesIO(img_bytes)).convert("RGBA")
        # РһРіСҖР°РҪРёСҮРёРІР°РөРј СҖР°Р·РјРөСҖ, СҮСӮРҫРұСӢ РҪРө Р»РҫРІРёСӮСҢ Р»РёРјРёСӮСӢ multipart, РҪРҫ СҒРҫС…СҖР°РҪСҸРөРј С…РҫСҖРҫСҲСғСҺ РҙРөСӮР°Р»РёР·Р°СҶРёСҺ.
        max_side = 1600
        if max(im.size) > max_side:
            im.thumbnail((max_side, max_side), Image.LANCZOS)
        bio = BytesIO()
        im.save(bio, format="PNG")
        return bio.getvalue(), "image.png", "image/png"
    except Exception:
        mime = sniff_image_mime(img_bytes)
        ext = "jpg" if mime == "image/jpeg" else "png"
        return img_bytes, f"image.{ext}", mime


async def _openai_image_edit_bytes(img_bytes: bytes, user_instruction: str) -> bytes | None:
    """Р РөР°Р»СҢРҪР°СҸ СҖРөСӮСғСҲСҢ СҮРөСҖРөР· OpenAI Images / gpt-image-1: /images/edits."""
    if not OPENAI_IMAGE_KEY:
        return None
    if OPENAI_IMAGE_KEY.startswith("sk-or-"):
        # OpenRouter РҪРө СғРјРөРөСӮ OpenAI Images edits.
        raise RuntimeError("OPENAI_IMAGE_KEY РҝРҫС…РҫР¶ РҪР° РәР»СҺСҮ OpenRouter. Р”Р»СҸ СҖРөСӮСғСҲРё РҪСғР¶РөРҪ РҫС„РёСҶРёР°Р»СҢРҪСӢР№ OpenAI key РёР»Рё СҒРҫРІРјРөСҒСӮРёРјСӢР№ image-edit proxy.")

    edit_bytes, filename, mime = _prepare_image_for_edit(img_bytes)
    prompt = _retouch_system_prompt(user_instruction)
    base = (IMAGES_BASE_URL or "https://api.openai.com/v1").rstrip("/")
    headers = {"Authorization": f"Bearer {OPENAI_IMAGE_KEY}"}

    # РқРөРәРҫСӮРҫСҖСӢРө РҝСҖРҫРәСҒРё РҪРө Р»СҺРұСҸСӮ size, РҝРҫСҚСӮРҫРјСғ РҙРөР»Р°РөРј РҙРІРө РҝРҫРҝСӢСӮРәРё.
    attempts = [
        {"model": IMAGES_MODEL, "prompt": prompt, "n": "1", "size": "1024x1024"},
        {"model": IMAGES_MODEL, "prompt": prompt, "n": "1"},
    ]
    last_err = ""
    async with httpx.AsyncClient(timeout=180.0, follow_redirects=True) as client:
        for data in attempts:
            try:
                files = {"image": (filename, edit_bytes, mime)}
                r = await client.post(f"{base}/images/edits", headers=headers, data=data, files=files)
                if r.status_code >= 400:
                    last_err = f"{r.status_code}: {_api_error_preview(r)}" if "_api_error_preview" in globals() else f"{r.status_code}: {r.text[:500]}"
                    log.warning("Image edit failed: %s", last_err)
                    continue
                js = r.json() or {}
                item = (js.get("data") or [{}])[0]
                b64 = item.get("b64_json")
                if b64:
                    return base64.b64decode(b64)
                url = item.get("url") or _extract_first_url(js) if "_extract_first_url" in globals() else item.get("url")
                if url:
                    rr = await client.get(url, timeout=180.0)
                    rr.raise_for_status()
                    return rr.content
                last_err = f"РҪРөСӮ b64_json/url РІ РҫСӮРІРөСӮРө: {json.dumps(js, ensure_ascii=False)[:500]}"
            except Exception as e:
                last_err = str(e)
                log.warning("Image edit exception: %s", e)
                continue
    raise RuntimeError(last_err or "image edit failed")


async def _edit_own_image_retouch(update: Update, context: ContextTypes.DEFAULT_TYPE, img_bytes: bytes, instruction: str):
    """РһСӮРҝСҖР°РІР»СҸРөСӮ С„РҫСӮРҫ РІ AI-СҖРөСӮСғСҲСҢ Рё РІРҫР·РІСҖР°СүР°РөСӮ СҖРөР·СғР»СҢСӮР°СӮ РҝРҫР»СҢР·РҫРІР°СӮРөР»СҺ."""
    instruction = (instruction or context.user_data.get("retouch_prompt") or "СғРұСҖР°СӮСҢ Р»РёСҲРҪСҺСҺ РҪР°РҙРҝРёСҒСҢ/РІРҫРҙСҸРҪРҫР№ Р·РҪР°Рә Рё РІРҫСҒСҒСӮР°РҪРҫРІРёСӮСҢ С„РҫРҪ").strip()
    if not OPENAI_IMAGE_KEY:
        await update.effective_message.reply_text("вқҢ Р РөСӮСғСҲСҢ РҪРөРҙРҫСҒСӮСғРҝРҪР°: РҪРө Р·Р°РҙР°РҪ OPENAI_IMAGE_KEY/OPENAI_API_KEY.")
        return
    if OPENAI_IMAGE_KEY.startswith("sk-or-"):
        await update.effective_message.reply_text(
            "вқҢ Р”Р»СҸ СҖРөСӮСғСҲРё РҪСғР¶РөРҪ РҫС„РёСҶРёР°Р»СҢРҪСӢР№ OpenAI image key. OpenRouter-РәР»СҺСҮ РҙР»СҸ image edit РҪРө РҝРҫРҙС…РҫРҙРёСӮ."
        )
        return

    await update.effective_message.reply_text(
        "рҹ§Ҫ Р—Р°РҝСғСҒРәР°СҺ СҖРөСӮСғСҲСҢ СҒРҫРұСҒСӮРІРөРҪРҪРҫРіРҫ РёР·РҫРұСҖР°Р¶РөРҪРёСҸ СҮРөСҖРөР· Images. "
        "РЈРұРөСҖСғ СӮРҫР»СҢРәРҫ СғРәР°Р·Р°РҪРҪСӢР№ Р»РёСҲРҪРёР№ СҚР»РөРјРөРҪСӮ Рё РҝРҫСҒСӮР°СҖР°СҺСҒСҢ РөСҒСӮРөСҒСӮРІРөРҪРҪРҫ РІРҫСҒСҒСӮР°РҪРҫРІРёСӮСҢ С„РҫРҪ."
    )
    with contextlib.suppress(Exception):
        await context.bot.send_chat_action(update.effective_chat.id, ChatAction.UPLOAD_PHOTO)
    try:
        out = await _openai_image_edit_bytes(img_bytes, instruction)
        if not out:
            await update.effective_message.reply_text("вқҢ РқРө СғРҙР°Р»РҫСҒСҢ РҝРҫР»СғСҮРёСӮСҢ СҖРөР·СғР»СҢСӮР°СӮ СҖРөСӮСғСҲРё.")
            return
        bio = BytesIO(out)
        bio.name = "retouched.png"
        await update.effective_message.reply_document(
            InputFile(bio),
            caption="вң… Р“РҫСӮРҫРІРҫ: СҖРөСӮСғСҲСҢ РІСӢРҝРҫР»РҪРөРҪР°. РҹСҖРҫРІРөСҖСҢСӮРө С„РҫРҪ Рё РҙРөСӮР°Р»Рё; РҝСҖРё РҪРөРҫРұС…РҫРҙРёРјРҫСҒСӮРё РјРҫР¶РҪРҫ РҫСӮРҝСҖР°РІРёСӮСҢ СғСӮРҫСҮРҪРөРҪРёРө, СҮСӮРҫ РҝРҫРҝСҖР°РІРёСӮСҢ."
        )
        return True
    except Exception as e:
        log.exception("image retouch error: %s", e)
        await update.effective_message.reply_text(
            "вқҢ РқРө СғРҙР°Р»РҫСҒСҢ РІСӢРҝРҫР»РҪРёСӮСҢ СҖРөСӮСғСҲСҢ РёР·РҫРұСҖР°Р¶РөРҪРёСҸ. "
            f"РўРөС…РҪРёСҮРөСҒРәР°СҸ РҝСҖРёСҮРёРҪР°: {str(e)[:700]}"
        )
        return False


async def _start_image_retouch(update: Update, context: ContextTypes.DEFAULT_TYPE, img_bytes: bytes, instruction: str):
    """РҹР»Р°СӮС‘Р¶РҪР°СҸ РҫРұС‘СҖСӮРәР° РҙР»СҸ СҖРөСӮСғСҲРё: РёСҒРҝРҫР»СҢР·СғРөРј СӮРҫСӮ Р¶Рө РұСҺРҙР¶РөСӮ Images."""
    if not OPENAI_IMAGE_KEY or OPENAI_IMAGE_KEY.startswith("sk-or-"):
        return await _edit_own_image_retouch(update, context, img_bytes, instruction)

    async def _go():
        return await _edit_own_image_retouch(update, context, img_bytes, instruction)

    await _try_pay_then_do(
        update, context, update.effective_user.id,
        "img", IMG_COST_USD, _go,
        remember_kind="image_retouch",
        remember_payload={"instruction": instruction},
    )


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ UI / СӮРөРәСҒСӮСӢ в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
START_TEXT = (
    "рҹ‘Ӣ РҹСҖРёРІРөСӮ! РҜ *Neyro-Bot GPT 5 Studio* вҖ” РјСғР»СҢСӮРёРјРҫРҙРөР»СҢРҪР°СҸ AI-СҒСӮСғРҙРёСҸ РІ Telegram РҙР»СҸ СӮРөРәСҒСӮР°, РҙРҫРәСғРјРөРҪСӮРҫРІ, С„РҫСӮРҫ, РІРёРҙРөРҫ, РјСғР·СӢРәРё, СҖРөСҮРё Рё live-РҝРҫРёСҒРәР°.\n"
    "Р Р°РұРҫСӮР°СҺ РІ РҝСҖРҫРҙР°РәСҲРҪ-СҖРөР¶РёРјР°С…: СҒР°Рј РҝРҫРҙРұРёСҖР°СҺ РҝРҫРҙС…РҫРҙСҸСүРёР№ РҙРІРёР¶РҫРә РҝРҫРҙ Р·Р°РҙР°СҮСғ РёР»Рё РҙР°СҺ РІСӢРұСҖР°СӮСҢ РөРіРҫ РІСҖСғСҮРҪСғСҺ СҮРөСҖРөР· В«рҹ§  Р”РІРёР¶РәРёВ».\n"
    "\n"
    "рҹҡҖ *Р§СӮРҫ СғРјРөСҺ:*\n"
    "вҖў рҹҺ“ *РЈСҮС‘РұР°* вҖ” РҫРұСҠСҸСҒРҪРөРҪРёРө СӮРөРј, Р·Р°РҙР°СҮРё, СҚСҒСҒРө/СҖРөС„РөСҖР°СӮСӢ, РәРҫРҪСҒРҝРөРәСӮСӢ РёР· PDF/EPUB/DOCX, РҝР»Р°РҪСӢ Рә СҚРәР·Р°РјРөРҪР°Рј, СҖРөСҮСҢ вҶ” СӮРөРәСҒСӮ.\n"
    "вҖў рҹ’ј *Р Р°РұРҫСӮР°/Р‘РёР·РҪРөСҒ* вҖ” РҝРёСҒСҢРјР°, РҡРҹ, РҙРҫРәСғРјРөРҪСӮСӢ, Р°РҪР°Р»РёСӮРёРәР°, ToDo, РұСҖРёС„СӢ, РҝСҖРөР·РөРҪСӮР°СҶРёРё, PDF-РәР°СӮР°Р»РҫРіРё, Р»РҫРіРҫСӮРёРҝСӢ, СғРҙР°Р»РөРҪРёРө РІРҫРҙСҸРҪСӢС… Р·РҪР°РәРҫРІ, РҙРёРәСӮРҫРІРәР° Рё РҫР·РІСғСҮРәР°.\n"
    "вҖў рҹ”Ҙ *Р Р°Р·РІР»РөСҮРөРҪРёСҸ* вҖ” РҫР¶РёРІР»РөРҪРёРө С„РҫСӮРҫ, РіРҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ, С„РҫСӮРҫвҶ’РІРёРҙРөРҫРәР»РёРҝ СҒ РјСғР·СӢРәРҫР№, РәР»РёРҝ СҒ РІРҫРәР°Р»РҫРј РҙР»СҸ 1 СҮРөР»РҫРІРөРәР°, РІРёРҙРөРҫ РҝРҫ СӮРөРәСҒСӮСғ/РіРҫР»РҫСҒСғ, Reels/Shorts, РјРёРҪРё-С„РёР»СҢРјСӢ, СҒСҶРөРҪР°СҖРёРё, Р·Р°РјРөРҪР° Р»РёСҶР°, СғРҙР°Р»РөРҪРёРө/Р·Р°РјРөРҪР° С„РҫРҪР°, СҖРөСӮСғСҲСҢ Рё РјСғР·СӢРәР° Suno.\n"
    "вҖў рҹ©ә *РңРөРҙРёСҶРёРҪР°* вҖ” СҒРҝСҖР°РІРҫСҮРҪСӢР№ СҖР°Р·РұРҫСҖ РІСӢРҝРёСҒРҫРә, Р°РҪР°Р»РёР·РҫРІ, Р·Р°РәР»СҺСҮРөРҪРёР№, РңР Рў/РҡРў/СҒРҪРёРјРәРҫРІ Рё РҝРҫРҙРіРҫСӮРҫРІРәР° РІРҫРҝСҖРҫСҒРҫРІ РІСҖР°СҮСғ. РӯСӮРҫ РҪРө РҙРёР°РіРҪРҫР· Рё РҪРө Р·Р°РјРөРҪР° РҫСҮРҪРҫР№ РәРҫРҪСҒСғР»СҢСӮР°СҶРёРё.\n"
    "вҖў рҹҢҗ *РҗРәСӮСғР°Р»СҢРҪР°СҸ РёРҪС„РҫСҖРјР°СҶРёСҸ* вҖ” live-РҝРҫРёСҒРә, РҪРҫРІРҫСҒСӮРё, РәСғСҖСҒСӢ, С„Р°РәСӮСӢ Рё РҙР°РҪРҪСӢРө РёР· РёРҪСӮРөСҖРҪРөСӮР°.\n"
    "вҖў рҹ’¬ *РңРҫРё СҮР°СӮСӢ* вҖ” РҙРҫ СҮРөСӮСӢСҖС‘С… РҫСӮРҙРөР»СҢРҪСӢС… РҙРёР°Р»РҫРіРҫРІ СҒ СҒРҫРұСҒСӮРІРөРҪРҪРҫР№ РҝР°РјСҸСӮСҢСҺ Рё РёСҒСӮРҫСҖРёРөР№.\n"
    "вҖў рҹ’і *Р‘Р°Р»Р°РҪСҒ Рё РҝРҫРҙРҝРёСҒРәР°* вҖ” СӮР°СҖРёС„СӢ, РҝРҫРҝРҫР»РҪРөРҪРёРө, РҝР»Р°СӮРҪСӢРө РіРөРҪРөСҖР°СҶРёРё Рё РәРҫРҪСӮСҖРҫР»СҢ СҖР°СҒС…РҫРҙРҫРІ.\n"
    "\n"
    "рҹ§© *РҹРҫРҙ РәР°РҝРҫСӮРҫРј:* РјСғР»СҢСӮРёРјРҫРҙРөР»СҢРҪР°СҸ СҒРІСҸР·РәР° GPT/OpenRouter, OpenAI Images/TTS, Deepgram, Tavily, Runway, Sora 2, Kling, Midjourney, Suno, Photoroom, PiAPI, Segmind Рё CometAPI.\n"
    "\n"
    "Р’СӢРұРөСҖРёСӮРө СҖРөР¶РёРј РәРҪРҫРҝРәРҫР№ РҪРёР¶Рө РёР»Рё РҝСҖРҫСҒСӮРҫ РҪР°РҝРёСҲРёСӮРө Р·Р°РҙР°СҮСғ СӮРөРәСҒСӮРҫРј/РіРҫР»РҫСҒРҫРј."
)

def engines_kb():
    # РһСӮРҙРөР»СҢРҪСӢР№ СҖРөР¶РёРј В«Р”РІРёР¶РәРёВ»: РҝРҫРәР°Р·СӢРІР°РөРј РІСҒРө СҒРөСҖРІРёСҒСӢ/РҪРөР№СҖРҫСҒРөСӮРё,
    # СҖРөР°Р»СҢРҪРҫ РҝРҫРҙРәР»СҺСҮРөРҪРҪСӢРө РІ СҚСӮРҫР№ СҒРұРҫСҖРәРө. РҡРҪРҫРҝРәРё вҖ” РҝРҫ РҫРҙРҪРҫР№ РІ СҒСӮСҖРҫРәРө,
    # СҮСӮРҫРұСӢ Telegram РҪР° Android РҪРө РҫРұСҖРөР·Р°Р» РҙР»РёРҪРҪСӢР№ СӮРөРәСҒСӮ.
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("рҹ§  GPT / OpenRouter", callback_data="engine:gpt")],
        [InlineKeyboardButton("рҹ–ј OpenAI Images", callback_data="engine:images")],
        [InlineKeyboardButton("рҹ”Ҡ OpenAI TTS", callback_data="engine:openai_tts")],
        [InlineKeyboardButton("рҹ—Ј Deepgram STT", callback_data="engine:deepgram")],
        [InlineKeyboardButton("рҹҢҗ Tavily РҝРҫРёСҒРә", callback_data="engine:tavily")],
        [InlineKeyboardButton("рҹҺҘ Runway РІРёРҙРөРҫ", callback_data="engine:runway")],
        [InlineKeyboardButton("рҹҺһ Sora 2 В· РұРөР· Р»СҺРҙРөР№", callback_data="engine:sora")],
        [InlineKeyboardButton("рҹҺ¬ Kling РІРёРҙРөРҫ", callback_data="engine:kling")],
        [InlineKeyboardButton("рҹҺЁ Midjourney РёР·РҫРұСҖР°Р¶РөРҪРёСҸ", callback_data="engine:midjourney")],
        [InlineKeyboardButton("рҹҺө Suno РјСғР·СӢРәР°", callback_data="engine:suno")],
        [InlineKeyboardButton("рҹ§ј Photoroom С„РҫРҪ", callback_data="engine:photoroom")],
        [InlineKeyboardButton("рҹҺӯ PiAPI Р»РёСҶРҫ", callback_data="engine:piapi")],
        [InlineKeyboardButton("рҹ’Һ Segmind Р»РёСҶРҫ", callback_data="engine:segmind")],
        [InlineKeyboardButton("рҹҢү CometAPI СҲР»СҺР·", callback_data="engine:comet")],
        [InlineKeyboardButton("рҹ’° РҰРөРҪСӢ РіРөРҪРөСҖР°СҶРёР№", callback_data="pricing:list")],
    ])

ENGINE_INFO_TEXT = {
    "gpt": (
        "рҹ§  GPT / OpenRouter: РҫСҒРҪРҫРІРҪРҫР№ СӮРөРәСҒСӮРҫРІСӢР№ РјРҫР·Рі РұРҫСӮР°.\n"
        "РҳСҒРҝРҫР»СҢР·СғРөСӮСҒСҸ РҙР»СҸ СҮР°СӮР°, РҙРҫРәСғРјРөРҪСӮРҫРІ, Р°РҪР°Р»РёР·Р°, СғСҮРөРұРҪСӢС… Рё СҖР°РұРҫСҮРёС… Р·Р°РҙР°СҮ.\n"
        "РҡР»СҺСҮРё: OPENAI_API_KEY / OPENAI_BASE_URL / OPENAI_MODEL."
    ),
    "images": (
        "рҹ–ј OpenAI Images: РіРөРҪРөСҖР°СҶРёСҸ Рё СҮР°СҒСӮСҢ СҖРөРҙР°РәСӮРёСҖРҫРІР°РҪРёСҸ РёР·РҫРұСҖР°Р¶РөРҪРёР№.\n"
        "Р—Р°РҝСғСҒРә: /img <РҫРҝРёСҒР°РҪРёРө> РёР»Рё РәРҪРҫРҝРәРё С„РҫСӮРҫ-РјР°СҒСӮРөСҖСҒРәРҫР№.\n"
        "РҡР»СҺСҮРё: OPENAI_IMAGE_KEY РёР»Рё OPENAI_API_KEY."
    ),
    "openai_tts": (
        "рҹ”Ҡ OpenAI TTS СӮРөРәСҒСӮР° РіРҫР»РҫСҒРҫРј.\n"
        "РҳСҒРҝРҫР»СҢР·СғРөСӮСҒСҸ РҙР»СҸ СҖРөР¶РёРјР° СҖРөСҮСҢ/РҫР·РІСғСҮРәР°.\n"
        "РҡР»СҺСҮРё: OPENAI_TTS_KEY / OPENAI_TTS_MODEL / OPENAI_TTS_VOICE."
    ),
    "deepgram": (
        "рҹ—Ј Deepgram STT: СҖР°СҒРҝРҫР·РҪР°РІР°РҪРёРө РіРҫР»РҫСҒРҫРІСӢС… СҒРҫРҫРұСүРөРҪРёР№ Рё Р°СғРҙРёРҫ РІ СӮРөРәСҒСӮ.\n"
        "РҡР»СҺСҮ: DEEPGRAM_API_KEY. Р РөР·РөСҖРІ: OPENAI_STT_KEY / Whisper."
    ),
    "stt_tts": (
        "рҹ—Ј Р РөСҮСҢ вҶ” СӮРөРәСҒСӮ: СҒРІСҸР·РәР° STT/TTS.\n"
        "STT: Deepgram РёР»Рё OpenAI Whisper. TTS: OpenAI TTS."
    ),
    "tavily": (
        "рҹҢҗ Tavily: live-РҝРҫРёСҒРә РІ РёРҪСӮРөСҖРҪРөСӮРө РҙР»СҸ СҒРІРөР¶РёС… РҙР°РҪРҪСӢС…, РҪРҫРІРҫСҒСӮРөР№, РәСғСҖСҒРҫРІ, Р·Р°РәРҫРҪРҫРІ Рё РёСҒСӮРҫСҮРҪРёРәРҫРІ.\n"
        "РҡР»СҺСҮ: TAVILY_API_KEY."
    ),
    "runway": (
        "рҹҺҘ Runway: РҫС„РёСҶРёР°Р»СҢРҪР°СҸ РіРөРҪРөСҖР°СҶРёСҸ РІРёРҙРөРҫ РҝРҫ СӮРөРәСҒСӮСғ Рё РҫР¶РёРІР»РөРҪРёРө С„РҫСӮРҫ.\n"
        "РһСҒРҪРҫРІРҪРҫР№ РјР°СҖСҲСҖСғСӮ вҖ” Runway Developer API; CometAPI Рё Kling РёСҒРҝРҫР»СҢР·СғСҺСӮСҒСҸ СӮРҫР»СҢРәРҫ РәР°Рә СҖРөР·РөСҖРІ.\n"
        "РҡР»СҺСҮ: RUNWAYML_API_SECRET. Powered by Runway вҖ” https://runwayml.com"
    ),
    "sora": (
        "рҹҺһ Sora 2: РіРөРҪРөСҖР°СҶРёСҸ РІРёРҙРөРҫ РҝРҫ СӮРөРәСҒСӮСғ Рё С„РҫСӮРҫ СҮРөСҖРөР· CometAPI.\n"
        "Р’Р°Р¶РҪРҫ: СҖРөР¶РёРј Sora 2 РұРөР· Р»СҺРҙРөР№ вҖ” РҙР»СҸ С„РҫСӮРҫ/РІРёРҙРөРҫ СҒ Р»СҺРҙСҢРјРё Р»СғСҮСҲРө Runway РёР»Рё Kling.\n"
        "РҡР»СҺСҮ: COMET_API_KEY."
    ),
    "kling": (
        "рҹҺ¬ Kling: РіРөРҪРөСҖР°СҶРёСҸ РІРёРҙРөРҫ РҝРҫ СӮРөРәСҒСӮСғ Рё С„РҫСӮРҫ СҮРөСҖРөР· CometAPI.\n"
        "РҳСҒРҝРҫР»СҢР·СғРөСӮСҒСҸ РҙР»СҸ РҙРёРҪР°РјРёСҮРҪСӢС… РәР»РёРҝРҫРІ, Reels/Shorts Рё fallback РҙР»СҸ РҫР¶РёРІР»РөРҪРёСҸ С„РҫСӮРҫ.\n"
        "РҡР»СҺСҮ: COMET_API_KEY."
    ),
    "midjourney": (
        "рҹҺЁ Midjourney: РіРөРҪРөСҖР°СҶРёСҸ РёР·РҫРұСҖР°Р¶РөРҪРёР№ СҮРөСҖРөР· CometAPI/СҒРҫРІРјРөСҒСӮРёРјСӢР№ РәР°РҪР°Р».\n"
        "РҳСҒРҝРҫР»СҢР·СғРөСӮСҒСҸ РҙР»СҸ РІРёР·СғР°Р»РҫРІ, РҫРұР»РҫР¶РөРә, Р»РҫРіРҫСӮРёРҝРҫРІ Рё РёР»Р»СҺСҒСӮСҖР°СҶРёР№.\n"
        "РҡР»СҺСҮ: COMET_API_KEY."
    ),
    "photoroom": (
        "рҹ§ј Photoroom: РҝСҖРҫРҙР°РәСҲРҪ-РёРҪСҒСӮСҖСғРјРөРҪСӮ РҙР»СҸ СғРҙР°Р»РөРҪРёСҸ Рё AI-Р·Р°РјРөРҪСӢ С„РҫРҪР°.\n"
        "РҳСҒРҝРҫР»СҢР·СғРөСӮСҒСҸ РІ РәРҪРҫРҝРәР°С… В«РЈРҙР°Р»РёСӮСҢ С„РҫРҪВ» Рё В«Р—Р°РјРөРҪРёСӮСҢ С„РҫРҪВ».\n"
        "РҡР»СҺСҮ: PHOTOROOM_API_KEY."
    ),
    "piapi": (
        "рҹҺӯ PiAPI: РұСӢСҒСӮСҖСӢР№ СҖРөР¶РёРј Р·Р°РјРөРҪСӢ Р»РёСҶР°.\n"
        "РҳСҒРҝРҫР»СҢР·СғРөСӮСҒСҸ РәР°Рә РҫСҒРҪРҫРІРҪРҫР№ РұСӢСҒСӮСҖСӢР№ FaceSwap-РҝСҖРҫРІР°Р№РҙРөСҖ.\n"
        "РҡР»СҺСҮ: PIAPI_API_KEY."
    ),
    "segmind": (
        "рҹ’Һ Segmind: РҝСҖРөРјРёСғРј/СҖРөР·РөСҖРІРҪСӢР№ FaceSwap.\n"
        "РҳСҒРҝРҫР»СҢР·СғРөСӮСҒСҸ РҙР»СҸ РұРҫР»РөРө РәР°СҮРөСҒСӮРІРөРҪРҪРҫР№ Р·Р°РјРөРҪСӢ Р»РёСҶР° Рё РәР°Рә fallback, РәРҫРіРҙР° РҙРҫСҒСӮР°СӮРҫСҮРҪРҫ РәСҖРөРҙРёСӮРҫРІ.\n"
        "РҡР»СҺСҮ: SEGMIND_API_KEY."
    ),
    "comet": (
        "рҹҢү CometAPI: РҫРұСүРёР№ СҲР»СҺР· РҙР»СҸ Sora 2, Kling, Runway/Midjourney-РәР°РҪР°Р»РҫРІ Рё РҙСҖСғРіРёС… РіРөРҪРөСҖР°СӮРёРІРҪСӢС… РјРҫРҙРөР»РөР№.\n"
        "РҡР»СҺСҮ: COMET_API_KEY."
    ),
}

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ MODES (РЈСҮС‘РұР° / Р Р°РұРҫСӮР° / Р Р°Р·РІР»РөСҮРөРҪРёСҸ) в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ

from telegram import InlineKeyboardMarkup, InlineKeyboardButton
from telegram.ext import CallbackQueryHandler, MessageHandler, filters

# РўРөРәСҒСӮ РәРҫСҖРҪРөРІРҫРіРҫ РјРөРҪСҺ СҖРөР¶РёРјРҫРІ
def _modes_root_text() -> str:
    return (
        "Р’СӢРұРөСҖРёСӮРө СҖРөР¶РёРј СҖР°РұРҫСӮСӢ. Р‘РҫСӮ СҖР°РұРҫСӮР°РөСӮ РәР°Рә РөРҙРёРҪСӢР№ РҝСҖРҫРҙР°РәСҲРҪ-СҶРөРҪСӮСҖ РҳРҳ:\n"
        "вҖў рҹҺ“ РЈСҮС‘РұР° вҖ” РҫРұСҠСҸСҒРҪРөРҪРёСҸ, Р·Р°РҙР°СҮРё, РәРҫРҪСҒРҝРөРәСӮСӢ, СҚРәР·Р°РјРөРҪСӢ, С„Р°Р№Р»СӢ Рё РіРҫР»РҫСҒ.\n"
        "вҖў рҹ’ј Р Р°РұРҫСӮР°/Р‘РёР·РҪРөСҒ вҖ” РҝРёСҒСҢРјР°, РҙРҫРәСғРјРөРҪСӮСӢ, Р°РҪР°Р»РёСӮРёРәР°, РҝСҖРөР·РөРҪСӮР°СҶРёРё, PDF-РәР°СӮР°Р»РҫРіРё, Р»РҫРіРҫСӮРёРҝСӢ, СҖРөСӮСғСҲСҢ Рё РҝР»Р°РҪСӢ.\n"
        "вҖў рҹ”Ҙ Р Р°Р·РІР»РөСҮРөРҪРёСҸ вҖ” С„РҫСӮРҫ, РІРёРҙРөРҫ, Reels/Shorts, РјРёРҪРё-С„РёР»СҢРјСӢ, С„РҫРҪ, Р»РёСҶРҫ, AI-СҒРөР»С„Рё, РІРҫРәР°Р»СҢРҪСӢРө РәР»РёРҝСӢ Рё РјСғР·СӢРәР°.\n"
        "вҖў рҹ©ә РңРөРҙРёСҶРёРҪР° вҖ” СҒРҝСҖР°РІРҫСҮРҪСӢР№ СҖР°Р·РұРҫСҖ Р°РҪР°Р»РёР·РҫРІ, РІСӢРҝРёСҒРҫРә, Р·Р°РәР»СҺСҮРөРҪРёР№ Рё СҒРҪРёРјРәРҫРІ.\n"
        "вҖў рҹ§  Р”РІРёР¶РәРё вҖ” СҖСғСҮРҪРҫР№ РІСӢРұРҫСҖ GPT/OpenRouter, OpenAI Images/TTS, Deepgram, Tavily, Runway, Sora 2, Kling, Midjourney, Suno, Photoroom, PiAPI, Segmind Рё CometAPI.\n"
        "вҖў рҹ’і Р‘Р°Р»Р°РҪСҒ/РҝРҫРҙРҝРёСҒРәР° вҖ” СӮР°СҖРёС„СӢ, РҝРҫРҝРҫР»РҪРөРҪРёРө Рё РәРҫРҪСӮСҖРҫР»СҢ РҝР»Р°СӮРҪСӢС… РіРөРҪРөСҖР°СҶРёР№.\n\n"
        "РңРҫР¶РҪРҫ РҝСҖРҫСҒСӮРҫ РҪР°РҝРёСҒР°СӮСҢ Р·Р°РҝСҖРҫСҒ вҖ” РұРҫСӮ СҒР°Рј РҫРҝСҖРөРҙРөР»РёСӮ Р·Р°РҙР°СҮСғ."
    )

def modes_root_kb() -> InlineKeyboardMarkup:
    # РқРө СҒСӮР°РІРёРј 3 РҙР»РёРҪРҪСӢРө РәРҪРҫРҝРәРё РІ РҫРҙРҪСғ СҒСӮСҖРҫРәСғ вҖ” РҪР° Android Telegram СҖРөР¶РөСӮ СӮРөРәСҒСӮ.
    return InlineKeyboardMarkup([
        [
            InlineKeyboardButton("рҹҺ“ РЈСҮС‘РұР°", callback_data="mode:study"),
            InlineKeyboardButton("рҹ’ј Р Р°РұРҫСӮР°/Р‘РёР·РҪРөСҒ", callback_data="mode:work"),
        ],
        [
            InlineKeyboardButton("рҹ”Ҙ Р Р°Р·РІР»РөСҮРөРҪРёСҸ", callback_data="mode:fun"),
            InlineKeyboardButton("рҹ©ә РңРөРҙРёСҶРёРҪР°", callback_data="mode:medicine"),
        ],
    ])

# в”Җв”Җ РһРҝРёСҒР°РҪРёРө Рё РҝРҫРҙРјРөРҪСҺ РҝРҫ СҖРөР¶РёРјР°Рј
def _mode_desc(key: str) -> str:
    if key == "study":
        return (
            "рҹҺ“ *РЈСҮС‘РұР°*\n"
            "Р“РёРұСҖРёРҙ: GPT-5 РҙР»СҸ РҫРұСҠСҸСҒРҪРөРҪРёР№/РәРҫРҪСҒРҝРөРәСӮРҫРІ, Vision РҙР»СҸ С„РҫСӮРҫ-Р·Р°РҙР°СҮ, "
            "STT/TTS РҙР»СҸ РіРҫР»РҫСҒРҫРІСӢС… Рё Р°СғРҙРёРҫ. Р“РөРҪРөСҖР°СӮРёРІРҪСӢРө С„РҫСӮРҫ/РІРёРҙРөРҫ-РҙРІРёР¶РәРё РІСӢРҪРөСҒРөРҪСӢ РІ РҫСӮРҙРөР»СҢРҪСӢР№ СҖРөР¶РёРј В«Р”РІРёР¶РәРёВ».\n\n"
            "Р‘СӢСҒСӮСҖСӢРө РҙРөР№СҒСӮРІРёСҸ РҪРёР¶Рө. РңРҫР¶РҪРҫ РҪР°РҝРёСҒР°СӮСҢ СҒРІРҫРұРҫРҙРҪСӢР№ Р·Р°РҝСҖРҫСҒ (РҪР°РҝСҖРёРјРөСҖ: "
            "В«СҒРҙРөР»Р°Р№ РәРҫРҪСҒРҝРөРәСӮ РёР· PDFВ», В«РҫРұСҠСҸСҒРҪРё РёРҪСӮРөРіСҖР°Р»СӢ СҒ РҝСҖРёРјРөСҖР°РјРёВ»)."
        )
    if key == "work":
        return (
            "рҹ’ј *Р Р°РұРҫСӮР°/Р‘РёР·РҪРөСҒ*\n"
            "Р“РёРұСҖРёРҙ: GPT-5 РҙР»СҸ РҙРҫРәСғРјРөРҪСӮРҫРІ, РҡРҹ, Р°РҪР°Р»РёСӮРёРәРё Рё СҒСӮСҖСғРәСӮСғСҖСӢ; Vision РҙР»СҸ СӮР°РұР»РёСҶ/СҒРәСҖРёРҪРҫРІ; "
            "OpenAI Images/Midjourney РҙР»СҸ Р»РҫРіРҫСӮРёРҝРҫРІ Рё РІРёР·СғР°Р»РҫРІ; image-edit РҙР»СҸ СғРҙР°Р»РөРҪРёСҸ РІРҫРҙСҸРҪСӢС… Р·РҪР°РәРҫРІ; "
            "РҝРҫР»РҪРҫСҶРөРҪРҪСӢР№ РјР°СҒСӮРөСҖ РҝСҖРөР·РөРҪСӮР°СҶРёР№/РәР°СӮР°Р»РҫРіРҫРІ СҒ Р»РҫРіРҫСӮРёРҝРҫРј, РјР°СҒСҒРҫРІРҫР№ Р·Р°РіСҖСғР·РәРҫР№ С„РҫСӮРҫ, AI-РІРёР·СғР°Р»Р°РјРё Рё СҚРәСҒРҝРҫСҖСӮРҫРј PDF+PPTX; STT/TTS РҙР»СҸ РҙРёРәСӮРҫРІРәРё Рё РҫР·РІСғСҮРәРё.\n\n"
            "Р‘СӢСҒСӮСҖСӢРө РҙРөР№СҒСӮРІРёСҸ РҪРёР¶Рө. РңРҫР¶РҪРҫ РҪР°РҝРёСҒР°СӮСҢ СҒРІРҫРұРҫРҙРҪСӢР№ Р·Р°РҝСҖРҫСҒ: В«СҒРҙРөР»Р°Р№ РҡРҹВ», В«СҒРҫР·РҙР°Р№ РҝСҖРөР·РөРҪСӮР°СҶРёСҺ РҪР° 10 СҒР»Р°Р№РҙРҫРІВ», "
            "В«СҒРҫРұРөСҖРё PDF-РәР°СӮР°Р»РҫРі РҫРұСҠРөРәСӮРҫРІВ», В«СҒРҫР·РҙР°Р№ Р»РҫРіРҫСӮРёРҝ РҙР»СҸ РұСҖРөРҪРҙР°В», В«СғРұРөСҖРё РІРҫРҙСҸРҪРҫР№ Р·РҪР°Рә СҒ С„РҫСӮРҫВ»."
        )
    if key == "fun":
        return (
            "рҹ”Ҙ *Р Р°Р·РІР»РөСҮРөРҪРёСҸ*\n"
            "Р“РёРұСҖРёРҙ: GPT-5 (РёРҙРөРё, СҒСҶРөРҪР°СҖРёРё, СҖР°СҒРәР°РҙСҖРҫРІРәР°), Vision (С„РҫСӮРҫ/СҖРөС„РөСҖРөРҪСҒСӢ), "
            "Sora 2 РұРөР· Р»СҺРҙРөР№, Kling Рё Runway (РІРёРҙРөРҫ РҝРҫ СӮРөРәСҒСӮСғ), Runway/Kling (РҫР¶РёРІР»РөРҪРёРө С„РҫСӮРҫ СҒ Р»СҺРҙСҢРјРё), Sora 2 вҖ” СӮРҫР»СҢРәРҫ СҒСҶРөРҪСӢ РұРөР· Р»СҺРҙРөР№, STT/TTS (РіРҫР»РҫСҒ Рё РҫР·РІСғСҮРәР°).\n\n"
            "Р“Р»Р°РІРҪСӢРө РұСӢСҒСӮСҖСӢРө РҙРөР№СҒСӮРІРёСҸ: РҫР¶РёРІРёСӮСҢ С„РҫСӮРҫ, РіРҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ, С„РҫСӮРҫвҶ’РІРёРҙРөРҫРәР»РёРҝ СҒ РјСғР·СӢРәРҫР№, РәР»РёРҝ СҒ РІРҫРәР°Р»РҫРј РҙР»СҸ 1 СҮРөР»РҫРІРөРәР°, РІРёРҙРөРҫ РҝРҫ СӮРөРәСҒСӮСғ/РіРҫР»РҫСҒСғ СҮРөСҖРөР· Kling/Runway, СҒРҙРөР»Р°СӮСҢ Reels/Shorts, СҒРҫР·РҙР°СӮСҢ РјРёРҪРё-С„РёР»СҢРј, Р·Р°РјРөРҪРёСӮСҢ Р»РёСҶРҫ, СғРҙР°Р»РёСӮСҢ РёР»Рё Р·Р°РјРөРҪРёСӮСҢ С„РҫРҪ РҪР° С„РҫСӮРҫ.\n"
            "РңРҫР¶РҪРҫ РҪР°РҝРёСҒР°СӮСҢ СҒРІРҫРұРҫРҙРҪСӢР№ Р·Р°РҝСҖРҫСҒ, РҪР°РҝСҖРёРјРөСҖ: В«СҒРҙРөР»Р°Р№ СҖРёР»СҒ 20 СҒРөРәСғРҪРҙ РҝСҖРҫ РІРёР»Р»Сғ РҪР° РЎР°РјСғРё, СҒСӮРёР»СҢ luxury, 9:16В»."
        )
    if key == "medicine":
        return _medical_menu_text()
    return "Р РөР¶РёРј РҪРө РҪР°Р№РҙРөРҪ."

def _mode_kb(key: str) -> InlineKeyboardMarkup:
    if key == "study":
        # Р’СҒРө РҙР»РёРҪРҪСӢРө РҝРҫРҙРҝРёСҒРё РёРҙСғСӮ РҝРҫ РҫРҙРҪРҫР№ РәРҪРҫРҝРәРө РІ СҒСӮСҖРҫРәРө вҖ” Telegram РҪР° Android РҪРө СҖРөР¶РөСӮ СӮРөРәСҒСӮ.
        return InlineKeyboardMarkup([
            [InlineKeyboardButton("рҹ“ҡ РҡРҫРҪСҒРҝРөРәСӮ PDF/EPUB/DOCX", callback_data="act:study:pdf_summary")],
            [InlineKeyboardButton("рҹ”Қ РһРұСҠСҸСҒРҪРёСӮСҢ СӮРөРјСғ", callback_data="act:study:explain")],
            [InlineKeyboardButton("рҹ§® Р РөСҲРёСӮСҢ Р·Р°РҙР°СҮРё", callback_data="act:study:tasks")],
            [InlineKeyboardButton("вңҚпёҸ РӯСҒСҒРө / СҖРөС„РөСҖР°СӮ", callback_data="act:study:essay")],
            [InlineKeyboardButton("рҹ“қ РҹР»Р°РҪ Рә СҚРәР·Р°РјРөРҪСғ", callback_data="act:study:exam_plan")],
            [InlineKeyboardButton("рҹ—Ј Р РөСҮСҢ вҶ” СӮРөРәСҒСӮ", callback_data="act:open:voice")],
            [InlineKeyboardButton("рҹ“қ РЎРІРҫРұРҫРҙРҪСӢР№ Р·Р°РҝСҖРҫСҒ", callback_data="act:free")],
            [InlineKeyboardButton("в¬…пёҸ РқР°Р·Р°Рҙ", callback_data="mode:root")],
        ])

    if key == "work":
        return InlineKeyboardMarkup([
            [InlineKeyboardButton("рҹ“„ РҹРёСҒСҢРјРҫ / РҙРҫРәСғРјРөРҪСӮ", callback_data="act:work:doc")],
            [InlineKeyboardButton("рҹ“Ҡ РҗРҪР°Р»РёСӮРёРәР° / СҒРІРҫРҙРәР°", callback_data="act:work:report")],
            [InlineKeyboardButton("рҹ—Ӯ РҹР»Р°РҪ / ToDo", callback_data="act:work:plan")],
            [InlineKeyboardButton("рҹ’Ў РҳРҙРөРё / РұСҖРёС„", callback_data="act:work:idea")],
            [InlineKeyboardButton("рҹ“Ҡ РҹСҖРөР·РөРҪСӮР°СҶРёСҸ PDF + PPTX", callback_data="act:work:presentation")],
            [InlineKeyboardButton("рҹ“• РҡР°СӮР°Р»РҫРі PDF + PPTX", callback_data="act:work:catalog_pdf")],
            [InlineKeyboardButton("рҹҺЁ РЎРҫР·РҙР°СӮСҢ Р»РҫРіРҫСӮРёРҝ", callback_data="act:work:logo")],
            [InlineKeyboardButton("рҹ§Ҫ РЈРҙР°Р»РёСӮСҢ РІРҫРҙСҸРҪРҫР№ Р·РҪР°Рә", callback_data="act:work:watermark")],
            [InlineKeyboardButton("рҹ—Ј Р РөСҮСҢ вҶ” СӮРөРәСҒСӮ", callback_data="act:open:voice")],
            [InlineKeyboardButton("рҹ“қ РЎРІРҫРұРҫРҙРҪСӢР№ Р·Р°РҝСҖРҫСҒ", callback_data="act:free")],
            [InlineKeyboardButton("в¬…пёҸ РқР°Р·Р°Рҙ", callback_data="mode:root")],
        ])

    if key == "fun":
        return InlineKeyboardMarkup([
            [InlineKeyboardButton("рҹӘ„ РһР¶РёРІРёСӮСҢ С„РҫСӮРҫ", callback_data="act:fun:revive")],
            [InlineKeyboardButton("рҹ—Ј Р“РҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ", callback_data="act:fun:avatar")],
            [InlineKeyboardButton("рҹҺө РӨРҫСӮРҫ вҶ’ РІРёРҙРөРҫРәР»РёРҝ", callback_data="act:fun:photoclip")],
            [InlineKeyboardButton("рҹҺӨ РҡР»РёРҝ СҒ РІРҫРәР°Р»РҫРј (1 СҮРөР»РҫРІРөРә)", callback_data="act:fun:vocalclip")],
            [InlineKeyboardButton("рҹҺ¬ Р’РёРҙРөРҫ РҝРҫ СӮРөРәСҒСӮСғ/РіРҫР»РҫСҒСғ", callback_data="act:fun:textvideo")],
            [InlineKeyboardButton("рҹӨі AI-СҒРөР»С„Рё СҒРҫ Р·РІРөР·РҙРҫР№", callback_data="act:fun:aiselfie")],
            [InlineKeyboardButton("рҹҺӯ Р—Р°РјРөРҪР° Р»РёСҶР° РҪР° С„РҫСӮРҫ", callback_data="act:fun:faceswap")],
            [InlineKeyboardButton("рҹ§ј РЈРҙР°Р»РёСӮСҢ С„РҫРҪ РҪР° С„РҫСӮРҫ", callback_data="act:fun:removebg")],
            [InlineKeyboardButton("рҹ–ј Р—Р°РјРөРҪРёСӮСҢ С„РҫРҪ РҪР° С„РҫСӮРҫ", callback_data="act:fun:replacebg")],
            [InlineKeyboardButton("рҹ“ұ Reels / Shorts", callback_data="act:fun:reels")],
            [InlineKeyboardButton("рҹҺһ РЎРҫР·РҙР°СӮСҢ РјРёРҪРё-С„РёР»СҢРј", callback_data="act:fun:film")],
            [InlineKeyboardButton("рҹҺ¬ РЎСҶРөРҪР°СҖРёР№ СҲРҫСҖСӮР°", callback_data="act:fun:shorts")],
            [InlineKeyboardButton("рҹҺө РңСғР·СӢРәР° / РҝРөСҒРҪСҸ", callback_data="act:fun:music")],
            [InlineKeyboardButton("рҹҺ® РҳРіСҖСӢ / РәРІРёР·", callback_data="act:fun:games")],
            [InlineKeyboardButton("рҹҺӯ РҳРҙРөРё РҙР»СҸ РҙРҫСҒСғРіР°", callback_data="act:fun:ideas")],
            [InlineKeyboardButton("рҹ“қ РЎРІРҫРұРҫРҙРҪСӢР№ Р·Р°РҝСҖРҫСҒ", callback_data="act:free")],
            [InlineKeyboardButton("в¬…пёҸ РқР°Р·Р°Рҙ", callback_data="mode:root")],
        ])

    if key == "medicine":
        return medicine_kb()

    return modes_root_kb()


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РҹРҫРҙРјРөРҪСҺ: РіРҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ / С„РҫСӮРҫвҶ’РІРёРҙРөРҫРәР»РёРҝ в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def _avatar_action_kb(prefix: str = "act") -> InlineKeyboardMarkup:
    """РҹРҫРҙРјРөРҪСҺ РіРҫРІРҫСҖСҸСүРөРіРҫ Р°РІР°СӮР°СҖР°. prefix='act' РҙР»СҸ _mode_kb, prefix='fun' РҙР»СҸ СҒСӮР°СҖРҫРіРҫ quick-menu."""
    base = "act:fun" if prefix == "act" else "fun"
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("рҹ“ё Р—Р°РіСҖСғР·РёСӮСҢ РҝРҫСҖСӮСҖРөСӮ", callback_data=f"{base}:avatar_upload")],
        [InlineKeyboardButton("вң… РҳСҒРҝРҫР»СҢР·РҫРІР°СӮСҢ РҝРҫСҒР»РөРҙРҪРөРө С„РҫСӮРҫ", callback_data=f"{base}:avatar_last")],
        [InlineKeyboardButton("рҹ‘© Р“РҫР»РҫСҒ Nova", callback_data=f"{base}:av_voice_nova"), InlineKeyboardButton("рҹ‘Ё Р“РҫР»РҫСҒ Onyx", callback_data=f"{base}:av_voice_onyx")],
        [InlineKeyboardButton("вҡӘ Р“РҫР»РҫСҒ Alloy", callback_data=f"{base}:av_voice_alloy"), InlineKeyboardButton("вңЁ Р“РҫР»РҫСҒ Shimmer", callback_data=f"{base}:av_voice_shimmer")],
        [InlineKeyboardButton("рҹ“– Р“РҫР»РҫСҒ Fable", callback_data=f"{base}:av_voice_fable")],
        [InlineKeyboardButton("рҹ“қ РўРөРәСҒСӮ вҶ’ РіРҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ", callback_data=f"{base}:avatar_text")],
        [InlineKeyboardButton("рҹҺҷ Voice/audio вҶ’ РіРҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ", callback_data=f"{base}:avatar_voice")],
        [InlineKeyboardButton("в¬…пёҸ РқР°Р·Р°Рҙ РІ Р Р°Р·РІР»РөСҮРөРҪРёСҸ", callback_data="mode:fun" if prefix == "act" else "fun:back")],
    ])


def _avatar_voice_choice_kb(prefix: str = "act") -> InlineKeyboardMarkup:
    """РЁР°Рі 2 РҙР»СҸ СӮРөРәСҒСӮРҫРІРҫРіРҫ Р°РІР°СӮР°СҖР°: РІСӢРұРҫСҖ TTS-РіРҫР»РҫСҒР° РҙРҫ РІРІРҫРҙР° СӮРөРәСҒСӮР°."""
    base = "act:fun" if prefix == "act" else "fun"
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("рҹ‘© Nova вҖ” РјСҸРіРәРёР№ Р¶РөРҪСҒРәРёР№", callback_data=f"{base}:av_voice_nova")],
        [InlineKeyboardButton("рҹ‘Ё Onyx вҖ” РҪРёР·РәРёР№ РјСғР¶СҒРәРҫР№", callback_data=f"{base}:av_voice_onyx")],
        [InlineKeyboardButton("вңЁ Shimmer вҖ” СҒРІРөСӮР»СӢР№ Р¶РөРҪСҒРәРёР№", callback_data=f"{base}:av_voice_shimmer")],
        [InlineKeyboardButton("рҹ“– Fable вҖ” СҒСӮРҫСҖРёСӮРөР»Р»РёРҪРі", callback_data=f"{base}:av_voice_fable")],
        [InlineKeyboardButton("вҡӘ Alloy вҖ” РҪРөР№СӮСҖР°Р»СҢРҪСӢР№", callback_data=f"{base}:av_voice_alloy")],
        [InlineKeyboardButton("рҹҺҷ Р’РјРөСҒСӮРҫ TTS РҝСҖРёСҒР»Р°СӮСҢ СҒРІРҫР№ voice/audio", callback_data=f"{base}:avatar_voice")],
        [InlineKeyboardButton("в¬…пёҸ РқР°Р·Р°Рҙ Рә Р°РІР°СӮР°СҖСғ", callback_data=f"{base}:avatar")],
    ])


def _avatar_voice_choice_text() -> str:
    return (
        "рҹҺҷ РЁР°Рі 2/3: РІСӢРұРөСҖРёСӮРө РіРҫР»РҫСҒ РҙР»СҸ СӮРөРәСҒСӮРҫРІРҫР№ РҫР·РІСғСҮРәРё Р°РІР°СӮР°СҖР°.\n\n"
        "РЎР°РјСӢР№ РөСҒСӮРөСҒСӮРІРөРҪРҪСӢР№ РІР°СҖРёР°РҪСӮ вҖ” РҝСҖРёСҒР»Р°СӮСҢ СҒРІРҫР№ MP3/WAV/M4A/AAC РёР»Рё Telegram voice СҮРөСҖРөР· РәРҪРҫРҝРәСғ В«Voice/audio вҶ’ РіРҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖВ».\n"
        "Р•СҒР»Рё РҪСғР¶РөРҪ СҒРёРҪСӮРөР· РҝРҫ СӮРөРәСҒСӮСғ, РІСӢРұРөСҖРёСӮРө РҫРҙРёРҪ РёР· РіРҫР»РҫСҒРҫРІ РҪРёР¶Рө, Р·Р°СӮРөРј РҝСҖРёСҲР»РёСӮРө С„СҖР°Р·Сғ РҙР»СҸ Р°РІР°СӮР°СҖР°."
    )

def _avatar_menu_text() -> str:
    return (
        "рҹ—Ј *Р“РҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ*\n"
        "РӨРҫСӮРҫ СҮРөР»РҫРІРөРәР° вҶ’ СӮРөРәСҒСӮ, voice РёР»Рё Р°СғРҙРёРҫ вҶ’ РІРёРҙРөРҫ СҒ СҒРёРҪС…СҖРҫРҪРёР·Р°СҶРёРөР№ РіСғРұ.\n\n"
        "РҹРҫСҒР»РөРҙРҫРІР°СӮРөР»СҢРҪРҫСҒСӮСҢ РҙР»СҸ СӮРөРәСҒСӮР°:\n"
        "1) Р·Р°РіСҖСғР·РёСӮСҢ РҝРҫСҖСӮСҖРөСӮ;\n"
        "2) РІСӢРұСҖР°СӮСҢ РіРҫР»РҫСҒ;\n"
        "3) РҝСҖРёСҒР»Р°СӮСҢ СӮРөРәСҒСӮ вҖ” РҝРҫСҒР»Рө СҚСӮРҫРіРҫ СҒСҖР°Р·Сғ Р·Р°РҝСғСҒРәР°РөСӮСҒСҸ Kling Avatar.\n\n"
        "Р”Р»СҸ СҒР°РјРҫРіРҫ РөСҒСӮРөСҒСӮРІРөРҪРҪРҫРіРҫ СҖРөР·СғР»СҢСӮР°СӮР° РјРҫР¶РҪРҫ РІРјРөСҒСӮРҫ TTS РҝСҖРёСҒР»Р°СӮСҢ СҒРІРҫР№ voice/audio вҖ” СӮРҫРіРҙР° Р°РІР°СӮР°СҖ РіРҫРІРҫСҖРёСӮ СҖРөР°Р»СҢРҪСӢРј Р·Р°РҝРёСҒР°РҪРҪСӢРј РіРҫР»РҫСҒРҫРј.\n\n"
        "РӣСғСҮСҲРө РІСҒРөРіРҫ СҖР°РұРҫСӮР°РөСӮ С„СҖРҫРҪСӮР°Р»СҢРҪСӢР№ РҝРҫСҖСӮСҖРөСӮ: Р»РёСҶРҫ РІРёРҙРҪРҫ РҝРҫР»РҪРҫСҒСӮСҢСҺ, РұРөР· СҒРёР»СҢРҪРҫРіРҫ РҝРҫРІРҫСҖРҫСӮР° Рё РұРөР· Р·Р°РәСҖСӢСӮРҫРіРҫ СҖСӮР°."
    )


def _photoclip_action_kb(prefix: str = "act") -> InlineKeyboardMarkup:
    """РҹРҫРҙРјРөРҪСҺ С„РҫСӮРҫвҶ’РІРёРҙРөРҫРәР»РёРҝ СҒ РҝСҖРөСҒРөСӮР°РјРё/РІР°СҖРёР°СҶРёСҸРјРё."""
    base = "act:fun" if prefix == "act" else "fun"
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("рҹ“ё Р—Р°РіСҖСғР·РёСӮСҢ С„РҫСӮРҫ", callback_data=f"{base}:photoclip_upload")],
        [InlineKeyboardButton("вң… РҳСҒРҝРҫР»СҢР·РҫРІР°СӮСҢ РҝРҫСҒР»РөРҙРҪРөРө С„РҫСӮРҫ", callback_data=f"{base}:photoclip_last")],
        [InlineKeyboardButton("рҹҺ¬ РҡРёРҪРөРјР°СӮРҫРіСҖР°С„РёСҮРҪСӢР№ РәР»РёРҝ", callback_data=f"{base}:pc_preset_cinematic")],
        [InlineKeyboardButton("рҹ“ұ Reels / Shorts 9:16", callback_data=f"{base}:pc_preset_reels")],
        [InlineKeyboardButton("рҹҸқ Travel / luxury РәР»РёРҝ", callback_data=f"{base}:pc_preset_luxury")],
        [InlineKeyboardButton("рҹ’ғ РңСғР·СӢРәР°Р»СҢРҪСӢР№ РәР»РёРҝ", callback_data=f"{base}:pc_preset_music")],
        [InlineKeyboardButton("рҹ“Ј Р РөРәР»Р°РјРҪСӢР№ РәР»РёРҝ", callback_data=f"{base}:pc_preset_ads")],
        [InlineKeyboardButton("рҹ“қ РЎРІРҫР№ РҝСҖРҫРјРҝСӮ", callback_data=f"{base}:photoclip_custom")],
        [InlineKeyboardButton("в¬…пёҸ РқР°Р·Р°Рҙ РІ Р Р°Р·РІР»РөСҮРөРҪРёСҸ", callback_data="mode:fun" if prefix == "act" else "fun:back")],
    ])


def _photoclip_menu_text() -> str:
    return (
        "рҹҺө *РӨРҫСӮРҫ вҶ’ РІРёРҙРөРҫРәР»РёРҝ*\n"
        "РӨРҫСӮРҫ СҮРөР»РҫРІРөРәР° РёР»Рё РҫРұСҠРөРәСӮР° вҶ’ РәРҫСҖРҫСӮРәРёР№ РәР»РёРҝ СҮРөСҖРөР· Kling imageвҶ’video.\n\n"
        "Р’СӢРұРөСҖРёСӮРө РІР°СҖРёР°РҪСӮ РҪРёР¶Рө: РјРҫР¶РҪРҫ СҒСҖР°Р·Сғ Р·Р°РҝСғСҒСӮРёСӮСҢ РіРҫСӮРҫРІСӢР№ РҝСҖРөСҒРөСӮ РёР»Рё РІСӢРұСҖР°СӮСҢ В«РЎРІРҫР№ РҝСҖРҫРјРҝСӮВ» Рё РҫРҝРёСҒР°СӮСҢ РҙРІРёР¶РөРҪРёРө, РјСғР·СӢРәСғ, РҪР°СҒСӮСҖРҫРөРҪРёРө, РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ Рё С„РҫСҖРјР°СӮ.\n\n"
        "РҹРҫ СғРјРҫР»СҮР°РҪРёСҺ: 9:16, РҙРёРҪР°РјРёСҮРҪСӢР№ РІРёРҙРөРҫРәР»РёРҝ, Р·РІСғРә Kling вҖ” РөСҒР»Рё РәР°РҪР°Р» Comet РҝРҫРҙРҙРөСҖР¶РёРІР°РөСӮ sound:on."
    )


def _photoclip_preset_prompt(kind: str) -> str:
    kind = (kind or "").strip().lower()
    presets = {
        "cinematic": "cinematic portrait music video, slow dolly-in camera movement, realistic face and body motion, soft dramatic light, shallow depth of field, premium film look, 5 seconds, vertical 9:16",
        "reels": "viral Reels/Shorts style clip, energetic camera movement, modern social media pacing, confident pose, dynamic background motion, clean realistic look, 5 seconds, vertical 9:16",
        "luxury": "luxury travel music video, tropical premium lifestyle mood, warm sunset light, smooth camera movement, elegant cinematic motion, realistic identity preservation, 5 seconds, vertical 9:16",
        "music": "music video performance, the person moves naturally to the beat, expressive face, subtle body movement, dynamic camera, stylish lighting, realistic cinematic clip, 5 seconds, vertical 9:16",
        "ads": "premium advertising video clip, product/commercial style, confident natural movement, clean composition, luxury brand look, smooth camera motion, realistic high quality, 5 seconds, vertical 9:16",
    }
    return presets.get(kind, presets["cinematic"])




def _vocal_clip_menu_text() -> str:
    return (
        "рҹҺӨ *РҡР»РёРҝ СҒ РІРҫРәР°Р»РҫРј / lip-sync (1 СҮРөР»РҫРІРөРә)*\n\n"
        "РӯСӮРҫСӮ СҖРөР¶РёРј РҪСғР¶РөРҪ, РәРҫРіРҙР° РҝРөСҖСҒРҫРҪР°Р¶ РҪР° С„РҫСӮРҫ РҙРҫР»Р¶РөРҪ *РҝРөСӮСҢ РҝРҫРҙ РІРҫРәР°Р»*. "
        "Р”Р»СҸ СҒСӮР°РұРёР»СҢРҪРҫРіРҫ СҖРөР·СғР»СҢСӮР°СӮР° СҖР°РұРҫСӮР°РөСӮ СӮРҫР»СҢРәРҫ СҒ *РҫРҙРҪРёРј СҮРөР»РҫРІРөРәРҫРј РІ РәР°РҙСҖРө*.\n\n"
        "РҹРҫСҒР»РөРҙРҫРІР°СӮРөР»СҢРҪРҫСҒСӮСҢ:\n"
        "1) Р·Р°РіСҖСғР·РёСӮРө С„СҖРҫРҪСӮР°Р»СҢРҪСӢР№ РҝРҫСҖСӮСҖРөСӮ РҫРҙРҪРҫРіРҫ СҮРөР»РҫРІРөРәР°;\n"
        "2) РҫРҝРёСҲРёСӮРө РҝРөСҒРҪСҺ: СҒСӮРёР»СҢ, РҪР°СҒСӮСҖРҫРөРҪРёРө, СҸР·СӢРә, РҝСҖРёРҝРөРІ, РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ;\n"
        "3) РұРҫСӮ РіРөРҪРөСҖРёСҖСғРөСӮ РІРҫРәР°Р» СҮРөСҖРөР· Suno;\n"
        "4) Р·Р°СӮРөРј РҙРөР»Р°РөСӮ lip-sync СҮРөСҖРөР· Kling Avatar.\n\n"
        "Р’Р°Р¶РҪРҫ: РөСҒР»Рё РҪР° С„РҫСӮРҫ РҙРІР° СҮРөР»РҫРІРөРәР°, РәР°СҮРөСҒСӮРІРҫ lip-sync СҖРөР·РәРҫ РҝР°РҙР°РөСӮ. Р”Р»СҸ РҙСғСҚСӮР° Р»СғСҮСҲРө РҙРөР»Р°СӮСҢ РҙРІРө РҫСӮРҙРөР»СҢРҪСӢРө СҒСҶРөРҪСӢ РҝРҫ РҫРҙРҪРҫРјСғ СҮРөР»РҫРІРөРәСғ."
    )


def _vocal_clip_action_kb(prefix: str = "act") -> InlineKeyboardMarkup:
    base = "act:fun" if prefix == "act" else "fun"
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("рҹ“ё Р—Р°РіСҖСғР·РёСӮСҢ РҝРҫСҖСӮСҖРөСӮ 1 СҮРөР»РҫРІРөРәР°", callback_data=f"{base}:vocalclip_upload")],
        [InlineKeyboardButton("вң… РҳСҒРҝРҫР»СҢР·РҫРІР°СӮСҢ РҝРҫСҒР»РөРҙРҪРөРө С„РҫСӮРҫ", callback_data=f"{base}:vocalclip_last")],
        [InlineKeyboardButton("рҹ“қ РһРҝРёСҒР°СӮСҢ РҝРөСҒРҪСҺ/РәР»РёРҝ", callback_data=f"{base}:vocalclip_prompt")],
        [InlineKeyboardButton("в¬…пёҸ РқР°Р·Р°Рҙ РІ Р Р°Р·РІР»РөСҮРөРҪРёСҸ", callback_data="mode:fun" if prefix == "act" else "fun:back")],
    ])


def _textvideo_menu_text() -> str:
    return (
        "рҹҺ¬ *Р’РёРҙРөРҫ РҝРҫ СӮРөРәСҒСӮСғ РёР»Рё РіРҫР»РҫСҒСғ*\n\n"
        "РһРҝРёСҲРёСӮРө СҒСҶРөРҪСғ СӮРөРәСҒСӮРҫРј РёР»Рё РҫСӮРҝСҖР°РІСҢСӮРө voice вҖ” РұРҫСӮ СҖР°СҒРҝРҫР·РҪР°РөСӮ СҖРөСҮСҢ Рё Р·Р°РҝСғСҒСӮРёСӮ РіРөРҪРөСҖР°СҶРёСҺ.\n\n"
        "Р”РІРёР¶РәРё:\n"
        "вҖў *Sora 2 вҖ” СӮРҫР»СҢРәРҫ РұРөР· Р»СҺРҙРөР№* вҖ” РҝСҖРөРҙРјРөСӮСӢ, Р¶РёРІРҫСӮРҪСӢРө, Р·РҙР°РҪРёСҸ, РҝРөР№Р·Р°Р¶Рё, РёРҪСӮРөСҖСҢРөСҖ.\n"
        "вҖў *Kling* вҖ” Р»СҺРҙРё, РҙРёРҪР°РјРёРәР°, Reels/Shorts Рё СғРҪРёРІРөСҖСҒР°Р»СҢРҪСӢРө СҒСҶРөРҪСӢ.\n"
        "вҖў *Runway* вҖ” Р»СҺРҙРё Рё РәРёРҪРөРјР°СӮРҫРіСҖР°С„РёСҮРҪСӢРө СҒСҶРөРҪСӢ; РҫС„РёСҶРёР°Р»СҢРҪСӢР№ Runway API РёСҒРҝРҫР»СҢР·СғРөСӮСҒСҸ РҝРөСҖРІСӢРј, Comet Рё Kling вҖ” СҖРөР·РөСҖРІ.\n\n"
        "РЎСӮРҫРёРјРҫСҒСӮСҢ РҝРҫРәР°Р·СӢРІР°РөСӮСҒСҸ РІ РәСҖРөРҙРёСӮР°С… Рё СҒРҝРёСҒСӢРІР°РөСӮСҒСҸ СӮРҫР»СҢРәРҫ РҝРҫСҒР»Рө СғСҒРҝРөСҲРҪРҫРіРҫ СҖРөР·СғР»СҢСӮР°СӮР°."
    )


def _textvideo_action_kb(prefix: str = "act") -> InlineKeyboardMarkup:
    base = "act:fun" if prefix == "act" else "fun"
    rows = [
        [InlineKeyboardButton(f"рҹҺһ Sora 2 В· РұРөР· Р»СҺРҙРөР№ В· {_video_price_credits('sora', 5)} РәСҖ.", callback_data=f"{base}:tv_engine_sora")],
        [InlineKeyboardButton(f"рҹҺ¬ Kling В· {_video_price_credits('kling', 5)} РәСҖ.", callback_data=f"{base}:tv_engine_kling")],
    ]
    if TEXT_VIDEO_ALLOW_RUNWAY:
        rows.append([InlineKeyboardButton(f"рҹҺҘ Runway В· {_video_price_credits('runway', 5)} РәСҖ.", callback_data=f"{base}:tv_engine_runway")])
    rows.extend([
        [InlineKeyboardButton("рҹ“қ Р’РІРөСҒСӮРё СӮРөРәСҒСӮ / рҹҺҷ РҫСӮРҝСҖР°РІРёСӮСҢ voice", callback_data=f"{base}:tv_prompt")],
        [InlineKeyboardButton("в¬…пёҸ РқР°Р·Р°Рҙ РІ Р Р°Р·РІР»РөСҮРөРҪРёСҸ", callback_data="mode:fun" if prefix == "act" else "fun:back")],
    ])
    return InlineKeyboardMarkup(rows)

def _ai_selfie_action_kb(prefix: str = "act") -> InlineKeyboardMarkup:
    """РҹРҫРҙРјРөРҪСҺ Nano Banana/Gemini style AI selfie."""
    base = "act:fun" if prefix == "act" else "fun"
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("рҹ“ё Р—Р°РіСҖСғР·РёСӮСҢ СҒРІРҫС‘ СҒРөР»С„Рё", callback_data=f"{base}:aiselfie_upload")],
        [InlineKeyboardButton("вң… РҳСҒРҝРҫР»СҢР·РҫРІР°СӮСҢ РҝРҫСҒР»РөРҙРҪРөРө С„РҫСӮРҫ", callback_data=f"{base}:aiselfie_last")],
        [InlineKeyboardButton("вӯҗ РЎРҫ Р·РҪР°РјРөРҪРёСӮРҫСҒСӮСҢСҺ", callback_data=f"{base}:aiselfie_custom")],
        [InlineKeyboardButton("рҹҰё РЎ РәРёРҪРҫРіРөСҖРҫРөРј / РҝРөСҖСҒРҫРҪР°Р¶РөРј", callback_data=f"{base}:as_preset_character")],
        [InlineKeyboardButton("рҹҺ¬ РҡРёРҪРҫСҒСҶРөРҪР° / РҝСҖРөРјСҢРөСҖР°", callback_data=f"{base}:as_preset_movie")],
        [InlineKeyboardButton("рҹҸқ Travel / luxury selfie", callback_data=f"{base}:as_preset_luxury")],
        [InlineKeyboardButton("рҹ“Ј Р РөРәР»Р°РјРҪСӢР№ РәР°РҙСҖ", callback_data=f"{base}:as_preset_ads")],
        [InlineKeyboardButton("рҹ“қ РЎРІРҫР№ РҝСҖРҫРјРҝСӮ", callback_data=f"{base}:aiselfie_custom")],
        [InlineKeyboardButton("в¬…пёҸ РқР°Р·Р°Рҙ РІ Р Р°Р·РІР»РөСҮРөРҪРёСҸ", callback_data="mode:fun" if prefix == "act" else "fun:back")],
    ])


def _ai_selfie_menu_text() -> str:
    return (
        "рҹӨі *AI-СҒРөР»С„Рё СҒРҫ Р·РІРөР·РҙРҫР№ / РҝРөСҖСҒРҫРҪР°Р¶РөРј*\n"
        "Р—Р°РіСҖСғР·РёСӮРө СҒРІРҫС‘ СҒРөР»С„Рё, Р·Р°СӮРөРј РјРҫРҙРөР»СҢ РҝРөСҖРөСҒРҫРұРөСҖС‘СӮ РҪРҫРІСғСҺ СҖРөР°Р»РёСҒСӮРёСҮРҪСғСҺ AI-С„РҫСӮРҫСҒСҶРөРҪСғ "
        "СҒ РҪСғР¶РҪРҫР№ Р·РҪР°РјРөРҪРёСӮРҫСҒСӮСҢСҺ, РҝРөСҖСҒРҫРҪР°Р¶РөРј, Р»РҫРәР°СҶРёРөР№ РёР»Рё СҖРөРәР»Р°РјРҪСӢРј СҒСҺР¶РөСӮРҫРј, СҒСӮР°СҖР°СҸСҒСҢ СҒРҫС…СҖР°РҪРёСӮСҢ РІР°СҲРө Р»РёСҶРҫ.\n\n"
        "РҡР°Рә РҝРҫР»СҢР·РҫРІР°СӮСҢСҒСҸ:\n"
        "1) Р·Р°РіСҖСғР·РёСӮРө СҒРІРҫС‘ СҒРөР»С„Рё РёР»Рё РёСҒРҝРҫР»СҢР·СғР№СӮРө РҝРҫСҒР»РөРҙРҪРөРө С„РҫСӮРҫ;\n"
        "2) РҪР°РҝРёСҲРёСӮРө: `СҒРөР»С„Рё СҒ ...`, РҪР°РҝСҖРёРјРөСҖ: `СҒРөР»С„Рё СҒ РёР·РІРөСҒСӮРҪСӢРј Р°РәСӮС‘СҖРҫРј РҪР° РәСҖР°СҒРҪРҫР№ РҙРҫСҖРҫР¶РәРө, iPhone selfie, 4:5`;\n"
        "3) РұРҫСӮ РІРөСҖРҪС‘СӮ РҪРҫРІСғСҺ AI-С„РҫСӮРҫРіСҖР°С„РёСҺ.\n\n"
        "Р”Р»СҸ РұРөР·РҫРҝР°СҒРҪРҫРіРҫ РёСҒРҝРҫР»СҢР·РҫРІР°РҪРёСҸ РҪРө РҝСҖРёРјРөРҪСҸР№СӮРө СҖРөР·СғР»СҢСӮР°СӮ РәР°Рә РҙРҫРәР°Р·Р°СӮРөР»СҢСҒСӮРІРҫ СҖРөР°Р»СҢРҪРҫР№ РІСҒСӮСҖРөСҮРё, СҖРөРәР»Р°РјСӢ, РҝРҫРҙРҙРөСҖР¶РәРё РёР»Рё РҪРҫРІРҫСҒСӮРё."
    )


def _ai_selfie_preset_prompt(kind: str) -> str:
    kind = (kind or "").strip().lower()
    presets = {
        "character": "Create a realistic fan selfie with a famous movie character specified by the user. If no character is specified, ask the user to name the character. iPhone selfie look, natural lighting, realistic skin, same user's face and identity, no text, no logos, 4:5.",
        "movie": "Create a realistic red-carpet movie premiere selfie. The user is standing next to the named celebrity or actor. Premium event lighting, shallow depth of field, iPhone selfie perspective, same user's face and identity, no text, no logos, 4:5.",
        "luxury": "Create a luxury travel selfie. The user is posing with the named celebrity or public figure in a premium travel location, yacht/hotel/resort mood, natural iPhone selfie, same user's face and identity, no text, no logos, 4:5.",
        "ads": "Create a premium advertising style AI selfie/poster scene with the user and the named celebrity/character. High-end commercial photography, clean composition, same user's face and identity, no fake endorsement text, no logos, 4:5.",
    }
    return presets.get(kind, "")


def _is_ai_selfie_intent(text: str) -> bool:
    t = (text or "").lower().replace("С‘", "Рө")
    return bool(re.search(r"(ai[-\s]?СҒРөР»С„Рё|Р°Рё[-\s]?СҒРөР»С„Рё|СҒРөР»С„Рё\s+СҒРҫ\s+Р·РІРөР·Рҙ|СҒРөР»С„Рё\s+СҒ\s+РёР·РІРөСҒСӮРҪ|СҒРөР»С„Рё\s+СҒ\s+Р°РәСӮ|СҒРөР»С„Рё\s+СҒ\s+РҝРөРІ|СҒРөР»С„Рё\s+СҒ\s+РҝРөСҖСҒРҫРҪР°Р¶|nano\s*banana|РҪР°РҪРҫ\s*РұР°РҪР°РҪ|С„РҫСӮРҫ\s+СҒРҫ\s+Р·РІРөР·Рҙ|С„РҫСӮРҫ\s+СҒ\s+Р·РҪР°РјРөРҪРёСӮ|celebrity\s+selfie|ai\s+selfie)", t, re.I))


def _clean_ai_selfie_prompt(text: str) -> str:
    t = (text or "").strip()
    t = re.sub(r"^(?:СҒРҙРөР»Р°Р№|СҒРҫР·РҙР°Р№|СҒРіРөРҪРөСҖРёСҖСғР№|РҪР°СҖРёСҒСғР№)?\s*(?:ai[-\s]?СҒРөР»С„Рё|Р°Рё[-\s]?СҒРөР»С„Рё|СҒРөР»С„Рё|С„РҫСӮРҫ|nano\s*banana|РҪР°РҪРҫ\s*РұР°РҪР°РҪ|celebrity\s*selfie|ai\s*selfie)\s*[:пјҡ,-]?\s*", "", t, flags=re.I)
    return t.strip()


def _set_ai_selfie_wait(context: ContextTypes.DEFAULT_TYPE):
    context.user_data["awaiting_ai_selfie_prompt"] = True
    context.user_data.pop("awaiting_avatar_script", None)
    context.user_data.pop("awaiting_photo_clip_prompt", None)




def _clear_vocal_clip_wait(context: ContextTypes.DEFAULT_TYPE):
    context.user_data.pop("awaiting_vocal_clip_photo", None)
    context.user_data.pop("awaiting_vocal_clip_prompt", None)
    context.user_data.pop("vocal_clip_preset_prompt", None)


def _set_vocal_clip_wait(context: ContextTypes.DEFAULT_TYPE):
    context.user_data["awaiting_vocal_clip_prompt"] = True


def _clear_text_video_wait(context: ContextTypes.DEFAULT_TYPE):
    context.user_data.pop("awaiting_text_video_prompt", None)


def _set_text_video_wait(context: ContextTypes.DEFAULT_TYPE, engine: str | None = None):
    context.user_data["awaiting_text_video_prompt"] = True
    if engine:
        context.user_data["text_video_engine"] = engine


async def _handle_vocalclip_upload_choice(update: Update, context: ContextTypes.DEFAULT_TYPE, q, prefix: str = "act"):
    _clear_transient_flows(context)
    _set_mode_clean(q.from_user.id, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "vocalclip")
    context.user_data["awaiting_vocal_clip_photo"] = True
    await q.message.reply_text(
        "рҹҺӨ РҹСҖРёСҲР»РёСӮРө С„СҖРҫРҪСӮР°Р»СҢРҪСӢР№ РҝРҫСҖСӮСҖРөСӮ *РҫРҙРҪРҫРіРҫ СҮРөР»РҫРІРөРәР°*. Р•СҒР»Рё РІ РәР°РҙСҖРө РҙРІР° СҮРөР»РҫРІРөРәР°, lip-sync РұСғРҙРөСӮ РҪРөСҒСӮР°РұРёР»СҢРҪСӢРј.",
        parse_mode="Markdown",
        reply_markup=_vocal_clip_action_kb(prefix),
    )


async def _handle_vocalclip_prompt_choice(update: Update, context: ContextTypes.DEFAULT_TYPE, q, prefix: str = "act"):
    _clear_transient_flows(context)
    _set_mode_clean(q.from_user.id, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "vocalclip")
    img = _get_cached_photo(q.from_user.id)
    if img:
        _set_vocal_clip_wait(context)
        await q.answer("Р“РҫСӮРҫРІРҫ")
        await q.message.reply_text(
            "рҹҺӨ РҳСҒРҝРҫР»СҢР·СғСҺ РҝРҫСҒР»РөРҙРҪРөРө С„РҫСӮРҫ. РўРөРҝРөСҖСҢ РҫРҝРёСҲРёСӮРө РҝРөСҒРҪСҺ/РәР»РёРҝ: СҒСӮРёР»СҢ, СҸР·СӢРә, РҪР°СҒСӮСҖРҫРөРҪРёРө, РҝСҖРёРҝРөРІ, РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ.\n\n"
            "Р’Р°Р¶РҪРҫ: СҖРөР¶РёРј СҖР°СҒСҒСҮРёСӮР°РҪ РҪР° РҫРҙРҪРҫРіРҫ СҮРөР»РҫРІРөРәР° РІ РәР°РҙСҖРө."
        )
    else:
        context.user_data["awaiting_vocal_clip_photo"] = True
        await q.message.reply_text(
            "РЎРҪР°СҮР°Р»Р° РҝСҖРёСҲР»РёСӮРө РҝРҫСҖСӮСҖРөСӮ РҫРҙРҪРҫРіРҫ СҮРөР»РҫРІРөРәР°. РҹРҫСҒР»Рө Р·Р°РіСҖСғР·РәРё СҸ РҝРҫРҝСҖРҫСҲСғ РҫРҝРёСҒР°РҪРёРө РҝРөСҒРҪРё/РәР»РёРҝР°.",
            reply_markup=_vocal_clip_action_kb(prefix),
        )


async def _handle_textvideo_engine_choice(update: Update, context: ContextTypes.DEFAULT_TYPE, q, engine: str, prefix: str = "act"):
    _clear_transient_flows(context)
    _set_mode_clean(q.from_user.id, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "textvideo")
    engine = (engine or TEXT_VIDEO_DEFAULT_ENGINE or "kling").strip().lower()
    if engine == "runway" and not TEXT_VIDEO_ALLOW_RUNWAY:
        engine = "kling"
    if engine not in ("sora", "kling", "runway"):
        engine = "kling"
    _set_text_video_wait(context, engine)
    labels = {"sora": "Sora 2 вҖ” СӮРҫР»СҢРәРҫ РұРөР· Р»СҺРҙРөР№", "kling": "Kling", "runway": "Runway"}
    await q.answer(labels[engine][:180])
    warning = "\nвҡ пёҸ Р’ РҝСҖРҫРјРҝСӮРө РҪРө РҙРҫР»Р¶РҪРҫ РұСӢСӮСҢ Р»СҺРҙРөР№, Р»РёСҶ РёР»Рё РҝРөСҖСҒРҫРҪР°Р¶РөР№." if engine == "sora" else ""
    await q.message.reply_text(
        f"вң… Р’СӢРұСҖР°РҪ РҙРІРёР¶РҫРә: {labels[engine]}.{warning}\n"
        "РўРөРҝРөСҖСҢ РҪР°РҝРёСҲРёСӮРө РҝСҖРҫРјРҝСӮ РёР»Рё РҫСӮРҝСҖР°РІСҢСӮРө voice. РЈРәР°Р¶РёСӮРө СҒСҶРөРҪСғ, СҒСӮРёР»СҢ, РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ Рё С„РҫСҖРјР°СӮ: РҪР°РҝСҖРёРјРөСҖ `10 СҒРөРәСғРҪРҙ, 16:9`."
    )


async def _handle_textvideo_prompt_choice(update: Update, context: ContextTypes.DEFAULT_TYPE, q, prefix: str = "act"):
    _clear_transient_flows(context)
    _set_mode_clean(q.from_user.id, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "textvideo")
    _set_text_video_wait(context, context.user_data.get("text_video_engine") or TEXT_VIDEO_DEFAULT_ENGINE)
    await q.answer("Р’РёРҙРөРҫ")
    await q.message.reply_text(
        "рҹҺ¬ РқР°РҝРёСҲРёСӮРө СӮРөРәСҒСӮРҫРІСӢР№ РҝСҖРҫРјРҝСӮ РёР»Рё РҫСӮРҝСҖР°РІСҢСӮРө voice. РҹРҫ СғРјРҫР»СҮР°РҪРёСҺ РёСҒРҝРҫР»СҢР·СғСҺ Kling. "
        "РңРҫР¶РҪРҫ СҒРҪР°СҮР°Р»Р° РІСӢРұСҖР°СӮСҢ Sora 2 РұРөР· Р»СҺРҙРөР№, Kling РёР»Рё Runway РәРҪРҫРҝРәРҫР№ РІСӢСҲРө."
    )

def _clear_ai_selfie_wait(context: ContextTypes.DEFAULT_TYPE):
    context.user_data.pop("awaiting_ai_selfie_prompt", None)
    context.user_data.pop("awaiting_ai_selfie_photo", None)
    context.user_data.pop("ai_selfie_preset_prompt", None)


async def _handle_avatar_upload_choice(update: Update, context: ContextTypes.DEFAULT_TYPE, q, prefix: str = "act"):
    _clear_transient_flows(context)
    _set_mode_clean(q.from_user.id, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "avatar")
    context.user_data["awaiting_avatar_photo"] = True
    await q.message.reply_text(
        "рҹ—Ј РЁР°Рі 1/3: РҝСҖРёСҲР»РёСӮРө РҝРҫСҖСӮСҖРөСӮ СҮРөР»РҫРІРөРәР°. РҹРҫСҒР»Рө Р·Р°РіСҖСғР·РәРё СҸ РҝСҖРөРҙР»РҫР¶Сғ РІСӢРұСҖР°СӮСҢ РіРҫР»РҫСҒ, Р·Р°СӮРөРј РҝРҫРҝСҖРҫСҲСғ СӮРөРәСҒСӮ. Р•СҒР»Рё РҪСғР¶РөРҪ СҖРөР°Р»СҢРҪСӢР№ РіРҫР»РҫСҒ вҖ” РІСӢРұРөСҖРёСӮРө voice/audio СҖРөР¶РёРј.",
        reply_markup=_avatar_action_kb(prefix),
    )


async def _handle_avatar_script_choice(update: Update, context: ContextTypes.DEFAULT_TYPE, q, prefix: str = "act", voice_mode: bool = False):
    _clear_transient_flows(context)
    _set_mode_clean(q.from_user.id, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "avatar")
    img = _get_cached_photo(q.from_user.id)
    if img:
        if voice_mode:
            _set_avatar_wait(context)
            await q.message.reply_text("рҹҺҷ РҳСҒРҝРҫР»СҢР·СғСҺ РҝРҫСҒР»РөРҙРҪРөРө Р·Р°РіСҖСғР¶РөРҪРҪРҫРө С„РҫСӮРҫ. РўРөРҝРөСҖСҢ РҝСҖРёСҲР»РёСӮРө voice/audio РҙР»СҸ СҖРөСҮРё Р°РІР°СӮР°СҖР°.")
        else:
            _set_avatar_voice_choice_wait(context)
            await q.message.reply_text(_avatar_voice_choice_text(), reply_markup=_avatar_voice_choice_kb(prefix))
        await q.answer("Р“РҫСӮРҫРІРҫ")
    else:
        context.user_data["awaiting_avatar_photo"] = True
        await q.message.reply_text(
            "РЎРҪР°СҮР°Р»Р° РҝСҖРёСҲР»РёСӮРө РҝРҫСҖСӮСҖРөСӮ СҮРөР»РҫРІРөРәР°. РҹРҫСҒР»Рө Р·Р°РіСҖСғР·РәРё СҸ РҝСҖРөРҙР»РҫР¶Сғ РІСӢРұСҖР°СӮСҢ РіРҫР»РҫСҒ, Р·Р°СӮРөРј РҝРҫРҝСҖРҫСҲСғ СӮРөРәСҒСӮ. Р”Р»СҸ СҖРөР°Р»СҢРҪРҫРіРҫ РіРҫР»РҫСҒР° РјРҫР¶РҪРҫ РІСӢРұСҖР°СӮСҢ voice/audio СҖРөР¶РёРј.",
            reply_markup=_avatar_action_kb(prefix),
        )


async def _handle_photoclip_upload_choice(update: Update, context: ContextTypes.DEFAULT_TYPE, q, prefix: str = "act"):
    _clear_transient_flows(context)
    _set_mode_clean(q.from_user.id, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "photoclip")
    context.user_data["awaiting_photo_clip_photo"] = True
    await q.message.reply_text(
        "рҹҺө РҹСҖРёСҲР»РёСӮРө С„РҫСӮРҫ СҮРөР»РҫРІРөРәР°/РҫРұСҠРөРәСӮР°. РҹРҫСҒР»Рө Р·Р°РіСҖСғР·РәРё СҸ РҝРҫРҝСҖРҫСҲСғ РҫРҝРёСҒР°РҪРёРө РәР»РёРҝР° РёР»Рё Р·Р°РҝСғСүСғ РІСӢРұСҖР°РҪРҪСӢР№ РҝСҖРөСҒРөСӮ.",
        reply_markup=_photoclip_action_kb(prefix),
    )


async def _handle_photoclip_prompt_choice(update: Update, context: ContextTypes.DEFAULT_TYPE, q, prefix: str = "act"):
    _clear_transient_flows(context)
    _set_mode_clean(q.from_user.id, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "photoclip")
    img = _get_cached_photo(q.from_user.id)
    if img:
        _set_photo_clip_wait(context)
        await q.message.reply_text("рҹҺө РҳСҒРҝРҫР»СҢР·СғСҺ РҝРҫСҒР»РөРҙРҪРөРө Р·Р°РіСҖСғР¶РөРҪРҪРҫРө С„РҫСӮРҫ. РһРҝРёСҲРёСӮРө СҒСӮРёР»СҢ РІРёРҙРөРҫРәР»РёРҝР°: РјСғР·СӢРәР°, РҙРІРёР¶РөРҪРёРө, РҪР°СҒСӮСҖРҫРөРҪРёРө, РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ Рё С„РҫСҖРјР°СӮ 9:16/16:9.")
        await q.answer("Р“РҫСӮРҫРІРҫ")
    else:
        context.user_data["awaiting_photo_clip_photo"] = True
        await q.message.reply_text(
            "РЎРҪР°СҮР°Р»Р° РҝСҖРёСҲР»РёСӮРө С„РҫСӮРҫ. РҹРҫСҒР»Рө Р·Р°РіСҖСғР·РәРё РҫРҝРёСҲРөСӮРө СҒСӮРёР»СҢ РІРёРҙРөРҫРәР»РёРҝР°.",
            reply_markup=_photoclip_action_kb(prefix),
        )


async def _handle_photoclip_preset_choice(update: Update, context: ContextTypes.DEFAULT_TYPE, q, kind: str, prefix: str = "act"):
    prompt = _photoclip_preset_prompt(kind)
    _clear_transient_flows(context)
    _set_mode_clean(q.from_user.id, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "photoclip")
    img = _get_cached_photo(q.from_user.id)
    if img:
        await q.answer("Р—Р°РҝСғСҒРәР°СҺ РәР»РёРҝ")
        await q.message.reply_text("рҹҺ¬ РҳСҒРҝРҫР»СҢР·СғСҺ РҝРҫСҒР»РөРҙРҪРөРө Р·Р°РіСҖСғР¶РөРҪРҪРҫРө С„РҫСӮРҫ Рё Р·Р°РҝСғСҒРәР°СҺ РІСӢРұСҖР°РҪРҪСӢР№ РҝСҖРөСҒРөСӮ РІРёРҙРөРҫРәР»РёРҝР°.")
        await _start_photo_music_clip(update, context, img, prompt)
    else:
        context.user_data["awaiting_photo_clip_photo"] = True
        context.user_data["photo_clip_preset_prompt"] = prompt
        await q.message.reply_text(
            "рҹҺ¬ РҹСҖРөСҒРөСӮ РІСӢРұСҖР°РҪ. РўРөРҝРөСҖСҢ РҝСҖРёСҲР»РёСӮРө С„РҫСӮРҫ вҖ” РҝРҫСҒР»Рө Р·Р°РіСҖСғР·РәРё РәР»РёРҝ Р·Р°РҝСғСҒСӮРёСӮСҒСҸ Р°РІСӮРҫРјР°СӮРёСҮРөСҒРәРё.",
            reply_markup=_photoclip_action_kb(prefix),
        )

# РҹРҫРәР°Р·Р°СӮСҢ РІСӢРұСҖР°РҪРҪСӢР№ СҖРөР¶РёРј (РёСҒРҝРҫР»СҢР·СғРөСӮСҒСҸ Рё РҙР»СҸ callback, Рё РҙР»СҸ СӮРөРәСҒСӮР°)
async def _send_mode_menu(update, context, key: str):
    text = _mode_desc(key)
    kb = _mode_kb(key)
    # Р•СҒР»Рё РҝСҖРёСҲР»Рё РёР· callback вҖ” СҖРөРҙР°РәСӮРёСҖСғРөРј; РөСҒР»Рё СӮРөРәСҒСӮРҫРј вҖ” СҲР»С‘Рј РҪРҫРІСӢРј СҒРҫРҫРұСүРөРҪРёРөРј
    if getattr(update, "callback_query", None):
        q = update.callback_query
        await q.message.reply_text(text, reply_markup=kb, parse_mode="Markdown")
        await q.answer()
    else:
        await update.effective_message.reply_text(text, reply_markup=kb, parse_mode="Markdown")

# РһРұСҖР°РұРҫСӮСҮРёРә callback РҝРҫ СҖРөР¶РёРјР°Рј
async def on_mode_cb(update, context):
    q = update.callback_query
    data = (q.data or "").strip()
    uid = q.from_user.id

    # РқР°РІРёРіР°СҶРёСҸ
    if data == "mode:root":
        await q.message.reply_text(_modes_root_text(), reply_markup=modes_root_kb())
        await q.answer(); return

    if data.startswith("mode:"):
        _, key = data.split(":", 1)
        _clear_transient_flows(context)
        if key == "medicine":
            _set_mode_clean(uid, "РңРөРҙРёСҶРёРҪР°", "")
        elif key == "study":
            _set_mode_clean(uid, "РЈСҮС‘РұР°", "")
        elif key == "work":
            _set_mode_clean(uid, "Р Р°РұРҫСӮР°/Р‘РёР·РҪРөСҒ", "")
        elif key == "fun":
            _set_mode_clean(uid, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "")
        await _send_mode_menu(update, context, key)
        return

    # РЎРІРҫРұРҫРҙРҪСӢР№ РІРІРҫРҙ РёР· РҝРҫРҙРјРөРҪСҺ
    if data == "act:free":
        await q.answer()
        await q.message.reply_text(
            "рҹ“қ РқР°РҝРёСҲРёСӮРө СҒРІРҫРұРҫРҙРҪСӢР№ Р·Р°РҝСҖРҫСҒ РҪРёР¶Рө СӮРөРәСҒСӮРҫРј РёР»Рё РіРҫР»РҫСҒРҫРј вҖ” СҸ РҝРҫРҙСҒСӮСҖРҫСҺСҒСҢ.",
            reply_markup=modes_root_kb(),
        )
        return

    # === РЈСҮС‘РұР°
    if data == "act:study:pdf_summary":
        await q.answer()
        _mode_track_set(uid, "pdf_summary")
        await q.message.reply_text(
            "рҹ“ҡ РҹСҖРёСҲР»РёСӮРө PDF/EPUB/DOCX/FB2/TXT вҖ” СҒРҙРөР»Р°СҺ СҒСӮСҖСғРәСӮСғСҖРёСҖРҫРІР°РҪРҪСӢР№ РәРҫРҪСҒРҝРөРәСӮ.\n"
            "РңРҫР¶РҪРҫ РІ РҝРҫРҙРҝРёСҒРё СғРәР°Р·Р°СӮСҢ СҶРөР»СҢ (РәРҫСҖРҫСӮРәРҫ/РҝРҫРҙСҖРҫРұРҪРҫ, СҸР·СӢРә Рё СӮ.Рҝ.).",
            reply_markup=_mode_kb("study"),
        )
        return

    if data == "act:study:explain":
        await q.answer()
        study_sub_set(uid, "explain")
        _mode_track_set(uid, "explain")
        await q.message.reply_text(
            "рҹ”Қ РқР°РҝРёСҲРёСӮРө СӮРөРјСғ + СғСҖРҫРІРөРҪСҢ (СҲРәРҫР»Р°/РІСғР·/РҝСҖРҫС„Рё). Р‘СғРҙРөСӮ РҫРұСҠСҸСҒРҪРөРҪРёРө СҒ РҝСҖРёРјРөСҖР°РјРё.",
            reply_markup=_mode_kb("study"),
        )
        return

    if data == "act:study:tasks":
        await q.answer()
        study_sub_set(uid, "tasks")
        _mode_track_set(uid, "tasks")
        await q.message.reply_text(
            "рҹ§® РҹСҖРёСҲР»РёСӮРө СғСҒР»РҫРІРёРө(СҸ) вҖ” СҖРөСҲСғ РҝРҫСҲР°РіРҫРІРҫ (С„РҫСҖРјСғР»СӢ, РҝРҫСҸСҒРҪРөРҪРёСҸ, РёСӮРҫРі).",
            reply_markup=_mode_kb("study"),
        )
        return

    if data == "act:study:essay":
        await q.answer()
        study_sub_set(uid, "essay")
        _mode_track_set(uid, "essay")
        await q.message.reply_text(
            "вңҚпёҸ РўРөРјР° + СӮСҖРөРұРҫРІР°РҪРёСҸ (РҫРұСҠС‘Рј/СҒСӮРёР»СҢ/СҸР·СӢРә) вҖ” РҝРҫРҙРіРҫСӮРҫРІР»СҺ СҚСҒСҒРө/СҖРөС„РөСҖР°СӮ.",
            reply_markup=_mode_kb("study"),
        )
        return

    if data == "act:study:exam_plan":
        await q.answer()
        study_sub_set(uid, "quiz")
        _mode_track_set(uid, "exam_plan")
        await q.message.reply_text(
            "рҹ“қ РЈРәР°Р¶РёСӮРө РҝСҖРөРҙРјРөСӮ Рё РҙР°СӮСғ СҚРәР·Р°РјРөРҪР° вҖ” СҒРҫСҒСӮР°РІР»СҺ РҝР»Р°РҪ РҝРҫРҙРіРҫСӮРҫРІРәРё СҒ РІРөС…Р°РјРё.",
            reply_markup=_mode_kb("study"),
        )
        return

    # === Р Р°РұРҫСӮР°
    if data == "act:work:doc":
        await q.answer()
        _mode_track_set(uid, "work_doc")
        await q.message.reply_text(
            "рҹ“„ Р§СӮРҫ Р·Р° РҙРҫРәСғРјРөРҪСӮ/Р°РҙСҖРөСҒР°СӮ/РәРҫРҪСӮРөРәСҒСӮ? РЎС„РҫСҖРјРёСҖСғСҺ СҮРөСҖРҪРҫРІРёРә РҝРёСҒСҢРјР°/РҙРҫРәСғРјРөРҪСӮР°.",
            reply_markup=_mode_kb("work"),
        )
        return

    if data == "act:work:report":
        await q.answer()
        _mode_track_set(uid, "work_report")
        await q.message.reply_text(
            "рҹ“Ҡ РҹСҖРёСҲР»РёСӮРө СӮРөРәСҒСӮ/С„Р°Р№Р»/СҒСҒСӢР»РәСғ вҖ” СҒРҙРөР»Р°СҺ Р°РҪР°Р»РёСӮРёСҮРөСҒРәСғСҺ РІСӢР¶РёРјРәСғ.",
            reply_markup=_mode_kb("work"),
        )
        return

    if data == "act:work:plan":
        await q.answer()
        _mode_track_set(uid, "work_plan")
        await q.message.reply_text(
            "рҹ—Ӯ РһРҝРёСҲРёСӮРө Р·Р°РҙР°СҮСғ/СҒСҖРҫРәРё вҖ” СҒРҫРұРөСҖСғ ToDo/РҝР»Р°РҪ СҒРҫ СҒСҖРҫРәР°РјРё Рё РҝСҖРёРҫСҖРёСӮРөСӮР°РјРё.",
            reply_markup=_mode_kb("work"),
        )
        return

    if data == "act:work:idea":
        await q.answer()
        _mode_track_set(uid, "work_idea")
        await q.message.reply_text(
            "рҹ’Ў Р Р°СҒСҒРәР°Р¶РёСӮРө РҝСҖРҫРҙСғРәСӮ/РҰРҗ/РәР°РҪР°Р»СӢ вҖ” РҝРҫРҙРіРҫСӮРҫРІР»СҺ РұСҖРёС„/РёРҙРөРё.",
            reply_markup=_mode_kb("work"),
        )
        return

    if data == "act:work:presentation":
        await q.answer()
        _clear_transient_flows(context)
        _set_mode_clean(uid, "Р Р°РұРҫСӮР°/Р‘РёР·РҪРөСҒ", "work_presentation")
        await _presentation_studio_get().start(update, context, "presentation")
        return

    if data == "act:work:catalog_pdf":
        await q.answer()
        _clear_transient_flows(context)
        _set_mode_clean(uid, "Р Р°РұРҫСӮР°/Р‘РёР·РҪРөСҒ", "work_catalog_pdf")
        await _presentation_studio_get().start(update, context, "catalog")
        return

    if data == "act:work:logo":
        await q.answer()
        _clear_transient_flows(context)
        _set_mode_clean(uid, "Р Р°РұРҫСӮР°/Р‘РёР·РҪРөСҒ", "work_logo")
        context.user_data["awaiting_work_logo_brief"] = True
        await q.message.reply_text(
            "рҹҺЁ РһРҝРёСҲРёСӮРө Р»РҫРіРҫСӮРёРҝ: РҪР°Р·РІР°РҪРёРө РұСҖРөРҪРҙР°, РҪРёСҲР°, СҒСӮРёР»СҢ, СҶРІРөСӮР°, СҒР»РҫРіР°РҪ, РіРҙРө РұСғРҙРөСӮ РёСҒРҝРҫР»СҢР·РҫРІР°СӮСҢСҒСҸ. РҜ СҒРіРөРҪРөСҖРёСҖСғСҺ РІРёР·СғР°Р».",
            reply_markup=_mode_kb("work"),
        )
        return

    if data == "act:work:watermark":
        await q.answer()
        _clear_medicine_wait(context)
        _set_mode_clean(uid, "Р Р°РұРҫСӮР°/Р‘РёР·РҪРөСҒ", "work_watermark")
        img = _get_cached_photo(uid)
        if img:
            _set_retouch_wait_text(context)
            await q.message.reply_text(
                "рҹ§Ҫ РҳСҒРҝРҫР»СҢР·СғСҺ РҝРҫСҒР»РөРҙРҪРөРө Р·Р°РіСҖСғР¶РөРҪРҪРҫРө С„РҫСӮРҫ. РқР°РҝРёСҲРёСӮРө, СҮСӮРҫ СғРұСҖР°СӮСҢ Рё РіРҙРө РҪР°С…РҫРҙРёСӮСҒСҸ СҚР»РөРјРөРҪСӮ: РІРҫРҙСҸРҪРҫР№ Р·РҪР°Рә, РҪР°РҙРҝРёСҒСҢ, Р»РҫРіРҫСӮРёРҝ РёР»Рё Р»РёСҲРҪРёР№ РҫРұСҠРөРәСӮ.\n"
                "РһСӮРҝСҖР°РІР»СҸСҸ РәРҫРјР°РҪРҙСғ, РІСӢ РҝРҫРҙСӮРІРөСҖР¶РҙР°РөСӮРө, СҮСӮРҫ СҚСӮРҫ РІР°СҲРө РёР·РҫРұСҖР°Р¶РөРҪРёРө РёР»Рё Сғ РІР°СҒ РөСҒСӮСҢ РҝСҖР°РІРҫ РөРіРҫ СҖРөРҙР°РәСӮРёСҖРҫРІР°СӮСҢ."
            )
        else:
            _set_waiting_image_retouch(update, context, "СғРұСҖР°СӮСҢ РІРҫРҙСҸРҪРҫР№ Р·РҪР°Рә/РҪР°РҙРҝРёСҒСҢ/Р»РҫРіРҫСӮРёРҝ Рё РІРҫСҒСҒСӮР°РҪРҫРІРёСӮСҢ С„РҫРҪ")
            await q.message.reply_text(
                "рҹ§Ҫ РҹСҖРёСҲР»РёСӮРө С„РҫСӮРҫ, Р·Р°СӮРөРј РҪР°РҝРёСҲРёСӮРө, СҮСӮРҫ СғРұСҖР°СӮСҢ: РІРҫРҙСҸРҪРҫР№ Р·РҪР°Рә, РҪР°РҙРҝРёСҒСҢ, Р»РҫРіРҫСӮРёРҝ РёР»Рё Р»РёСҲРҪРёР№ РҫРұСҠРөРәСӮ.",
                reply_markup=_mode_kb("work"),
            )
        return

    # === Р Р°Р·РІР»РөСҮРөРҪРёСҸ: РіРҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ Рё С„РҫСӮРҫвҶ’РІРёРҙРөРҫРәР»РёРҝ
    if data == "act:fun:avatar":
        await q.answer("Р“РҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ")
        _clear_transient_flows(context)
        _set_mode_clean(uid, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "avatar")
        await q.message.reply_text(_avatar_menu_text(), parse_mode="Markdown", reply_markup=_avatar_action_kb("act"))
        return

    if data.startswith("act:fun:av_voice_"):
        voice = data.rsplit("_", 1)[-1].strip()
        context.user_data["avatar_tts_voice"] = voice
        context.user_data.pop("awaiting_avatar_voice_choice", None)
        if _get_cached_photo(q.from_user.id):
            _set_avatar_wait(context)
            await q.answer(f"Р“РҫР»РҫСҒ: {voice}")
            await q.message.reply_text(f"вң… Р”Р»СҸ Р°РІР°СӮР°СҖР° РІСӢРұСҖР°РҪ РіРҫР»РҫСҒ: {_avatar_tts_voice_label(voice)}. РЁР°Рі 3/3: РҝСҖРёСҲР»РёСӮРө СӮРөРәСҒСӮ, РәРҫСӮРҫСҖСӢР№ РҙРҫР»Р¶РөРҪ РҝСҖРҫРёР·РҪРөСҒСӮРё Р°РІР°СӮР°СҖ.")
        else:
            context.user_data["awaiting_avatar_photo"] = True
            await q.answer(f"Р“РҫР»РҫСҒ: {voice}")
            await q.message.reply_text(f"вң… Р“РҫР»РҫСҒ РІСӢРұСҖР°РҪ: {_avatar_tts_voice_label(voice)}. РўРөРҝРөСҖСҢ РҝСҖРёСҲР»РёСӮРө РҝРҫСҖСӮСҖРөСӮ СҮРөР»РҫРІРөРәР°.")
        return

    if data == "act:fun:avatar_upload":
        await q.answer("Р—Р°РіСҖСғР·РёСӮРө РҝРҫСҖСӮСҖРөСӮ")
        await _handle_avatar_upload_choice(update, context, q, prefix="act")
        return

    if data in ("act:fun:avatar_last", "act:fun:avatar_text"):
        await _handle_avatar_script_choice(update, context, q, prefix="act", voice_mode=False)
        return

    if data == "act:fun:avatar_voice":
        await _handle_avatar_script_choice(update, context, q, prefix="act", voice_mode=True)
        return

    if data == "act:fun:vocalclip":
        await q.answer("РҡР»РёРҝ СҒ РІРҫРәР°Р»РҫРј")
        _clear_transient_flows(context)
        _set_mode_clean(uid, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "vocalclip")
        await q.message.reply_text(_vocal_clip_menu_text(), parse_mode="Markdown", reply_markup=_vocal_clip_action_kb("act"))
        return

    if data == "act:fun:vocalclip_upload":
        await q.answer("Р—Р°РіСҖСғР·РёСӮРө РҝРҫСҖСӮСҖРөСӮ")
        await _handle_vocalclip_upload_choice(update, context, q, prefix="act")
        return

    if data in ("act:fun:vocalclip_last", "act:fun:vocalclip_prompt"):
        await _handle_vocalclip_prompt_choice(update, context, q, prefix="act")
        return

    if data == "act:fun:textvideo":
        await q.answer("Р’РёРҙРөРҫ РҝРҫ СӮРөРәСҒСӮСғ/РіРҫР»РҫСҒСғ")
        _clear_transient_flows(context)
        _set_mode_clean(uid, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "textvideo")
        await q.message.reply_text(_textvideo_menu_text(), parse_mode="Markdown", reply_markup=_textvideo_action_kb("act"))
        return

    if data == "act:fun:tv_engine_sora":
        await _handle_textvideo_engine_choice(update, context, q, "sora", prefix="act")
        return

    if data == "act:fun:tv_engine_kling":
        await _handle_textvideo_engine_choice(update, context, q, "kling", prefix="act")
        return

    if data == "act:fun:tv_engine_runway":
        await _handle_textvideo_engine_choice(update, context, q, "runway", prefix="act")
        return

    if data == "act:fun:tv_prompt":
        await _handle_textvideo_prompt_choice(update, context, q, prefix="act")
        return

    if data == "act:fun:photoclip":
        await q.answer("РӨРҫСӮРҫ вҶ’ РІРёРҙРөРҫРәР»РёРҝ")
        _clear_transient_flows(context)
        _set_mode_clean(uid, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "photoclip")
        await q.message.reply_text(_photoclip_menu_text(), parse_mode="Markdown", reply_markup=_photoclip_action_kb("act"))
        return

    if data == "act:fun:photoclip_upload":
        await q.answer("Р—Р°РіСҖСғР·РёСӮРө С„РҫСӮРҫ")
        await _handle_photoclip_upload_choice(update, context, q, prefix="act")
        return

    if data in ("act:fun:photoclip_last", "act:fun:photoclip_custom"):
        await _handle_photoclip_prompt_choice(update, context, q, prefix="act")
        return

    if data.startswith("act:fun:pc_preset_"):
        kind = data.rsplit("_", 1)[-1]
        await _handle_photoclip_preset_choice(update, context, q, kind, prefix="act")
        return

    if data == "act:fun:aiselfie":
        await q.answer("AI-СҒРөР»С„Рё")
        _clear_transient_flows(context)
        _set_mode_clean(uid, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "aiselfie")
        await q.message.reply_text(_ai_selfie_menu_text(), parse_mode="Markdown", reply_markup=_ai_selfie_action_kb("act"))
        return

    if data == "act:fun:aiselfie_upload":
        await q.answer("Р—Р°РіСҖСғР·РёСӮРө СҒРөР»С„Рё")
        _clear_transient_flows(context)
        _set_mode_clean(uid, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "aiselfie")
        context.user_data["awaiting_ai_selfie_photo"] = True
        await q.message.reply_text("рҹӨі РҹСҖРёСҲР»РёСӮРө СҒРІРҫС‘ СҒРөР»С„Рё. РҹРҫСҒР»Рө Р·Р°РіСҖСғР·РәРё СҸ РҝРҫРҝСҖРҫСҲСғ РҪР°РҝРёСҒР°СӮСҢ, СҒ РәРөРј/РіРҙРө СҒРҙРөР»Р°СӮСҢ AI-С„РҫСӮРҫ.", reply_markup=_ai_selfie_action_kb("act"))
        return

    if data in ("act:fun:aiselfie_last", "act:fun:aiselfie_custom"):
        _clear_transient_flows(context)
        _set_mode_clean(uid, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "aiselfie")
        img = _get_cached_photo(uid)
        if img:
            _set_ai_selfie_wait(context)
            await q.answer("Р“РҫСӮРҫРІРҫ")
            await q.message.reply_text("рҹӨі РҳСҒРҝРҫР»СҢР·СғСҺ РҝРҫСҒР»РөРҙРҪРөРө С„РҫСӮРҫ. РқР°РҝРёСҲРёСӮРө СҒСҶРөРҪСғ: РҪР°РҝСҖРёРјРөСҖ В«СҒРөР»С„Рё СҒ РёР·РІРөСҒСӮРҪСӢРј Р°РәСӮС‘СҖРҫРј РҪР° РәСҖР°СҒРҪРҫР№ РҙРҫСҖРҫР¶РәРө, iPhone selfie, 4:5В».")
        else:
            context.user_data["awaiting_ai_selfie_photo"] = True
            await q.message.reply_text("РЎРҪР°СҮР°Р»Р° РҝСҖРёСҲР»РёСӮРө СҒРІРҫС‘ СҒРөР»С„Рё. РҹРҫСҒР»Рө Р·Р°РіСҖСғР·РәРё РҪР°РҝРёСҲРөСӮРө, СҒ РәРөРј/РіРҙРө СҒРҙРөР»Р°СӮСҢ AI-С„РҫСӮРҫ.", reply_markup=_ai_selfie_action_kb("act"))
        return

    if data.startswith("act:fun:as_preset_"):
        kind = data.rsplit("_", 1)[-1]
        preset = _ai_selfie_preset_prompt(kind)
        _clear_transient_flows(context)
        _set_mode_clean(uid, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "aiselfie")
        img = _get_cached_photo(uid)
        if img:
            context.user_data["ai_selfie_preset_prompt"] = preset
            _set_ai_selfie_wait(context)
            await q.answer("РҹСҖРөСҒРөСӮ РІСӢРұСҖР°РҪ")
            await q.message.reply_text("рҹӨі РҹСҖРөСҒРөСӮ РІСӢРұСҖР°РҪ. РўРөРҝРөСҖСҢ РҪР°РҝРёСҲРёСӮРө РёРјСҸ Р·РҪР°РјРөРҪРёСӮРҫСҒСӮРё/РҝРөСҖСҒРҫРҪР°Р¶Р° Рё РҙРөСӮР°Р»Рё СҒСҶРөРҪСӢ.")
        else:
            context.user_data["awaiting_ai_selfie_photo"] = True
            context.user_data["ai_selfie_preset_prompt"] = preset
            await q.message.reply_text("рҹӨі РҹСҖРөСҒРөСӮ РІСӢРұСҖР°РҪ. РўРөРҝРөСҖСҢ РҝСҖРёСҲР»РёСӮРө СҒРІРҫС‘ СҒРөР»С„Рё вҖ” РҝРҫСҒР»Рө Р·Р°РіСҖСғР·РәРё СҸ РҝРҫРҝСҖРҫСҲСғ РёРјСҸ Р·РҪР°РјРөРҪРёСӮРҫСҒСӮРё/РҝРөСҖСҒРҫРҪР°Р¶Р°.", reply_markup=_ai_selfie_action_kb("act"))
        return

    # === Р Р°Р·РІР»РөСҮРөРҪРёСҸ (РәР°Рә РұСӢР»Рҫ)
    if data == "act:fun:ideas":
        await q.answer()
        await q.message.reply_text(
            "рҹ”Ҙ Р’СӢРұРөСҖРөРј С„РҫСҖРјР°СӮ: РҙРҫРј/СғР»РёСҶР°/РіРҫСҖРҫРҙ/РІ РҝРҫРөР·РҙРәРө. РқР°РҝРёСҲРёСӮРө РұСҺРҙР¶РөСӮ/РҪР°СҒСӮСҖРҫРөРҪРёРө.",
            reply_markup=_mode_kb("fun"),
        )
        return
    if data == "act:fun:shorts":
        await q.answer()
        await q.message.reply_text(
            "рҹҺ¬ РўРөРјР°, РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ (15вҖ“30 СҒРөРә), СҒСӮРёР»СҢ вҖ” СҒРҙРөР»Р°СҺ СҒСҶРөРҪР°СҖРёР№ СҲРҫСҖСӮР° + РҝРҫРҙСҒРәР°Р·РәРё РҙР»СҸ РҫР·РІСғСҮРәРё.",
            reply_markup=_mode_kb("fun"),
        )
        return
    if data == "act:fun:games":
        await q.answer()
        await q.message.reply_text(
            "рҹҺ® РўРөРјР°СӮРёРәР° РәРІРёР·Р°/РёРіСҖСӢ? РЎРіРөРҪРөСҖРёСҖСғСҺ РұСӢСҒСӮСҖСғСҺ РІРёРәСӮРҫСҖРёРҪСғ РёР»Рё РјРёРҪРё-РёРіСҖСғ РІ СҮР°СӮРө.",
            reply_markup=_mode_kb("fun"),
        )
        return

    if data == "act:fun:faceswap":
        await q.answer()
        _clear_transient_flows(context)
        _set_mode_clean(uid, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "faceswap")
        await _start_faceswap_flow(update, context, None, use_cached=False)
        return

    if data == "act:fun:removebg":
        await q.answer()
        _clear_transient_flows(context)
        _set_mode_clean(uid, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "removebg")
        img = _get_cached_photo(uid)
        if img:
            await q.message.reply_text("рҹ§ј РҳСҒРҝРҫР»СҢР·СғСҺ РҝРҫСҒР»РөРҙРҪРөРө Р·Р°РіСҖСғР¶РөРҪРҪРҫРө С„РҫСӮРҫ Рё СғРҙР°Р»СҸСҺ С„РҫРҪ.")
            await _pedit_removebg(update, context, img)
        else:
            _set_waiting_removebg(context)
            await q.message.reply_text("рҹ§ј РҹСҖРёСҲР»РёСӮРө С„РҫСӮРҫ вҖ” СғРҙР°Р»СҺ С„РҫРҪ Рё РІРөСҖРҪСғ PNG СҒ РҝСҖРҫР·СҖР°СҮРҪРҫР№ РҝРҫРҙР»РҫР¶РәРҫР№.", reply_markup=_mode_kb("fun"))
        return

    if data == "act:fun:replacebg":
        await q.answer()
        _clear_transient_flows(context)
        _set_mode_clean(uid, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "replacebg")
        img = _get_cached_photo(uid)
        if img:
            await q.message.reply_text("рҹ–ј Р’СӢРұРөСҖРёСӮРө РҪРҫРІСӢР№ С„РҫРҪ РҙР»СҸ РҝРҫСҒР»РөРҙРҪРөРіРҫ Р·Р°РіСҖСғР¶РөРҪРҪРҫРіРҫ С„РҫСӮРҫ:", reply_markup=background_presets_kb())
        else:
            context.user_data["photo_flow"] = "replacebg_menu"
            await q.message.reply_text("рҹ–ј РҹСҖРёСҲР»РёСӮРө С„РҫСӮРҫ. РҹРҫСҒР»Рө Р·Р°РіСҖСғР·РәРё СҸ РҝРҫРәР°Р¶Сғ РІР°СҖРёР°РҪСӮСӢ С„РҫРҪР°: РҝР»СҸР¶, РіРҫСҖСӢ, РҝСҖРёСҖРҫРҙР°, РіРҫСҖРҫРҙ РёР»Рё СҒРІРҫР№ СӮРөРәСҒСӮ.", reply_markup=_mode_kb("fun"))
        return

    if data == "act:fun:revive":
        await q.answer()
        _set_waiting_photo_revival(update, context)
        await q.message.reply_text(
            _fun_revive_help_text(),
            parse_mode="Markdown",
            reply_markup=_mode_kb("fun"),
        )
        return

    if data == "act:fun:reels":
        await q.answer()
        _clear_transient_flows(context)
        _set_mode_clean(uid, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "fun_reels")
        context.user_data["awaiting_reels_material"] = True
        await q.message.reply_text(
            _fun_reels_help_text(),
            parse_mode="Markdown",
            reply_markup=_mode_kb("fun"),
        )
        return

    if data == "act:fun:film":
        await q.answer()
        _clear_transient_flows(context)
        _set_mode_clean(uid, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "fun_film")
        context.user_data["awaiting_film_material"] = True
        await q.message.reply_text(
            _fun_film_help_text(),
            parse_mode="Markdown",
            reply_markup=_mode_kb("fun"),
        )
        return

    if data == "act:fun:music":
        await q.answer("РңСғР·СӢРәР° / Suno")
        _clear_transient_flows(context)
        _set_mode_clean(uid, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "suno_music")
        context.user_data["awaiting_suno_brief"] = True
        await _show_suno_help_from_callback(q, context, reply_markup=_suno_menu_kb(), submenu=True)
        return

    # === РңРөРҙРёСҶРёРҪР°
    if data.startswith("act:med:"):
        await q.answer()
        action = data.split(":", 2)[2]
        mapping = {
            "extract": "med_extract",
            "scan": "med_scan",
            "conclusion": "med_conclusion",
            "mri": "med_mri",
            "ct": "med_ct",
            "labs": "med_labs",
            "free": "med_free",
        }
        track = mapping.get(action, "med_free")
        _set_mode_clean(uid, "РңРөРҙРёСҶРёРҪР°", track)
        _clear_transient_flows(context)
        context.user_data["medicine_waiting_for_material"] = True
        context.user_data["pending_med_task"] = track
        await q.message.reply_text(_medical_menu_text(track), reply_markup=medicine_kb())
        return

    # === РңРҫРҙСғР»Рё (РәР°Рә РұСӢР»Рҫ)
    if data == "act:open:runway":
        await q.answer()
        await q.message.reply_text(
            "рҹҺ¬ РңРҫРҙСғР»СҢ Runway: РҝСҖРёСҲР»РёСӮРө РёРҙРөСҺ/СҖРөС„РөСҖРөРҪСҒ вҖ” РҝРҫРҙРіРҫСӮРҫРІР»СҺ РҝСҖРҫРјРҝСӮ Рё РұСҺРҙР¶РөСӮ.",
            reply_markup=modes_root_kb(),
        )
        return
    if data == "act:open:mj":
        await q.answer()
        await q.message.reply_text(
            "рҹҺЁ РңРҫРҙСғР»СҢ Midjourney: РҫРҝРёСҲРёСӮРө РәР°СҖСӮРёРҪРәСғ вҖ” РҝСҖРөРҙР»РҫР¶Сғ 3 РҝСҖРҫРјРҝСӮР° Рё СҒРөСӮРәСғ СҒСӮРёР»РөР№.",
            reply_markup=modes_root_kb(),
        )
        return
    if data == "act:open:voice":
        await q.answer()
        await q.message.reply_text(
            "рҹ—Ј Р“РҫР»РҫСҒ: /voice_on вҖ” РҫР·РІСғСҮРәР° РҫСӮРІРөСӮРҫРІ, /voice_off вҖ” РІСӢРәР»СҺСҮРёСӮСҢ. "
            "РңРҫР¶РөСӮРө РҝСҖРёСҒР»Р°СӮСҢ РіРҫР»РҫСҒРҫРІРҫРө вҖ” СҖР°СҒРҝРҫР·РҪР°СҺ Рё РҫСӮРІРөСҮСғ.",
            reply_markup=modes_root_kb(),
        )
        return

    await q.answer()

# Fallback вҖ” РөСҒР»Рё РҝРҫР»СҢР·РҫРІР°СӮРөР»СҢ РҪР°Р¶РјС‘СӮ В«РЈСҮС‘РұР°/Р Р°РұРҫСӮР°/Р Р°Р·РІР»РөСҮРөРҪРёСҸВ» РҫРұСӢСҮРҪРҫР№ РәРҪРҫРҝРәРҫР№/СӮРөРәСҒСӮРҫРј
async def on_mode_text(update, context):
    text = (update.effective_message.text or "").strip().lower()
    mapping = {
        "СғСҮС‘РұР°": "study", "СғСҮРөРұР°": "study",
        "СҖР°РұРҫСӮР°": "work", "СҖР°РұРҫСӮР°/РұРёР·РҪРөСҒ": "work", "РұРёР·РҪРөСҒ": "work",
        "СҖР°Р·РІР»РөСҮРөРҪРёСҸ": "fun", "СҖР°Р·РІР»РөСҮРөРҪРёРө": "fun",
        "РјРөРҙРёСҶРёРҪР°": "medicine", "РјРөРҙРёСҶРёРҪa": "medicine",
    }
    key = mapping.get(text)
    if key:
        await _send_mode_menu(update, context, key)
        
def main_keyboard():
    # 2 РәРҪРҫРҝРәРё РІ СҒСӮСҖРҫРәРө вҖ” СӮР°Рә Telegram РҪР° РјРҫРұРёР»СҢРҪСӢС… РҪРө РҫРұСҖРөР·Р°РөСӮ РҝРҫРҙРҝРёСҒРё.
    return ReplyKeyboardMarkup(
        [
            [KeyboardButton("рҹҺ“ РЈСҮС‘РұР°"), KeyboardButton("рҹ’ј Р Р°РұРҫСӮР°/Р‘РёР·РҪРөСҒ")],
            [KeyboardButton("рҹ”Ҙ Р Р°Р·РІР»РөСҮРөРҪРёСҸ"), KeyboardButton("рҹ©ә РңРөРҙРёСҶРёРҪР°")],
            [KeyboardButton("рҹ’¬ РңРҫРё СҮР°СӮСӢ"), KeyboardButton("вһ• РқРҫРІСӢР№ СҮР°СӮ")],
            [KeyboardButton("рҹ§  Р”РІРёР¶РәРё"), KeyboardButton("рҹ§ҫ Р‘Р°Р»Р°РҪСҒ")],
            [KeyboardButton("вӯҗ РҹРҫРҙРҝРёСҒРәР° В· РҹРҫРјРҫСүСҢ")],
        ],
        resize_keyboard=True,
        one_time_keyboard=False,
        selective=False,
        input_field_placeholder="Р’СӢРұРөСҖРёСӮРө СҖРөР¶РёРј РёР»Рё РҪР°РҝРёСҲРёСӮРө Р·Р°РҝСҖРҫСҒвҖҰ",
    )

main_kb = main_keyboard()

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ /start в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
async def cmd_start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    with contextlib.suppress(Exception):
        _chat_ensure_active(update.effective_user.id, update.effective_chat.id)
    await update.effective_chat.send_message(
        START_TEXT,
        reply_markup=main_kb,
        disable_web_page_preview=True,
    )


def _chat_dt_label(ts: int) -> str:
    if not ts:
        return "вҖ”"
    dt = datetime.fromtimestamp(int(ts), tz=timezone.utc)
    now = datetime.now(timezone.utc)
    if dt.date() == now.date():
        return "СҒРөРіРҫРҙРҪСҸ " + dt.strftime("%H:%M")
    if dt.date() == (now - timedelta(days=1)).date():
        return "РІСҮРөСҖР° " + dt.strftime("%H:%M")
    return dt.strftime("%d.%m.%Y")


def _chat_list_kb(user_id: int, telegram_chat_id: int) -> InlineKeyboardMarkup:
    rows = []
    for item in _chat_list(user_id, telegram_chat_id):
        mark = "вң…" if item["active"] else "рҹ’¬"
        title = f"{mark} {item['title']} В· {_chat_dt_label(item['updated_ts'])}"
        rows.append([InlineKeyboardButton(title[:60], callback_data=f"chat:open:{item['id']}")])
        rows.append([
            InlineKeyboardButton("рҹ“– РҳСҒСӮРҫСҖРёСҸ", callback_data=f"chat:history:{item['id']}:0"),
            InlineKeyboardButton("вңҸпёҸ РҳРјСҸ", callback_data=f"chat:rename:{item['id']}"),
            InlineKeyboardButton("рҹ—‘ РЈРҙР°Р»РёСӮСҢ", callback_data=f"chat:delete:{item['id']}"),
        ])
    rows.append([InlineKeyboardButton("вһ• РЎРҫР·РҙР°СӮСҢ РҪРҫРІСӢР№ СҮР°СӮ", callback_data="chat:new")])
    return InlineKeyboardMarkup(rows)


async def cmd_chats(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id; tgid = update.effective_chat.id
    _chat_ensure_active(uid, tgid)
    chats = _chat_list(uid, tgid)
    active = next((x for x in chats if x["active"]), None)
    text = (
        f"рҹ’¬ Р’Р°СҲРё СҮР°СӮСӢ: {len(chats)}/{CHAT_MAX_CONVERSATIONS}.\n"
        "РҡР°Р¶РҙСӢР№ СҮР°СӮ С…СҖР°РҪРёСӮ РҫСӮРҙРөР»СҢРҪСӢР№ РәРҫРҪСӮРөРәСҒСӮ GPT. РқР°Р¶РјРёСӮРө РҪР° СҮР°СӮ, СҮСӮРҫРұСӢ РҝСҖРҫРҙРҫР»Р¶РёСӮСҢ, РёР»Рё РҫСӮРәСҖРҫР№СӮРө РөРіРҫ РёСҒСӮРҫСҖРёСҺ."
    )
    if active:
        text += f"\n\nРЎРөР№СҮР°СҒ Р°РәСӮРёРІРөРҪ: В«{active['title']}В»."
    await update.effective_message.reply_text(text, reply_markup=_chat_list_kb(uid, tgid))


async def cmd_newchat(update: Update, context: ContextTypes.DEFAULT_TYPE):
    uid = update.effective_user.id; tgid = update.effective_chat.id
    cid, err = _chat_create(uid, tgid, "РқРҫРІСӢР№ СҮР°СӮ")
    if not cid:
        await update.effective_message.reply_text(err, reply_markup=_chat_list_kb(uid, tgid))
        return
    _clear_transient_flows(context)
    await update.effective_message.reply_text(
        "вһ• РқРҫРІСӢР№ СҮР°СӮ СҒРҫР·РҙР°РҪ Рё РІСӢРұСҖР°РҪ. РқР°РҝРёСҲРёСӮРө РҝРөСҖРІСӢР№ Р·Р°РҝСҖРҫСҒ вҖ” РҪР°Р·РІР°РҪРёРө РҝРҫСҸРІРёСӮСҒСҸ Р°РІСӮРҫРјР°СӮРёСҮРөСҒРәРё.",
        reply_markup=main_kb,
    )


def _chat_history_page(user_id: int, telegram_chat_id: int, ai_chat_id: int, page: int) -> tuple[list[dict], int, str]:
    page = max(0, int(page or 0)); offset = page * CHAT_HISTORY_PAGE_SIZE
    con = sqlite3.connect(DB_PATH); cur = con.cursor()
    cur.execute("SELECT title FROM ai_chats WHERE id=? AND user_id=? AND telegram_chat_id=?", (int(ai_chat_id), int(user_id), int(telegram_chat_id)))
    row = cur.fetchone()
    if not row:
        con.close(); return [], 0, "Р§Р°СӮ"
    title = row[0] or "Р§Р°СӮ"
    cur.execute("SELECT COUNT(*) FROM ai_chat_messages WHERE ai_chat_id=?", (int(ai_chat_id),))
    total = int((cur.fetchone() or [0])[0])
    # Page 0 is the newest page; inside the page messages are chronological.
    cur.execute(
        "SELECT role, content, created_ts FROM ai_chat_messages WHERE ai_chat_id=? ORDER BY id DESC LIMIT ? OFFSET ?",
        (int(ai_chat_id), CHAT_HISTORY_PAGE_SIZE, offset),
    )
    rows = cur.fetchall(); con.close(); rows.reverse()
    return [{"role": r[0], "content": r[1] or "", "created_ts": int(r[2] or 0)} for r in rows], total, title


async def _send_chat_history(update: Update, context: ContextTypes.DEFAULT_TYPE, ai_chat_id: int, page: int = 0):
    uid = update.effective_user.id; tgid = update.effective_chat.id
    items, total, title = _chat_history_page(uid, tgid, ai_chat_id, page)
    if total <= 0:
        await update.effective_message.reply_text(f"рҹ“– Р’ СҮР°СӮРө В«{title}В» РҝРҫРәР° РҪРөСӮ СҒРҫРҫРұСүРөРҪРёР№.", reply_markup=_chat_list_kb(uid, tgid))
        return
    lines = [f"рҹ“– РҳСҒСӮРҫСҖРёСҸ: {title}", f"РЎРҫРҫРұСүРөРҪРёР№: {total} В· СҒСӮСҖР°РҪРёСҶР° {page + 1}", ""]
    for item in items:
        who = "рҹ‘Ө Р’СӢ" if item["role"] == "user" else "рҹӨ– Р‘РҫСӮ"
        stamp = datetime.fromtimestamp(item["created_ts"], tz=timezone.utc).strftime("%d.%m %H:%M") if item["created_ts"] else ""
        content = (item["content"] or "").strip()
        lines.append(f"{who} В· {stamp}\n{content}\n")
    text = "\n".join(lines)
    # Telegram limit: split long pages safely.
    for i in range(0, len(text), 3800):
        await update.effective_message.reply_text(text[i:i+3800])
    max_page = max(0, (total - 1) // CHAT_HISTORY_PAGE_SIZE)
    nav = []
    if page < max_page:
        nav.append(InlineKeyboardButton("в¬…пёҸ РЎСӮР°СҖРөРө", callback_data=f"chat:history:{ai_chat_id}:{page+1}"))
    if page > 0:
        nav.append(InlineKeyboardButton("РқРҫРІРөРө вһЎпёҸ", callback_data=f"chat:history:{ai_chat_id}:{page-1}"))
    rows = [nav] if nav else []
    rows.append([InlineKeyboardButton("в–¶пёҸ РҹСҖРҫРҙРҫР»Р¶РёСӮСҢ СҚСӮРҫСӮ СҮР°СӮ", callback_data=f"chat:open:{ai_chat_id}")])
    rows.append([InlineKeyboardButton("рҹ’¬ Р’СҒРө СҮР°СӮСӢ", callback_data="chat:list")])
    await update.effective_message.reply_text("Р”РөР№СҒСӮРІРёСҸ СҒ РёСҒСӮРҫСҖРёРөР№:", reply_markup=InlineKeyboardMarkup(rows))


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ СҒРҫС…СҖР°РҪРөРҪРёРө РІСӢРұСҖР°РҪРҪРҫРіРҫ СҖРөР¶РёРјР°/РҝРҫРҙСҖРөР¶РёРјР° (SQLite kv) в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def _mode_set(user_id: int, mode: str):
    kv_set(f"mode:{user_id}", mode)

def _mode_get(user_id: int) -> str:
    return (kv_get(f"mode:{user_id}", "none") or "none")

def _mode_track_set(user_id: int, track: str):
    kv_set(f"mode_track:{user_id}", track)

def _mode_track_get(user_id: int) -> str:
    return kv_get(f"mode_track:{user_id}", "") or ""


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ Р‘РөР·РҫРҝР°СҒРҪР°СҸ РјР°СҖСҲСҖСғСӮРёР·Р°СҶРёСҸ СҖРөР¶РёРјРҫРІ в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def _set_mode_clean(user_id: int, mode: str, track: str = ""):
    """Р•РҙРёРҪСӢР№ РІС…РҫРҙ РҙР»СҸ РҝРөСҖРөРәР»СҺСҮРөРҪРёСҸ СҖРөР¶РёРјР°: РҪРҫРІСӢР№ СҖРөР¶РёРј РІСҒРөРіРҙР° СҒРұСҖР°СҒСӢРІР°РөСӮ СҒСӮР°СҖСӢР№ РҝРҫРҙСҖРөР¶РёРј."""
    with contextlib.suppress(Exception):
        _mode_set(user_id, mode)
    with contextlib.suppress(Exception):
        _mode_track_set(user_id, track or "")


def _clear_transient_flows(context):
    """РЎРұСҖР°СҒСӢРІР°РөСӮ РәСҖР°СӮРәРҫР¶РёРІСғСүРёРө РҫР¶РёРҙР°РҪРёСҸ, РәРҫСӮРҫСҖСӢРө РҪРө РҙРҫР»Р¶РҪСӢ СӮСҸРҪСғСӮСҢСҒСҸ РјРөР¶РҙСғ СҖРөР¶РёРјР°РјРё."""
    if not context:
        return
    for key in (
        "awaiting_photo_for",
        "photo_flow",
        "retouch_prompt",
        "retouch_wait_text",
        "awaiting_med_file",
        "awaiting_med_photo",
        "awaiting_med_document",
        "pending_med_task",
        "medicine_waiting_for_material",
        "awaiting_reels_material",
        "awaiting_film_material",
        "awaiting_suno_brief",
        "suno_preset_kind",
        "awaiting_avatar_photo",
        "awaiting_avatar_voice_choice",
        "awaiting_avatar_script",
        "awaiting_photo_clip_photo",
        "awaiting_photo_clip_prompt",
        "photo_clip_preset_prompt",
        "awaiting_ai_selfie_photo",
        "awaiting_ai_selfie_prompt",
        "ai_selfie_preset_prompt",
        "awaiting_vocal_clip_photo",
        "awaiting_vocal_clip_prompt",
        "vocal_clip_preset_prompt",
        "awaiting_text_video_prompt",
        "text_video_engine",
        "presentation_studio_active",
    ):
        with contextlib.suppress(Exception):
            context.user_data.pop(key, None)


def _set_waiting_photo_revival(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """РҹРҫСҒР»Рө СӮРөРәСҒСӮР°/РіРҫР»РҫСҒР° В«РҫР¶РёРІРёСӮСҢ С„РҫСӮРҫВ» СҒР»РөРҙСғСҺСүР°СҸ С„РҫСӮРҫРіСҖР°С„РёСҸ РҙРҫР»Р¶РҪР° РёРҙСӮРё РІ С„РҫСӮРҫ-РјР°СҒСӮРөСҖСҒРәСғСҺ, Р° РҪРө РІ РјРөРҙРёСҶРёРҪСғ."""
    uid = update.effective_user.id
    _set_mode_clean(uid, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "")
    _clear_transient_flows(context)
    context.user_data["awaiting_photo_for"] = "revive"
    context.user_data["photo_flow"] = "revive"


def _is_waiting_photo_revival(context) -> bool:
    if not context:
        return False
    return (
        context.user_data.get("awaiting_photo_for") == "revive"
        or context.user_data.get("photo_flow") == "revive"
    )


def _clear_photo_revival_wait(context):
    if not context:
        return
    for key in ("awaiting_photo_for", "photo_flow"):
        with contextlib.suppress(Exception):
            context.user_data.pop(key, None)


# Р РөСӮСғСҲСҢ СҒРҫРұСҒСӮРІРөРҪРҪРҫРіРҫ РёР·РҫРұСҖР°Р¶РөРҪРёСҸ РёСҒРҝРҫР»СҢР·СғРөСӮ РҫСӮРҙРөР»СҢРҪСӢР№ photo_flow, СҮСӮРҫРұСӢ РҪРө РәРҫРҪС„Р»РёРәСӮРҫРІР°СӮСҢ СҒ РјРөРҙРёСҶРёРҪРҫР№.
# РӨСғРҪРәСҶРёРё _set_waiting_image_retouch/_is_waiting_image_retouch РҫРұСҠСҸРІР»РөРҪСӢ РІСӢСҲРө РІ РұР»РҫРәРө Image retouch.


def _set_medical_waiting(update: Update, context: ContextTypes.DEFAULT_TYPE, track: str = ""):
    """РңРөРҙРёСҶРёРҪР° СҒСӮР°РҪРҫРІРёСӮСҒСҸ Р°РәСӮРёРІРҪРҫР№ СӮРҫР»СҢРәРҫ РәР°Рә РјРөРҪСҺ/РәРҫРҪРәСҖРөСӮРҪР°СҸ РјРөРҙ. Р·Р°РҙР°СҮР°, РҪРҫ РҪРө РҙРҫР»Р¶РҪР° РҝРөСҖРөС…РІР°СӮСӢРІР°СӮСҢ С„РҫСӮРҫ РҝРҫСҒР»Рө РҙСҖСғРіРёС… Р·Р°РҝСҖРҫСҒРҫРІ."""
    uid = update.effective_user.id
    _set_mode_clean(uid, "РңРөРҙРёСҶРёРҪР°", track or "")
    _clear_transient_flows(context)
    if track:
        context.user_data["medicine_waiting_for_material"] = True
        context.user_data["pending_med_task"] = track


def _clear_medicine_wait(context):
    if not context:
        return
    for key in ("medicine_waiting_for_material", "pending_med_task", "awaiting_med_file", "awaiting_med_photo", "awaiting_med_document"):
        with contextlib.suppress(Exception):
            context.user_data.pop(key, None)


def _should_route_medical(context: ContextTypes.DEFAULT_TYPE, user_id: int, caption_or_text: str = "", filename: str = "") -> bool:
    """РңР°СҖСҲСҖСғСӮРёР·РёСҖСғРөСӮ РІ РјРөРҙРёСҶРёРҪСғ СӮРҫР»СҢРәРҫ СҸРІРҪСӢР№ РјРөРҙ. РҝРҫРҙСҖРөР¶РёРј/РҫР¶РёРҙР°РҪРёРө РёР»Рё СҸРІРҪСӢРө РјРөРҙ. СҒР»РҫРІР°.
    РЎР°Рј С„Р°РәСӮ, СҮСӮРҫ РҝРҫР»СҢР·РҫРІР°СӮРөР»СҢ РәРҫРіРҙР°-СӮРҫ РҪР°Р¶Р°Р» В«РңРөРҙРёСҶРёРҪР°В», РұРҫР»СҢСҲРө РҪРө РҙРөР»Р°РөСӮ РІСҒРө СҒР»РөРҙСғСҺСүРёРө С„РҫСӮРҫ РјРөРҙРёСҶРёРҪСҒРәРёРјРё.
    """
    track = ""
    with contextlib.suppress(Exception):
        track = _mode_track_get(user_id)
    combined = f"{caption_or_text or ''} {filename or ''}"
    # Р’ РјРөРҙРёСҶРёРҪСғ РҫСӮРҝСҖР°РІР»СҸРөРј СӮРҫР»СҢРәРҫ СҸРІРҪСӢР№ РјРөРҙРёСҶРёРҪСҒРәРёР№ РјР°СӮРөСҖРёР°Р» РёР»Рё РјР°СӮРөСҖРёР°Р»,
    # РәРҫСӮРҫСҖСӢР№ РҝРҫР»СҢР·РҫРІР°СӮРөР»СҢ РҝСҖСҸРјРҫ Р¶РҙРөСӮ РҝРҫСҒР»Рө РҪР°Р¶Р°СӮРёСҸ РјРөРҙ. РҝРҫРҙРјРөРҪСҺ.
    # РЎСӮР°СҖСӢР№ mode_track=med_* СҒР°Рј РҝРҫ СҒРөРұРө РұРҫР»СҢСҲРө РҪРө РҙРҫР»Р¶РөРҪ РҝРөСҖРөС…РІР°СӮСӢРІР°СӮСҢ РІСҒРө РҝРҫРҙСҖСҸРҙ.
    return (
        bool(_MEDICAL_TERMS_RE.search(combined))
        or bool(context and context.user_data.get("medicine_waiting_for_material") and (track or "").startswith("med_"))
    )


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РўРөРәСҒСӮСӢ РұСӢСҒСӮСҖСӢС… РҙРөР№СҒСӮРІРёР№ В«Р Р°Р·РІР»РөСҮРөРҪРёСҸВ» в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def _fun_revive_help_text() -> str:
    return (
        "рҹӘ„ *РһР¶РёРІРёСӮСҢ С„РҫСӮРҫ / С„РҫСӮРҫвҶ’РІРёРҙРөРҫ*\n"
        "РңРҫР¶РҪРҫ СҒРҙРөР»Р°СӮСҢ РәРҫСҖРҫСӮРәРҫРө РІРёРҙРөРҫ РёР· РҫРҙРҪРҫР№ С„РҫСӮРҫРіСҖР°С„РёРё, РіРҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ РёР»Рё РәР»РёРҝРҫРІСғСҺ Р°РҪРёРјР°СҶРёСҺ.\n\n"
        "РҡР°Рә Р·Р°РҝСғСҒСӮРёСӮСҢ:\n"
        "1) РҫСӮРҝСҖР°РІСҢСӮРө С„РҫСӮРҫ РҫРұСӢСҮРҪРҫР№ РәР°СҖСӮРёРҪРәРҫР№;\n"
        "2) РҝРҫСҒР»Рө Р·Р°РіСҖСғР·РәРё РҝРҫСҸРІСҸСӮСҒСҸ РәРҪРҫРҝРәРё вңЁ РһР¶РёРІРёСӮСҢ, рҹ—Ј Р“РҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ Рё рҹҺө РӨРҫСӮРҫвҶ’РІРёРҙРөРҫРәР»РёРҝ;\n"
        "3) Р»РёРұРҫ СҒСҖР°Р·Сғ РҫСӮРҝСҖР°РІСҢСӮРө С„РҫСӮРҫ СҒ РҝРҫРҙРҝРёСҒСҢСҺ: `РҫР¶РёРІРё С„РҫСӮРҫ: Р»С‘РіРәР°СҸ СғР»СӢРұРәР°, РҝР»Р°РІРҪРҫРө РҙРІРёР¶РөРҪРёРө РәР°РјРөСҖСӢ, 5 СҒРөРәСғРҪРҙ, 9:16`;\n"
        "4) РҙР»СҸ РіРҫРІРҫСҖСҸСүРөРіРҫ Р°РІР°СӮР°СҖР° РҪР°Р¶РјРёСӮРө рҹ—Ј Рё РҝСҖРёСҲР»РёСӮРө СӮРөРәСҒСӮ, voice РёР»Рё Р°СғРҙРёРҫС„Р°Р№Р» 2вҖ“60 СҒРөРәСғРҪРҙ.\n\n"
        "Р РөРәРҫРјРөРҪРҙР°СҶРёСҸ: Runway вҖ” СҒСӮР°РұРёР»СҢРҪРҫРө РҫР¶РёРІР»РөРҪРёРө РҝРҫСҖСӮСҖРөСӮРҫРІ; Kling вҖ” РҙРёРҪР°РјРёРәР°, Р°РІР°СӮР°СҖСӢ Рё РәР»РёРҝСӢ; Sora 2 вҖ” СӮРҫР»СҢРәРҫ РҙР»СҸ СҒСҶРөРҪ РұРөР· Р»СҺРҙРөР№, РөСҒР»Рё РәР°РҪР°Р» РІ Comet РҙРҫСҒСӮСғРҝРөРҪ."
    )


def _fun_reels_help_text() -> str:
    return (
        "рҹ“ұ *РЎРҙРөР»Р°СӮСҢ Reels / Shorts*\n"
        "РҜ РјРҫРіСғ СҒРҫРұСҖР°СӮСҢ РёРҙРөСҺ, С…СғРә, СҒСҶРөРҪР°СҖРёР№, СӮР°Р№Рј-РәРҫРҙСӢ, РҝРҫРҙРҝРёСҒРё, РҝСҖРҫРјРҝСӮСӢ РҙР»СҸ РІРёРҙРөРҫРіРөРҪРөСҖР°СҶРёРё Рё РҝР»Р°РҪ РјРҫРҪСӮР°Р¶Р°.\n\n"
        "РҡР°Рә Р·Р°РҝСғСҒСӮРёСӮСҢ:\n"
        "вҖў РҪР°РҝРёСҲРёСӮРө СӮРөРјСғ: `СҒРҙРөР»Р°Р№ СҖРёР»СҒ 20 СҒРөРәСғРҪРҙ РҝСҖРҫ РІРёР»Р»Сғ РҪР° РЎР°РјСғРё, luxury, 9:16`;\n"
        "вҖў РёР»Рё РҝСҖРёСҲР»РёСӮРө РёСҒС…РҫРҙРҪРҫРө РІРёРҙРөРҫ/С„РҫСӮРҫ Рё РҝРҫРҙРҝРёСҒСҢ: СҮСӮРҫ РҫСҒСӮР°РІРёСӮСҢ, СҒСӮРёР»СҢ, РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ, РҰРҗ;\n"
        "вҖў РҙР»СҸ AI-РІСҒСӮР°РІРҫРә Р»СғСҮСҲРө РёСҒРҝРҫР»СҢР·РҫРІР°СӮСҢ Runway/Kling; Sora 2 вҖ” СӮРҫР»СҢРәРҫ РөСҒР»Рё РҙРҫСҒСӮСғРҝРөРҪ РәР°РҪР°Р» РҝСҖРҫРІР°Р№РҙРөСҖР°.\n\n"
        "РӨРҫСҖРјР°СӮ СҖРөР·СғР»СҢСӮР°СӮР°: hook вҶ’ СҒСҶРөРҪСӢ вҶ’ СӮРөРәСҒСӮ РҪР° СҚРәСҖР°РҪРө вҶ’ voice-over вҶ’ CTA вҶ’ РҝСҖРҫРјРҝСӮСӢ РҙР»СҸ РІСӢРұСҖР°РҪРҪРҫРіРҫ РҙРІРёР¶РәР°."
    )


def _fun_film_help_text() -> str:
    return (
        "рҹҺһ *РЎРҫР·РҙР°СӮСҢ С„РёР»СҢРј / РјРёРҪРё-С„РёР»СҢРј*\n"
        "РҹРҫРҙС…РҫРҙРёСӮ РҙР»СҸ СҖРҫР»РёРәР° 30вҖ“90 СҒРөРәСғРҪРҙ РёР»Рё СҒРөСҖРёРё СҒСҶРөРҪ: СҖРөРәР»Р°РјР°, РёСҒСӮРҫСҖРёСҸ, СӮСҖРөР№Р»РөСҖ, РҝСҖРҫРјРҫ, РәР»РёРҝ.\n\n"
        "РҡР°Рә Р·Р°РҝСғСҒСӮРёСӮСҢ:\n"
        "1) РҪР°РҝРёСҲРёСӮРө РёРҙРөСҺ, Р¶Р°РҪСҖ, РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ, СҒСӮРёР»СҢ Рё С„РҫСҖРјР°СӮ РәР°РҙСҖР°;\n"
        "2) СҸ СҒРҙРөР»Р°СҺ СҒСҶРөРҪР°СҖРёР№, СҒРҝРёСҒРҫРә СҒСҶРөРҪ, СҖР°СҒРәР°РҙСҖРҫРІРәСғ Рё РҝСҖРҫРјРҝСӮСӢ;\n"
        "3) СҒСҶРөРҪСӢ РіРөРҪРөСҖРёСҖСғСҺСӮСҒСҸ РәРҫСҖРҫСӮРәРёРјРё С„СҖР°РіРјРөРҪСӮР°РјРё СҮРөСҖРөР· Sora 2 РұРөР· Р»СҺРҙРөР№, Kling РёР»Рё Runway, Р·Р°СӮРөРј СҒРәР»РөРёРІР°СҺСӮСҒСҸ.\n\n"
        "Р РөРәРҫРјРөРҪРҙР°СҶРёСҸ: Runway вҖ” РҙР»СҸ РәРҫРҪСӮСҖРҫР»СҸ Рё РәР°СҮРөСҒСӮРІР°, Kling вҖ” РҙР»СҸ РҙРёРҪР°РјРёСҮРҪСӢС… СҒСҶРөРҪ, Sora 2 вҖ” РөСҒР»Рё РҙРҫСҒСӮСғРҝРөРҪ СҮРөСҖРөР· РІР°СҲ Comet-РәР°РҪР°Р» Рё РІ РәР°РҙСҖРө РҪРөСӮ Р»СҺРҙРөР№."
    )

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РҹРҫРҙРјРөРҪСҺ СҖРөР¶РёРјРҫРІ в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def _school_kb():
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("рҹ”Һ РһРұСҠСҸСҒРҪРөРҪРёРө", callback_data="school:explain")],
        [InlineKeyboardButton("рҹ§® Р—Р°РҙР°СҮРё", callback_data="school:tasks")],
        [InlineKeyboardButton("вңҚпёҸ РӯСҒСҒРө / СҖРөС„РөСҖР°СӮ", callback_data="school:essay")],
        [InlineKeyboardButton("рҹ“қ РӯРәР·Р°РјРөРҪ / РәРІРёР·", callback_data="school:quiz")],
    ])

def _work_kb():
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("рҹ“§ РҹРёСҒСҢРјРҫ / РҙРҫРәСғРјРөРҪСӮ", callback_data="work:doc")],
        [InlineKeyboardButton("рҹ“Ҡ РҗРҪР°Р»РёСӮРёРәР° / СҒРІРҫРҙРәР°", callback_data="work:report")],
        [InlineKeyboardButton("рҹ—Ӯ РҹР»Р°РҪ / ToDo", callback_data="work:plan")],
        [InlineKeyboardButton("рҹ’Ў РҳРҙРөРё / РұСҖРёС„", callback_data="work:idea")],
    ])

def _fun_quick_kb():
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("рҹӘ„ РһР¶РёРІРёСӮСҢ С„РҫСӮРҫ", callback_data="fun:revive")],
        [InlineKeyboardButton("рҹ—Ј Р“РҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ", callback_data="fun:avatar")],
        [InlineKeyboardButton("рҹҺө РӨРҫСӮРҫ вҶ’ РІРёРҙРөРҫРәР»РёРҝ", callback_data="fun:photoclip")],
        [InlineKeyboardButton("рҹҺӨ РҡР»РёРҝ СҒ РІРҫРәР°Р»РҫРј (1 СҮРөР»РҫРІРөРә)", callback_data="fun:vocalclip")],
        [InlineKeyboardButton("рҹҺ¬ Р’РёРҙРөРҫ РҝРҫ СӮРөРәСҒСӮСғ/РіРҫР»РҫСҒСғ", callback_data="fun:textvideo")],
        [InlineKeyboardButton("рҹӨі AI-СҒРөР»С„Рё СҒРҫ Р·РІРөР·РҙРҫР№", callback_data="fun:aiselfie")],
        [InlineKeyboardButton("рҹҺӯ Р—Р°РјРөРҪР° Р»РёСҶР° РҪР° С„РҫСӮРҫ", callback_data="fun:faceswap")],
        [InlineKeyboardButton("рҹ§ј РЈРҙР°Р»РёСӮСҢ С„РҫРҪ РҪР° С„РҫСӮРҫ", callback_data="fun:removebg")],
        [InlineKeyboardButton("рҹ–ј Р—Р°РјРөРҪРёСӮСҢ С„РҫРҪ РҪР° С„РҫСӮРҫ", callback_data="fun:replacebg")],
        [InlineKeyboardButton("рҹ“ұ Reels / Shorts", callback_data="fun:reels")],
        [InlineKeyboardButton("рҹҺһ РЎРҫР·РҙР°СӮСҢ РјРёРҪРё-С„РёР»СҢРј", callback_data="fun:film")],
        [InlineKeyboardButton("рҹҺ¬ Р Р°СҒРәР°РҙСҖРҫРІРәР° Reels", callback_data="fun:storyboard")],
        [InlineKeyboardButton("рҹ–Ң РЎРҫР·РҙР°СӮСҢ РёР·РҫРұСҖР°Р¶РөРҪРёРө", callback_data="fun:img")],
    ])

def _fun_kb():
    # РҫСҒСӮР°РІРёРј Рё СҒСӮР°СҖРҫРө РҝРҫРҙРјРөРҪСҺ вҖ” РҪРө РёСҒРҝРҫР»СҢР·СғРөСӮСҒСҸ СҒРөР№СҮР°СҒ
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("рҹ–ј РӨРҫСӮРҫ-РјР°СҒСӮРөСҖСҒРәР°СҸ", callback_data="fun:photo"),
         InlineKeyboardButton("рҹҺ¬ Р’РёРҙРөРҫ-РёРҙРөРё",      callback_data="fun:video")],
        [InlineKeyboardButton("рҹҺІ РҡРІРёР·СӢ/РёРіСҖСӢ",      callback_data="fun:quiz"),
         InlineKeyboardButton("рҹҳҶ РңРөРјСӢ/СҲСғСӮРәРё",      callback_data="fun:meme")],
    ])


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РҡРҫРјР°РҪРҙСӢ/РәРҪРҫРҝРәРё СҖРөР¶РёРјРҫРІ в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
async def cmd_mode_school(update: Update, context: ContextTypes.DEFAULT_TYPE):
    _clear_transient_flows(context)
    _set_mode_clean(update.effective_user.id, "РЈСҮС‘РұР°", "")
    # РҝРҫРәР°Р·СӢРІР°РөРј РқРһР’РһР• РҝРҫРҙРјРөРҪСҺ В«РЈСҮС‘РұР°В»
    await _send_mode_menu(update, context, "study")

async def cmd_mode_work(update: Update, context: ContextTypes.DEFAULT_TYPE):
    _clear_transient_flows(context)
    _set_mode_clean(update.effective_user.id, "Р Р°РұРҫСӮР°/Р‘РёР·РҪРөСҒ", "")
    # РҝРҫРәР°Р·СӢРІР°РөРј РқРһР’РһР• РҝРҫРҙРјРөРҪСҺ В«Р Р°РұРҫСӮР°В»
    await _send_mode_menu(update, context, "work")

async def cmd_mode_fun(update: Update, context: ContextTypes.DEFAULT_TYPE):
    _clear_transient_flows(context)
    _set_mode_clean(update.effective_user.id, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "")
    await update.effective_message.reply_text(
        "рҹ”Ҙ Р Р°Р·РІР»РөСҮРөРҪРёСҸ вҖ” РұСӢСҒСӮСҖСӢРө РҙРөР№СҒСӮРІРёСҸ: РІРёРҙРөРҫ, С„РҫСӮРҫ, С„РҫРҪ Рё Р·Р°РјРөРҪР° Р»РёСҶР°.",
        reply_markup=_fun_quick_kb()
    )

async def cmd_mode_medicine(update: Update, context: ContextTypes.DEFAULT_TYPE):
    _set_medical_waiting(update, context, "")
    await update.effective_message.reply_text(_medical_menu_text(), reply_markup=medicine_kb())


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РҡРҫР»Р»РұСҚРәРё РҝРҫРҙСҖРөР¶РёРјРҫРІ в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
async def on_cb_mode(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    data = (q.data or "")
    try:
        if any(data.startswith(p) for p in ("school:", "work:", "fun:")):
            # РұР°Р·РҫРІСӢР№ СӮСҖРөРәРёРҪРі СҒСӮР°СҖСӢС… РІРөСӮРҫРә (photo/video/quiz/meme)
            if data in ("fun:revive","fun:clip","fun:img","fun:storyboard"):
                # СҚСӮРё РҫРұСҖР°РұР°СӮСӢРІР°СҺСӮСҒСҸ РҫСӮРҙРөР»СҢРҪСӢРј С…РөРҪРҙР»РөСҖРҫРј on_cb_fun
                return
            _, track = data.split(":", 1)
            _mode_track_set(update.effective_user.id, track)
            mode = _mode_get(update.effective_user.id)
            await q.message.reply_text(f"{mode} вҶ’ {track}. РқР°РҝРёСҲРёСӮРө Р·Р°РҙР°РҪРёРө/СӮРөРјСғ вҖ” СҒРҙРөР»Р°СҺ.")
            return
    finally:
        with contextlib.suppress(Exception):
            await q.answer()

# РұСӢСҒСӮСҖСӢРө РҙРөР№СҒСӮРІРёСҸ В«Р Р°Р·РІР»РөСҮРөРҪРёСҸВ»
async def on_cb_fun(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    await q.answer()
    data = q.data or ""
    if data == "fun:img":
        return await q.message.reply_text("РҹСҖРёСҲР»Рё РҝСҖРҫРјРҝСӮ РёР»Рё РёСҒРҝРҫР»СҢР·СғР№ РәРҫРјР°РҪРҙСғ /img <РҫРҝРёСҒР°РҪРёРө> вҖ” СҒРіРөРҪРөСҖРёСҖСғСҺ РёР·РҫРұСҖР°Р¶РөРҪРёРө.")
    if data == "fun:revive":
        return await q.message.reply_text("Р—Р°РіСҖСғР·Рё С„РҫСӮРҫ (РәР°Рә РәР°СҖСӮРёРҪРәСғ) Рё РҪР°РҝРёСҲРё, СҮСӮРҫ РҫР¶РёРІРёСӮСҢ/РәР°Рә РҙРІРёРіР°СӮСҢСҒСҸ. РЎРҙРөР»Р°СҺ Р°РҪРёРјР°СҶРёСҺ.")
    if data == "fun:avatar":
        context.user_data["awaiting_avatar_photo"] = True
        return await q.message.reply_text("Р—Р°РіСҖСғР·Рё РҝРҫСҖСӮСҖРөСӮ СҮРөР»РҫРІРөРәР°, Р·Р°СӮРөРј РҪР°Р¶РјРё рҹ—Ј Р“РҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ. РҹРҫСҒР»Рө СҚСӮРҫРіРҫ РҝСҖРёСҲР»Рё СӮРөРәСҒСӮ, voice РёР»Рё Р°СғРҙРёРҫ РҙР»СҸ СҖРөСҮРё.")
    if data == "fun:photoclip":
        context.user_data["awaiting_photo_clip_photo"] = True
        return await q.message.reply_text("Р—Р°РіСҖСғР·Рё С„РҫСӮРҫ СҮРөР»РҫРІРөРәР°, Р·Р°СӮРөРј РҪР°Р¶РјРё рҹҺө РӨРҫСӮРҫ вҶ’ РІРёРҙРөРҫРәР»РёРҝ Рё РҫРҝРёСҲРё СҒСӮРёР»СҢ: РјСғР·СӢРәР°, РҙРІРёР¶РөРҪРёРө, РҪР°СҒСӮСҖРҫРөРҪРёРө, 5/10 СҒРөРә, 9:16 РёР»Рё 16:9.")
    if data == "fun:aiselfie":
        context.user_data["awaiting_ai_selfie_photo"] = True
        return await q.message.reply_text("рҹӨі Р—Р°РіСҖСғР·Рё СҒРІРҫС‘ СҒРөР»С„Рё, Р·Р°СӮРөРј РҪР°РҝРёСҲРё, СҒ РәРөРј/РіРҙРө СҒРҙРөР»Р°СӮСҢ AI-С„РҫСӮРҫ: Р·РҪР°РјРөРҪРёСӮРҫСҒСӮСҢ, РҝРөСҖСҒРҫРҪР°Р¶, РҝСҖРөРјСҢРөСҖР°, СҖРөРәР»Р°РјР°, travel/luxury.")
    if data == "fun:clip":
        return await q.message.reply_text("РҹСҖРёСҲР»Рё СӮРөРәСҒСӮ/РіРҫР»РҫСҒ Рё С„РҫСҖРјР°СӮ (Reels/Shorts), РјСғР·СӢРәСғ/СҒСӮРёР»СҢ вҖ” СҒРҫРұРөСҖСғ РәР»РёРҝ. Р”Р»СҸ РіРөРҪРөСҖР°СҶРёРё РІРёРҙРөРҫ РҙРҫСҒСӮСғРҝРҪСӢ Sora 2 РұРөР· Р»СҺРҙРөР№, Kling Рё Runway.")
    if data == "fun:storyboard":
        return await q.message.reply_text("РҹСҖРёСҲР»Рё С„РҫСӮРҫ РёР»Рё РҫРҝРёСҲРё РёРҙРөСҺ СҖРҫР»РёРәР° вҖ” РІРөСҖРҪСғ СҖР°СҒРәР°РҙСҖРҫРІРәСғ РҝРҫРҙ Reels СҒ СӮР°Р№Рј-РәРҫРҙР°РјРё.")

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РЎСӮР°СҖСӮ / Р”РІРёР¶РәРё / РҹРҫРјРҫСүСҢ в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
async def cmd_start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    with contextlib.suppress(Exception):
        _chat_ensure_active(update.effective_user.id, update.effective_chat.id)
    welcome_url = kv_get("welcome_url", BANNER_URL)
    if welcome_url:
        with contextlib.suppress(Exception):
            await update.effective_message.reply_photo(welcome_url)
    await update.effective_message.reply_text(START_TEXT, reply_markup=main_kb, disable_web_page_preview=True)

async def cmd_engines(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.effective_message.reply_text("рҹ§  Р’СӢРұРөСҖРёСӮРө РҪРөР№СҖРҫСҒРөСӮСҢ РёР»Рё РҝСҖРҫРІР°Р№РҙРөСҖР°:", reply_markup=engines_kb())

async def cmd_subs_help(update: Update, context: ContextTypes.DEFAULT_TYPE):
    kb = InlineKeyboardMarkup([
        [InlineKeyboardButton("вӯҗ РһСӮРәСҖСӢСӮСҢ СӮР°СҖРёС„СӢ", web_app=WebAppInfo(url=TARIFF_URL))],
        [InlineKeyboardButton("рҹҡҖ PRO РҪР° РјРөСҒСҸСҶ", callback_data="buyinv:pro:1")],
    ])
    await update.effective_message.reply_text("вӯҗ РўР°СҖРёС„СӢ Рё РҝРҫРјРҫСүСҢ.\n\n" + HELP_TEXT, reply_markup=kb, disable_web_page_preview=True)

async def cmd_help(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.effective_message.reply_text(HELP_TEXT, disable_web_page_preview=True)

async def cmd_examples(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.effective_message.reply_text(EXAMPLES_TEXT, disable_web_page_preview=True)

async def cmd_version(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.effective_message.reply_text(
        f"вң… РҡРҫРҙ Р·Р°РҝСғСүРөРҪ: {PATCH_VERSION}\n"
        f"РӨР°Р№Р» РҙРҫР»Р¶РөРҪ РұСӢСӮСҢ РёРјРөРҪРҪРҫ main.py РҪР° Render. Start Command: python -u main.py"
    )


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ Р”РёР°РіРҪРҫСҒСӮРёРәР°/Р»РёРјРёСӮСӢ в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
async def cmd_diag_limits(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    tier = get_subscription_tier(user_id)
    lim = _limits_for(user_id)
    row = _usage_row(user_id, _today_ymd())
    lines = [
        f"рҹ‘Ө РўР°СҖРёС„: {tier}",
        f"вҖў РўРөРәСҒСӮСӢ СҒРөРіРҫРҙРҪСҸ: {row['text_count']} / {lim['text_per_day']}",
        f"вҖў Р‘РөСҒРҝР»Р°СӮРҪСӢРө РәР°СҖСӮРёРҪРәРё: {row['free_img_gen_count']} / {FREE_IMAGE_GENERATIONS_PER_DAY}",
        f"вҖў Р‘РөСҒРҝР»Р°СӮРҪСӢРө РҫРұСҖР°РұРҫСӮРәРё С„РҫСӮРҫ: {row['free_img_proc_count']} / {FREE_IMAGE_PROCESSINGS_PER_DAY}",
        f"вҖў РЎРөРұРөСҒСӮРҫРёРјРҫСҒСӮСҢ РІРёРҙРөРҫ СҒРөРіРҫРҙРҪСҸ: {_credits_fmt_from_usd(row['runway_usd'])} (РІРҪСғСӮСҖРөРҪРҪСҸСҸ РҙРёР°РіРҪРҫСҒСӮРёРәР°)",
        f"вҖў РЎРөРұРөСҒСӮРҫРёРјРҫСҒСӮСҢ РёР·РҫРұСҖР°Р¶РөРҪРёР№ СҒРөРіРҫРҙРҪСҸ: {_credits_fmt_from_usd(row['img_usd'])} (РІРҪСғСӮСҖРөРҪРҪСҸСҸ РҙРёР°РіРҪРҫСҒСӮРёРәР°)",
    ]
    await update.effective_message.reply_text("\n".join(lines))


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ Capability Q&A в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
_CAP_PDF   = re.compile(r"(pdf|РҙРҫРәСғРјРөРҪСӮ(СӢ)?|С„Р°Р№Р»(СӢ)?)", re.I)
_CAP_EBOOK = re.compile(r"(ebook|e-?book|СҚР»РөРәСӮСҖРҫРҪРҪ(Р°СҸ|СӢРө)\s+РәРҪРёРі|epub|fb2|docx|txt|mobi|azw)", re.I)
_CAP_AUDIO = re.compile(r"(Р°СғРҙРёРҫ ?РәРҪРёРі|audiobook|audio ?book|mp3|m4a|wav|ogg|webm|voice)", re.I)
_CAP_IMAGE = re.compile(r"(РёР·РҫРұСҖР°Р¶РөРҪ|РәР°СҖСӮРёРҪРә|С„РҫСӮРҫ|image|picture|img)", re.I)
_CAP_VIDEO = re.compile(r"(РІРёРҙРөРҫ|СҖРҫР»РёРә|shorts?|reels?|clip)", re.I)

def _is_photo_revival_question(text: str) -> bool:
    """Р–С‘СҒСӮРәРёР№ РҝРөСҖРөС…РІР°СӮ РІРҫРҝСҖРҫСҒРҫРІ/С„СҖР°Р· РҝСҖРҫ РІРҫР·РјРҫР¶РҪРҫСҒСӮСҢ РҫР¶РёРІРёСӮСҢ С„РҫСӮРҫ.
    РқРө РҫСӮРҙР°С‘Рј СӮР°РәРёРө С„СҖР°Р·СӢ РІ GPT, РҝРҫСӮРҫРјСғ СҮСӮРҫ РјРҫРҙРөР»СҢ РјРҫР¶РөСӮ РҫСӮРІРөСӮРёСӮСҢ РҫРұСүРёРј РҫСӮРәР°Р·РҫРј.
    """
    tl = (text or "").strip().lower()
    if not tl:
        return False
    has_photo = bool(re.search(r"(С„РҫСӮРҫ|С„РҫСӮРҫРіСҖР°С„|РәР°СҖСӮРёРҪРә|РёР·РҫРұСҖР°Р¶РөРҪ|image|picture|photo)", tl, re.I))
    has_revival = bool(re.search(r"(РҫР¶РёРІ|Р°РҪРёРјРёСҖ|РҙРІРёР¶РөРҪ|image\s*to\s*video|i2v|revive|animate)", tl, re.I))
    has_ability = bool(re.search(r"(РјРҫР¶(РөСҲСҢ|РөСӮРө|РҪРҫ)|СғРјРө(РөСҲСҢ|РөСӮРө)|СҒРҝРҫСҒРҫРұРөРҪ|РҝРҫРҙРҙРөСҖР¶РёРІР°РөСҲСҢ|РҙРөР»Р°РөСҲСҢ|РҝРҫР»СғСҮРёСӮСҒСҸ|РјРҫР¶РөСӮ\s+Р»Рё)", tl, re.I))
    # РӣРҫРІРёРј Рё РҝСҖСҸРјРҫР№ РІРҫРҝСҖРҫСҒ В«РјРҫР¶РөСҲСҢ РҫР¶РёРІРёСӮСҢ С„РҫСӮРҫ?В», Рё С„СҖР°Р·СӢ РІРёРҙР° В«РҫР¶РёРІР»РөРҪРёРө С„РҫСӮРҫ РІРҫР·РјРҫР¶РҪРҫ?В»
    return has_photo and has_revival and (has_ability or "?" in tl)

def _is_photo_revival_intent(text: str) -> bool:
    """РӣРҫРІРёСӮ РҪРө СӮРҫР»СҢРәРҫ РІРҫРҝСҖРҫСҒ В«РјРҫР¶РөСҲСҢ РҫР¶РёРІРёСӮСҢ С„РҫСӮРҫ?В», РҪРҫ Рё РәРҫРјР°РҪРҙСғ/РҪР°РјРөСҖРөРҪРёРө В«РҫР¶РёРІРё СҚСӮРҫ С„РҫСӮРҫВ»."""
    tl = (text or "").strip().lower().replace("С‘", "Рө")
    if not tl:
        return False
    has_photo = bool(re.search(r"(С„РҫСӮРҫ|С„РҫСӮРҫРіСҖР°С„|РәР°СҖСӮРёРҪРә|РёР·РҫРұСҖР°Р¶РөРҪ|image|picture|photo)", tl, re.I))
    has_revival = bool(re.search(r"(РҫР¶РёРІ|Р°РҪРёРјРёСҖ|РҙРІРёР¶РөРҪ|image\s*to\s*video|i2v|revive|animate|СҒРҙРөР»Р°Р№\s+РІРёРҙРөРҫ)", tl, re.I))
    return has_photo and has_revival


def _revival_engine_from_text(text: str, default: str = "runway") -> str:
    tl = (text or "").lower().replace("С‘", "Рө")
    if "luma" in tl or "Р»СғРјР°" in tl:
        return "luma"
    if "sora" in tl or "СҒРҫСҖР°" in tl:
        return "sora"
    if "kling" in tl or "РәР»РёРҪРі" in tl:
        return "kling"
    if "runway" in tl or "СҖР°РҪРІРөР№" in tl or "СҖР°РҪРІСҚР№" in tl:
        return "runway"
    return default


def _clean_revival_prompt(text: str) -> str:
    cleaned = re.sub(
        r"\b(РҫР¶РёРІРё|РҫР¶РёРІРёСӮСҢ|Р°РҪРёРјРёСҖСғР№|Р°РҪРёРјРёСҖРҫРІР°СӮСҢ|СҒРҙРөР»Р°Р№\s+РІРёРҙРөРҫ|revive|animate|image\s*to\s*video|i2v|runway|luma|sora|kling|Р»СғРјР°|СҒРҫСҖР°|РәР»РёРҪРі|СҖР°РҪРІРөР№|СҖР°РҪРІСҚР№)\b",
        "",
        text or "",
        flags=re.I,
    )
    cleaned = re.sub(
        r"\b(С„РҫСӮРҫ|С„РҫСӮРҫРіСҖР°С„Рё(?:СҸ|СҺ|Рё|РөР№)?|РәР°СҖСӮРёРҪРә(?:Р°|Сғ|Рё|РҫР№)?|РёР·РҫРұСҖР°Р¶РөРҪРё(?:Рө|СҸ|СҺ)?|photo|image|picture)\b",
        "",
        cleaned,
        flags=re.I,
    ).strip(" ,.:-вҖ”")
    # РқРө РҫСӮРҝСҖР°РІР»СҸРөРј РІ РҙРІРёР¶РҫРә РјСғСҒРҫСҖРҪСӢР№ РҝСҖРҫРјРҝСӮ РІСҖРҫРҙРө В«СҚСӮРҫВ» / В«РҝРҫР¶Р°Р»СғР№СҒСӮР°В».
    if len(cleaned) < 4:
        return ""
    return cleaned


def _photo_revival_capability_text() -> str:
    return (
        "Р”Р°, РјРҫРіСғ РҫР¶РёРІРёСӮСҢ С„РҫСӮРҫРіСҖР°С„РёСҺ Рё СҒРҙРөР»Р°СӮСҢ РёР· РҪРөС‘ РәРҫСҖРҫСӮРәРҫРө РІРёРҙРөРҫ.\n\n"
        "РҡР°Рә Р·Р°РҝСғСҒСӮРёСӮСҢ:\n"
        "1) Р·Р°РіСҖСғР·РёСӮРө С„РҫСӮРҫ;\n"
        "2) РҪР°Р¶РјРёСӮРө РәРҪРҫРҝРәСғ вңЁ РһР¶РёРІРёСӮСҢ: Runway / Kling / Sora 2 РұРөР· Р»СҺРҙРөР№;\n"
        "3) Р»РёРұРҫ РҫСӮРҝСҖР°РІСҢСӮРө С„РҫСӮРҫ СҒ РҝРҫРҙРҝРёСҒСҢСҺ: В«РҫР¶РёРІРё С„РҫСӮРҫ: Р»С‘РіРәР°СҸ СғР»СӢРұРәР°, РҙРІРёР¶РөРҪРёРө РәР°РјРөСҖСӢ, 5 СҒРөРәСғРҪРҙ, 9:16В».\n\n"
        "Р’Р°Р¶РҪРҫ: Sora 2 СҮР°СҒСӮРҫ РұР»РҫРәРёСҖСғРөСӮ Р·Р°РіСҖСғР¶РөРҪРҪСӢРө С„РҫСӮРҫ СҒ Р»СҺРҙСҢРјРё РёР·-Р·Р° РјРҫРҙРөСҖР°СҶРёРё. "
        "Р”Р»СҸ РҝРҫСҖСӮСҖРөСӮРҫРІ Рё С„РҫСӮРҫ Р»СҺРҙРөР№ Р»СғСҮСҲРө РІСӢРұРёСҖР°СӮСҢ Runway РёР»Рё Kling. "
        "Sora 2 РҫСҒСӮР°РІР»РөРҪР° РҙР»СҸ РёР·РҫРұСҖР°Р¶РөРҪРёР№ РұРөР· Р»СҺРҙРөР№: РҝСҖРөРҙРјРөСӮСӢ, Р¶РёРІРҫСӮРҪСӢРө, Р·РҙР°РҪРёСҸ, РҝРөР№Р·Р°Р¶Рё, РёРҪСӮРөСҖСҢРөСҖ.\n\n"
        "Р•СҒР»Рё РҙРІРёР¶РҫРә РІРөСҖРҪС‘СӮ РҫСҲРёРұРәСғ, СҸ РҝРҫРәР°Р¶Сғ СӮРөС…РҪРёСҮРөСҒРәСғСҺ РҝСҖРёСҮРёРҪСғ: РәР»СҺСҮ, Р»РёРјРёСӮ, РәСҖРөРҙРёСӮСӢ, РјРҫРҙРөСҖР°СҶРёСҸ РёР»Рё С„РҫСҖРјР°СӮ Р·Р°РҝСҖРҫСҒР°."
    )

MEDICAL_DISCLAIMER = (
    "\n\nвҡ пёҸ Р’Р°Р¶РҪРҫ: СҚСӮРҫ СҒРҝСҖР°РІРҫСҮРҪСӢР№ СҖР°Р·РұРҫСҖ Рё РҝРҫРҙРіРҫСӮРҫРІРәР° РІРҫРҝСҖРҫСҒРҫРІ Рә РІСҖР°СҮСғ, "
    "Р° РҪРө РҫС„РёСҶРёР°Р»СҢРҪСӢР№ РҙРёР°РіРҪРҫР·, РҪРө РјРөРҙРёСҶРёРҪСҒРәРҫРө Р·Р°РәР»СҺСҮРөРҪРёРө Рё РҪРө Р·Р°РјРөРҪР° РҫСҮРҪРҫР№ РәРҫРҪСҒСғР»СҢСӮР°СҶРёРё/РҫРұСҒР»РөРҙРҫРІР°РҪРёСҸ."
)

_MEDICAL_TERMS_RE = re.compile(
    # Р’РҗР–РқРһ: РҪРө РёСҒРҝРҫР»СҢР·СғРөРј РіРҫР»СӢР№ РәРҫСҖРөРҪСҢ В«Р°РҪР°Р»РёР·В»,
    # РёРҪР°СҮРө С„СҖР°Р·СӢ В«Р°РҪР°Р»РёР· PDF/СҚР»РөРәСӮСҖРҫРҪРҪСӢС… РәРҪРёРі/С„РҫСӮРҫВ» РҫСҲРёРұРҫСҮРҪРҫ СғС…РҫРҙСҸСӮ РІ РјРөРҙРёСҶРёРҪСғ.
    r"(РјРөРҙРёСҶРёРҪ|РјРөРҙРәР°СҖСӮ|РёСҒСӮРҫСҖРё[СҸРё]\s+РұРҫР»РөР·РҪ|РІСӢРҝРёСҒРә|Р°РҪР°РјРҪРөР·|"
    r"СҖРөР·СғР»СҢСӮР°СӮ(?:СӢ)?\s+(?:Р»Р°РұРҫСҖР°СӮРҫСҖРҪ(?:СӢС…|РҫРіРҫ)?\s+)?Р°РҪР°Р»РёР·(?:РҫРІ|СӢ)?|"
    r"Р»Р°РұРҫСҖР°СӮРҫСҖРҪ(?:СӢРө|СӢС…|РҫРіРҫ)?\s+Р°РҪР°Р»РёР·(?:СӢ|РҫРІ)?|"
    r"Р°РҪР°Р»РёР·(?:СӢ|РҫРІ)?\s+(?:РәСҖРҫРІРё|РјРҫСҮРё|РәР°Р»[Р°-СҸ]*|РіРҫСҖРјРҫРҪ(?:СӢ|РҫРІ)?|РұРёРҫС…РёРјРё[СҸРё]|С„РөСҖСҖРёСӮРёРҪ|РіР»СҺРәРҫР·[Р°-СҸ]*|С…РҫР»РөСҒСӮРөСҖРёРҪ)|"
    r"Р·Р°РәР»СҺСҮРөРҪРё|РІСҖР°СҮРөРұРҪ|РҙРёР°РіРҪРҫР·|СҚРҝРёРәСҖРёР·|РҪР°Р·РҪР°СҮРөРҪРё|СҖРөСҶРөРҝСӮ|"
    r"СҖРөРҪСӮРіРөРҪ|СғР·Рё|РјСҖСӮ|РәСӮ|mri|ct|СӮРҫРјРҫРіСҖР°С„|"
    r"(?:РјРөРҙРёСҶРёРҪСҒРә(?:РёР№|РҫРіРҫ|РҫРјСғ|РёРј|РҫРј|Р°СҸ|СғСҺ|РҫРө|РёРө|РёС…)\s+)?СҒРҪРёРјРҫРә|"
    r"РҝРөСҖРөР»РҫРј|СӮСҖР°РІРј|РұРҫР»СҢ|РұРҫР»РёСӮ|РҫСӮ[РөС‘]Рә|СҒРёРјРҝСӮРҫРј|Р»РөСҮРөРҪРёРө|РІСҖР°СҮ|РҝР°СҶРёРөРҪСӮ|"
    r"blood\s*test|medical|x[-\s]?ray|ultrasound|scan)",
    re.I,
)
_MEDICAL_ABILITY_RE = re.compile(
    r"(РјРҫР¶(РөСҲСҢ|РөСӮРө|РҪРҫ)|СғРјРө(РөСҲСҢ|РөСӮРө)|СҒРҝРҫСҒРҫРұРөРҪ|РҝРҫРҙРҙРөСҖР¶РёРІР°РөСҲСҢ|Р°РҪР°Р»РёР·РёСҖСғРөСҲСҢ|"
    r"СҖР°РұРҫСӮР°РөСҲСҢ|СҖР°Р·РұРёСҖР°РөСҲСҢ|РјРҫР¶РөСӮ\s+Р»Рё|РјРҫР¶РҪРҫ\s+Р»Рё|РҝРҫР»СғСҮРёСӮСҒСҸ)",
    re.I,
)

def _is_medical_capability_question(text: str) -> bool:
    tl = (text or "").strip().lower()
    if not tl:
        return False
    return bool(_MEDICAL_TERMS_RE.search(tl) and (_MEDICAL_ABILITY_RE.search(tl) or "?" in tl))

def _medical_capability_text() -> str:
    return (
        "Р”Р°, РјРҫРіСғ РҝРҫРјРҫСҮСҢ СҒ РјРөРҙРёСҶРёРҪСҒРәРёРјРё С„Р°Р№Р»Р°РјРё Рё РҙРҫРәСғРјРөРҪСӮР°РјРё.\n\n"
        "Р§СӮРҫ РјРҫРіСғ СҖР°Р·РҫРұСҖР°СӮСҢ:\n"
        "вҖў РІСӢРҝРёСҒРәСғ РёР· РјРөРҙРёСҶРёРҪСҒРәРҫР№ РәР°СҖСӮСӢ Рё Р°РҪР°РјРҪРөР·;\n"
        "вҖў РІСҖР°СҮРөРұРҪРҫРө Р·Р°РәР»СҺСҮРөРҪРёРө/СҚРҝРёРәСҖРёР·/РҪР°Р·РҪР°СҮРөРҪРёСҸ;\n"
        "вҖў СҖРөР·СғР»СҢСӮР°СӮСӢ Р»Р°РұРҫСҖР°СӮРҫСҖРҪСӢС… Р°РҪР°Р»РёР·РҫРІ;\n"
        "вҖў СҒРҪРёРјРҫРә, С„РҫСӮРҫ РҙРҫРәСғРјРөРҪСӮР°, РЈР—Рҳ/СҖРөРҪСӮРіРөРҪ;\n"
        "вҖў РҫРҝРёСҒР°РҪРёРө РңР Рў Рё РҡРў, Р° СӮР°РәР¶Рө РёР·РҫРұСҖР°Р¶РөРҪРёРө, РөСҒР»Рё РөРіРҫ РјРҫР¶РҪРҫ РҝСҖРҫСҮРёСӮР°СӮСҢ РҝРҫ С„РҫСӮРҫ.\n\n"
        "Р§СӮРҫ РІСӢ РҝРҫР»СғСҮРёСӮРө: РҝРҫРҪСҸСӮРҪСғСҺ РІСӢР¶РёРјРәСғ, СҖР°СҒСҲРёС„СҖРҫРІРәСғ СӮРөСҖРјРёРҪРҫРІ, РәР»СҺСҮРөРІСӢРө РҝРҫРәР°Р·Р°СӮРөР»Рё, "
        "РәСҖР°СҒРҪСӢРө С„Р»Р°РіРё, СҒРҝРёСҒРҫРә РІРҫРҝСҖРҫСҒРҫРІ РІСҖР°СҮСғ Рё РҝР»Р°РҪ, СҮСӮРҫ СғСӮРҫСҮРҪРёСӮСҢ РҙР°Р»СҢСҲРө."
        + MEDICAL_DISCLAIMER
    )

def _medical_menu_text(track: str = "") -> str:
    title = {
        "med_extract": "РІСӢРҝРёСҒРәСғ РёР· РјРөРҙРёСҶРёРҪСҒРәРҫР№ РәР°СҖСӮСӢ",
        "med_scan": "СҒРҪРёРјРҫРә/С„РҫСӮРҫ РјРөРҙРёСҶРёРҪСҒРәРҫРіРҫ РҙРҫРәСғРјРөРҪСӮР°",
        "med_conclusion": "РІСҖР°СҮРөРұРҪРҫРө Р·Р°РәР»СҺСҮРөРҪРёРө",
        "med_mri": "РңР Рў",
        "med_ct": "РҡРў",
        "med_labs": "СҖРөР·СғР»СҢСӮР°СӮСӢ Р°РҪР°Р»РёР·РҫРІ",
    }.get(track, "РјРөРҙРёСҶРёРҪСҒРәРёР№ РҙРҫРәСғРјРөРҪСӮ РёР»Рё СҒРҪРёРјРҫРә")
    return (
        f"рҹ©ә РңРөРҙРёСҶРёРҪР° вҖ” РіРҫСӮРҫРІ СҖР°Р·РҫРұСҖР°СӮСҢ {title}.\n\n"
        "РҡР°Рә Р·Р°РіСҖСғР·РёСӮСҢ:\n"
        "1) РҪР°Р¶РјРёСӮРө РҪСғР¶РҪСғСҺ РәРҪРҫРҝРәСғ РҪРёР¶Рө;\n"
        "2) РҫСӮРҝСҖР°РІСҢСӮРө PDF/DOCX/TXT РёР»Рё С„РҫСӮРҫ/СҒРәСҖРёРҪ/РёР·РҫРұСҖР°Р¶РөРҪРёРө;\n"
        "3) РІ РҝРҫРҙРҝРёСҒРё РјРҫР¶РҪРҫ РҪР°РҝРёСҒР°СӮСҢ СҶРөР»СҢ: В«РәРҫСҖРҫСӮРәРҫВ», В«РҝРҫРҙСҖРҫРұРҪРҫВ», В«СҮСӮРҫ СҒРҝСҖРҫСҒРёСӮСҢ Сғ РІСҖР°СҮР°В», "
        "В«РҫРұСҠСҸСҒРҪРё РҝСҖРҫСҒСӮСӢРјРё СҒР»РҫРІР°РјРёВ».\n\n"
        "Р РөР·СғР»СҢСӮР°СӮ: РәСҖР°СӮРәРҫРө СҖРөР·СҺРјРө, СҖР°СҒСҲРёС„СҖРҫРІРәР° СӮРөСҖРјРёРҪРҫРІ, РІР°Р¶РҪСӢРө СҶРёС„СҖСӢ/РҪР°С…РҫРҙРәРё, "
        "РІРҫР·РјРҫР¶РҪСӢРө СҖРёСҒРәРё, РІРҫРҝСҖРҫСҒСӢ РІСҖР°СҮСғ Рё СҮСӮРҫ РҝСҖРҫРІРөСҖРёСӮСҢ РҙРҫРҝРҫР»РҪРёСӮРөР»СҢРҪРҫ."
        + MEDICAL_DISCLAIMER
    )

def medicine_kb() -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("рҹ“„ Р Р°Р·РұРҫСҖ РІСӢРҝРёСҒРәРё", callback_data="act:med:extract")],
        [InlineKeyboardButton("рҹ–ј Р Р°Р·РұРҫСҖ СҒРҪРёРјРәР°", callback_data="act:med:scan")],
        [InlineKeyboardButton("рҹ§ҫ Р—Р°РәР»СҺСҮРөРҪРёРө РІСҖР°СҮР°", callback_data="act:med:conclusion")],
        [InlineKeyboardButton("рҹ§І РңР Рў", callback_data="act:med:mri"),
         InlineKeyboardButton("рҹ§  РҡРў", callback_data="act:med:ct")],
        [InlineKeyboardButton("рҹ§Ә РҗРҪР°Р»РёР·СӢ", callback_data="act:med:labs")],
        [InlineKeyboardButton("рҹ“қ РңРөРҙ. РІРҫРҝСҖРҫСҒ", callback_data="act:med:free")],
        [InlineKeyboardButton("в¬…пёҸ РқР°Р·Р°Рҙ", callback_data="mode:root")],
    ])

def _medical_track_title(track: str) -> str:
    return {
        "med_extract": "РІСӢРҝРёСҒРәР° / Р°РҪР°РјРҪРөР·",
        "med_scan": "РјРөРҙРёСҶРёРҪСҒРәРёР№ СҒРҪРёРјРҫРә / С„РҫСӮРҫ",
        "med_conclusion": "РІСҖР°СҮРөРұРҪРҫРө Р·Р°РәР»СҺСҮРөРҪРёРө",
        "med_mri": "РңР Рў",
        "med_ct": "РҡРў",
        "med_labs": "Р»Р°РұРҫСҖР°СӮРҫСҖРҪСӢРө Р°РҪР°Р»РёР·СӢ",
        "med_free": "РјРөРҙРёСҶРёРҪСҒРәРёР№ РІРҫРҝСҖРҫСҒ",
    }.get(track or "", "РјРөРҙРёСҶРёРҪСҒРәРёР№ РјР°СӮРөСҖРёР°Р»")

def _is_medical_context(user_id: int, caption_or_text: str = "", filename: str = "") -> bool:
    # РЎРҫРІРјРөСҒСӮРёРјРҫСҒСӮСҢ РҙР»СҸ СҒСӮР°СҖСӢС… РјРөСҒСӮ РІСӢР·РҫРІР° РұРөР· context.
    # Р’РҗР–РқРһ: РҫРҙРёРҪ СӮРҫР»СҢРәРҫ mode == "РңРөРҙРёСҶРёРҪР°" РұРҫР»СҢСҲРө РҪРө СҒСҮРёСӮР°РөСӮСҒСҸ РҝСҖРёСҮРёРҪРҫР№
    # РҫСӮРҝСҖР°РІР»СҸСӮСҢ Р»СҺРұРҫРө С„РҫСӮРҫ/РҙРҫРәСғРјРөРҪСӮ РІ РјРөРҙРёСҶРёРҪСҒРәСғСҺ РІРөСӮРәСғ.
    track = ""
    with contextlib.suppress(Exception):
        track = _mode_track_get(user_id)
    combined = f"{caption_or_text or ''} {filename or ''}"
    return bool(_MEDICAL_TERMS_RE.search(combined))

def _medical_doc_prompt(text: str, track: str = "", goal: str | None = None) -> str:
    task = _medical_track_title(track)
    extra = f"\nРҰРөР»СҢ РҝРҫР»СҢР·РҫРІР°СӮРөР»СҸ: {goal}" if goal else ""
    return (
        f"РўСӢ Р°РҪР°Р»РёР·РёСҖСғРөСҲСҢ РјРөРҙРёСҶРёРҪСҒРәРёР№ РјР°СӮРөСҖРёР°Р»: {task}. "
        "РқРө СҒСӮР°РІСҢ РҙРёР°РіРҪРҫР· Рё РҪРө РҪР°Р·РҪР°СҮР°Р№ Р»РөСҮРөРҪРёРө. Р”Р°Р№ СҒРҝСҖР°РІРҫСҮРҪСӢР№ СҖР°Р·РұРҫСҖ РҪР° СҖСғСҒСҒРәРҫРј. "
        "РЎСӮСҖСғРәСӮСғСҖР° РҫСӮРІРөСӮР°:\n"
        "1) Р§СӮРҫ СҚСӮРҫ Р·Р° РҙРҫРәСғРјРөРҪСӮ/РёСҒСҒР»РөРҙРҫРІР°РҪРёРө.\n"
        "2) РҡСҖР°СӮРәРҫРө СҖРөР·СҺРјРө РҝСҖРҫСҒСӮСӢРјРё СҒР»РҫРІР°РјРё.\n"
        "3) РҡР»СҺСҮРөРІСӢРө РҝРҫРәР°Р·Р°СӮРөР»Рё/РҪР°С…РҫРҙРәРё Рё СҮСӮРҫ РҫРҪРё РјРҫРіСғСӮ РҫР·РҪР°СҮР°СӮСҢ РІ РҫРұСүРөРј СҒРјСӢСҒР»Рө.\n"
        "4) Р’РҫР·РјРҫР¶РҪСӢРө СӮСҖРөРІРҫР¶РҪСӢРө РҝСғРҪРәСӮСӢ, РәРҫСӮРҫСҖСӢРө СҒСӮРҫРёСӮ РҫРұСҒСғРҙРёСӮСҢ СҒ РІСҖР°СҮРҫРј.\n"
        "5) РЎРҝРёСҒРҫРә РІРҫРҝСҖРҫСҒРҫРІ РІСҖР°СҮСғ.\n"
        "6) Р§СӮРҫ РҝСҖРҫРІРөСҖРёСӮСҢ/РәР°РәРёРө РҙР°РҪРҪСӢРө РҪСғР¶РҪСӢ РҙР»СҸ РҝРҫР»РҪРҫСҶРөРҪРҪРҫР№ РҫСҶРөРҪРәРё.\n"
        "Р’ РәРҫРҪСҶРө РҫРұСҸР·Р°СӮРөР»СҢРҪРҫ РҙРҫРұР°РІСҢ: СҚСӮРҫ РҪРө РҙРёР°РіРҪРҫР· Рё РҪРө Р·Р°РјРөРҪР° РәРҫРҪСҒСғР»СҢСӮР°СҶРёРё РІСҖР°СҮР°."
        f"{extra}\n\nРўРөРәСҒСӮ РҙРҫРәСғРјРөРҪСӮР°:\n{text[:24000]}"
    )

async def _medical_analyze_text(update: Update, context: ContextTypes.DEFAULT_TYPE, text: str, goal: str | None = None):
    user_id = update.effective_user.id
    track = _mode_track_get(user_id)
    await update.effective_message.reply_text("рҹ©ә РҗРҪР°Р»РёР·РёСҖСғСҺ РјРөРҙРёСҶРёРҪСҒРәРёР№ РјР°СӮРөСҖРёР°Р». РӯСӮРҫ СҒРҝСҖР°РІРҫСҮРҪСӢР№ СҖР°Р·РұРҫСҖ, РҪРө РҙРёР°РіРҪРҫР·вҖҰ")
    ans = await ask_openai_text(_medical_doc_prompt(text, track=track, goal=goal))
    if MEDICAL_DISCLAIMER.strip() not in ans:
        ans = ans.rstrip() + MEDICAL_DISCLAIMER
    await update.effective_message.reply_text(ans[:3900])
    if len(ans) > 3900:
        await update.effective_message.reply_text(ans[3900:7800])
    await maybe_tts_reply(update, context, ans[:TTS_MAX_CHARS])

async def _medical_analyze_image(update: Update, context: ContextTypes.DEFAULT_TYPE, img_bytes: bytes, goal: str | None = None):
    user_id = update.effective_user.id
    track = _mode_track_get(user_id)
    await update.effective_message.reply_text("рҹ©ә РҗРҪР°Р»РёР·РёСҖСғСҺ РјРөРҙРёСҶРёРҪСҒРәРҫРө РёР·РҫРұСҖР°Р¶РөРҪРёРө/СҒРҪРёРјРҫРә. РӯСӮРҫ СҒРҝСҖР°РІРҫСҮРҪСӢР№ СҖР°Р·РұРҫСҖ, РҪРө РҙРёР°РіРҪРҫР·вҖҰ")
    prompt = (
        f"Р Р°Р·РұРөСҖРё РјРөРҙРёСҶРёРҪСҒРәРҫРө РёР·РҫРұСҖР°Р¶РөРҪРёРө/РҙРҫРәСғРјРөРҪСӮ: {_medical_track_title(track)}. "
        "Р•СҒР»Рё СҚСӮРҫ С„РҫСӮРҫ РІСӢРҝРёСҒРәРё/Р°РҪР°Р»РёР·РҫРІ вҖ” РҝСҖРҫСҮРёСӮР°Р№ РІРёРҙРёРјСӢР№ СӮРөРәСҒСӮ Рё СҒСӮСҖСғРәСӮСғСҖРёСҖСғР№. "
        "Р•СҒР»Рё СҚСӮРҫ СҒРҪРёРјРҫРә РңР Рў/РҡРў/СҖРөРҪСӮРіРөРҪ/РЈР—Рҳ вҖ” РҫРҝРёСҲРё СӮРҫР»СҢРәРҫ РІРёРҙРёРјСӢРө СҚР»РөРјРөРҪСӮСӢ Рё РҫРіСҖР°РҪРёСҮРөРҪРёСҸ: "
        "РҝРҫ РҫРҙРҪРҫРјСғ РёР·РҫРұСҖР°Р¶РөРҪРёСҺ РҪРөР»СҢР·СҸ СҒСӮР°РІРёСӮСҢ РҙРёР°РіРҪРҫР·, РҪСғР¶РөРҪ РҫС„РёСҶРёР°Р»СҢРҪСӢР№ РҝСҖРҫСӮРҫРәРҫР» Рё РІСҖР°СҮ. "
        "Р”Р°Р№: 1) СҮСӮРҫ РІРёРҙРҪРҫ, 2) РІРҫР·РјРҫР¶РҪСӢР№ СҒРјСӢСҒР» СӮРөСҖРјРёРҪРҫРІ/РҝРҫРәР°Р·Р°СӮРөР»РөР№, 3) СҮСӮРҫ СғСӮРҫСҮРҪРёСӮСҢ Сғ РІСҖР°СҮР°, "
        "4) РәР°РәРёРө РҙРҫРҝРҫР»РҪРёСӮРөР»СҢРҪСӢРө РҙР°РҪРҪСӢРө РҪСғР¶РҪСӢ. РқРө РҪР°Р·РҪР°СҮР°Р№ Р»РөСҮРөРҪРёРө. "
        "Р’ РәРҫРҪСҶРө РҙРҫРұР°РІСҢ РҝСҖРөРҙСғРҝСҖРөР¶РҙРөРҪРёРө, СҮСӮРҫ СҚСӮРҫ РҪРө РҫС„РёСҶРёР°Р»СҢРҪРҫРө РјРөРҙРёСҶРёРҪСҒРәРҫРө Р·Р°РәР»СҺСҮРөРҪРёРө."
    )
    if goal:
        prompt += f"\nРҰРөР»СҢ РҝРҫР»СҢР·РҫРІР°СӮРөР»СҸ: {goal}"
    b64 = base64.b64encode(img_bytes).decode("ascii")
    ans = await ask_openai_vision(prompt, b64, sniff_image_mime(img_bytes))
    if MEDICAL_DISCLAIMER.strip() not in ans:
        ans = ans.rstrip() + MEDICAL_DISCLAIMER
    await update.effective_message.reply_text(ans[:3900])
    if len(ans) > 3900:
        await update.effective_message.reply_text(ans[3900:7800])
    await maybe_tts_reply(update, context, ans[:TTS_MAX_CHARS])

async def on_photo_revival_capability(update: Update, context: ContextTypes.DEFAULT_TYPE):
    _set_waiting_photo_revival(update, context)
    await update.effective_message.reply_text(_photo_revival_capability_text(), reply_markup=main_kb)

def capability_answer(text: str) -> str | None:
    """
    РҡРҫСҖРҫСӮРәРёРө РҫСӮРІРөСӮСӢ РҪР° РІРҫРҝСҖРҫСҒСӢ РІРёРҙР°:
    - В«СӮСӢ РјРҫР¶РөСҲСҢ Р°РҪР°Р»РёР·РёСҖРҫРІР°СӮСҢ PDF?В»
    - В«СӮСӢ СғРјРөРөСҲСҢ СҖР°РұРҫСӮР°СӮСҢ СҒ СҚР»РөРәСӮСҖРҫРҪРҪСӢРјРё РәРҪРёРіР°РјРё?В»
    - В«СӮСӢ РјРҫР¶РөСҲСҢ СҒРҫР·РҙР°РІР°СӮСҢ РІРёРҙРөРҫ?В» Рё СӮ.Рҝ.

    Р’Р°Р¶РҪРҫ: РҪРө РҝРөСҖРөС…РІР°СӮСӢРІР°РөРј СҖРөР°Р»СҢРҪСӢРө РәРҫРјР°РҪРҙСӢ
    В«СҒРҙРөР»Р°Р№ РІРёРҙРөРҫвҖҰВ», В«СҒРіРөРҪРөСҖРёСҖСғР№ РәР°СҖСӮРёРҪРәСғвҖҰВ» Рё СӮ.Рҙ.
    """
    tl = (text or "").strip().lower()
    if not tl:
        return None

    if _is_general_capability_query(tl):
        return (
            "Р”Р°. РҜ вҖ” Neyro-Bot GPT 5 Studio, РјСғР»СҢСӮРёРјРҫРҙРөР»СҢРҪР°СҸ AI-СҒСӮСғРҙРёСҸ РІ Telegram РҙР»СҸ СӮРөРәСҒСӮР°, РҙРҫРәСғРјРөРҪСӮРҫРІ, С„РҫСӮРҫ, РІРёРҙРөРҫ, РјСғР·СӢРәРё, СҖРөСҮРё Рё live-РҙР°РҪРҪСӢС….\n\n"
            "РһСҒРҪРҫРІРҪСӢРө СҖРөР¶РёРјСӢ:\n"
            "вҖў рҹҺ“ РЈСҮС‘РұР° вҖ” РҫРұСҠСҸСҒРҪРөРҪРёРө СӮРөРј, Р·Р°РҙР°СҮРё, РәРҫРҪСҒРҝРөРәСӮСӢ РёР· PDF/EPUB/DOCX, СҚСҒСҒРө Рё РҝР»Р°РҪСӢ Рә СҚРәР·Р°РјРөРҪСғ.\n"
            "вҖў рҹ’ј Р Р°РұРҫСӮР°/Р‘РёР·РҪРөСҒ вҖ” РҝРёСҒСҢРјР°, РҡРҹ, РҙРҫРәСғРјРөРҪСӮСӢ, Р°РҪР°Р»РёСӮРёРәР°, РҝСҖРөР·РөРҪСӮР°СҶРёРё, PDF-РәР°СӮР°Р»РҫРіРё, Р»РҫРіРҫСӮРёРҝСӢ, ToDo, РұСҖРёС„СӢ Рё СҖРөСӮСғСҲСҢ РұРёР·РҪРөСҒ-С„РҫСӮРҫ.\n"
            "вҖў рҹ”Ҙ Р Р°Р·РІР»РөСҮРөРҪРёСҸ вҖ” РҫР¶РёРІР»РөРҪРёРө С„РҫСӮРҫ, Р·Р°РјРөРҪР° Р»РёСҶР°, СғРҙР°Р»РөРҪРёРө/Р·Р°РјРөРҪР° С„РҫРҪР°, СғРҙР°Р»РөРҪРёРө РІРҫРҙСҸРҪСӢС… Р·РҪР°РәРҫРІ, Reels/Shorts, РјРёРҪРё-С„РёР»СҢРјСӢ, СҒСҶРөРҪР°СҖРёРё.\n"
            "вҖў рҹ©ә РңРөРҙРёСҶРёРҪР° вҖ” СҒРҝСҖР°РІРҫСҮРҪСӢР№ СҖР°Р·РұРҫСҖ РІСӢРҝРёСҒРҫРә, Р°РҪР°Р»РёР·РҫРІ, РңР Рў/РҡРў/СҒРҪРёРјРәРҫРІ РұРөР· РҝРҫСҒСӮР°РҪРҫРІРәРё РҙРёР°РіРҪРҫР·Р°.\n"
            "вҖў рҹ§  Р”РІРёР¶РәРё вҖ” РҝСҖРёРҪСғРҙРёСӮРөР»СҢРҪСӢР№ РІСӢРұРҫСҖ GPT-5, OpenAI Images, Runway, Sora 2, Kling, Midjourney, Deepgram/STT/TTS, Photoroom, PiAPI/Segmind.\n\n"
            "РңРҫР¶РҪРҫ РҝСҖРҫСҒСӮРҫ РҪР°РҝРёСҒР°СӮСҢ Р·Р°РҙР°СҮСғ РёР»Рё РҫСӮРҝСҖР°РІРёСӮСҢ С„РҫСӮРҫ/РҙРҫРәСғРјРөРҪСӮ/РіРҫР»РҫСҒРҫРІРҫРө. Р”Р»СҸ Р·Р°РјРөРҪСӢ Р»РёСҶР°: рҹ”Ҙ Р Р°Р·РІР»РөСҮРөРҪРёСҸ вҶ’ рҹҺӯ Р—Р°РјРөРҪР° Р»РёСҶР° РҪР° С„РҫСӮРҫ. "
            "Р”Р»СҸ С„РҫРҪР°: рҹ”Ҙ Р Р°Р·РІР»РөСҮРөРҪРёСҸ вҶ’ рҹ§ј РЈРҙР°Р»РёСӮСҢ С„РҫРҪ РёР»Рё рҹ–ј Р—Р°РјРөРҪРёСӮСҢ С„РҫРҪ. Р”Р»СҸ РҫР·РІСғСҮРәРё РҫСӮРІРөСӮРҫРІ: /voice_on."
        )

    if _is_medical_capability_question(tl):
        return _medical_capability_text()

    if _is_photo_revival_question(tl):
        return _photo_revival_capability_text()

    # --- Р—Р°РјРөРҪР° Р»РёСҶ / FaceSwap ---
    if (
        re.search(r"(РјРҫР¶(РөСҲСҢ|РөСӮРө|РҪРҫ)|СғРјРө(РөСҲСҢ|РөСӮРө)|СҒРҝРҫСҒРҫРұРөРҪ|РҝРҫРҙРҙРөСҖР¶РёРІР°РөСҲСҢ|РҙРөР»Р°РөСҲСҢ|Р·Р°РјРөРҪСҸРөСҲСҢ|РјРҫР¶РөСӮ\s+Р»Рё)", tl, re.I)
        and re.search(r"(Р·Р°РјРөРҪ|РҝРҫРјРөРҪСҸ|РҝРҫРҙРјРөРҪ|face\s*swap|faceswap)", tl, re.I)
        and re.search(r"(Р»РёСҶ|Р»РёСҶР°|Р»РёСҶРҫ|face|СҮРөР»РҫРІРөРәР°|С„РҫСӮРҫ|С„РҫСӮРҫРіСҖР°С„|РәР°СҖСӮРёРҪРә|РёР·РҫРұСҖР°Р¶РөРҪ)", tl, re.I)
    ):
        return (
            "Р”Р°, СҸ СғРјРөСҺ Р·Р°РјРөРҪСҸСӮСҢ Р»РёСҶР° РҪР° С„РҫСӮРҫРіСҖР°С„РёСҸС…. РңРҫР¶РҪРҫ СҖР°РұРҫСӮР°СӮСҢ СҒ С„РҫСӮРҫ, РіРҙРө РҫРҙРёРҪ, РҙРІР° РёР»Рё РҪРөСҒРәРҫР»СҢРәРҫ СҮРөР»РҫРІРөРә: "
            "РұРҫСӮ РҝРҫРәР°Р¶РөСӮ РҪР°Р№РҙРөРҪРҪСӢРө Р»РёСҶР° СҒ РҪРҫРјРөСҖР°РјРё, РІСӢ РІСӢРұРөСҖРөСӮРө, РәР°РәРҫРө РёРјРөРҪРҪРҫ Р»РёСҶРҫ Р·Р°РјРөРҪРёСӮСҢ, Р·Р°СӮРөРј Р·Р°РіСҖСғР·РёСӮРө Р»РёСҶРҫ-РёСҒСӮРҫСҮРҪРёРә.\n\n"
            "РҡР°Рә Р·Р°РҝСғСҒСӮРёСӮСҢ:\n"
            "1) РҪР°Р¶РјРёСӮРө рҹ”Ҙ Р Р°Р·РІР»РөСҮРөРҪРёСҸ вҶ’ рҹҺӯ Р—Р°РјРөРҪР° Р»РёСҶР° РҪР° С„РҫСӮРҫ;\n"
            "2) Р·Р°РіСҖСғР·РёСӮРө РҫСҒРҪРҫРІРҪРҫРө С„РҫСӮРҫ;\n"
            "3) РІСӢРұРөСҖРёСӮРө СҶРөР»РөРІРҫРө Р»РёСҶРҫ РҝРҫ РҪРҫРјРөСҖСғ вҖ” СҒР»РөРІР°/СҶРөРҪСӮСҖ/СҒРҝСҖР°РІР°;\n"
            "4) Р·Р°РіСҖСғР·РёСӮРө С„РҫСӮРҫ Р»РёСҶР°, РәРҫСӮРҫСҖРҫРө РҪСғР¶РҪРҫ РІСҒСӮР°РІРёСӮСҢ;\n"
            "5) РІСӢРұРөСҖРёСӮРө РәР°СҮРөСҒСӮРІРҫ: вҡЎ Р‘СӢСҒСӮСҖРҫ СҮРөСҖРөР· PiAPI РёР»Рё рҹ’Һ РҹСҖРөРјРёСғРј СҮРөСҖРөР· Segmind, РөСҒР»Рё РҪР° РұР°Р»Р°РҪСҒРө Segmind РөСҒСӮСҢ РәСҖРөРҙРёСӮСӢ.\n\n"
            "РўР°РәР¶Рө РјРҫР¶РҪРҫ РҝСҖРҫСҒСӮРҫ Р·Р°РіСҖСғР·РёСӮСҢ С„РҫСӮРҫ Рё РҪР°Р¶Р°СӮСҢ РәРҪРҫРҝРәСғ рҹҺӯ Р—Р°РјРөРҪР° Р»РёСҶР°. "
            "РҜ СҒРҫС…СҖР°РҪСҸСҺ СӮРөР»Рҫ, РҫРҙРөР¶РҙСғ, С„РҫРҪ Рё РәРҫРјРҝРҫР·РёСҶРёСҺ РёСҒС…РҫРҙРҪРҫРіРҫ С„РҫСӮРҫ, РјРөРҪСҸРөСӮСҒСҸ СӮРҫР»СҢРәРҫ РІСӢРұСҖР°РҪРҪРҫРө Р»РёСҶРҫ."
        )

    # --- РһР¶РёРІР»РөРҪРёРө С„РҫСӮРҫ / image-to-video ---
    if (
        re.search(r"(РјРҫР¶(РөСҲСҢ|РөСӮРө)|СғРјРө(РөСҲСҢ|РөСӮРө)|РјРҫР¶РөСӮ\s+Р»Рё|СҒРҝРҫСҒРҫРұРөРҪ|РҝРҫРҙРҙРөСҖР¶РёРІР°РөСҲСҢ)", tl)
        and re.search(r"(РҫР¶РёРІ|Р°РҪРёРјРёСҖ|РҙРІРёР¶РөРҪРё|image\s*to\s*video|i2v)", tl)
        and re.search(r"(С„РҫСӮРҫ|С„РҫСӮРҫРіСҖР°С„|РәР°СҖСӮРёРҪРә|РёР·РҫРұСҖР°Р¶РөРҪ|image|picture)", tl)
    ):
        return (
            "Р”Р°, РјРҫРіСғ РҫР¶РёРІРёСӮСҢ С„РҫСӮРҫРіСҖР°С„РёСҺ Рё СҒРҙРөР»Р°СӮСҢ РёР· РҪРөС‘ РәРҫСҖРҫСӮРәРҫРө РІРёРҙРөРҫ. "
            "Р—Р°РіСҖСғР·РёСӮРө С„РҫСӮРҫ вҖ” СҸ РҝРҫРәР°Р¶Сғ РәРҪРҫРҝРәРё Runway, Kling Рё Sora 2 РұРөР· Р»СҺРҙРөР№. "
            "РңРҫР¶РҪРҫ СӮР°РәР¶Рө РҫСӮРҝСҖР°РІРёСӮСҢ С„РҫСӮРҫ СҒ РҝРҫРҙРҝРёСҒСҢСҺ: В«РҫР¶РёРІРё С„РҫСӮРҫ: Р»С‘РіРәР°СҸ СғР»СӢРұРәР°, РҙРІРёР¶РөРҪРёРө РәР°РјРөСҖСӢ, 5 СҒРөРәСғРҪРҙ, 9:16В»."
        )

    # --- РңРөРҙРёСҶРёРҪР° / РјРөРҙ. РҙРҫРәСғРјРөРҪСӮСӢ / Р°РҪР°Р»РёР·СӢ / РңР Рў / РҡРў ---
    if _is_medical_capability_question(tl):
        return _medical_capability_text()

    # --- Р”РҫРәСғРјРөРҪСӮСӢ / С„Р°Р№Р»СӢ ---
    if re.search(r"\b(pdf|docx|epub|fb2|txt|mobi|azw)\b", tl) and "?" in tl:
        return (
            "Р”Р°, СҸ РјРҫРіСғ РҝРҫРјРҫСҮСҢ СҒ Р°РҪР°Р»РёР·РҫРј РҙРҫРәСғРјРөРҪСӮРҫРІ Рё СҚР»РөРәСӮСҖРҫРҪРҪСӢС… РәРҪРёРі. "
            "РһСӮРҝСҖР°РІСҢ С„Р°Р№Р» (PDF, EPUB, DOCX, FB2, TXT, MOBI/AZW вҖ“ РҝРҫ РІРҫР·РјРҫР¶РҪРҫСҒСӮРё), "
            "Р° РІ СҒРҫРҫРұСүРөРҪРёРё РҪР°РҝРёСҲРё, СҮСӮРҫ РҪСғР¶РҪРҫ: РәРҫРҪСҒРҝРөРәСӮ, РҝР»Р°РҪ, СҖР°Р·РұРҫСҖ Рё СӮ.Рҝ."
        )

    # --- РҗСғРҙРёРҫ / СҖРөСҮСҢ ---
    if "Р°СғРҙРёРҫ" in tl or "РіРҫР»РҫСҒРҫРІ" in tl or "speech" in tl:
        if "?" in tl or "РјРҫР¶РөСҲСҢ" in tl or "СғРјРөРөСҲСҢ" in tl:
            return (
                "Р”Р°, СҸ РјРҫРіСғ СҖР°СҒРҝРҫР·РҪР°РІР°СӮСҢ СҖРөСҮСҢ РёР· РіРҫР»РҫСҒРҫРІСӢС… Рё Р°СғРҙРёРҫ. "
                "РҹСҖРҫСҒСӮРҫ РҝСҖРёСҲР»Рё РіРҫР»РҫСҒРҫРІРҫРө СҒРҫРҫРұСүРөРҪРёРө вҖ” СҸ СҖР°СҒСҲРёС„СҖСғСҺ РөРіРҫ РІ СӮРөРәСҒСӮ "
                "Рё РҫСӮРІРөСҮСғ РәР°Рә РҪР° РҫРұСӢСҮРҪСӢР№ Р·Р°РҝСҖРҫСҒ."
            )

    # --- РңСғР·СӢРәР° / Suno ---
    if (
        re.search(r"(РјСғР·СӢРә|РҝРөСҒРҪ|СӮСҖРөРә|РҙР¶РёРҪРіР»|РјРёРҪСғСҒРҫРІРә|suno|music|song)", tl, re.I)
        and re.search(r"(РјРҫР¶(РөСҲСҢ|РөСӮРө|РҪРҫ)|СғРјРө(РөСҲСҢ|РөСӮРө)|СҒРҫР·РҙР°[РөС‘]СҲСҢ|РҙРөР»Р°РөСҲСҢ|СҒРіРөРҪРөСҖРёСҖСғРөСҲСҢ|РҪР°РҝРёСҒР°СӮСҢ|РҝРҫР»СғСҮРёСӮСҒСҸ)", tl, re.I)
    ):
        return (
            "Р”Р°, РјРҫРіСғ Р·Р°РҝСғСҒРәР°СӮСҢ РіРөРҪРөСҖР°СҶРёСҺ РјСғР·СӢРәРё СҮРөСҖРөР· Suno. "
            "РҹРҫРҙС…РҫРҙРёСӮ РҙР»СҸ РҝРөСҒРөРҪ СҒ СӮРөРәСҒСӮРҫРј, РҙР¶РёРҪРіР»РҫРІ, С„РҫРҪРҫРІРҫР№ РјСғР·СӢРәРё, РёРҪСӮСҖРҫ/Р°СғСӮСҖРҫ, РјРёРҪСғСҒРҫРІРҫРә Рё РҙРөРјРҫ-СӮСҖРөРәРҫРІ.\n\n"
            "РҡР°Рә Р·Р°РҝСғСҒСӮРёСӮСҢ: рҹ”Ҙ Р Р°Р·РІР»РөСҮРөРҪРёСҸ вҶ’ рҹҺө РңСғР·СӢРәР° / РҝРөСҒРҪСҸ РёР»Рё РәРҫРјР°РҪРҙР° /music <РҫРҝРёСҒР°РҪРёРө>. "
            "Р’ РҫРҝРёСҒР°РҪРёРё СғРәР°Р¶РёСӮРө Р¶Р°РҪСҖ, РҪР°СҒСӮСҖРҫРөРҪРёРө, СҸР·СӢРә, РҝСҖРёРјРөСҖРҪСӢР№ СӮРөРјРҝ/BPM, РҪСғР¶РөРҪ Р»Рё РІРҫРәР°Р», СӮРөРәСҒСӮ РәСғРҝР»РөСӮР°/РҝСҖРёРҝРөРІР° РёР»Рё instrumental."
        )

    # --- Р’РёРҙРөРҫ / СҖРҫР»РёРәРё / Reels / Shorts ---
    # Р’Р°Р¶РҪРҫ: РіРҫР»РҫСҒРҫРІСӢРө СҮР°СҒСӮРҫ РҝСҖРёС…РҫРҙСҸСӮ РұРөР· РІРҫРҝСҖРҫСҒРёСӮРөР»СҢРҪРҫРіРҫ Р·РҪР°РәР°: В«СӮСӢ РјРҫР¶РөСҲСҢ СҒРҫР·РҙР°СӮСҢ РІРёРҙРөРҫ ...В».
    # РўР°РәРҫР№ Р·Р°РҝСҖРҫСҒ РҪРө РҙРҫР»Р¶РөРҪ СғС…РҫРҙРёСӮСҢ РІ РҫРұСүРёР№ GPT-СҮР°СӮ, РіРҙРө РјРҫРҙРөР»СҢ РјРҫР¶РөСӮ РҫСӮРІРөСӮРёСӮСҢ РҫСӮРәР°Р·РҫРј.
    if (
        re.search(r"(РІРёРҙРөРҫ|СҖРҫР»РёРә|РәР»РёРҝ|reels?|shorts?|video|clip)", tl, re.I)
        and re.search(r"(СӮСӢ\s+)?(РјРҫР¶(РөСҲСҢ|РөСӮРө|РҪРҫ)|СғРјРө(РөСҲСҢ|РөСӮРө)|СҒРҝРҫСҒРҫРұРөРҪ|РҝРҫРҙРҙРөСҖР¶РёРІР°РөСҲСҢ|РҙРөР»Р°РөСҲСҢ|СҒРҫР·РҙР°[РөС‘]СҲСҢ|РјРҫР¶РөСӮ\s+Р»Рё|РјРҫР¶РҪРҫ\s+Р»Рё|РҝРҫР»СғСҮРёСӮСҒСҸ\s+Р»Рё)", tl, re.I)
    ):
        return (
            "Р”Р°, РјРҫРіСғ СҒРҫР·РҙР°РІР°СӮСҢ РәРҫСҖРҫСӮРәРёР№ РІРёРҙРөРҫРәРҫРҪСӮРөРҪСӮ РҝРҫ СӮРөРәСҒСӮСғ РёР»Рё РіРҫР»РҫСҒРҫРІРҫРјСғ Р·Р°РҝСҖРҫСҒСғ. "
            "РқР°РҝРёСҲРёСӮРө РёР»Рё СҒРәР°Р¶РёСӮРө РәРҫРјР°РҪРҙРҫР№: В«СҒРҫР·РҙР°Р№ РІРёРҙРөРҫ: РјРөРҙРІРөРҙСҢ Р»РөСӮРёСӮ РҪР° РІРҫР·РҙСғСҲРҪРҫРј СҲР°СҖРө, 5 СҒРөРәСғРҪРҙ, 16:9В». "
            "РҹРҫСҒР»Рө СҚСӮРҫРіРҫ СҸ РҝРҫРәР°Р¶Сғ РІСӢРұРҫСҖ РҙРІРёР¶РәР°: Sora 2 РұРөР· Р»СҺРҙРөР№, Kling РёР»Рё Runway. "
            "Luma СҒРөР№СҮР°СҒ РІСҖРөРјРөРҪРҪРҫ СҒРәСҖСӢСӮР° Рё РҪРө РёСҒРҝРҫР»СҢР·СғРөСӮСҒСҸ."
        )

    # --- РҡР°СҖСӮРёРҪРәРё / РёР·РҫРұСҖР°Р¶РөРҪРёСҸ ---
    if re.search(r"(РәР°СҖСӮРёРҪРә|РёР·РҫРұСҖР°Р¶РөРҪ|С„РҫСӮРҫ|С„РҫСӮРҫРіСҖР°С„|image|picture|Р»РҫРіРҫСӮРёРҝ|РұР°РҪРҪРөСҖ)", tl) and "?" in tl:
        return (
            "Р”Р°, РјРҫРіСғ СҖР°РұРҫСӮР°СӮСҢ СҒ РёР·РҫРұСҖР°Р¶РөРҪРёСҸРјРё: Р°РҪР°Р»РёР·РёСҖРҫРІР°СӮСҢ С„РҫСӮРҫ, СғРҙР°Р»СҸСӮСҢ/Р·Р°РјРөРҪСҸСӮСҢ С„РҫРҪ, СҖР°СҒСҲРёСҖСҸСӮСҢ РәР°РҙСҖ, "
            "РҙРөР»Р°СӮСҢ СҖРөСӮСғСҲСҢ СҒРҫРұСҒСӮРІРөРҪРҪСӢС… РёР·РҫРұСҖР°Р¶РөРҪРёР№, СғРұРёСҖР°СӮСҢ Р»РёСҲРҪРёРө РҪР°РҙРҝРёСҒРё/РҫРұСҠРөРәСӮСӢ, "
            "СҒРҫР·РҙР°РІР°СӮСҢ Р»РҫРіРҫСӮРёРҝСӢ/РәР°СҖСӮРёРҪРәРё РҝРҫ РҫРҝРёСҒР°РҪРёСҺ Рё РҫР¶РёРІР»СҸСӮСҢ С„РҫСӮРҫ РІ РәРҫСҖРҫСӮРәРҫРө РІРёРҙРөРҫ. "
            "Р—Р°РіСҖСғР·РёСӮРө С„РҫСӮРҫ вҖ” РҝРҫСҸРІСҸСӮСҒСҸ РұСӢСҒСӮСҖСӢРө РәРҪРҫРҝРәРё РҙРөР№СҒСӮРІРёР№."
        )

    # РқРёСҮРөРіРҫ РҝРҫРҙС…РҫРҙСҸСүРөРіРҫ вҖ” РҝСғСҒСӮСҢ РҙР°Р»СҢСҲРө РҫРұСҖР°РұР°СӮСӢРІР°РөСӮСҒСҸ РҫРұСӢСҮРҪРҫР№ Р»РҫРіРёРәРҫР№
    return None


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РңРҫРҙСӢ/РҙРІРёР¶РәРё РҙР»СҸ study в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def _uk(user_id: int, name: str) -> str: return f"user:{user_id}:{name}"
def mode_set(user_id: int, mode: str):     kv_set(_uk(user_id, "mode"), (mode or "default"))
def mode_get(user_id: int) -> str:         return kv_get(_uk(user_id, "mode"), "default") or "default"
def engine_set(user_id: int, engine: str): kv_set(_uk(user_id, "engine"), (engine or "gpt"))
def engine_get(user_id: int) -> str:       return kv_get(_uk(user_id, "engine"), "gpt") or "gpt"
def study_sub_set(user_id: int, sub: str): kv_set(_uk(user_id, "study_sub"), (sub or "explain"))
def study_sub_get(user_id: int) -> str:    return kv_get(_uk(user_id, "study_sub"), "explain") or "explain"

def modes_kb():
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("рҹҺ“ РЈСҮС‘РұР°", callback_data="mode:set:study"),
         InlineKeyboardButton("рҹ–ј РӨРҫСӮРҫ",  callback_data="mode:set:photo")],
        [InlineKeyboardButton("рҹ“„ Р”РҫРәСғРјРөРҪСӮСӢ", callback_data="mode:set:docs"),
         InlineKeyboardButton("рҹҺҷ Р“РҫР»РҫСҒ",     callback_data="mode:set:voice")],
        [InlineKeyboardButton("рҹ§  Р”РІРёР¶РәРё", callback_data="mode:engines")]
    ])

def study_kb():
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("рҹ”Қ РһРұСҠСҸСҒРҪРөРҪРёРө",          callback_data="study:set:explain"),
         InlineKeyboardButton("рҹ§® Р—Р°РҙР°СҮРё",              callback_data="study:set:tasks")],
        [InlineKeyboardButton("вңҚпёҸ РӯСҒСҒРө / СҖРөС„РөСҖР°СӮ", callback_data="study:set:essay")],
        [InlineKeyboardButton("рҹ“қ РӯРәР·Р°РјРөРҪ/РәРІРёР·",        callback_data="study:set:quiz")]
    ])

async def study_process_text(update: Update, context: ContextTypes.DEFAULT_TYPE, text: str):
    sub = study_sub_get(update.effective_user.id)
    if sub == "explain":
        prompt = f"РһРұСҠСҸСҒРҪРё РҝСҖРҫСҒСӮСӢРјРё СҒР»РҫРІР°РјРё, СҒ 2вҖ“3 РҝСҖРёРјРөСҖР°РјРё Рё РјРёРҪРё-РёСӮРҫРіРҫРј:\n\n{text}"
    elif sub == "tasks":
        prompt = ("Р РөСҲРё Р·Р°РҙР°СҮСғ(Рё) РҝРҫСҲР°РіРҫРІРҫ: С„РҫСҖРјСғР»СӢ, РҝРҫСҸСҒРҪРөРҪРёСҸ, РёСӮРҫРіРҫРІСӢР№ РҫСӮРІРөСӮ. "
                  "Р•СҒР»Рё РҪРө С…РІР°СӮР°РөСӮ РҙР°РҪРҪСӢС… вҖ” СғСӮРҫСҮРҪСҸСҺСүРёРө РІРҫРҝСҖРҫСҒСӢ РІ РәРҫРҪСҶРө.\n\n" + text)
    elif sub == "essay":
        prompt = ("РқР°РҝРёСҲРё СҒСӮСҖСғРәСӮСғСҖРёСҖРҫРІР°РҪРҪСӢР№ СӮРөРәСҒСӮ 400вҖ“600 СҒР»РҫРІ (СҚСҒСҒРө/СҖРөС„РөСҖР°СӮ/РҙРҫРәР»Р°Рҙ): "
                  "РІРІРөРҙРөРҪРёРө, 3вҖ“5 СӮРөР·РёСҒРҫРІ СҒ С„Р°РәСӮР°РјРё, РІСӢРІРҫРҙ, СҒРҝРёСҒРҫРә РёР· 3 РёСҒСӮРҫСҮРҪРёРәРҫРІ (РөСҒР»Рё СғРјРөСҒСӮРҪРҫ).\n\nРўРөРјР°:\n" + text)
    elif sub == "quiz":
        prompt = ("РЎРҫСҒСӮР°РІСҢ РјРёРҪРё-РәРІРёР· РҝРҫ СӮРөРјРө: 10 РІРҫРҝСҖРҫСҒРҫРІ, Сғ РәР°Р¶РҙРҫРіРҫ 4 РІР°СҖРёР°РҪСӮР° AвҖ“D; "
                  "РІ РәРҫРҪСҶРө РҙР°Р№ РәР»СҺСҮ РҫСӮРІРөСӮРҫРІ (РҪРҫРјРөСҖвҶ’РұСғРәРІР°). РўРөРјР°:\n\n" + text)
    else:
        prompt = text
    user_id = update.effective_user.id
    chat_id = update.effective_chat.id if update.effective_chat else 0
    ans = await ask_openai_text(prompt, user_id=user_id, chat_id=chat_id)
    await update.effective_message.reply_text(ans)
    _chat_memory_add(user_id, chat_id, "user", text)
    _chat_memory_add(user_id, chat_id, "assistant", ans)
    await maybe_tts_reply(update, context, ans[:TTS_MAX_CHARS])


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РҡРҪРҫРҝРәР° РҝСҖРёРІРөСӮСҒСӮРІРөРҪРҪРҫР№ РәР°СҖСӮРёРҪРәРё в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
async def cmd_set_welcome(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != OWNER_ID:
        await update.effective_message.reply_text("РҡРҫРјР°РҪРҙР° РҙРҫСҒСӮСғРҝРҪР° СӮРҫР»СҢРәРҫ РІР»Р°РҙРөР»СҢСҶСғ.")
        return
    if not context.args:
        await update.effective_message.reply_text("РӨРҫСҖРјР°СӮ: /set_welcome <url_РәР°СҖСӮРёРҪРәРё>")
        return
    url = " ".join(context.args).strip()
    kv_set("welcome_url", url)
    await update.effective_message.reply_text("РҡР°СҖСӮРёРҪРәР° РҝСҖРёРІРөСӮСҒСӮРІРёСҸ РҫРұРҪРҫРІР»РөРҪР°. РһСӮРҝСҖР°РІСҢСӮРө /start РҙР»СҸ РҝСҖРҫРІРөСҖРәРё.")

async def cmd_show_welcome(update: Update, context: ContextTypes.DEFAULT_TYPE):
    url = kv_get("welcome_url", BANNER_URL)
    if url:
        await update.effective_message.reply_photo(url, caption="РўРөРәСғСүР°СҸ РәР°СҖСӮРёРҪРәР° РҝСҖРёРІРөСӮСҒСӮРІРёСҸ")
    else:
        await update.effective_message.reply_text("РҡР°СҖСӮРёРҪРәР° РҝСҖРёРІРөСӮСҒСӮРІРёСҸ РҪРө Р·Р°РҙР°РҪР°.")



# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РҡРҫРјРјРөСҖСҮРөСҒРәРёР№ РҝСҖР°Р№СҒ в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def _pricing_catalog_text() -> str:
    lines = [
        "рҹ’° РҰРөРҪСӢ РҝР»Р°СӮРҪСӢС… РіРөРҪРөСҖР°СҶРёР№",
        "Р’СҒРө СҶРөРҪСӢ СғР¶Рө РІРәР»СҺСҮР°СҺСӮ СҒРөСҖРІРёСҒРҪСғСҺ РјР°СҖР¶Сғ Рё РҫРәСҖСғРіР»РөРҪСӢ РҙРҫ 5 РәСҖРөРҙРёСӮРҫРІ.",
        "1 РәСҖРөРҙРёСӮ = 1 вӮҪ; СҒРҝРёСҒР°РҪРёРө вҖ” СӮРҫР»СҢРәРҫ РҝРҫСҒР»Рө СғСҒРҝРөСҲРҪРҫРіРҫ СҖРөР·СғР»СҢСӮР°СӮР°.",
        "",
        "рҹ–ј РҳР·РҫРұСҖР°Р¶РөРҪРёСҸ",
        f"вҖў OpenAI Images, 1024Г—1024 medium вҖ” {_retail_credits(IMG_COST_USD)} РәСҖ.",
        f"вҖў Midjourney Fast вҖ” {_retail_credits(MIDJOURNEY_UNIT_COST_USD)} РәСҖ.",
        f"вҖў РһРұСҖР°РұРҫСӮРәР°/СҖРөСӮСғСҲСҢ/С„РҫРҪ вҖ” РҫСӮ {_retail_credits(IMG_PROCESS_COST_USD)} РәСҖ.",
        f"вҖў AI-СҒРөР»С„Рё вҖ” {_retail_credits(AI_SELFIE_UNIT_COST_USD)} РәСҖ.",
        "",
        "рҹҺ¬ Р’РёРҙРөРҫ",
        f"вҖў Sora 2 РұРөР· Р»СҺРҙРөР№, 5 СҒРөРә вҖ” {_video_price_credits('sora', 5)} РәСҖ.; 10 СҒРөРә вҖ” {_video_price_credits('sora', 10)} РәСҖ.",
        f"вҖў Kling, 5 СҒРөРә вҖ” {_video_price_credits('kling', 5)} РәСҖ.; 10 СҒРөРә вҖ” {_video_price_credits('kling', 10)} РәСҖ.",
        f"вҖў Runway Gen-4.5, 5 СҒРөРә вҖ” {_video_price_credits('runway', 5)} РәСҖ.; 10 СҒРөРә вҖ” {_video_price_credits('runway', 10)} РәСҖ.",
        "",
        "рҹҺө РҹСҖРҫРҙР°РәСҲРҪ-С„СғРҪРәСҶРёРё",
        f"вҖў Suno РјСғР·СӢРәР° вҖ” {_retail_credits(SUNO_COST_USD)} РәСҖ.",
        f"вҖў Р“РҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ вҖ” {_retail_credits(AVATAR_UNIT_COST_USD)} РәСҖ.",
        f"вҖў РӨРҫСӮРҫ вҶ’ РІРёРҙРөРҫРәР»РёРҝ СҒ РјСғР·СӢРәРҫР№ вҖ” {_retail_credits(PHOTO_CLIP_UNIT_COST_USD)} РәСҖ.",
        f"вҖў РҡР»РёРҝ СҒ РІРҫРәР°Р»РҫРј / lip-sync вҖ” {_retail_credits(VOCAL_CLIP_UNIT_COST_USD)} РәСҖ.",
        "",
        "рҹҺӯ Р‘РёР·РҪРөСҒ Рё С„РҫСӮРҫ",
        f"вҖў FaceSwap РұСӢСҒСӮСҖРҫ вҖ” {_retail_credits(FACESWAP_FAST_COST_USD)} РәСҖ.; РҝСҖРөРјРёСғРј вҖ” {_retail_credits(FACESWAP_PREMIUM_COST_USD)} РәСҖ.",
        f"вҖў РЎРұРҫСҖРәР° РҝСҖРөР·РөРҪСӮР°СҶРёРё PDF + PPTX вҖ” {_retail_credits(PRESENTATION_RENDER_COST_USD)} РәСҖ. + РІСӢРұСҖР°РҪРҪСӢРө AI-РёР·РҫРұСҖР°Р¶РөРҪРёСҸ.",
        "",
        "РӨРёРҪР°Р»СҢРҪР°СҸ СҶРөРҪР° РІСҒРөРіРҙР° РҝРҫРәР°Р·СӢРІР°РөСӮСҒСҸ РҙРҫ Р·Р°РҝСғСҒРәР° Р·Р°РҙР°СҮРё.",
    ]
    return "\n".join(lines)

async def cmd_prices(update: Update, context: ContextTypes.DEFAULT_TYPE):
    await update.effective_message.reply_text(
        _pricing_catalog_text(),
        reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("вһ• РҹРҫРҝРҫР»РҪРёСӮСҢ РұР°Р»Р°РҪСҒ", callback_data="topup")]]),
    )

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ Р‘Р°Р»Р°РҪСҒ / РҝРҫРҝРҫР»РҪРөРҪРёРө в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
async def cmd_balance(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    w = _wallet_get(user_id)
    total = _wallet_total_get(user_id)
    row = _usage_row(user_id)
    lim = _limits_for(user_id)
    credits = int(round(_usd_to_credits(total)))
    reserved = int(round(_usd_to_credits(_wallet_reserved_total_get(user_id))))
    available = max(0, credits - reserved)
    msg = (
        "рҹ§ҫ РҡСҖРөРҙРёСӮРҪСӢР№ РұР°Р»Р°РҪСҒ:\n"
        f"вҖў Р’СҒРөРіРҫ: {credits} РәСҖ.\n"
        f"вҖў Р”РҫСҒСӮСғРҝРҪРҫ: {available} РәСҖ.\n"
        + (f"вҖў Р—Р°СҖРөР·РөСҖРІРёСҖРҫРІР°РҪРҫ Р°РәСӮРёРІРҪСӢРјРё Р·Р°РҙР°СҮР°РјРё: {reserved} РәСҖ.\n" if reserved else "")
        + "\n1 РәСҖРөРҙРёСӮ = 1 вӮҪ. РҰРөРҪР° РәР°Р¶РҙРҫР№ РҝР»Р°СӮРҪРҫР№ РіРөРҪРөСҖР°СҶРёРё СғР¶Рө РІРәР»СҺСҮР°РөСӮ РјР°СҖР¶Сғ СҒРөСҖРІРёСҒР°; СҒРҝРёСҒР°РҪРёРө РҝСҖРҫРёСҒС…РҫРҙРёСӮ СӮРҫР»СҢРәРҫ РҝРҫСҒР»Рө СғСҒРҝРөСҲРҪРҫРіРҫ СҖРөР·СғР»СҢСӮР°СӮР°.\n\n"
        "РҗРәСӮСғР°Р»СҢРҪСӢР№ РҝСҖР°Р№СҒ: /prices\n\n"
        "РӣРёРјРёСӮСӢ СҒРөРіРҫРҙРҪСҸ:\n"
        f"вҖў РўРөРәСҒСӮ: {row['text_count']} / {lim['text_per_day']}\n"
        f"вҖў Р‘РөСҒРҝР»Р°СӮРҪР°СҸ РіРөРҪРөСҖР°СҶРёСҸ РәР°СҖСӮРёРҪРҫРә: {row['free_img_gen_count']} / {FREE_IMAGE_GENERATIONS_PER_DAY}\n"
        f"вҖў Р‘РөСҒРҝР»Р°СӮРҪР°СҸ РҫРұСҖР°РұРҫСӮРәР° С„РҫСӮРҫ: {row['free_img_proc_count']} / {FREE_IMAGE_PROCESSINGS_PER_DAY}\n"
    )
    kb = InlineKeyboardMarkup([[InlineKeyboardButton("вһ• РҹРҫРҝРҫР»РҪРёСӮСҢ РұР°Р»Р°РҪСҒ", callback_data="topup"), InlineKeyboardButton("рҹ’° РҰРөРҪСӢ", callback_data="pricing:list")]])
    await update.effective_message.reply_text(msg, reply_markup=kb)

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РҹРҫРҙРҝРёСҒРәР° / СӮР°СҖРёС„СӢ вҖ” UI Рё РҫРҝР»Р°СӮСӢ (PATCH) в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
# Р—Р°РІРёСҒРёРјРҫСҒСӮРё РҫРәСҖСғР¶РөРҪРёСҸ:
#  - YOOKASSA_PROVIDER_TOKEN  (РҝР»Р°СӮС‘Р¶РҪСӢР№ СӮРҫРәРөРҪ Telegram Payments РҫСӮ Р®Kassa)
#  - YOOKASSA_CURRENCY        (РҝРҫ СғРјРҫР»СҮР°РҪРёСҺ "RUB")
#  - CRYPTO_PAY_API_TOKEN     (https://pay.crypt.bot вҖ” СӮРҫРәРөРҪ РҝСҖРҫРҙР°РІСҶР°)
#  - CRYPTO_ASSET             (РҪР°РҝСҖРёРјРөСҖ "USDT", РҝРҫ СғРјРҫР»СҮР°РҪРёСҺ "USDT")
#  - PRICE_START_RUB, PRICE_PRO_RUB, PRICE_ULT_RUB  (СҶРөР»РҫРө СҮРёСҒР»Рҫ, вӮҪ)
#  - PRICE_START_USD, PRICE_PRO_USD, PRICE_ULT_USD  (СҮРёСҒР»Рҫ СҒ СӮРҫСҮРәРҫР№, $)
#
# РҘСҖР°РҪРёР»РёСүРө РҝРҫРҙРҝРёСҒРәРё Рё РәРҫСҲРөР»СҢРәР° РёСҒРҝРҫР»СҢР·СғРөСӮСҒСҸ РҪР° kv_*:
#   sub:tier:{user_id}   -> "start" | "pro" | "ultimate"
#   sub:until:{user_id}  -> ISO-СҒСӮСҖРҫРәР° РҙР°СӮСӢ РҫРәРҫРҪСҮР°РҪРёСҸ
#   wallet:usd:{user_id} -> РұР°Р»Р°РҪСҒ РІ USD (float)

YOOKASSA_PROVIDER_TOKEN = os.environ.get("YOOKASSA_PROVIDER_TOKEN", "").strip()
YOOKASSA_CURRENCY = (os.environ.get("YOOKASSA_CURRENCY") or "RUB").upper()

# v80: РҝСҖСҸРјСӢРө РҝР»Р°СӮРөР¶Рё Р®Kassa РҝРҫ СҒСҒСӢР»РәРө/QR/РұР°РҪРәРҫРІСҒРәРёРј РҝСҖРёР»РҫР¶РөРҪРёСҸРј.
# РӯСӮРҫ РқР• Telegram Payments invoice, Р° API Р®Kassa /v3/payments + polling СҒСӮР°СӮСғСҒР°.
#
# Render РҪР° РјРҫРұРёР»СҢРҪРҫРј РёРҪРҫРіРҙР° РҪРө РҙР°С‘СӮ РҙРҫРұР°РІРёСӮСҢ РҪРҫРІСӢРө ENV. РҹРҫСҚСӮРҫРјСғ РәР»СҺСҮРё Р®Kassa РјРҫР¶РҪРҫ
# РҝРөСҖРөРҙР°СӮСҢ СӮСҖРөРјСҸ РҝСғСӮСҸРјРё, РІ РҝРҫСҖСҸРҙРәРө РҝСҖРёРҫСҖРёСӮРөСӮР°:
# 1) ENV: YOO_SHOP_ID/YOO_SECRET_KEY РёР»Рё YOOKASSA_SHOP_ID/YOOKASSA_SECRET_KEY
# 2) РәРҫСҖРҫСӮРәРёРө ENV: YK_ID/YK_KEY
# 3) Render Secret File: /etc/secrets/yookassa.env СҒ СҒРҫРҙРөСҖР¶РёРјСӢРј:
#    YK_ID=1185809
#    YK_KEY=СҒРөРәСҖРөСӮРҪСӢР№_РәР»СҺСҮ
def _read_simple_secret_file(path: str) -> dict:
    data = {}
    try:
        with open(path, "r", encoding="utf-8") as f:
            for raw in f.read().splitlines():
                line = raw.strip()
                if not line or line.startswith("#") or "=" not in line:
                    continue
                k, v = line.split("=", 1)
                data[k.strip()] = v.strip().strip('"').strip("'")
    except Exception:
        pass
    return data

_YOO_SECRET_FILE_VALUES = {}
for _p in (
    "/etc/secrets/yookassa.env",
    "/etc/secrets/yookassa.txt",
    "/etc/secrets/yk.env",
    "/etc/secrets/yk.txt",
):
    _YOO_SECRET_FILE_VALUES.update(_read_simple_secret_file(_p))

def _secret_value(*names: str) -> str:
    for name in names:
        val = (os.environ.get(name) or "").strip()
        if val:
            return val
    for name in names:
        val = (_YOO_SECRET_FILE_VALUES.get(name) or "").strip()
        if val:
            return val
    return ""

YOO_SHOP_ID = _secret_value("YOO_SHOP_ID", "YOOKASSA_SHOP_ID", "YK_ID", "PAY_ID")
YOO_SECRET_KEY = _secret_value("YOO_SECRET_KEY", "YOOKASSA_SECRET_KEY", "YK_KEY", "PAY_KEY")
YOO_DIRECT_ENABLED = os.environ.get("YOO_DIRECT_ENABLED", "1").lower() in ("1", "true", "yes", "on")
YOO_PAYMENT_RETURN_URL = os.environ.get("YOO_PAYMENT_RETURN_URL", "").strip() or f"{PUBLIC_URL.rstrip('/')}/payment_return"
YOO_SBP_ENABLED = os.environ.get("YOO_SBP_ENABLED", "1").lower() in ("1", "true", "yes", "on")
YOO_SBERPAY_ENABLED = os.environ.get("YOO_SBERPAY_ENABLED", "1").lower() in ("1", "true", "yes", "on")
YOO_TPAY_ENABLED = os.environ.get("YOO_TPAY_ENABLED", "1").lower() in ("1", "true", "yes", "on")
YOO_MIRPAY_ENABLED = os.environ.get("YOO_MIRPAY_ENABLED", "1").lower() in ("1", "true", "yes", "on")
YOO_CARD_ENABLED = os.environ.get("YOO_CARD_ENABLED", "1").lower() in ("1", "true", "yes", "on")
YOO_PAYMENT_POLL_SECONDS = int(os.environ.get("YOO_PAYMENT_POLL_SECONDS", "900") or 900)
YOO_PAYMENT_POLL_INTERVAL_S = float(os.environ.get("YOO_PAYMENT_POLL_INTERVAL_S", "5") or 5)
YOO_DEFAULT_METHOD = (os.environ.get("YOO_DEFAULT_METHOD", "sbp") or "sbp").strip().lower()
# Р•СҒР»Рё РІ Р®Kassa РІРәР»СҺСҮРөРҪР° РҫРҪР»Р°Р№РҪ-РәР°СҒСҒР°/Р°РІСӮРҫРҫСӮРҝСҖР°РІРәР° СҮРөРәРҫРІ, API РјРҫР¶РөСӮ СӮСҖРөРұРҫРІР°СӮСҢ receipt.
# РӯСӮРё Р·РҪР°СҮРөРҪРёСҸ РјРҫР¶РҪРҫ РҝРҫР»РҫР¶РёСӮСҢ РІ Render Secret File yookassa.env РұРөР· РҙРҫРұР°РІР»РөРҪРёСҸ ENV-РҝРөСҖРөРјРөРҪРҪСӢС…:
# YK_RECEIPT_EMAIL=owner@example.com
# YK_VAT_CODE=1
YOO_RECEIPT_EMAIL = _secret_value("YOO_RECEIPT_EMAIL", "YOOKASSA_RECEIPT_EMAIL", "YK_RECEIPT_EMAIL", "PAY_RECEIPT_EMAIL")
YOO_VAT_CODE = int((_secret_value("YOO_VAT_CODE", "YK_VAT_CODE") or os.environ.get("YOO_VAT_CODE") or "1").strip() or 1)
YOO_DEBUG_PAY_ERRORS = os.environ.get("YOO_DEBUG_PAY_ERRORS", "1").lower() in ("1", "true", "yes", "on")

CRYPTO_PAY_API_TOKEN = os.environ.get("CRYPTO_PAY_API_TOKEN", "").strip()
CRYPTO_ASSET = (os.environ.get("CRYPTO_ASSET") or "USDT").upper()

# === COMPAT with existing vars/DB in your main.py ===
# 1) Р®Kassa: РөСҒР»Рё СғР¶Рө РөСҒСӮСҢ PROVIDER_TOKEN (РёР· PROVIDER_TOKEN_YOOKASSA), РёСҒРҝРҫР»СҢР·СғРөРј РөРіРҫ:
if not YOOKASSA_PROVIDER_TOKEN and 'PROVIDER_TOKEN' in globals() and PROVIDER_TOKEN:
    YOOKASSA_PROVIDER_TOKEN = PROVIDER_TOKEN

# 2) РҡРҫСҲРөР»С‘Рә: РёСҒРҝРҫР»СҢР·СғРөРј СӮРІРҫР№ РөРҙРёРҪСӢР№ USD-РәРҫСҲРөР»С‘Рә (wallet table) РІРјРөСҒСӮРҫ kv:
def _user_balance_get(user_id: int) -> float:
    return _wallet_total_get(user_id)

def _user_balance_add(user_id: int, delta: float) -> float:
    if delta > 0:
        _wallet_total_add(user_id, delta)
    elif delta < 0:
        _wallet_total_take(user_id, -delta)
    return _wallet_total_get(user_id)

def _user_balance_debit(user_id: int, amount: float) -> bool:
    return _wallet_total_take(user_id, amount)

# 3) РҹРҫРҙРҝРёСҒРәР°: Р°РәСӮРёРІРёСҖСғРөРј СҮРөСҖРөР· СӮРІРҫРё С„СғРҪРәСҶРёРё СҒ Р‘Р”, Р° РҪРө kv:
def _sub_activate(user_id: int, tier_key: str, months: int = 1) -> str:
    dt = activate_subscription_with_tier(user_id, tier_key, months)
    return dt.isoformat()

def _sub_info_text(user_id: int) -> str:
    tier = get_subscription_tier(user_id)
    dt = get_subscription_until(user_id)
    human_until = dt.strftime("%d.%m.%Y") if dt else ""
    bal = _user_balance_get(user_id)
    line_until = f"\nвҸі РҗРәСӮРёРІРҪР° РҙРҫ: {human_until}" if tier != "free" and human_until else ""
    return f"рҹ§ҫ РўРөРәСғСүР°СҸ РҝРҫРҙРҝРёСҒРәР°: {tier.upper() if tier!='free' else 'РҪРөСӮ'}{line_until}\nрҹӘҷ Р‘Р°Р»Р°РҪСҒ: {_credits_fmt_from_usd(bal)}"

# РҰРөРҪСӢ вҖ” РёР· env СҒ РҫСҒРјСӢСҒР»РөРҪРҪСӢРјРё РҙРөС„РҫР»СӮР°РјРё
PRICE_START_RUB = int(os.environ.get("PRICE_START_RUB", "599"))
PRICE_PRO_RUB = int(os.environ.get("PRICE_PRO_RUB", "1990"))
PRICE_ULT_RUB = int(os.environ.get("PRICE_ULT_RUB", "4990"))

PRICE_START_USD = _env_float("PRICE_START_USD", 5.99)
PRICE_PRO_USD   = _env_float("PRICE_PRO_USD", 19.90)
PRICE_ULT_USD   = _env_float("PRICE_ULT_USD", 49.90)

SUBS_TIERS = {
    "start": {
        "title": "START",
        "rub": PRICE_START_RUB,
        "usd": PRICE_START_USD,
        "credits": SUBSCRIPTION_CREDITS.get("start", 200),
        "features": [
            "рҹ’¬ GPT-СҮР°СӮ, РҙРҫРәСғРјРөРҪСӮСӢ Рё РҝРөСҖРөРІРҫРҙСҮРёРә",
            "рҹҺ§ РһР·РІСғСҮРәР° РҫСӮРІРөСӮРҫРІ Рё СҖР°СҒРҝРҫР·РҪР°РІР°РҪРёРө СҖРөСҮРё",
            "рҹ–ј РӨРҫСӮРҫ-РјР°СҒСӮРөСҖСҒРәР°СҸ, РәР°СҖСӮРёРҪРәРё Рё AI-С„СғРҪРәСҶРёРё вҖ” СҮРөСҖРөР· РәСҖРөРҙРёСӮСӢ",
            "рҹҺ¬ Р’РёРҙРөРҫ, РјСғР·СӢРәР° Рё Р°РІР°СӮР°СҖСӢ РјРҫР¶РҪРҫ Р·Р°РҝСғСҒРәР°СӮСҢ СҮРөСҖРөР· РәСҖРөРҙРёСӮСӢ",
            "вҡҷпёҸ РһРұСӢСҮРҪР°СҸ РҫСҮРөСҖРөРҙСҢ",
            f"рҹӘҷ {SUBSCRIPTION_CREDITS.get('start', 200)} РәСҖРөРҙРёСӮРҫРІ РәР°Р¶РҙСӢР№ РјРөСҒСҸСҶ",
        ],
    },
    "pro": {
        "title": "PRO",
        "rub": PRICE_PRO_RUB,
        "usd": PRICE_PRO_USD,
        "credits": SUBSCRIPTION_CREDITS.get("pro", 1200),
        "features": [
            "рҹ’¬ РҹРҫРІСӢСҲРөРҪРҪСӢРө GPT-Р»РёРјРёСӮСӢ",
            "рҹ“ҡ Р“Р»СғРұРҫРәРёР№ СҖР°Р·РұРҫСҖ PDF/DOCX/EPUB",
            "рҹҺ¬ Reels/Shorts, С„РҫСӮРҫвҶ’РІРёРҙРөРҫ, Suno, AI-С„РҫСӮРҫ вҖ” СҮРөСҖРөР· РәСҖРөРҙРёСӮСӢ",
            "рҹ–ј Outpaint, С„РҫРҪ, FaceSwap, AI-СҒРөР»С„Рё",
            "рҹҺӨ Р“РҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ Рё РәР»РёРҝ СҒ РІРҫРәР°Р»РҫРј РҙРҫСҒСӮСғРҝРҪСӢ СҮРөСҖРөР· РәСҖРөРҙРёСӮСӢ",
            "вҡЎ РЈСҒРәРҫСҖРөРҪРҪР°СҸ РҫСҮРөСҖРөРҙСҢ",
            f"рҹӘҷ {SUBSCRIPTION_CREDITS.get('pro', 1200)} РәСҖРөРҙРёСӮРҫРІ РәР°Р¶РҙСӢР№ РјРөСҒСҸСҶ",
        ],
    },
    "ultimate": {
        "title": "ULTIMATE",
        "rub": PRICE_ULT_RUB,
        "usd": PRICE_ULT_USD,
        "credits": SUBSCRIPTION_CREDITS.get("ultimate", 3500),
        "features": [
            "рҹ’¬ РңР°РәСҒРёРјР°Р»СҢРҪСӢРө GPT-Р»РёРјРёСӮСӢ",
            "рҹҡҖ Р‘РҫР»СҢСҲРө РәСҖРөРҙРёСӮРҫРІ РҙР»СҸ Runway / Kling / Sora 2 / Suno",
            "рҹҺӨ РҡР»РёРҝ СҒ РІРҫРәР°Р»РҫРј, РіРҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ, РІРёРҙРөРҫ РҝРҫ СӮРөРәСҒСӮСғ/РіРҫР»РҫСҒСғ",
            "рҹҺ¬ Р‘РҫР»СҢСҲРө РІРёРҙРөРҫ Рё РҝСҖРөРјРёСғРј-СҖРөРҪРҙРөСҖРҫРІ",
            "рҹ§  РҹСҖРёРҫСҖРёСӮРөСӮРҪР°СҸ РҫСҮРөСҖРөРҙСҢ",
            "рҹӣ  PRO-РёРҪСҒСӮСҖСғРјРөРҪСӮСӢ Рё СҖР°СҒСҲРёСҖРөРҪРҪСӢРө СҒСҶРөРҪР°СҖРёРё",
            f"рҹӘҷ {SUBSCRIPTION_CREDITS.get('ultimate', 3500)} РәСҖРөРҙРёСӮРҫРІ РәР°Р¶РҙСӢР№ РјРөСҒСҸСҶ",
        ],
    },
}

def _money_fmt_rub(v: int) -> str:
    return f"{v:,}".replace(",", " ") + " вӮҪ"

def _money_fmt_usd(v: float) -> str:
    return f"${v:.2f}"

def _user_balance_get(user_id: int) -> float:
    return _wallet_total_get(user_id)

def _user_balance_add(user_id: int, delta: float) -> float:
    if delta >= 0:
        _wallet_total_add(user_id, delta)
    else:
        _wallet_total_take(user_id, -delta)
    return _wallet_total_get(user_id)

def _user_balance_debit(user_id: int, amount: float) -> bool:
    return _wallet_total_take(user_id, amount)

def _sub_activate(user_id: int, tier_key: str, months: int = 1) -> str:
    until_dt = activate_subscription_with_tier(user_id, tier_key, months)
    kv_set(f"sub:tier:{user_id}", tier_key)
    kv_set(f"sub:until:{user_id}", until_dt.isoformat())
    return until_dt.isoformat()

def _sub_info_text(user_id: int) -> str:
    tier = get_subscription_tier(user_id)
    until_dt = get_subscription_until(user_id)
    human_until = until_dt.strftime("%d.%m.%Y") if until_dt else ""
    bal = _user_balance_get(user_id)
    line_until = f"\nвҸі РҗРәСӮРёРІРҪР° РҙРҫ: {human_until}" if tier != "free" and human_until else ""
    title = tier.upper() if tier != "free" else "РҪРөСӮ"
    return f"рҹ§ҫ РўРөРәСғСүР°СҸ РҝРҫРҙРҝРёСҒРәР°: {title}{line_until}\nрҹӘҷ Р‘Р°Р»Р°РҪСҒ: {_credits_fmt_from_usd(bal)}"

def _plan_card_text(key: str) -> str:
    p = SUBS_TIERS[key]
    fs = "\n".join("вҖў " + f for f in p["features"])
    return (
        f"вӯҗ РўР°СҖРёС„ {p['title']}\n"
        f"РҰРөРҪР°: {_money_fmt_rub(p['rub'])} РІ РјРөСҒСҸСҶ.\n"
        f"Р’РәР»СҺСҮРөРҪРҫ: {p.get('credits', 0)} РәСҖРөРҙРёСӮРҫРІ РҪР° РҝР»Р°СӮРҪСӢРө С„СғРҪРәСҶРёРё.\n\n"
        f"{fs}\n"
    )

def _plans_overview_text(user_id: int) -> str:
    parts = [
        "вӯҗ РҹРҫРҙРҝРёСҒРәР° Рё СӮР°СҖРёС„СӢ",
        "Р’СҒРө СӮР°СҖРёС„СӢ РҫСӮРәСҖСӢРІР°СҺСӮ РҙРҫСҒСӮСғРҝ Рә AI-С„СғРҪРәСҶРёСҸРј Neyro-Bot GPT 5 Studio. 1 РәСҖРөРҙРёСӮ = 1 вӮҪ. РҰРөРҪСӢ РіРөРҪРөСҖР°СҶРёР№ СғР¶Рө РІРәР»СҺСҮР°СҺСӮ РјР°СҖР¶Сғ СҒРөСҖРІРёСҒР°; РәСҖРөРҙРёСӮСӢ РҫРәРҫРҪСҮР°СӮРөР»СҢРҪРҫ СҒРҝРёСҒСӢРІР°СҺСӮСҒСҸ СӮРҫР»СҢРәРҫ РҝРҫСҒР»Рө СғСҒРҝРөСҲРҪРҫР№ РІСӢРҙР°СҮРё СҖРөР·СғР»СҢСӮР°СӮР°. Р§РөРј РІСӢСҲРө СӮР°СҖРёС„ вҖ” СӮРөРј РұРҫР»СҢСҲРө РәСҖРөРҙРёСӮРҫРІ, РІСӢСҲРө GPT-Р»РёРјРёСӮСӢ Рё РұСӢСҒСӮСҖРөРө РҫСҮРөСҖРөРҙСҢ.",
        _sub_info_text(user_id),
        "вҖ” вҖ” вҖ”",
        _plan_card_text("start"),
        _plan_card_text("pro"),
        _plan_card_text("ultimate"),
        "Р’СӢРұРөСҖРёСӮРө СӮР°СҖРёС„ РәРҪРҫРҝРәРҫР№ РҪРёР¶Рө.",
    ]
    return "\n".join(parts)

def plans_root_kb() -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup([
        [
            InlineKeyboardButton("вӯҗ START",    callback_data="plan:start"),
            InlineKeyboardButton("рҹҡҖ PRO",      callback_data="plan:pro"),
            InlineKeyboardButton("рҹ‘‘ ULTIMATE", callback_data="plan:ultimate"),
        ],
        [InlineKeyboardButton("рҹӘҷ РҡСғРҝРёСӮСҢ РәСҖРөРҙРёСӮСӢ", callback_data="topup")],
    ])

def plan_pay_kb(plan_key: str) -> InlineKeyboardMarkup:
    rows = []
    if YOO_SBP_ENABLED:
        rows.append([InlineKeyboardButton("вҡЎ РЎР‘Рҹ / QR", callback_data=f"pay:yoo_sbp:{plan_key}")])
    app_row = []
    if YOO_SBERPAY_ENABLED:
        app_row.append(InlineKeyboardButton("рҹҹў SberPay", callback_data=f"pay:yoo_sberpay:{plan_key}"))
    if YOO_TPAY_ENABLED:
        app_row.append(InlineKeyboardButton("рҹҹЎ T-Pay", callback_data=f"pay:yoo_tpay:{plan_key}"))
    if app_row:
        rows.append(app_row)
    app_row2 = []
    if YOO_MIRPAY_ENABLED:
        app_row2.append(InlineKeyboardButton("рҹ’ҷ Mir Pay", callback_data=f"pay:yoo_mirpay:{plan_key}"))
    if YOO_CARD_ENABLED:
        app_row2.append(InlineKeyboardButton("рҹ’і РҡР°СҖСӮР° Telegram", callback_data=f"pay:yookassa:{plan_key}"))
    if app_row2:
        rows.append(app_row2)
    rows.append([InlineKeyboardButton("рҹҢҗ Р’СҒРө СҒРҝРҫСҒРҫРұСӢ Р®Kassa", callback_data=f"pay:yoo_all:{plan_key}")])
    rows.append([InlineKeyboardButton("рҹ’  CryptoBot", callback_data=f"pay:cryptobot:{plan_key}")])
    rows.append([InlineKeyboardButton("рҹ§ҫ РЎ РұР°Р»Р°РҪСҒР°", callback_data=f"pay:balance:{plan_key}")])
    rows.append([InlineKeyboardButton("в¬…пёҸ Рҡ СӮР°СҖРёС„Р°Рј", callback_data="plan:root")])
    return InlineKeyboardMarkup(rows)


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ Р®Kassa direct API: РЎР‘Рҹ / QR, SberPay, T-Pay, Mir Pay в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
YOO_DIRECT_METHODS = {
    "yoo_sbp":      {"type": "sbp",          "confirmation": "redirect", "label": "вҡЎ РЎР‘Рҹ / QR"},
    "yoo_sberpay": {"type": "sberbank",     "confirmation": "redirect", "label": "рҹҹў SberPay"},
    "yoo_tpay":    {"type": "tinkoff_bank", "confirmation": "redirect", "label": "рҹҹЎ T-Pay"},
    # Р’ API Р®Kassa Mir Pay РҝСҖРҫС…РҫРҙРёСӮ РәР°Рә bank_card; РҪР° С„РҫСҖРјРө/РјРҫРұ. СғСҒСӮСҖРҫР№СҒСӮРІРө РёСҒСӮРҫСҮРҪРёРә РұСғРҙРөСӮ mir_pay.
    "yoo_mirpay":  {"type": "bank_card",    "confirmation": "redirect", "label": "рҹ’ҷ Mir Pay"},
    # РЈРҪРёРІРөСҖСҒР°Р»СҢРҪР°СҸ СҒСҒСӢР»РәР° Р®Kassa: РұРөР· Р¶РөСҒСӮРәРҫРіРҫ payment_method_data, С„РҫСҖРјР° РҝРҫРәР°Р¶РөСӮ РҙРҫСҒСӮСғРҝРҪСӢРө РјРөСӮРҫРҙСӢ РјР°РіР°Р·РёРҪР°.
    "yoo_all":     {"type": None,           "confirmation": "redirect", "label": "рҹҢҗ Р’СҒРө СҒРҝРҫСҒРҫРұСӢ Р®Kassa"},
}

def _yoo_direct_configured() -> bool:
    return bool(YOO_DIRECT_ENABLED and YOO_SHOP_ID and YOO_SECRET_KEY)

def _yoo_auth():
    return (YOO_SHOP_ID, YOO_SECRET_KEY)

async def _yoo_create_direct_payment(user_id: int, plan_key: str, months: int, method_key: str) -> dict:
    """РЎРҫР·РҙР°С‘СӮ direct payment РІ Р®Kassa Рё РІРҫР·РІСҖР°СүР°РөСӮ JSON РҝР»Р°СӮРөР¶Р°."""
    plan = SUBS_TIERS[plan_key]
    method = YOO_DIRECT_METHODS[method_key]
    amount_rub = int(plan["rub"]) * max(1, int(months or 1))
    title = f"РҹРҫРҙРҝРёСҒРәР° {plan['title']} вҖў {months} РјРөСҒ."
    confirmation_type = method.get("confirmation") or "redirect"
    payload = {
        "amount": {"value": f"{amount_rub:.2f}", "currency": "RUB"},
        "capture": True,
        "description": title,
        "metadata": {
            "kind": "subscription",
            "user_id": str(user_id),
            "tier": plan_key,
            "months": str(months),
            "credits": str(plan.get("credits", 0) * max(1, int(months or 1))),
            "method": method_key,
        },
        "confirmation": {"type": confirmation_type},
    }
    if method.get("type"):
        payload["payment_method_data"] = {"type": method["type"]}
    if confirmation_type in ("redirect", "external"):
        payload["confirmation"]["return_url"] = YOO_PAYMENT_RETURN_URL

    # Р•СҒР»Рё РјР°РіР°Р·РёРҪ СҖР°РұРҫСӮР°РөСӮ СҒ РҫРҪР»Р°Р№РҪ-РәР°СҒСҒРҫР№, Р®Kassa РјРҫР¶РөСӮ РҪРө СҒРҫР·РҙР°СӮСҢ РҝР»Р°СӮРөР¶ РұРөР· receipt.
    # Р”РҫРұР°РІР»СҸРөРј СҮРөРә, РөСҒР»Рё РІ secret file СғРәР°Р·Р°РҪ YK_RECEIPT_EMAIL.
    if YOO_RECEIPT_EMAIL:
        payload["receipt"] = {
            "customer": {"email": YOO_RECEIPT_EMAIL},
            "items": [{
                "description": title[:128],
                "quantity": "1.00",
                "amount": {"value": f"{amount_rub:.2f}", "currency": "RUB"},
                "vat_code": YOO_VAT_CODE,
                "payment_mode": "full_payment",
                "payment_subject": "service",
            }],
        }
    headers = {"Idempotence-Key": str(uuid.uuid4())}
    async with httpx.AsyncClient(timeout=30.0) as client:
        r = await client.post("https://api.yookassa.ru/v3/payments", auth=_yoo_auth(), headers=headers, json=payload)
        if r.status_code >= 400:
            log.error("YooKassa create payment failed status=%s method=%s body=%s", r.status_code, method_key, r.text[:1200])
            raise RuntimeError(f"YooKassa {r.status_code}: {r.text[:800]}")
        return r.json()



def _credit_pack_resolve(requested_credits: int = 0, requested_rub: int = 0) -> tuple[int, int] | None:
    """Р Р°Р·СҖРөСҲР°РөСӮ СӮРҫР»СҢРәРҫ Р·Р°С„РёРәСҒРёСҖРҫРІР°РҪРҪСӢРө РҝР°РәРөСӮСӢ; Р·РҪР°СҮРөРҪРёСҸ РёР· WebApp РҪРө СҒСҮРёСӮР°СҺСӮСҒСҸ РҙРҫРІРөСҖРөРҪРҪСӢРјРё."""
    for credits_raw, rub_raw in CREDIT_PACKAGES_RUB.items():
        credits, rub = int(credits_raw), int(rub_raw)
        if requested_credits and requested_credits == credits:
            return credits, rub
        if requested_rub and requested_rub == rub:
            return credits, rub
    return None


async def _yoo_create_credit_payment(user_id: int, credits: int, amount_rub: int, method_key: str = "yoo_all") -> dict:
    resolved = _credit_pack_resolve(int(credits), int(amount_rub))
    if not resolved:
        raise ValueError("Unknown credit package")
    credits, amount_rub = resolved
    method_key = method_key if method_key in YOO_DIRECT_METHODS else "yoo_all"
    method = YOO_DIRECT_METHODS[method_key]
    title = f"РҹРҫРҝРҫР»РҪРөРҪРёРө РұР°Р»Р°РҪСҒР°: {credits} РәСҖРөРҙРёСӮРҫРІ"
    confirmation_type = method.get("confirmation") or "redirect"
    payload = {
        "amount": {"value": f"{amount_rub:.2f}", "currency": "RUB"},
        "capture": True,
        "description": title,
        "metadata": {
            "kind": "credit_topup",
            "user_id": str(user_id),
            "credits": str(credits),
            "amount_rub": str(amount_rub),
            "method": method_key,
        },
        "confirmation": {"type": confirmation_type},
    }
    if method.get("type"):
        payload["payment_method_data"] = {"type": method["type"]}
    if confirmation_type in ("redirect", "external"):
        payload["confirmation"]["return_url"] = YOO_PAYMENT_RETURN_URL
    if YOO_RECEIPT_EMAIL:
        payload["receipt"] = {
            "customer": {"email": YOO_RECEIPT_EMAIL},
            "items": [{
                "description": title[:128],
                "quantity": "1.00",
                "amount": {"value": f"{amount_rub:.2f}", "currency": "RUB"},
                "vat_code": YOO_VAT_CODE,
                "payment_mode": "full_payment",
                "payment_subject": "service",
            }],
        }
    headers = {"Idempotence-Key": str(uuid.uuid4())}
    async with httpx.AsyncClient(timeout=30.0) as client:
        r = await client.post("https://api.yookassa.ru/v3/payments", auth=_yoo_auth(), headers=headers, json=payload)
        if r.status_code >= 400:
            log.error("YooKassa credit payment failed status=%s body=%s", r.status_code, r.text[:1200])
            raise RuntimeError(f"YooKassa {r.status_code}: {r.text[:800]}")
        return r.json()


async def _poll_yoo_credit_payment(
    context: ContextTypes.DEFAULT_TYPE,
    chat_id: int,
    message_id: int,
    user_id: int,
    payment_id: str,
    credits: int,
    amount_rub: int,
):
    deadline = time.time() + max(60, YOO_PAYMENT_POLL_SECONDS)
    try:
        while time.time() < deadline:
            p = await _yoo_get_payment(payment_id)
            st = (p or {}).get("status", "").lower()
            if st == "succeeded":
                paid_key = f"yoo:credit_paid:{payment_id}"
                if kv_get(paid_key) != "1":
                    _wallet_total_add(user_id, _credits_to_usd(credits))
                    kv_set(paid_key, "1")
                with contextlib.suppress(Exception):
                    await context.bot.edit_message_text(
                        chat_id=chat_id,
                        message_id=message_id,
                        text=(
                            "вң… РһРҝР»Р°СӮР° Р®Kassa РҝРҫРҙСӮРІРөСҖР¶РҙРөРҪР°.\n"
                            f"РқР°СҮРёСҒР»РөРҪРҫ: {credits} РәСҖРөРҙРёСӮРҫРІ Р·Р° {amount_rub} вӮҪ.\n"
                            f"РўРөРәСғСүРёР№ РұР°Р»Р°РҪСҒ: {_credits_fmt_from_usd(_user_balance_get(user_id))}."
                        ),
                    )
                return
            if st in ("canceled", "cancelled", "failed"):
                with contextlib.suppress(Exception):
                    await context.bot.edit_message_text(
                        chat_id=chat_id,
                        message_id=message_id,
                        text=f"вқҢ РһРҝР»Р°СӮР° РҝР°РәРөСӮР° РәСҖРөРҙРёСӮРҫРІ РҪРө Р·Р°РІРөСҖСҲРөРҪР°. РЎСӮР°СӮСғСҒ: {st}.",
                    )
                return
            await asyncio.sleep(max(2.0, YOO_PAYMENT_POLL_INTERVAL_S))
        with contextlib.suppress(Exception):
            await context.bot.edit_message_text(
                chat_id=chat_id,
                message_id=message_id,
                text="вҢӣ Р’СҖРөРјСҸ РҫР¶РёРҙР°РҪРёСҸ РҫРҝР»Р°СӮСӢ РІСӢСҲР»Рҫ. Р•СҒР»Рё РҝР»Р°СӮС‘Р¶ Р·Р°РІРөСҖСҲС‘РҪ вҖ” РҫРұСҖР°СӮРёСӮРөСҒСҢ РІ РҝРҫРҙРҙРөСҖР¶РәСғ Рё СғРәР°Р¶РёСӮРө РІСҖРөРјСҸ РҫРҝР»Р°СӮСӢ.",
            )
    except Exception as e:
        log.exception("YooKassa credit poll error: %s", e)

async def _yoo_get_payment(payment_id: str) -> dict | None:
    if not payment_id or not _yoo_direct_configured():
        return None
    try:
        async with httpx.AsyncClient(timeout=20.0) as client:
            r = await client.get(f"https://api.yookassa.ru/v3/payments/{payment_id}", auth=_yoo_auth())
            if r.status_code >= 400:
                log.warning("YooKassa get payment failed %s: %s", r.status_code, r.text[:300])
                return None
            return r.json()
    except Exception as e:
        log.exception("YooKassa get payment error: %s", e)
        return None

async def _poll_yoo_subscription_payment(context: ContextTypes.DEFAULT_TYPE, chat_id: int, message_id: int, user_id: int, payment_id: str, plan_key: str, months: int):
    """РһР¶РёРҙР°РөСӮ РҫРҝР»Р°СӮСғ direct YooKassa Рё Р°РәСӮРёРІРёСҖСғРөСӮ СӮР°СҖРёС„."""
    deadline = time.time() + max(60, YOO_PAYMENT_POLL_SECONDS)
    try:
        while time.time() < deadline:
            p = await _yoo_get_payment(payment_id)
            st = (p or {}).get("status", "").lower()
            if st == "succeeded":
                until = activate_subscription_with_tier(user_id, plan_key, months)
                credits = SUBSCRIPTION_CREDITS.get((plan_key or "").lower(), 0) * int(months)
                kv_set(f"yoo:paid:{payment_id}", "1")
                with contextlib.suppress(Exception):
                    await context.bot.edit_message_text(
                        chat_id=chat_id,
                        message_id=message_id,
                        text=(
                            "вң… РһРҝР»Р°СӮР° Р®Kassa РҝРҫРҙСӮРІРөСҖР¶РҙРөРҪР°.\n"
                            f"РҹРҫРҙРҝРёСҒРәР° {plan_key.upper()} Р°РәСӮРёРІРҪР° РҙРҫ {until.strftime('%Y-%m-%d')}.\n"
                            f"рҹӘҷ РҡСҖРөРҙРёСӮСӢ РҪР°СҮРёСҒР»РөРҪСӢ: {credits} РәСҖ."
                        ),
                    )
                return
            if st in ("canceled", "cancelled", "failed"):
                with contextlib.suppress(Exception):
                    await context.bot.edit_message_text(chat_id=chat_id, message_id=message_id, text=f"вқҢ РһРҝР»Р°СӮР° РҪРө Р·Р°РІРөСҖСҲРөРҪР°. РЎСӮР°СӮСғСҒ: {st}.")
                return
            await asyncio.sleep(max(2.0, YOO_PAYMENT_POLL_INTERVAL_S))
        with contextlib.suppress(Exception):
            await context.bot.edit_message_text(
                chat_id=chat_id,
                message_id=message_id,
                text="вҢӣ Р’СҖРөРјСҸ РҫР¶РёРҙР°РҪРёСҸ РҫРҝР»Р°СӮСӢ РІСӢСҲР»Рҫ. Р•СҒР»Рё РІСӢ РҫРҝР»Р°СӮРёР»Рё вҖ” РҪР°Р¶РјРёСӮРө В«рҹ”Һ РҹСҖРҫРІРөСҖРёСӮСҢ РҫРҝР»Р°СӮСғВ» РІ РҪРҫРІРҫРј СҒСҮС‘СӮРө РёР»Рё РҪР°РҝРёСҲРёСӮРө РІ РҝРҫРҙРҙРөСҖР¶РәСғ.",
            )
    except Exception as e:
        log.exception("YooKassa subscription poll error: %s", e)

# РҡРҪРҫРҝРәР° В«вӯҗ РҹРҫРҙРҝРёСҒРәР° В· РҹРҫРјРҫСүСҢВ»
async def on_btn_plans(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    text = _plans_overview_text(user_id)
    await update.effective_chat.send_message(text, reply_markup=plans_root_kb())

# РһРұСҖР°РұРҫСӮСҮРёРә РҪР°СҲРёС… РәРҫР»РұСҚРәРҫРІ РҝРҫ РҝРҫРҙРҝРёСҒРәРө/РҫРҝР»Р°СӮР°Рј (Р·Р°СҖРөРіРёСҒСӮСҖРёСҖРҫРІР°СӮСҢ Р”Рһ РҫРұСүРөРіРҫ on_cb!)
async def on_cb_plans(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    data = q.data or ""
    user_id = q.from_user.id
    chat_id = q.message.chat.id  # FIX: РәРҫСҖСҖРөРәСӮРҪРҫРө РҝРҫР»Рө РІ PTB v21+

    # РқР°РІРёРіР°СҶРёСҸ РјРөР¶РҙСғ СӮР°СҖРёС„Р°РјРё
    if data.startswith("plan:"):
        _, arg = data.split(":", 1)
        if arg == "root":
            await q.message.reply_text(_plans_overview_text(user_id), reply_markup=plans_root_kb())
            await q.answer()
            return
        if arg in SUBS_TIERS:
            await q.message.reply_text(
                _plan_card_text(arg) + "\nР’СӢРұРөСҖРёСӮРө СҒРҝРҫСҒРҫРұ РҫРҝР»Р°СӮСӢ:",
                reply_markup=plan_pay_kb(arg)
            )
            await q.answer()
            return

    # РҹР»Р°СӮРөР¶Рё
    if data.startswith("pay:"):
        # РұРөР·РҫРҝР°СҒРҪСӢР№ РҝР°СҖСҒРёРҪРі
        try:
            _, method, plan_key = data.split(":", 2)
        except ValueError:
            await q.answer("РқРөРәРҫСҖСҖРөРәСӮРҪСӢРө РҙР°РҪРҪСӢРө РәРҪРҫРҝРәРё.", show_alert=True)
            return

        plan = SUBS_TIERS.get(plan_key)
        if not plan:
            await q.answer("РқРөРёР·РІРөСҒСӮРҪСӢР№ СӮР°СҖРёС„.", show_alert=True)
            return

        # Р®Kassa direct API: РЎР‘Рҹ/QR, SberPay, T-Pay, Mir Pay
        if method in YOO_DIRECT_METHODS:
            await q.answer("РЎРҫР·РҙР°СҺ СҒСҒСӢР»РәСғ РҪР° РҫРҝР»Р°СӮСғвҖҰ")
            if not _yoo_direct_configured():
                await q.message.reply_text(
                    "вҡ пёҸ Р‘СӢСҒСӮСҖР°СҸ РҫРҝР»Р°СӮР° Р®Kassa РҝРҫРәР° РҪРө РҪР°СҒСӮСҖРҫРөРҪР°: РҪСғР¶РҪСӢ YOO_SHOP_ID/YOO_SECRET_KEY РёР»Рё Secret File yookassa.env СҒ YK_ID/YK_KEY."
                )
                return
            try:
                pay = await _yoo_create_direct_payment(user_id, plan_key, 1, method)
                payment_id = str(pay.get("id") or "")
                conf = pay.get("confirmation") or {}
                pay_url = conf.get("confirmation_url") or conf.get("confirmation_data") or conf.get("external_url") or ""
                label = YOO_DIRECT_METHODS[method]["label"]
                if not pay_url:
                    raise RuntimeError(f"YooKassa did not return confirmation url: {pay}")
                kb = InlineKeyboardMarkup([
                    [InlineKeyboardButton(f"{label} вҖ” РҫРҝР»Р°СӮРёСӮСҢ", url=pay_url)],
                    [InlineKeyboardButton("в¬…пёҸ Рҡ СӮР°СҖРёС„Сғ", callback_data=f"plan:{plan_key}")],
                ])
                msg = await q.message.reply_text(
                    _plan_card_text(plan_key)
                    + f"\nРЎРҝРҫСҒРҫРұ РҫРҝР»Р°СӮСӢ: {label}\n"
                    + "РһСӮРәСҖРҫР№СӮРө СҒСҒСӢР»РәСғ. РҹРҫСҒР»Рө РҫРҝР»Р°СӮСӢ РұРҫСӮ Р°РәСӮРёРІРёСҖСғРөСӮ РҝРҫРҙРҝРёСҒРәСғ Р°РІСӮРҫРјР°СӮРёСҮРөСҒРәРё.",
                    reply_markup=kb,
                )
                kv_set(f"yoo:pending:{payment_id}", json.dumps({"user_id": user_id, "tier": plan_key, "months": 1, "method": method}, ensure_ascii=False))
                context.application.create_task(_poll_yoo_subscription_payment(
                    context, msg.chat.id, msg.message_id, user_id, payment_id, plan_key, 1
                ))
            except Exception as e:
                log.exception("YooKassa direct payment create failed: %s", e)
                err = str(e)[:700]
                user_msg = "вҡ пёҸ РқРө СғРҙР°Р»РҫСҒСҢ СҒРҫР·РҙР°СӮСҢ РұСӢСҒСӮСҖСғСҺ РҫРҝР»Р°СӮСғ Р®Kassa. РҹРҫРҝСҖРҫРұСғР№СӮРө РәР°СҖСӮСғ Telegram, CryptoBot РёР»Рё РҝРҫР·Р¶Рө."
                if YOO_DEBUG_PAY_ERRORS:
                    user_msg += "\n\nР”РёР°РіРҪРҫСҒСӮРёРәР° Р®Kassa: " + err
                await q.message.reply_text(user_msg)
            return

        # Р®Kassa СҮРөСҖРөР· Telegram Payments / РұР°РҪРәРҫРІСҒРәР°СҸ РәР°СҖСӮР°
        if method == "yookassa":
            if not YOOKASSA_PROVIDER_TOKEN:
                await q.answer("Р®Kassa РҪРө РҝРҫРҙРәР»СҺСҮРөРҪР° (РҪРөСӮ YOOKASSA_PROVIDER_TOKEN).", show_alert=True)
                return

            title = f"РҹРҫРҙРҝРёСҒРәР° {plan['title']} вҖў 1 РјРөСҒСҸСҶ"
            desc = "Р”РҫСҒСӮСғРҝ Рә С„СғРҪРәСҶРёСҸРј РұРҫСӮР° СҒРҫРіР»Р°СҒРҪРҫ РІСӢРұСҖР°РҪРҪРҫРјСғ СӮР°СҖРёС„Сғ. РҹРҫРҙРҝРёСҒРәР° Р°РәСӮРёРІРёСҖСғРөСӮСҒСҸ СҒСҖР°Р·Сғ РҝРҫСҒР»Рө РҫРҝР»Р°СӮСӢ."
            payload = json.dumps({"tier": plan_key, "months": 1})

            # Telegram РҫР¶РёРҙР°РөСӮ СҒСғРјРјСғ РІ РјРёРҪРҫСҖРҪСӢС… РөРҙРёРҪРёСҶР°С… (РәРҫРҝРөР№РәРё/СҶРөРҪСӮСӢ)
            if YOOKASSA_CURRENCY == "RUB":
                total_minor = int(round(float(plan["rub"]) * 100))
            else:
                total_minor = int(round(float(plan["usd"]) * 100))

            prices = [LabeledPrice(label=f"{plan['title']} 1 РјРөСҒ.", amount=total_minor)]
            await context.bot.send_invoice(
                chat_id=chat_id,
                title=title,
                description=desc,
                payload=payload,
                provider_token=YOOKASSA_PROVIDER_TOKEN,
                currency=YOOKASSA_CURRENCY,
                prices=prices,
                need_email=True,
                is_flexible=False,
            )
            await q.answer("РЎСҮС‘СӮ РІСӢСҒСӮР°РІР»РөРҪ вң…")
            return

        # CryptoBot (Crypto Pay API: СҒРҫР·РҙР°С‘Рј РёРҪРІРҫР№СҒ Рё РҫСӮРҙР°С‘Рј СҒСҒСӢР»РәСғ)
        if method == "cryptobot":  # FIX: РІСӢСҖРҫРІРҪРөРҪ РҫСӮСҒСӮСғРҝ
            if not CRYPTO_PAY_API_TOKEN:
                await q.answer("CryptoBot РҪРө РҝРҫРҙРәР»СҺСҮС‘РҪ (РҪРөСӮ CRYPTO_PAY_API_TOKEN).", show_alert=True)
                return
            try:
                amount = float(plan["usd"])
                async with httpx.AsyncClient(timeout=20) as client:
                    r = await client.post(
                        "https://pay.crypt.bot/api/createInvoice",
                        headers={"Crypto-Pay-API-Token": CRYPTO_PAY_API_TOKEN},
                        json={
                            "asset": CRYPTO_ASSET,
                            "amount": f"{amount:.2f}",
                            "description": f"Subscription {plan['title']} вҖў 1 month",
                            "allow_comments": False,
                            "allow_anonymous": True,
                        },
                    )
                    data = r.json()
                    if not data.get("ok"):
                        raise RuntimeError(str(data))
                    res = data["result"]
                    pay_url = res["pay_url"]
                    inv_id = str(res["invoice_id"])

                kb = InlineKeyboardMarkup([
                    [InlineKeyboardButton("рҹ’  CryptoBot", url=pay_url)],
                    [InlineKeyboardButton("в¬…пёҸ Рҡ СӮР°СҖРёС„Сғ", callback_data=f"plan:{plan_key}")],
                ])
                msg = await q.message.reply_text(
                    _plan_card_text(plan_key) + "\nРһСӮРәСҖРҫР№СӮРө СҒСҒСӢР»РәСғ РҙР»СҸ РҫРҝР»Р°СӮСӢ:",
                    reply_markup=kb
                )
                # Р°РІСӮРҫРҝСғР» СҒСӮР°СӮСғСҒР° РёРјРөРҪРҪРҫ РҙР»СҸ РҹРһР”РҹРҳРЎРҡРҳ
                context.application.create_task(_poll_crypto_sub_invoice(
                    context, msg.chat.id, msg.message_id, user_id, inv_id, plan_key, 1  # FIX: msg.chat.id
                ))
                await q.answer()
            except Exception as e:
                await q.answer("РқРө СғРҙР°Р»РҫСҒСҢ СҒРҫР·РҙР°СӮСҢ СҒСҮС‘СӮ РІ CryptoBot.", show_alert=True)
                log.exception("CryptoBot invoice error: %s", e)
            return

        # РЎРҝРёСҒР°РҪРёРө СҒ РІРҪСғСӮСҖРөРҪРҪРөРіРҫ РәСҖРөРҙРёСӮРҪРҫРіРҫ РұР°Р»Р°РҪСҒР°
        if method == "balance":
            price_credits = int(plan["rub"])
            price_internal = _credits_to_usd(price_credits)
            if not _user_balance_debit(user_id, price_internal):
                await q.answer("РқРөРҙРҫСҒСӮР°СӮРҫСҮРҪРҫ РәСҖРөРҙРёСӮРҫРІ РҪР° РұР°Р»Р°РҪСҒРө.", show_alert=True)
                return
            until = _sub_activate(user_id, plan_key, months=1)
            await q.message.reply_text(
                f"вң… РҹРҫРҙРҝРёСҒРәР° {plan['title']} Р°РәСӮРёРІРёСҖРҫРІР°РҪР° РҙРҫ {until[:10]}.\n"
                f"рҹӘҷ РЎРҝРёСҒР°РҪРҫ: {price_credits} РәСҖ. "
                f"РўРөРәСғСүРёР№ РұР°Р»Р°РҪСҒ: {_credits_fmt_from_usd(_user_balance_get(user_id))}",
                reply_markup=plans_root_kb(),
            )
            await q.answer()
            return

    # Р•СҒР»Рё РәРҫР»РұСҚРә РҪРө РҪР°СҲ вҖ” РҝСҖРҫРҝСғСҒРәР°РөРј РҙР°Р»СҢСҲРө
    await q.answer()
    return


# Р•СҒР»Рё Сғ СӮРөРұСҸ СғР¶Рө РөСҒСӮСҢ on_precheckout / on_successful_payment вҖ” РҫСҒСӮР°РІСҢ РёС….
# Р•СҒР»Рё РҪРөСӮ, РјРҫР¶РөСҲСҢ РёСҒРҝРҫР»СҢР·РҫРІР°СӮСҢ СҚСӮРё РҝСҖРҫСҒСӮСӢРө СҖРөР°Р»РёР·Р°СҶРёРё:

async def on_precheckout(update: Update, context: ContextTypes.DEFAULT_TYPE):
    try:
        await update.pre_checkout_query.answer(ok=True)
    except Exception as e:
        log.exception("precheckout error: %s", e)

async def on_successful_payment(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """
    РЈРҪРёРІРөСҖСҒР°Р»СҢРҪСӢР№ РҫРұСҖР°РұРҫСӮСҮРёРә Telegram Payments:
    - РҹРҫРҙРҙРөСҖР¶РёРІР°РөСӮ payload РІ РҙРІСғС… С„РҫСҖРјР°СӮР°С…:
        1) JSON: {"tier":"pro","months":1}
        2) РЎСӮСҖРҫРәР°: "sub:pro:1"
    - РҳРҪР°СҮРө СӮСҖР°РәСӮСғРөСӮ РәР°Рә РҝРҫРҝРҫР»РҪРөРҪРёРө РөРҙРёРҪРҫРіРҫ USD-РәРҫСҲРөР»СҢРәР°.
    """
    try:
        sp = update.message.successful_payment
        payload_raw = sp.invoice_payload or ""
        total_minor = sp.total_amount or 0
        rub = total_minor / 100.0
        uid = update.effective_user.id

        # 1) РҹСӢСӮР°РөРјСҒСҸ СҖР°СҒРҝР°СҖСҒРёСӮСҢ JSON
        tier, months = None, None
        try:
            if payload_raw.strip().startswith("{"):
                obj = json.loads(payload_raw)
                tier = (obj.get("tier") or "").strip().lower() or None
                months = int(obj.get("months") or 1)
        except Exception:
            pass

        # 2) РҹСӢСӮР°РөРјСҒСҸ СҖР°СҒРҝР°СҖСҒРёСӮСҢ СҒСӮСҖРҫРәРҫРІСӢР№ С„РҫСҖРјР°СӮ "sub:tier:months"
        if not tier and payload_raw.startswith("sub:"):
            try:
                _, t, m = payload_raw.split(":", 2)
                tier = (t or "pro").strip().lower()
                months = int(m or 1)
            except Exception:
                tier, months = None, None

        if tier and months:
            until = activate_subscription_with_tier(uid, tier, months)
            await update.effective_message.reply_text(
                f"рҹҺү РһРҝР»Р°СӮР° РҝСҖРҫСҲР»Р° СғСҒРҝРөСҲРҪРҫ!\n"
                f"вң… РҹРҫРҙРҝРёСҒРәР° {tier.upper()} Р°РәСӮРёРІРёСҖРҫРІР°РҪР° РҙРҫ {until.strftime('%Y-%m-%d')}.\nрҹӘҷ РҡСҖРөРҙРёСӮСӢ РҪР°СҮРёСҒР»РөРҪСӢ: {SUBSCRIPTION_CREDITS.get((tier or "").lower(), 0) * int(months)} РәСҖ."
            )
            return

        if payload_raw.startswith("topup:"):
            try:
                _, credits_s, rub_s = payload_raw.split(":", 2)
                resolved = _credit_pack_resolve(int(credits_s), int(rub_s))
                if resolved:
                    credits, expected_rub = resolved
                    _wallet_total_add(uid, _credits_to_usd(credits))
                    await update.effective_message.reply_text(
                        f"вң… РһРҝР»Р°СӮР° РҝСҖРҫСҲР»Р° СғСҒРҝРөСҲРҪРҫ. РқР°СҮРёСҒР»РөРҪРҫ: {credits} РәСҖРөРҙРёСӮРҫРІ Р·Р° {expected_rub} вӮҪ."
                    )
                    return
            except Exception:
                log.exception("Failed to parse topup payload: %s", payload_raw)

        # РҳРҪР°СҮРө СҒСҮРёСӮР°РөРј, СҮСӮРҫ СҚСӮРҫ РҝРҫРҝРҫР»РҪРөРҪРёРө РәРҫСҲРөР»СҢРәР° РІ СҖСғРұР»СҸС…
        usd = _credits_to_usd(rub)
        _wallet_total_add(uid, usd)
        await update.effective_message.reply_text(
            f"рҹ’і РҹРҫРҝРҫР»РҪРөРҪРёРө: {rub:.0f} вӮҪ. РқР°СҮРёСҒР»РөРҪРҫ: {_credits_fmt_from_usd(usd)}."
        )

    except Exception as e:
        log.exception("successful_payment handler error: %s", e)
        with contextlib.suppress(Exception):
            await update.effective_message.reply_text("вҡ пёҸ РһСҲРёРұРәР° РҫРұСҖР°РұРҫСӮРәРё РҝР»Р°СӮРөР¶Р°. Р•СҒР»Рё РҙРөРҪСҢРіРё СҒРҝРёСҒР°Р»РёСҒСҢ вҖ” РҪР°РҝРёСҲРёСӮРө РІ РҝРҫРҙРҙРөСҖР¶РәСғ.")
# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РҡРҫРҪРөСҶ PATCH в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
        
# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РҡРҫРјР°РҪРҙР° /img в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
async def cmd_img(update: Update, context: ContextTypes.DEFAULT_TYPE):
    prompt = " ".join(context.args).strip() if context.args else ""
    if not prompt:
        await update.effective_message.reply_text("РӨРҫСҖРјР°СӮ: /img <РҫРҝРёСҒР°РҪРёРө>")
        return

    async def _go():
        return await _do_img_generate(update, context, prompt)

    user_id = update.effective_user.id
    await _try_pay_then_do(
        update, context, user_id,
        "img", IMG_COST_USD, _go,
        remember_kind="img_generate", remember_payload={"prompt": prompt}
    )


async def cmd_midjourney(update: Update, context: ContextTypes.DEFAULT_TYPE):
    prompt = " ".join(context.args).strip() if context.args else ""
    await _start_midjourney_image(update, context, prompt)


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ Work/Business generators + Suno music в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def _safe_filename(name: str, ext: str) -> str:
    base = re.sub(r"[^a-zA-ZР°-СҸРҗ-РҜ0-9_.-]+", "_", (name or "file")).strip("._")[:48] or "file"
    return f"{base}.{ext.lstrip('.')}"

async def _reply_long_text(update: Update, text: str, reply_markup=None):
    text = (text or "").strip()
    if not text:
        return
    first = True
    for i in range(0, len(text), 3800):
        await update.effective_message.reply_text(text[i:i+3800], reply_markup=reply_markup if first else None)
        first = False

async def _generate_business_presentation(update: Update, context: ContextTypes.DEFAULT_TYPE, brief: str):
    prompt = (
        "РўСӢ РұРёР·РҪРөСҒ-РәРҫРҪСҒСғР»СҢСӮР°РҪСӮ Рё РҝСҖРөР·РөРҪСӮР°СҶРёРҫРҪРҪСӢР№ СҖРөРҙР°РәСӮРҫСҖ. РқР° СҖСғСҒСҒРәРҫРј СҒРҙРөР»Р°Р№ СҒСӮСҖСғРәСӮСғСҖСғ РҝСҖРөР·РөРҪСӮР°СҶРёРё РҝРҫ РўР— РҝРҫР»СҢР·РҫРІР°СӮРөР»СҸ. "
        "РӨРҫСҖРјР°СӮ СҒСӮСҖРҫРіРҫ: Р·Р°РіРҫР»РҫРІРҫРә, Р·Р°СӮРөРј 8-12 СҒР»Р°Р№РҙРҫРІ. Р”Р»СҸ РәР°Р¶РҙРҫРіРҫ СҒР»Р°Р№РҙР°: РқР°Р·РІР°РҪРёРө, 3-5 РұСғР»Р»РөСӮРҫРІ, РІРёР·СғР°Р»СҢРҪР°СҸ РёРҙРөСҸ, Р·Р°РјРөСӮРәР° СҒРҝРёРәРөСҖР°. "
        "РқРө РҝРёСҲРё Р»РёСҲРҪРёС… РІСҒСӮСғРҝР»РөРҪРёР№. РўР—:\n" + (brief or "")
    )
    reply = await ask_openai_text(prompt, user_id=update.effective_user.id, chat_id=update.effective_chat.id)
    await _reply_long_text(update, "рҹ“Ҡ Р§РөСҖРҪРҫРІРёРә РҝСҖРөР·РөРҪСӮР°СҶРёРё РіРҫСӮРҫРІ. РқРёР¶Рө СҒСӮСҖСғРәСӮСғСҖР°, РҫСӮРҙРөР»СҢРҪРҫ РҝСҖРёРәСҖРөРҝР»СҸСҺ PPTX-С„Р°Р№Р».\n\n" + reply)
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation()
        title = (brief or "РҹСҖРөР·РөРҪСӮР°СҶРёСҸ")[:80]
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = "Р§РөСҖРҪРҫРІРёРә СҒРҫР·РҙР°РҪ Neyro-Bot GPT 5 Studio"
        chunks = re.split(r"\n(?=\s*(?:РЎР»Р°Р№Рҙ\s*\d+|\d+[\).]|#{1,3}\s+))", reply)
        slide_count = 0
        for ch in chunks:
            lines = [re.sub(r"^[#\s]*", "", x).strip(" -вҖў\t") for x in ch.splitlines() if x.strip()]
            if not lines:
                continue
            if slide_count >= 12:
                break
            ttl = lines[0][:90] or f"РЎР»Р°Р№Рҙ {slide_count+1}"
            body_lines = [x for x in lines[1:8] if x][:6]
            sl = prs.slides.add_slide(prs.slide_layouts[1])
            sl.shapes.title.text = ttl
            tf = sl.placeholders[1].text_frame
            tf.clear()
            if not body_lines:
                body_lines = [ttl]
            for j, line in enumerate(body_lines):
                p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
                p.text = line[:240]
                p.level = 0
                with contextlib.suppress(Exception):
                    p.font.size = Pt(20)
            slide_count += 1
        if slide_count == 0:
            sl = prs.slides.add_slide(prs.slide_layouts[1])
            sl.shapes.title.text = "РЎСӮСҖСғРәСӮСғСҖР°"
            sl.placeholders[1].text = reply[:1500]
        bio = BytesIO()
        prs.save(bio)
        bio.seek(0)
        bio.name = _safe_filename("presentation", "pptx")
        await update.effective_message.reply_document(InputFile(bio), caption="рҹ“Ҡ РҹСҖРөР·РөРҪСӮР°СҶРёСҸ PPTX РіРҫСӮРҫРІР°.")
    except Exception as e:
        log.exception("pptx generation failed: %s", e)
        bio = BytesIO(reply.encode("utf-8")); bio.name = _safe_filename("presentation_outline", "txt")
        await update.effective_message.reply_document(InputFile(bio), caption="вҡ пёҸ PPTX РҪРө СҒРҫРұСҖР°Р»СҒСҸ РҪР° СҒРөСҖРІРөСҖРө, РҫСӮРҝСҖР°РІР»СҸСҺ СҒСӮСҖСғРәСӮСғСҖСғ TXT.")

async def _generate_business_catalog_pdf(update: Update, context: ContextTypes.DEFAULT_TYPE, brief: str):
    prompt = (
        "РўСӢ РјР°СҖРәРөСӮРҫР»РҫРі Рё СҖРөРҙР°РәСӮРҫСҖ РәРҫРјРјРөСҖСҮРөСҒРәРёС… РәР°СӮР°Р»РҫРіРҫРІ. РқР° СҖСғСҒСҒРәРҫРј РҝРҫРҙРіРҫСӮРҫРІСҢ СӮРөРәСҒСӮ PDF-РәР°СӮР°Р»РҫРіР° РҝРҫ РўР— РҝРҫР»СҢР·РҫРІР°СӮРөР»СҸ. "
        "РЎСӮСҖСғРәСӮСғСҖР°: РҫРұР»РҫР¶РәР°, РәСҖР°СӮРәРҫРө РЈРўРҹ, РұР»РҫРәРё РҫРұСҠРөРәСӮРҫРІ/СӮРҫРІР°СҖРҫРІ/СғСҒР»СғРі, С…Р°СҖР°РәСӮРөСҖРёСҒСӮРёРәРё, РҝСҖРөРёРјСғСүРөСҒСӮРІР°, СғСҒР»РҫРІРёСҸ, CTA, РәРҫРҪСӮР°РәСӮСӢ. "
        "Р•СҒР»Рё РҙР°РҪРҪСӢС… РҪРө С…РІР°СӮР°РөСӮ вҖ” СҒРҙРөР»Р°Р№ Р°РәРәСғСҖР°СӮРҪСӢР№ СҲР°РұР»РҫРҪ СҒ РјРөСҒСӮР°РјРё РҙР»СҸ Р·Р°РҝРҫР»РҪРөРҪРёСҸ. РўР—:\n" + (brief or "")
    )
    reply = await ask_openai_text(prompt, user_id=update.effective_user.id, chat_id=update.effective_chat.id)
    await _reply_long_text(update, "рҹ“• Р§РөСҖРҪРҫРІРёРә PDF-РәР°СӮР°Р»РҫРіР° РіРҫСӮРҫРІ. РқРёР¶Рө СӮРөРәСҒСӮ, РҫСӮРҙРөР»СҢРҪРҫ РҝСҖРёРәСҖРөРҝР»СҸСҺ PDF-С„Р°Р№Р».\n\n" + reply)
    try:
        from reportlab.lib.pagesizes import A4
        from reportlab.pdfgen import canvas
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        from reportlab.lib.units import mm
        # РҹСӢСӮР°РөРјСҒСҸ РҝРҫРҙРәР»СҺСҮРёСӮСҢ DejaVu РҙР»СҸ РәРёСҖРёР»Р»РёСҶСӢ, РөСҒР»Рё РҫРҪ РөСҒСӮСҢ РІ РҫРәСҖСғР¶РөРҪРёРё Render.
        font_name = "Helvetica"
        for fp in ("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", "/usr/share/fonts/dejavu/DejaVuSans.ttf"):
            if os.path.exists(fp):
                pdfmetrics.registerFont(TTFont("DejaVuSans", fp))
                font_name = "DejaVuSans"
                break
        bio = BytesIO()
        c = canvas.Canvas(bio, pagesize=A4)
        width, height = A4
        x, y = 18*mm, height - 20*mm
        c.setFont(font_name, 16)
        c.drawString(x, y, "PDF-РәР°СӮР°Р»РҫРі")
        y -= 12*mm
        c.setFont(font_name, 10)
        for raw in reply.splitlines():
            line = raw.strip()
            if not line:
                y -= 5*mm
                continue
            # РҹСҖРҫСҒСӮР°СҸ РҝРөСҖРөРҪРҫСҒРәР° СҒСӮСҖРҫРә.
            while line:
                part = line[:105]
                line = line[105:]
                if y < 18*mm:
                    c.showPage(); c.setFont(font_name, 10); y = height - 20*mm
                c.drawString(x, y, part)
                y -= 5.2*mm
        c.save()
        bio.seek(0); bio.name = _safe_filename("catalog", "pdf")
        await update.effective_message.reply_document(InputFile(bio), caption="рҹ“• PDF-РәР°СӮР°Р»РҫРі РіРҫСӮРҫРІ.")
    except Exception as e:
        log.exception("pdf catalog generation failed: %s", e)
        bio = BytesIO(reply.encode("utf-8")); bio.name = _safe_filename("catalog_text", "txt")
        await update.effective_message.reply_document(InputFile(bio), caption="вҡ пёҸ PDF РҪРө СҒРҫРұСҖР°Р»СҒСҸ РҪР° СҒРөСҖРІРөСҖРө, РҫСӮРҝСҖР°РІР»СҸСҺ СӮРөРәСҒСӮ РәР°СӮР°Р»РҫРіР° TXT.")

def _extract_logo_brand_name(brief: str) -> str:
    s = (brief or "").strip()
    patterns = [
        r"(?:РҪР°Р·РІР°РҪРёРө|РұСҖРөРҪРҙ|РәРҫРјРҝР°РҪ(?:РёСҸ|РёРё)|РҝСҖРҫРөРәСӮ)\s*[:вҖ”-]\s*['\"В«]?([^'\"В»\n,.;]{2,40})",
        r"['\"В«]([^'\"В»]{2,40})['\"В»]",
        r"(?:Р»РҫРіРҫСӮРёРҝ|Р»РҫРіРҫ)\s+(?:РҙР»СҸ|РұСҖРөРҪРҙР°|РәРҫРјРҝР°РҪРёРё)?\s*['\"В«]?([^'\"В»\n,.;]{2,40})",
    ]
    for pat in patterns:
        m = re.search(pat, s, flags=re.I)
        if m:
            name = re.sub(r"\s+", " ", m.group(1)).strip(" -вҖ”:;,.В«В»\"'")
            if 2 <= len(name) <= 40:
                return name
    words = re.findall(r"[A-Za-zРҗ-РҜР°-СҸРҒС‘0-9]+", s)
    stop = {"СҒРҫР·РҙР°Р№", "СҒРҙРөР»Р°Р№", "Р»РҫРіРҫСӮРёРҝ", "Р»РҫРіРҫ", "РҙР»СҸ", "РұСҖРөРҪРҙ", "РұСҖРөРҪРҙР°", "РәРҫРјРҝР°РҪРёРё", "РҪРөР№СҖРҫ", "РјСғР»СҢСӮРёРјРҫРҙР°Р»СҢРҪСӢР№", "РұРҫСӮ", "СҒСӮРёР»СҢ", "СҶРІРөСӮ", "СҶРІРөСӮР°"}
    picked = [w for w in words[:8] if w.lower() not in stop]
    return " ".join(picked[:2])[:32] or "Brand"


def _logo_initials(name: str) -> str:
    parts = re.findall(r"[A-Za-zРҗ-РҜР°-СҸРҒС‘0-9]+", name or "Brand")
    if not parts:
        return "B"
    return "".join(p[0].upper() for p in parts[:2])[:3]


def _load_logo_font(size: int, bold: bool = True):
    if ImageFont is None:
        return None
    candidates = [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
        "/usr/share/fonts/dejavu/DejaVuSans-Bold.ttf" if bold else "/usr/share/fonts/dejavu/DejaVuSans.ttf",
    ]
    for fp in candidates:
        if fp and os.path.exists(fp):
            with contextlib.suppress(Exception):
                return ImageFont.truetype(fp, size)
    return ImageFont.load_default()


def _draw_centered_text(draw, xy, text_value, font, fill):
    x, y = xy
    try:
        bbox = draw.textbbox((0, 0), text_value, font=font)
        w, h = bbox[2] - bbox[0], bbox[3] - bbox[1]
    except Exception:
        w, h = draw.textlength(text_value, font=font), 40
    draw.text((x - w / 2, y - h / 2), text_value, font=font, fill=fill)


def _make_local_logo_png(brief: str) -> bytes | None:
    if not (LOGO_LOCAL_FALLBACK and Image and ImageDraw):
        return None
    try:
        brand = _extract_logo_brand_name(brief)
        initials = _logo_initials(brand)
        W, H = 1024, 1024
        img = Image.new("RGB", (W, H), (12, 18, 33))
        draw = ImageDraw.Draw(img)
        # simple premium gradient
        for y in range(H):
            r = int(12 + 18 * y / H)
            g = int(18 + 30 * y / H)
            b = int(33 + 55 * y / H)
            draw.line([(0, y), (W, y)], fill=(r, g, b))
        # soft glow circles
        for radius, alpha_color, cx, cy in [(360, (47, 128, 255), 240, 180), (300, (124, 224, 255), 840, 820)]:
            overlay = Image.new("RGBA", (W, H), (0, 0, 0, 0))
            od = ImageDraw.Draw(overlay)
            for i in range(radius, 0, -8):
                a = int(55 * (i / radius) ** 2)
                c = alpha_color + (a,)
                od.ellipse((cx-i, cy-i, cx+i, cy+i), fill=c)
            img = Image.alpha_composite(img.convert("RGBA"), overlay).convert("RGB")
            draw = ImageDraw.Draw(img)
        # emblem
        draw.rounded_rectangle((248, 168, 776, 696), radius=120, outline=(124, 224, 255), width=10, fill=(17, 28, 51))
        draw.ellipse((332, 246, 692, 606), outline=(47, 128, 255), width=18)
        f_big = _load_logo_font(170, True)
        f_brand = _load_logo_font(58, True)
        f_sub = _load_logo_font(30, False)
        _draw_centered_text(draw, (512, 430), initials, f_big, (248, 251, 255))
        _draw_centered_text(draw, (512, 790), brand[:32], f_brand, (248, 251, 255))
        tagline = "AI вҖў BRAND вҖў PRODUCT" if re.search(r"ai|РёРё|РҪРөР№СҖРҫ|РұРҫСӮ|tech|digital", brief or "", re.I) else "PREMIUM BRAND IDENTITY"
        _draw_centered_text(draw, (512, 854), tagline, f_sub, (199, 215, 235))
        out = BytesIO()
        img.save(out, format="PNG", optimize=True)
        return out.getvalue()
    except Exception as e:
        log.exception("local logo fallback failed: %s", e)
        return None


async def _download_image_url_bytes(url: str, timeout_s: float = 180.0) -> bytes | None:
    if not url:
        return None
    try:
        async with httpx.AsyncClient(timeout=timeout_s, follow_redirects=True) as client:
            r = await client.get(url)
            r.raise_for_status()
            ct = (r.headers.get("content-type") or "").lower()
            if r.content and len(r.content) > 1000 and ("image" in ct or url.lower().split("?")[0].endswith((".png", ".jpg", ".jpeg", ".webp"))):
                return bytes(r.content)
    except Exception as e:
        log.warning("download image url failed: %s", e)
    return None


async def _comet_generate_image_bytes(prompt: str) -> bytes | None:
    if not (COMET_API_KEY and COMET_BASE_URL and COMET_IMAGE_GEN_PATH):
        return None
    url = f"{COMET_BASE_URL}{COMET_IMAGE_GEN_PATH}"
    headers = {"Authorization": f"Bearer {COMET_API_KEY}", "Content-Type": "application/json", "Accept": "application/json"}
    payloads = [
        {"model": COMET_IMAGE_GEN_MODEL, "prompt": prompt, "size": "1024x1024", "n": 1},
        {"model": COMET_IMAGE_GEN_MODEL, "prompt": prompt, "response_format": "b64_json", "size": "1024x1024", "n": 1},
    ]
    try:
        async with httpx.AsyncClient(timeout=float(COMET_IMAGE_GEN_TIMEOUT_S), follow_redirects=True) as client:
            last_body = ""
            for payload in payloads:
                r = await client.post(url, headers=headers, json=payload)
                last_body = r.text[:500]
                if r.status_code >= 400:
                    log.warning("Comet image gen HTTP %s: %s", r.status_code, last_body)
                    continue
                j = r.json() or {}
                candidates = []
                if isinstance(j.get("data"), list):
                    candidates.extend(j.get("data") or [])
                if isinstance(j.get("images"), list):
                    candidates.extend(j.get("images") or [])
                if isinstance(j.get("output"), list):
                    candidates.extend(j.get("output") or [])
                candidates.append(j)
                for item in candidates:
                    if not isinstance(item, dict):
                        continue
                    b64 = item.get("b64_json") or item.get("image_base64") or item.get("base64")
                    if b64:
                        if "," in b64 and b64.strip().startswith("data:"):
                            b64 = b64.split(",", 1)[1]
                        with contextlib.suppress(Exception):
                            return base64.b64decode(b64)
                    url2 = item.get("url") or item.get("image_url")
                    if isinstance(url2, dict):
                        url2 = url2.get("url")
                    if url2:
                        got = await _download_image_url_bytes(str(url2), timeout_s=float(COMET_IMAGE_GEN_TIMEOUT_S))
                        if got:
                            return got
            log.warning("Comet image generation did not return image. Last body: %s", last_body)
    except Exception as e:
        log.exception("Comet image generation failed: %s", e)
    return None


async def _generate_logo_image_bytes(prompt: str, brief: str) -> bytes | None:
    # 1) OpenAI/Luma route from older builds
    img = await _luma_generate_image_bytes(prompt)
    if img:
        return img
    # 2) Comet OpenAI-compatible image generation route
    img = await _comet_generate_image_bytes(prompt)
    if img:
        return img
    # 3) Guaranteed local fallback, so the user does not get a dead end
    return await asyncio.to_thread(_make_local_logo_png, brief)


async def _generate_business_logo(update: Update, context: ContextTypes.DEFAULT_TYPE, brief: str):
    prompt = (
        "Professional logo design, clean vector style, high-end brand identity, flat vector mark, no mockup, no watermark. "
        "Create a modern logo based on this Russian brief. Keep brand name readable if provided. "
        "Avoid small unreadable text. Square composition, transparent-background style if possible. Brief: " + (brief or "")
    )
    await update.effective_message.reply_text(
        "рҹҺЁ Р—Р°РҝСғСҒРәР°СҺ СҒРҫР·РҙР°РҪРёРө Р»РҫРіРҫСӮРёРҝР°. РҹСҖРҫРұСғСҺ image-РҝСҖРҫРІР°Р№РҙРөСҖ, Р° РөСҒР»Рё РҫРҪ РҪРөРҙРҫСҒСӮСғРҝРөРҪ вҖ” СҒРҫРұРөСҖСғ Р°РәРәСғСҖР°СӮРҪСӢР№ Р»РҫРәР°Р»СҢРҪСӢР№ PNG-Р»РҫРіРҫСӮРёРҝ, СҮСӮРҫРұСӢ С„СғРҪРәСҶРёСҸ РҪРө РҝР°РҙР°Р»Р°."
    )

    async def _go():
        img = await _generate_logo_image_bytes(prompt, brief)
        if not img:
            await update.effective_message.reply_text("РқРө СғРҙР°Р»РҫСҒСҢ СҒРҫР·РҙР°СӮСҢ РёР·РҫРұСҖР°Р¶РөРҪРёРө. РҹСҖРҫРІРөСҖСҢСӮРө OPENAI_IMAGE_KEY РёР»Рё COMET_API_KEY.")
            return False
        try:
            await update.effective_message.reply_photo(photo=img, caption="рҹҺЁ РӣРҫРіРҫСӮРёРҝ РіРҫСӮРҫРІ вң…")
        except Exception:
            bio = BytesIO(img); bio.name = _safe_filename("logo", "png")
            await update.effective_message.reply_document(InputFile(bio), caption="рҹҺЁ РӣРҫРіРҫСӮРёРҝ РіРҫСӮРҫРІ вң…")
        return True

    await _try_pay_then_do(
        update, context, update.effective_user.id,
        "img", IMG_COST_USD, _go,
        remember_kind="business_logo",
        remember_payload={"brief": brief},
    )



# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ Presentation / Catalog Studio v86 integration в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
_PRESENTATION_STUDIO_INSTANCE = None


async def _presentation_llm_call(update: Update, prompt: str) -> str:
    return await ask_openai_text(
        prompt,
        user_id=update.effective_user.id if update.effective_user else None,
        chat_id=update.effective_chat.id if update.effective_chat else None,
        extra_system=(
            "РўСӢ СҖР°РұРҫСӮР°РөСҲСҢ РІРҪСғСӮСҖРё РјР°СҒСӮРөСҖР° РҝСҖРөР·РөРҪСӮР°СҶРёР№. Р’РҫР·РІСҖР°СүР°Р№ СҒСӮСҖСғРәСӮСғСҖРёСҖРҫРІР°РҪРҪСӢРө РҙР°РҪРҪСӢРө СӮРҫСҮРҪРҫ РІ Р·Р°РҝСҖРҫСҲРөРҪРҪРҫРј С„РҫСҖРјР°СӮРө, "
            "РҪРө РҙРҫРұР°РІР»СҸР№ РІСӢРјСӢСҲР»РөРҪРҪСӢРө СҶРөРҪСӢ, СӮРөР»РөС„РҫРҪСӢ, СҒРөСҖСӮРёС„РёРәР°СӮСӢ Рё С„Р°РәСӮСӢ."
        ),
    )


def _presentation_local_visual(prompt: str, index: int = 0) -> bytes:
    """Guaranteed local visual fallback so a long presentation workflow never ends with an empty asset."""
    try:
        if Image is None or ImageDraw is None:
            return b""
        palettes = [
            ((12, 18, 32), (58, 125, 255), (124, 224, 255)),
            ((28, 25, 22), (205, 168, 83), (245, 236, 210)),
            ((233, 239, 235), (74, 122, 101), (199, 184, 145)),
            ((245, 247, 251), (47, 105, 235), (150, 192, 255)),
        ]
        bg, accent, accent2 = palettes[index % len(palettes)]
        w, h = 1536, 1024
        im = Image.new("RGB", (w, h), bg)
        d = ImageDraw.Draw(im)
        for y in range(h):
            ratio = y / max(1, h - 1)
            fill = tuple(int(bg[i] * (1 - ratio * 0.25) + accent[i] * ratio * 0.25) for i in range(3))
            d.line((0, y, w, y), fill=fill)
        d.ellipse((850, -220, 1630, 560), fill=accent)
        d.ellipse((-320, 480, 520, 1320), fill=accent2)
        d.rounded_rectangle((150, 150, 910, 820), radius=80, outline=(255, 255, 255), width=5)
        d.line((220, 705, 760, 705), fill=(255, 255, 255), width=8)
        # Keep the fallback intentionally text-free: the renderer adds approved copy itself.
        out = BytesIO(); im.save(out, format="JPEG", quality=92, optimize=True)
        return out.getvalue()
    except Exception:
        return b""


async def _presentation_paid_runner(
    update: Update,
    context: ContextTypes.DEFAULT_TYPE,
    engine: str,
    feature: str,
    provider_cost_usd: float,
    action,
):
    """Run a presentation subtask and return its result.

    Logo concepts are included in the presentation workflow and use a silent provider attempt;
    deterministic local concepts are supplied by presentation_studio.py when the provider is unavailable.
    This prevents partial logo cards followed by a misleading global error message.
    """
    if feature == "presentation_logo_variants":
        try:
            return await action()
        except Exception as e:
            log.exception("Presentation logo provider attempt failed silently: %s", e)
            return None

    box = {"result": None}

    async def _go():
        result = await action()
        box["result"] = result
        return bool(result)

    await _try_pay_then_do(
        update,
        context,
        update.effective_user.id,
        engine if engine in ("luma", "runway", "img") else "img",
        max(0.0, float(provider_cost_usd or 0.0)),
        _go,
        remember_kind=feature,
        remember_payload={"presentation_studio": True, "feature": feature},
    )
    return box.get("result")


async def _presentation_image_batch(
    update: Update,
    context: ContextTypes.DEFAULT_TYPE,
    prompts: list[str],
    engine: str,
    feature: str,
) -> list[bytes] | None:
    prompts = [re.sub(r"\s+", " ", (p or "").strip())[:1800] for p in prompts if (p or "").strip()]
    if not prompts:
        return []

    resolved_engines: list[str] = []
    text_sensitive_re = re.compile(
        r"(label|etiket|packag|bottle|box|poster|menu|sign|screen|document|book|magazine|logo|brand|"
        r"СҚСӮРёРәРөСӮ|СғРҝР°РәРҫРІ|РұСғСӮСӢР»|РәРҫСҖРҫРұ|РҝР»Р°РәР°СӮ|РјРөРҪСҺ|РІСӢРІРөСҒРә|СҚРәСҖР°РҪ|РҙРҫРәСғРјРөРҪСӮ|РәРҪРёРі|Р¶СғСҖРҪР°Р»|Р»РҫРіРҫСӮРёРҝ|РұСҖРөРҪРҙ)", re.I
    )
    safe_prompts: list[str] = []
    for prompt in prompts:
        text_sensitive = bool(text_sensitive_re.search(prompt))
        safe_prompt = prompt
        if PRESENTATION_TEXT_SAFE_VISUALS and text_sensitive:
            safe_prompt += (
                " STRICT TYPOGRAPHY SAFETY: all labels, packaging, signs and printed areas must be completely blank and unprinted; "
                "no letters, pseudo-letters, glyphs, numbers or fake brand marks. Exact approved copy will be overlaid later by the renderer."
            )
        safe_prompts.append(safe_prompt)
        if PRESENTATION_FORCE_OPENAI_FOR_TEXT and text_sensitive:
            # Midjourney is excellent for atmosphere but unreliable for exact lettering.
            resolved_engines.append("openai")
        elif engine == "auto":
            premium = bool(re.search(r"luxury|premium|editorial|lifestyle|cinematic|Р°СӮРјРҫСҒС„РөСҖ|РҝСҖРөРјРёСғРј|Р»СҺРәСҒ", prompt, re.I))
            resolved_engines.append("midjourney" if premium and MIDJOURNEY_ENABLED and COMET_API_KEY else "openai")
        else:
            resolved_engines.append(engine)

    provider_cost = 0.0
    for e in resolved_engines:
        provider_cost += MIDJOURNEY_UNIT_COST_USD if e == "midjourney" else IMG_COST_USD

    async def _action():
        result: list[bytes] = []
        for idx, (prompt, selected_engine) in enumerate(zip(safe_prompts, resolved_engines)):
            img = None
            try:
                if selected_engine == "midjourney":
                    img, _task = await _midjourney_generate_image_bytes(prompt)
                    if not img:
                        img = await _luma_generate_image_bytes(prompt)
                else:
                    img = await _luma_generate_image_bytes(prompt)
                    if not img:
                        img = await _comet_generate_image_bytes(prompt)
                    # Logo concepts remain text-safe: do not silently switch them to Midjourney.
                    if not img and feature != "presentation_logo_variants" and MIDJOURNEY_ENABLED and COMET_API_KEY:
                        with contextlib.suppress(Exception):
                            img, _task = await _midjourney_generate_image_bytes(prompt)
            except Exception as e:
                log.warning("Presentation image %s/%s failed: %s", idx + 1, len(prompts), e)
            if not img and feature != "presentation_logo_variants":
                img = _presentation_local_visual(prompt, idx)
            if not img:
                # For logos, keep any successful partial results; the studio fills missing concepts locally.
                continue
            result.append(img)
        return result

    return await _presentation_paid_runner(update, context, "img", feature, provider_cost, _action)


def _presentation_studio_get() -> PresentationStudio:
    global _PRESENTATION_STUDIO_INSTANCE
    if _PRESENTATION_STUDIO_INSTANCE is None:
        _PRESENTATION_STUDIO_INSTANCE = PresentationStudio(
            StudioConfig(
                db_path=DB_PATH,
                data_dir=PRESENTATION_DATA_DIR,
                max_uploads=max(5, PRESENTATION_MAX_UPLOADS),
                max_generated_images=max(1, PRESENTATION_MAX_GENERATED_IMAGES),
                render_cost_usd=max(0.0, PRESENTATION_RENDER_COST_USD),
            ),
            llm_call=_presentation_llm_call,
            image_batch_call=_presentation_image_batch,
            paid_runner=_presentation_paid_runner,
        )
    return _PRESENTATION_STUDIO_INSTANCE


async def cmd_presentation(update: Update, context: ContextTypes.DEFAULT_TYPE):
    _clear_transient_flows(context)
    await _presentation_studio_get().start(update, context, "presentation")


async def cmd_catalog(update: Update, context: ContextTypes.DEFAULT_TYPE):
    _clear_transient_flows(context)
    await _presentation_studio_get().start(update, context, "catalog")


async def cmd_diag_presentation(update: Update, context: ContextTypes.DEFAULT_TYPE):
    studio = _presentation_studio_get()
    project = studio._active_project(update.effective_user.id, update.effective_chat.id)
    lines = [
        f"рҹ§Ә Presentation Studio / {PATCH_VERSION}",
        f"data_dir={studio.data_dir}",
        f"db_path={DB_PATH}",
        f"max_uploads={PRESENTATION_MAX_UPLOADS}",
        f"max_generated_images={PRESENTATION_MAX_GENERATED_IMAGES}",
        f"render_cost_usd={PRESENTATION_RENDER_COST_USD}",
        f"text_safe_visuals={PRESENTATION_TEXT_SAFE_VISUALS}",
        f"force_openai_for_text={PRESENTATION_FORCE_OPENAI_FOR_TEXT}",
        f"pptx={'on' if shutil.which('python') or True else 'off'}",
        f"active_project={project.get('id') if project else '-'}",
        f"state={project.get('state') if project else '-'}",
    ]
    await update.effective_message.reply_text("\n".join(lines))


async def _reply_audio_from_url(update: Update, client: httpx.AsyncClient, url: str, caption: str):
    try:
        r = await client.get(url, timeout=240.0, follow_redirects=True)
        r.raise_for_status()
        if r.content and len(r.content) > 1024 and "json" not in (r.headers.get("content-type") or "").lower():
            bio = BytesIO(r.content)
            bio.name = "suno_track.mp3"
            try:
                await update.effective_message.reply_audio(audio=InputFile(bio), caption=caption)
                return
            except Exception:
                bio.seek(0)
                await update.effective_message.reply_document(document=InputFile(bio), caption=caption)
                return
    except Exception as e:
        log.warning("audio download failed: %s", e)
    await update.effective_message.reply_text(f"{caption}\nРЎСҒСӢР»РәР° РҪР° СҖРөР·СғР»СҢСӮР°СӮ: {url}", disable_web_page_preview=False)

async def _poll_suno_task(update: Update, client: httpx.AsyncClient, headers: dict, task_id: str):
    started = time.time()
    status_paths = [
        SUNO_STATUS_PATH,
        "/suno/fetch/{id}",
        "/suno/v1/music/{id}",
        "/api/v1/task/{id}",
        "/v1/tasks/{id}",
    ]
    last = ""
    while time.time() - started < SUNO_TIMEOUT_S:
        for path in status_paths:
            try:
                rs = await client.get(f"{SUNO_BASE_URL}{path}".format(id=task_id), headers=headers, timeout=60.0)
                if rs.status_code >= 400:
                    last = f"{rs.status_code}: {_api_error_preview(rs)}"
                    continue
                js = rs.json() or {}
                data_obj = js.get("data")
                nested = data_obj if isinstance(data_obj, dict) else {}
                nested_data = nested.get("data") if isinstance(nested, dict) else None
                url = (
                    _extract_first_url(js.get("audio_url"))
                    or _extract_first_url(js.get("audio"))
                    or _extract_first_url(js.get("output"))
                    or _extract_first_url(data_obj)
                    or _extract_first_url(nested_data)
                    or _extract_first_url(js)
                )
                st = str(js.get("status") or js.get("state") or js.get("task_status") or nested.get("status") or "").lower()
                if url and (st in ("", "completed", "succeeded", "success", "finished", "done", "ready") or not st):
                    await _reply_audio_from_url(update, client, url, "рҹҺө РңСғР·СӢРәР° Suno РіРҫСӮРҫРІР° вң…")
                    return True
                if st in ("failed", "fail", "error", "canceled", "cancelled", "rejected"):
                    await update.effective_message.reply_text(f"вқҢ Suno: РҫСҲРёРұРәР° РіРөРҪРөСҖР°СҶРёРё.\n{json.dumps(js, ensure_ascii=False)[:1200]}")
                    return True
                last = json.dumps(js, ensure_ascii=False)[:700]
            except Exception as e:
                last = str(e)
        await asyncio.sleep(SUNO_POLL_DELAY_S)
    await update.effective_message.reply_text(f"вҢӣ Suno: РІСҖРөРјСҸ РҫР¶РёРҙР°РҪРёСҸ РІСӢСҲР»Рҫ. РҹРҫСҒР»РөРҙРҪРёР№ РҫСӮРІРөСӮ: {last[:700]}")
    return False

async def _run_suno_music(update: Update, context: ContextTypes.DEFAULT_TYPE, brief: str):
    brief = (brief or "").strip()
    if not brief:
        context.user_data["awaiting_suno_brief"] = True
        await update.effective_message.reply_text(_suno_submenu_text(), reply_markup=_suno_menu_kb())
        return
    if not SUNO_ENABLED:
        await update.effective_message.reply_text("вҡ пёҸ Р РөР¶РёРј Suno РҫСӮРәР»СҺСҮС‘РҪ РІ ENV: SUNO_ENABLED=0.")
        return
    if not SUNO_API_KEY:
        await update.effective_message.reply_text("вқҢ Р”Р»СҸ РјСғР·СӢРәРё РҪСғР¶РөРҪ SUNO_API_KEY РёР»Рё COMET_API_KEY РІ Environment.")
        return

    async def _go():
        await update.effective_message.reply_text("рҹҺө Р—Р°РҝСғСҒРәР°СҺ Suno. РһРұСӢСҮРҪРҫ СӮСҖРөРә РіРҫСӮРҫРІРёСӮСҒСҸ 1вҖ“5 РјРёРҪСғСӮвҖҰ")
        headers = {"Authorization": f"Bearer {SUNO_API_KEY}", "Content-Type": "application/json", "Accept": "application/json"}
        instrumental = bool(re.search(r"РёРҪСҒСӮСҖСғРјРөРҪСӮР°Р»|instrumental|РұРөР· РІРҫРәР°Р»Р°|РұРөР· РіРҫР»РҫСҒР°|РјРёРҪСғСҒРҫРІ|background music|С„РҫРҪРҫРІР°", brief, re.I))
        # Documented CometAPI Suno endpoint: /suno/submit/music.
        # For free-form user prompts we use Inspiration Mode; for instrumental prompts
        # we add make_instrumental=True and empty prompt, as Comet expects.
        base_payload = {"mv": SUNO_MODEL, "gpt_description_prompt": brief}
        if instrumental:
            base_payload.update({"prompt": "", "make_instrumental": True})
        payloads = [
            (SUNO_CREATE_PATH, base_payload),
            ("/suno/submit/music", base_payload),
        ]
        last_err = ""
        async with httpx.AsyncClient(timeout=90.0) as client:
            for path, payload in payloads:
                try:
                    r = await client.post(f"{SUNO_BASE_URL}{path}", headers=headers, json=payload)
                    if r.status_code >= 400:
                        last_err = f"POST {path} вҶ’ {r.status_code}: {_api_error_preview(r)}"
                        log.warning("Suno create failed: %s", last_err)
                        continue
                    js = r.json() or {}
                    url = _extract_first_url(js.get("audio_url")) or _extract_first_url(js.get("audio")) or _extract_first_url(js.get("output")) or _extract_first_url(js.get("data")) or _extract_first_url(js)
                    if url:
                        await _reply_audio_from_url(update, client, url, "рҹҺө РңСғР·СӢРәР° Suno РіРҫСӮРҫРІР° вң…")
                        return True
                    task_id = str(js.get("id") or js.get("task_id") or js.get("taskId") or js.get("request_id") or "").strip()
                    data_obj = js.get("data")
                    if not task_id and isinstance(data_obj, str):
                        # Comet returns {"code":"success","data":"<task_id>"} for /suno/submit/music.
                        if re.fullmatch(r"[A-Za-z0-9_.:-]{8,}", data_obj.strip()):
                            task_id = data_obj.strip()
                    if not task_id and isinstance(data_obj, dict):
                        d = data_obj or {}
                        task_id = str(d.get("id") or d.get("task_id") or d.get("taskId") or d.get("request_id") or "").strip()
                    if task_id:
                        await update.effective_message.reply_text(f"вҸі Suno: Р·Р°РҙР°СҮР° РҝСҖРёРҪСҸСӮР°, id={task_id}. Р–РҙСғ СҖРөР·СғР»СҢСӮР°СӮвҖҰ")
                        return bool(await _poll_suno_task(update, client, headers, task_id))
                    last_err = f"POST {path}: РҪРөСӮ audio_url/id РІ РҫСӮРІРөСӮРө {json.dumps(js, ensure_ascii=False)[:700]}"
                except Exception as e:
                    last_err = f"POST {path}: {e}"
                    log.warning("Suno create exception: %s", e)
            await update.effective_message.reply_text(f"вқҢ Suno: РҪРө СғРҙР°Р»РҫСҒСҢ СҒРҫР·РҙР°СӮСҢ Р·Р°РҙР°СҮСғ.\n{last_err[:1500]}")
            return False
    await _try_pay_then_do(update, context, update.effective_user.id, SUNO_BILLING_ENGINE, SUNO_COST_USD, _go, remember_kind="suno_music", remember_payload={"prompt": brief})


async def cmd_diag_yookassa(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Р‘РөР·РҫРҝР°СҒРҪР°СҸ РҙРёР°РіРҪРҫСҒСӮРёРәР° Р®Kassa: РҪРө РҝРҫРәР°Р·СӢРІР°РөСӮ СҒРөРәСҖРөСӮСӢ, СӮРҫР»СҢРәРҫ РҪР°Р»РёСҮРёРө РҪР°СҒСӮСҖРҫРөРә."""
    def mask(v: str) -> str:
        v = (v or "").strip()
        if not v:
            return "РҪРөСӮ"
        if len(v) <= 8:
            return "РөСҒСӮСҢ"
        return v[:5] + "вҖҰ" + v[-4:]

    source = []
    for _p in ("/etc/secrets/yookassa.env", "/etc/secrets/yookassa.txt", "/etc/secrets/yk.env", "/etc/secrets/yk.txt"):
        try:
            if os.path.exists(_p):
                source.append(_p)
        except Exception:
            pass
    provider_ok = bool(YOOKASSA_PROVIDER_TOKEN)
    provider_hint = "РөСҒСӮСҢ" if provider_ok else "РҪРөСӮ"
    if YOOKASSA_PROVIDER_TOKEN and not (":" in YOOKASSA_PROVIDER_TOKEN and "LIVE" in YOOKASSA_PROVIDER_TOKEN.upper()):
        provider_hint = "РөСҒСӮСҢ, РҪРҫ С„РҫСҖРјР°СӮ РҪРө РҝРҫС…РҫР¶ РҪР° Telegram provider token"

    lines = [
        f"рҹ§ҫ YooKassa diagnostic / {PATCH_VERSION}",
        f"Direct API enabled: {YOO_DIRECT_ENABLED}",
        f"ShopID: {mask(YOO_SHOP_ID)}",
        f"API Secret: {mask(YOO_SECRET_KEY)}",
        f"Secret files found: {', '.join(source) if source else 'РҪРөСӮ'}",
        f"Telegram provider token: {provider_hint}",
        f"Return URL: {YOO_PAYMENT_RETURN_URL}",
        "",
        f"РЎР‘Рҹ/QR: {YOO_SBP_ENABLED}",
        f"SberPay: {YOO_SBERPAY_ENABLED}",
        f"T-Pay: {YOO_TPAY_ENABLED}",
        f"Mir Pay: {YOO_MIRPAY_ENABLED}",
        f"Telegram card: {YOO_CARD_ENABLED}",
    ]
    if _yoo_direct_configured():
        lines.append("вң… Direct API Р®Kassa РҪР°СҒСӮСҖРҫРөРҪ. РңРҫР¶РҪРҫ СӮРөСҒСӮРёСҖРҫРІР°СӮСҢ РЎР‘Рҹ/QR.")
    else:
        lines.append("вқҢ Direct API Р®Kassa РҪРө РҪР°СҒСӮСҖРҫРөРҪ: РҪСғР¶РөРҪ yookassa.env СҒ YK_ID Рё YK_KEY.")
    await update.effective_message.reply_text("\n".join(lines))

async def cmd_diag_suno(update: Update, context: ContextTypes.DEFAULT_TYPE):
    lines = [
        f"рҹ§Ә Suno diagnostic / {PATCH_VERSION}",
        f"SUNO_ENABLED={SUNO_ENABLED}",
        f"SUNO_API_KEY={'on' if bool(SUNO_API_KEY) else 'off'}",
        f"SUNO_BASE_URL={SUNO_BASE_URL}",
        f"SUNO_MODEL/mv={SUNO_MODEL}",
        f"SUNO_CREATE_PATH={SUNO_CREATE_PATH}",
        f"SUNO_STATUS_PATH={SUNO_STATUS_PATH}",
        f"SUNO_TIMEOUT_S={SUNO_TIMEOUT_S} poll={SUNO_POLL_DELAY_S}",
        "Expected Comet endpoints: POST /suno/submit/music, GET /suno/fetch/{task_id}",
    ]
    await update.effective_message.reply_text("\n".join(lines)[:3900])


def _suno_submenu_text() -> str:
    # Plain text intentionally: Telegram Markdown is fragile in callback-edited messages.
    return (
        "рҹҺө РңСғР·СӢРәР° / РҝРөСҒРҪРё вҖ” Suno\n"
        "РңРҫР¶РҪРҫ СҒРҫР·РҙР°РІР°СӮСҢ РіРҫСӮРҫРІСӢРө РҝРөСҒРҪРё СҒ РІРҫРәР°Р»РҫРј, СҖРөРәР»Р°РјРҪСӢРө РҙР¶РёРҪРіР»СӢ, РёРҪСӮСҖРҫ/Р°СғСӮСҖРҫ РҙР»СҸ Reels, "
        "С„РҫРҪРҫРІСғСҺ РјСғР·СӢРәСғ, РјРёРҪСғСҒРҫРІРәРё, РҙРөРјРҫ-СӮСҖРөРәРё Рё РәРҫСҖРҫСӮРәРёРө РұСҖРөРҪРҙ-Р°СғРҙРёРҫ.\n\n"
        "РҡР°Рә Р·Р°РҝСғСҒРәР°СӮСҢ:\n"
        "1) РҪР°Р¶РјРёСӮРө РҝСҖРөСҒРөСӮ РҪРёР¶Рө;\n"
        "2) РҪР°РҝРёСҲРёСӮРө РҙРөСӮР°Р»Рё: СӮРөРјР°, РұСҖРөРҪРҙ/РҪРёСҲР°, СҒСӮРёР»СҢ, СҸР·СӢРә, РіРҫР»РҫСҒ, РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ;\n"
        "3) СҸ СҒРҫРұРөСҖСғ РҝСҖР°РІРёР»СҢРҪСӢР№ Suno-РҝСҖРҫРјРҝСӮ Рё РҫСӮРҝСҖР°РІР»СҺ Р·Р°РҙР°СҮСғ.\n\n"
        "РңРҫР¶РҪРҫ РІСӢРұСҖР°СӮСҢ В«РЎРІРҫРұРҫРҙРҪСӢР№ Р·Р°РҝСҖРҫСҒВ» Рё РҪР°РҝРёСҒР°СӮСҢ РІСҒС‘ РҫРҙРҪРёРј СҒРҫРҫРұСүРөРҪРёРөРј. РһРұСӢСҮРҪРҫ РіРҫСӮРҫРІРҪРҫСҒСӮСҢ 1вҖ“5 РјРёРҪСғСӮ."
    )


def _suno_menu_kb() -> InlineKeyboardMarkup:
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("вңҚпёҸ РЎРІРҫРұРҫРҙРҪСӢР№ Р·Р°РҝСҖРҫСҒ", callback_data="suno:free")],
        [InlineKeyboardButton("рҹҸў Р”Р¶РёРҪРіР» РҙР»СҸ РұРёР·РҪРөСҒР°", callback_data="suno:preset:jingle")],
        [InlineKeyboardButton("рҹ“ұ РҳРҪСӮСҖРҫ РҙР»СҸ Reels", callback_data="suno:preset:reels_intro")],
        [InlineKeyboardButton("рҹҺ¬ РӨРҫРҪРҫРІР°СҸ РјСғР·СӢРәР°", callback_data="suno:preset:background")],
        [InlineKeyboardButton("рҹҺӨ РҹРөСҒРҪСҸ СҒ РІРҫРәР°Р»РҫРј", callback_data="suno:preset:vocal_song")],
        [InlineKeyboardButton("рҹҺ§ РңРёРҪСғСҒРҫРІРәР° / instrumental", callback_data="suno:preset:instrumental")],
        [InlineKeyboardButton("рҹҸқ Luxury tropical", callback_data="suno:preset:tropical_luxury")],
        [InlineKeyboardButton("в¬…пёҸ РқР°Р·Р°Рҙ Рә СҖР°Р·РІР»РөСҮРөРҪРёСҸРј", callback_data="suno:back_fun")],
    ])


def _suno_preset_instruction(kind: str) -> str:
    presets = {
        "jingle": (
            "рҹҸў Р”Р¶РёРҪРіР» РҙР»СҸ РұРёР·РҪРөСҒР°\n"
            "РқР°РҝРёСҲРёСӮРө: РҪРёСҲР°/РұСҖРөРҪРҙ, РіРҫСҖРҫРҙ РёР»Рё РіРөРҫРіСҖР°С„РёСҸ, РҪР°СҒСӮСҖРҫРөРҪРёРө, СҸР·СӢРә, РҪСғР¶РөРҪ Р»Рё РІРҫРәР°Р», РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ.\n\n"
            "РҹСҖРёРјРөСҖ: Р°РіРөРҪСӮСҒСӮРІРҫ РҪРөРҙРІРёР¶РёРјРҫСҒСӮРё РҪР° РЎР°РјСғРё, 20 СҒРөРәСғРҪРҙ, tropical house, luxury, РұРөР· РІРҫРәР°Р»Р°, Р·Р°РҝРҫРјРёРҪР°СҺСүРёР№СҒСҸ РұСҖРөРҪРҙРҫРІСӢР№ Р·РІСғРә."
        ),
        "reels_intro": (
            "рҹ“ұ РҳРҪСӮСҖРҫ / Р°СғСӮСҖРҫ РҙР»СҸ Reels\n"
            "РқР°РҝРёСҲРёСӮРө СӮРөРјСғ Reels, РҪР°СҒСӮСҖРҫРөРҪРёРө, РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ 5вҖ“20 СҒРөРә, РҪСғР¶РөРҪ Р»Рё РіРҫР»РҫСҒ/РІРҫРәР°Р».\n\n"
            "РҹСҖРёРјРөСҖ: РёРҪСӮСҖРҫ 12 СҒРөРәСғРҪРҙ РҙР»СҸ СҖРёР»СҒРҫРІ РҝСҖРҫ РІРёР»Р»СӢ РҪР° РЎР°РјСғРё, luxury tropical, СҚРҪРөСҖРіРёСҮРҪРҫ, РұРөР· РІРҫРәР°Р»Р°, СҒРҫРІСҖРөРјРөРҪРҪСӢР№ РұРёСӮ."
        ),
        "background": (
            "рҹҺ¬ РӨРҫРҪРҫРІР°СҸ РјСғР·СӢРәР°\n"
            "РқР°РҝРёСҲРёСӮРө РҙР»СҸ СҮРөРіРҫ С„РҫРҪ: РІРёРҙРөРҫ, РҝСҖРөР·РөРҪСӮР°СҶРёСҸ, СҒСӮРҫСҖРёСҒ, СҲРҫСғСҖСғРј, СҖРөСҒСӮРҫСҖР°РҪ, РҪРөРҙРІРёР¶РёРјРҫСҒСӮСҢ. РЈРәР°Р¶РёСӮРө РҪР°СҒСӮСҖРҫРөРҪРёРө Рё РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ.\n\n"
            "РҹСҖРёРјРөСҖ: С„РҫРҪРҫРІР°СҸ РјСғР·СӢРәР° 60 СҒРөРәСғРҪРҙ РҙР»СҸ РҝСҖРөР·РөРҪСӮР°СҶРёРё РҝСҖРөРјРёР°Р»СҢРҪРҫР№ РІРёР»Р»СӢ, cinematic tropical, СҒРҝРҫРәРҫР№РҪР°СҸ, РұРөР· РІРҫРәР°Р»Р°."
        ),
        "vocal_song": (
            "рҹҺӨ РҹРөСҒРҪСҸ СҒ РІРҫРәР°Р»РҫРј\n"
            "РқР°РҝРёСҲРёСӮРө СӮРөРјСғ РҝРөСҒРҪРё, СҸР·СӢРә, РјСғР¶СҒРәРҫР№/Р¶РөРҪСҒРәРёР№ РІРҫРәР°Р», Р¶Р°РҪСҖ, РҪР°СҒСӮСҖРҫРөРҪРёРө, РјРҫР¶РҪРҫ РҙРҫРұР°РІРёСӮСҢ СӮРөРәСҒСӮ РҝСҖРёРҝРөРІР° РёР»Рё РәСғРҝР»РөСӮР°.\n\n"
            "РҹСҖРёРјРөСҖ: РҝРөСҒРҪСҸ РҪР° СҖСғСҒСҒРәРҫРј РҝСҖРҫ Р¶РёР·РҪСҢ Сғ РјРҫСҖСҸ, РјСғР¶СҒРәРҫР№ РІРҫРәР°Р», pop house, РҝСҖРёРҝРөРІ РҙРҫР»Р¶РөРҪ РұСӢСӮСҢ Р·Р°РҝРҫРјРёРҪР°СҺСүРёРјСҒСҸ."
        ),
        "instrumental": (
            "рҹҺ§ РңРёРҪСғСҒРҫРІРәР° / instrumental\n"
            "РқР°РҝРёСҲРёСӮРө Р¶Р°РҪСҖ, РҪР°СҒСӮСҖРҫРөРҪРёРө, СӮРөРјРҝ Рё РіРҙРө РұСғРҙРөСӮ РёСҒРҝРҫР»СҢР·РҫРІР°СӮСҢСҒСҸ СӮСҖРөРә.\n\n"
            "РҹСҖРёРјРөСҖ: РёРҪСҒСӮСҖСғРјРөРҪСӮР°Р» 90 СҒРөРәСғРҪРҙ, deep house, luxury, Р°СӮРјРҫСҒС„РөСҖРҪРҫ, РҙР»СҸ РІРёРҙРөРҫ СӮСғСҖР° РҝРҫ РІРёР»Р»Рө, РұРөР· РІРҫРәР°Р»Р°."
        ),
        "tropical_luxury": (
            "рҹҸқ Luxury tropical\n"
            "РҹСҖРөСҒРөСӮ РҝРҫРҙ РҪРөРҙРІРёР¶РёРјРҫСҒСӮСҢ, РІРёР»Р»СӢ, РҝСғСӮРөСҲРөСҒСӮРІРёСҸ, РҝР»СҸР¶РҪСӢР№ РҝСҖРөРјРёСғРј-РәРҫРҪСӮРөРҪСӮ. РқР°РҝРёСҲРёСӮРө РҫРұСҠРөРәСӮ/РұСҖРөРҪРҙ Рё РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ.\n\n"
            "РҹСҖРёРјРөСҖ: 30 СҒРөРәСғРҪРҙ РҙР»СҸ Р°РіРөРҪСӮСҒСӮРІР° Cozy Asia, РЎР°РјСғРё, premium tropical house, РҙРҫСҖРҫРіРҫР№ РјСҸРіРәРёР№ Р·РІСғРә, РұРөР· РІРҫРәР°Р»Р°."
        ),
    }
    return presets.get(kind, _suno_submenu_text())


def _suno_preset_prefix(kind: str) -> str:
    prefixes = {
        "jingle": "Create a polished business advertising jingle. Duration 15-30 seconds unless user specifies otherwise. Catchy, brand-friendly, high production value.",
        "reels_intro": "Create a short intro/outro music bed for Reels/Shorts. Duration 5-20 seconds unless user specifies otherwise. Modern, hooky, social-media ready.",
        "background": "Create background music for video/presentation. Instrumental unless user explicitly asks for vocals. Smooth mix, no harsh lead vocal.",
        "vocal_song": "Create a full song with vocals. Follow requested language, vocal gender, genre, lyrics/theme and mood.",
        "instrumental": "Create an instrumental track, no vocals. Follow requested genre, tempo, mood, and use-case.",
        "tropical_luxury": "Create luxury tropical house / premium island real-estate music. Warm, expensive, elegant, suitable for villa and travel content. Instrumental unless user asks for vocals.",
    }
    return prefixes.get(kind, "")


def _prepare_suno_brief_from_context(context: ContextTypes.DEFAULT_TYPE, user_text: str) -> str:
    kind = ""
    with contextlib.suppress(Exception):
        kind = context.user_data.pop("suno_preset_kind", "") or ""
    user_text = (user_text or "").strip()
    prefix = _suno_preset_prefix(kind)
    if prefix:
        return f"{prefix}\nUser details: {user_text}".strip()
    return user_text

def _suno_help_text() -> str:
    # Plain text intentionally: Telegram legacy Markdown can crash callback editing
    # on slash/backtick-heavy help texts on some clients.
    return (
        "рҹҺө РңСғР·СӢРәР° / Suno\n"
        "РңРҫР¶РҪРҫ СҒРҫР·РҙР°РІР°СӮСҢ: РҝРөСҒРҪРё СҒ РІРҫРәР°Р»РҫРј, РҙР¶РёРҪРіР»СӢ, РёРҪСӮСҖРҫ/Р°СғСӮСҖРҫ РҙР»СҸ Reels, С„РҫРҪРҫРІСғСҺ РјСғР·СӢРәСғ, РјРёРҪСғСҒРҫРІРәРё Рё РҙРөРјРҫ-СӮСҖРөРәРё.\n\n"
        "Р§СӮРҫ РҪР°РҝРёСҒР°СӮСҢ:\n"
        "вҖў Р¶Р°РҪСҖ: pop, rap, house, cinematic, acoustic;\n"
        "вҖў РҪР°СҒСӮСҖРҫРөРҪРёРө: luxury, СҚРҪРөСҖРіРёСҮРҪРҫ, РҙСҖР°РјР°СӮРёСҮРҪРҫ, СҖРҫРјР°РҪСӮРёСҮРҪРҫ;\n"
        "вҖў СҸР·СӢРә Рё РІРҫРәР°Р»: СҖСғСҒСҒРәРёР№ РјСғР¶СҒРәРҫР№ РІРҫРәР°Р» / Р°РҪРіР»РёР№СҒРәРёР№ Р¶РөРҪСҒРәРёР№ РІРҫРәР°Р» / instrumental;\n"
        "вҖў СӮРөРәСҒСӮ РәСғРҝР»РөСӮР°/РҝСҖРёРҝРөРІР° РёР»Рё СӮРөРјСғ РҝРөСҒРҪРё;\n"
        "вҖў РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ/С„РҫСҖРјР°СӮ: intro 15 СҒРөРә, jingle, full song.\n\n"
        "РҹСҖРёРјРөСҖ: СҒРҙРөР»Р°Р№ РҙР¶РёРҪРіР» 20 СҒРөРә РҙР»СҸ Р°РіРөРҪСӮСҒСӮРІР° РҪРөРҙРІРёР¶РёРјРҫСҒСӮРё РҪР° РЎР°РјСғРё, tropical house, luxury, РұРөР· РІРҫРәР°Р»Р°.\n\n"
        "РҹРҫСҒР»Рө СҚСӮРҫРіРҫ СҸ РҫСӮРҝСҖР°РІР»СҺ Р·Р°РҙР°СҮСғ РІ Suno. РһРұСӢСҮРҪРҫ РіРҫСӮРҫРІРҪРҫСҒСӮСҢ 1вҖ“5 РјРёРҪСғСӮ."
    )

async def _show_suno_help_from_callback(q, context, reply_markup=None, submenu: bool = False):
    """Show Suno instructions safely from any callback.

    Callback buttons should never fall into the global error handler only because
    Telegram refused to edit an old/unchanged message or parse markdown.
    """
    context.user_data["awaiting_suno_brief"] = True
    text = _suno_submenu_text() if submenu else _suno_help_text()
    kb = reply_markup or (_suno_menu_kb() if submenu else _mode_kb("fun"))
    try:
        await q.message.reply_text(text, reply_markup=kb)
    except Exception as e:
        log.warning("Suno menu edit failed, fallback to reply: %s", e)
        with contextlib.suppress(Exception):
            await q.message.reply_text(text, reply_markup=kb)


async def _safe_suno_callback_reply(q, text: str, reply_markup=None, *, edit: bool = True):
    """Р‘РөР·РҫРҝР°СҒРҪСӢР№ РІСӢРІРҫРҙ Suno-РјРөРҪСҺ.

    РқР° СҮР°СҒСӮРё Р°РәРәР°СғРҪСӮРҫРІ Telegram РјРҫР¶РөСӮ РҫСӮРәР°Р·Р°СӮСҢ РІ edit_message_text РҙР»СҸ СҒСӮР°СҖРҫРіРҫ/РәСҚСҲРёСҖРҫРІР°РҪРҪРҫРіРҫ
    СҒРҫРҫРұСүРөРҪРёСҸ РёР»Рё РҝСҖРё РҝРҫРІСӮРҫСҖРҪРҫРј РҪР°Р¶Р°СӮРёРё inline-РәРҪРҫРҝРәРё. Р Р°РҪСҢСҲРө СҚСӮРҫ СғС…РҫРҙРёР»Рҫ РІ РҫРұСүРёР№ error-handler
    Рё РҝРҫР»СҢР·РҫРІР°СӮРөР»СҢ РІРёРҙРөР» В«РЈРҝСҒВ». РўРөРҝРөСҖСҢ Suno-РәРҪРҫРҝРәРё РҪРө РҝР°РҙР°СҺСӮ: РөСҒР»Рё edit РҪРө РҝСҖРҫСҲС‘Р», РҫСӮРҝСҖР°РІР»СҸРөРј
    РҪРҫРІРҫРө СҒРҫРҫРұСүРөРҪРёРө.
    """
    try:
        if edit:
            await q.message.reply_text(text, reply_markup=reply_markup)
            return
    except Exception as e:
        log.warning("Suno callback edit failed, fallback to reply. data=%s err=%s", getattr(q, "data", ""), e)
    try:
        await q.message.reply_text(text, reply_markup=reply_markup)
    except Exception as e:
        log.warning("Suno callback reply failed: %s", e)


async def on_cb_suno(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    data = (q.data or "").strip()
    uid = q.from_user.id

    try:
        if data == "suno:back_fun":
            context.user_data.pop("awaiting_suno_brief", None)
            context.user_data.pop("suno_preset_kind", None)
            _set_mode_clean(uid, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "")
            with contextlib.suppress(Exception):
                await q.answer("РқР°Р·Р°Рҙ")
            await _safe_suno_callback_reply(q, _mode_desc("fun"), reply_markup=_mode_kb("fun"))
            return

        if data == "suno:free":
            _clear_transient_flows(context)
            _set_mode_clean(uid, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "suno_music")
            context.user_data["awaiting_suno_brief"] = True
            with contextlib.suppress(Exception):
                await q.answer("РЎРІРҫРұРҫРҙРҪСӢР№ Р·Р°РҝСҖРҫСҒ")
            await _safe_suno_callback_reply(
                q,
                "вңҚпёҸ РЎРІРҫРұРҫРҙРҪСӢР№ Р·Р°РҝСҖРҫСҒ РҙР»СҸ Suno\n\n"
                "РқР°РҝРёСҲРёСӮРө РҫРҙРҪРёРј СҒРҫРҫРұСүРөРҪРёРөРј, РәР°РәСғСҺ РјСғР·СӢРәСғ СҒРҙРөР»Р°СӮСҢ: Р¶Р°РҪСҖ, РҪР°СҒСӮСҖРҫРөРҪРёРө, СҸР·СӢРә/РІРҫРәР°Р», РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ Рё РҪР°Р·РҪР°СҮРөРҪРёРө.\n\n"
                "РҹСҖРёРјРөСҖ: СҒРҙРөР»Р°Р№ РҙР¶РёРҪРіР» 20 СҒРөРәСғРҪРҙ РҙР»СҸ Р°РіРөРҪСӮСҒСӮРІР° РҪРөРҙРІРёР¶РёРјРҫСҒСӮРё РҪР° РЎР°РјСғРё, tropical house, luxury, РұРөР· РІРҫРәР°Р»Р°.",
                reply_markup=_suno_menu_kb(),
            )
            return

        if data.startswith("suno:preset:"):
            kind = data.split(":", 2)[2]
            _clear_transient_flows(context)
            _set_mode_clean(uid, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "suno_music")
            context.user_data["awaiting_suno_brief"] = True
            context.user_data["suno_preset_kind"] = kind
            with contextlib.suppress(Exception):
                await q.answer("РҹСҖРөСҒРөСӮ Suno")
            # Р”Р»СҸ РҝСҖРөСҒРөСӮРҫРІ Р»СғСҮСҲРө РҫСӮРҝСҖР°РІР»СҸСӮСҢ РҪРҫРІРҫРө СҒРҫРҫРұСүРөРҪРёРө, Р° РҪРө СҖРөРҙР°РәСӮРёСҖРҫРІР°СӮСҢ СҒСӮР°СҖРҫРө: СӮР°Рә РҪРө Р»РҫРІРёРј
            # Telegram BadRequest РҪР° СҒСӮРҫСҖРҫРҪРҪРёС… Р°РәРәР°СғРҪСӮР°С…/СҒСӮР°СҖСӢС… inline-СҒРҫРҫРұСүРөРҪРёСҸС….
            await _safe_suno_callback_reply(q, _suno_preset_instruction(kind), reply_markup=_suno_menu_kb(), edit=False)
            return

        with contextlib.suppress(Exception):
            await q.answer()
    except Exception as e:
        log.exception("Suno callback failed safely. data=%s uid=%s err=%s", data, uid, e)
        with contextlib.suppress(Exception):
            await q.answer("РһСҲРёРұРәР° Suno-РјРөРҪСҺ", show_alert=False)
        with contextlib.suppress(Exception):
            await q.message.reply_text(
                "вҡ пёҸ РқРө СғРҙР°Р»РҫСҒСҢ РҫСӮРәСҖСӢСӮСҢ СҚСӮРҫСӮ РҝСғРҪРәСӮ Suno. РқР°Р¶РјРёСӮРө В«рҹҺө РңСғР·СӢРәР° / РҝРөСҒРҪСҸВ» РөСүС‘ СҖР°Р· РёР»Рё РҫСӮРҝСҖР°РІСҢСӮРө РҫРҝРёСҒР°РҪРёРө РҝРөСҒРҪРё СӮРөРәСҒСӮРҫРј.",
                reply_markup=_suno_menu_kb(),
            )


async def cmd_music(update: Update, context: ContextTypes.DEFAULT_TYPE):
    brief = " ".join(context.args).strip() if context.args else ""
    if not brief:
        _set_mode_clean(update.effective_user.id, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "suno_music")
        context.user_data["awaiting_suno_brief"] = True
        await update.effective_message.reply_text(_suno_submenu_text(), reply_markup=_suno_menu_kb())
        return
    await _run_suno_music(update, context, brief)

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ Photo quick actions в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def photo_quick_actions_kb():
    # Р”Р»РёРҪРҪСӢРө РҝРҫРҙРҝРёСҒРё вҖ” РҫСӮРҙРөР»СҢРҪСӢРјРё СҒСӮСҖРҫРәР°РјРё, СҮСӮРҫРұСӢ РҪРө РҫРұСҖРөР·Р°Р»РёСҒСҢ РІ Telegram.
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("вңЁ РһР¶РёРІРёСӮСҢ СҮРөСҖРөР· Runway", callback_data="pedit:revive_runway")],
        [InlineKeyboardButton("вңЁ РһР¶РёРІРёСӮСҢ СҮРөСҖРөР· Kling", callback_data="pedit:revive_kling")],
        [InlineKeyboardButton("вңЁ Sora 2 РұРөР· Р»СҺРҙРөР№", callback_data="pedit:revive_sora")],
        [InlineKeyboardButton("рҹ—Ј Р“РҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ", callback_data="pedit:avatar")],
        [InlineKeyboardButton("рҹҺө РӨРҫСӮРҫ вҶ’ РІРёРҙРөРҫРәР»РёРҝ", callback_data="pedit:photoclip")],
        [InlineKeyboardButton("рҹҺӨ РҡР»РёРҝ СҒ РІРҫРәР°Р»РҫРј (1 СҮРөР»РҫРІРөРә)", callback_data="pedit:vocalclip")],
        [InlineKeyboardButton("рҹӨі AI-СҒРөР»С„Рё СҒРҫ Р·РІРөР·РҙРҫР№", callback_data="pedit:aiselfie")],
        [InlineKeyboardButton("рҹҺӯ Р—Р°РјРөРҪР° Р»РёСҶР° РҪР° С„РҫСӮРҫ", callback_data="pedit:faceswap")],
        [InlineKeyboardButton("рҹ§ј РЈРҙР°Р»РёСӮСҢ С„РҫРҪ РҪР° С„РҫСӮРҫ", callback_data="pedit:removebg")],
        [InlineKeyboardButton("рҹ–ј Р—Р°РјРөРҪРёСӮСҢ С„РҫРҪ РҪР° С„РҫСӮРҫ", callback_data="pedit:replacebg")],
        [InlineKeyboardButton("рҹ§Ҫ РЈРҙР°Р»РёСӮСҢ РІРҫРҙСҸРҪРҫР№ Р·РҪР°Рә", callback_data="pedit:retouch")],
        [InlineKeyboardButton("рҹ§ӯ Р Р°СҒСҲРёСҖРёСӮСҢ РәР°РҙСҖ", callback_data="pedit:outpaint")],
        [InlineKeyboardButton("рҹ“Ҫ Р Р°СҒРәР°РҙСҖРҫРІРәР° С„РҫСӮРҫ", callback_data="pedit:story")],
        [InlineKeyboardButton("рҹ–Ң РҡР°СҖСӮРёРҪРәР° РёР· СӮРөРәСҒСӮР°", callback_data="pedit:lumaimg")],
        [InlineKeyboardButton("рҹ‘Ғ РҗРҪР°Р»РёР· С„РҫСӮРҫ", callback_data="pedit:vision")],
    ])

_photo_cache = {}      # user_id -> bytes
_photo_url_cache = {}  # user_id -> Telegram file URL when available

def _cache_photo(user_id: int, data: bytes, file_url: str | None = None):
    try:
        _photo_cache[user_id] = data
        if file_url:
            _photo_url_cache[user_id] = str(file_url)
    except Exception:
        pass

def _get_cached_photo(user_id: int) -> bytes | None:
    return _photo_cache.get(user_id)

def _get_cached_photo_url(user_id: int) -> str:
    return _photo_url_cache.get(user_id, "") or ""


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ Kling Avatar / PhotoвҶ’music clip в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def _is_avatar_intent(text: str) -> bool:
    t = (text or "").lower()
    return bool(re.search(r"(РіРҫРІРҫСҖСҸСү\w*\s+Р°РІР°СӮР°СҖ|Р°РІР°СӮР°СҖ\w*\s+РіРҫРІРҫСҖ|talking\s*avatar|avatar\s*video|lip[-\s]?sync|Р»РёРҝ[-\s]?СҒРёРҪРә|СҒРёРҪС…СҖРҫРҪ\w*\s+РіСғРұ)", t, re.I))


def _is_photo_clip_intent(text: str) -> bool:
    t = (text or "").lower()
    return bool(re.search(r"(С„РҫСӮРҫ\s*(?:РІ|вҶ’|-)\s*РІРёРҙРөРҫРәР»РёРҝ|С„РҫСӮРҫ\s*(?:РІ|вҶ’|-)\s*РәР»РёРҝ|РІРёРҙРөРҫРәР»РёРҝ\s+РёР·\s+С„РҫСӮРҫ|РәР»РёРҝ\s+РёР·\s+С„РҫСӮРҫ|music\s*video|photo\s*clip)", t, re.I))


def _clean_avatar_script(text: str) -> str:
    t = (text or "").strip()
    t = re.sub(r"^(?:СҒРҙРөР»Р°Р№|СҒРҫР·РҙР°Р№|СҒРіРөРҪРөСҖРёСҖСғР№)?\s*(?:РіРҫРІРҫСҖСҸСү\w*\s+Р°РІР°СӮР°СҖ|Р°РІР°СӮР°СҖ|talking\s*avatar|lip[-\s]?sync)\s*[:пјҡ,-]?\s*", "", t, flags=re.I)
    return t.strip()


def _clean_photo_clip_prompt(text: str) -> str:
    t = (text or "").strip()
    t = re.sub(r"^(?:СҒРҙРөР»Р°Р№|СҒРҫР·РҙР°Р№|СҒРіРөРҪРөСҖРёСҖСғР№)?\s*(?:С„РҫСӮРҫ\s*(?:РІ|вҶ’|-)\s*(?:РІРёРҙРөРҫРәР»РёРҝ|РәР»РёРҝ)|РІРёРҙРөРҫРәР»РёРҝ\s+РёР·\s+С„РҫСӮРҫ|РәР»РёРҝ\s+РёР·\s+С„РҫСӮРҫ|music\s*video|photo\s*clip)\s*[:пјҡ,-]?\s*", "", t, flags=re.I)
    return t.strip()


def _set_avatar_wait(context: ContextTypes.DEFAULT_TYPE):
    context.user_data["awaiting_avatar_script"] = True
    context.user_data.pop("awaiting_avatar_voice_choice", None)
    context.user_data.pop("awaiting_photo_clip_prompt", None)


def _set_avatar_voice_choice_wait(context: ContextTypes.DEFAULT_TYPE):
    context.user_data["awaiting_avatar_voice_choice"] = True
    context.user_data.pop("awaiting_avatar_script", None)
    context.user_data.pop("awaiting_photo_clip_prompt", None)


def _set_photo_clip_wait(context: ContextTypes.DEFAULT_TYPE):
    context.user_data["awaiting_photo_clip_prompt"] = True
    context.user_data.pop("awaiting_avatar_script", None)


def _clear_avatar_wait(context: ContextTypes.DEFAULT_TYPE):
    context.user_data.pop("awaiting_avatar_script", None)
    context.user_data.pop("awaiting_avatar_voice_choice", None)
    context.user_data.pop("awaiting_avatar_photo", None)
    context.user_data.pop("avatar_pending_script", None)


def _clear_photo_clip_wait(context: ContextTypes.DEFAULT_TYPE):
    context.user_data.pop("awaiting_photo_clip_prompt", None)
    context.user_data.pop("awaiting_photo_clip_photo", None)


def _telegram_file_public_url(file_path: str) -> str:
    fp = (file_path or "").strip()
    if not fp:
        return ""
    if fp.startswith("http://") or fp.startswith("https://"):
        return fp
    return f"https://api.telegram.org/file/bot{BOT_TOKEN}/{fp.lstrip('/')}"


def _avatar_audio_supported_filename(filename: str, mime: str = "") -> bool:
    f = (filename or "").lower()
    m = (mime or "").lower()
    return (
        f.endswith((".mp3", ".wav", ".m4a", ".aac"))
        or "mpeg" in m or "mp3" in m or "wav" in m or "m4a" in m or "aac" in m
    )


def _avatar_tts_voice_get(context: ContextTypes.DEFAULT_TYPE | None = None) -> str:
    try:
        v = ((context.user_data.get("avatar_tts_voice") if context else "") or "").strip()
    except Exception:
        v = ""
    return v or AVATAR_TTS_DEFAULT_VOICE or OPENAI_TTS_VOICE or "alloy"


def _avatar_tts_voice_label(voice: str) -> str:
    labels = {
        "nova": "Nova вҖ” РјСҸРіРәРёР№ Р¶РөРҪСҒРәРёР№",
        "alloy": "Alloy вҖ” РҪРөР№СӮСҖР°Р»СҢРҪСӢР№",
        "onyx": "Onyx вҖ” РҪРёР·РәРёР№ РјСғР¶СҒРәРҫР№",
        "shimmer": "Shimmer вҖ” СҒРІРөСӮР»СӢР№ Р¶РөРҪСҒРәРёР№",
        "fable": "Fable вҖ” СҒСӮРҫСҖРёСӮРөР»Р»РёРҪРі",
        "echo": "Echo вҖ” РјСғР¶СҒРәРҫР№",
        "sage": "Sage вҖ” СҒРҝРҫРәРҫР№РҪСӢР№",
        "ash": "Ash вҖ” РҝР»РҫСӮРҪСӢР№",
        "coral": "Coral вҖ” СҸСҖРәРёР№",
        "verse": "Verse вҖ” РІСӢСҖР°Р·РёСӮРөР»СҢРҪСӢР№",
        "ballad": "Ballad вҖ” РјСҸРіРәРёР№",
    }
    return labels.get((voice or "").strip().lower(), (voice or "alloy").strip())


async def _upload_bytes_to_telegram_file_url(update: Update, context: ContextTypes.DEFAULT_TYPE, raw: bytes, filename: str, caption: str = "") -> str:
    if not raw:
        return ""
    bio = BytesIO(raw)
    bio.seek(0)
    bio.name = filename
    sent = await update.effective_message.reply_document(
        document=InputFile(bio),
        caption=caption or "РӨР°Р№Р» РҝРҫРҙРіРҫСӮРҫРІР»РөРҪ РҙР»СҸ РіРөРҪРөСҖР°СҶРёРё.",
    )
    media = getattr(sent, "document", None) or getattr(sent, "audio", None) or getattr(sent, "voice", None)
    if not media:
        return ""
    tg_file = await context.bot.get_file(media.file_id)
    return _telegram_file_public_url(getattr(tg_file, "file_path", "") or "")


async def _text_to_public_mp3_url(update: Update, context: ContextTypes.DEFAULT_TYPE, text: str) -> str:
    text = (text or "").strip()
    if not text or not OPENAI_TTS_KEY:
        return ""
    mp3 = await asyncio.to_thread(_tts_bytes_sync, text[:TTS_MAX_CHARS], "mp3", _avatar_tts_voice_get(context))
    if not mp3:
        return ""
    return await _upload_bytes_to_telegram_file_url(
        update, context, mp3, "avatar_speech.mp3", "рҹ”Ҡ РһР·РІСғСҮРәР° РҙР»СҸ РіРҫРІРҫСҖСҸСүРөРіРҫ Р°РІР°СӮР°СҖР° РҝРҫРҙРіРҫСӮРҫРІР»РөРҪР°."
    )


def _extract_task_status(js: object) -> str:
    if not isinstance(js, dict):
        return ""
    for k in ("status", "state", "task_status"):
        v = js.get(k)
        if v is not None:
            return str(v).lower().strip()
    for k in ("data", "result", "response", "payload"):
        v = js.get(k)
        if isinstance(v, dict):
            st = _extract_task_status(v)
            if st:
                return st
    return ""


def _extract_task_id_from_response(js: object) -> str:
    if not isinstance(js, dict):
        return ""
    for k in ("id", "task_id", "generation_id", "video_id", "audio_id"):
        v = js.get(k)
        if v:
            return str(v).strip()
    for k in ("data", "result", "response", "payload"):
        v = js.get(k)
        if isinstance(v, dict):
            tid = _extract_task_id_from_response(v)
            if tid:
                return tid
    return ""


async def _create_kling_tts_audio_ref(text: str) -> tuple[str, str, str]:
    """Returns (audio_id, linked_task_id, sound_url). Prefer OpenAI TTS for RU; this is fallback."""
    if not (KLING_API_KEY or COMET_API_KEY):
        return "", "", ""
    text = (text or "").strip()[:1000]
    if not text:
        return "", "", ""
    headers = {"Authorization": f"Bearer {KLING_API_KEY or COMET_API_KEY}", "Content-Type": "application/json", "Accept": "application/json"}
    payload = {
        "text": text,
        "voice_id": KLING_TTS_VOICE_ID,
        "voice_language": KLING_TTS_LANGUAGE,
        "voice_speed": max(0.8, min(2.0, float(KLING_TTS_SPEED or 1.0))),
    }
    async with httpx.AsyncClient(timeout=90.0) as client:
        r = await client.post(f"{COMET_BASE_URL}{KLING_TTS_CREATE_PATH}", headers=headers, json=payload)
        if r.status_code >= 400:
            log.warning("Kling TTS create failed: %s", _api_error_preview(r))
            return "", "", ""
        js = r.json() or {}
        sound_url = _extract_first_url(js) or ""
        task_id = _extract_task_id_from_response(js)
        audio_id = str((js.get("data") or {}).get("audio_id") or js.get("audio_id") or task_id or "").strip()
        if sound_url or audio_id:
            return audio_id, task_id, sound_url
        return "", task_id, ""


def _avatar_prompt_from_text(script_text: str = "") -> str:
    base = KLING_AVATAR_PROMPT or "The person talks naturally to camera, realistic facial motion, accurate lip sync."
    script_text = (script_text or "").strip()
    if script_text:
        return (base + " Speech content: " + script_text[:500]).strip()
    return base


async def _run_kling_avatar(
    update: Update,
    context: ContextTypes.DEFAULT_TYPE,
    img_bytes: bytes,
    script_text: str = "",
    audio_bytes: bytes | None = None,
    audio_filename: str = "",
    audio_file_url: str = "",
    audio_mime: str = "",
    avatar_prompt_override: str = "",
    max_wait_s: int | None = None,
):
    if not (KLING_API_KEY or COMET_API_KEY):
        await update.effective_message.reply_text("вқҢ Kling Avatar: COMET_API_KEY/KLING_API_KEY РҪРө Р·Р°РҙР°РҪ РІ Render ENV.")
        return False

    await context.bot.send_chat_action(update.effective_chat.id, ChatAction.RECORD_VIDEO)
    raw_b64 = base64.b64encode(img_bytes).decode("ascii")
    sound_file = ""
    audio_id = ""
    linked_task_id = ""
    script_text = (script_text or "").strip()

    if audio_file_url and _avatar_audio_supported_filename(audio_filename, audio_mime):
        sound_file = _telegram_file_public_url(audio_file_url)

    if not sound_file and audio_bytes and not script_text:
        try:
            script_text = await _stt_transcribe_bytes(audio_filename or "voice.ogg", audio_bytes)
        except Exception as e:
            log.warning("Avatar voice STT failed: %s", e)
            script_text = ""

    if not sound_file and script_text and AVATAR_TTS_PROVIDER in ("openai", "auto", ""):
        sound_file = await _text_to_public_mp3_url(update, context, script_text)

    if not sound_file and script_text:
        audio_id, linked_task_id, tts_url = await _create_kling_tts_audio_ref(script_text)
        if tts_url:
            sound_file = tts_url

    if not sound_file and not audio_id:
        await update.effective_message.reply_text(
            "вқҢ РқРө РҝРҫР»СғСҮРёР»РҫСҒСҢ РҝРҫРҙРіРҫСӮРҫРІРёСӮСҢ Р°СғРҙРёРҫ РҙР»СҸ Р°РІР°СӮР°СҖР°. РҹСҖРёСҲР»РёСӮРө MP3/WAV/M4A/AAC 2вҖ“60 СҒРөРәСғРҪРҙ РёР»Рё РҙРҫРұР°РІСҢСӮРө OPENAI_TTS_KEY РІ Render ENV."
        )
        return False

    prompt = (avatar_prompt_override or "").strip() or _avatar_prompt_from_text(script_text)
    payload = {
        "image": raw_b64,
        "prompt": prompt,
        "mode": KLING_AVATAR_MODE,
    }
    if sound_file:
        payload["sound_file"] = sound_file
    else:
        payload["audio_id"] = audio_id
        if linked_task_id:
            payload["task_id"] = linked_task_id

    return await _create_and_poll_i2v(
        update,
        COMET_BASE_URL,
        KLING_API_KEY or COMET_API_KEY,
        [(KLING_AVATAR_CREATE_PATH, payload), ("/kling/v1/videos/avatar/image2video", payload)],
        [KLING_AVATAR_STATUS_PATH, "/kling/v1/videos/avatar/image2video/{id}", "/kling/v1/videos/{id}", "/v1/tasks/{id}"],
        "Kling talking avatar",
        max_wait_s=max_wait_s,
    )


async def _start_talking_avatar(
    update: Update,
    context: ContextTypes.DEFAULT_TYPE,
    img_bytes: bytes,
    script_text: str = "",
    audio_bytes: bytes | None = None,
    audio_filename: str = "",
    audio_file_url: str = "",
    audio_mime: str = "",
    avatar_prompt_override: str = "",
):
    async def _go():
        ok = await _run_kling_avatar(update, context, img_bytes, script_text, audio_bytes, audio_filename, audio_file_url, audio_mime, avatar_prompt_override)
        return bool(ok)

    await _try_pay_then_do(
        update, context, update.effective_user.id,
        "runway", AVATAR_UNIT_COST_USD, _go,
        remember_kind="kling_talking_avatar",
        remember_payload={"script": (script_text or "")[:500], "audio_filename": audio_filename},
    )


def _photo_clip_prompt(user_prompt: str, base_seconds: int = 10) -> str:
    user_prompt = (user_prompt or "").strip()
    if not user_prompt:
        user_prompt = "СҚРҪРөСҖРіРёСҮРҪСӢР№ РјСғР·СӢРәР°Р»СҢРҪСӢР№ РәР»РёРҝ, РіРөСҖРҫР№ РІ РәР°РҙСҖРө, РҝР»Р°РІРҪРҫРө РҙРІРёР¶РөРҪРёРө РәР°РјРөСҖСӢ, cinematic, social media, high quality"
    return (
        "Create ONLY the visual part of a short cinematic music-video style clip from this person photo. "
        "Keep the same identity and face, avoid deformation, natural body/head motion, dynamic but clean camera, premium lighting. "
        "Do not create subtitles, do not add text overlays, do not add fake logos. "
        f"Base visual duration requested from provider: {base_seconds} seconds. "
        "Music/audio will be added separately in post-production, so focus on visual motion only. "
        "User direction: " + user_prompt[:600]
    )


def _photo_clip_target_duration(user_prompt: str) -> int:
    txt = (user_prompt or "").lower()
    m = re.search(r"(\d{1,2})\s*(?:СҒРөРә|СҒРөРәСғРҪРҙ|second|seconds|s)\b", txt, re.I)
    if m:
        try:
            val = int(m.group(1))
        except Exception:
            val = PHOTO_CLIP_DEFAULT_DURATION_S
    else:
        val = PHOTO_CLIP_DEFAULT_DURATION_S
    return max(5, min(int(PHOTO_CLIP_MAX_DURATION_S or 30), val))


def _ffmpeg_exe() -> str:
    try:
        import imageio_ffmpeg  # type: ignore
        exe = imageio_ffmpeg.get_ffmpeg_exe()
        if exe:
            return exe
    except Exception:
        pass
    return shutil.which("ffmpeg") or "ffmpeg"


async def _download_binary_from_url(client: httpx.AsyncClient, url: str, accept: str = "*/*", timeout_s: float = 300.0) -> bytes | None:
    if not url:
        return None
    try:
        headers = {"User-Agent": "Mozilla/5.0 (compatible; GPT5ProBot/1.0)", "Accept": accept}
        r = await client.get(url, headers=headers, timeout=timeout_s, follow_redirects=True)
        if r.status_code >= 400:
            log.warning("download_binary failed %s: %s", r.status_code, _api_error_preview(r))
            return None
        ctype = (r.headers.get("content-type") or "").lower()
        if not r.content or len(r.content) < 512:
            log.warning("download_binary empty response url=%s bytes=%s", url[:120], len(r.content or b""))
            return None
        if "text/html" in ctype or "application/json" in ctype:
            log.warning("download_binary non-binary content-type=%s body=%s", ctype, r.text[:300])
            return None
        return bytes(r.content)
    except Exception as e:
        log.warning("download_binary exception url=%s: %s", url[:120], e)
        return None


async def _poll_video_task_for_bytes(
    client: httpx.AsyncClient,
    headers: dict,
    base_url: str,
    status_paths: list[str],
    task_id: str,
    caption: str,
    max_wait_s: int = 1200,
) -> bytes | None:
    started = time.time()
    last_body = ""
    while time.time() - started < max_wait_s:
        for path in status_paths:
            try:
                url = f"{base_url}{path}".format(id=task_id)
                rs = await client.get(url, headers=headers, timeout=60.0, follow_redirects=True)
                if rs.status_code >= 400:
                    last_body = f"{rs.status_code}: {_api_error_preview(rs)}"
                    continue
                try:
                    js = rs.json() or {}
                except Exception:
                    js = {}
                st = str(js.get("status") or js.get("state") or js.get("task_status") or "").lower()
                ready_url = _extract_first_url(js.get("output")) or _extract_first_url(js.get("outputs")) or _extract_first_url(js.get("assets")) or _extract_first_url(js.get("data")) or _extract_first_url(js)
                if ready_url and (st in ("", "completed", "succeeded", "success", "finished", "ready", "done", "succeed") or not st):
                    content = await _download_binary_from_url(client, ready_url, accept="video/mp4,video/*,*/*;q=0.8", timeout_s=300.0)
                    if content:
                        return content
                    last_body = f"ready_url download failed: {ready_url[:200]}"
                if st in ("failed", "fail", "error", "canceled", "cancelled", "rejected"):
                    raise RuntimeError(f"{caption}: render failed {json.dumps(js, ensure_ascii=False)[:900]}")
                last_body = json.dumps(js, ensure_ascii=False)[:700]
            except Exception as e:
                last_body = str(e)
                continue
        await asyncio.sleep(VIDEO_POLL_DELAY_S)
    raise TimeoutError(f"{caption}: timeout, last={last_body[:700]}")


async def _create_and_poll_i2v_bytes(
    base_url: str,
    api_key: str,
    create_payloads: list[tuple[str, dict]],
    status_paths: list[str],
    caption: str,
    max_wait_s: int = 1200,
) -> bytes | None:
    if not api_key:
        raise RuntimeError(f"{caption}: API key missing")
    auth_headers = {"Authorization": f"Bearer {api_key}", "Accept": "application/json"}
    last_err = ""
    async with httpx.AsyncClient(timeout=120.0, follow_redirects=True) as client:
        # v61: for photoвҶ’music clip do not try multiple create paths after success and do not send direct links.
        for path, payload in create_payloads[:1]:
            try:
                headers = dict(auth_headers)
                headers["Content-Type"] = "application/json"
                if str(path).startswith("/runwayml/"):
                    headers["X-Runway-Version"] = RUNWAY_API_VERSION or "2024-11-06"
                r = await client.post(f"{base_url}{path}", headers=headers, json=payload)
                if r.status_code >= 400:
                    last_err = f"POST {path} вҶ’ {r.status_code}: {_api_error_preview(r)}"
                    log.warning("%s create failed: %s", caption, last_err)
                    continue
                try:
                    js = r.json() or {}
                except Exception:
                    js = {}
                ready_url = _extract_first_url(js.get("output")) or _extract_first_url(js.get("outputs")) or _extract_first_url(js.get("assets")) or _extract_first_url(js.get("data")) or _extract_first_url(js)
                if ready_url:
                    content = await _download_binary_from_url(client, ready_url, accept="video/mp4,video/*,*/*;q=0.8", timeout_s=300.0)
                    if content:
                        return content
                task_id = str(
                    js.get("id") or js.get("task_id") or js.get("generation_id") or js.get("video_id") or js.get("taskId") or js.get("taskID") or js.get("request_id") or js.get("uuid") or ""
                ).strip()
                if not task_id and isinstance(js.get("data"), dict):
                    d = js.get("data") or {}
                    task_id = str(d.get("id") or d.get("task_id") or d.get("generation_id") or d.get("video_id") or d.get("taskId") or d.get("taskID") or d.get("request_id") or d.get("uuid") or "").strip()
                if not task_id and isinstance(js.get("result"), dict):
                    d = js.get("result") or {}
                    task_id = str(d.get("id") or d.get("task_id") or d.get("generation_id") or d.get("video_id") or d.get("taskId") or d.get("taskID") or d.get("request_id") or d.get("uuid") or "").strip()
                if not task_id:
                    last_err = f"POST {path}: no task id in {json.dumps(js, ensure_ascii=False)[:700]}"
                    continue
                log.info("%s accepted task_id=%s", caption, task_id)
                return await _poll_video_task_for_bytes(client, headers, base_url, status_paths, task_id, caption, max_wait_s=max_wait_s)
            except Exception as e:
                last_err = str(e)
                log.warning("%s create/poll exception: %s", caption, e)
                continue
    raise RuntimeError(last_err or f"{caption}: no result")


async def _run_kling_photo_clip_result(img_bytes: bytes, prompt: str, duration_s: int, aspect: str) -> bytes | None:
    if not (KLING_API_KEY or COMET_API_KEY):
        raise RuntimeError("Kling photoвҶ’clip: COMET_API_KEY/KLING_API_KEY РҪРө Р·Р°РҙР°РҪ")
    raw_b64 = base64.b64encode(img_bytes).decode("ascii")
    base_duration = str(_duration_for_engine("kling", min(10, duration_s)))
    payload = {
        "image": raw_b64,
        "prompt": _photo_clip_prompt(prompt, int(base_duration)),
        "model_name": KLING_MODEL,
        "model": KLING_MODEL,
        "mode": PHOTO_CLIP_MODE,
        "duration": base_duration,
        "aspect_ratio": aspect,
        # v61: music is added by Suno+ffmpeg. Native provider sound is not reliable here.
        "sound": "off",
    }
    return await _create_and_poll_i2v_bytes(
        COMET_BASE_URL,
        KLING_API_KEY or COMET_API_KEY,
        [(KLING_CREATE_PATH, payload), ("/kling/v1/videos/image2video", payload)],
        ["/kling/v1/videos/image2video/{id}", KLING_STATUS_PATH, "/kling/v1/videos/{id}", "/v1/tasks/{id}"],
        "Kling photoвҶ’music clip",
        max_wait_s=max(LUMA_MAX_WAIT_S, RUNWAY_MAX_WAIT_S),
    )


async def _run_suno_music_result_bytes(update: Update, brief: str) -> bytes | None:
    brief = (brief or "").strip() or "dynamic catchy music video song, cinematic, social media ready"
    if not (SUNO_AUTO_FOR_PHOTO_CLIP and SUNO_ENABLED and SUNO_API_KEY):
        return None
    headers = {"Authorization": f"Bearer {SUNO_API_KEY}", "Content-Type": "application/json", "Accept": "application/json"}
    instrumental = bool(re.search(r"РёРҪСҒСӮСҖСғРјРөРҪСӮР°Р»|instrumental|РұРөР· РІРҫРәР°Р»Р°|РұРөР· РіРҫР»РҫСҒР°|РјРёРҪСғСҒРҫРІ|background music|С„РҫРҪРҫРІР°", brief, re.I))
    base_payload = {"mv": SUNO_MODEL, "gpt_description_prompt": brief}
    if instrumental:
        base_payload.update({"prompt": "", "make_instrumental": True})
    status_paths = [SUNO_STATUS_PATH, "/suno/fetch/{id}", "/suno/v1/music/{id}", "/api/v1/task/{id}", "/v1/tasks/{id}"]
    async with httpx.AsyncClient(timeout=120.0, follow_redirects=True) as client:
        last_err = ""
        for path in (SUNO_CREATE_PATH, "/suno/submit/music"):
            try:
                r = await client.post(f"{SUNO_BASE_URL}{path}", headers=headers, json=base_payload)
                if r.status_code >= 400:
                    last_err = f"POST {path} вҶ’ {r.status_code}: {_api_error_preview(r)}"
                    continue
                try:
                    js = r.json() or {}
                except Exception:
                    js = {}
                url = _extract_first_url(js.get("audio_url")) or _extract_first_url(js.get("audio")) or _extract_first_url(js.get("output")) or _extract_first_url(js.get("data")) or _extract_first_url(js)
                if url:
                    audio = await _download_binary_from_url(client, url, accept="audio/mpeg,audio/*,*/*;q=0.8", timeout_s=300.0)
                    if audio:
                        return audio
                task_id = str(js.get("id") or js.get("task_id") or js.get("taskId") or js.get("request_id") or "").strip()
                data_obj = js.get("data")
                if not task_id and isinstance(data_obj, str) and re.fullmatch(r"[A-Za-z0-9_.:-]{8,}", data_obj.strip()):
                    task_id = data_obj.strip()
                if not task_id and isinstance(data_obj, dict):
                    d = data_obj or {}
                    task_id = str(d.get("id") or d.get("task_id") or d.get("taskId") or d.get("request_id") or "").strip()
                if task_id:
                    started = time.time()
                    while time.time() - started < max(30, min(SUNO_TIMEOUT_S, PHOTO_CLIP_SUNO_FAST_TIMEOUT_S)):
                        for sp in status_paths:
                            try:
                                rs = await client.get(f"{SUNO_BASE_URL}{sp}".format(id=task_id), headers=headers, timeout=60.0)
                                if rs.status_code >= 400:
                                    continue
                                try:
                                    sj = rs.json() or {}
                                except Exception:
                                    sj = {}
                                data_nested = sj.get("data") if isinstance(sj, dict) else None
                                nested = data_nested if isinstance(data_nested, dict) else {}
                                url = _extract_first_url(sj.get("audio_url")) or _extract_first_url(sj.get("audio")) or _extract_first_url(sj.get("output")) or _extract_first_url(data_nested) or _extract_first_url(nested.get("data") if isinstance(nested, dict) else None) or _extract_first_url(sj)
                                st = str(sj.get("status") or sj.get("state") or sj.get("task_status") or nested.get("status") or "").lower()
                                if url and (st in ("", "completed", "succeeded", "success", "finished", "done", "ready") or not st):
                                    audio = await _download_binary_from_url(client, url, accept="audio/mpeg,audio/*,*/*;q=0.8", timeout_s=300.0)
                                    if audio:
                                        return audio
                                if st in ("failed", "fail", "error", "canceled", "cancelled", "rejected"):
                                    raise RuntimeError(json.dumps(sj, ensure_ascii=False)[:900])
                            except Exception as e:
                                last_err = str(e)
                                continue
                        await asyncio.sleep(SUNO_POLL_DELAY_S)
            except Exception as e:
                last_err = str(e)
                continue
        log.warning("Suno result bytes failed: %s", last_err)
    return None


def _mux_video_audio_sync(video_bytes: bytes, audio_bytes: bytes | None, target_duration_s: int) -> bytes | None:
    """Build one Telegram-safe MP4 from Kling video + Suno audio.
    v65: copy-first if small, then compressed fallback under FFMPEG_MUX_MAX_MB.
    """
    if not video_bytes or not audio_bytes:
        return None
    target_duration_s = max(5, min(int(PHOTO_CLIP_MAX_DURATION_S or 30), int(target_duration_s or PHOTO_CLIP_DEFAULT_DURATION_S or 15)))
    timeout_s = max(30, int(FFMPEG_MUX_TIMEOUT_S or 180))
    max_bytes = max(5, int(FFMPEG_MUX_MAX_MB or 45)) * 1024 * 1024
    started = time.time()
    try:
        ffmpeg = _ffmpeg_exe()
        with tempfile.TemporaryDirectory() as td:
            vin = os.path.join(td, "input.mp4")
            a_in = os.path.join(td, "audio.mp3")
            out_fast = os.path.join(td, "final_fast.mp4")
            out_compact = os.path.join(td, "final_compact.mp4")
            out_small = os.path.join(td, "final_small.mp4")
            with open(vin, "wb") as f:
                f.write(video_bytes)
            with open(a_in, "wb") as f:
                f.write(audio_bytes)

            attempts: list[tuple[str, list[str], str, bool]] = []
            if FFMPEG_MUX_COPY_FIRST:
                attempts.append((
                    "copy-first",
                    [
                        ffmpeg, "-y",
                        "-stream_loop", "-1", "-fflags", "+genpts", "-i", vin,
                        "-stream_loop", "-1", "-i", a_in,
                        "-t", str(target_duration_s),
                        "-map", "0:v:0", "-map", "1:a:0",
                        "-c:v", "copy",
                        "-c:a", "aac", "-b:a", FFMPEG_MUX_AUDIO_BITRATE,
                        "-shortest", "-movflags", "+faststart",
                        out_fast,
                    ],
                    out_fast,
                    True,
                ))
            # Telegram-safe compact encode. This is the primary fallback after copy-first.
            attempts.append((
                "compact-720p",
                [
                    ffmpeg, "-y",
                    "-stream_loop", "-1", "-fflags", "+genpts", "-i", vin,
                    "-stream_loop", "-1", "-i", a_in,
                    "-t", str(target_duration_s),
                    "-map", "0:v:0", "-map", "1:a:0",
                    "-vf", f"scale=-2:{int(FFMPEG_MUX_SCALE_HEIGHT or 720)},fps={int(FFMPEG_MUX_FPS or 24)}",
                    "-c:v", "libx264", "-preset", FFMPEG_MUX_REENCODE_PRESET,
                    "-crf", str(FFMPEG_MUX_CRF or "32"), "-pix_fmt", "yuv420p",
                    "-c:a", "aac", "-b:a", FFMPEG_MUX_AUDIO_BITRATE,
                    "-shortest", "-movflags", "+faststart",
                    out_compact,
                ],
                out_compact,
                False,
            ))
            # Last-resort smaller encode if the first compact file is still too large.
            attempts.append((
                "small-540p",
                [
                    ffmpeg, "-y",
                    "-stream_loop", "-1", "-fflags", "+genpts", "-i", vin,
                    "-stream_loop", "-1", "-i", a_in,
                    "-t", str(target_duration_s),
                    "-map", "0:v:0", "-map", "1:a:0",
                    "-vf", "scale=-2:540,fps=20",
                    "-c:v", "libx264", "-preset", "ultrafast",
                    "-crf", "35", "-pix_fmt", "yuv420p",
                    "-c:a", "aac", "-b:a", "96k",
                    "-shortest", "-movflags", "+faststart",
                    out_small,
                ],
                out_small,
                False,
            ))

            last_err = ""
            best_too_large: bytes | None = None
            for name, cmd, out_path, accept_only_if_small in attempts:
                left = max(20, timeout_s - int(time.time() - started))
                log.info("ffmpeg mux attempt=%s video=%s audio=%s duration=%s max_bytes=%s timeout_left=%s", name, len(video_bytes), len(audio_bytes), target_duration_s, max_bytes, left)
                try:
                    res = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=left)
                except subprocess.TimeoutExpired:
                    last_err = f"{name}: timeout after {left}s"
                    log.warning("ffmpeg mux timeout: %s", last_err)
                    continue
                if res.returncode != 0:
                    last_err = f"{name}: rc={res.returncode} stderr={res.stderr.decode('utf-8', 'ignore')[-1500:]}"
                    log.warning("ffmpeg mux failed: %s", last_err)
                    continue
                if not os.path.exists(out_path):
                    last_err = f"{name}: output missing"
                    continue
                with open(out_path, "rb") as f:
                    final = f.read()
                if not final or len(final) <= 1024:
                    last_err = f"{name}: empty output"
                    continue
                log.info("ffmpeg mux ok attempt=%s output=%s elapsed=%.1fs", name, len(final), time.time() - started)
                if len(final) <= max_bytes:
                    return final
                best_too_large = final
                last_err = f"{name}: output too large {len(final)} > {max_bytes}"
                log.warning("ffmpeg mux output too large: %s", last_err)
                if accept_only_if_small:
                    continue
            if best_too_large and PHOTO_CLIP_SEND_BASE_IF_MUX_FAILS:
                return best_too_large
            log.warning("ffmpeg mux all attempts failed/too-large: %s", last_err[-1500:])
            return None
    except Exception as e:
        log.warning("mux_video_audio exception: %s", e)
        return None





def _trim_audio_for_vocal_clip_sync(audio_bytes: bytes, max_seconds: int = 65) -> bytes:
    """Telegram/Suno often returns a full 2-4 minute song, while Kling Avatar is much more reliable
    on short audio. Trim/re-encode to a safe MP3 fragment for lip-sync.
    """
    if not audio_bytes:
        return audio_bytes
    max_seconds = max(8, min(90, int(max_seconds or 65)))
    try:
        ffmpeg = _ffmpeg_exe()
        with tempfile.TemporaryDirectory() as td:
            src = os.path.join(td, "suno_full.mp3")
            out = os.path.join(td, "suno_lipsync_safe.mp3")
            with open(src, "wb") as f:
                f.write(audio_bytes)
            cmd = [
                ffmpeg, "-y", "-hide_banner", "-loglevel", "error",
                "-i", src,
                "-t", str(max_seconds),
                "-vn", "-ac", "2", "-ar", "44100",
                "-c:a", "libmp3lame", "-b:a", "128k",
                out,
            ]
            res = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=90)
            if res.returncode == 0 and os.path.exists(out) and os.path.getsize(out) > 2048:
                with open(out, "rb") as f:
                    return f.read()
            log.warning("vocal clip audio trim failed: rc=%s err=%s", res.returncode, res.stderr.decode("utf-8", "ignore")[-500:])
    except Exception as e:
        log.warning("vocal clip audio trim exception: %s", e)
    return audio_bytes

async def _trim_audio_for_vocal_clip(audio_bytes: bytes, max_seconds: int = 65) -> bytes:
    return await asyncio.to_thread(_trim_audio_for_vocal_clip_sync, audio_bytes, max_seconds)

async def _start_vocal_clip(update: Update, context: ContextTypes.DEFAULT_TYPE, img_bytes: bytes, user_prompt: str):
    """Premium one-person vocal/lip-sync clip: Suno audio -> Kling Avatar."""
    prompt = (user_prompt or "").strip()
    if not prompt:
        await update.effective_message.reply_text("РһРҝРёСҲРёСӮРө РҝРөСҒРҪСҺ/РәР»РёРҝ: СҒСӮРёР»СҢ, СҸР·СӢРә, РҪР°СҒСӮСҖРҫРөРҪРёРө, РҝСҖРёРҝРөРІ, РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ.")
        return
    try:
        faces = _detect_faces_for_choice(img_bytes) if FACESWAP_FACE_DETECTION_ENABLED else []
    except Exception:
        faces = []
    if len(faces) > 1:
        await update.effective_message.reply_text(
            "вқҢ Р”Р»СҸ СҖРөР¶РёРјР° В«РҡР»РёРҝ СҒ РІРҫРәР°Р»РҫРј / lip-syncВ» РҪСғР¶РөРҪ РҝРҫСҖСӮСҖРөСӮ РҫРҙРҪРҫРіРҫ СҮРөР»РҫРІРөРәР°. РқР° С„РҫСӮРҫ РҪР°Р№РҙРөРҪРҫ РҪРөСҒРәРҫР»СҢРәРҫ Р»РёСҶ.\n"
            "РЎРҙРөР»Р°Р№СӮРө РҫСӮРҙРөР»СҢРҪСӢР№ РәР°РҙСҖ РҫРҙРҪРҫРіРҫ РіРөСҖРҫСҸ РәСҖСғРҝРҪРөРө Рё Р·Р°РҝСғСҒСӮРёСӮРө СҖРөР¶РёРј СҒРҪРҫРІР°."
        )
        return
    user_id = update.effective_user.id
    img_digest = hashlib.sha1((img_bytes or b"")[:256000]).hexdigest()[:16]
    job_key = f"vocal:{user_id}:{img_digest}:{hashlib.sha1(prompt.encode('utf-8')).hexdigest()[:16]}"
    if job_key in _vocal_clip_background_jobs:
        await update.effective_message.reply_text("вҸі РўР°РәРҫР№ РІРҫРәР°Р»СҢРҪСӢР№ РәР»РёРҝ СғР¶Рө РҫРұСҖР°РұР°СӮСӢРІР°РөСӮСҒСҸ. Р”РҫР¶РҙРёСӮРөСҒСҢ СҖРөР·СғР»СҢСӮР°СӮР°, СҮСӮРҫРұСӢ РҪРө РҝРҫР»СғСҮРёСӮСҢ РҙСғРұР»Рё.")
        return

    async def _job():
        _vocal_clip_background_jobs.add(job_key)
        try:
            if not (SUNO_ENABLED and SUNO_API_KEY):
                raise RuntimeError("Р”Р»СҸ РІРҫРәР°Р»СҢРҪРҫРіРҫ РәР»РёРҝР° РҪСғР¶РөРҪ SUNO_ENABLED=1 Рё SUNO_API_KEY/COMET_API_KEY.")
            await update.effective_message.reply_text(
                "рҹҺӨ РҡР»РёРҝ СҒ РІРҫРәР°Р»РҫРј РҝСҖРёРҪСҸСӮ. Р РөР¶РёРј СҖР°СҒСҒСҮРёСӮР°РҪ РҪР° *РҫРҙРҪРҫРіРҫ СҮРөР»РҫРІРөРәР°*: СҒРҪР°СҮР°Р»Р° РіРөРҪРөСҖРёСҖСғСҺ РІРҫРәР°Р» Suno, Р·Р°СӮРөРј РҙРөР»Р°СҺ lip-sync СҮРөСҖРөР· Kling Avatar.",
                parse_mode="Markdown",
            )
            await context.bot.send_chat_action(update.effective_chat.id, ChatAction.RECORD_VIDEO)
            audio_bytes = await _run_suno_music_result_bytes(update, prompt)
            if not audio_bytes:
                raise RuntimeError("Suno РҪРө РІРөСҖРҪСғР» РІРҫРәР°Р»/РјСғР·СӢРәСғ. РҹРҫРҝСҖРҫРұСғР№СӮРө РәРҫСҖРҫСҮРө РҫРҝРёСҒР°СӮСҢ РҝРөСҒРҪСҺ РёР»Рё СғРІРөР»РёСҮРёСӮСҢ SUNO_TIMEOUT_S.")
            original_audio_size = len(audio_bytes)
            safe_audio_bytes = await _trim_audio_for_vocal_clip(audio_bytes, VOCAL_CLIP_MAX_AUDIO_S)
            trim_note = (
                f"рҹҺ§ Р’РҫРәР°Р»/СӮСҖРөРә РҝРҫР»СғСҮРөРҪ. Р”Р»СҸ СҒСӮР°РұРёР»СҢРҪРҫРіРҫ lip-sync РұРөСҖСғ РҝРөСҖРІСӢРө ~{VOCAL_CLIP_MAX_AUDIO_S} СҒРөРә. Р—Р°РҝСғСҒРәР°СҺ РіРөСҖРҫСҸвҖҰ"
                if len(safe_audio_bytes or b"") != original_audio_size
                else "рҹҺ§ Р’РҫРәР°Р»/СӮСҖРөРә РҝРҫР»СғСҮРөРҪ. Р—Р°РҝСғСҒРәР°СҺ lip-sync РіРөСҖРҫСҸвҖҰ"
            )
            audio_url = await _upload_bytes_to_telegram_file_url(
                update, context, safe_audio_bytes, "vocal_clip_lipsync_safe.mp3", trim_note
            )
            if not audio_url:
                raise RuntimeError("РқРө СғРҙР°Р»РҫСҒСҢ РҝРҫРҙРіРҫСӮРҫРІРёСӮСҢ РҝСғРұР»РёСҮРҪСӢР№ URL Р°СғРҙРёРҫ РҙР»СҸ lip-sync.")
            avatar_prompt = (
                "A single person performs as a music video singer. Accurate lip sync to the provided vocal audio, "
                "natural mouth shapes, expressive singing to camera, subtle rhythmic head and shoulder motion, "
                "premium music-video lighting, stable identity, no extra people, no subtitles, no text overlays. "
                f"Music direction: {prompt[:500]}"
            )
            ok = await _run_kling_avatar(
                update, context, img_bytes,
                script_text="singing performance",
                audio_file_url=audio_url,
                audio_filename="vocal_clip_lipsync_safe.mp3",
                audio_mime="audio/mpeg",
                avatar_prompt_override=avatar_prompt,
                max_wait_s=VOCAL_CLIP_KLING_MAX_WAIT_S,
            )
            return bool(ok)
        except Exception as e:
            log.exception("vocal clip failed: %s", e)
            with contextlib.suppress(Exception):
                await update.effective_message.reply_text(f"вқҢ РҡР»РёРҝ СҒ РІРҫРәР°Р»РҫРј РҪРө РҝРҫР»СғСҮРёР»СҒСҸ. РҹСҖРёСҮРёРҪР°: {str(e)[:900]}")
            return False
        finally:
            _vocal_clip_background_jobs.discard(job_key)

    async def _paid_start():
        # Billing must wait for the real provider result; do not detach paid jobs.
        return await _job()

    await _try_pay_then_do(
        update, context, update.effective_user.id,
        "runway", VOCAL_CLIP_UNIT_COST_USD, _paid_start,
        remember_kind="vocal_lipsync_clip",
        remember_payload={"prompt": prompt[:500]},
    )


async def _start_text_video(update: Update, context: ContextTypes.DEFAULT_TYPE, prompt: str):
    prompt = (prompt or "").strip()
    if not prompt:
        await update.effective_message.reply_text("РһРҝРёСҲРёСӮРө СҒСҶРөРҪСғ РҙР»СҸ РІРёРҙРөРҫ.")
        return False
    duration, aspect = parse_video_opts(prompt)
    engine = (context.user_data.get("text_video_engine") or TEXT_VIDEO_DEFAULT_ENGINE or "kling").strip().lower()
    if engine == "runway" and not TEXT_VIDEO_ALLOW_RUNWAY:
        engine = "kling"
    if engine not in ("sora", "runway", "kling"):
        engine = "kling"
    if engine == "sora" and _prompt_likely_has_people(prompt):
        await update.effective_message.reply_text(
            "вҡ пёҸ Sora 2 РІ РұРҫСӮРө РёСҒРҝРҫР»СҢР·СғРөСӮСҒСҸ СӮРҫР»СҢРәРҫ РҙР»СҸ СҒСҶРөРҪ РұРөР· Р»СҺРҙРөР№. Р’ Р·Р°РҝСҖРҫСҒРө РҫРұРҪР°СҖСғР¶РөРҪ СҮРөР»РҫРІРөРә/РҝРөСҖСҒРҫРҪР°Р¶. Р’СӢРұРөСҖРёСӮРө Kling РёР»Рё Runway.",
            reply_markup=_textvideo_action_kb("act"),
        )
        return False
    provider_cost = _video_provider_cost_usd(engine, duration)
    label = {"sora": "Sora 2 В· РұРөР· Р»СҺРҙРөР№", "runway": "Runway", "kling": "Kling"}[engine]

    async def _go():
        await update.effective_message.reply_text(f"рҹҺ¬ Р—Р°РҝСғСҒРәР°СҺ {label}: {duration} СҒРөРә В· {aspect}.")
        if engine == "runway":
            return await _run_runway_video(update, context, prompt, duration, aspect)
        return await _run_comet_text_video(update, context, engine, prompt, duration, aspect)

    await _try_pay_then_do(
        update, context, update.effective_user.id,
        "runway", provider_cost, _go,
        remember_kind=f"text_video_{engine}",
        remember_payload={"prompt": prompt[:500], "duration": duration, "aspect": aspect, "engine": engine},
    )
    return True

async def _run_kling_photo_clip(update: Update, context: ContextTypes.DEFAULT_TYPE, img_bytes: bytes, prompt: str, duration_s: int, aspect: str):
    """Legacy wrapper kept for external callers; v61 clean pipeline uses _start_photo_music_clip."""
    await context.bot.send_chat_action(update.effective_chat.id, ChatAction.RECORD_VIDEO)
    try:
        video = await _run_kling_photo_clip_result(img_bytes, prompt, duration_s, aspect)
        if not video:
            raise RuntimeError("empty video result")
        await _reply_video_bytes(update, video, "Kling photoвҶ’music clip вң…")
        return True
    except Exception as e:
        await update.effective_message.reply_text(f"вқҢ Kling photoвҶ’clip: {str(e)[:700]}")
        return False


async def _start_photo_music_clip(update: Update, context: ContextTypes.DEFAULT_TYPE, img_bytes: bytes, user_prompt: str):
    _, aspect = parse_video_opts(user_prompt or "")
    if not any(a in (user_prompt or "") for a in _ASPECTS):
        aspect = "9:16"
    elif aspect not in _ASPECTS:
        aspect = "9:16"
    target_duration = _photo_clip_target_duration(user_prompt or "")
    base_duration = min(10, target_duration)
    user_id = update.effective_user.id

    # РҡР»СҺСҮ Р·Р°СүРёСӮСӢ РҫСӮ РҙСғРұР»РөР№: РҙР»РёРҪРҪСӢРө webhook-Р·Р°РҙР°СҮРё Telegram РјРҫР¶РөСӮ РҝРҫРІСӮРҫСҖРҪРҫ РҙРҫСҒСӮР°РІР»СҸСӮСҢ.
    img_digest = hashlib.sha1((img_bytes or b"")[:256000]).hexdigest()[:16]
    job_key = f"{user_id}:{img_digest}:{hashlib.sha1((user_prompt or '').encode('utf-8')).hexdigest()[:16]}:{target_duration}:{aspect}"
    if job_key in _photo_clip_background_jobs:
        await update.effective_message.reply_text("вҸі РўР°РәРҫР№ С„РҫСӮРҫвҶ’РІРёРҙРөРҫРәР»РёРҝ СғР¶Рө РҫРұСҖР°РұР°СӮСӢРІР°РөСӮСҒСҸ. Р”РҫР¶РҙРёСӮРөСҒСҢ СҖРөР·СғР»СҢСӮР°СӮР°, СҮСӮРҫРұСӢ РҪРө РҝРҫР»СғСҮРёСӮСҢ РҙСғРұР»Рё.")
        return

    async def _photo_clip_job():
        _photo_clip_background_jobs.add(job_key)
        try:
            if not (PHOTO_CLIP_MUX_AUDIO and SUNO_AUTO_FOR_PHOTO_CLIP and SUNO_ENABLED and SUNO_API_KEY):
                raise RuntimeError("Р”Р»СҸ РҝРҫР»РҪРҫРіРҫ РәР»РёРҝР° СҒ РјСғР·СӢРәРҫР№ РІРәР»СҺСҮРёСӮРө SUNO_ENABLED=1, SUNO_AUTO_FOR_PHOTO_CLIP=1 Рё Р·Р°РҙР°Р№СӮРө SUNO_API_KEY (РёР»Рё COMET_API_KEY).")

            await update.effective_message.reply_text(
                f"рҹҺ¬ РӨРҫСӮРҫвҶ’РІРёРҙРөРҫРәР»РёРҝ РҝСҖРёРҪСҸСӮ. Р”РөР»Р°СҺ РһР”РҳРқ РёСӮРҫРіРҫРІСӢР№ MP4 СҒ РјСғР·СӢРәРҫР№. "
                f"Р’РёРҙРөРҫ: Kling; РјСғР·СӢРәР°: Suno; РҙР»РёРҪР°: ~{target_duration} СҒРөРә. РһРұСӢСҮРҪРҫ РҫР¶РёРҙР°РҪРёРө 5вҖ“15 РјРёРҪСғСӮ."
            )
            await update.effective_message.reply_text("рҹҺө Р“РөРҪРөСҖРёСҖСғСҺ РјСғР·СӢРәСғ/РҝРөСҒРҪСҺ СҮРөСҖРөР· SunoвҖҰ")
            await update.effective_message.reply_text("рҹҺһпёҸ Р“РөРҪРөСҖРёСҖСғСҺ РІРёРҙРөРҫСҖСҸРҙ СҮРөСҖРөР· KlingвҖҰ")
            await context.bot.send_chat_action(update.effective_chat.id, ChatAction.RECORD_VIDEO)

            started_at = time.time()
            video_task = asyncio.create_task(_run_kling_photo_clip_result(img_bytes, user_prompt, base_duration, aspect))
            audio_task = asyncio.create_task(_run_suno_music_result_bytes(update, user_prompt))

            try:
                video_bytes = await asyncio.wait_for(video_task, timeout=max(300, PHOTO_CLIP_TOTAL_USER_WAIT_S))
            except asyncio.TimeoutError:
                audio_task.cancel()
                raise RuntimeError(f"Kling РҪРө РІРөСҖРҪСғР» РІРёРҙРөРҫ Р·Р° {PHOTO_CLIP_TOTAL_USER_WAIT_S} СҒРөРә. РҹРҫРҝСҖРҫРұСғР№СӮРө РұРҫР»РөРө РәРҫСҖРҫСӮРәРёР№ prompt РёР»Рё РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ 10вҖ“15 СҒРөРә.")
            if not video_bytes:
                audio_task.cancel()
                raise RuntimeError("Kling РҪРө РІРөСҖРҪСғР» РІРёРҙРөРҫ")

            await update.effective_message.reply_text("вң… Р’РёРҙРөРҫСҖСҸРҙ РіРҫСӮРҫРІ. РһР¶РёРҙР°СҺ РјСғР·СӢРәСғ Suno Рё Р·Р°СӮРөРј СҒРҫРұРёСҖР°СҺ С„РёРҪР°Р»СҢРҪСӢР№ MP4вҖҰ")

            remaining_wait = max(30, PHOTO_CLIP_TOTAL_USER_WAIT_S - int(time.time() - started_at))
            audio_wait = max(30, min(max(PHOTO_CLIP_AUDIO_AFTER_VIDEO_WAIT_S, remaining_wait), SUNO_TIMEOUT_S))
            audio_bytes = None
            if audio_task.done():
                with contextlib.suppress(Exception):
                    audio_bytes = audio_task.result()
            else:
                try:
                    audio_bytes = await asyncio.wait_for(audio_task, timeout=audio_wait)
                except asyncio.TimeoutError:
                    audio_task.cancel()
                    audio_bytes = None

            if not audio_bytes:
                raise RuntimeError(
                    "Suno РҪРө РІРөСҖРҪСғР» РјСғР·СӢРәСғ РІ РҫСӮРІРөРҙС‘РҪРҪРҫРө РІСҖРөРјСҸ. РҳСӮРҫРіРҫРІСӢР№ РәР»РёРҝ РұРөР· РјСғР·СӢРәРё РҪРө РҫСӮРҝСҖР°РІР»СҸСҺ. "
                    "РҹРҫРҝСҖРҫРұСғР№СӮРө РөСүС‘ СҖР°Р· РёР»Рё СғРІРөР»РёСҮСҢСӮРө SUNO_TIMEOUT_S / PHOTO_CLIP_TOTAL_USER_WAIT_S."
                )

            await update.effective_message.reply_text(f"рҹҺ§ РңСғР·СӢРәР° РҝРҫР»СғСҮРөРҪР°. РЎРәР»РөРёРІР°СҺ Рё СҒР¶РёРјР°СҺ MP4 Р»РҫРәР°Р»СҢРҪРҫ СҮРөСҖРөР· ffmpeg, Р»РёРјРёСӮ ~{FFMPEG_MUX_TIMEOUT_S} СҒРөРә, СҖР°Р·РјРөСҖ РҙРҫ ~{FFMPEG_MUX_MAX_MB}MBвҖҰ")

            final_bytes = None
            if PHOTO_CLIP_PIPELINE:
                final_bytes = await asyncio.wait_for(asyncio.to_thread(_mux_video_audio_sync, video_bytes, audio_bytes, target_duration), timeout=max(60, FFMPEG_MUX_TIMEOUT_S + 30))
            if not final_bytes and PHOTO_CLIP_SEND_BASE_IF_MUX_FAILS:
                final_bytes = video_bytes
            if not final_bytes:
                raise RuntimeError("ffmpeg РҪРө СғСҒРҝРөР» РёР»Рё РҪРө СҒРјРҫРі СҒРҫРұСҖР°СӮСҢ РёСӮРҫРіРҫРІСӢР№ MP4 СҒ РјСғР·СӢРәРҫР№. РҹСҖРҫРІРөСҖСҢСӮРө Render CPU/Р»РҫРіРё ffmpeg РёР»Рё СғРІРөР»РёСҮСҢСӮРө FFMPEG_MUX_TIMEOUT_S.")

            caption = f"РӨРҫСӮРҫвҶ’РІРёРҙРөРҫРәР»РёРҝ СҒ РјСғР·СӢРәРҫР№ вң… РһРҙРёРҪ MP4 В· ~{target_duration} СҒРөРә"
            await _reply_video_bytes(update, final_bytes, caption)
            return True
        except Exception as e:
            log.exception("photo music clip background pipeline failed: %s", e)
            with contextlib.suppress(Exception):
                await update.effective_message.reply_text(f"вқҢ РӨРҫСӮРҫвҶ’РІРёРҙРөРҫРәР»РёРҝ РҪРө РҝРҫР»СғСҮРёР»СҒСҸ. РҹСҖРёСҮРёРҪР°: {str(e)[:900]}")
            return False
        finally:
            _photo_clip_background_jobs.discard(job_key)

    async def _go():
        # РҹР»Р°СӮРҪР°СҸ Р·Р°РҙР°СҮР° РҫСҒСӮР°С‘СӮСҒСҸ РҝСҖРёРІСҸР·Р°РҪРҪРҫР№ РҙРҫ СғСҒРҝРөСҲРҪРҫР№ РҫСӮРҝСҖР°РІРәРё РёСӮРҫРіРҫРІРҫРіРҫ MP4.
        return await _photo_clip_job()

    await _try_pay_then_do(
        update, context, update.effective_user.id,
        "runway", PHOTO_CLIP_UNIT_COST_USD, _go,
        remember_kind="kling_photo_music_clip_v62_background",
        remember_payload={"prompt": user_prompt, "duration": target_duration, "aspect": aspect},
    )

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ AI selfie / Nano Banana style image fusion в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def _prepare_reference_image_for_gemini(img_bytes: bytes, max_side: int | None = None) -> tuple[str, str]:
    """Return (base64, mime). JPEG is used to keep payload smaller."""
    max_side = int(max_side or AI_SELFIE_MAX_SIDE or 1536)
    if Image is None:
        return base64.b64encode(img_bytes).decode("ascii"), sniff_image_mime(img_bytes)
    try:
        im = Image.open(BytesIO(img_bytes)).convert("RGB")
        if max(im.size) > max_side:
            im.thumbnail((max_side, max_side), Image.LANCZOS)
        bio = BytesIO()
        im.save(bio, format="JPEG", quality=92, optimize=True)
        return base64.b64encode(bio.getvalue()).decode("ascii"), "image/jpeg"
    except Exception:
        return base64.b64encode(img_bytes).decode("ascii"), sniff_image_mime(img_bytes)


def _ai_selfie_final_prompt(user_prompt: str, preset_prompt: str = "") -> str:
    user_prompt = (user_prompt or "").strip()
    preset_prompt = (preset_prompt or "").strip()
    if not user_prompt and not preset_prompt:
        user_prompt = "realistic iPhone selfie with a celebrity, natural light, social media style, 4:5"
    if not AI_SELFIE_ALLOW_PUBLIC_FIGURES:
        safety_line = "Use a generic fictional celebrity-like person instead of a real public figure."
    else:
        safety_line = "This is a fictional AI-generated fan/selfie scene, not documentary evidence, endorsement, news, or political support."
    parts = [
        "Edit the uploaded user selfie into a new realistic AI photo.",
        "Preserve the user's face identity and natural likeness.",
        "Create one coherent iPhone/selfie-style scene with natural lighting, correct perspective and realistic skin.",
        "No text, no captions, no watermarks, no logos, no UI overlays, no distorted hands, no duplicate faces.",
        safety_line,
    ]
    if preset_prompt:
        parts.append("Preset: " + preset_prompt[:350])
    if user_prompt:
        parts.append("User request: " + user_prompt[:650])
    parts.append(f"Aspect ratio {AI_SELFIE_DEFAULT_ASPECT}. Output image size {AI_SELFIE_IMAGE_SIZE}.")
    return " ".join(parts).strip()


def _extract_image_b64_from_gemini(obj) -> str:
    """Robust extraction from Interactions API responses: output_image or nested steps/content image blocks."""
    if isinstance(obj, dict):
        oi = obj.get("output_image")
        if isinstance(oi, dict):
            data = oi.get("data") or oi.get("b64_json") or oi.get("base64")
            if isinstance(data, str) and len(data) > 100:
                return data
        # Common REST structures and SDK-like serialization.
        if obj.get("type") == "image":
            data = obj.get("data") or obj.get("b64_json") or obj.get("base64")
            if isinstance(data, str) and len(data) > 100:
                return data
        for key in ("image", "generated_image", "output", "content", "steps", "candidates", "data", "result", "response"):
            val = obj.get(key)
            found = _extract_image_b64_from_gemini(val)
            if found:
                return found
        for val in obj.values():
            found = _extract_image_b64_from_gemini(val)
            if found:
                return found
    elif isinstance(obj, (list, tuple)):
        for item in obj:
            found = _extract_image_b64_from_gemini(item)
            if found:
                return found
    return ""


async def _run_comet_ai_selfie_bytes(img_bytes: bytes, user_prompt: str, preset_prompt: str = "") -> bytes | None:
    """AI selfie via Comet Nano Banana/Gemini generateContent.
    v59: one primary Comet route first, longer timeout, smaller reference image, then optional compact retry.
    """
    if not COMET_API_KEY:
        return None

    prompt = _ai_selfie_final_prompt(user_prompt, preset_prompt)
    img_b64, mime = _prepare_reference_image_for_gemini(img_bytes, AI_SELFIE_MAX_SIDE)
    mime = mime or "image/jpeg"

    models: list[str] = []
    for m in COMET_IMAGE_EDIT_FALLBACK_MODELS:
        m = (m or "").strip()
        if m and m not in models:
            models.append(m)
    if AI_SELFIE_FAST_MODE and models:
        # In production do not try many slow models before returning; first model must be the enabled Comet channel.
        models = models[:2]

    base_paths: list[str] = []
    for p in (COMET_IMAGE_EDIT_PATH, "/v1beta/models/{model}:generateContent"):
        p = (p or "").strip()
        if p and p not in base_paths:
            base_paths.append(p)

    def _mk_generate_content_payload(style: str, compact: bool = False) -> dict:
        req_prompt = prompt
        if compact:
            req_prompt = (
                "Use the uploaded selfie as identity reference. Create a realistic AI selfie scene. "
                f"User request: {(user_prompt or preset_prompt or 'celebrity selfie')[:450]}. "
                "Preserve face identity. No text, no logos, no watermark."
            )
        if style == "camel":
            image_part = {"inlineData": {"mimeType": mime, "data": img_b64}}
            gen_cfg = {"responseModalities": ["TEXT", "IMAGE"]}
        else:
            image_part = {"inline_data": {"mime_type": mime, "data": img_b64}}
            gen_cfg = {"response_modalities": ["TEXT", "IMAGE"]}
        return {
            "contents": [{
                "role": "user",
                "parts": [
                    {"text": req_prompt},
                    image_part,
                ],
            }],
            "generationConfig": gen_cfg,
        }

    # Main Comet route. Bearer is the normal gateway auth; avoid slow duplicate auth variants unless explicitly disabled fast mode.
    header_variants = [
        {"Authorization": f"Bearer {COMET_API_KEY}", "Accept": "application/json", "Content-Type": "application/json"},
    ]
    if not AI_SELFIE_FAST_MODE:
        header_variants.extend([
            {"x-goog-api-key": COMET_API_KEY, "Accept": "application/json", "Content-Type": "application/json"},
            {"Authorization": COMET_API_KEY, "Accept": "application/json", "Content-Type": "application/json"},
        ])

    last_err = ""
    timeout_seen = False
    async with httpx.AsyncClient(timeout=httpx.Timeout(COMET_IMAGE_EDIT_TIMEOUT_S, connect=40.0, read=COMET_IMAGE_EDIT_TIMEOUT_S, write=120.0), follow_redirects=True) as client:
        for model in models:
            for path_tmpl in base_paths:
                path = path_tmpl.replace("{model}", model)
                if "/images/" in path and "generateContent" not in path:
                    continue
                url = f"{COMET_BASE_URL}{path}"
                payload_styles = ("camel", "snake") if AI_SELFIE_FAST_MODE else ("camel", "snake")
                for payload_style in payload_styles:
                    payload = _mk_generate_content_payload(payload_style, compact=False)
                    for headers in header_variants:
                        try:
                            r = await client.post(url, headers=headers, json=payload)
                            if r.status_code >= 400:
                                last_err = f"{r.status_code}: {_api_error_preview(r)}" if "_api_error_preview" in globals() else f"{r.status_code}: {r.text[:500]}"
                                log.warning("Comet generateContent AI selfie failed model=%s path=%s style=%s: %s", model, path, payload_style, last_err)
                                continue
                            out = await _image_bytes_from_response(r, client)
                            if out:
                                return out
                            try:
                                js = r.json() or {}
                            except Exception:
                                js = {}
                            last_err = f"Comet generateContent: РҪРөСӮ РёР·РҫРұСҖР°Р¶РөРҪРёСҸ РІ РҫСӮРІРөСӮРө {json.dumps(js, ensure_ascii=False)[:800]}"
                        except (httpx.ReadTimeout, httpx.ConnectTimeout, httpx.TimeoutException) as e:
                            timeout_seen = True
                            last_err = f"Comet image generation timeout after {COMET_IMAGE_EDIT_TIMEOUT_S}s: {type(e).__name__}"
                            log.warning("Comet generateContent AI selfie timeout model=%s path=%s: %s", model, path, e)
                            # Do not spend all time on additional variants when the model is just slow.
                            break
                        except Exception as e:
                            last_err = str(e)
                            log.warning("Comet generateContent AI selfie exception model=%s path=%s: %s", model, path, e)
                            continue
                    if timeout_seen and AI_SELFIE_FAST_MODE:
                        break
                if timeout_seen and AI_SELFIE_FAST_MODE:
                    break
            if timeout_seen and AI_SELFIE_FAST_MODE:
                break

        if timeout_seen and AI_SELFIE_RETRY_ON_TIMEOUT:
            # One compact retry with smaller image. This often helps when the gateway times out on upload+generation+download.
            retry_b64, retry_mime = _prepare_reference_image_for_gemini(img_bytes, 768)
            img_b64, mime = retry_b64, retry_mime or "image/jpeg"
            model = models[0] if models else COMET_IMAGE_EDIT_MODEL
            path = (COMET_IMAGE_EDIT_PATH or "/v1beta/models/{model}:generateContent").replace("{model}", model)
            url = f"{COMET_BASE_URL}{path}"
            payload = _mk_generate_content_payload("camel", compact=True)
            try:
                r = await client.post(url, headers={"Authorization": f"Bearer {COMET_API_KEY}", "Accept": "application/json", "Content-Type": "application/json"}, json=payload)
                if r.status_code < 400:
                    out = await _image_bytes_from_response(r, client)
                    if out:
                        return out
                    try:
                        js = r.json() or {}
                    except Exception:
                        js = {}
                    last_err = f"retry: РҪРөСӮ РёР·РҫРұСҖР°Р¶РөРҪРёСҸ РІ РҫСӮРІРөСӮРө {json.dumps(js, ensure_ascii=False)[:800]}"
                else:
                    last_err = f"retry {r.status_code}: {_api_error_preview(r)}" if "_api_error_preview" in globals() else f"retry {r.status_code}: {r.text[:500]}"
            except Exception as e:
                last_err = f"retry exception: {type(e).__name__}: {e}"

        if COMET_IMAGE_EDIT_OPENAI_FALLBACK:
            headers = {"Authorization": f"Bearer {COMET_API_KEY}", "Accept": "application/json"}
            openai_paths = ["/v1/images/edits", "/v1/images/generations"]
            edit_bytes, filename, edit_mime = _prepare_image_for_edit(img_bytes)
            attempts = [
                {"prompt": prompt, "n": "1", "response_format": "b64_json", "size": "1024x1024"},
                {"prompt": prompt, "n": "1", "response_format": "b64_json"},
            ]
            for path in openai_paths:
                for model in models:
                    for data in attempts:
                        try:
                            payload = {"model": model, **data}
                            files = {"image": ("image.png", edit_bytes, edit_mime or "image/png")}
                            r = await client.post(f"{COMET_BASE_URL}{path}", headers=headers, data=payload, files=files)
                            if r.status_code >= 400:
                                last_err = f"{r.status_code}: {_api_error_preview(r)}" if "_api_error_preview" in globals() else f"{r.status_code}: {r.text[:500]}"
                                continue
                            out = await _image_bytes_from_response(r, client)
                            if out:
                                return out
                        except Exception as e:
                            last_err = str(e)
                            continue

    raise RuntimeError(last_err or "Comet AI selfie generateContent failed")


async def _run_openai_ai_selfie_fallback(img_bytes: bytes, user_prompt: str, preset_prompt: str = "") -> bytes | None:
    if not OPENAI_IMAGE_KEY or OPENAI_IMAGE_KEY.startswith("sk-or-"):
        return None
    return await _openai_image_edit_bytes(img_bytes, _ai_selfie_final_prompt(user_prompt, preset_prompt))


async def _run_ai_selfie_image(update: Update, context: ContextTypes.DEFAULT_TYPE, img_bytes: bytes, user_prompt: str, preset_prompt: str = "") -> bool:
    await context.bot.send_chat_action(update.effective_chat.id, ChatAction.UPLOAD_PHOTO)
    if not img_bytes:
        await update.effective_message.reply_text("вқҢ РқРөСӮ РёР·РҫРұСҖР°Р¶РөРҪРёСҸ РҙР»СҸ AI-СҒРөР»С„Рё.")
        return False
    try:
        with contextlib.suppress(Exception):
            await update.effective_message.reply_text("вҸі AI-СҒРөР»С„Рё РҝСҖРёРҪСҸСӮРҫ. Comet/Nano Banana РјРҫР¶РөСӮ РҫРұСҖР°РұР°СӮСӢРІР°СӮСҢ 1вҖ“5 РјРёРҪСғСӮ, РҫР¶РёРҙР°СҺ СҖРөР·СғР»СҢСӮР°СӮвҖҰ")
        out = None
        if AI_SELFIE_PROVIDER in ("comet", "auto", "nano", "nanobanana", "nano_banana", "gemini") and COMET_API_KEY:
            out = await _run_comet_ai_selfie_bytes(img_bytes, user_prompt, preset_prompt)
        if not out and AI_SELFIE_PROVIDER in ("openai", "auto", "gpt", "gptimages", "gpt_images") and OPENAI_IMAGE_KEY and not OPENAI_IMAGE_KEY.startswith("sk-or-"):
            out = await _run_openai_ai_selfie_fallback(img_bytes, user_prompt, preset_prompt)
        if not out:
            await update.effective_message.reply_text(
                "вқҢ AI-СҒРөР»С„Рё РҪРөРҙРҫСҒСӮСғРҝРҪРҫ: РҝСҖРҫРІРөСҖСҢСӮРө COMET_API_KEY, COMET_IMAGE_EDIT_MODEL Рё COMET_IMAGE_EDIT_PATH. "
                "РўРөРәСғСүРёР№ РјР°СҖСҲСҖСғСӮ: Comet image-edit."
            )
            return False
        bio = BytesIO(out)
        bio.name = "ai_selfie.png"
        caption = "рҹӨі AI-СҒРөР»С„Рё РіРҫСӮРҫРІРҫ вң…\nРҹРҫРјРөСӮРәР°: РёР·РҫРұСҖР°Р¶РөРҪРёРө СҒРіРөРҪРөСҖРёСҖРҫРІР°РҪРҫ РҳРҳ; РҪРө РёСҒРҝРҫР»СҢР·СғР№СӮРө РәР°Рә РҙРҫРәР°Р·Р°СӮРөР»СҢСҒСӮРІРҫ СҖРөР°Р»СҢРҪРҫР№ РІСҒСӮСҖРөСҮРё/РҝРҫРҙРҙРөСҖР¶РәРё."
        if AI_SELFIE_SEND_AS_DOCUMENT:
            await update.effective_message.reply_document(InputFile(bio), caption=caption)
        else:
            await update.effective_message.reply_photo(photo=out, caption=caption)
        return True
    except Exception as e:
        log.exception("AI selfie error: %s", e)
        err_txt = str(e)[:1200]
        if "timeout" in err_txt.lower() or "timed out" in err_txt.lower():
            err_txt = (
                "Comet/Nano Banana РҪРө СғСҒРҝРөР» РІРөСҖРҪСғСӮСҢ РёР·РҫРұСҖР°Р¶РөРҪРёРө РІ РҫСӮРІРөРҙС‘РҪРҪРҫРө РІСҖРөРјСҸ. "
                "РҹРҫРҝСҖРҫРұСғР№СӮРө РөСүС‘ СҖР°Р· СҒ 1K/РәРҫСҖРҫСӮРәРёРј РҝСҖРҫРјРҝСӮРҫРј РёР»Рё СғРІРөР»РёСҮСҢСӮРө COMET_IMAGE_EDIT_TIMEOUT_S РҙРҫ 600. "
                f"РўРөС…РҪРёСҮРөСҒРәРё: {err_txt}"
            )
        await update.effective_message.reply_text(f"вқҢ AI-СҒРөР»С„Рё РҪРө РҝРҫР»СғСҮРёР»РҫСҒСҢ. РҹСҖРёСҮРёРҪР°: {err_txt}")
        return False


async def _start_ai_selfie(update: Update, context: ContextTypes.DEFAULT_TYPE, img_bytes: bytes, user_prompt: str, preset_prompt: str = ""):
    async def _go():
        return bool(await _run_ai_selfie_image(update, context, img_bytes, user_prompt, preset_prompt))

    await _try_pay_then_do(
        update, context, update.effective_user.id,
        "img", AI_SELFIE_UNIT_COST_USD, _go,
        remember_kind="ai_selfie_gemini_image_edit",
        remember_payload={"prompt": user_prompt, "preset": preset_prompt[:400]},
    )

def background_presets_kb():
    # РҹСҖРөСҒРөСӮСӢ РҝРҫ РҫРҙРҪРҫРјСғ РІ СҒСӮСҖРҫРәРө: РҪР° РјР°Р»РөРҪСҢРәРёС… СҚРәСҖР°РҪР°С… СӮРөРәСҒСӮ РҪРө РҫРұСҖРөР·Р°РөСӮСҒСҸ.
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("рҹҢ« Р Р°Р·РјСӢСӮСҢ РёСҒС…РҫРҙРҪСӢР№ С„РҫРҪ", callback_data="pedit:bg:blur")],
        [InlineKeyboardButton("рҹҸ– РҹР»СҸР¶ / РјРҫСҖРө", callback_data="pedit:bg:beach")],
        [InlineKeyboardButton("вӣ° Р“РҫСҖСӢ / Р°Р»СҢРҝСӢ", callback_data="pedit:bg:mountains")],
        [InlineKeyboardButton("рҹҢҝ РҹСҖРёСҖРҫРҙР° / РҝР°СҖРә", callback_data="pedit:bg:nature")],
        [InlineKeyboardButton("рҹҸҷ РҡСҖСӢСҲР° / РіРҫСҖРҫРҙ", callback_data="pedit:bg:roof")],
        [InlineKeyboardButton("рҹҸў РһС„РёСҒ / РұРёР·РҪРөСҒ", callback_data="pedit:bg:office")],
        [InlineKeyboardButton("вҡӘ Р‘РөР»СӢР№ СҒСӮСғРҙРёР№РҪСӢР№ С„РҫРҪ", callback_data="pedit:bg:white")],
        [InlineKeyboardButton("вңҚпёҸ РЎРІРҫР№ С„РҫРҪ", callback_data="pedit:bg:custom")],
        [InlineKeyboardButton("в¬…пёҸ РқР°Р·Р°Рҙ", callback_data="pedit:back")],
    ])


def _set_waiting_removebg(context):
    context.user_data["awaiting_photo_for"] = "removebg"
    context.user_data["photo_flow"] = "removebg"


def _is_waiting_removebg(context) -> bool:
    return bool(context and (context.user_data.get("awaiting_photo_for") == "removebg" or context.user_data.get("photo_flow") == "removebg"))


def _clear_removebg_wait(context):
    if not context:
        return
    if context.user_data.get("awaiting_photo_for") == "removebg":
        context.user_data.pop("awaiting_photo_for", None)
    if context.user_data.get("photo_flow") == "removebg":
        context.user_data.pop("photo_flow", None)


def _set_waiting_replacebg(context, prompt: str = ""):
    context.user_data["awaiting_photo_for"] = "replacebg"
    context.user_data["photo_flow"] = "replacebg"
    if prompt:
        context.user_data["replacebg_prompt"] = prompt.strip()


def _set_replacebg_wait_text(context):
    context.user_data["replacebg_wait_text"] = "1"


def _is_replacebg_wait_text(context) -> bool:
    return bool(context and context.user_data.get("replacebg_wait_text"))


def _is_waiting_replacebg(context) -> bool:
    return bool(context and (context.user_data.get("awaiting_photo_for") == "replacebg" or context.user_data.get("photo_flow") == "replacebg"))


def _clear_replacebg_wait(context):
    if not context:
        return
    for key in ("replacebg_wait_text", "replacebg_prompt"):
        context.user_data.pop(key, None)
    if context.user_data.get("awaiting_photo_for") == "replacebg":
        context.user_data.pop("awaiting_photo_for", None)
    if context.user_data.get("photo_flow") == "replacebg":
        context.user_data.pop("photo_flow", None)


def _is_remove_bg_request(text: str) -> bool:
    tl = (text or "").lower()
    return any(k in tl for k in ("СғРҙР°Р»Рё С„РҫРҪ", "СғРҙР°Р»РёСӮСҢ С„РҫРҪ", "СғРұРөСҖРё С„РҫРҪ", "СғРұСҖР°СӮСҢ С„РҫРҪ", "removebg", "remove background", "РҝСҖРҫР·СҖР°СҮРҪСӢР№ С„РҫРҪ", "РұРөР· С„РҫРҪР°"))


def _is_replace_bg_request(text: str) -> bool:
    tl = (text or "").lower()
    return any(k in tl for k in ("Р·Р°РјРөРҪРё С„РҫРҪ", "Р·Р°РјРөРҪРёСӮСҢ С„РҫРҪ", "РҝРҫРјРөРҪСҸР№ С„РҫРҪ", "РҝРҫРјРөРҪСҸСӮСҢ С„РҫРҪ", "replacebg", "replace background", "С„РҫРҪ РҪР°", "РҙСҖСғРіРҫР№ С„РҫРҪ"))


def _bg_kind_from_text(text: str) -> tuple[str, str]:
    tl = (text or "").lower()
    if any(k in tl for k in ("РҝР»СҸР¶", "РјРҫСҖРө", "РҫРәРөР°РҪ", "beach", "coast", "shore")):
        return "beach", text
    if any(k in tl for k in ("РіРҫСҖСӢ", "РіРҫСҖР°", "mountain", "alps", "Р°Р»СҢРҝ")):
        return "mountains", text
    if any(k in tl for k in ("РәСҖСӢСҲР°", "РіРҫСҖРҫРҙ", "РҪРөРұРҫСҒРәСҖРөРұ", "rooftop", "city", "skyline", "СӮРөСҖСҖР°СҒР°")):
        return "roof", text
    if any(k in tl for k in ("РҫС„РёСҒ", "РәР°РұРёРҪРөСӮ", "business", "office", "coworking", "РҝРөСҖРөРіРҫРІРҫСҖ")):
        return "office", text
    if any(k in tl for k in ("РҝСҖРёСҖРҫРҙР°", "Р»РөСҒ", "Р·РөР»РөРҪСҢ", "nature", "forest", "park", "РҝР°СҖРә")):
        return "nature", text
    if any(k in tl for k in ("РұРөР»СӢР№", "white", "studio", "СҒСӮСғРҙ")):
        return "white", text
    if any(k in tl for k in ("СҮРөСҖРҪСӢР№", "СҮС‘СҖРҪСӢР№", "black")):
        return "black", text
    if any(k in tl for k in ("СҖР°Р·РјСӢ", "blur", "РұР»СҺСҖ")):
        return "blur", text
    return "custom", text


def _safe_b64decode_image(value: str) -> bytes | None:
    if not value or not isinstance(value, str):
        return None
    s = value.strip()
    if s.startswith("data:") and "," in s:
        s = s.split(",", 1)[1]
    try:
        return base64.b64decode(s, validate=False)
    except Exception:
        return None


def _find_first_image_b64(obj) -> bytes | None:
    if isinstance(obj, dict):
        # Gemini / Nano Banana generateContent returns candidates[].content.parts[].inlineData/inline_data.data.
        for key in ("inline_data", "inlineData"):
            v = obj.get(key)
            if isinstance(v, dict):
                mime = str(v.get("mime_type") or v.get("mimeType") or "").lower()
                if mime.startswith("image/") or v.get("data"):
                    out = _safe_b64decode_image(v.get("data") or v.get("bytes_base64") or v.get("bytesBase64"))
                    if out:
                        return out
        for key in ("b64_json", "image_b64", "image", "png", "result_b64", "bytes_base64", "bytesBase64"):
            if key in obj:
                out = _safe_b64decode_image(obj.get(key))
                if out:
                    return out
        for v in obj.values():
            out = _find_first_image_b64(v)
            if out:
                return out
    elif isinstance(obj, list):
        for v in obj:
            out = _find_first_image_b64(v)
            if out:
                return out
    return None


def _find_first_image_url(obj) -> str:
    if isinstance(obj, dict):
        for key in ("url", "image_url", "output_url", "result_url", "download_url", "file_url"):
            val = obj.get(key)
            if isinstance(val, str) and val.startswith(("http://", "https://")):
                return val
            if isinstance(val, dict):
                nested = val.get("url")
                if isinstance(nested, str) and nested.startswith(("http://", "https://")):
                    return nested
        for v in obj.values():
            url = _find_first_image_url(v)
            if url:
                return url
    elif isinstance(obj, list):
        for v in obj:
            url = _find_first_image_url(v)
            if url:
                return url
    return ""


async def _image_bytes_from_response(resp: httpx.Response, client: httpx.AsyncClient) -> bytes | None:
    ctype = (resp.headers.get("content-type") or "").lower()
    if ctype.startswith("image/") and resp.content:
        return bytes(resp.content)
    try:
        obj = resp.json()
    except Exception:
        return None
    out = _find_first_image_b64(obj)
    if out:
        return out
    url = _find_first_image_url(obj)
    if url:
        try:
            rr = await client.get(url, timeout=BG_REMOVE_TIMEOUT_S)
            rr.raise_for_status()
            if rr.content:
                return bytes(rr.content)
        except Exception as e:
            log.warning("background result url download failed: %s", e)
    return None


def _prepare_bytes_for_rembg(img_bytes: bytes) -> bytes:
    """
    Render Starter РјРҫР¶РөСӮ СғРҝРёСҖР°СӮСҢСҒСҸ РІ RAM РҪР° РұРҫР»СҢСҲРёС… С„РҫСӮРҫ.
    Р”Р»СҸ local rembg СғРјРөРҪСҢСҲР°РөРј СҒР»РёСҲРәРҫРј РәСҖСғРҝРҪСӢРө РёР·РҫРұСҖР°Р¶РөРҪРёСҸ РҙРҫ REMBG_MAX_SIDE.
    """
    if Image is None or not REMBG_MAX_SIDE or REMBG_MAX_SIDE <= 0:
        return img_bytes
    try:
        im = Image.open(BytesIO(img_bytes))
        with contextlib.suppress(Exception):
            im = ImageOps.exif_transpose(im)
        w, h = im.size
        mx = max(w, h)
        if mx <= REMBG_MAX_SIDE:
            return img_bytes
        scale = REMBG_MAX_SIDE / float(mx)
        nw, nh = max(1, int(w * scale)), max(1, int(h * scale))
        resample = getattr(Image, "Resampling", Image).LANCZOS
        im = im.convert("RGB").resize((nw, nh), resample)
        out = BytesIO()
        im.save(out, format="JPEG", quality=94, optimize=True)
        _bg_note_error(f"local rembg input resized {w}x{h}->{nw}x{nh}")
        return out.getvalue()
    except Exception as e:
        _bg_note_error(f"local rembg resize skipped: {e}")
        return img_bytes


def _get_local_rembg_session():
    global _REMBG_SESSION
    if rembg_remove is None:
        _bg_note_error(f"rembg import failed: {REMBG_IMPORT_ERROR}")
        return None
    if rembg_new_session is None:
        _bg_note_error("rembg.new_session is not available")
        return None
    if _REMBG_SESSION is not None:
        return _REMBG_SESSION
    with _REMBG_SESSION_LOCK:
        if _REMBG_SESSION is not None:
            return _REMBG_SESSION
        with contextlib.suppress(Exception):
            os.environ.setdefault("U2NET_HOME", U2NET_HOME)
            os.makedirs(U2NET_HOME, exist_ok=True)
            os.makedirs(XDG_CACHE_HOME, exist_ok=True)
        last_exc = ""
        for model_name in REMBG_MODEL_FALLBACKS:
            try:
                log.info("Initializing local rembg session: model=%s U2NET_HOME=%s", model_name, os.environ.get("U2NET_HOME"))
                _REMBG_SESSION = rembg_new_session(model_name)
                log.info("local rembg session ready: model=%s", model_name)
                return _REMBG_SESSION
            except Exception as e:
                last_exc = repr(e)
                log.warning("local rembg session init failed model=%s: %s", model_name, e)
                _bg_note_error(f"local rembg session init failed model={model_name}: {e}")
        _bg_note_error(f"local rembg session unavailable: {last_exc}")
    return None



async def _local_rembg_remove_subprocess(img_bytes: bytes) -> bytes | None:
    """
    РҳР·РҫР»РёСҖРҫРІР°РҪРҪСӢР№ local rembg worker. Р•СҒР»Рё rembg/onnxruntime Р·Р°РІРёСҒРҪРөСӮ РҝСҖРё СҒРәР°СҮРёРІР°РҪРёРё
    РёР»Рё РёРҪРёСҶРёР°Р»РёР·Р°СҶРёРё РјРҫРҙРөР»Рё, РҫСҒРҪРҫРІРҪРҫР№ РұРҫСӮ РҪРө Р·Р°РІРёСҒР°РөСӮ: РҝСҖРҫСҶРөСҒСҒ СғРұРёРІР°РөРј РҝРҫ timeout.
    """
    if not LOCAL_REMBG_ENABLED or rembg_remove is None:
        return None
    if Image is not None:
        img_bytes = _prepare_bytes_for_rembg(img_bytes)
    timeout_s = max(20.0, float(LOCAL_REMBG_TIMEOUT_S or 180))
    env = os.environ.copy()
    env["U2NET_HOME"] = env.get("U2NET_HOME") or U2NET_HOME or "/tmp/.u2net"
    env["XDG_CACHE_HOME"] = env.get("XDG_CACHE_HOME") or XDG_CACHE_HOME or "/tmp/.cache"
    env.setdefault("OMP_NUM_THREADS", "1")
    env.setdefault("MPLCONFIGDIR", "/tmp/.matplotlib")
    for d in (env["U2NET_HOME"], env["XDG_CACHE_HOME"], env["MPLCONFIGDIR"]):
        with contextlib.suppress(Exception):
            os.makedirs(d, exist_ok=True)

    worker_code = r'''
import os, sys
from pathlib import Path
try:
    from rembg import remove, new_session
except Exception as e:
    print(f"IMPORT_ERROR: {e}", file=sys.stderr)
    sys.exit(11)
model = os.environ.get("REMBG_MODEL", "u2netp") or "u2netp"
inp, outp = sys.argv[1], sys.argv[2]
try:
    Path(os.environ.get("U2NET_HOME", "/tmp/.u2net")).mkdir(parents=True, exist_ok=True)
    Path(os.environ.get("XDG_CACHE_HOME", "/tmp/.cache")).mkdir(parents=True, exist_ok=True)
    data = Path(inp).read_bytes()
    sess = new_session(model)
    try:
        out = remove(data, session=sess, force_return_bytes=True)
    except TypeError:
        out = remove(data, session=sess)
    Path(outp).write_bytes(out)
except Exception as e:
    print(f"REMBG_WORKER_ERROR: {type(e).__name__}: {e}", file=sys.stderr)
    sys.exit(12)
'''
    in_path = out_path = None
    try:
        with tempfile.NamedTemporaryFile(prefix="rembg_in_", suffix=".jpg", delete=False) as f:
            f.write(img_bytes)
            in_path = f.name
        fd, out_path = tempfile.mkstemp(prefix="rembg_out_", suffix=".png")
        os.close(fd)
        proc = await asyncio.create_subprocess_exec(
            sys.executable, "-c", worker_code, in_path, out_path,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
            env=env,
        )
        try:
            stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=timeout_s)
        except asyncio.TimeoutError:
            with contextlib.suppress(Exception):
                proc.kill()
            with contextlib.suppress(Exception):
                await proc.wait()
            _bg_note_error(f"local rembg subprocess timeout after {timeout_s:.0f}s")
            return None
        if proc.returncode != 0:
            msg = (stderr or stdout or b"").decode("utf-8", "replace")[:800]
            _bg_note_error(f"local rembg subprocess failed rc={proc.returncode}: {msg}")
            return None
        if not out_path or not os.path.exists(out_path) or os.path.getsize(out_path) < 512:
            _bg_note_error("local rembg subprocess produced empty output")
            return None
        with open(out_path, "rb") as f:
            return f.read()
    except Exception as e:
        _bg_note_error(f"local rembg subprocess exception: {e}")
        log.warning("local rembg subprocess exception: %s", e)
        return None
    finally:
        for path in (in_path, out_path):
            if path:
                with contextlib.suppress(Exception):
                    os.remove(path)

def _local_rembg_remove_sync(img_bytes: bytes) -> bytes:
    img_bytes = _prepare_bytes_for_rembg(img_bytes)
    session = _get_local_rembg_session()
    if session is not None:
        try:
            return rembg_remove(img_bytes, session=session, force_return_bytes=True)
        except TypeError:
            return rembg_remove(img_bytes, session=session)
    try:
        return rembg_remove(img_bytes, force_return_bytes=True)
    except TypeError:
        return rembg_remove(img_bytes)


async def _local_rembg_remove_bytes(img_bytes: bytes) -> bytes | None:
    if not LOCAL_REMBG_ENABLED:
        return None
    if rembg_remove is None:
        log.warning("local rembg is not available: %s", REMBG_IMPORT_ERROR)
        return None

    if LOCAL_REMBG_SUBPROCESS:
        return await _local_rembg_remove_subprocess(img_bytes)

    try:
        return await asyncio.wait_for(
            asyncio.to_thread(_local_rembg_remove_sync, img_bytes),
            timeout=LOCAL_REMBG_TIMEOUT_S,
        )
    except Exception as e:
        log.warning("local rembg failed: %s", e)
        _bg_note_error(f"local rembg failed: {e}")
        return None




def _resize_image_bytes_for_bg_api(img_bytes: bytes, max_side: int | None = None) -> tuple[bytes, str, str]:
    """РЎР¶РёРјР°РөСӮ РІС…РҫРҙ РҝРөСҖРөРҙ Photoroom, СҮСӮРҫРұСӢ РҪРө Р»РҫРІРёСӮСҢ API/Telegram timeout РҪР° РұРҫР»СҢСҲРёС… С„РҫСӮРҫ."""
    max_side = int(max_side or PHOTOROOM_INPUT_MAX_SIDE or 1600)
    mime = sniff_image_mime(img_bytes) or "image/jpeg"
    if Image is None or max_side <= 0:
        ext = ".jpg" if mime == "image/jpeg" else (".png" if mime == "image/png" else ".webp")
        return img_bytes, f"image{ext}", mime
    try:
        im = Image.open(BytesIO(img_bytes))
        # Р”Р»СҸ РёСҒС…РҫРҙРҪСӢС… С„РҫСӮРҫ РҝСҖРҫР·СҖР°СҮРҪРҫСҒСӮСҢ РҪРө РҪСғР¶РҪР°, Photoroom СҒР°Рј РІРөСҖРҪС‘СӮ PNG/RGBA.
        im = ImageOps.exif_transpose(im) if ImageOps else im
        if max(im.size) > max_side:
            im.thumbnail((max_side, max_side), Image.LANCZOS)
        # JPEG СҒРёР»СҢРҪРҫ СғРјРөРҪСҢСҲР°РөСӮ СҖР°Р·РјРөСҖ Р·Р°РҝСҖРҫСҒР° Рё СғСҒРәРҫСҖСҸРөСӮ РұРөСҒРҝР»Р°СӮРҪСӢР№ API.
        if im.mode not in ("RGB", "L"):
            im = im.convert("RGB")
        bio = BytesIO()
        im.save(bio, format="JPEG", quality=92, optimize=True, progressive=True)
        return bio.getvalue(), "image.jpg", "image/jpeg"
    except Exception as e:
        _bg_note_error(f"input resize for Photoroom failed: {e}")
        ext = ".jpg" if mime == "image/jpeg" else (".png" if mime == "image/png" else ".webp")
        return img_bytes, f"image{ext}", mime

async def _photoroom_api_remove_bytes(img_bytes: bytes) -> bytes | None:
    """Photoroom Remove Background API. Returns transparent PNG/RGBA bytes."""
    if not PHOTOROOM_API_KEY:
        _bg_note_error("Photoroom API key missing: set PHOTOROOM_API_KEY in Render Environment")
        return None
    upload_bytes, upload_name, mime = _resize_image_bytes_for_bg_api(img_bytes, PHOTOROOM_INPUT_MAX_SIDE)
    url = f"{PHOTOROOM_BASE_URL}{PHOTOROOM_REMOVE_PATH}"
    headers = {"x-api-key": PHOTOROOM_API_KEY}
    # Photoroom /v1/segment accepts remove.bg-compatible fields.
    data = {
        "format": PHOTOROOM_FORMAT or "png",
        "channels": PHOTOROOM_CHANNELS or "rgba",
        "size": PHOTOROOM_SIZE or "hd",
        "crop": PHOTOROOM_CROP or "false",
        "despill": PHOTOROOM_DESPILL or "false",
    }
    timeout = httpx.Timeout(PHOTOROOM_TIMEOUT_S, connect=20.0)
    try:
        async with httpx.AsyncClient(timeout=timeout, follow_redirects=True) as client:
            files = {"image_file": (upload_name, upload_bytes, mime)}
            r = await client.post(url, headers=headers, data=data, files=files)
            if r.status_code >= 400:
                body = (r.text or "")[:900]
                _bg_note_error(f"Photoroom API failed status={r.status_code} body={body}")
                log.warning("Photoroom API failed status=%s body=%s", r.status_code, body)
                return None
            ctype = (r.headers.get("content-type") or "").lower()
            if not r.content or len(r.content) < 300:
                _bg_note_error(f"Photoroom API returned empty/short content: {len(r.content or b'')} bytes")
                return None
            if (not ctype.startswith("image/")) and (r.content[:1] not in (b"\x89", b"\xff")):
                _bg_note_error(f"Photoroom API returned non-image content-type={ctype} body={(r.text or '')[:700]}")
                return None
            return bytes(r.content)
    except Exception as e:
        _bg_note_error(f"Photoroom API exception: {type(e).__name__}: {e}")
        log.warning("Photoroom API exception: %s", e)
        return None


def _normalize_photo_result(img_bytes: bytes) -> bytes:
    if Image is None:
        return img_bytes
    try:
        im = Image.open(BytesIO(img_bytes)).convert("RGB")
        max_side = int(BG_OUTPUT_MAX_SIDE or 1600)
        if max_side > 0 and max(im.size) > max_side:
            im.thumbnail((max_side, max_side), Image.LANCZOS)
        bio = BytesIO()
        im.save(bio, format="JPEG", quality=93, optimize=True, progressive=True)
        return bio.getvalue()
    except Exception as e:
        _bg_note_error(f"normalize photo result failed: {e}")
        return img_bytes


def _clean_bg_prompt_text(prompt: str) -> str:
    """Remove wording that makes AI draw a visible phone/mirror instead of just selfie-like perspective."""
    s = (prompt or "").strip()
    replacements = {
        "РҪР° СӮРөР»РөС„РҫРҪ": "",
        "СҒ СӮРөР»РөС„РҫРҪР°": "",
        "РҪР° СҒРјР°СҖСӮС„РҫРҪ": "",
        "СҒРҫ СҒРјР°СҖСӮС„РҫРҪР°": "",
        "СӮРөР»РөС„РҫРҪ": "РәР°РјРөСҖР°",
        "СҒРјР°СҖСӮС„РҫРҪ": "РәР°РјРөСҖР°",
        "phone": "camera",
        "smartphone": "camera",
        "mirror": "",
        "Р·РөСҖРәР°Р»Рҫ": "",
        "selfie stick": "",
        "СҒРөР»С„Рё-РҝР°Р»РәР°": "",
    }
    for a, b in replacements.items():
        s = re.sub(re.escape(a), b, s, flags=re.IGNORECASE)
    s = re.sub(r"\s{2,}", " ", s).strip(" ,.;")
    return s


def _build_background_scene_prompt(kind: str, prompt: str = "") -> str:
    kind = (kind or "custom").lower().strip()
    user_prompt = _clean_bg_prompt_text(prompt)
    preset_map = {
        "beach": (
            "photorealistic tropical beach background, natural sandy shore, clean sea horizon, calm blue water, "
            "soft waves, realistic daylight, believable travel-photo atmosphere, open air, uncluttered foreground, "
            "no buildings unless naturally distant"
        ),
        "mountains": (
            "photorealistic mountain landscape, scenic alpine or green mountains, open outdoor view, realistic sky, "
            "natural daylight, atmospheric perspective, clean travel-photo composition, no indoor elements"
        ),
        "nature": (
            "photorealistic park or nature background, greenery, trees, soft depth, realistic outdoor daylight, "
            "calm natural environment, clean believable background"
        ),
        "roof": (
            "photorealistic rooftop terrace or city skyline, elegant urban atmosphere, realistic architecture in the distance, "
            "natural perspective, outdoor light, clean modern background"
        ),
        "office": (
            "photorealistic premium office or business-lounge background, modern interior, clean lines, "
            "soft daylight, realistic depth, elegant professional atmosphere, no random gadgets in foreground"
        ),
        "white": "clean white studio background with soft even light and realistic subtle shadow",
        "black": "clean dark studio background with soft controlled light and realistic subtle shadow",
        "blur": "soft natural bokeh background derived from the original scene, realistic lens blur, no extra objects",
        "custom": user_prompt,
    }
    scene = user_prompt or preset_map.get(kind) or "photorealistic clean environment"
    return (
        "Create only the BACKGROUND SCENE with no people. This is step 2 of a two-stage workflow: "
        "the subject will be composited later, so leave clean free space for one adult person in the center foreground. "
        "The result must be photorealistic, believable, natural perspective, eye-level camera angle, realistic light, "
        "real-world textures, coherent depth and a clean uncluttered composition. "
        f"Scene request: {scene}."
    ).strip()


def _build_background_negative_prompt(kind: str = "", prompt: str = "") -> str:
    parts = [
        "no people", "no portraits", "no selfie", "no face", "no hands", "no body",
        "no phone", "no smartphone", "no screen", "no mirror", "no selfie stick",
        "no black panel", "no kiosk", "no wall device", "no random indoor artifact",
        "no duplicate objects", "no text", "no watermark", "no logo", "no cartoon",
        "no painting", "no illustration", "no anime", "no CGI", "no 3d render"
    ]
    kind = (kind or "").lower().strip()
    if kind == "beach":
        parts += ["no snow", "no mountains in foreground", "no office"]
    elif kind == "mountains":
        parts += ["no beach", "no ocean", "no tropical palm trees"]
    elif kind == "office":
        parts += ["no beach", "no mountains", "no random device close to camera"]
    return ", ".join(dict.fromkeys(parts))


def _build_selfie_background_prompt(kind: str, prompt: str = "") -> str:
    # Legacy helper kept for compatibility/debug.
    scene = _build_background_scene_prompt(kind, prompt)
    negative = _build_background_negative_prompt(kind, prompt)
    return (
        "Generate only a new background and keep the original main subject completely unchanged: "
        "do not alter the face, hair, clothes, body, pose, proportions or skin texture. "
        + scene + " Avoid: " + negative
    ).strip()


def _should_use_photoroom_ai_background(kind: str, prompt: str = "") -> bool:
    # In production we now prefer a strict two-stage pipeline: remove background -> build/select new background -> composite.
    # Photoroom edit is kept only as an optional debug path and is disabled for presets by default.
    if not PHOTOROOM_EDIT_ENABLED or not PHOTOROOM_API_KEY:
        return False
    return False


async def _photoroom_api_edit_background_bytes(img_bytes: bytes, kind: str = "custom", prompt: str = "") -> bytes | None:
    """Legacy direct AI edit path. Kept for optional diagnostics, not used by the default two-stage production pipeline."""
    if not PHOTOROOM_API_KEY:
        _bg_note_error("Photoroom API key missing for AI background: set PHOTOROOM_API_KEY in Render Environment")
        return None
    if not PHOTOROOM_EDIT_ENABLED:
        _bg_note_error("Photoroom AI background is disabled by PHOTOROOM_EDIT_ENABLED=0")
        return None
    upload_bytes, upload_name, mime = _resize_image_bytes_for_bg_api(img_bytes, PHOTOROOM_INPUT_MAX_SIDE)
    url = f"{PHOTOROOM_EDIT_BASE_URL}{PHOTOROOM_EDIT_PATH}"
    headers = {"x-api-key": PHOTOROOM_API_KEY}
    bg_prompt = _build_selfie_background_prompt(kind, prompt)
    data = {
        "referenceBox": "originalImage",
        "background.prompt": bg_prompt,
        "background.expandPrompt.mode": PHOTOROOM_EDIT_EXPAND_PROMPT_MODE or "ai.never",
    }
    if PHOTOROOM_EDIT_NEGATIVE_PROMPT:
        data["background.negativePrompt"] = PHOTOROOM_EDIT_NEGATIVE_PROMPT
    timeout = httpx.Timeout(PHOTOROOM_EDIT_TIMEOUT_S, connect=20.0)
    try:
        async with httpx.AsyncClient(timeout=timeout, follow_redirects=True) as client:
            files = {"imageFile": (upload_name, upload_bytes, mime)}
            r = await client.post(url, headers=headers, data=data, files=files)
            if r.status_code >= 400:
                body = (r.text or "")[:900]
                _bg_note_error(f"Photoroom edit API failed status={r.status_code} body={body}")
                log.warning("Photoroom edit API failed status=%s body=%s", r.status_code, body)
                return None
            ctype = (r.headers.get("content-type") or "").lower()
            if ctype.startswith("image/") and r.content and len(r.content) > 300:
                return bytes(r.content)
            out = await _image_bytes_from_response(r, client)
            if out:
                return out
            _bg_note_error(f"Photoroom edit API returned non-image content-type={ctype} body={(r.text or '')[:700]}")
            return None
    except Exception as e:
        _bg_note_error(f"Photoroom edit API exception: {type(e).__name__}: {e}")
        log.warning("Photoroom edit API exception: %s", e)
        return None


async def _removebg_api_remove_bytes(img_bytes: bytes) -> bytes | None:
    """Official remove.bg-compatible API primary path. Returns transparent PNG bytes."""
    if not REMOVE_BG_API_KEY:
        _bg_note_error("remove.bg API key missing: set REMOVE_BG_API_KEY in Render Environment")
        return None
    mime = sniff_image_mime(img_bytes) or "image/jpeg"
    url = f"{REMOVE_BG_BASE_URL}{REMOVE_BG_PATH}"
    headers = {"X-Api-Key": REMOVE_BG_API_KEY}
    data = {
        "size": REMOVE_BG_SIZE,
        "format": REMOVE_BG_FORMAT,
        "type": "auto",
    }
    timeout = httpx.Timeout(REMOVE_BG_TIMEOUT_S, connect=20.0)
    try:
        async with httpx.AsyncClient(timeout=timeout, follow_redirects=True) as client:
            files = {"image_file": ("image.jpg", img_bytes, mime)}
            r = await client.post(url, headers=headers, data=data, files=files)
            if r.status_code >= 400:
                body = (r.text or "")[:700]
                _bg_note_error(f"remove.bg API failed status={r.status_code} body={body}")
                log.warning("remove.bg API failed status=%s body=%s", r.status_code, body)
                return None
            if not r.content or len(r.content) < 300:
                _bg_note_error(f"remove.bg API returned empty/short content: {len(r.content or b'')} bytes")
                return None
            return bytes(r.content)
    except Exception as e:
        _bg_note_error(f"remove.bg API exception: {type(e).__name__}: {e}")
        log.warning("remove.bg API exception: %s", e)
        return None


async def _comet_bria_remove_bg_bytes(img_bytes: bytes) -> bytes | None:
    """Remote fallback only. Main production path should be local rembg."""
    if not COMET_API_KEY:
        return None
    headers = {"Authorization": f"Bearer {COMET_API_KEY}"}
    mime = sniff_image_mime(img_bytes)
    paths: list[str] = []
    for p in (BG_COMET_REMOVE_PATH, "/v1/images/edits", "/v1/images/generations"):
        if p and p not in paths:
            paths.append(p)
    timeout = httpx.Timeout(BG_REMOVE_TIMEOUT_S, connect=20.0)
    async with httpx.AsyncClient(timeout=timeout, follow_redirects=True) as client:
        for path in paths:
            url = f"{COMET_BASE_URL}{path}"
            try:
                files = {"image": ("image.png", img_bytes, mime or "application/octet-stream")}
                data = {
                    "model": BG_COMET_MODEL,
                    "response_format": "b64_json",
                    "transparent_background": "true",
                }
                r = await client.post(url, headers=headers, files=files, data=data)
                if r.status_code >= 400:
                    err = f"Comet BG remove failed path={path} status={r.status_code} body={r.text[:500]}"
                    log.warning(err)
                    _bg_note_error(err)
                    continue
                out = await _image_bytes_from_response(r, client)
                if out:
                    return out
            except Exception as e:
                msg = f"Comet background remove exception path={path}: {e}"
                log.warning(msg)
                _bg_note_error(msg)
    return None


async def _bria_direct_remove_bg_bytes(img_bytes: bytes) -> bytes | None:
    """Direct Bria fallback, only if BRIA_API_KEY is configured."""
    if not BRIA_API_KEY:
        return None
    mime = sniff_image_mime(img_bytes)
    headers_variants = [
        {"Authorization": f"Bearer {BRIA_API_KEY}"},
        {"api_token": BRIA_API_KEY},
        {"Authorization": f"Bearer {BRIA_API_KEY}", "api_token": BRIA_API_KEY},
    ]
    paths: list[str] = []
    for p in (BRIA_REMOVE_PATH, "/v1/background/remove", "/v1/remove_background", "/background/remove"):
        if p and p not in paths:
            paths.append(p)
    timeout = httpx.Timeout(BG_REMOVE_TIMEOUT_S, connect=20.0)
    async with httpx.AsyncClient(timeout=timeout, follow_redirects=True) as client:
        for path in paths:
            url = f"{BRIA_BASE_URL}{path}"
            for headers in headers_variants:
                try:
                    files = {"image": ("image.png", img_bytes, mime or "application/octet-stream")}
                    r = await client.post(url, headers=headers, files=files)
                    if r.status_code >= 400:
                        _bg_note_error(f"Bria direct failed path={path} status={r.status_code} body={r.text[:400]}")
                        continue
                    out = await _image_bytes_from_response(r, client)
                    if out:
                        return out
                except Exception as e:
                    msg = f"Bria direct remove exception path={path}: {e}"
                    log.warning(msg)
                    _bg_note_error(msg)
    return None


def _normalize_transparent_png(img_bytes: bytes) -> bytes:
    if Image is None:
        return img_bytes
    try:
        im = Image.open(BytesIO(img_bytes)).convert("RGBA")
        max_side = int(BG_OUTPUT_MAX_SIDE or 1600)
        if max_side > 0 and max(im.size) > max_side:
            im.thumbnail((max_side, max_side), Image.LANCZOS)
        bio = BytesIO()
        im.save(bio, format="PNG", optimize=True, compress_level=6)
        return bio.getvalue()
    except Exception as e:
        _bg_note_error(f"normalize transparent png failed: {e}")
        return img_bytes

async def _remove_bg_bytes_primary(img_bytes: bytes) -> bytes | None:
    """
    Production background pipeline.
    Primary: Photoroom Remove Background API.
    Optional fallbacks are used only in provider=multi/auto. Local rembg remains disabled on Render Starter.
    """
    provider = (BG_PROVIDER or "photoroom-api-only").lower().strip()
    photoroom_only = provider in (
        "photoroom", "photoroom-api", "photoroom-api-only",
        "photoroom-only", "api", "api-only"
    )

    # 1) Stable production path: Photoroom API.
    if provider in ("auto", "multi", "photoroom", "photoroom-api", "photoroom-api-only", "photoroom-only", "api", "api-only"):
        out = await _photoroom_api_remove_bytes(img_bytes)
        if out:
            return _normalize_transparent_png(out)
        if photoroom_only:
            return None

    # 2) Optional legacy remove.bg-compatible fallback if explicitly configured.
    if provider in ("auto", "multi", "removebg", "remove.bg", "removebg-api", "removebg-api-only", "removebg-only"):
        out = await _removebg_api_remove_bytes(img_bytes)
        if out:
            return _normalize_transparent_png(out)
        if provider in ("removebg", "remove.bg", "removebg-api", "removebg-api-only", "removebg-only"):
            return None

    # 3) Optional Comet/Bria fallback only when provider allows remote fallback.
    remote_allowed = provider not in (
        "local-only", "rembg-only", "local", "rembg",
        "photoroom-api-only", "photoroom-only", "api-only"
    )
    if remote_allowed and provider in ("auto", "multi", "comet", "cometapi", "bria-comet", "bria/remove-background"):
        out = await _comet_bria_remove_bg_bytes(img_bytes)
        if out:
            return _normalize_transparent_png(out)

    if remote_allowed and provider in ("auto", "multi", "bria", "direct-bria", "bria-direct"):
        out = await _bria_direct_remove_bg_bytes(img_bytes)
        if out:
            return _normalize_transparent_png(out)

    # 4) Last-resort local rembg only if explicitly enabled; disabled in production on Render Starter.
    local_allowed = bool((not BG_DISABLE_LOCAL_REMBG) and LOCAL_REMBG_ENABLED and rembg_remove is not None and provider not in ("photoroom-api-only", "photoroom-only", "removebg-api-only", "removebg-only", "api-only"))
    if local_allowed:
        out = await _local_rembg_remove_bytes(img_bytes)
        if out:
            return _normalize_transparent_png(out)

    if not PHOTOROOM_API_KEY and provider in ("photoroom", "photoroom-api", "photoroom-api-only", "photoroom-only", "api", "api-only"):
        _bg_note_error("PHOTOROOM_API_KEY is not configured")
    if rembg_remove is None and not BG_DISABLE_LOCAL_REMBG:
        _bg_note_error(f"local rembg unavailable: import failed {REMBG_IMPORT_ERROR}")
    return None

def _fit_cover(im, size: tuple[int, int]):
    if ImageOps:
        return ImageOps.fit(im, size, method=Image.LANCZOS, centering=(0.5, 0.5))
    return im.resize(size, Image.LANCZOS)


def _gradient_background(size: tuple[int, int], top: tuple[int, int, int], bottom: tuple[int, int, int]):
    w, h = size
    bg = Image.new("RGB", size, top)
    if ImageDraw is None:
        return bg
    draw = ImageDraw.Draw(bg)
    for y in range(max(1, h)):
        t = y / max(1, h - 1)
        c = tuple(int(top[i] * (1 - t) + bottom[i] * t) for i in range(3))
        draw.line([(0, y), (w, y)], fill=c)
    return bg



_REAL_BG_URLS = {
    # Р РөР°Р»СҢРҪСӢРө С„РҫСӮРҫСҒСҶРөРҪСӢ, Р° РҪРө СҖРёСҒРҫРІР°РҪРҪСӢРө Р·Р°РіР»СғСҲРәРё. Р•СҒР»Рё СҒРөСӮСҢ РҪРөРҙРҫСҒСӮСғРҝРҪР° вҖ” РҪРёР¶Рө РҫСҒСӮР°РҪРөСӮСҒСҸ Р»РҫРәР°Р»СҢРҪСӢР№ fallback.
    "beach": "https://images.unsplash.com/photo-1507525428034-b723cf961d3e?auto=format&fit=crop&w=1800&q=85",
    "mountains": "https://images.unsplash.com/photo-1506905925346-21bda4d32df4?auto=format&fit=crop&w=1800&q=85",
    "nature": "https://images.unsplash.com/photo-1448375240586-882707db888b?auto=format&fit=crop&w=1800&q=85",
    "roof": "https://images.unsplash.com/photo-1480714378408-67cf0d13bc1f?auto=format&fit=crop&w=1800&q=85",
    "office": "https://images.unsplash.com/photo-1497366754035-f200968a6e72?auto=format&fit=crop&w=1800&q=85",
}


def _safe_cache_name(key: str) -> str:
    return re.sub(r"[^a-zA-Z0-9_.-]+", "_", key)[:80] or "bg"


def _real_photo_background(size: tuple[int, int], kind: str):
    """Р‘РөСҖС‘СӮ РҪР°СҒСӮРҫСҸСүРёР№ С„РҫСӮРҫС„РҫРҪ РёР· РәСҚСҲР°/URL Рё РәСҖРҫРҝРёСӮ РҝРҫРҙ СҖР°Р·РјРөСҖ РҫРұСҠРөРәСӮР°."""
    if Image is None or not BG_REALISTIC_BACKGROUNDS:
        return None
    kind = (kind or "").lower().strip()
    url = _REAL_BG_URLS.get(kind)
    if not url:
        return None
    try:
        os.makedirs(BG_CACHE_DIR, exist_ok=True)
        cache_path = os.path.join(BG_CACHE_DIR, _safe_cache_name(kind) + ".jpg")
        data = None
        if os.path.exists(cache_path) and os.path.getsize(cache_path) > 1024:
            with open(cache_path, "rb") as f:
                data = f.read()
        else:
            with httpx.Client(timeout=BG_BACKGROUND_TIMEOUT_S, follow_redirects=True) as client:
                r = client.get(url, headers={"User-Agent": "GPT5ProBot/1.0"})
                r.raise_for_status()
                data = r.content
            with contextlib.suppress(Exception):
                with open(cache_path, "wb") as f:
                    f.write(data)
        if not data:
            return None
        bg = Image.open(BytesIO(data)).convert("RGB")
        bg = _fit_cover(bg, size)
        # РӣС‘РіРәРҫРө СҖР°Р·РјСӢСӮРёРө Рё Р·Р°СӮРөРјРҪРөРҪРёРө, СҮСӮРҫРұСӢ РҫРұСҠРөРәСӮ РҪРө РІСӢРіР»СҸРҙРөР» РІРәР»РөРөРҪРҪСӢРј СҒР»РёСҲРәРҫРј СҖРөР·РәРҫ.
        if ImageFilter:
            bg = bg.filter(ImageFilter.GaussianBlur(radius=max(1.2, min(size) / 450)))
        return bg
    except Exception as e:
        _bg_note_error(f"real background download failed kind={kind}: {type(e).__name__}: {e}")
        log.warning("real background failed kind=%s: %s", kind, e)
        return None

def _make_local_background(size: tuple[int, int], kind: str, original_bytes: bytes | None = None, prompt: str = ""):
    w, h = size
    kind = (kind or "blur").lower()
    if BG_REPLACE_USE_STOCK_BACKGROUNDS and kind in ("beach", "mountains", "nature", "roof", "office"):
        real_bg = _real_photo_background(size, kind)
        if real_bg is not None:
            return real_bg
    if kind == "blur" and original_bytes and ImageFilter:
        base = Image.open(BytesIO(original_bytes)).convert("RGB")
        bg = _fit_cover(base, size)
        return bg.filter(ImageFilter.GaussianBlur(radius=max(18, min(w, h) // 18)))
    if kind == "white":
        return Image.new("RGB", size, (255, 255, 255))
    if kind == "black":
        return Image.new("RGB", size, (18, 18, 18))

    if kind == "beach":
        bg = _gradient_background(size, (114, 194, 237), (245, 219, 162))
        if ImageDraw:
            d = ImageDraw.Draw(bg)
            d.rectangle([0, int(h * 0.64), w, h], fill=(233, 203, 147))
            d.rectangle([0, int(h * 0.48), w, int(h * 0.66)], fill=(58, 161, 202))
        return bg

    if kind == "mountains":
        bg = _gradient_background(size, (126, 176, 229), (230, 237, 244))
        if ImageDraw:
            d = ImageDraw.Draw(bg)
            d.polygon([(0, h), (int(w*0.26), int(h*0.40)), (int(w*0.58), h)], fill=(97, 114, 126))
            d.polygon([(int(w*0.34), h), (int(w*0.67), int(h*0.28)), (w, h)], fill=(78, 93, 108))
        return bg

    if kind == "roof":
        bg = _gradient_background(size, (96, 123, 170), (34, 38, 56))
        if ImageDraw:
            d = ImageDraw.Draw(bg)
            for i in range(8):
                x0 = int(w * (i / 8.0))
                bw = max(18, w // 12)
                bh = int(h * (0.15 + (i % 4) * 0.08))
                d.rectangle([x0, h - bh, x0 + bw, h], fill=(40, 46, 62))
        return bg

    if kind == "office":
        bg = _gradient_background(size, (242, 244, 247), (212, 219, 228))
        if ImageDraw:
            d = ImageDraw.Draw(bg)
            d.rectangle([0, int(h*0.70), w, h], fill=(210, 196, 178))
            for x in range(0, w, max(40, w // 6)):
                d.rectangle([x, int(h*0.18), min(w, x + max(24, w // 12)), int(h*0.58)], fill=(190, 205, 219))
        return bg

    # nature/custom fallback: СҒРҝРҫРәРҫР№РҪСӢР№ Р·РөР»С‘РҪСӢР№ С„РҫРҪ СҒ РіР»СғРұРёРҪРҫР№.
    bg = _gradient_background(size, (125, 190, 140), (35, 80, 50))
    if ImageDraw:
        d = ImageDraw.Draw(bg)
        for i in range(12):
            x = int(w * i / 11)
            y = int(h * (0.45 + (i % 3) * 0.07))
            r = max(30, w // 12)
            d.ellipse([x-r, y-r, x+r, y+r], fill=(45, 115, 65))
        d.rectangle([0, int(h*0.70), w, h], fill=(42, 95, 55))
    return bg


def _compose_subject_on_background(bg_rgb, fg_rgba):
    bg = bg_rgb.convert("RGBA") if getattr(bg_rgb, "mode", "") != "RGBA" else bg_rgb.copy()
    fg = fg_rgba.convert("RGBA")
    alpha = fg.getchannel("A")
    if ImageFilter:
        try:
            shadow_mask = alpha.filter(ImageFilter.GaussianBlur(radius=max(6, min(fg.size)//90)))
            shadow = Image.new("RGBA", fg.size, (0, 0, 0, 0))
            shadow.putalpha(shadow_mask.point(lambda p: int(p * 0.20)))
            dx = max(2, fg.size[0] // 100)
            dy = max(2, fg.size[1] // 90)
            bg.alpha_composite(shadow, dest=(dx, dy))
        except Exception:
            pass
    bg.alpha_composite(fg)
    return bg


async def _generate_background_only_bytes(size: tuple[int, int], kind: str = "custom", prompt: str = "") -> bytes | None:
    if Image is None:
        return None
    kind = (kind or "custom").lower().strip()
    want_gen = (kind == "custom" and BG_REPLACE_GENERATE_CUSTOM) or (kind != "custom" and BG_REPLACE_GENERATE_PRESETS)
    if not want_gen:
        return None
    full_prompt = _build_background_scene_prompt(kind, prompt)
    negative = _build_background_negative_prompt(kind, prompt)
    final_prompt = f"{full_prompt} Avoid: {negative}."
    try:
        data = await _luma_generate_image_bytes(final_prompt)
        if not data:
            _bg_note_error(f"background-only generation returned empty output for kind={kind}")
            return None
        bg = Image.open(BytesIO(data)).convert("RGB")
        bg = _fit_cover(bg, size)
        bio = BytesIO()
        bg.save(bio, format="JPEG", quality=max(86, min(98, BG_REPLACE_JPEG_QUALITY)), optimize=True, progressive=True)
        return bio.getvalue()
    except Exception as e:
        _bg_note_error(f"background-only generation failed kind={kind}: {type(e).__name__}: {e}")
        log.warning("background-only generation failed kind=%s: %s", kind, e)
        return None


async def _compose_replace_bg(img_bytes: bytes, bg_kind: str = "blur", prompt: str = "") -> bytes | None:
    if Image is None:
        return None

    # Strict production pipeline: 1) remove background, 2) build/select new background, 3) composite subject back.
    cutout_bytes = await _remove_bg_bytes_primary(img_bytes)
    if not cutout_bytes:
        return None

    fg = Image.open(BytesIO(cutout_bytes)).convert("RGBA")
    target_size = fg.size

    bg_rgb = None
    gen_bytes = await _generate_background_only_bytes(target_size, kind=bg_kind, prompt=prompt)
    if gen_bytes:
        try:
            bg_rgb = Image.open(BytesIO(gen_bytes)).convert("RGB")
        except Exception:
            bg_rgb = None

    if bg_rgb is None:
        bg_rgb = _make_local_background(target_size, bg_kind, original_bytes=img_bytes, prompt=prompt).convert("RGB")

    composed = _compose_subject_on_background(bg_rgb, fg)
    bio = BytesIO()
    composed.convert("RGB").save(
        bio,
        format="JPEG",
        quality=max(86, min(98, BG_REPLACE_JPEG_QUALITY)),
        optimize=True,
        progressive=True,
    )
    return bio.getvalue()


async def _pedit_removebg(update: Update, context: ContextTypes.DEFAULT_TYPE, img_bytes: bytes):
    if not context.user_data.pop("_image_processing_quota_ok", False):
        async def _go():
            context.user_data["_image_processing_quota_ok"] = True
            return await _pedit_removebg(update, context, img_bytes)
        await _try_pay_then_do(
            update, context, update.effective_user.id,
            "img", IMG_PROCESS_COST_USD, _go,
            remember_kind="removebg",
        )
        return
    try:
        await update.effective_message.reply_text("рҹ§ј РЈРҙР°Р»СҸСҺ С„РҫРҪ. Р’РөСҖРҪСғ PNG СҒ РҝСҖРҫР·СҖР°СҮРҪРҫР№ РҝРҫРҙР»РҫР¶РәРҫР№.")
        out = await asyncio.wait_for(_remove_bg_bytes_primary(img_bytes), timeout=BG_ACTION_TIMEOUT_S)
        if not out:
            if rembg_remove is None:
                await update.effective_message.reply_text(
                    "вқҢ РқРө СғРҙР°Р»РҫСҒСҢ СғРҙР°Р»РёСӮСҢ С„РҫРҪ СҮРөСҖРөР· Photoroom API. "
                    "РҹСҖРҫРІРөСҖСҢСӮРө PHOTOROOM_API_KEY Рё Р»РёРјРёСӮСӢ СӮРөСҒСӮРҫРІРҫРіРҫ РәР»СҺСҮР°. РҹРҫСҒР»РөРҙРҪРёРө РҫСҲРёРұРәРё:\n" + _bg_last_errors_text()
                )
            else:
                await update.effective_message.reply_text("вқҢ РқРө СғРҙР°Р»РҫСҒСҢ СғРҙР°Р»РёСӮСҢ С„РҫРҪ. Р”Р»СҸ СӮРҫСҮРҪРҫР№ РҝСҖРёСҮРёРҪСӢ Р·Р°РҝСғСҒСӮРёСӮРө /diag_bg. РҹРҫСҒР»РөРҙРҪРёРө РҫСҲРёРұРәРё:\n" + _bg_last_errors_text())
            return False
        bio = BytesIO(out)
        bio.name = "no_bg.png"
        await update.effective_message.reply_document(
            InputFile(bio),
            caption="РӨРҫРҪ СғРҙР°Р»С‘РҪ вң… PNG СҒ РҝСҖРҫР·СҖР°СҮРҪРҫР№ РҝРҫРҙР»РҫР¶РәРҫР№.",
            read_timeout=120,
            write_timeout=120,
            connect_timeout=30,
            pool_timeout=30,
        )
        return True
    except TimedOut as e:
        _bg_note_error(f"telegram upload timeout: {type(e).__name__}: {e}")
        await update.effective_message.reply_text(
            "вқҢ РӨРҫРҪ СғРҙР°Р»С‘РҪ, РҪРҫ Telegram РҪРө СғСҒРҝРөР» РҝСҖРёРҪСҸСӮСҢ PNG-С„Р°Р№Р». "
            "РҜ СғРјРөРҪСҢСҲРёР» СҖР°Р·РјРөСҖ С„Р°Р№Р»РҫРІ РІ РҪРҫРІРҫР№ РІРөСҖСҒРёРё; РҝРҫРҝСҖРҫРұСғР№СӮРө РөСүС‘ СҖР°Р· РёР»Рё РҫСӮРҝСҖР°РІСҢСӮРө С„РҫСӮРҫ РјРөРҪСҢСҲРөРіРҫ СҖР°Р·РјРөСҖР°."
        )
        return False
    except asyncio.TimeoutError:
        _bg_note_error(f"removebg action timeout after {BG_ACTION_TIMEOUT_S:.0f}s")
        await update.effective_message.reply_text("вқҢ РЈРҙР°Р»РөРҪРёРө С„РҫРҪР° Р·Р°РҪСҸР»Рҫ СҒР»РёСҲРәРҫРј РјРҪРҫРіРҫ РІСҖРөРјРөРҪРё Рё РұСӢР»Рҫ РҫСҒСӮР°РҪРҫРІР»РөРҪРҫ. РҹРҫРҝСҖРҫРұСғР№СӮРө С„РҫСӮРҫ РјРөРҪСҢСҲРөРіРҫ СҖР°Р·РјРөСҖР° РёР»Рё РҝРҫРІСӮРҫСҖРёСӮРө РҝРҫР·Р¶Рө. РҹРҫСҒР»РөРҙРҪРёРө РҫСҲРёРұРәРё:\n" + _bg_last_errors_text())
        return False
    except Exception as e:
        log.exception("removebg error: %s", e)
        _bg_note_error(f"removebg exception: {type(e).__name__}: {e}")
        await update.effective_message.reply_text("РқРө СғРҙР°Р»РҫСҒСҢ СғРҙР°Р»РёСӮСҢ С„РҫРҪ. РҹРҫСҒР»РөРҙРҪРёРө РҫСҲРёРұРәРё:\n" + _bg_last_errors_text())
        return False


async def _pedit_replacebg(update: Update, context: ContextTypes.DEFAULT_TYPE, img_bytes: bytes, kind: str = "blur", prompt: str = ""):
    if not context.user_data.pop("_image_processing_quota_ok", False):
        async def _go():
            context.user_data["_image_processing_quota_ok"] = True
            return await _pedit_replacebg(update, context, img_bytes, kind=kind, prompt=prompt)
        await _try_pay_then_do(
            update, context, update.effective_user.id,
            "img", IMG_PROCESS_COST_USD, _go,
            remember_kind="replacebg",
        )
        return
    if Image is None:
        await update.effective_message.reply_text("Pillow РҪРө СғСҒСӮР°РҪРҫРІР»РөРҪ.")
        return False
    try:
        await update.effective_message.reply_text("рҹ–ј Р—Р°РҝСғСҒРәР°СҺ РҙРІСғС…СҚСӮР°РҝРҪСғСҺ Р·Р°РјРөРҪСғ С„РҫРҪР°: 1) Р°РәРәСғСҖР°СӮРҪРҫ РІСӢСҖРөР·Р°СҺ СҮРөР»РҫРІРөРәР°/РҫРұСҠРөРәСӮ, 2) РҫСӮРҙРөР»СҢРҪРҫ РҝРҫРҙРұРёСҖР°СҺ РёР»Рё РіРөРҪРөСҖРёСҖСғСҺ РҪРҫРІСӢР№ С„РҫРҪ, 3) СҒРҫРұРёСҖР°СҺ РёСӮРҫРі РұРөР· РҝРөСҖРөСҖРёСҒРҫРІРәРё Р»РёСҶР°, РҫРҙРөР¶РҙСӢ Рё РҝРҫР·СӢ.")
        out = await asyncio.wait_for(_compose_replace_bg(img_bytes, bg_kind=kind, prompt=prompt), timeout=BG_ACTION_TIMEOUT_S)
        if not out:
            await update.effective_message.reply_text("вқҢ РқРө СғРҙР°Р»РҫСҒСҢ РҫСӮРҙРөР»РёСӮСҢ РҫРұСҠРөРәСӮ РҫСӮ С„РҫРҪР°. Р”Р»СҸ СӮРҫСҮРҪРҫР№ РҝСҖРёСҮРёРҪСӢ Р·Р°РҝСғСҒСӮРёСӮРө /diag_bg. РҹРҫСҒР»РөРҙРҪРёРө РҫСҲРёРұРәРё:\n" + _bg_last_errors_text())
            return False
        bio = BytesIO(out)
        bio.name = f"replace_bg_{kind or 'custom'}.jpg"
        cap = "РӨРҫРҪ Р·Р°РјРөРҪС‘РҪ вң… РһРұСҠРөРәСӮ РҫСҒСӮР°РІР»РөРҪ РёР· РёСҒС…РҫРҙРҪРҫРіРҫ С„РҫСӮРҫ."
        if prompt:
            cap += f"\nР—Р°РҝСҖРҫСҒ С„РҫРҪР°: {prompt[:250]}"
        await update.effective_message.reply_photo(
            InputFile(bio),
            caption=cap,
            read_timeout=180,
            write_timeout=180,
            connect_timeout=30,
            pool_timeout=30,
        )
        return True
    except TimedOut as e:
        _bg_note_error(f"telegram replacebg upload timeout: {type(e).__name__}: {e}")
        await update.effective_message.reply_text(
            "вҡ пёҸ Telegram СҒР»РёСҲРәРҫРј РҙРҫР»РіРҫ РҝСҖРёРҪРёРјР°Р» РёР·РҫРұСҖР°Р¶РөРҪРёРө РҝРҫСҒР»Рө Р·Р°РјРөРҪСӢ С„РҫРҪР°. "
            "Р•СҒР»Рё С„РҫСӮРҫ СғР¶Рө РҝРҫСҸРІРёР»РҫСҒСҢ РІ СҮР°СӮРө вҖ” СҖРөР·СғР»СҢСӮР°СӮ СғСҒРҝРөСҲРҪРҫ РҙРҫСҒСӮР°РІР»РөРҪ. "
            "Р•СҒР»Рё С„РҫСӮРҫ РҪРө РҝРҫСҸРІРёР»РҫСҒСҢ, РҝРҫРІСӮРҫСҖРёСӮРө РҝРҫРҝСӢСӮРәСғ РөСүС‘ СҖР°Р· РёР»Рё РҫСӮРҝСҖР°РІСҢСӮРө С„РҫСӮРҫ РјРөРҪСҢСҲРөРіРҫ СҖР°Р·РјРөСҖР°."
        )
        return False
    except asyncio.TimeoutError:
        _bg_note_error(f"replacebg action timeout after {BG_ACTION_TIMEOUT_S:.0f}s")
        await update.effective_message.reply_text("вқҢ Р—Р°РјРөРҪР° С„РҫРҪР° Р·Р°РҪСҸР»Р° СҒР»РёСҲРәРҫРј РјРҪРҫРіРҫ РІСҖРөРјРөРҪРё Рё РұСӢР»Р° РҫСҒСӮР°РҪРҫРІР»РөРҪР°. РҹРҫРҝСҖРҫРұСғР№СӮРө С„РҫСӮРҫ РјРөРҪСҢСҲРөРіРҫ СҖР°Р·РјРөСҖР° РёР»Рё РҝРҫРІСӮРҫСҖРёСӮРө РҝРҫР·Р¶Рө. РҹРҫСҒР»РөРҙРҪРёРө РҫСҲРёРұРәРё:\n" + _bg_last_errors_text())
        return False
    except Exception as e:
        log.exception("replacebg error: %s", e)
        _bg_note_error(f"replacebg exception: {type(e).__name__}: {e}")
        await update.effective_message.reply_text("РқРө СғРҙР°Р»РҫСҒСҢ Р·Р°РјРөРҪРёСӮСҢ С„РҫРҪ. РҹРҫСҒР»РөРҙРҪРёРө РҫСҲРёРұРәРё:\n" + _bg_last_errors_text())
        return False


async def _pedit_outpaint(update: Update, context: ContextTypes.DEFAULT_TYPE, img_bytes: bytes):
    if not context.user_data.pop("_image_processing_quota_ok", False):
        async def _go():
            context.user_data["_image_processing_quota_ok"] = True
            return await _pedit_outpaint(update, context, img_bytes)
        await _try_pay_then_do(
            update, context, update.effective_user.id,
            "img", IMG_PROCESS_COST_USD, _go,
            remember_kind="outpaint",
        )
        return
    if Image is None:
        await update.effective_message.reply_text("Pillow РҪРө СғСҒСӮР°РҪРҫРІР»РөРҪ.")
        return False
    try:
        im = Image.open(BytesIO(img_bytes)).convert("RGB")
        pad = max(64, min(256, max(im.size)//6))
        big = Image.new("RGB", (im.width + 2*pad, im.height + 2*pad))
        bg = im.resize(big.size, Image.LANCZOS).filter(ImageFilter.GaussianBlur(radius=24)) if ImageFilter else im.resize(big.size)
        big.paste(bg, (0, 0)); big.paste(im, (pad, pad))
        bio = BytesIO(); big.save(bio, format="JPEG", quality=92); bio.seek(0); bio.name = "outpaint.jpg"
        await update.effective_message.reply_photo(InputFile(bio), caption="РҹСҖРҫСҒСӮРҫР№ outpaint: СҖР°СҒСҲРёСҖРёР» РҝРҫР»РҫСӮРҪРҫ СҒ РјСҸРіРәРёРјРё РәСҖР°СҸРјРё.")
        return True
    except Exception as e:
        log.exception("outpaint error: %s", e)
        await update.effective_message.reply_text("РқРө СғРҙР°Р»РҫСҒСҢ СҒРҙРөР»Р°СӮСҢ outpaint.")
        return False

async def _pedit_storyboard(update: Update, context: ContextTypes.DEFAULT_TYPE, img_bytes: bytes):
    try:
        b64 = base64.b64encode(img_bytes).decode("ascii")
        desc = await ask_openai_vision("РһРҝРёСҲРё РәР»СҺСҮРөРІСӢРө СҚР»РөРјРөРҪСӮСӢ РәР°РҙСҖР° РҫСҮРөРҪСҢ РәСҖР°СӮРәРҫ.", b64, sniff_image_mime(img_bytes))
        plan = await ask_openai_text(
            "РЎРҙРөР»Р°Р№ СҖР°СҒРәР°РҙСҖРҫРІРәСғ (6 РәР°РҙСҖРҫРІ) РҝРҫРҙ 6вҖ“10 СҒРөРәСғРҪРҙРҪСӢР№ РәР»РёРҝ. "
            "РҡР°Р¶РҙСӢР№ РәР°РҙСҖ вҖ” 1 СҒСӮСҖРҫРәР°: РәР°РҙСҖ/РҙРөР№СҒСӮРІРёРө/СҖР°РәСғСҖСҒ/СҒРІРөСӮ. РһСҒРҪРҫРІР°:\n" + (desc or "")
        )
        await update.effective_message.reply_text("Р Р°СҒРәР°РҙСҖРҫРІРәР°:\РҪ" + plan)
    except Exception as e:
        log.exception("storyboard error: %s", e)
        await update.effective_message.reply_text("РқРө СғРҙР°Р»РҫСҒСҢ РҝРҫСҒСӮСҖРҫРёСӮСҢ СҖР°СҒРәР°РҙСҖРҫРІРәСғ.")



# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ Face swap production helpers в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
_faceswap_target_cache = {}   # user_id -> bytes: photo where face must be replaced
_faceswap_source_cache = {}   # user_id -> bytes: face/reference photo
_faceswap_target_face_index_cache = {}  # user_id -> Segmind/API face index in target image
_faceswap_source_face_index_cache = {}  # user_id -> Segmind/API face index in source image
_faceswap_target_face_count_cache = {}  # user_id -> detected target faces count
_faceswap_source_face_count_cache = {}  # user_id -> detected source faces count
_faceswap_target_faces_cache = {}  # user_id -> detected target faces list with boxes
_faceswap_source_faces_cache = {}  # user_id -> detected source faces list with boxes
_faceswap_errors = []


def _faceswap_note_error(msg: str):
    try:
        msg = str(msg).strip()
        if not msg:
            return
        _faceswap_errors.append(msg[:1200])
        del _faceswap_errors[:-8]
        log.warning("faceswap: %s", msg)
    except Exception:
        pass


def _faceswap_last_errors_text() -> str:
    if not _faceswap_errors:
        return "РҫСҲРёРұРҫРә РҝРҫРәР° РҪРөСӮ"
    return "\n".join("вҖў " + e for e in _faceswap_errors[-5:])


def _faceswap_cv2_status() -> str:
    if not FACESWAP_FACE_DETECTION_ENABLED:
        return "disabled_by_env"
    try:
        import cv2  # type: ignore
        ver = getattr(cv2, "__version__", "")
        if not hasattr(cv2, "CascadeClassifier"):
            return f"BROKEN {ver}: no CascadeClassifier".strip()
        if not hasattr(cv2, "data") or not getattr(cv2.data, "haarcascades", ""):
            return f"BROKEN {ver}: no haarcascades".strip()
        return f"ok {ver}".strip()
    except Exception as e:
        return f"FAILED {type(e).__name__}: {e}"


async def cmd_diag_face(update: Update, context: ContextTypes.DEFAULT_TYPE):
    lines = [
        f"рҹ§Ә FaceSwap diagnostic / {PATCH_VERSION}",
        f"FACESWAP_ENABLED={FACESWAP_ENABLED} provider={FACESWAP_PROVIDER} fallback={FACESWAP_FALLBACK_PROVIDER}",
        f"fast_provider={FACESWAP_FAST_PROVIDER} premium_provider={FACESWAP_PREMIUM_PROVIDER}",
        f"target_choice={FACESWAP_ASK_TARGET_FACE} source_choice={FACESWAP_ASK_SOURCE_FACE} strict_selected={FACESWAP_STRICT_SELECTED_FACE}",
        f"manual_choice_if_detection_fail={FACESWAP_MANUAL_CHOICE_IF_DETECTION_FAIL}",
        f"face_detection={FACESWAP_FACE_DETECTION_ENABLED} cv2={_faceswap_cv2_status()} detect_max={FACESWAP_DETECTION_MAX_SIDE} preview_max={FACESWAP_PREVIEW_MAX_SIDE}",
        f"precise_composite={FACESWAP_PRECISE_COMPOSITE} force_segmind_multi={FACESWAP_FORCE_SEGMIND_FOR_MULTI} segmind_fallback={FACESWAP_GROUP_ALLOW_SEGMIND_FALLBACK}",
        f"face_filter_ratio={FACESWAP_FACE_BOX_FILTER_RATIO} source_crop={FACESWAP_SOURCE_CROP_MARGIN} hide_margin={FACESWAP_TARGET_HIDE_MARGIN}",
        f"PIAPI_API_KEY={'on' if PIAPI_API_KEY else 'off'} base={PIAPI_BASE_URL} model={PIAPI_FACE_MODEL} task={PIAPI_FACE_TASK_TYPE}",
        f"SEGMIND_API_KEY={'on' if SEGMIND_API_KEY else 'off'} base={SEGMIND_BASE_URL} fast={SEGMIND_FACESWAP_MODEL_FAST} premium={SEGMIND_FACESWAP_MODEL_PREMIUM}",
        f"timeout={FACESWAP_TIMEOUT_S}s poll={FACESWAP_POLL_DELAY_S}s input_max={FACESWAP_INPUT_MAX_SIDE} output_max={FACESWAP_OUTPUT_MAX_SIDE}",
        "РҹРҫСҒР»РөРҙРҪРёРө РҫСҲРёРұРәРё:",
        _faceswap_last_errors_text(),
    ]
    await update.effective_message.reply_text("\n".join(lines)[:3900])


def face_swap_quality_kb():
    return InlineKeyboardMarkup([
        [InlineKeyboardButton(f"вҡЎ Р‘СӢСҒСӮСҖРҫ В· {_retail_credits(FACESWAP_FAST_COST_USD)} РәСҖ.", callback_data="faceswap:run:fast")],
        [InlineKeyboardButton(f"рҹ’Һ РҹСҖРөРјРёСғРј В· {_retail_credits(FACESWAP_PREMIUM_COST_USD)} РәСҖ.", callback_data="faceswap:run:premium")],
        [InlineKeyboardButton("в¬…пёҸ РқР°Р·Р°Рҙ", callback_data="pedit:back")],
    ])


def _is_face_swap_request(text: str) -> bool:
    tl = (text or "").strip().lower().replace("С‘", "Рө")
    return any(k in tl for k in (
        "Р·Р°РјРөРҪРё Р»РёСҶРҫ", "Р·Р°РјРөРҪРёСӮСҢ Р»РёСҶРҫ", "РҝРҫРјРөРҪСҸР№ Р»РёСҶРҫ", "РҝРҫРјРөРҪСҸСӮСҢ Р»РёСҶРҫ", "РҝРҫРҙСҒСӮР°РІСҢ Р»РёСҶРҫ", "РІСҒСӮР°РІСҢ Р»РёСҶРҫ",
        "face swap", "faceswap", "swap face", "replace face", "СҒРјРөРҪР° Р»РёСҶР°", "Р·Р°РјРөРҪР° Р»РёСҶР°"
    ))


def _set_faceswap_wait_source(context, target_bytes: bytes):
    # Р РөР°Р»СҢРҪСӢР№ user_id РҝСҖРҫРәРёРҙСӢРІР°РөРј РҫСӮРҙРөР»СҢРҪРҫ РІ РІСӢР·РҫРІР°С… СҮРөСҖРөР· _faceswap_target_cache.
    context.user_data["faceswap_flow"] = "await_source"


def _clear_faceswap_flow(context):
    for k in ("faceswap_flow", "awaiting_photo_for"):
        with contextlib.suppress(Exception):
            context.user_data.pop(k, None)


def _clear_faceswap_user_cache(user_id: int):
    for d in (
        _faceswap_target_cache, _faceswap_source_cache,
        _faceswap_target_face_index_cache, _faceswap_source_face_index_cache,
        _faceswap_target_face_count_cache, _faceswap_source_face_count_cache,
        _faceswap_target_faces_cache, _faceswap_source_faces_cache,
    ):
        with contextlib.suppress(Exception):
            d.pop(user_id, None)


def _faceswap_target_for(user_id: int) -> bytes | None:
    return _faceswap_target_cache.get(user_id)


def _faceswap_source_for(user_id: int) -> bytes | None:
    return _faceswap_source_cache.get(user_id)


def _faceswap_selected_target_index(user_id: int) -> int:
    return int(_faceswap_target_face_index_cache.get(user_id, 0) or 0)


def _faceswap_selected_source_index(user_id: int) -> int:
    return int(_faceswap_source_face_index_cache.get(user_id, 0) or 0)


def _resize_image_bytes_for_faceswap_api(img_bytes: bytes, max_side: int | None = None) -> tuple[bytes, str, str]:
    max_side = int(max_side or FACESWAP_INPUT_MAX_SIDE or 1600)
    mime = sniff_image_mime(img_bytes) or "image/jpeg"
    if Image is None or max_side <= 0:
        ext = ".jpg" if mime == "image/jpeg" else (".png" if mime == "image/png" else ".webp")
        return img_bytes, f"faceswap{ext}", mime
    try:
        im = Image.open(BytesIO(img_bytes))
        im = ImageOps.exif_transpose(im) if ImageOps else im
        if max(im.size) > max_side:
            im.thumbnail((max_side, max_side), Image.LANCZOS)
        if im.mode not in ("RGB", "L"):
            im = im.convert("RGB")
        bio = BytesIO()
        im.save(bio, format="JPEG", quality=93, optimize=True, progressive=True)
        return bio.getvalue(), "faceswap.jpg", "image/jpeg"
    except Exception as e:
        _faceswap_note_error(f"faceswap input resize failed: {type(e).__name__}: {e}")
        ext = ".jpg" if mime == "image/jpeg" else (".png" if mime == "image/png" else ".webp")
        return img_bytes, f"faceswap{ext}", mime


def _b64_for_faceswap(img_bytes: bytes) -> str:
    b, _name, mime = _resize_image_bytes_for_faceswap_api(img_bytes, FACESWAP_INPUT_MAX_SIDE)
    raw = base64.b64encode(b).decode("ascii")
    if FACESWAP_IMAGE_DATA_URL:
        return f"data:{mime};base64,{raw}"
    return raw


def _normalize_output_image_bytes(img_bytes: bytes) -> bytes:
    if Image is None or not img_bytes or FACESWAP_OUTPUT_MAX_SIDE <= 0:
        return img_bytes
    try:
        im = Image.open(BytesIO(img_bytes))
        im = ImageOps.exif_transpose(im) if ImageOps else im
        if max(im.size) > FACESWAP_OUTPUT_MAX_SIDE:
            im.thumbnail((FACESWAP_OUTPUT_MAX_SIDE, FACESWAP_OUTPUT_MAX_SIDE), Image.LANCZOS)
        out = BytesIO()
        if im.mode in ("RGBA", "LA"):
            im.save(out, format="PNG", optimize=True)
        else:
            im = im.convert("RGB")
            im.save(out, format="JPEG", quality=94, optimize=True, progressive=True)
        return out.getvalue()
    except Exception:
        return img_bytes


def _maybe_resize_output_image(img_bytes: bytes) -> bytes:
    """Backward-compatible output resize helper used by face swap pipeline."""
    return _normalize_output_image_bytes(img_bytes)


def _detect_faces_for_choice(img_bytes: bytes) -> list[dict]:
    """РһРҝСҖРөРҙРөР»РөРҪРёРө Р»РёСҶ РҙР»СҸ UI Рё РІСӢРұРҫСҖР° РёРҪРҙРөРәСҒРҫРІ.

    Р’Р°Р¶РҪСӢР№ РҝСҖР°РәСӮРёСҮРөСҒРәРёР№ РјРҫРјРөРҪСӮ: РІ СҖРөР°Р»СҢРҪСӢС… СӮРөСҒСӮР°С… РҝСҖРҫРІР°Р№РҙРөСҖСӢ РІРөР»Рё СҒРөРұСҸ СҒСӮР°РұРёР»СҢРҪРөРө,
    РәРҫРіРҙР° РёРҪРҙРөРәСҒСӢ Р»РёСҶ РҝРөСҖРөРҙР°РІР°Р»РёСҒСҢ РІ РІРёР·СғР°Р»СҢРҪРҫРј РҝРҫСҖСҸРҙРәРө СҒР»РөРІР°вҶ’РҪР°РҝСҖР°РІРҫ.
    РҹРҫСҚСӮРҫРјСғ api_index Р·РҙРөСҒСҢ СҒРёРҪС…СҖРҫРҪРёР·РёСҖРҫРІР°РҪ СҒ display_index, Р° РҪРө СҒ СҒРҫСҖСӮРёСҖРҫРІРәРҫР№ РҝРҫ СҖР°Р·РјРөСҖСғ.
    Р”РҫРҝРҫР»РҪРёСӮРөР»СҢРҪРҫ СҖРөР¶РөРј Р»РҫР¶РҪСӢРө СҒСҖР°РұР°СӮСӢРІР°РҪРёСҸ (РјРөР»РәРёРө РұРҫРәСҒСӢ РҪР° РҫРҙРөР¶РҙРө/С„РҫРҪРө).
    """
    if not FACESWAP_FACE_DETECTION_ENABLED or Image is None:
        return []
    try:
        import cv2  # type: ignore
        import numpy as np  # type: ignore

        def _iou(a: dict, b: dict) -> float:
            ax1, ay1, ax2, ay2 = a["x"], a["y"], a["x"] + a["w"], a["y"] + a["h"]
            bx1, by1, bx2, by2 = b["x"], b["y"], b["x"] + b["w"], b["y"] + b["h"]
            ix1, iy1 = max(ax1, bx1), max(ay1, by1)
            ix2, iy2 = min(ax2, bx2), min(ay2, by2)
            iw, ih = max(0, ix2 - ix1), max(0, iy2 - iy1)
            inter = iw * ih
            if inter <= 0:
                return 0.0
            union = a["area"] + b["area"] - inter
            return (inter / union) if union > 0 else 0.0

        im = Image.open(BytesIO(img_bytes))
        im = ImageOps.exif_transpose(im) if ImageOps else im
        im = im.convert("RGB")
        orig_w, orig_h = im.size
        img_area = max(1, orig_w * orig_h)
        scale = 1.0
        max_side = int(FACESWAP_DETECTION_MAX_SIDE or 1200)
        if max_side > 0 and max(orig_w, orig_h) > max_side:
            scale = max_side / float(max(orig_w, orig_h))
            im_det = im.resize((max(1, int(orig_w * scale)), max(1, int(orig_h * scale))), Image.LANCZOS)
        else:
            im_det = im
        arr = np.array(im_det)
        gray = cv2.cvtColor(arr, cv2.COLOR_RGB2GRAY)
        cascade_path = os.path.join(cv2.data.haarcascades, "haarcascade_frontalface_default.xml")
        cascade = cv2.CascadeClassifier(cascade_path)
        if cascade.empty():
            _faceswap_note_error("cv2 haarcascade_frontalface_default.xml not loaded")
            return []
        faces = cascade.detectMultiScale(gray, scaleFactor=1.08, minNeighbors=5, minSize=(40, 40))
        items = []
        for (x, y, w, h) in faces:
            ox = int(round(x / scale)); oy = int(round(y / scale)); ow = int(round(w / scale)); oh = int(round(h / scale))
            if ow < 24 or oh < 24:
                continue
            area = ow * oh
            # СҒР»РёСҲРәРҫРј РјР°Р»РөРҪСҢРәРёРө РұРҫРәСҒСӢ РҝРҫСҮСӮРё РІСҒРөРіРҙР° Р»РҫР¶РҪСӢРө СҒСҖР°РұР°СӮСӢРІР°РҪРёСҸ
            if area < max(1600, int(img_area * 0.004)):
                continue
            ratio = (ow / float(max(1, oh)))
            # СҒР»РёСҲРәРҫРј РІСӢСӮСҸРҪСғСӮСӢРө РұРҫРәСҒСӢ РҙР»СҸ РІСӢРұРҫСҖР° Р»РёСҶР° РҝРҫСҮСӮРё РІСҒРөРіРҙР° СҲСғРј
            if ratio < 0.65 or ratio > 1.45:
                continue
            cx, cy = ox + ow / 2, oy + oh / 2
            items.append({"x": ox, "y": oy, "w": ow, "h": oh, "area": area, "cx": cx, "cy": cy})

        if not items:
            return []

        # NMS / РҙРөРҙСғРҝР»РёРәР°СҶРёСҸ СҒРёР»СҢРҪРҫ РҝРөСҖРөРәСҖСӢРІР°СҺСүРёС…СҒСҸ РұРҫРәСҒРҫРІ
        dedup = []
        for f in sorted(items, key=lambda z: z["area"], reverse=True):
            if any(_iou(f, kept) >= 0.35 for kept in dedup):
                continue
            dedup.append(f)
        items = dedup

        if not items:
            return []

        # v35: Р»РҫР¶РҪСӢРө РұРҫРәСҒСӢ РҪР° Р¶РёР»РөСӮРө/РҫРҙРөР¶РҙРө/СҒСӮРөРҪР°С… СҮР°СҒСӮРҫ РұСӢР»Рё РјРөРҪСҢСҲРө СҖРөР°Р»СҢРҪСӢС… Р»РёСҶ.
        # Р”Р»СҸ РІСӢРұРҫСҖР° Р»РёСҶР° РІ РҝСҖРҫРҙР°РәСҲРҪРө Р»СғСҮСҲРө РҝСҖРҫРҝСғСҒСӮРёСӮСҢ РҫСҮРөРҪСҢ РјРөР»РәРёРө РұРҫРәСҒСӢ, СҮРөРј РҙР°СӮСҢ РҝРҫР»СҢР·РҫРІР°СӮРөР»СҺ
        # РІСӢРұСҖР°СӮСҢ "Р»РёСҶРҫ", РәРҫСӮРҫСҖРҫРө РҝСҖРҫРІР°Р№РҙРөСҖ Р·Р°СӮРөРј РҪРө СҒРјРҫР¶РөСӮ СҒРҫРҝРҫСҒСӮР°РІРёСӮСҢ.
        if len(items) > 1:
            largest_area = max(z["area"] for z in items)
            min_keep = max(int(largest_area * max(0.10, min(0.80, FACESWAP_FACE_BOX_FILTER_RATIO))), int(img_area * 0.006))
            items = [z for z in items if z["area"] >= min_keep]
            if not items:
                return []

        # Р•СҒР»Рё СҸРІРҪРҫ РҙРҫРјРёРҪРёСҖСғРөСӮ РҫРҙРҪРҫ Р»РёСҶРҫ (СӮРёРҝРёСҮРҪСӢР№ СҒРөР»С„Рё-РёСҒСӮРҫСҮРҪРёРә),
        # РҫСӮРұСҖР°СҒСӢРІР°РөРј РјРөР»РәРёРө Р»РҫР¶РҪСӢРө РұРҫРәСҒСӢ РҪР° РҝР»РөСҮРө/С„РҫРҪРө.
        largest = max(items, key=lambda z: z["area"])
        second_area = max([z["area"] for z in items if z is not largest] or [0])
        dominant = largest["area"] >= max(2.0 * second_area, int(img_area * 0.03))
        if dominant:
            filt = []
            for f in items:
                if f is largest:
                    filt.append(f)
                    continue
                if f["area"] >= largest["area"] * 0.45:
                    filt.append(f)
            items = filt

        # РҳСӮРҫРіРҫРІСӢР№ РҝРҫСҖСҸРҙРҫРә: СҒР»РөРІР° РҪР°РҝСҖР°РІРҫ, Р·Р°СӮРөРј СҒРІРөСҖС…Сғ РІРҪРёР·.
        # РўР°РәРҫР№ Р¶Рө РёРҪРҙРөРәСҒ РҝРөСҖРөРҙР°С‘Рј РҝСҖРҫРІР°Р№РҙРөСҖСғ.
        display = sorted(items, key=lambda f: (f["cx"], f["cy"]))[:8]
        for i, f in enumerate(display, 1):
            f["display_index"] = i
            f["api_index"] = i - 1
            if len(display) == 2:
                f["pos_label"] = "СҒР»РөРІР°" if i == 1 else "СҒРҝСҖР°РІР°"
            elif len(display) == 3:
                f["pos_label"] = ["СҒР»РөРІР°", "СҶРөРҪСӮСҖ", "СҒРҝСҖР°РІР°"][i - 1]
            else:
                f["pos_label"] = f"Р»РёСҶРҫ {i}"
        return display
    except Exception as e:
        _faceswap_note_error(f"face detection failed: {type(e).__name__}: {e}")
        return []


def _face_choice_preview_bytes(img_bytes: bytes, faces: list[dict], title: str = "Р’СӢРұРөСҖРёСӮРө Р»РёСҶРҫ") -> bytes | None:
    if Image is None or ImageDraw is None or not faces:
        return None
    try:
        im = Image.open(BytesIO(img_bytes))
        im = ImageOps.exif_transpose(im) if ImageOps else im
        im = im.convert("RGB")
        scale = 1.0
        max_side = int(FACESWAP_PREVIEW_MAX_SIDE or 1200)
        if max_side > 0 and max(im.size) > max_side:
            scale = max_side / float(max(im.size))
            im = im.resize((max(1, int(im.width * scale)), max(1, int(im.height * scale))), Image.LANCZOS)
        draw = ImageDraw.Draw(im)
        try:
            font_big = ImageFont.truetype("DejaVuSans-Bold.ttf", max(24, int(42 * scale)))
            font_small = ImageFont.truetype("DejaVuSans-Bold.ttf", max(14, int(22 * scale)))
        except Exception:
            font_big = None
            font_small = None
        # title strip
        draw.rectangle([0, 0, im.width, min(im.height, 44)], fill=(0, 0, 0))
        draw.text((12, 8), title, fill=(255, 255, 255), font=font_small)
        for f in faces:
            x = int(f["x"] * scale); y = int(f["y"] * scale); w = int(f["w"] * scale); h = int(f["h"] * scale)
            label = str(f.get("display_index") or f.get("api_index") or 0)
            # Р–С‘Р»СӮР°СҸ СҖР°РјРәР° + СҮС‘СҖРҪР°СҸ РҝРҫРҙР»РҫР¶РәР° РҪРҫРјРөСҖР° С…РҫСҖРҫСҲРҫ РІРёРҙРҪСӢ РІ Telegram.
            for off in range(4):
                draw.rectangle([x-off, y-off, x+w+off, y+h+off], outline=(255, 221, 0))
            lx, ly = x, max(44, y - 42)
            draw.ellipse([lx, ly, lx + 42, ly + 42], fill=(255, 221, 0), outline=(0, 0, 0), width=2)
            draw.text((lx + 14, ly + 6), label, fill=(0, 0, 0), font=font_big)
        bio = BytesIO()
        im.save(bio, format="JPEG", quality=92, optimize=True)
        return bio.getvalue()
    except Exception as e:
        _faceswap_note_error(f"face choice preview failed: {type(e).__name__}: {e}")
        return None


def _face_choice_kb(stage: str, faces: list[dict]) -> InlineKeyboardMarkup:
    buttons = []
    row = []
    for f in faces[:8]:
        disp = int(f.get("display_index") or 0)
        api_idx = int(f.get("api_index") or 0)
        pos = str(f.get("pos_label") or "")
        text = f"{disp} вҖ” {pos}" if pos and not pos.startswith("Р»РёСҶРҫ") else f"РӣРёСҶРҫ {disp}"
        row.append(InlineKeyboardButton(text, callback_data=f"faceswap:{stage}:{api_idx}"))
        if len(row) == 2:
            buttons.append(row); row = []
    if row:
        buttons.append(row)
    buttons.append([InlineKeyboardButton("в¬…пёҸ РһСӮРјРөРҪР°", callback_data="pedit:back")])
    return InlineKeyboardMarkup(buttons)


def _manual_face_choice_kb(stage: str) -> InlineKeyboardMarkup:
    # Р СғСҮРҪРҫР№ РІСӢРұРҫСҖ РҪСғР¶РөРҪ, РөСҒР»Рё РҙРөСӮРөРәСӮРҫСҖ Р»РёСҶ РҪРөРҙРҫСҒСӮСғРҝРөРҪ/СҒР»Р°РұСӢР№ РёР»Рё Telegram РҝСҖРёСҒР»Р°Р» СҒР»РҫР¶РҪРҫРө РіСҖСғРҝРҝРҫРІРҫРө С„РҫСӮРҫ.
    # РҳРҪРҙРөРәСҒСӢ РҝРөСҖРөРҙР°СҺСӮСҒСҸ РҝСҖРҫРІР°Р№РҙРөСҖР°Рј РәР°Рә target_face_index/source_face_index.
    return InlineKeyboardMarkup([
        [InlineKeyboardButton("рҹ‘Ҳ РӣРёСҶРҫ СҒР»РөРІР°", callback_data=f"faceswap:{stage}:0")],
        [InlineKeyboardButton("рҹҺҜ РӣРёСҶРҫ РІ СҶРөРҪСӮСҖРө", callback_data=f"faceswap:{stage}:1")],
        [InlineKeyboardButton("рҹ‘ү РӣРёСҶРҫ СҒРҝСҖР°РІР°", callback_data=f"faceswap:{stage}:2")],
        [InlineKeyboardButton("рҹӨ– РҗРІСӮРҫ / РҝРөСҖРІРҫРө Р»РёСҶРҫ", callback_data=f"faceswap:{stage}:0")],
        [InlineKeyboardButton("в¬…пёҸ РһСӮРјРөРҪР°", callback_data="pedit:back")],
    ])


async def _ask_faceswap_source_photo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    context.user_data["faceswap_flow"] = "await_source"
    await update.effective_message.reply_text(
        "рҹҺӯ РўРөРҝРөСҖСҢ РҝСҖРёСҲР»РёСӮРө С„РҫСӮРҫ Р»РёСҶР°, РәРҫСӮРҫСҖРҫРө РҪСғР¶РҪРҫ РІСҒСӮР°РІРёСӮСҢ.\n\n"
        "Р’Р°Р¶РҪРҫ: РёСҒРҝРҫР»СҢР·СғР№СӮРө СӮРҫР»СҢРәРҫ СҒРІРҫРё РёР·РҫРұСҖР°Р¶РөРҪРёСҸ РёР»Рё РёР·РҫРұСҖР°Р¶РөРҪРёСҸ, РҪР° РәРҫСӮРҫСҖСӢРө Сғ РІР°СҒ РөСҒСӮСҢ РҝСҖР°РІРҫ. "
        "РқРөР»СҢР·СҸ РёСҒРҝРҫР»СҢР·РҫРІР°СӮСҢ С„СғРҪРәСҶРёСҺ РҙР»СҸ РҫРұРјР°РҪР°, СҲР°РҪСӮР°Р¶Р°, РҙРҫРәСғРјРөРҪСӮРҫРІ, РёРҪСӮРёРјРҪРҫРіРҫ РәРҫРҪСӮРөРҪСӮР° Рё РІСҖРөРҙРҫРҪРҫСҒРҪСӢС… РҝРҫРҙРҙРөР»РҫРә."
    )


async def _maybe_choose_target_face(update: Update, context: ContextTypes.DEFAULT_TYPE, user_id: int, img: bytes):
    _faceswap_target_cache[user_id] = img
    faces = _detect_faces_for_choice(img)
    _faceswap_target_faces_cache[user_id] = faces
    _faceswap_target_face_count_cache[user_id] = len(faces)
    _faceswap_target_face_index_cache[user_id] = 0
    if FACESWAP_ASK_TARGET_FACE and len(faces) > 1:
        context.user_data["faceswap_flow"] = "choose_target_face"
        preview = _face_choice_preview_bytes(img, faces, "РҡРҫРіРҫ Р·Р°РјРөРҪРёСӮСҢ? Р’СӢРұРөСҖРёСӮРө РҪРҫРјРөСҖ Р»РёСҶР°")
        text = "рҹҺӯ РқР° С„РҫСӮРҫ РҪР°Р№РҙРөРҪРҫ РҪРөСҒРәРҫР»СҢРәРҫ Р»РёСҶ. Р’СӢРұРөСҖРёСӮРө, Сғ РәР°РәРҫРіРҫ СҮРөР»РҫРІРөРәР° Р·Р°РјРөРҪРёСӮСҢ Р»РёСҶРҫ."
        if preview:
            bio = BytesIO(preview); bio.name = "faces_target_choice.jpg"
            await update.effective_message.reply_photo(InputFile(bio), caption=text, reply_markup=_face_choice_kb("target", faces))
        else:
            await update.effective_message.reply_text(text, reply_markup=_face_choice_kb("target", faces))
        return
    if FACESWAP_ASK_TARGET_FACE and not faces and FACESWAP_MANUAL_CHOICE_IF_DETECTION_FAIL:
        # Р•СҒР»Рё OpenCV/РҙРөСӮРөРәСӮРҫСҖ РҪРө РҪР°СҲС‘Р» Р»РёСҶР°, РҪРө РұРөСҖС‘Рј СҖР°РҪРҙРҫРјРҪСӢР№ index 0 РјРҫР»СҮР°.
        # Р”Р°С‘Рј РҝРҫР»СҢР·РҫРІР°СӮРөР»СҺ СҖСғСҮРҪРҫР№ РІСӢРұРҫСҖ: Р»РөРІРҫРө/СҶРөРҪСӮСҖ/РҝСҖР°РІРҫРө/Р°РІСӮРҫ.
        context.user_data["faceswap_flow"] = "choose_target_face"
        _faceswap_target_face_count_cache[user_id] = 3
        await update.effective_message.reply_text(
            "рҹҺӯ РқРө СҒРјРҫРі СғРІРөСҖРөРҪРҪРҫ РҫРҝСҖРөРҙРөР»РёСӮСҢ Р»РёСҶР° РҪР° С„РҫСӮРҫ. Р’СӢРұРөСҖРёСӮРө РІСҖСғСҮРҪСғСҺ, Сғ РәР°РәРҫРіРҫ СҮРөР»РҫРІРөРәР° Р·Р°РјРөРҪРёСӮСҢ Р»РёСҶРҫ:",
            reply_markup=_manual_face_choice_kb("target"),
        )
        return
    await _ask_faceswap_source_photo(update, context)


async def _maybe_choose_source_face(update: Update, context: ContextTypes.DEFAULT_TYPE, user_id: int, img: bytes):
    _faceswap_source_cache[user_id] = img
    faces = _detect_faces_for_choice(img)
    _faceswap_source_faces_cache[user_id] = faces
    _faceswap_source_face_count_cache[user_id] = len(faces)
    _faceswap_source_face_index_cache[user_id] = 0
    if FACESWAP_ASK_SOURCE_FACE and len(faces) > 1:
        context.user_data["faceswap_flow"] = "choose_source_face"
        preview = _face_choice_preview_bytes(img, faces, "Р§СҢС‘ Р»РёСҶРҫ РІР·СҸСӮСҢ? Р’СӢРұРөСҖРёСӮРө РҪРҫРјРөСҖ")
        text = "рҹҺӯ РқР° С„РҫСӮРҫ-РёСҒСӮРҫСҮРҪРёРәРө РҪР°Р№РҙРөРҪРҫ РҪРөСҒРәРҫР»СҢРәРҫ Р»РёСҶ. Р’СӢРұРөСҖРёСӮРө, СҮСҢС‘ Р»РёСҶРҫ РІСҒСӮР°РІРёСӮСҢ."
        if preview:
            bio = BytesIO(preview); bio.name = "faces_source_choice.jpg"
            await update.effective_message.reply_photo(InputFile(bio), caption=text, reply_markup=_face_choice_kb("source", faces))
        else:
            await update.effective_message.reply_text(text, reply_markup=_face_choice_kb("source", faces))
        return
    if FACESWAP_ASK_SOURCE_FACE and not faces and FACESWAP_MANUAL_CHOICE_IF_DETECTION_FAIL:
        context.user_data["faceswap_flow"] = "choose_source_face"
        _faceswap_source_face_count_cache[user_id] = 3
        await update.effective_message.reply_text(
            "рҹҺӯ РқРө СҒРјРҫРі СғРІРөСҖРөРҪРҪРҫ РҫРҝСҖРөРҙРөР»РёСӮСҢ Р»РёСҶРҫ-РёСҒСӮРҫСҮРҪРёРә. Р’СӢРұРөСҖРёСӮРө РІСҖСғСҮРҪСғСҺ, СҮСҢС‘ Р»РёСҶРҫ РІР·СҸСӮСҢ:",
            reply_markup=_manual_face_choice_kb("source"),
        )
        return
    context.user_data["faceswap_flow"] = "ready"
    await update.effective_message.reply_text("рҹҺӯ РӨРҫСӮРҫ Р»РёСҶР° РҝРҫР»СғСҮРөРҪРҫ. Р’СӢРұРөСҖРёСӮРө РәР°СҮРөСҒСӮРІРҫ Р·Р°РјРөРҪСӢ:", reply_markup=face_swap_quality_kb())


async def _extract_image_bytes_from_json_or_response(resp: httpx.Response, client: httpx.AsyncClient) -> bytes | None:
    out = await _image_bytes_from_response(resp, client)
    if out:
        return out
    try:
        obj = resp.json()
    except Exception:
        return None
    url = _find_first_image_url(obj)
    if url:
        rr = await client.get(url, timeout=60.0)
        rr.raise_for_status()
        return bytes(rr.content)
    return None



def _faceswap_image_size(img_bytes: bytes) -> tuple[int, int]:
    try:
        im = Image.open(BytesIO(img_bytes))
        im = ImageOps.exif_transpose(im) if ImageOps else im
        return int(im.width), int(im.height)
    except Exception:
        return (0, 0)


def _faceswap_scaled_faces(faces: list[dict] | None, src_bytes: bytes, dst_bytes: bytes) -> list[dict]:
    """РҹРөСҖРөСҒСҮРёСӮР°СӮСҢ РұРҫРәСҒСӢ Р»РёСҶ РёР· РёСҒС…РҫРҙРҪРҫРіРҫ СҖР°Р·РјРөСҖР° РІ СҖР°Р·РјРөСҖ РәР°СҖСӮРёРҪРәРё, СҖРөР°Р»СҢРҪРҫ РҫСӮРҝСҖР°РІР»СҸРөРјРҫР№ РІ API."""
    if not faces:
        return []
    sw, sh = _faceswap_image_size(src_bytes)
    dw, dh = _faceswap_image_size(dst_bytes)
    if sw <= 0 or sh <= 0 or dw <= 0 or dh <= 0:
        return [dict(f) for f in faces]
    sx, sy = dw / float(sw), dh / float(sh)
    out = []
    for f in faces:
        g = dict(f)
        g["x"] = int(round(float(f.get("x", 0)) * sx))
        g["y"] = int(round(float(f.get("y", 0)) * sy))
        g["w"] = max(1, int(round(float(f.get("w", 1)) * sx)))
        g["h"] = max(1, int(round(float(f.get("h", 1)) * sy)))
        g["cx"] = g["x"] + g["w"] / 2.0
        g["cy"] = g["y"] + g["h"] / 2.0
        g["area"] = g["w"] * g["h"]
        out.append(g)
    return out


def _faceswap_get_face_by_index(faces: list[dict] | None, idx: int) -> dict | None:
    if not faces:
        return None
    idx = int(idx or 0)
    for f in faces:
        if int(f.get("api_index", -999)) == idx or int(f.get("display_index", 0)) - 1 == idx:
            return dict(f)
    if 0 <= idx < len(faces):
        return dict(faces[idx])
    return None


def _faceswap_expand_box(box: dict, width: int, height: int, margin: float = 1.5,
                         margin_x: float | None = None, margin_y_up: float | None = None,
                         margin_y_down: float | None = None) -> tuple[int, int, int, int]:
    x, y, w, h = int(box.get("x", 0)), int(box.get("y", 0)), int(box.get("w", 1)), int(box.get("h", 1))
    cx = x + w / 2.0
    mx = float(margin_x if margin_x is not None else margin)
    myu = float(margin_y_up if margin_y_up is not None else margin)
    myd = float(margin_y_down if margin_y_down is not None else margin)
    x1 = int(round(cx - (w * mx) / 2.0))
    x2 = int(round(cx + (w * mx) / 2.0))
    y1 = int(round(y - h * (myu - 1.0)))
    y2 = int(round(y + h + h * (myd - 1.0)))
    return max(0, x1), max(0, y1), min(width, x2), min(height, y2)


def _faceswap_crop_source_face(img_bytes: bytes, face: dict | None) -> bytes:
    """РһСҒСӮР°РІРёСӮСҢ РІ source СӮРҫР»СҢРәРҫ РІСӢРұСҖР°РҪРҪРҫРө Р»РёСҶРҫ. РӯСӮРҫ СғРұРёСҖР°РөСӮ Р»РҫР¶РҪСӢРө РёРҪРҙРөРәСҒСӢ РҪР° СҒРөР»С„Рё/РҫРҙРөР¶РҙРө."""
    if Image is None or not face:
        return img_bytes
    try:
        im = Image.open(BytesIO(img_bytes))
        im = ImageOps.exif_transpose(im) if ImageOps else im
        im = im.convert("RGB")
        x1, y1, x2, y2 = _faceswap_expand_box(face, im.width, im.height, margin=FACESWAP_SOURCE_CROP_MARGIN)
        if x2 <= x1 or y2 <= y1:
            return img_bytes
        crop = im.crop((x1, y1, x2, y2))
        bio = BytesIO()
        crop.save(bio, format="JPEG", quality=94, optimize=True)
        return bio.getvalue()
    except Exception as e:
        _faceswap_note_error(f"source face crop failed: {type(e).__name__}: {e}")
        return img_bytes


def _faceswap_hide_other_faces(target_bytes: bytes, faces: list[dict] | None, selected_idx: int) -> bytes:
    """РЎРәСҖСӢСӮСҢ РІСҒРө Р»РёСҶР°, РәСҖРҫРјРө РІСӢРұСҖР°РҪРҪРҫРіРҫ, РҝРөСҖРөРҙ РҫСӮРҝСҖР°РІРәРҫР№ РҝСҖРҫРІР°Р№РҙРөСҖСғ.

    РӨРёРҪР°Р»СҢРҪСӢР№ СҖРөР·СғР»СҢСӮР°СӮ РұРөСҖС‘СӮСҒСҸ РҪРө СҶРөР»РёРәРҫРј: РҪРёР¶Рө РјСӢ РІРәР»РөРёРІР°РөРј СӮРҫР»СҢРәРҫ РҫРұР»Р°СҒСӮСҢ РІСӢРұСҖР°РҪРҪРҫРіРҫ Р»РёСҶР°
    РҫРұСҖР°СӮРҪРҫ РІ РёСҒС…РҫРҙРҪСӢР№ РәР°РҙСҖ. РҹРҫСҚСӮРҫРјСғ РІСҖРөРјРөРҪРҪРҫРө СҒРәСҖСӢСӮРёРө СҒРҫСҒРөРҙРҪРёС… Р»РёСҶ РұРөР·РҫРҝР°СҒРҪРҫ Рё СҖРөСҲР°РөСӮ
    РҝСҖРҫРұР»РөРјСғ, РәРҫРіРҙР° API РјРөРҪСҸРөСӮ РҪРө СӮРҫРіРҫ СҮРөР»РҫРІРөРәР°.
    """
    if Image is None or ImageFilter is None or not faces:
        return target_bytes
    try:
        im = Image.open(BytesIO(target_bytes))
        im = ImageOps.exif_transpose(im) if ImageOps else im
        im = im.convert("RGB")
        selected = _faceswap_get_face_by_index(faces, selected_idx)
        if not selected:
            return target_bytes
        blurred = im.filter(ImageFilter.GaussianBlur(radius=max(18, int(max(im.size) * 0.025))))
        for f in faces:
            if int(f.get("api_index", -999)) == int(selected.get("api_index", -888)):
                continue
            x1, y1, x2, y2 = _faceswap_expand_box(f, im.width, im.height, margin=FACESWAP_TARGET_HIDE_MARGIN)
            if x2 <= x1 or y2 <= y1:
                continue
            patch = blurred.crop((x1, y1, x2, y2))
            im.paste(patch, (x1, y1))
        bio = BytesIO()
        im.save(bio, format="JPEG", quality=94, optimize=True)
        return bio.getvalue()
    except Exception as e:
        _faceswap_note_error(f"hide other faces failed: {type(e).__name__}: {e}")
        return target_bytes


def _faceswap_composite_selected_region(original_target_bytes: bytes, provider_output: bytes, selected_face: dict | None) -> bytes:
    """Р’РөСҖРҪСғСӮСҢ РІ РёСҒС…РҫРҙРҪСӢР№ РәР°РҙСҖ СӮРҫР»СҢРәРҫ РёР·РјРөРҪС‘РҪРҪСғСҺ РҫРұР»Р°СҒСӮСҢ РІСӢРұСҖР°РҪРҪРҫРіРҫ Р»РёСҶР°."""
    if Image is None or ImageFilter is None or not selected_face:
        return provider_output
    try:
        base = Image.open(BytesIO(original_target_bytes))
        base = ImageOps.exif_transpose(base) if ImageOps else base
        base = base.convert("RGB")
        out = Image.open(BytesIO(provider_output))
        out = ImageOps.exif_transpose(out) if ImageOps else out
        out = out.convert("RGB")
        if out.size != base.size:
            out = out.resize(base.size, Image.LANCZOS)
        x1, y1, x2, y2 = _faceswap_expand_box(
            selected_face, base.width, base.height,
            margin_x=FACESWAP_COMPOSITE_MARGIN_X,
            margin_y_up=FACESWAP_COMPOSITE_MARGIN_Y_UP,
            margin_y_down=FACESWAP_COMPOSITE_MARGIN_Y_DOWN,
        )
        if x2 <= x1 or y2 <= y1:
            return provider_output
        patch = out.crop((x1, y1, x2, y2))
        mask = Image.new("L", (x2 - x1, y2 - y1), 0)
        md = ImageDraw.Draw(mask)
        pad = max(2, int(min(mask.size) * 0.04))
        md.rounded_rectangle([pad, pad, mask.width - pad, mask.height - pad], radius=max(12, int(min(mask.size) * 0.28)), fill=255)
        mask = mask.filter(ImageFilter.GaussianBlur(radius=max(6, int(min(mask.size) * 0.06))))
        base.paste(patch, (x1, y1), mask)
        bio = BytesIO()
        base.save(bio, format="JPEG", quality=94, optimize=True)
        return bio.getvalue()
    except Exception as e:
        _faceswap_note_error(f"selected face composite failed: {type(e).__name__}: {e}")
        return provider_output


async def _piapi_faceswap(target_img: bytes, source_face: bytes, quality: str = "fast", target_index: int = 0, source_index: int = 0) -> bytes | None:
    if not PIAPI_API_KEY:
        _faceswap_note_error("PIAPI_API_KEY missing")
        return None
    target_b64 = _b64_for_faceswap(target_img)
    source_b64 = _b64_for_faceswap(source_face)
    url = f"{PIAPI_BASE_URL}{PIAPI_FACE_CREATE_PATH}"
    headers = {"x-api-key": PIAPI_API_KEY, "X-API-Key": PIAPI_API_KEY, "Content-Type": "application/json", "Accept": "application/json"}
    payload = {
        "model": PIAPI_FACE_MODEL,
        "task_type": PIAPI_FACE_TASK_TYPE,
        "input": {
            "target_image": target_b64,
            "swap_image": source_b64,
            # Р•СҒР»Рё PiAPI РҪР°СҮРҪС‘СӮ РҝРҫРҙРҙРөСҖР¶РёРІР°СӮСҢ СҸРІРҪСӢР№ РёРҪРҙРөРәСҒ, СҚСӮРё РҝРҫР»СҸ СғР¶Рө РұСғРҙСғСӮ РҝРөСҖРөРҙР°РҪСӢ.
            "target_face_index": int(target_index or 0),
            "source_face_index": int(source_index or 0),
        },
    }
    async with httpx.AsyncClient(timeout=60.0, follow_redirects=True) as client:
        r = await client.post(url, headers=headers, json=payload)
        if r.status_code >= 400:
            _faceswap_note_error(f"PiAPI create failed status={r.status_code} body={r.text[:900]}")
            return None
        try:
            obj = r.json()
        except Exception:
            obj = {}
        direct = await _extract_image_bytes_from_json_or_response(r, client)
        if direct:
            return _normalize_output_image_bytes(direct)
        data = obj.get("data") if isinstance(obj, dict) else None
        task_id = ""
        if isinstance(data, dict):
            task_id = str(data.get("task_id") or data.get("id") or "")
        task_id = task_id or str(obj.get("task_id") or obj.get("id") or "") if isinstance(obj, dict) else ""
        if not task_id:
            _faceswap_note_error(f"PiAPI create returned no task_id body={str(obj)[:900]}")
            return None
        deadline = time.monotonic() + max(30.0, FACESWAP_TIMEOUT_S)
        status_url = f"{PIAPI_BASE_URL}{PIAPI_FACE_STATUS_PATH.format(task_id=task_id)}"
        last_obj = None
        while time.monotonic() < deadline:
            await asyncio.sleep(max(1.0, FACESWAP_POLL_DELAY_S))
            rr = await client.get(status_url, headers=headers)
            if rr.status_code >= 400:
                _faceswap_note_error(f"PiAPI status failed status={rr.status_code} body={rr.text[:900]}")
                return None
            try:
                last_obj = rr.json()
            except Exception:
                last_obj = {}
            img = await _extract_image_bytes_from_json_or_response(rr, client)
            if img:
                return _normalize_output_image_bytes(img)
            data = last_obj.get("data") if isinstance(last_obj, dict) else None
            st = ""
            if isinstance(data, dict):
                st = str(data.get("status") or data.get("state") or "").lower()
            st = st or str(last_obj.get("status") or last_obj.get("state") or "").lower() if isinstance(last_obj, dict) else ""
            if st in ("failed", "error", "canceled", "cancelled"):
                _faceswap_note_error(f"PiAPI task failed status={st} body={str(last_obj)[:900]}")
                return None
        _faceswap_note_error(f"PiAPI task timeout after {FACESWAP_TIMEOUT_S:.0f}s task_id={task_id}")
        return None


async def _segmind_faceswap_v2(target_img: bytes, source_face: bytes, target_index: int = 0, source_index: int = 0) -> bytes | None:
    if not SEGMIND_API_KEY:
        _faceswap_note_error("SEGMIND_API_KEY missing")
        return None
    url = f"{SEGMIND_BASE_URL}/v1/{SEGMIND_FACESWAP_MODEL_FAST}"
    payload = {
        "source_img": _b64_for_faceswap(source_face),
        "target_img": _b64_for_faceswap(target_img),
        "input_faces_index": str(int(target_index or 0)),
        "source_faces_index": str(int(source_index or 0)),
        "face_restore": SEGMIND_FACE_RESTORE,
        "base64": False,
    }
    headers = {"x-api-key": SEGMIND_API_KEY, "Content-Type": "application/json", "Accept": "application/json"}
    async with httpx.AsyncClient(timeout=FACESWAP_TIMEOUT_S, follow_redirects=True) as client:
        r = await client.post(url, headers=headers, json=payload)
        if r.status_code >= 400:
            _faceswap_note_error(f"Segmind v2 failed status={r.status_code} body={r.text[:900]}")
            return None
        img = await _extract_image_bytes_from_json_or_response(r, client)
        if img:
            return _normalize_output_image_bytes(img)
        _faceswap_note_error(f"Segmind v2 no output body={r.text[:900]}")
        return None


async def _segmind_faceswap_v4(target_img: bytes, source_face: bytes, quality: str = "premium") -> bytes | None:
    if not SEGMIND_API_KEY:
        _faceswap_note_error("SEGMIND_API_KEY missing")
        return None
    url = f"{SEGMIND_BASE_URL}/v1/{SEGMIND_FACESWAP_MODEL_PREMIUM}"
    payload = {
        "source_image": _b64_for_faceswap(source_face),
        "target_image": _b64_for_faceswap(target_img),
        "model_type": "quality" if quality == "premium" else "speed",
        "swap_type": SEGMIND_FACE_SWAP_TYPE,
        "style_type": SEGMIND_FACE_STYLE_TYPE,
        "seed": int(time.time()) % 100000000,
        "image_format": "png",
        "image_quality": 95 if quality == "premium" else 90,
        "hardware": "fast",
        "base64": False,
    }
    headers = {"x-api-key": SEGMIND_API_KEY, "Content-Type": "application/json", "Accept": "application/json"}
    async with httpx.AsyncClient(timeout=FACESWAP_TIMEOUT_S, follow_redirects=True) as client:
        r = await client.post(url, headers=headers, json=payload)
        if r.status_code >= 400:
            _faceswap_note_error(f"Segmind v4 failed status={r.status_code} body={r.text[:900]}")
            return None
        img = await _extract_image_bytes_from_json_or_response(r, client)
        if img:
            return _normalize_output_image_bytes(img)
        _faceswap_note_error(f"Segmind v4 no output body={r.text[:900]}")
        return None


def _faceswap_provider_supports_indices(provider: str) -> bool:
    provider = (provider or "").lower().strip()
    # РҹСҖР°РәСӮРёРәР° СӮРөСҒСӮРҫРІ РҝРҫРәР°Р·Р°Р»Р°:
    # - Segmind v2 РҙРөР№СҒСӮРІРёСӮРөР»СҢРҪРҫ СғРјРөРөСӮ Р°РҙСҖРөСҒРҪСӢР№ РІСӢРұРҫСҖ Р»РёСҶ РҝРҫ РёРҪРҙРөРәСҒР°Рј.
    # - PiAPI/Qubico РјРҫР¶РөСӮ СғСҒРҝРөСҲРҪРҫ СҒРІР°РҝР°СӮСҢ Р»РёСҶР°, РҪРҫ РҪРө РҙР°С‘СӮ РҪР°РҙС‘Р¶РҪРҫР№ РіР°СҖР°РҪСӮРёРё,
    #   СҮСӮРҫ target_face_index/source_face_index РұСғРҙСғСӮ СҒРҫРұР»СҺРҙРөРҪСӢ РҙР»СҸ РіСҖСғРҝРҝРҫРІСӢС… С„РҫСӮРҫ.
    # РҹРҫСҚСӮРҫРјСғ РҙР»СҸ СҒСӮСҖРҫРіРҫРіРҫ СҖРөР¶РёРјР° РІСӢРұРҫСҖР° Р»РёСҶР° СҒСҮРёСӮР°РөРј РёРҪРҙРөРәСҒРёСҖСғРөРјСӢРј СӮРҫР»СҢРәРҫ Segmind v2.
    return provider in ("segmind", "segmind-v2", "segmind2")


async def _run_faceswap_provider(provider: str, target_img: bytes, source_face: bytes, quality: str, target_index: int = 0, source_index: int = 0) -> bytes | None:
    provider = (provider or "").lower().strip()
    if provider in ("piapi", "pi", "qubico"):
        return await _piapi_faceswap(target_img, source_face, quality=quality, target_index=target_index, source_index=source_index)
    if provider in ("segmind", "segmind-v2", "segmind2"):
        return await _segmind_faceswap_v2(target_img, source_face, target_index=target_index, source_index=source_index)
    if provider in ("segmind-v4", "segmind4"):
        return await _segmind_faceswap_v4(target_img, source_face, quality=quality)
    _faceswap_note_error(f"unknown faceswap provider={provider}")
    return None


def _faceswap_provider_order(quality: str, user_id: int) -> list[str]:
    if quality == "premium":
        base = [FACESWAP_PREMIUM_PROVIDER, FACESWAP_FALLBACK_PROVIDER, "segmind-v2", FACESWAP_PROVIDER]
    else:
        base = [FACESWAP_FAST_PROVIDER, FACESWAP_PROVIDER, FACESWAP_FALLBACK_PROVIDER, "segmind-v2"]
    seen = []
    order = [p for p in base if p and not (p in seen or seen.append(p))]
    selected_multi = (_faceswap_target_face_count_cache.get(user_id, 0) or 0) > 1 or (_faceswap_source_face_count_cache.get(user_id, 0) or 0) > 1

    # v35: РөСҒР»Рё РҝРҫР»СҢР·РҫРІР°СӮРөР»СҢ РІСӢРұРёСҖР°Р» Р»РёСҶРҫ РҪР° РіСҖСғРҝРҝРҫРІРҫРј С„РҫСӮРҫ, РҪРө РҫСӮРҙР°С‘Рј Р·Р°РҙР°СҮСғ PiAPI.
    # PiAPI С…РҫСҖРҫСҲРёР№ РұСӢСҒСӮСҖСӢР№ РҝСҖРҫРІР°Р№РҙРөСҖ, РҪРҫ РІ СӮРөСҒСӮР°С… РҫРҪ РёРҪРҫРіРҙР° РёРіРҪРҫСҖРёСҖРҫРІР°Р» РІСӢРұСҖР°РҪРҪСӢР№ РёРҪРҙРөРәСҒ.
    # Р”Р»СҸ СӮРҫСҮРҪРҫСҒСӮРё РҫСҒСӮР°РІР»СҸРөРј СӮРҫР»СҢРәРҫ Segmind; v4 РјРҫР¶РөСӮ РёСҒРҝРҫР»СҢР·РҫРІР°СӮСҢСҒСҸ СҮРөСҖРөР· isolate+composite,
    # v2 РҫСҒСӮР°С‘СӮСҒСҸ РёРҪРҙРөРәСҒРёСҖСғРөРјСӢРј СҖРөР·РөСҖРІРҫРј, РөСҒР»Рё СҖР°Р·СҖРөСҲС‘РҪ.
    if FACESWAP_FORCE_SEGMIND_FOR_MULTI and selected_multi:
        seg = ["segmind-v2"]
        if quality == "premium" and FACESWAP_GROUP_ALLOW_SEGMIND_FALLBACK:
            seg.append("segmind-v4")
        elif quality != "premium" and FACESWAP_GROUP_ALLOW_SEGMIND_FALLBACK:
            seg.append("segmind-v4")
        seen2 = []
        return [p for p in seg if p and not (p in seen2 or seen2.append(p))]

    if FACESWAP_STRICT_SELECTED_FACE and selected_multi:
        indexed = [p for p in order if _faceswap_provider_supports_indices(p)]
        if indexed:
            return indexed
        _faceswap_note_error("strict target/source face choice needs Segmind faceswap-v2, but no indexed provider configured")
    return order

async def _faceswap_process(update: Update, context: ContextTypes.DEFAULT_TYPE, quality: str = "fast"):
    user_id = update.effective_user.id
    target = _faceswap_target_for(user_id)
    source = _faceswap_source_for(user_id)
    target_index = _faceswap_selected_target_index(user_id)
    source_index = _faceswap_selected_source_index(user_id)
    if not target or not source:
        await update.effective_message.reply_text("РқСғР¶РҪСӢ 2 С„РҫСӮРҫ: СҒРҪР°СҮР°Р»Р° С„РҫСӮРҫ, РіРҙРө Р·Р°РјРөРҪРёСӮСҢ Р»РёСҶРҫ, Р·Р°СӮРөРј С„РҫСӮРҫ Р»РёСҶР° РҙР»СҸ РІСҒСӮР°РІРәРё.", reply_markup=main_kb)
        return
    if not FACESWAP_ENABLED:
        await update.effective_message.reply_text("рҹҺӯ Р—Р°РјРөРҪР° Р»РёСҶР° РІСҖРөРјРөРҪРҪРҫ РҫСӮРәР»СҺСҮРөРҪР° РІ РҪР°СҒСӮСҖРҫР№РәР°С… СҒРөСҖРІРөСҖР°.")
        return

    original_target = target
    target_api, _target_name, _target_mime = _resize_image_bytes_for_faceswap_api(target, FACESWAP_INPUT_MAX_SIDE)
    target_faces_api = _faceswap_scaled_faces(_faceswap_target_faces_cache.get(user_id), target, target_api)
    source_faces = _faceswap_source_faces_cache.get(user_id) or []
    selected_target_face = _faceswap_get_face_by_index(target_faces_api, target_index)
    selected_source_face = _faceswap_get_face_by_index(source_faces, source_index)

    selected_multi = (_faceswap_target_face_count_cache.get(user_id, 0) or 0) > 1 or (_faceswap_source_face_count_cache.get(user_id, 0) or 0) > 1
    provider_order = _faceswap_provider_order(quality, user_id)
    indexed_available = any(_faceswap_provider_supports_indices(p) for p in provider_order)
    precise = bool(FACESWAP_PRECISE_COMPOSITE and selected_multi and selected_target_face and not indexed_available)
    provider_target = _faceswap_hide_other_faces(target_api, target_faces_api, target_index) if precise else target_api
    provider_source = _faceswap_crop_source_face(source, selected_source_face) if selected_source_face else source
    provider_target_index = 0 if precise else target_index
    provider_source_index = 0 if selected_source_face else source_index

    await update.effective_message.reply_text(
        f"рҹҺӯ Р—Р°РҝСғСҒРәР°СҺ Р·Р°РјРөРҪСғ Р»РёСҶР° ({'РҹСҖРөРјРёСғРј' if quality == 'premium' else 'Р‘СӢСҒСӮСҖРҫ'}). "
        f"РҰРөР»РөРІРҫРө Р»РёСҶРҫ в„–{target_index + 1}. РӣРёСҶРҫ-РёСҒСӮРҫСҮРҪРёРә в„–{source_index + 1}. "
        f"РҹСҖРҫРІР°Р№РҙРөСҖСӢ: {', '.join(provider_order)}. "
        "РЎРҫС…СҖР°РҪСҸСҺ РёСҒС…РҫРҙРҪРҫРө СӮРөР»Рҫ, РҫРҙРөР¶РҙСғ, С„РҫРҪ Рё РәРҫРјРҝРҫР·РёСҶРёСҺ. РһРұСӢСҮРҪРҫ 10вҖ“180 СҒРөРәСғРҪРҙвҖҰ"
    )

    async def _go():
        out = None
        used_provider = ""
        for provider in provider_order:
            try:
                out = await asyncio.wait_for(
                    _run_faceswap_provider(provider, provider_target, provider_source, quality, target_index=provider_target_index, source_index=provider_source_index),
                    timeout=FACESWAP_TIMEOUT_S + 30,
                )
                if out:
                    used_provider = provider
                    break
            except asyncio.TimeoutError:
                _faceswap_note_error(f"{provider} timeout after {FACESWAP_TIMEOUT_S:.0f}s")
            except Exception as e:
                _faceswap_note_error(f"{provider} exception: {type(e).__name__}: {e}")
        if not out:
            await update.effective_message.reply_text("вқҢ РқРө СғРҙР°Р»РҫСҒСҢ Р·Р°РјРөРҪРёСӮСҢ Р»РёСҶРҫ. РҹРҫСҒР»РөРҙРҪРёРө РҫСҲРёРұРәРё:\n" + _faceswap_last_errors_text())
            return False
        if precise:
            # РЎРҫС…СҖР°РҪСҸРөРј РҫСҒСӮР°Р»СҢРҪСӢРө Р»РёСҶР° Рё С„РҫРҪ РёР· РёСҒС…РҫРҙРҪРҫРіРҫ РәР°РҙСҖР°; РјРөРҪСҸРөРј СӮРҫР»СҢРәРҫ РІСӢРұСҖР°РҪРҪСғСҺ РҫРұР»Р°СҒСӮСҢ.
            out = _faceswap_composite_selected_region(target_api, out, selected_target_face)
        out = _maybe_resize_output_image(out)
        _cache_photo(user_id, out)
        bio = BytesIO(out)
        cap = (
            f"рҹҺӯ РӣРёСҶРҫ Р·Р°РјРөРҪРөРҪРҫ вң… Р РөР¶РёРј: {'РҹСҖРөРјРёСғРј' if quality == 'premium' else 'Р‘СӢСҒСӮСҖРҫ'} В· РҝСҖРҫРІР°Р№РҙРөСҖ: {used_provider}. "
            f"РҰРөР»РөРІРҫРө Р»РёСҶРҫ в„–{target_index + 1}, РёСҒСӮРҫСҮРҪРёРә в„–{source_index + 1}. "
            "Р РөР·СғР»СҢСӮР°СӮ СҒРҫС…СҖР°РҪС‘РҪ РәР°Рә РёСҒС…РҫРҙРҪРҫРө С„РҫСӮРҫ РҙР»СҸ РҙР°Р»СҢРҪРөР№СҲРёС… РҙРөР№СҒСӮРІРёР№."
        )
        if FACESWAP_RESULT_AS_DOCUMENT:
            bio.name = "face_swap.png"
            await update.effective_message.reply_document(InputFile(bio), caption=cap, reply_markup=photo_quick_actions_kb())
        else:
            bio.name = "face_swap.jpg"
            await update.effective_message.reply_photo(InputFile(bio), caption=cap, reply_markup=photo_quick_actions_kb())
        _clear_faceswap_flow(context)
        _clear_faceswap_user_cache(user_id)
        return True

    est = FACESWAP_PREMIUM_COST_USD if quality == "premium" else FACESWAP_FAST_COST_USD
    await _try_pay_then_do(update, context, user_id, "img", est, _go, remember_kind=f"faceswap_{quality}")

async def _start_faceswap_flow(update: Update, context: ContextTypes.DEFAULT_TYPE, img_bytes: bytes | None = None, use_cached: bool = True):
    user_id = update.effective_user.id
    target = img_bytes or (_get_cached_photo(user_id) if use_cached else None)
    if not target:
        _clear_faceswap_user_cache(user_id)
        context.user_data["faceswap_flow"] = "await_target"
        await update.effective_message.reply_text(
            "рҹҺӯ Р—Р°РјРөРҪР° Р»РёСҶР°. РҹСҖРёСҲР»РёСӮРө РқРһР’РһР• С„РҫСӮРҫ, РіРҙРө РҪСғР¶РҪРҫ Р·Р°РјРөРҪРёСӮСҢ Р»РёСҶРҫ. Р•СҒР»Рё РҪР° С„РҫСӮРҫ РҪРөСҒРәРҫР»СҢРәРҫ Р»СҺРҙРөР№, СҸ РҝРҫРәР°Р¶Сғ РҪРҫРјРөСҖР° Рё РҝРҫРҝСҖРҫСҲСғ РІСӢРұСҖР°СӮСҢ РҪСғР¶РҪРҫРіРҫ СҮРөР»РҫРІРөРәР°. Р—Р°СӮРөРј СҸ РҝРҫРҝСҖРҫСҲСғ С„РҫСӮРҫ Р»РёСҶР° РҙР»СҸ РІСҒСӮР°РІРәРё."
        )
        return
    await _maybe_choose_target_face(update, context, user_id, target)

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ WebApp data (СӮР°СҖРёС„СӢ/РҝРҫРҝРҫР»РҪРөРҪРёСҸ) в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
async def on_webapp_data(update: Update, context: ContextTypes.DEFAULT_TYPE):
    try:
        wad = update.effective_message.web_app_data
        raw = wad.data if wad else ""
        data = {}
        try:
            data = json.loads(raw)
        except Exception:
            for part in (raw or "").split("&"):
                if "=" in part:
                    k, v = part.split("=", 1)
                    data[k] = v

        typ = (data.get("type") or data.get("action") or "").lower()
        immediate = str(data.get("immediate") or "").lower() in ("1", "true", "yes", "on")
        user_id = update.effective_user.id

        if typ in ("subscribe", "buy", "buy_sub", "sub"):
            tier = (data.get("tier") or "pro").lower()
            if tier not in SUBS_TIERS:
                await update.effective_message.reply_text("РқРөРёР·РІРөСҒСӮРҪСӢР№ СӮР°СҖРёС„. РһСӮРәСҖРҫР№СӮРө СҒСӮСҖР°РҪРёСҶСғ СӮР°СҖРёС„РҫРІ Р·Р°РҪРҫРІРҫ.")
                return
            months = max(1, min(12, int(data.get("months") or 1)))
            method = (data.get("method") or "yoo_all").lower()
            if method not in YOO_DIRECT_METHODS:
                method = "yoo_all"

            if immediate and _yoo_direct_configured():
                try:
                    pay = await _yoo_create_direct_payment(user_id, tier, months, method)
                    payment_id = str(pay.get("id") or "")
                    conf = pay.get("confirmation") or {}
                    pay_url = conf.get("confirmation_url") or conf.get("confirmation_data") or conf.get("external_url") or ""
                    if not pay_url:
                        raise RuntimeError("YooKassa did not return confirmation URL")
                    label = YOO_DIRECT_METHODS[method]["label"]
                    msg = await update.effective_message.reply_text(
                        f"вӯҗ РўР°СҖРёС„ {tier.upper()} РҪР° {months} РјРөСҒ.\n"
                        f"РЎРҝРҫСҒРҫРұ: {label}. РқР°Р¶РјРёСӮРө РәРҪРҫРҝРәСғ РҙР»СҸ РҫРҝР»Р°СӮСӢ; РҝРҫСҒР»Рө РҝРҫРҙСӮРІРөСҖР¶РҙРөРҪРёСҸ РҝРҫРҙРҝРёСҒРәР° Р°РәСӮРёРІРёСҖСғРөСӮСҒСҸ Р°РІСӮРҫРјР°СӮРёСҮРөСҒРәРё.",
                        reply_markup=InlineKeyboardMarkup([
                            [InlineKeyboardButton(f"{label} вҖ” РҝРөСҖРөР№СӮРё Рә РҫРҝР»Р°СӮРө", url=pay_url)],
                            [InlineKeyboardButton("Р”СҖСғРіРҫР№ СҒРҝРҫСҒРҫРұ РҫРҝР»Р°СӮСӢ", callback_data=f"plan:{tier}")],
                        ]),
                    )
                    kv_set(
                        f"yoo:pending:{payment_id}",
                        json.dumps({"user_id": user_id, "tier": tier, "months": months, "method": method}, ensure_ascii=False),
                    )
                    context.application.create_task(
                        _poll_yoo_subscription_payment(context, msg.chat.id, msg.message_id, user_id, payment_id, tier, months)
                    )
                    return
                except Exception as exc:
                    log.exception("WebApp immediate subscription payment failed: %s", exc)

            await update.effective_message.reply_text(
                f"РһС„РҫСҖРјР»РөРҪРёРө РҝРҫРҙРҝРёСҒРәРё {tier.upper()} РҪР° {months} РјРөСҒ.\nР’СӢРұРөСҖРёСӮРө СҒРҝРҫСҒРҫРұ РҫРҝР»Р°СӮСӢ:",
                reply_markup=plan_pay_kb(tier),
            )
            return

        if typ in ("topup_rub", "rub_topup", "buy_credits", "credit_pack"):
            requested_rub = int(data.get("amount") or data.get("rub") or 0)
            requested_credits = int(data.get("credits") or 0)
            resolved = _credit_pack_resolve(requested_credits, requested_rub)
            if not resolved:
                await update.effective_message.reply_text("РқРөРёР·РІРөСҒСӮРҪСӢР№ РҝР°РәРөСӮ РәСҖРөРҙРёСӮРҫРІ. РһСӮРәСҖРҫР№СӮРө СҒСӮСҖР°РҪРёСҶСғ СӮР°СҖРёС„РҫРІ Р·Р°РҪРҫРІРҫ.")
                return
            credits, amount_rub = resolved
            method = (data.get("method") or "yoo_all").lower()
            if method not in YOO_DIRECT_METHODS:
                method = "yoo_all"

            if immediate and _yoo_direct_configured():
                try:
                    pay = await _yoo_create_credit_payment(user_id, credits, amount_rub, method)
                    payment_id = str(pay.get("id") or "")
                    conf = pay.get("confirmation") or {}
                    pay_url = conf.get("confirmation_url") or conf.get("confirmation_data") or conf.get("external_url") or ""
                    if not pay_url:
                        raise RuntimeError("YooKassa did not return confirmation URL")
                    label = YOO_DIRECT_METHODS[method]["label"]
                    msg = await update.effective_message.reply_text(
                        f"рҹӘҷ РҹР°РәРөСӮ: {credits} РәСҖРөРҙРёСӮРҫРІ Р·Р° {amount_rub} вӮҪ.\n"
                        f"РЎРҝРҫСҒРҫРұ: {label}. РҹРҫСҒР»Рө РҫРҝР»Р°СӮСӢ РәСҖРөРҙРёСӮСӢ РҪР°СҮРёСҒР»СҸСӮСҒСҸ Р°РІСӮРҫРјР°СӮРёСҮРөСҒРәРё.",
                        reply_markup=InlineKeyboardMarkup([
                            [InlineKeyboardButton(f"{label} вҖ” РҝРөСҖРөР№СӮРё Рә РҫРҝР»Р°СӮРө", url=pay_url)],
                            [InlineKeyboardButton("Р”СҖСғРіРёРө РҝР°РәРөСӮСӢ", callback_data="topup")],
                        ]),
                    )
                    kv_set(
                        f"yoo:credit_pending:{payment_id}",
                        json.dumps({"user_id": user_id, "credits": credits, "amount_rub": amount_rub, "method": method}, ensure_ascii=False),
                    )
                    context.application.create_task(
                        _poll_yoo_credit_payment(context, msg.chat.id, msg.message_id, user_id, payment_id, credits, amount_rub)
                    )
                    return
                except Exception as exc:
                    log.exception("WebApp immediate credit payment failed: %s", exc)

            await _send_invoice_rub(
                f"{credits} РәСҖРөРҙРёСӮРҫРІ",
                f"РҹРҫРҝРҫР»РҪРөРҪРёРө РұР°Р»Р°РҪСҒР° Neyro-Bot: {credits} РәСҖРөРҙРёСӮРҫРІ.",
                amount_rub,
                f"topup:{credits}:{amount_rub}",
                update,
            )
            return

        if typ in ("topup_crypto", "crypto_topup"):
            if not CRYPTO_PAY_API_TOKEN:
                await update.effective_message.reply_text("CryptoBot РҪРө РҪР°СҒСӮСҖРҫРөРҪ.")
                return
            usd = float(data.get("usd") or 0)
            inv_id, pay_url, usd_amount, asset = await _crypto_create_invoice(usd, asset="USDT")
            if not inv_id or not pay_url:
                await update.effective_message.reply_text("РқРө СғРҙР°Р»РҫСҒСҢ СҒРҫР·РҙР°СӮСҢ СҒСҮС‘СӮ РІ CryptoBot.")
                return
            msg = await update.effective_message.reply_text(
                f"РһРҝР»Р°СӮРёСӮРө СҮРөСҖРөР· CryptoBot: {usd_amount:.2f} {asset} вҶ’ {_credits_fmt_from_usd(usd_amount)}.",
                reply_markup=InlineKeyboardMarkup([
                    [InlineKeyboardButton("рҹ’  CryptoBot", url=pay_url)],
                    [InlineKeyboardButton("рҹ”Һ РҹСҖРҫРІРөСҖРёСӮСҢ", callback_data=f"crypto:check:{inv_id}")],
                ]),
            )
            context.application.create_task(
                _poll_crypto_invoice(context, msg.chat_id, msg.message_id, user_id, inv_id, usd_amount)
            )
            return

        await update.effective_message.reply_text("РҹРҫР»СғСҮРөРҪСӢ РҙР°РҪРҪСӢРө РёР· РјРёРҪРё-РҝСҖРёР»РҫР¶РөРҪРёСҸ, РҪРҫ РәРҫРјР°РҪРҙР° РҪРө СҖР°СҒРҝРҫР·РҪР°РҪР°.")
    except Exception as exc:
        log.exception("on_webapp_data error: %s", exc)
        await update.effective_message.reply_text("РһСҲРёРұРәР° РҫРұСҖР°РұРҫСӮРәРё РҙР°РҪРҪСӢС… РјРёРҪРё-РҝСҖРёР»РҫР¶РөРҪРёСҸ.")


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ v93: server-side checkout bridge for Inline/Menu Mini Apps в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def _validate_telegram_webapp_init_data(init_data: str, max_age_s: int = 86400) -> dict:
    """Validate Telegram Mini App initData and return parsed fields.

    Inline/Menu Mini Apps receive a query_id and must communicate through a
    server-side bridge. Never trust user id or purchase parameters from JS
    without validating initData using the bot token.
    """
    raw = (init_data or "").strip()
    if not raw:
        raise ValueError("Telegram initData is empty")
    pairs = urllib.parse.parse_qsl(raw, keep_blank_values=True)
    data = {str(k): str(v) for k, v in pairs}
    received_hash = data.pop("hash", "")
    # For bot-token HMAC validation, all received fields except hash remain
    # in the data-check string, including the optional signature field.
    if not received_hash:
        raise ValueError("Telegram initData hash is missing")
    data_check_string = "\n".join(f"{k}={data[k]}" for k in sorted(data))
    secret_key = hmac.new(b"WebAppData", BOT_TOKEN.encode("utf-8"), hashlib.sha256).digest()
    calculated_hash = hmac.new(secret_key, data_check_string.encode("utf-8"), hashlib.sha256).hexdigest()
    if not hmac.compare_digest(calculated_hash, received_hash):
        raise ValueError("Telegram initData signature is invalid")
    auth_date = int(data.get("auth_date") or 0)
    now_ts = int(time.time())
    if not auth_date or auth_date > now_ts + 60 or now_ts - auth_date > max_age_s:
        raise ValueError("Telegram initData has expired")
    user_raw = data.get("user") or "{}"
    try:
        user = json.loads(user_raw)
    except Exception as exc:
        raise ValueError("Telegram initData user is invalid") from exc
    user_id = int(user.get("id") or 0)
    if user_id <= 0:
        raise ValueError("Telegram user id is missing")
    data["user_obj"] = user
    data["user_id"] = user_id
    return data


def _install_webapp_checkout_bridge(application):
    """Extend PTB's Tornado webhook app with /webapp/checkout.

    PTB 21.6 does not expose custom routes in run_webhook, so this replaces only
    the small internal WebhookAppClass while preserving TelegramHandler.
    """
    try:
        import tornado.web
        import telegram.ext._updater as ptb_updater
        from telegram.ext._utils.webhookhandler import TelegramHandler
    except Exception as exc:
        log.exception("Checkout bridge dependencies unavailable: %s", exc)
        return False

    class CheckoutBridgeHandler(tornado.web.RequestHandler):
        def set_default_headers(self):
            self.set_header("Content-Type", "application/json; charset=utf-8")
            self.set_header("Access-Control-Allow-Origin", "*")
            self.set_header("Access-Control-Allow-Headers", "Content-Type")
            self.set_header("Access-Control-Allow-Methods", "POST, OPTIONS")
            self.set_header("Cache-Control", "no-store")

        async def options(self):
            self.set_status(204)
            self.finish()

        async def post(self):
            try:
                if not _yoo_direct_configured():
                    self.set_status(503)
                    self.finish(json.dumps({"ok": False, "error": "YooKassa direct API is not configured"}, ensure_ascii=False))
                    return
                try:
                    body = json.loads(self.request.body.decode("utf-8") or "{}")
                except Exception:
                    self.set_status(400)
                    self.finish(json.dumps({"ok": False, "error": "Invalid JSON"}, ensure_ascii=False))
                    return
                init_data = str(body.get("init_data") or "")
                purchase = body.get("purchase") or {}
                if not isinstance(purchase, dict):
                    raise ValueError("Invalid purchase payload")
                validated = _validate_telegram_webapp_init_data(init_data)
                user_id = int(validated["user_id"])
                typ = str(purchase.get("type") or purchase.get("action") or "").lower()
                method = str(purchase.get("method") or "yoo_all").lower()
                if method not in YOO_DIRECT_METHODS:
                    method = "yoo_all"

                if typ in ("subscribe", "buy", "buy_sub", "sub"):
                    tier = str(purchase.get("tier") or "pro").lower()
                    if tier not in SUBS_TIERS:
                        raise ValueError("Unknown subscription tier")
                    months = max(1, min(12, int(purchase.get("months") or 1)))
                    pay = await _yoo_create_direct_payment(user_id, tier, months, method)
                    payment_id = str(pay.get("id") or "")
                    conf = pay.get("confirmation") or {}
                    pay_url = str(conf.get("confirmation_url") or conf.get("confirmation_data") or conf.get("external_url") or "")
                    if not payment_id or not pay_url:
                        raise RuntimeError("YooKassa did not return confirmation URL")
                    label = YOO_DIRECT_METHODS[method]["label"]
                    msg = await application.bot.send_message(
                        chat_id=user_id,
                        text=(
                            f"вӯҗ РўР°СҖРёС„ {tier.upper()} РҪР° {months} РјРөСҒ.\n"
                            f"РЎРҝРҫСҒРҫРұ: {label}. РқР°Р¶РјРёСӮРө РәРҪРҫРҝРәСғ РҙР»СҸ РҫРҝР»Р°СӮСӢ; РҝРҫСҒР»Рө РҝРҫРҙСӮРІРөСҖР¶РҙРөРҪРёСҸ РҝРҫРҙРҝРёСҒРәР° Р°РәСӮРёРІРёСҖСғРөСӮСҒСҸ Р°РІСӮРҫРјР°СӮРёСҮРөСҒРәРё."
                        ),
                        reply_markup=InlineKeyboardMarkup([
                            [InlineKeyboardButton(f"{label} вҖ” РҝРөСҖРөР№СӮРё Рә РҫРҝР»Р°СӮРө", url=pay_url)],
                            [InlineKeyboardButton("Р”СҖСғРіРҫР№ СҒРҝРҫСҒРҫРұ РҫРҝР»Р°СӮСӢ", callback_data=f"plan:{tier}")],
                        ]),
                    )
                    kv_set(
                        f"yoo:pending:{payment_id}",
                        json.dumps({"user_id": user_id, "tier": tier, "months": months, "method": method}, ensure_ascii=False),
                    )
                    simple_context = types.SimpleNamespace(bot=application.bot)
                    application.create_task(
                        _poll_yoo_subscription_payment(simple_context, msg.chat.id, msg.message_id, user_id, payment_id, tier, months)
                    )
                    result = {"ok": True, "url": pay_url, "payment_id": payment_id, "kind": "subscription"}

                elif typ in ("topup_rub", "rub_topup", "buy_credits", "credit_pack"):
                    requested_rub = int(purchase.get("amount") or purchase.get("rub") or 0)
                    requested_credits = int(purchase.get("credits") or 0)
                    resolved = _credit_pack_resolve(requested_credits, requested_rub)
                    if not resolved:
                        raise ValueError("Unknown credit package")
                    credits, amount_rub = resolved
                    pay = await _yoo_create_credit_payment(user_id, credits, amount_rub, method)
                    payment_id = str(pay.get("id") or "")
                    conf = pay.get("confirmation") or {}
                    pay_url = str(conf.get("confirmation_url") or conf.get("confirmation_data") or conf.get("external_url") or "")
                    if not payment_id or not pay_url:
                        raise RuntimeError("YooKassa did not return confirmation URL")
                    label = YOO_DIRECT_METHODS[method]["label"]
                    msg = await application.bot.send_message(
                        chat_id=user_id,
                        text=(
                            f"рҹӘҷ РҹР°РәРөСӮ: {credits} РәСҖРөРҙРёСӮРҫРІ Р·Р° {amount_rub} вӮҪ.\n"
                            f"РЎРҝРҫСҒРҫРұ: {label}. РҹРҫСҒР»Рө РҫРҝР»Р°СӮСӢ РәСҖРөРҙРёСӮСӢ РҪР°СҮРёСҒР»СҸСӮСҒСҸ Р°РІСӮРҫРјР°СӮРёСҮРөСҒРәРё."
                        ),
                        reply_markup=InlineKeyboardMarkup([
                            [InlineKeyboardButton(f"{label} вҖ” РҝРөСҖРөР№СӮРё Рә РҫРҝР»Р°СӮРө", url=pay_url)],
                            [InlineKeyboardButton("Р”СҖСғРіРёРө РҝР°РәРөСӮСӢ", callback_data="topup")],
                        ]),
                    )
                    kv_set(
                        f"yoo:credit_pending:{payment_id}",
                        json.dumps({"user_id": user_id, "credits": credits, "amount_rub": amount_rub, "method": method}, ensure_ascii=False),
                    )
                    simple_context = types.SimpleNamespace(bot=application.bot)
                    application.create_task(
                        _poll_yoo_credit_payment(simple_context, msg.chat.id, msg.message_id, user_id, payment_id, credits, amount_rub)
                    )
                    result = {"ok": True, "url": pay_url, "payment_id": payment_id, "kind": "credits"}
                else:
                    raise ValueError("Unknown purchase type")

                self.set_status(200)
                self.finish(json.dumps(result, ensure_ascii=False))
            except ValueError as exc:
                log.warning("WebApp checkout rejected: %s", exc)
                self.set_status(400)
                self.finish(json.dumps({"ok": False, "error": str(exc)}, ensure_ascii=False))
            except Exception as exc:
                log.exception("WebApp checkout bridge failed: %s", exc)
                self.set_status(500)
                message = str(exc)[:500] if YOO_DEBUG_PAY_ERRORS else "РқРө СғРҙР°Р»РҫСҒСҢ СҒРҫР·РҙР°СӮСҢ РҫРҝР»Р°СӮСғ"
                self.finish(json.dumps({"ok": False, "error": message}, ensure_ascii=False))

    class CheckoutWebhookApp(tornado.web.Application):
        def __init__(self, webhook_path, bot, update_queue, secret_token=None):
            shared = {"bot": bot, "update_queue": update_queue, "secret_token": secret_token}
            handlers = [
                (rf"{webhook_path}/?", TelegramHandler, shared),
                (r"/webapp/checkout/?", CheckoutBridgeHandler),
                (r"/healthz/?", HealthBridgeHandler),
            ]
            super().__init__(handlers)

        def log_request(self, handler):
            status = handler.get_status()
            if status >= 400:
                log.warning("HTTP %s %s %s", status, handler.request.method, handler.request.uri)

    class HealthBridgeHandler(tornado.web.RequestHandler):
        async def get(self):
            self.set_header("Content-Type", "application/json; charset=utf-8")
            self.finish(json.dumps({"ok": True, "version": PATCH_VERSION}, ensure_ascii=False))

    ptb_updater.WebhookAppClass = CheckoutWebhookApp
    log.info("WebApp checkout bridge installed: %s/webapp/checkout", PUBLIC_URL.rstrip("/"))
    return True


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ CallbackQuery (РІСҒС‘ РҫСҒСӮР°Р»СҢРҪРҫРө) в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
_pending_actions = {}

def _new_aid() -> str:
    return uuid.uuid4().hex[:12]

async def on_cb(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    data = (q.data or "").strip()
    try:
        # Presentation/Catalog Studio v86
        if data.startswith("ps:"):
            await _presentation_studio_get().handle_callback(update, context)
            return

        # Persistent virtual chats
        if data == "chat:list":
            await q.answer()
            await cmd_chats(update, context)
            return
        if data == "chat:new":
            await q.answer()
            await cmd_newchat(update, context)
            return
        if data.startswith("chat:open:"):
            await q.answer()
            try: cid = int(data.split(":", 2)[2])
            except Exception: cid = 0
            if not cid or not _chat_set_active(q.from_user.id, q.message.chat_id, cid):
                await q.message.reply_text("Р§Р°СӮ РҪРө РҪР°Р№РҙРөРҪ.")
                return
            _clear_transient_flows(context)
            title = next((x["title"] for x in _chat_list(q.from_user.id, q.message.chat_id) if x["id"] == cid), "Р§Р°СӮ")
            await q.message.reply_text(f"в–¶пёҸ Р§Р°СӮ В«{title}В» РІСӢРұСҖР°РҪ. РҹСҖРҫРҙРҫР»Р¶Р°Р№СӮРө СҖР°Р·РіРҫРІРҫСҖ.", reply_markup=main_kb)
            return
        if data.startswith("chat:history:"):
            await q.answer()
            parts = data.split(":")
            try: cid = int(parts[2]); page = int(parts[3]) if len(parts) > 3 else 0
            except Exception: cid, page = 0, 0
            if cid:
                await _send_chat_history(update, context, cid, page)
            return
        if data.startswith("chat:rename:"):
            await q.answer()
            try: cid = int(data.split(":", 2)[2])
            except Exception: cid = 0
            if cid:
                context.user_data["awaiting_chat_rename"] = cid
                await q.message.reply_text("вңҸпёҸ РһСӮРҝСҖР°РІСҢСӮРө РҪРҫРІРҫРө РҪР°Р·РІР°РҪРёРө СҮР°СӮР° РҫРҙРҪРёРј СҒРҫРҫРұСүРөРҪРёРөРј (РҙРҫ 60 СҒРёРјРІРҫР»РҫРІ).")
            return
        if data.startswith("chat:delete_confirm:"):
            await q.answer()
            try: cid = int(data.split(":", 2)[2])
            except Exception: cid = 0
            if cid and _chat_delete(q.from_user.id, q.message.chat_id, cid):
                await q.message.reply_text("рҹ—‘ Р§Р°СӮ Рё РөРіРҫ РёСҒСӮРҫСҖРёСҸ СғРҙР°Р»РөРҪСӢ.", reply_markup=_chat_list_kb(q.from_user.id, q.message.chat_id))
            else:
                await q.message.reply_text("Р§Р°СӮ РҪРө РҪР°Р№РҙРөРҪ.")
            return
        if data.startswith("chat:delete:"):
            await q.answer()
            try: cid = int(data.split(":", 2)[2])
            except Exception: cid = 0
            await q.message.reply_text(
                "РЈРҙР°Р»РёСӮСҢ СҚСӮРҫСӮ СҮР°СӮ РІРјРөСҒСӮРө СҒРҫ РІСҒРөР№ РёСҒСӮРҫСҖРёРөР№?",
                reply_markup=InlineKeyboardMarkup([
                    [InlineKeyboardButton("Р”Р°, СғРҙР°Р»РёСӮСҢ", callback_data=f"chat:delete_confirm:{cid}")],
                    [InlineKeyboardButton("РһСӮРјРөРҪР°", callback_data="chat:list")],
                ]),
            )
            return

        if data == "pricing:list":
            await q.answer()
            await q.message.reply_text(
                _pricing_catalog_text(),
                reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("вһ• РҹРҫРҝРҫР»РҪРёСӮСҢ РұР°Р»Р°РҪСҒ", callback_data="topup")]]),
            )
            return

        # TOPUP РјРөРҪСҺ
        if data == "topup":
            await q.answer()
            await _send_topup_menu(update, context)
            return

        # TOPUP RUB
        if data.startswith("topup:rub:"):
            await q.answer()
            try:
                amount_rub = int((data.split(":", 2)[-1] or "0").strip() or "0")
            except Exception:
                amount_rub = 0
            resolved = _credit_pack_resolve(0, amount_rub)
            if not resolved:
                await q.message.reply_text("РқРөРёР·РІРөСҒСӮРҪСӢР№ РҝР°РәРөСӮ РәСҖРөРҙРёСӮРҫРІ. РһСӮРәСҖРҫР№СӮРө РјРөРҪСҺ РҝРҫРҝРҫР»РҪРөРҪРёСҸ Р·Р°РҪРҫРІРҫ.")
                return
            credits, amount_rub = resolved
            ok = await _send_invoice_rub(
                f"{credits} РәСҖРөРҙРёСӮРҫРІ",
                f"РҹРҫРҝРҫР»РҪРөРҪРёРө РұР°Р»Р°РҪСҒР° Neyro-Bot: {credits} РәСҖРөРҙРёСӮРҫРІ.",
                amount_rub,
                f"topup:{credits}:{amount_rub}",
                update,
            )
            await q.answer("Р’СӢСҒСӮР°РІР»СҸСҺ СҒСҮС‘СӮвҖҰ" if ok else "РқРө СғРҙР°Р»РҫСҒСҢ РІСӢСҒСӮР°РІРёСӮСҢ СҒСҮС‘СӮ", show_alert=not ok)
            return

        # TOPUP CRYPTO
        if data.startswith("topup:crypto:"):
            await q.answer()
            if not CRYPTO_PAY_API_TOKEN:
                await q.message.reply_text("РқР°СҒСӮСҖРҫР№СӮРө CRYPTO_PAY_API_TOKEN РҙР»СҸ РҫРҝР»Р°СӮСӢ СҮРөСҖРөР· CryptoBot.")
                return
            try:
                usd = float((data.split(":", 2)[-1] or "0").strip() or "0")
            except Exception:
                usd = 0.0
            if usd <= 0.0:
                await q.message.reply_text("РқРөРІРөСҖРҪР°СҸ СҒСғРјРјР°.")
                return
            inv_id, pay_url, usd_amount, asset = await _crypto_create_invoice(usd, asset="USDT", description="Wallet top-up")
            if not inv_id or not pay_url:
                await q.message.reply_text("РқРө СғРҙР°Р»РҫСҒСҢ СҒРҫР·РҙР°СӮСҢ СҒСҮС‘СӮ РІ CryptoBot. РҹРҫРҝСҖРҫРұСғР№СӮРө РҝРҫР·Р¶Рө.")
                return
            msg = await update.effective_message.reply_text(
                f"РһРҝР»Р°СӮРёСӮРө СҮРөСҖРөР· CryptoBot: {usd_amount:.2f} {asset} вҶ’ {_credits_fmt_from_usd(usd_amount)}.\nРҹРҫСҒР»Рө РҫРҝР»Р°СӮСӢ РәСҖРөРҙРёСӮСӢ РҝРҫРҝРҫР»РҪСҸСӮСҒСҸ Р°РІСӮРҫРјР°СӮРёСҮРөСҒРәРё.",
                reply_markup=InlineKeyboardMarkup([
                    [InlineKeyboardButton("рҹ’  CryptoBot", url=pay_url)],
                    [InlineKeyboardButton("рҹ”Һ РҹСҖРҫРІРөСҖРёСӮСҢ", callback_data=f"crypto:check:{inv_id}")]
                ])
            )
            context.application.create_task(_poll_crypto_invoice(
                context, msg.chat_id, msg.message_id, update.effective_user.id, inv_id, usd_amount
            ))
            return

        if data.startswith("crypto:check:"):
            await q.answer()
            inv_id = data.split(":", 2)[-1]
            inv = await _crypto_get_invoice(inv_id)
            if not inv:
                await q.message.reply_text("РқРө РҪР°СҲС‘Р» СҒСҮС‘СӮ. РЎРҫР·РҙР°Р№СӮРө РҪРҫРІСӢР№.")
                return
            st = (inv.get("status") or "").lower()
            if st == "paid":
                usd_amount = float(inv.get("amount", 0.0))
                if (inv.get("asset") or "").upper() == "TON":
                    usd_amount *= TON_USD_RATE
                _wallet_total_add(update.effective_user.id, usd_amount)
                await q.message.reply_text(f"рҹ’і РһРҝР»Р°СӮР° РҝРҫР»СғСҮРөРҪР°. РқР°СҮРёСҒР»РөРҪРҫ: {_credits_fmt_from_usd(usd_amount)}.")
            elif st == "active":
                await q.answer("РҹР»Р°СӮС‘Р¶ РөСүС‘ РҪРө РҝРҫРҙСӮРІРөСҖР¶РҙС‘РҪ", show_alert=True)
            else:
                await q.message.reply_text(f"РЎСӮР°СӮСғСҒ СҒСҮС‘СӮР°: {st}")
            return

        # РҹРҫРҙРҝРёСҒРәР°: РІСӢРұРҫСҖ СҒРҝРҫСҒРҫРұР°
        if data.startswith("buy:"):
            await q.answer()
            _, tier, months = data.split(":", 2)
            months = int(months)
            desc = f"РҹРҫРҙРҝРёСҒРәР° {tier.upper()} РҪР° {months} РјРөСҒ."
            await q.message.reply_text(
                f"{desc}\nР’СӢРұРөСҖРёСӮРө СҒРҝРҫСҒРҫРұ РҫРҝР»Р°СӮСӢ:",
                reply_markup=InlineKeyboardMarkup([
                    [InlineKeyboardButton("рҹ’і РҡР°СҖСӮРҫР№ Р®Kassa", callback_data=f"buyinv:{tier}:{months}")],
                    [InlineKeyboardButton("рҹӘҷ РЎ РәСҖРөРҙРёСӮРҪРҫРіРҫ РұР°Р»Р°РҪСҒР°", callback_data=f"buywallet:{tier}:{months}")],
                ])
            )
            return

        # РҹРҫРҙРҝРёСҒРәР° СҮРөСҖРөР· Р®Kassa direct API РёР· СҒСӮР°СҖРҫРіРҫ buy-РјРөРҪСҺ
        if data.startswith("buyyoo:"):
            await q.answer("РЎРҫР·РҙР°СҺ СҒСҒСӢР»РәСғ РҪР° РҫРҝР»Р°СӮСғвҖҰ")
            try:
                _, m, tier, months = data.split(":", 3)
                months = int(months)
                method_map = {"sbp": "yoo_sbp", "sberpay": "yoo_sberpay", "tpay": "yoo_tpay", "mirpay": "yoo_mirpay"}
                method_key = method_map.get(m)
                if not method_key:
                    await q.message.reply_text("РқРөРёР·РІРөСҒСӮРҪСӢР№ СҒРҝРҫСҒРҫРұ РҫРҝР»Р°СӮСӢ.")
                    return
                if not _yoo_direct_configured():
                    await q.message.reply_text("вҡ пёҸ Р‘СӢСҒСӮСҖР°СҸ РҫРҝР»Р°СӮР° Р®Kassa РҝРҫРәР° РҪРө РҪР°СҒСӮСҖРҫРөРҪР°: РҪСғР¶РҪСӢ YOO_SHOP_ID/YOO_SECRET_KEY РёР»Рё Secret File yookassa.env СҒ YK_ID/YK_KEY.")
                    return
                pay = await _yoo_create_direct_payment(update.effective_user.id, tier, months, method_key)
                payment_id = str(pay.get("id") or "")
                conf = pay.get("confirmation") or {}
                pay_url = conf.get("confirmation_url") or conf.get("confirmation_data") or conf.get("external_url") or ""
                if not pay_url:
                    raise RuntimeError(f"YooKassa did not return confirmation url: {pay}")
                label = YOO_DIRECT_METHODS[method_key]["label"]
                msg = await q.message.reply_text(
                    f"РҹРҫРҙРҝРёСҒРәР° {tier.upper()} РҪР° {months} РјРөСҒ.\nРЎРҝРҫСҒРҫРұ РҫРҝР»Р°СӮСӢ: {label}\nРһСӮРәСҖРҫР№СӮРө СҒСҒСӢР»РәСғ РҙР»СҸ РҫРҝР»Р°СӮСӢ:",
                    reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton(f"{label} вҖ” РҫРҝР»Р°СӮРёСӮСҢ", url=pay_url)]])
                )
                kv_set(f"yoo:pending:{payment_id}", json.dumps({"user_id": update.effective_user.id, "tier": tier, "months": months, "method": method_key}, ensure_ascii=False))
                context.application.create_task(_poll_yoo_subscription_payment(context, msg.chat_id, msg.message_id, update.effective_user.id, payment_id, tier, months))
            except Exception as e:
                log.exception("buyyoo payment failed: %s", e)
                err = str(e)[:700]
                user_msg = "вҡ пёҸ РқРө СғРҙР°Р»РҫСҒСҢ СҒРҫР·РҙР°СӮСҢ РҫРҝР»Р°СӮСғ Р®Kassa. РҹРҫРҝСҖРҫРұСғР№СӮРө РҙСҖСғРіРҫР№ СҒРҝРҫСҒРҫРұ."
                if YOO_DEBUG_PAY_ERRORS:
                    user_msg += "\n\nР”РёР°РіРҪРҫСҒСӮРёРәР° Р®Kassa: " + err
                await q.message.reply_text(user_msg)
            return

        # РҹРҫРҙРҝРёСҒРәР° СҮРөСҖРөР· Р®Kassa
        if data.startswith("buyinv:"):
            await q.answer()
            _, tier, months = data.split(":", 2)
            months = int(months)
            payload, amount_rub, title = _plan_payload_and_amount(tier, months)
            desc = f"РһС„РҫСҖРјР»РөРҪРёРө РҝРҫРҙРҝРёСҒРәРё {tier.upper()} РҪР° {months} РјРөСҒ."
            ok = await _send_invoice_rub(title, desc, amount_rub, payload, update)
            if not ok:
                await q.answer("РқРө СғРҙР°Р»РҫСҒСҢ РІСӢСҒСӮР°РІРёСӮСҢ СҒСҮС‘СӮ", show_alert=True)
            return

        # РҹРҫРҙРҝРёСҒРәР° СҒРҝРёСҒР°РҪРёРөРј РёР· РәСҖРөРҙРёСӮРҪРҫРіРҫ РұР°Р»Р°РҪСҒР°
        if data.startswith("buywallet:"):
            await q.answer()
            _, tier, months = data.split(":", 2)
            months = int(months)
            amount_rub = _plan_rub(tier, {1: "month", 3: "quarter", 12: "year"}[months])
            need_usd = _credits_to_usd(amount_rub)
            if _wallet_total_take(update.effective_user.id, need_usd):
                until = activate_subscription_with_tier(update.effective_user.id, tier, months)
                await q.message.reply_text(
                    f"вң… РҹРҫРҙРҝРёСҒРәР° {tier.upper()} Р°РәСӮРёРІРёСҖРҫРІР°РҪР° РҙРҫ {until.strftime('%Y-%m-%d')}.\n"
                    f"РЎРҝРёСҒР°РҪРҫ СҒ РұР°Р»Р°РҪСҒР°: {int(round(_usd_to_credits(need_usd)))} РәСҖ."
                )
            else:
                await q.message.reply_text(
                    "РқРөРҙРҫСҒСӮР°СӮРҫСҮРҪРҫ СҒСҖРөРҙСҒСӮРІ РҪР° РөРҙРёРҪРҫРј РұР°Р»Р°РҪСҒРө.\nРҹРҫРҝРҫР»РҪРёСӮРө РұР°Р»Р°РҪСҒ Рё РҝРҫРІСӮРҫСҖРёСӮРө.",
                    reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("вһ• РҹРҫРҝРҫР»РҪРёСӮСҢ РұР°Р»Р°РҪСҒ", callback_data="topup")]])
                )
            return

        # Р’СӢРұРҫСҖ РҙРІРёР¶РәР°
        if data.startswith("engine:"):
            await q.answer()
            engine = data.split(":", 1)[1]
            if engine == "luma" and LUMA_TEMP_DISABLED:
                await q.message.reply_text("вҡ пёҸ Luma РІСҖРөРјРөРҪРҪРҫ РҫСӮРәР»СҺСҮРөРҪР° Рё СҒРәСҖСӢСӮР° РёР· РјРөРҪСҺ. Р”Р»СҸ video РёСҒРҝРҫР»СҢР·СғР№СӮРө Sora 2 РұРөР· Р»СҺРҙРөР№, Kling РёР»Рё Runway.")
                return
            if engine == "midjourney":
                context.user_data["awaiting_midjourney_prompt"] = True
                await q.message.reply_text(
                    ENGINE_INFO_TEXT["midjourney"] +
                    f"\n\nРЎСӮРҫРёРјРҫСҒСӮСҢ РҫРҙРҪРҫР№ РіРөРҪРөСҖР°СҶРёРё: {_retail_credits(MIDJOURNEY_UNIT_COST_USD)} РәСҖ. РқР°РҝРёСҲРёСӮРө РҝСҖРҫРјРҝСӮ СҒР»РөРҙСғСҺСүРёРј СҒРҫРҫРұСүРөРҪРёРөРј РёР»Рё РёСҒРҝРҫР»СҢР·СғР№СӮРө /mj <РҫРҝРёСҒР°РҪРёРө>."
                )
                return
            if engine in ENGINE_INFO_TEXT:
                price_suffix = ""
                if engine == "images":
                    price_suffix = f"\n\nРҰРөРҪР° РіРөРҪРөСҖР°СҶРёРё: {_retail_credits(IMG_COST_USD)} РәСҖ."
                elif engine == "runway":
                    price_suffix = f"\n\nРҰРөРҪР° Gen-4.5: {_video_price_credits('runway', 5)} РәСҖ. Р·Р° 5 СҒРөРә."
                elif engine == "sora":
                    price_suffix = f"\n\nРҰРөРҪР° Sora 2: {_video_price_credits('sora', 5)} РәСҖ. Р·Р° 5 СҒРөРә."
                elif engine == "kling":
                    price_suffix = f"\n\nРҰРөРҪР° Kling: {_video_price_credits('kling', 5)} РәСҖ. Р·Р° 5 СҒРөРә."
                elif engine == "suno":
                    price_suffix = f"\n\nРҰРөРҪР° РҫРҙРҪРҫР№ РіРөРҪРөСҖР°СҶРёРё: {_retail_credits(SUNO_COST_USD)} РәСҖ."
                await q.message.reply_text(ENGINE_INFO_TEXT[engine] + price_suffix, disable_web_page_preview=True)
                return
            username = (update.effective_user.username or "")
            if engine == "runway":
                await q.message.reply_text(
                    "вң… Runway РҙРҫСҒСӮСғРҝРөРҪ РҙР»СҸ РІРёРҙРөРҫ РҝРҫ СӮРөРәСҒСӮСғ Рё РҙР»СҸ РҫР¶РёРІР»РөРҪРёСҸ С„РҫСӮРҫ.\n"
                    "Р”Р»СҸ РҫР¶РёРІР»РөРҪРёСҸ Р·Р°РіСҖСғР·РёСӮРө С„РҫСӮРҫРіСҖР°С„РёСҺ Рё РҪР°Р¶РјРёСӮРө вңЁ РһР¶РёРІРёСӮСҢ (Runway) РёР»Рё РҫСӮРҝСҖР°РІСҢСӮРө С„РҫСӮРҫ СҒ РҝРҫРҙРҝРёСҒСҢСҺ: "
                    "В«РҫР¶РёРІРё С„РҫСӮРҫ: Р»С‘РіРәР°СҸ СғР»СӢРұРәР°, РҙРІРёР¶РөРҪРёРө РәР°РјРөСҖСӢ, 5 СҒРөРәСғРҪРҙ, 9:16В».\n\n"
                    "Р”Р»СҸ СҒРҫР·РҙР°РҪРёСҸ РІРёРҙРөРҫ РҝРҫ СӮРөРәСҒСӮСғ/РіРҫР»РҫСҒСғ РёСҒРҝРҫР»СҢР·СғР№СӮРө Sora 2 РұРөР· Р»СҺРҙРөР№, Kling РёР»Рё Runway."
                )
                return
            if is_unlimited(update.effective_user.id, username):
                await q.message.reply_text(
                    f"вң… Р”РІРёР¶РҫРә В«{engine}В» РҙРҫСҒСӮСғРҝРөРҪ РұРөР· РҫРіСҖР°РҪРёСҮРөРҪРёР№.\n"
                    f"Р”Р»СҸ textвҶ’video РҙРҫСҒСӮСғРҝРҪСӢ Sora 2 РұРөР· Р»СҺРҙРөР№, Kling Рё Runway."
                )
                return

            if engine in ("gpt", "stt_tts", "midjourney", "sora", "kling"):
                await q.message.reply_text(
                    f"вң… Р’СӢРұСҖР°РҪ В«{engine}В». РһСӮРҝСҖР°РІСҢСӮРө Р·Р°РҝСҖРҫСҒ СӮРөРәСҒСӮРҫРј/С„РҫСӮРҫ. "
                    f"Р”Р»СҸ РІРёРҙРөРҫ РҪР°РҝРёСҲРёСӮРө: В«СҒРҫР·РҙР°Р№ РІРёРҙРөРҫ вҖҰ 5 СҒРөРәСғРҪРҙ 16:9В» вҖ” СҸ РҝСҖРөРҙР»РҫР¶Сғ Sora 2 РұРөР· Р»СҺРҙРөР№, Kling Рё Runway."
                )
                return

            est_cost = IMG_COST_USD if engine == "images" else (0.40 if engine == "luma" else max(1.0, RUNWAY_UNIT_COST_USD))
            map_engine = {"images": "img", "luma": "luma", "runway": "runway"}[engine]
            ok, offer = _can_spend_or_offer(update.effective_user.id, username, map_engine, est_cost)

            if ok:
                await q.message.reply_text(
                    "вң… Р”РҫСҒСӮСғРҝРҪРҫ. " +
                    ("Р—Р°РҝСғСҒСӮРёСӮРө: /img РәРҫСӮ РІ РҫСҮРәР°С…" if engine == "images"
                     else "Р”Р»СҸ РІРёРҙРөРҫ РҝРҫ СӮРөРәСҒСӮСғ РҙРҫСҒСӮСғРҝРҪСӢ Sora 2 РұРөР· Р»СҺРҙРөР№, Kling Рё Runway.")
                )
                return

            if offer == "ASK_SUBSCRIBE":
                await q.message.reply_text(
                    "Р”Р»СҸ СҚСӮРҫРіРҫ РҙРІРёР¶РәР° РҪСғР¶РҪР° Р°РәСӮРёРІРҪР°СҸ РҝРҫРҙРҝРёСҒРәР° РёР»Рё РөРҙРёРҪСӢР№ РұР°Р»Р°РҪСҒ. РһСӮРәСҖРҫР№СӮРө /plans РёР»Рё РҝРҫРҝРҫР»РҪРёСӮРө В«рҹ§ҫ Р‘Р°Р»Р°РҪСҒВ».",
                    reply_markup=InlineKeyboardMarkup(
                        [[InlineKeyboardButton("вӯҗ РўР°СҖРёС„СӢ", web_app=WebAppInfo(url=TARIFF_URL))],
                         [InlineKeyboardButton("вһ• РҹРҫРҝРҫР»РҪРёСӮСҢ РұР°Р»Р°РҪСҒ", callback_data="topup")]]
                    ),
                )
                return

            try:
                need_usd = float(offer.split(":", 1)[-1])
            except Exception:
                need_usd = est_cost
            amount_rub = _calc_oneoff_price_rub(map_engine, need_usd)
            await q.message.reply_text(
                f"Р’Р°СҲ РҙРҪРөРІРҪРҫР№ Р»РёРјРёСӮ РҝРҫ В«{engine}В» РёСҒСҮРөСҖРҝР°РҪ. Р Р°Р·РҫРІР°СҸ РҝРҫРәСғРҝРәР° вүҲ {amount_rub} вӮҪ "
                f"РёР»Рё РҝРҫРҝРҫР»РҪРёСӮРө РұР°Р»Р°РҪСҒ РІ В«рҹ§ҫ Р‘Р°Р»Р°РҪСҒВ».",
                reply_markup=InlineKeyboardMarkup(
                    [
                        [InlineKeyboardButton("вӯҗ РўР°СҖРёС„СӢ", web_app=WebAppInfo(url=TARIFF_URL))],
                        [InlineKeyboardButton("вһ• РҹРҫРҝРҫР»РҪРёСӮСҢ РұР°Р»Р°РҪСҒ", callback_data="topup")],
                    ]
                ),
            )
            return

        # Р РөР¶РёРјСӢ / Р”РІРёР¶РәРё
        if data == "mode:engines":
            await q.answer()
            await q.message.reply_text("Р”РІРёР¶РәРё:", reply_markup=engines_kb())
            return

        if data.startswith("mode:set:"):
            await q.answer()
            mode = data.split(":")[-1]
            mode_set(update.effective_user.id, mode)
            if mode == "study":
                study_sub_set(update.effective_user.id, "explain")
                await q.message.reply_text("Р РөР¶РёРј В«РЈСҮС‘РұР°В» РІРәР»СҺСҮС‘РҪ. Р’СӢРұРөСҖРёСӮРө РҝРҫРҙСҖРөР¶РёРј:", reply_markup=study_kb())
            elif mode == "photo":
                await q.message.reply_text("Р РөР¶РёРј В«РӨРҫСӮРҫВ» РІРәР»СҺСҮС‘РҪ. РҹСҖРёСҲР»РёСӮРө РёР·РҫРұСҖР°Р¶РөРҪРёРө вҖ” РҝРҫСҸРІСҸСӮСҒСҸ РұСӢСҒСӮСҖСӢРө РәРҪРҫРҝРәРё.", reply_markup=photo_quick_actions_kb())
            elif mode == "docs":
                await q.message.reply_text("Р РөР¶РёРј В«Р”РҫРәСғРјРөРҪСӮСӢВ». РҹСҖРёСҲР»РёСӮРө PDF/DOCX/EPUB/TXT вҖ” СҒРҙРөР»Р°СҺ РәРҫРҪСҒРҝРөРәСӮ.")
            elif mode == "voice":
                await q.message.reply_text("Р РөР¶РёРј В«Р“РҫР»РҫСҒВ». РһСӮРҝСҖР°РІСҢСӮРө voice/audio. РһР·РІСғСҮРәР° РҫСӮРІРөСӮРҫРІ: /voice_on")
            else:
                await q.message.reply_text(f"Р РөР¶РёРј В«{mode}В» Р°РәСӮРёРІРёСҖРҫРІР°РҪ.")
            return

        if data.startswith("study:set:"):
            await q.answer()
            sub = data.split(":")[-1]
            study_sub_set(update.effective_user.id, sub)
            await q.message.reply_text(f"РЈСҮС‘РұР° вҶ’ {sub}. РқР°РҝРёСҲРёСӮРө СӮРөРјСғ/Р·Р°РҙР°РҪРёРө.", reply_markup=study_kb())
            return

        # Photo edits require cached image
        if data.startswith("pedit:"):
            await q.answer()
            img = _get_cached_photo(update.effective_user.id)
            if not img:
                await q.message.reply_text("РЎРҪР°СҮР°Р»Р° РҝСҖРёСҲР»РёСӮРө С„РҫСӮРҫ, Р·Р°СӮРөРј РІСӢРұРөСҖРёСӮРө РҙРөР№СҒСӮРІРёРө.", reply_markup=photo_quick_actions_kb())
                return
            if data == "pedit:avatar":
                _clear_medicine_wait(context)
                if context.user_data.get("avatar_tts_voice"):
                    _set_avatar_wait(context)
                    await q.message.reply_text(
                        f"рҹ—Ј РҹРҫСҖСӮСҖРөСӮ РІСӢРұСҖР°РҪ. Р“РҫР»РҫСҒ СғР¶Рө РІСӢРұСҖР°РҪ: {_avatar_tts_voice_label(_avatar_tts_voice_get(context))}. РўРөРҝРөСҖСҢ РҝСҖРёСҲР»РёСӮРө СӮРөРәСҒСӮ, voice РёР»Рё Р°СғРҙРёРҫС„Р°Р№Р» MP3/WAV/M4A/AAC РҙР»СҸ СҖРөСҮРё Р°РІР°СӮР°СҖР°."
                    )
                else:
                    _set_avatar_voice_choice_wait(context)
                    await q.message.reply_text(_avatar_voice_choice_text(), reply_markup=_avatar_voice_choice_kb("act"))
                return
            if data == "pedit:photoclip":
                _clear_medicine_wait(context)
                _set_photo_clip_wait(context)
                await q.message.reply_text("рҹҺө РӨРҫСӮРҫ РІСӢРұСҖР°РҪРҫ. РһРҝРёСҲРёСӮРө СҒСӮРёР»СҢ РІРёРҙРөРҫРәР»РёРҝР°: РјСғР·СӢРәР°, РҙРІРёР¶РөРҪРёРө, РҪР°СҒСӮСҖРҫРөРҪРёРө, РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ Рё С„РҫСҖРјР°СӮ 9:16/16:9.")
                return
            if data == "pedit:vocalclip":
                _clear_medicine_wait(context)
                _clear_transient_flows(context)
                _set_mode_clean(q.from_user.id, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "vocalclip")
                _set_vocal_clip_wait(context)
                await q.message.reply_text(
                    "рҹҺӨ РҹРҫСҖСӮСҖРөСӮ РІСӢРұСҖР°РҪ РҙР»СҸ РәР»РёРҝР° СҒ РІРҫРәР°Р»РҫРј. РўРөРҝРөСҖСҢ РҫРҝРёСҲРёСӮРө РҝРөСҒРҪСҺ/РәР»РёРҝ: СҒСӮРёР»СҢ, СҸР·СӢРә, РҪР°СҒСӮСҖРҫРөРҪРёРө, РҝСҖРёРҝРөРІ, РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ.\n\n"
                    "Р’Р°Р¶РҪРҫ: СҖРөР¶РёРј СҖР°СҒСҒСҮРёСӮР°РҪ РҪР° РҫРҙРҪРҫРіРҫ СҮРөР»РҫРІРөРәР° РІ РәР°РҙСҖРө."
                )
                return
            if data == "pedit:aiselfie":
                _clear_medicine_wait(context)
                _set_ai_selfie_wait(context)
                await q.message.reply_text("рҹӨі РӨРҫСӮРҫ РІСӢРұСҖР°РҪРҫ. РқР°РҝРёСҲРёСӮРө СҒСҶРөРҪСғ: СҒ РәР°РәРҫР№ Р·РҪР°РјРөРҪРёСӮРҫСҒСӮСҢСҺ/РҝРөСҖСҒРҫРҪР°Р¶РөРј, РіРҙРө, СҒСӮРёР»СҢ, С„РҫСҖРјР°СӮ. РқР°РҝСҖРёРјРөСҖ: В«СҒРөР»С„Рё СҒ РёР·РІРөСҒСӮРҪСӢРј Р°РәСӮС‘СҖРҫРј РҪР° РәСҖР°СҒРҪРҫР№ РҙРҫСҖРҫР¶РәРө, iPhone selfie, 4:5В».")
                return
            if data == "pedit:retouch":
                _clear_medicine_wait(context)
                _set_retouch_wait_text(context)
                await q.message.reply_text(
                    "рҹ§Ҫ Р§СӮРҫ СғРұСҖР°СӮСҢ Рё РіРҙРө РҪР°С…РҫРҙРёСӮСҒСҸ СҚР»РөРјРөРҪСӮ?\n\n"
                    "РқР°РҝСҖРёРјРөСҖ: В«РІРҫРҙСҸРҪРҫР№ Р·РҪР°Рә СҒРҝСҖР°РІР° СҒРҪРёР·СғВ», В«РҪР°РҙРҝРёСҒСҢ РҝРҫ СҶРөРҪСӮСҖСғВ», В«Р»РҫРіРҫСӮРёРҝ РІ Р»РөРІРҫРј РІРөСҖС…РҪРөРј СғРіР»СғВ».\n"
                    "РһСӮРҝСҖР°РІР»СҸСҸ РәРҫРјР°РҪРҙСғ, РІСӢ РҝРҫРҙСӮРІРөСҖР¶РҙР°РөСӮРө, СҮСӮРҫ СҚСӮРҫ РІР°СҲРө РёР·РҫРұСҖР°Р¶РөРҪРёРө РёР»Рё Сғ РІР°СҒ РөСҒСӮСҢ РҝСҖР°РІРҫ РөРіРҫ СҖРөРҙР°РәСӮРёСҖРҫРІР°СӮСҢ."
                )
                return
            if data == "pedit:back":
                context.user_data.pop("photo_flow", None)
                await q.message.reply_text("РӨРҫСӮРҫ-РјР°СҒСӮРөСҖСҒРәР°СҸ:", reply_markup=photo_quick_actions_kb()); return
            if data == "pedit:removebg":
                await _pedit_removebg(update, context, img); return
            if data == "pedit:replacebg":
                await q.message.reply_text("Р’СӢРұРөСҖРёСӮРө РҪРҫРІСӢР№ С„РҫРҪ. РўРөРҝРөСҖСҢ Р·Р°РјРөРҪР° СҖР°РұРҫСӮР°РөСӮ РІ 2 СҚСӮР°РҝР°: 1) Р°РәРәСғСҖР°СӮРҪРҫ РІСӢСҖРөР·Р°СҺ СҮРөР»РҫРІРөРәР°/РҫРұСҠРөРәСӮ, 2) РҝРҫРҙСҒСӮР°РІР»СҸСҺ РҪРҫРІСӢР№ С„РҫРҪ РұРөР· РҝРөСҖРөСҖРёСҒРҫРІРәРё СҒР°РјРҫРіРҫ СҮРөР»РҫРІРөРәР°. РңРҫР¶РҪРҫ РІСӢРұСҖР°СӮСҢ РҝСҖРөСҒРөСӮ РёР»Рё РҪР°РҝРёСҒР°СӮСҢ СҒРІРҫР№ РІР°СҖРёР°РҪСӮ.", reply_markup=background_presets_kb()); return
            if data == "pedit:faceswap":
                await _start_faceswap_flow(update, context, img); return
            if data.startswith("pedit:bg:"):
                kind = data.split(":")[-1]
                if kind == "custom":
                    _set_replacebg_wait_text(context)
                    await q.message.reply_text("РқР°РҝРёСҲРёСӮРө, РәР°РәРҫР№ С„РҫРҪ РҝРҫСҒСӮР°РІРёСӮСҢ. РҜ СҒРҫС…СҖР°РҪСҺ СҮРөР»РҫРІРөРәР° РёР· РёСҒС…РҫРҙРҪРҫРіРҫ С„РҫСӮРҫ, РҫСӮРҙРөР»СҢРҪРҫ РҝРҫРҙРұРөСҖСғ/СҒРіРөРҪРөСҖРёСҖСғСҺ РҪРҫРІСӢР№ С„РҫРҪ Рё Р·Р°СӮРөРј Р°РәРәСғСҖР°СӮРҪРҫ СҒРҫРұРөСҖСғ РёСӮРҫРі. РҹСҖРёРјРөСҖСӢ: В«РҙРҫСҖРҫРіРҫР№ РҫС„РёСҒ СҒ РҝР°РҪРҫСҖР°РјРҪСӢРјРё РҫРәРҪР°РјРёВ», В«РҝР»СҸР¶ РЎР°РјСғРё РҪР° Р·Р°РәР°СӮРөВ», В«Р°Р»СҢРҝРёР№СҒРәРёРө РіРҫСҖСӢ Р»РөСӮРҫРјВ», В«СӮРөСҖСҖР°СҒР° РҪРөРұРҫСҒРәСҖС‘РұР° РҪРҫСҮСҢСҺВ». ")
                    return
                await _pedit_replacebg(update, context, img, kind=kind); return
            if data == "pedit:outpaint":
                await _pedit_outpaint(update, context, img); return
            if data == "pedit:story":
                await _pedit_storyboard(update, context, img); return
            if data in ("pedit:revive", "pedit:revive_runway", "pedit:revive_luma", "pedit:revive_sora", "pedit:revive_kling"):
                engine = {
                    "pedit:revive": "runway",
                    "pedit:revive_runway": "runway",
                    "pedit:revive_luma": "luma",
                    "pedit:revive_sora": "sora",
                    "pedit:revive_kling": "kling",
                }.get(data, "runway")
                if engine == "luma" and LUMA_TEMP_DISABLED:
                    await q.message.reply_text("вҡ пёҸ Luma РІСҖРөРјРөРҪРҪРҫ РҫСӮРәР»СҺСҮРөРҪР° Рё СҒРәСҖСӢСӮР° РёР· РјРөРҪСҺ. РҳСҒРҝРҫР»СҢР·СғР№СӮРө Runway, Kling РёР»Рё Sora 2 РұРөР· Р»СҺРҙРөР№.")
                    return
                # Р’РёРҙРёРјСӢР№ ACK СҒСҖР°Р·Сғ РҝРҫСҒР»Рө РәР»РёРәР°. РўСҸР¶С‘Р»Р°СҸ РіРөРҪРөСҖР°СҶРёСҸ РёРҙС‘СӮ РҝРҫСҒР»Рө РұСӢСҒСӮСҖРҫРіРҫ РҫСӮРІРөСӮР° Telegram.
                with contextlib.suppress(Exception):
                    shown_engine = "Runway СҒ Р°РІСӮРҫ-СҖРөР·РөСҖРІРҫРј Kling" if engine == "runway" else engine.upper()
                    await q.message.reply_text(f"рҹҹў Р—Р°РҝСғСҒРәР°СҺ РҫР¶РёРІР»РөРҪРёРө: {shown_engine}. Р•СҒР»Рё РҫСҒРҪРҫРІРҪРҫР№ РҙРІРёР¶РҫРә РҪРөРҙРҫСҒСӮСғРҝРөРҪ, РҝРөСҖРөРәР»СҺСҮСғСҒСҢ РҪР° СҖРөР·РөСҖРІРҪСӢР№.")
                try:
                    await _start_photo_revival(update, context, engine=engine, img_bytes=img, prompt="")
                except Exception as e:
                    log.exception("pedit revive failed: %s", e)
                    await update.effective_message.reply_text("вҡ пёҸ РқРө СғРҙР°Р»РҫСҒСҢ Р·Р°РҝСғСҒСӮРёСӮСҢ РҫСҒРҪРҫРІРҪРҫР№ РҙРІРёР¶РҫРә. РһСӮРәСҖРҫР№СӮРө РјРөРҪСҺ Рё РҝРҫРҝСҖРҫРұСғР№СӮРө Kling РёР»Рё РҝРҫРІСӮРҫСҖРёСӮРө РҝРҫР·Р¶Рө.")
                return

            if data == "pedit:lumaimg":
                _mode_track_set(update.effective_user.id, "lumaimg_wait_text")
                await q.message.reply_text("РқР°РҝРёСҲРёСӮРө РҫРҙРҪРҫ РҝСҖРөРҙР»РҫР¶РөРҪРёРө вҖ” СҮСӮРҫ СҒРіРөРҪРөСҖРёСҖРҫРІР°СӮСҢ. РҜ СҒРҙРөР»Р°СҺ РәР°СҖСӮРёРҪРәСғ.")
                return
            if data == "pedit:vision":
                b64 = base64.b64encode(img).decode("ascii")
                mime = sniff_image_mime(img)
                ans = await ask_openai_vision("РһРҝРёСҲРё С„РҫСӮРҫ Рё СӮРөРәСҒСӮ РҪР° РҪС‘Рј РәСҖР°СӮРәРҫ.", b64, mime)
                await update.effective_message.reply_text(ans or "Р“РҫСӮРҫРІРҫ.")
                return

        if data.startswith("faceswap:target:"):
            await q.answer()
            try:
                idx = int(data.split(":")[-1])
            except Exception:
                idx = 0
            user_id = update.effective_user.id
            _faceswap_target_face_index_cache[user_id] = max(0, idx)
            await q.message.reply_text(f"вң… Р’СӢРұСҖР°РҪРҫ СҶРөР»РөРІРҫРө Р»РёСҶРҫ в„–{max(0, idx) + 1}. РўРөРҝРөСҖСҢ РҝСҖРёСҲР»РёСӮРө С„РҫСӮРҫ Р»РёСҶР°, РәРҫСӮРҫСҖРҫРө РҪСғР¶РҪРҫ РІСҒСӮР°РІРёСӮСҢ.")
            context.user_data["faceswap_flow"] = "await_source"
            return

        if data.startswith("faceswap:source:"):
            await q.answer()
            try:
                idx = int(data.split(":")[-1])
            except Exception:
                idx = 0
            user_id = update.effective_user.id
            _faceswap_source_face_index_cache[user_id] = max(0, idx)
            context.user_data["faceswap_flow"] = "ready"
            await q.message.reply_text(f"вң… Р’СӢРұСҖР°РҪРҫ Р»РёСҶРҫ-РёСҒСӮРҫСҮРҪРёРә в„–{max(0, idx) + 1}. Р’СӢРұРөСҖРёСӮРө РәР°СҮРөСҒСӮРІРҫ Р·Р°РјРөРҪСӢ:", reply_markup=face_swap_quality_kb())
            return

        if data.startswith("faceswap:run:"):
            await q.answer()
            quality = data.split(":")[-1]
            await _faceswap_process(update, context, quality=quality)
            return

        if data.startswith("chooseimg:"):
            await q.answer()
            try:
                _, engine, aid = data.split(":", 2)
            except Exception:
                await q.message.reply_text("РқРө СғРҙР°Р»РҫСҒСҢ СҖР°СҒРҝРҫР·РҪР°СӮСҢ РІСӢРұРҫСҖ РҙРІРёР¶РәР°.")
                return
            meta = _pending_actions.pop(aid, None)
            if not meta or meta.get("kind") != "image_generate":
                await q.answer("Р—Р°РҙР°СҮР° СғСҒСӮР°СҖРөР»Р°", show_alert=True)
                return
            prompt = (meta.get("prompt") or "").strip()
            if not prompt:
                await q.message.reply_text("РҹСҖРҫРјРҝСӮ РҪРө РҪР°Р№РҙРөРҪ. Р—Р°РҝСғСҒСӮРёСӮРө Р·Р°РҝСҖРҫСҒ РөСүС‘ СҖР°Р·.")
                return
            await _run_selected_image_generation(update, context, prompt, engine)
            return

        # РҹРҫРҙСӮРІРөСҖР¶РҙРөРҪРёРө РІСӢРұРҫСҖР° РҙРІРёР¶РәР° РҙР»СҸ РІРёРҙРөРҫ
        if data.startswith("choose:"):
            await q.answer()
            _, engine, aid = data.split(":", 2)
            if engine == "luma" and LUMA_TEMP_DISABLED:
                _pending_actions.pop(aid, None)
                await q.message.reply_text("вҡ пёҸ Luma РІСҖРөРјРөРҪРҪРҫ РҫСӮРәР»СҺСҮРөРҪР°. Р’СӢРұРөСҖРёСӮРө Sora 2 РұРөР· Р»СҺРҙРөР№, Kling РёР»Рё Runway.")
                return
            meta = _pending_actions.pop(aid, None)
            if not meta:
                await q.answer("Р—Р°РҙР°СҮР° СғСҒСӮР°СҖРөР»Р°", show_alert=True); return
            prompt, duration, aspect = meta["prompt"], meta["duration"], meta["aspect"]
            engine = (engine or "").lower()
            if engine not in ("sora", "kling", "runway"):
                await update.effective_message.reply_text("вқҢ Р”РҫСҒСӮСғРҝРҪСӢ Sora 2 РұРөР· Р»СҺРҙРөР№, Kling Рё Runway."); return
            if engine == "runway" and not TEXT_VIDEO_ALLOW_RUNWAY:
                await update.effective_message.reply_text("вҡ пёҸ Runway РІСҖРөРјРөРҪРҪРҫ РҫСӮРәР»СҺСҮС‘РҪ РҪР°СҒСӮСҖРҫР№РәРҫР№ TEXT_VIDEO_ALLOW_RUNWAY."); return
            if engine == "sora" and _prompt_likely_has_people(prompt):
                await update.effective_message.reply_text("вҡ пёҸ Sora 2 РёСҒРҝРҫР»СҢР·СғРөСӮСҒСҸ СӮРҫР»СҢРәРҫ РұРөР· Р»СҺРҙРөР№. Р”Р»СҸ СҚСӮРҫРіРҫ Р·Р°РҝСҖРҫСҒР° РІСӢРұРөСҖРёСӮРө Kling РёР»Рё Runway."); return
            provider_cost = _video_provider_cost_usd(engine, duration)

            async def _start_real_render():
                if engine == "runway":
                    return await _run_runway_video(update, context, prompt, duration, aspect)
                return await _run_comet_text_video(update, context, engine, prompt, duration, aspect)

            await _try_pay_then_do(
                update, context, update.effective_user.id,
                "runway", provider_cost, _start_real_render,
                remember_kind=f"video_{engine}",
                remember_payload={"prompt": prompt, "duration": duration, "aspect": aspect, "engine": engine},
            )
            return

        await q.answer("РқРөРёР·РІРөСҒСӮРҪР°СҸ РәРҫРјР°РҪРҙР°", show_alert=True)

    except Exception as e:
        msg = str(e)
        if "query is too old" in msg.lower() or "query id is invalid" in msg.lower():
            log.warning("stale callback ignored: %s", msg)
            return
        log.exception("on_cb error: %s", e)
        with contextlib.suppress(Exception):
            await update.effective_message.reply_text("вҡ пёҸ РҡРҪРҫРҝРәР° СғСҒСӮР°СҖРөР»Р°. РһСӮРәСҖРҫР№СӮРө РјРөРҪСҺ Р·Р°РҪРҫРІРҫ Рё РҝРҫРІСӮРҫСҖРёСӮРө РҙРөР№СҒСӮРІРёРө.")
    finally:
        with contextlib.suppress(Exception):
            await q.answer()


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ STT в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def _mime_from_filename(fn: str) -> str:
    fnl = (fn or "").lower()
    if fnl.endswith((".ogg", ".oga")): return "audio/ogg"
    if fnl.endswith(".mp3"):           return "audio/mpeg"
    if fnl.endswith((".m4a", ".mp4")): return "audio/mp4"
    if fnl.endswith(".wav"):           return "audio/wav"
    if fnl.endswith(".webm"):          return "audio/webm"
    return "application/octet-stream"

async def transcribe_audio(buf: BytesIO, filename_hint: str = "audio.ogg") -> str:
    data = buf.getvalue()
    if DEEPGRAM_API_KEY:
        try:
            async with httpx.AsyncClient(timeout=60.0) as client:
                params = {"model": "nova-2", "language": "ru", "smart_format": "true", "punctuate": "true"}
                headers = {"Authorization": f"Token {DEEPGRAM_API_KEY}", "Content-Type": _mime_from_filename(filename_hint)}
                r = await client.post("https://api.deepgram.com/v1/listen", params=params, headers=headers, content=data)
                r.raise_for_status()
                dg = r.json()
                text = (dg.get("results", {}).get("channels", [{}])[0].get("alternatives", [{}])[0].get("transcript", "")).strip()
                if text: return text
        except Exception as e:
            log.exception("Deepgram STT error: %s", e)
    if oai_stt:
        try:
            buf2 = BytesIO(data); buf2.seek(0); setattr(buf2, "name", filename_hint)
            tr = oai_stt.audio.transcriptions.create(model=TRANSCRIBE_MODEL, file=buf2)
            return (tr.text or "").strip()
        except Exception as e:
            log.exception("Whisper STT error: %s", e)
    return ""


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ Р”РёР°РіРҪРҫСҒСӮРёРәР° РҙРІРёР¶РәРҫРІ в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
async def cmd_diag_stt(update: Update, context: ContextTypes.DEFAULT_TYPE):
    lines = []
    lines.append("рҹ”Һ STT РҙРёР°РіРҪРҫСҒСӮРёРәР°:")
    lines.append(f"вҖў OpenAI Whisper: {'вң… РәР»РёРөРҪСӮ Р°РәСӮРёРІРөРҪ' if oai_stt else 'вқҢ РҪРөРҙРҫСҒСӮСғРҝРөРҪ'}")
    lines.append(f"вҖў РңРҫРҙРөР»СҢ Whisper: {TRANSCRIBE_MODEL}")
    lines.append("вҖў РҹРҫРҙРҙРөСҖР¶РәР° С„РҫСҖРјР°СӮРҫРІ: ogg/oga, mp3, m4a/mp4, wav, webm")
    await update.effective_message.reply_text("\n".join(lines))

async def cmd_diag_images(update: Update, context: ContextTypes.DEFAULT_TYPE):
    key_env  = os.environ.get("OPENAI_IMAGE_KEY", "").strip()
    key_used = key_env or OPENAI_API_KEY
    base     = IMAGES_BASE_URL
    lines = [
        "рҹ§Ә Images (OpenAI) РҙРёР°РіРҪРҫСҒСӮРёРәР°:",
        f"вҖў OPENAI_IMAGE_KEY: {'вң… РҪР°Р№РҙРөРҪ' if key_used else 'вқҢ РҪРөСӮ'}",
        f"вҖў BASE_URL: {base}",
        f"вҖў MODEL: {IMAGES_MODEL}",
    ]
    if "openrouter" in (base or "").lower():
        lines.append("вҡ пёҸ BASE_URL СғРәР°Р·СӢРІР°РөСӮ РҪР° OpenRouter вҖ” СӮР°Рј РҪРөСӮ gpt-image-1.")
        lines.append("   РЈРәР°Р¶Рё https://api.openai.com/v1 (РёР»Рё СҒРІРҫР№ РҝСҖРҫРәСҒРё) РІ OPENAI_IMAGE_BASE_URL.")
    await update.effective_message.reply_text("\n".join(lines))

async def cmd_diag_video(update: Update, context: ContextTypes.DEFAULT_TYPE):
    lines = [
        f"рҹҺ¬ Р’РёРҙРөРҫ-РҙРІРёР¶РәРё / {PATCH_VERSION}:",
        f"вҖў Luma key: {'рҹҡ« СҒРәСҖСӢСӮРҫ РІСҖРөРјРөРҪРҪРҫ' if LUMA_TEMP_DISABLED else ('вң…' if bool(LUMA_API_KEY) else 'вқҢ')}  base={LUMA_BASE_URL}",
        f"  create={LUMA_CREATE_PATH}  status={LUMA_STATUS_PATH}  model={LUMA_MODEL}",
        f"вҖў Runway official: enabled={'вң…' if RUNWAY_DIRECT_ENABLED else 'вқҢ'} key={'вң…' if bool(RUNWAY_API_KEY) else 'вқҢ'} source={RUNWAY_KEY_SOURCE} fingerprint={runway_safe_key_fingerprint(RUNWAY_API_KEY)}",
        f"  base={RUNWAY_BASE_URL} version={RUNWAY_API_VERSION} text={RUNWAY_TEXT_CREATE_PATH} i2v={RUNWAY_I2V_PATH}",
        f"  uploads={RUNWAY_UPLOAD_PATH} org={RUNWAY_ORGANIZATION_PATH} tasks={RUNWAY_STATUS_PATH}",
        f"  text_models={','.join(_runway_direct_text_model_candidates())} i2v_models={','.join(_runway_direct_i2v_model_candidates())} poll={RUNWAY_DIRECT_POLL_INTERVAL_S:.1f}-{RUNWAY_DIRECT_POLL_MAX_INTERVAL_S:.1f}s upload_attempts={RUNWAY_DIRECT_UPLOAD_ATTEMPTS}",
        f"вҖў Comet key: {'вң…' if bool(COMET_API_KEY) else 'вқҢ'}  base={COMET_BASE_URL}",
        f"  Runway/Comet create={RUNWAY_COMET_CREATE_PATH}  status={RUNWAY_COMET_STATUS_PATH}",
        f"  Runway i2v enabled={'вң…' if RUNWAY_IMAGE2VIDEO_ENABLED else 'вқҢ'} models={','.join(_runway_i2v_model_candidates())} cooldown={_provider_cooldown_left('runway_i2v')}s hide_errors={'вң…' if RUNWAY_HIDE_TECH_ERRORS else 'вҖ”'}",
        f"вҖў BG remove: provider={BG_PROVIDER} photoroom={'вң…' if bool(PHOTOROOM_API_KEY) else 'вқҢ'} local_rembg={'вң…' if (LOCAL_REMBG_ENABLED and rembg_remove is not None) else 'вқҢ'}",
        f"вҖў Sora key: {'вң…' if bool(SORA_API_KEY) else 'вқҢ'}  model={SORA_MODEL}  create={SORA_CREATE_PATH}",
        f"вҖў Kling key: {'вң…' if bool(KLING_API_KEY) else 'вқҢ'}  model={KLING_MODEL}  create={KLING_CREATE_PATH}",
        f"вҖў Kling Avatar: create={KLING_AVATAR_CREATE_PATH}  status={KLING_AVATAR_STATUS_PATH}  mode={KLING_AVATAR_MODE}  avatar_voice_default={AVATAR_TTS_DEFAULT_VOICE} cost=${AVATAR_UNIT_COST_USD:.2f}",
        f"вҖў PhotoвҶ’clip pipeline: {'вң… on' if PHOTO_CLIP_PIPELINE else 'вҖ” off'}  engine={PHOTO_CLIP_VIDEO_ENGINE}  native_sound={'on' if PHOTO_CLIP_SOUND else 'off'}  mode={PHOTO_CLIP_MODE}  default={PHOTO_CLIP_DEFAULT_DURATION_S}s max={PHOTO_CLIP_MAX_DURATION_S}s mux_audio={'вң…' if PHOTO_CLIP_MUX_AUDIO else 'вҖ”'} cost=${PHOTO_CLIP_UNIT_COST_USD:.2f}",
        f"вҖў ffmpeg mux: timeout={FFMPEG_MUX_TIMEOUT_S}s copy_first={'вң…' if FFMPEG_MUX_COPY_FIRST else 'вҖ”'} preset={FFMPEG_MUX_REENCODE_PRESET} crf={FFMPEG_MUX_CRF} scale_h={FFMPEG_MUX_SCALE_HEIGHT} fps={FFMPEG_MUX_FPS} audio={FFMPEG_MUX_AUDIO_BITRATE} max={FFMPEG_MUX_MAX_MB}MB",
        f"вҖў Suno for photoвҶ’clip: {'вң… auto' if SUNO_AUTO_FOR_PHOTO_CLIP else 'вҖ” off'}  enabled={'вң…' if SUNO_ENABLED else 'вҖ”'} key={'вң…' if bool(SUNO_API_KEY) else 'вқҢ'} create={SUNO_CREATE_PATH} model={SUNO_MODEL}",
        f"вҖў AI selfie: provider={AI_SELFIE_PROVIDER} comet_key={'on' if bool(COMET_API_KEY) else 'off'} model={COMET_IMAGE_EDIT_MODEL} fallbacks={','.join(COMET_IMAGE_EDIT_FALLBACK_MODELS)} path={COMET_IMAGE_EDIT_PATH} timeout={COMET_IMAGE_EDIT_TIMEOUT_S}s max_side={AI_SELFIE_MAX_SIDE} size={AI_SELFIE_IMAGE_SIZE} fast={AI_SELFIE_FAST_MODE} cost=${AI_SELFIE_UNIT_COST_USD:.2f}",
        f"вҖў РқРҫСҖРјР°Р»РёР·Р°СҶРёСҸ duration: Kling 5/10 СҒРөРә; Sora 4/8/12 СҒРөРә РұРөР· Р»СҺРҙРөР№; Runway textвҶ’video Рё imageвҶ’video; Luma РІСҖРөРјРөРҪРҪРҫ СҒРәСҖСӢСӮР°",
        f"вҖў РҹРҫР»Р»РёРҪРі РәР°Р¶РҙСӢРө {VIDEO_POLL_DELAY_S:.1f} c",
    ]
    await update.effective_message.reply_text("\n".join(lines))

async def cmd_diag_runway(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Runway/Comet health diagnostic.

    Р‘РөР· Р°СҖРіСғРјРөРҪСӮРҫРІ вҖ” РұРөР·РҫРҝР°СҒРҪР°СҸ РҙРёР°РіРҪРҫСҒСӮРёРәР° РәРҫРҪС„РёРіСғСҖР°СҶРёРё Рё circuit breaker.
    /diag_runway auth вҖ” РұРөР·РҫРҝР°СҒРҪРҫ РҝСҖРҫРІРөСҖСҸРөСӮ РҫС„РёСҶРёР°Р»СҢРҪСӢР№ РәР»СҺСҮ Рё API-РәСҖРөРҙРёСӮСӢ.
    /diag_runway reset вҖ” СҒРұСҖРҫСҒ cooldown.
    /diag_runway test вҖ” РҝСҖРҫРұСғРөСӮ СҖРөР°Р»СҢРҪСғСҺ i2v Р·Р°РҙР°СҮСғ РҝРҫ РҝРҫСҒР»РөРҙРҪРөРјСғ С„РҫСӮРҫ.
    """
    args = [str(a).lower() for a in (context.args or [])]
    if "reset" in args:
        _provider_reset("runway_i2v")
        _provider_reset("runway_text_comet")
        _provider_reset("runway_direct")
        await update.effective_message.reply_text("вң… Runway provider cooldown СҒРұСҖРҫСҲРөРҪ РҙР»СҸ direct, imageвҶ’video Рё textвҶ’video.")
        return

    if "auth" in args:
        ok, detail = await _runway_direct_org_info()
        prefix = "вң… РһС„РёСҶРёР°Р»СҢРҪСӢР№ Runway API РҙРҫСҒСӮСғРҝРөРҪ" if ok else "вқҢ РһС„РёСҶРёР°Р»СҢРҪСӢР№ Runway API РҪРө РіРҫСӮРҫРІ"
        await update.effective_message.reply_text(prefix + "\n" + detail)
        return

    lines = [
        f"рҹ§Ә Runway/Comet РҙРёР°РіРҪРҫСҒСӮРёРәР° / {PATCH_VERSION}",
        f"вҖў enabled: {'вң…' if RUNWAY_IMAGE2VIDEO_ENABLED else 'вқҢ'}  use_comet={'вң…' if RUNWAY_USE_COMET else 'вҖ”'}",
        f"вҖў comet_key: {'вң…' if bool(COMET_API_KEY) else 'вқҢ'}  base={COMET_BASE_URL}",
        f"вҖў create={RUNWAY_COMET_CREATE_PATH}  status={RUNWAY_COMET_STATUS_PATH}",
        f"вҖў version={RUNWAY_API_VERSION or '2024-11-06'}",
        f"вҖў models={', '.join(_runway_i2v_model_candidates())}",
        f"вҖў direct_enabled: {'вң…' if RUNWAY_DIRECT_ENABLED else 'вқҢ'}  direct_key: {'вң…' if bool(RUNWAY_API_KEY) else 'вқҢ'}  direct_first={'вң…' if RUNWAY_DIRECT_FIRST else 'вҖ”'}",
        f"вҖў key_source={RUNWAY_KEY_SOURCE}  fingerprint={runway_safe_key_fingerprint(RUNWAY_API_KEY)}",
        f"вҖў official_base={RUNWAY_BASE_URL} version={RUNWAY_API_VERSION}",
        f"вҖў official_endpoints: text={RUNWAY_TEXT_CREATE_PATH}  i2v={RUNWAY_I2V_PATH}  upload={RUNWAY_UPLOAD_PATH}  org={RUNWAY_ORGANIZATION_PATH}",
        f"вҖў direct text models={', '.join(_runway_direct_text_model_candidates())}  i2v models={', '.join(_runway_direct_i2v_model_candidates())}",
        f"вҖў polling={RUNWAY_DIRECT_POLL_INTERVAL_S:.1f}-{RUNWAY_DIRECT_POLL_MAX_INTERVAL_S:.1f}s  retries={RUNWAY_DIRECT_RETRY_ATTEMPTS}  upload_attempts={RUNWAY_DIRECT_UPLOAD_ATTEMPTS}  data_uri_fallback={'вң…' if RUNWAY_DIRECT_DATA_URI_FALLBACK else 'вҖ”'}",
        f"вҖў hide_tech_errors={'вң…' if RUNWAY_HIDE_TECH_ERRORS else 'вҖ”'}  fallback_kling={'вң…' if RUNWAY_TEXT_FALLBACK_KLING and RUNWAY_AUTO_FALLBACK_KLING else 'вҖ”'}",
        f"вҖў i2v_cooldown={_provider_cooldown_left('runway_i2v')}s  text_cooldown={_provider_cooldown_left('runway_text_comet')}s",
    ]
    if _provider_last_error.get("runway_i2v"):
        lines.append("вҖў last_error=" + _provider_last_error.get("runway_i2v", "")[:700])

    if "test" not in args:
        lines.append("")
        lines.append("РҹСҖРҫРІРөСҖРёСӮСҢ РҫС„РёСҶРёР°Р»СҢРҪСӢР№ РәР»СҺСҮ Рё API-РәСҖРөРҙРёСӮСӢ: /diag_runway auth")
        lines.append("Р”Р»СҸ СҖРөР°Р»СҢРҪРҫРіРҫ СӮРөСҒСӮР° РҫСӮРҝСҖР°РІСҢСӮРө СҮРёСҒСӮРҫРө С„РҫСӮРҫ, Р·Р°СӮРөРј: /diag_runway test")
        await update.effective_message.reply_text("\n".join(lines)[:3900])
        return

    img = _get_cached_photo(update.effective_user.id)
    if not img:
        lines.append("вқҢ РқРөСӮ РҝРҫСҒР»РөРҙРҪРөРіРҫ С„РҫСӮРҫ. РЎРҪР°СҮР°Р»Р° РҫСӮРҝСҖР°РІСҢСӮРө С„РҫСӮРҫ РІ РұРҫСӮ, Р·Р°СӮРөРј /diag_runway test")
        await update.effective_message.reply_text("\n".join(lines)[:3900])
        return

    await update.effective_message.reply_text("\n".join(lines)[:2500] + "\n\nв–¶пёҸ Р—Р°РҝСғСҒРәР°СҺ СҖРөР°Р»СҢРҪСӢР№ РәРҫСҖРҫСӮРәРёР№ СӮРөСҒСӮ РҫС„РёСҶРёР°Р»СҢРҪРҫРіРҫ Runway. РӯСӮРҫ СҒРҝРёСҲРөСӮ API-РәСҖРөРҙРёСӮСӢ, РөСҒР»Рё Р·Р°РҙР°СҮР° РұСғРҙРөСӮ РҝСҖРёРҪСҸСӮР°.")
    if not RUNWAY_API_KEY:
        await update.effective_message.reply_text("вқҢ RUNWAYML_API_SECRET РҪРө РҪР°Р№РҙРөРҪ. Р”РҫРұР°РІСҢСӮРө РәР»СҺСҮ РІ Render Secret File runway.env РёР»Рё РІ Environment.")
        return
    ok = await _run_runway_direct_animate_photo(
        update, context, img,
        "subtle portrait animation, keep identity, small natural motion",
        5, "9:16",
    )
    if ok:
        await update.effective_message.reply_text("вң… Official Runway test: Р·Р°РҙР°СҮР° РҫСӮСҖР°РұРҫСӮР°Р»Р°.")
    else:
        await update.effective_message.reply_text("вҡ пёҸ Official Runway test РҪРө РҝСҖРҫСҲС‘Р». Р’СӢРҝРҫР»РҪРёСӮРө /diag_runway auth Рё РҝСҖРҫРІРөСҖСҢСӮРө API-РәСҖРөРҙРёСӮСӢ РІ Developer Portal.")

async def cmd_provider_status(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Short provider health dashboard for production checks."""
    args = [str(a).lower() for a in (context.args or [])]
    if "reset" in args:
        for name in ("runway_i2v",):
            _provider_reset(name)
        await update.effective_message.reply_text("вң… Provider cooldown СҒРұСҖРҫСҲРөРҪ.")
        return
    lines = [
        f"рҹ“Ҡ Provider status / {PATCH_VERSION}",
        f"вҖў Runway i2v: {'вң… available' if _provider_is_available('runway_i2v') else 'рҹҹЎ cooldown'} / cooldown={_provider_cooldown_left('runway_i2v')}s / fails={_provider_fail_counts.get('runway_i2v', 0)}",
        f"вҖў Runway models: {', '.join(_runway_i2v_model_candidates())}",
        f"вҖў Kling fallback: {'вң…' if (RUNWAY_AUTO_FALLBACK_KLING and bool(COMET_API_KEY)) else 'вқҢ'} / path={KLING_CREATE_PATH}",
        f"вҖў Hide tech errors: {'вң…' if RUNWAY_HIDE_TECH_ERRORS else 'вқҢ'}",
        f"вҖў I2V preprocess: {'вң…' if I2V_PREPROCESS_ENABLED else 'вқҢ'} / max_side={I2V_MAX_SOURCE_SIDE}",
    ]
    if _provider_last_error.get("runway_i2v"):
        lines.append("вҖў Runway last_error: " + _provider_last_error.get("runway_i2v", "")[:900])
    await update.effective_message.reply_text("\n".join(lines)[:3900])


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ MIME РҙР»СҸ РёР·РҫРұСҖР°Р¶РөРҪРёР№ в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def sniff_image_mime(data: bytes) -> str:
    if not data or len(data) < 12:
        return "application/octet-stream"
    b = data[:12]
    if b.startswith(b"\x89PNG\r\n\x1a\n"):
        return "image/png"
    if b[0:3] == b"\xff\xd8\xff":
        return "image/jpeg"
    if b[0:4] == b"RIFF" and b[8:12] == b"WEBP":
        return "image/webp"
    return "application/octet-stream"

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РҹР°СҖСҒ РҫРҝСҶРёР№ РІРёРҙРөРҫ в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
_ASPECTS = {"9:16", "16:9", "1:1", "4:5", "3:4", "4:3"}

def parse_video_opts(text: str) -> tuple[int, str]:
    tl = (text or "").lower()
    m = re.search(r"(\d+)\s*(?:СҒРөРә|СҒ)\b", tl)
    duration = int(m.group(1)) if m else LUMA_DURATION_S
    duration = max(3, min(20, duration))
    asp = None
    for a in _ASPECTS:
        if a in tl:
            asp = a
            break
    aspect = asp or (LUMA_ASPECT if LUMA_ASPECT in _ASPECTS else "16:9")
    return duration, aspect


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ Luma video в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
async def _run_luma_video(
    update: Update,
    context: ContextTypes.DEFAULT_TYPE,
    prompt: str,
    duration_s: int,
    aspect: str,
):
    await context.bot.send_chat_action(update.effective_chat.id, ChatAction.RECORD_VIDEO)
    try:
        async with httpx.AsyncClient(timeout=60.0) as client:
            base = await _pick_luma_base(client)
            create_url = f"{base}{LUMA_CREATE_PATH}"

            headers = {
                "Authorization": f"Bearer {LUMA_API_KEY}",
                "Accept": "application/json",
            }
            payload = {
                "model": LUMA_MODEL,
                "prompt": prompt,
                "duration": f"{duration_s}s",
                "aspect_ratio": aspect,
            }

            # СҒРҫР·РҙР°С‘Рј Р·Р°РҙР°СҮСғ
            r = await client.post(create_url, headers=headers, json=payload)
            if r.status_code >= 400:
                await update.effective_message.reply_text(
                    f"вҡ пёҸ Luma РҫСӮРәР»РҫРҪРёР»Р° Р·Р°РҙР°СҮСғ ({r.status_code})."
                )
                return

            data = r.json() or {}
            rid = data.get("id") or data.get("generation_id")
            if not rid:
                log.error("Luma: no generation id in response: %s", data)
                await update.effective_message.reply_text("вҡ пёҸ Luma РҪРө РІРөСҖРҪСғР»Р° id РіРөРҪРөСҖР°СҶРёРё.")
                return

            await update.effective_message.reply_text(
                "вҸі Luma СҖРөРҪРҙРөСҖРёСӮвҖҰ РҜ СҒРҫРҫРұСүСғ, РәРҫРіРҙР° РІРёРҙРөРҫ РұСғРҙРөСӮ РіРҫСӮРҫРІРҫ."
            )

            status_url = f"{base}{LUMA_STATUS_PATH}".format(id=rid)
            started = time.time()

            while True:
                rs = await client.get(status_url, headers=headers)
                try:
                    js = rs.json() or {}
                except Exception:
                    js = {}

                st = (js.get("state") or js.get("status") or "").lower()

                if st in ("completed", "succeeded", "finished", "ready"):
                    # --- РқРһР’Р«Рҷ РҪР°РҙС‘Р¶РҪСӢР№ РҝРҫРёСҒРә СҒСҒСӢР»РәРё РҪР° РІРёРҙРөРҫ ---
                    url = None
                    assets = js.get("assets")

                    def _extract_urls_from_assets(a):
                        urls = []
                        if isinstance(a, str):
                            urls.append(a)
                        elif isinstance(a, dict):
                            # СӮРёРҝРёСҮРҪСӢР№ С„РҫСҖРјР°СӮ: {"video": "https://..."} РёР»Рё {"video": {"url": "..."}}
                            for v in a.values():
                                urls.extend(_extract_urls_from_assets(v))
                        elif isinstance(a, (list, tuple)):
                            for item in a:
                                urls.extend(_extract_urls_from_assets(item))
                        return urls

                    if assets is not None:
                        for u in _extract_urls_from_assets(assets):
                            if isinstance(u, str) and u.startswith("http"):
                                url = u
                                break

                    # Р·Р°РҝР°СҒРҪСӢРө РәР»СҺСҮРё РҪР° РІСҒСҸРәРёР№ СҒР»СғСҮР°Р№
                    if not url:
                        for k in ("output_url", "video_url", "url"):
                            val = js.get(k)
                            if isinstance(val, str) and val.startswith("http"):
                                url = val
                                break

                    if not url:
                        log.error("Luma: РҫСӮРІРөСӮ РұРөР· СҒСҒСӢР»РәРё РҪР° РІРёРҙРөРҫ: %s", js)
                        await update.effective_message.reply_text(
                            "вқҢ Luma: РҫСӮРІРөСӮ РҝСҖРёСҲС‘Р» РұРөР· СҒСҒСӢР»РәРё РҪР° РІРёРҙРөРҫ."
                        )
                        return

                    # РЎРәР°СҮРёРІР°РөРј Рё РҫСӮРҝСҖР°РІР»СҸРөРј С„Р°Р№Р» РәР°Рә РІРёРҙРөРҫ
                    try:
                        v = await client.get(url, timeout=120.0)
                        v.raise_for_status()
                        bio = BytesIO(v.content)
                        bio.name = "luma.mp4"
                        await update.effective_message.reply_video(
                            InputFile(bio),
                            caption="рҹҺ¬ Luma: РіРҫСӮРҫРІРҫ вң…",
                        )
                    except Exception:
                        # РөСҒР»Рё РҪРө РҝРҫР»СғСҮРёР»РҫСҒСҢ СҒРәР°СҮР°СӮСҢ вҖ” С…РҫСӮСҸ РұСӢ РҙР°С‘Рј РҝСҖСҸРјСғСҺ СҒСҒСӢР»РәСғ
                        await update.effective_message.reply_text(
                            f"рҹҺ¬ Luma: РіРҫСӮРҫРІРҫ вң…\n{url}"
                        )
                    return

                if st in ("failed", "error", "canceled", "cancelled"):
                    log.error("Luma returned error state: %s", js)
                    await update.effective_message.reply_text("вқҢ Luma: РҫСҲРёРұРәР° СҖРөРҪРҙРөСҖР°.")
                    return

                if time.time() - started > LUMA_MAX_WAIT_S:
                    await update.effective_message.reply_text(
                        "вҢӣ Luma: РІСҖРөРјСҸ РҫР¶РёРҙР°РҪРёСҸ РІСӢСҲР»Рҫ."
                    )
                    return

                await asyncio.sleep(VIDEO_POLL_DELAY_S)

    except Exception as e:
        log.exception("Luma error: %s", e)
        await update.effective_message.reply_text(
            "вқҢ Luma: РҪРө СғРҙР°Р»РҫСҒСҢ Р·Р°РҝСғСҒСӮРёСӮСҢ/РҝРҫР»СғСҮРёСӮСҢ РІРёРҙРөРҫ."
        )
# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ Runway video в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def _dedupe_models(*items: str) -> list[str]:
    out: list[str] = []
    for m in items:
        m = (m or "").strip()
        if m and m not in out:
            out.append(m)
    return out

def _runway_i2v_model_candidates() -> list[str]:
    """РңРҫРҙРөР»Рё Runway РҙР»СҸ РҫР¶РёРІР»РөРҪРёСҸ С„РҫСӮРҫ/imageвҶ’video.

    v77: СҒРҝРёСҒРҫРә РұРөСҖС‘СӮСҒСҸ РёР· RUNWAY_IMAGE2VIDEO_MODELS. РҹРөСҖРІСӢРј СҒСӮРҫРёСӮ РҝСғРұР»РёСҮРҪСӢР№
    Р°Р»РёР°СҒ Comet runwayml-image-to-video, Р·Р°СӮРөРј backend-Р°Р»РёР°СҒСӢ. gen4.5 РҪР°РјРөСҖРөРҪРҪРҫ
    РҪРө РІРәР»СҺСҮС‘РҪ РІ РҙРөС„РҫР»СӮ, РҝРҫСӮРҫРјСғ СҮСӮРҫ РІ Р»РҫРіР°С… СҮР°СҒСӮРҫ РҫСӮРІРөСҮР°Р» model_not_found/no available channel.
    """
    env_models = [m.strip() for m in (RUNWAY_IMAGE2VIDEO_MODELS_ENV or "").split(",") if m.strip()]
    return _dedupe_models(*(env_models or ["runwayml-image-to-video", "gen4_turbo", "gen3a_turbo", "veo3.1_fast", "veo3.1", "veo3"]))

def _runway_direct_text_model_candidates() -> list[str]:
    """Official Runway textвҶ’video candidates. Gen-4.5 supports text-only input."""
    env_models = [m.strip() for m in (RUNWAY_DIRECT_TEXT_MODELS_ENV or "").split(",") if m.strip()]
    return _dedupe_models(*(env_models or [RUNWAY_TEXT_MODEL, "gen4.5"]))


def _runway_direct_i2v_model_candidates() -> list[str]:
    """Official Runway imageвҶ’video candidates only; no Comet aliases or Veo models."""
    env_models = [m.strip() for m in (RUNWAY_DIRECT_I2V_MODELS_ENV or "").split(",") if m.strip()]
    return _dedupe_models(*(env_models or ["gen4.5", "gen4_turbo"]))


def _runway_comet_text_model_candidates() -> list[str]:
    """Small fail-fast Comet candidate list; no payload/model storm when the channel is unavailable."""
    env_models = [m.strip() for m in (RUNWAY_COMET_TEXT_MODELS_ENV or "").split(",") if m.strip()]
    return _dedupe_models(*(env_models or ["runway-video", "gen4.5"]))


def _runway_direct_base_candidates() -> list[str]:
    """
    Direct Runway API СҒРөР№СҮР°СҒ СҖР°РұРҫСӮР°РөСӮ СҮРөСҖРөР· https://api.dev.runwayml.com/v1/... .
    Р•СҒР»Рё РІ ENV РҫСҒСӮР°Р»СҒСҸ СҒСӮР°СҖСӢР№ base_url, РІСҒС‘ СҖР°РІРҪРҫ РҙРҫРұР°РІР»СҸРөРј РҫС„РёСҶРёР°Р»СҢРҪСӢР№ base РәР°Рә fallback.
    """
    raw = (RUNWAY_BASE_URL or "").strip().rstrip("/")
    out: list[str] = []
    for base in (raw, "https://api.dev.runwayml.com"):
        if not base:
            continue
        if base.endswith("/v1"):
            base = base[:-3].rstrip("/")
        if base and base not in out:
            out.append(base)
    return out


def _runway_direct_headers() -> dict[str, str]:
    return {
        "Authorization": f"Bearer {RUNWAY_API_KEY}",
        "Accept": "application/json",
        "Content-Type": "application/json",
        "X-Runway-Version": RUNWAY_API_VERSION or "2024-11-06",
    }


def _runway_direct_ratio(aspect: str) -> str:
    """Current Runway Gen-4.5/Gen-4 Turbo landscape or portrait ratio."""
    return "720:1280" if (aspect or "").strip() in {"9:16", "3:4", "4:5"} else "1280:720"


def _runway_official_client() -> RunwayOfficialClient:
    return RunwayOfficialClient(
        RUNWAY_API_KEY,
        base_url=RUNWAY_BASE_URL,
        api_version=RUNWAY_API_VERSION or "2024-11-06",
        retry_attempts=RUNWAY_DIRECT_RETRY_ATTEMPTS,
        retry_base_s=RUNWAY_DIRECT_RETRY_BASE_S,
        poll_interval_s=RUNWAY_DIRECT_POLL_INTERVAL_S,
        poll_max_interval_s=RUNWAY_DIRECT_POLL_MAX_INTERVAL_S,
        upload_attempts=RUNWAY_DIRECT_UPLOAD_ATTEMPTS,
        data_uri_fallback=RUNWAY_DIRECT_DATA_URI_FALLBACK,
    )


def _runway_user_error_text(exc: Exception) -> str:
    if isinstance(exc, RunwayTaskTimeout):
        return "вҢӣ Runway РҪРө Р·Р°РІРөСҖСҲРёР» Р·Р°РҙР°СҮСғ Р·Р° РҫСӮРІРөРҙС‘РҪРҪРҫРө РІСҖРөРјСҸ. РҡСҖРөРҙРёСӮСӢ РұРҫСӮР° РҪРө СҒРҝРёСҒР°РҪСӢ."
    if isinstance(exc, RunwayAPIError):
        code = (exc.failure_code or "").upper()
        text = str(exc).lower()
        if code.startswith("SAFETY") or "safety" in text or "moderation" in text:
            return "вҡ пёҸ Runway РҫСӮРәР»РҫРҪРёР» Р·Р°РҝСҖРҫСҒ РҝРҫ РҝСҖР°РІРёР»Р°Рј РұРөР·РҫРҝР°СҒРҪРҫСҒСӮРё. РҳР·РјРөРҪРёСӮРө СҒСҶРөРҪСғ РёР»Рё С„РҫСҖРјСғР»РёСҖРҫРІРәСғ вҖ” РәСҖРөРҙРёСӮСӢ РҪРө СҒРҝРёСҒР°РҪСӢ."
        if exc.status_code in {401, 403}:
            return "вқҢ Runway РҪРө РҝСҖРёРҪСҸР» API-РәР»СҺСҮ. РҹСҖРҫРІРөСҖСҢСӮРө RUNWAYML_API_SECRET РІ Render Environment РёР»Рё Secret File runway.env."
        if exc.status_code == 402 or "credit" in text and ("insufficient" in text or "not enough" in text):
            return "вқҢ РқР° API-РұР°Р»Р°РҪСҒРө Runway РҪРөРҙРҫСҒСӮР°СӮРҫСҮРҪРҫ РәСҖРөРҙРёСӮРҫРІ. РҹРҫРҝРҫР»РҪРёСӮРө Billing РІ Runway Developer Portal."
        if exc.status_code == 429:
            return "вҡ пёҸ Р”РҫСҒСӮРёРіРҪСғСӮ Р»РёРјРёСӮ Runway РҙР»СҸ СӮРөРәСғСүРөРіРҫ API-tier. Р—Р°РҙР°СҮР° РұСғРҙРөСӮ РҪР°РҝСҖР°РІР»РөРҪР° РІ СҖРөР·РөСҖРІРҪСӢР№ РҙРІРёР¶РҫРә."
        if exc.status_code == 400:
            return "вҡ пёҸ Runway РҪРө РҝСҖРёРҪСҸР» РҝР°СҖР°РјРөСӮСҖСӢ Р·Р°РҙР°СҮРё. РҹСҖРҫРІРөСҖСҢСӮРө РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ, С„РҫСҖРјР°СӮ РёР·РҫРұСҖР°Р¶РөРҪРёСҸ Рё СӮРөРәСҒСӮ Р·Р°РҝСҖРҫСҒР°."
    return "вҡ пёҸ РһС„РёСҶРёР°Р»СҢРҪСӢР№ Runway РІСҖРөРјРөРҪРҪРҫ РҪРө РҫСӮРІРөСӮРёР». РҳСҒРҝРҫР»СҢР·СғСҺ СҖРөР·РөСҖРІРҪСӢР№ РјР°СҖСҲСҖСғСӮ."


async def _runway_request_with_retries(client: httpx.AsyncClient, method: str, url: str, *, headers: dict, json_body: dict | None = None):
    """Retry only transient Runway responses with exponential backoff and jitter."""
    transient = {429, 502, 503, 504}
    last = None
    for attempt in range(RUNWAY_DIRECT_RETRY_ATTEMPTS):
        try:
            r = await client.request(method, url, headers=headers, json=json_body)
            last = r
            if r.status_code not in transient or attempt >= RUNWAY_DIRECT_RETRY_ATTEMPTS - 1:
                return r
            delay = RUNWAY_DIRECT_RETRY_BASE_S * (2 ** attempt)
            delay += random.uniform(0, delay * 0.5)
            log.warning("Runway transient HTTP %s; retry %s/%s in %.1fs", r.status_code, attempt + 1, RUNWAY_DIRECT_RETRY_ATTEMPTS, delay)
            await asyncio.sleep(delay)
        except (httpx.TimeoutException, httpx.TransportError) as e:
            last = e
            if attempt >= RUNWAY_DIRECT_RETRY_ATTEMPTS - 1:
                raise
            delay = RUNWAY_DIRECT_RETRY_BASE_S * (2 ** attempt)
            delay += random.uniform(0, delay * 0.5)
            log.warning("Runway transport error; retry %s/%s in %.1fs: %s", attempt + 1, RUNWAY_DIRECT_RETRY_ATTEMPTS, delay, e)
            await asyncio.sleep(delay)
    return last


async def _runway_direct_org_info() -> tuple[bool, str]:
    """Read-only official Runway authentication, organization and credit check."""
    if not (RUNWAY_DIRECT_ENABLED and RUNWAY_API_KEY):
        return False, "RUNWAYML_API_SECRET РҪРө РҪР°Р№РҙРөРҪ РҪРё РІ Render Environment, РҪРё РІ Secret Files."
    format_ok, format_note = runway_key_format_hint(RUNWAY_API_KEY)
    source_name = RUNWAY_KEY_SOURCE
    try:
        async with _runway_official_client() as rw:
            org = await rw.organization(endpoint=RUNWAY_ORGANIZATION_PATH)
        payload = org.get("data") if isinstance(org.get("data"), dict) else org
        credits = next((payload.get(k) for k in ("creditBalance", "credits", "balance", "availableCredits") if k in payload), None)
        tier = next((payload.get(k) for k in ("tier", "usageTier", "rateLimitTier") if k in payload), None)
        parts = [
            "РәР»СҺСҮ РҝСҖРёРҪСҸСӮ",
            f"РҝРөСҖРөРјРөРҪРҪР°СҸ: {source_name}",
            f"fingerprint: {runway_safe_key_fingerprint(RUNWAY_API_KEY)}",
            format_note,
        ]
        if credits is not None:
            parts.append(f"API-РәСҖРөРҙРёСӮСӢ: {credits}")
        if tier is not None:
            parts.append(f"tier: {tier}")
        return True, "\n".join(parts)
    except Exception as e:
        log.warning("Runway organization auth failed: %s", e)
        return False, _runway_user_error_text(e) + f"\nFingerprint: {runway_safe_key_fingerprint(RUNWAY_API_KEY)}"

# v76 provider circuit breaker / health cache
_provider_fail_counts: dict[str, int] = {}
_provider_cooldown_until: dict[str, float] = {}
_provider_last_error: dict[str, str] = {}

def _provider_is_available(name: str) -> bool:
    return time.time() >= float(_provider_cooldown_until.get(name, 0) or 0)

def _provider_cooldown_left(name: str) -> int:
    return max(0, int(float(_provider_cooldown_until.get(name, 0) or 0) - time.time()))

def _provider_mark_success(name: str) -> None:
    _provider_fail_counts[name] = 0
    _provider_cooldown_until.pop(name, None)
    _provider_last_error.pop(name, None)

def _provider_mark_failure(name: str, reason: str = "") -> None:
    _provider_last_error[name] = (reason or "")[:700]
    n = int(_provider_fail_counts.get(name, 0) or 0) + 1
    _provider_fail_counts[name] = n
    if n >= max(1, RUNWAY_PROVIDER_FAIL_THRESHOLD):
        _provider_cooldown_until[name] = time.time() + max(30, RUNWAY_PROVIDER_COOLDOWN_S)

def _provider_reset(name: str) -> None:
    _provider_fail_counts.pop(name, None)
    _provider_cooldown_until.pop(name, None)
    _provider_last_error.pop(name, None)

def _is_runway_unavailable_text(s: str) -> bool:
    t = (s or "").lower()
    needles = (
        "model_not_found", "no available channel", "invalid url",
        "жӯӨжЁЎеһӢе·ІдёӢжһ¶", "model has been removed", "model is removed",
        "not found", "channel", "invalid_request_error"
    )
    return any(x in t for x in needles)

def _looks_like_screenshot_or_bad_i2v_source(img_bytes: bytes) -> str:
    """Soft heuristic only: warn user if input looks like a phone screenshot/frame."""
    if Image is None:
        return ""
    try:
        im = Image.open(BytesIO(img_bytes)).convert("RGB")
        w, h = im.size
        # РһСҮРөРҪСҢ РІСӢСҒРҫРәРёР№/СҲРёСҖРҫРәРёР№ РәР°РҙСҖ СҮР°СүРө РІСҒРөРіРҫ СҸРІР»СҸРөСӮСҒСҸ СҒРәСҖРёРҪСҲРҫСӮРҫРј СӮРөР»РөС„РҫРҪР°/СҚРәСҖР°РҪР°.
        if h > w * 1.55 or w > h * 1.55:
            return "РӨРҫСӮРҫ РҝРҫС…РҫР¶Рө РҪР° СҒРәСҖРёРҪСҲРҫСӮ/РәР°РҙСҖ СҒ РұРҫР»СҢСҲРёРјРё РҝРҫР»СҸРјРё. Р”Р»СҸ Р»СғСҮСҲРөРіРҫ РҫР¶РёРІР»РөРҪРёСҸ Р·Р°РіСҖСғР·РёСӮРө СҮРёСҒСӮСӢР№ РҝРҫСҖСӮСҖРөСӮ РұРөР· РёРҪСӮРөСҖС„РөР№СҒР° СӮРөР»РөС„РҫРҪР° Рё СҮС‘СҖРҪСӢС… СҖР°РјРҫРә."
        # Р‘РҫР»СҢСҲРёРө СӮС‘РјРҪСӢРө РҫРұР»Р°СҒСӮРё РҝРҫ РәСҖР°СҸРј вҖ” СҮР°СҒСӮСӢР№ РҝСҖРёР·РҪР°Рә С„РҫСӮРҫ СҚРәСҖР°РҪР°/РІРёРҙРөРҫ-РҝР»РөРөСҖР°.
        try:
            small = im.resize((64, 64))
            px = list(small.getdata())
            dark = sum(1 for r, g, b in px if max(r, g, b) < 32) / max(1, len(px))
            if dark > 0.35:
                return "Р’ РәР°РҙСҖРө РјРҪРҫРіРҫ СҮС‘СҖРҪСӢС… РҝРҫР»РөР№/СҚР»РөРјРөРҪСӮРҫРІ РёРҪСӮРөСҖС„РөР№СҒР°. РңРҫРҙРөР»СҢ РјРҫР¶РөСӮ РҫР¶РёРІРёСӮСҢ СҖР°РјРәСғ РёР»Рё СҚРәСҖР°РҪ РІРјРөСҒСӮРҫ СҮРөР»РҫРІРөРәР°. РӣСғСҮСҲРө Р·Р°РіСҖСғР·РёСӮСҢ СҮРёСҒСӮРҫРө С„РҫСӮРҫ."
        except Exception:
            pass
    except Exception:
        return ""
    return ""

def _prepare_i2v_source_image(img_bytes: bytes, aspect: str = "9:16") -> tuple[bytes, str]:
    """
    Production-safe preparation for imageвҶ’video providers.
    Р’РҫР·РІСҖР°СүР°РөСӮ (bytes, note). РқРө РҙРөР»Р°РөСӮ Р°РіСҖРөСҒСҒРёРІРҪСӢР№ face-crop, СҮСӮРҫРұСӢ РҪРө РёСҒРҝРҫСҖСӮРёСӮСҢ С„РҫСӮРҫ,
    РҪРҫ СғРұРёСҖР°РөСӮ СҸРІРҪСӢРө СҮС‘СҖРҪСӢРө СҖР°РјРәРё Рё РҪРҫСҖРјР°Р»РёР·СғРөСӮ СҖР°Р·РјРөСҖ/С„РҫСҖРјР°СӮ.
    """
    if not I2V_PREPROCESS_ENABLED or Image is None:
        return img_bytes, ""
    try:
        im = Image.open(BytesIO(img_bytes)).convert("RGB")
        w, h = im.size
        note_parts = []

        # 1) РЈРҙР°Р»РөРҪРёРө СҸРІРҪСӢС… СҮС‘СҖРҪСӢС… СҖР°РјРҫРә. РқРө СӮСҖРҫРіР°РөРј, РөСҒР»Рё crop СҒР»РёСҲРәРҫРј РјР°Р»/СҖРёСҒРәРҫРІР°РҪРҪСӢР№.
        if I2V_AUTOCROP_BLACK_BORDERS and min(w, h) >= 200:
            gray = im.convert("L")
            # РҹРёРәСҒРөР»Рё СҸСҖСҮРө РҝРҫСҖРҫРіР° СҒСҮРёСӮР°РөРј СҒРҫРҙРөСҖР¶РёРјСӢРј.
            mask = gray.point(lambda p: 255 if p > 28 else 0)
            bbox = mask.getbbox()
            if bbox:
                x1, y1, x2, y2 = bbox
                bw, bh = x2 - x1, y2 - y1
                area_ratio = (bw * bh) / max(1, w * h)
                # crop СӮРҫР»СҢРәРҫ РөСҒР»Рё РҫРҪ Р·Р°РјРөСӮРҪРҫ СғРұРёСҖР°РөСӮ РәСҖР°СҸ, РҪРҫ РҪРө РҝСҖРөРІСҖР°СүР°РөСӮ РәР°СҖСӮРёРҪРәСғ РІ РәСҖРҫСҲРөСҮРҪСӢР№ С„СҖР°РіРјРөРҪСӮ
                margin_removed = (x1 > w * 0.04 or y1 > h * 0.04 or x2 < w * 0.96 or y2 < h * 0.96)
                if margin_removed and 0.20 <= area_ratio <= 0.96:
                    pad = int(max(bw, bh) * 0.04)
                    x1 = max(0, x1 - pad); y1 = max(0, y1 - pad)
                    x2 = min(w, x2 + pad); y2 = min(h, y2 + pad)
                    im = im.crop((x1, y1, x2, y2))
                    w, h = im.size
                    note_parts.append("СғРұСҖР°Р» Р»РёСҲРҪРёРө СӮС‘РјРҪСӢРө РҝРҫР»СҸ")

        # 2) РқРҫСҖРјР°Р»РёР·Р°СҶРёСҸ СҖР°Р·РјРөСҖР°, СҮСӮРҫРұСӢ РҪРө РҫСӮРҝСҖР°РІР»СҸСӮСҢ РҫРіСҖРҫРјРҪСӢРө СҒРәСҖРёРҪСҲРҫСӮСӢ РҝСҖРҫРІР°Р№РҙРөСҖР°Рј.
        max_side = max(512, int(I2V_MAX_SOURCE_SIDE or 1280))
        if max(w, h) > max_side:
            im.thumbnail((max_side, max_side), getattr(Image, "Resampling", Image).LANCZOS)
            note_parts.append(f"СҒР¶Р°Р» РёСҒС…РҫРҙРҪРёРә РҙРҫ {max_side}px")

        out = BytesIO()
        im.save(out, format="JPEG", quality=92, optimize=True)
        return out.getvalue(), ", ".join(note_parts)
    except Exception as e:
        log.warning("i2v source preprocess failed: %s", e)
        return img_bytes, ""

async def _run_runway_video(update: Update, context: ContextTypes.DEFAULT_TYPE, prompt: str, duration_s: int, aspect: str):
    """Production Runway textвҶ’video: official API first, Comet second, Kling fallback.

    Official route follows current Runway documentation:
    POST /v1/text_to_video -> GET /v1/tasks/{id}.
    For backwards compatibility only, a 404/405 can fall back to
    /v1/image_to_video without promptImage.
    """
    await context.bot.send_chat_action(update.effective_chat.id, ChatAction.RECORD_VIDEO)
    prompt = (prompt or "").strip()
    if not prompt:
        await update.effective_message.reply_text("вқҢ Runway: РҝСғСҒСӮРҫР№ Р·Р°РҝСҖРҫСҒ РҙР»СҸ РІРёРҙРөРҫ.")
        return False

    duration = max(2, min(10, int(_duration_for_engine("runway", duration_s))))
    ratio = _runway_direct_ratio(aspect)
    errors: list[str] = []
    hard_stop = False

    async def try_direct() -> bool:
        nonlocal hard_stop
        if not (RUNWAY_DIRECT_ENABLED and RUNWAY_API_KEY):
            return False
        try:
            async with _runway_official_client() as rw:
                task_id = await rw.create_text_to_video(
                    prompt_text=prompt,
                    model=_runway_direct_text_model_candidates()[0],
                    ratio=ratio,
                    duration=duration,
                    endpoint=RUNWAY_TEXT_CREATE_PATH,
                    compatibility_endpoint=RUNWAY_TEXT_COMPAT_PATH,
                )
                await update.effective_message.reply_text(
                    f"вҸі Runway Gen-4.5: Р·Р°РҙР°СҮР° РҝСҖРёРҪСҸСӮР° ({duration} СҒ, {aspect}). РһР¶РёРҙР°СҺ СҖРөР·СғР»СҢСӮР°СӮвҖҰ"
                )
                result = await rw.wait_for_task(task_id, timeout_s=RUNWAY_MAX_WAIT_S)

            async with httpx.AsyncClient(timeout=240.0, follow_redirects=True) as dl_client:
                await _reply_video_from_url(
                    update, dl_client, result.first_output,
                    "Runway textвҶ’video вң… В· Powered by Runway",
                    task_id=task_id,
                )
            _provider_mark_success("runway_direct")
            return True
        except RunwayAPIError as e:
            errors.append(str(e)); _provider_mark_failure("runway_direct", str(e))
            log.warning("Official Runway text-to-video failed: %s", e)
            code = (e.failure_code or "").upper()
            if code.startswith("SAFETY") or e.status_code in {400, 401, 403, 402}:
                hard_stop = True
                await update.effective_message.reply_text(_runway_user_error_text(e))
            return False
        except Exception as e:
            errors.append(str(e)); _provider_mark_failure("runway_direct", str(e))
            log.exception("Official Runway text route failed: %s", e)
            return False

    async def try_comet() -> bool:
        provider_name = "runway_text_comet"
        if not (RUNWAY_USE_COMET and COMET_API_KEY):
            return False
        if not _provider_is_available(provider_name):
            log.warning("Runway text/Comet skipped: cooldown %ss", _provider_cooldown_left(provider_name))
            return False
        headers = {
            "Authorization": f"Bearer {COMET_API_KEY}",
            "Accept": "application/json",
            "Content-Type": "application/json",
            "X-Runway-Version": RUNWAY_API_VERSION or "2024-11-06",
        }
        async with httpx.AsyncClient(timeout=90.0) as client:
            for model in _runway_comet_text_model_candidates():
                payload = {"model": model, "promptText": prompt, "duration": duration, "ratio": ratio}
                try:
                    r = await client.post(f"{COMET_BASE_URL}{RUNWAY_COMET_CREATE_PATH}", headers=headers, json=payload)
                    if r.status_code >= 400:
                        err = f"Comet Runway {r.status_code}: {_api_error_preview(r)}"
                        errors.append(err); log.warning(err)
                        if _is_runway_unavailable_text(err) or r.status_code == 503:
                            _provider_mark_failure(provider_name, err)
                            break
                        continue
                    js = r.json() or {}
                    ready_url = _extract_first_url(js.get("output")) or _extract_first_url(js.get("data")) or _extract_first_url(js)
                    if ready_url:
                        _provider_mark_success(provider_name)
                        await _reply_video_from_url(update, client, ready_url, "Runway/Comet textвҶ’video вң… В· Powered by Runway")
                        return True
                    task_id = str(js.get("id") or js.get("task_id") or js.get("generation_id") or ((js.get("data") or {}).get("id") if isinstance(js.get("data"), dict) else "") or "").strip()
                    if not task_id:
                        err = f"Comet Runway: no task id: {json.dumps(js, ensure_ascii=False)[:500]}"
                        errors.append(err); _provider_mark_failure(provider_name, err)
                        continue
                    await update.effective_message.reply_text("вҸі Runway/Comet: Р·Р°РҙР°СҮР° РҝСҖРёРҪСҸСӮР°, РҫР¶РёРҙР°СҺ СҖРөР·СғР»СҢСӮР°СӮвҖҰ")
                    ok = bool(await _poll_video_task_generic(
                        update, client, headers, COMET_BASE_URL,
                        [RUNWAY_COMET_STATUS_PATH, "/runwayml/v1/tasks/{id}", "/v1/tasks/{id}"],
                        task_id, "Runway/Comet textвҶ’video В· Powered by Runway", RUNWAY_MAX_WAIT_S,
                    ))
                    if ok:
                        _provider_mark_success(provider_name)
                    else:
                        _provider_mark_failure(provider_name, "polling failed")
                    return ok
                except Exception as e:
                    err = f"Comet Runway exception: {e}"
                    errors.append(err); log.warning(err); _provider_mark_failure(provider_name, err)
        return False

    routes = (try_direct, try_comet) if RUNWAY_DIRECT_FIRST else (try_comet, try_direct)
    for route in routes:
        if hard_stop:
            return False
        try:
            if await route():
                return True
        except Exception as e:
            errors.append(str(e)); log.exception("Runway text route failed: %s", e)

    if hard_stop:
        return False

    if RUNWAY_TEXT_FALLBACK_KLING and RUNWAY_AUTO_FALLBACK_KLING and COMET_API_KEY:
        await update.effective_message.reply_text(RUNWAY_PUBLIC_FALLBACK_TEXT)
        try:
            return bool(await _run_comet_text_video(update, context, "kling", prompt, duration, aspect))
        except Exception as e:
            errors.append(f"Kling fallback: {e}"); log.exception("RunwayвҶ’Kling fallback failed: %s", e)

    if RUNWAY_HIDE_TECH_ERRORS:
        await update.effective_message.reply_text(
            "вҡ пёҸ Runway СҒРөР№СҮР°СҒ РҪРө РҝСҖРёРҪСҸР» Р·Р°РҙР°СҮСғ. РҡСҖРөРҙРёСӮСӢ Р·Р° РҪРөСғСҒРҝРөСҲРҪСғСҺ РіРөРҪРөСҖР°СҶРёСҺ РҪРө СҒРҝРёСҒСӢРІР°СҺСӮСҒСҸ. "
            "РҹСҖРҫРІРөСҖСҢСӮРө /diag_runway auth РёР»Рё РІСӢРұРөСҖРёСӮРө Kling."
        )
    else:
        details = "\n".join(errors[-3:]) or "API РҪРө РІРөСҖРҪСғР» РҝРҫРҙСҖРҫРұРҪРҫСҒСӮРё."
        await update.effective_message.reply_text(f"вқҢ Runway: Р·Р°РҙР°СҮР° РҪРө РІСӢРҝРҫР»РҪРөРҪР°.\n{details[:1600]}")
    return False

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ ImageвҶ’Video helpers в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def _api_error_preview(resp, limit: int = 900) -> str:
    try:
        body = json.dumps(resp.json(), ensure_ascii=False)
    except Exception:
        body = getattr(resp, "text", "") or ""
    body = re.sub(r"\s+", " ", body).strip()
    return body[:limit] if body else "РұРөР· СӮРөР»Р° РҫСӮРІРөСӮР°"


def _sora_people_moderation_text() -> str:
    return (
        "вҡ пёҸ Sora 2 Р·Р°РұР»РҫРәРёСҖРҫРІР°Р»Р° СҚСӮРҫ РёР·РҫРұСҖР°Р¶РөРҪРёРө РҪР° РјРҫРҙРөСҖР°СҶРёРё, РҝРҫСӮРҫРјСғ СҮСӮРҫ РҪР° С„РҫСӮРҫ РөСҒСӮСҢ СҮРөР»РҫРІРөРә/Р»СҺРҙРё.\n\n"
        "РӯСӮРҫ РҫРіСҖР°РҪРёСҮРөРҪРёРө Sora/Comet, Р° РҪРө РҫСҲРёРұРәР° РҙРөРҝР»РҫСҸ Рё РҪРө РҫСҲРёРұРәР° РәР»СҺСҮР°.\n\n"
        "Р”Р»СҸ РҫР¶РёРІР»РөРҪРёСҸ С„РҫСӮРҫ СҒ Р»СҺРҙСҢРјРё РёСҒРҝРҫР»СҢР·СғР№СӮРө:\n"
        "вҖў вңЁ РһР¶РёРІРёСӮСҢ (Runway)\n"
        "вҖў вңЁ РһР¶РёРІРёСӮСҢ (Kling)\n\n"
        "Sora 2 РҫСҒСӮР°РІР»РөРҪР° РҙР»СҸ РёР·РҫРұСҖР°Р¶РөРҪРёР№ РұРөР· Р»СҺРҙРөР№: РҝСҖРөРҙРјРөСӮСӢ, Р¶РёРІРҫСӮРҪСӢРө, Р·РҙР°РҪРёСҸ, РҝРөР№Р·Р°Р¶Рё, РёРҪСӮРөСҖСҢРөСҖ."
    )


def is_sora_people_moderation_error(err: object) -> bool:
    try:
        text = json.dumps(err, ensure_ascii=False).lower()
    except Exception:
        text = str(err).lower()
    return (
        "people-in-user-uploads" in text
        or "blocked by our moderation system" in text
        or ("moderation system" in text and "sora" in text)
        or ("request is blocked" in text and "people" in text)
    )

def _extract_first_url(obj) -> str | None:
    if isinstance(obj, str):
        if obj.startswith("http://") or obj.startswith("https://"):
            return obj
        return None
    if isinstance(obj, dict):
        preferred = ("video", "video_url", "output_url", "url", "download_url", "file", "asset_url")
        for k in preferred:
            if k in obj:
                found = _extract_first_url(obj.get(k))
                if found:
                    return found
        for v in obj.values():
            found = _extract_first_url(v)
            if found:
                return found
    if isinstance(obj, (list, tuple)):
        for item in obj:
            found = _extract_first_url(item)
            if found:
                return found
    return None



def _cleanup_sent_video_keys():
    now = time.time()
    stale = [k for k, ts in _SENT_VIDEO_KEYS.items() if (now - ts) > VIDEO_RESULT_DEDUPE_TTL_S]
    for k in stale:
        _SENT_VIDEO_KEYS.pop(k, None)


def _mark_video_sent_once(key: str) -> bool:
    if not key:
        return False
    _cleanup_sent_video_keys()
    if key in _SENT_VIDEO_KEYS:
        return True
    _SENT_VIDEO_KEYS[key] = time.time()
    return False


def _video_result_key(chat_id: int | str, task_id: str = "", url: str = "", content: bytes | None = None) -> str:
    base = f"{chat_id}|{task_id or ''}|{url or ''}"
    if content:
        try:
            digest = hashlib.sha1(content).hexdigest()
        except Exception:
            digest = ""
        base += f"|{digest}"
    return base


def _compress_video_for_telegram_sync(video_bytes: bytes, max_mb: int = 48) -> bytes | None:
    """Re-encode a provider MP4 to a Telegram-safe document/video size.
    Used only as a fallback when Telegram rejects the original file or it is too large.
    """
    if not video_bytes:
        return None
    max_bytes = max(5, int(max_mb or 48)) * 1024 * 1024
    try:
        ffmpeg = _ffmpeg_exe()
        with tempfile.TemporaryDirectory() as td:
            src = os.path.join(td, "input.mp4")
            out = os.path.join(td, "tg_safe.mp4")
            with open(src, "wb") as f:
                f.write(video_bytes)
            cmd = [
                ffmpeg, "-y", "-hide_banner", "-loglevel", "error",
                "-i", src,
                "-vf", "scale='min(720,iw)':-2,fps=24",
                "-c:v", "libx264", "-preset", "ultrafast", "-crf", "34", "-pix_fmt", "yuv420p",
                "-c:a", "aac", "-b:a", "96k",
                "-movflags", "+faststart",
                out,
            ]
            res = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, timeout=180)
            if res.returncode != 0:
                log.warning("telegram video compress failed rc=%s err=%s", res.returncode, res.stderr.decode("utf-8", "ignore")[-500:])
                return None
            if os.path.exists(out) and os.path.getsize(out) > 1024:
                with open(out, "rb") as f:
                    data = f.read()
                if len(data) <= max_bytes:
                    return data
                log.warning("telegram video compress too large: %s > %s", len(data), max_bytes)
    except Exception as e:
        log.warning("telegram video compress exception: %s", e)
    return None

async def _compress_video_for_telegram(video_bytes: bytes, max_mb: int = 48) -> bytes | None:
    return await asyncio.to_thread(_compress_video_for_telegram_sync, video_bytes, max_mb)

async def _reply_video_from_url(update: Update, client: httpx.AsyncClient, url: str, caption: str, task_id: str = ""):
    """
    РһСӮРҝСҖР°РІР»СҸРөСӮ РһР”РҳРқ СҖРөР·СғР»СҢСӮР°СӮ РІ Telegram.
    РҹРҫ СғРјРҫР»СҮР°РҪРёСҺ вҖ” MP4 РәР°Рә document, СҮСӮРҫРұСӢ Telegram РҪРө РјР°СҖРәРёСҖРҫРІР°Р» РәРҫСҖРҫСӮРәРёР№ СҖРҫР»РёРә РәР°Рә GIF.
    """
    headers = {
        "User-Agent": "Mozilla/5.0 (compatible; GPT5ProBot/1.0)",
        "Accept": "video/mp4,video/*,*/*;q=0.8",
    }

    downloaded: bytes | None = None
    try:
        r = await client.get(url, headers=headers, timeout=240.0, follow_redirects=True)
        r.raise_for_status()
        content_type = (r.headers.get("content-type") or "").lower()
        if not r.content or len(r.content) < 512:
            raise RuntimeError(f"empty video response: {len(r.content)} bytes")
        if "text/html" in content_type or "application/json" in content_type:
            raise RuntimeError(f"not a video response: {content_type}; {r.text[:300]}")
        downloaded = r.content
    except Exception as e:
        log.warning("reply_video_from_url: local download failed: %s", e)

    chat_id = getattr(getattr(update, "effective_chat", None), "id", "na")
    dedupe_key = _video_result_key(chat_id, task_id=task_id, url=url, content=downloaded)
    if _mark_video_sent_once(dedupe_key):
        log.info("reply_video_from_url: duplicate suppressed task_id=%s", task_id)
        return

    if downloaded:
        # v70: for Telegram rejection/size issues, try a compact MP4 before falling back to a raw link.
        if TELEGRAM_VIDEO_COMPRESS_ON_FAIL and len(downloaded) > max(5, int(TELEGRAM_RESULT_MAX_MB or 48)) * 1024 * 1024:
            compact = await _compress_video_for_telegram(downloaded, TELEGRAM_RESULT_MAX_MB)
            if compact:
                downloaded = compact
        if VIDEO_RESULT_SEND_AS_DOCUMENT:
            try:
                bio = BytesIO(downloaded)
                bio.name = "result.mp4"
                await update.effective_message.reply_document(document=InputFile(bio), caption=caption)
                return
            except Exception as e:
                log.warning("reply_video_from_url: document send failed: %s", e)
        try:
            bio = BytesIO(downloaded)
            bio.name = "result.mp4"
            await update.effective_message.reply_video(video=InputFile(bio), caption=caption, supports_streaming=True)
            return
        except Exception as e:
            log.warning("reply_video_from_url: video send failed: %s", e)
        try:
            bio = BytesIO(downloaded)
            bio.name = "result.mp4"
            await update.effective_message.reply_document(document=InputFile(bio), caption=caption)
            return
        except Exception as e:
            log.warning("reply_video_from_url: document send fallback failed: %s", e)
        if TELEGRAM_VIDEO_COMPRESS_ON_FAIL:
            compact = await _compress_video_for_telegram(downloaded, TELEGRAM_RESULT_MAX_MB)
            if compact and compact != downloaded:
                try:
                    bio = BytesIO(compact)
                    bio.name = "result_tg_safe.mp4"
                    await update.effective_message.reply_document(document=InputFile(bio), caption=caption + "\nрҹ“Ұ Р’РёРҙРөРҫ СҒР¶Р°СӮРҫ РҙР»СҸ РҫСӮРҝСҖР°РІРәРё РІ Telegram.")
                    return
                except Exception as e:
                    log.warning("reply_video_from_url: compressed document send failed: %s", e)

    if not VIDEO_RESULT_SEND_AS_DOCUMENT:
        try:
            await update.effective_message.reply_video(video=url, caption=caption, supports_streaming=True)
            return
        except Exception as e:
            log.warning("reply_video_from_url: telegram URL video send failed: %s", e)
    try:
        await update.effective_message.reply_document(document=url, caption=caption)
        return
    except Exception as e:
        log.warning("reply_video_from_url: telegram URL document send failed: %s", e)

    safe_url = (url or "")[:3500]
    await update.effective_message.reply_text(
        f"{caption}\nвҡ пёҸ Telegram РҪРө РҝСҖРёРҪСҸР» РІРёРҙРөРҫС„Р°Р№Р» РҪР°РҝСҖСҸРјСғСҺ, РҫСҒСӮР°РІР»СҸСҺ СҒСҒСӢР»РәСғ:\n{safe_url}",
        disable_web_page_preview=False,
    )

async def _reply_video_bytes(update: Update, content: bytes, caption: str, task_id: str = ""):
    if not content or len(content) < 512:
        raise RuntimeError(f"empty video bytes: {len(content or b'')} bytes")
    chat_id = getattr(getattr(update, "effective_chat", None), "id", "na")
    dedupe_key = _video_result_key(chat_id, task_id=task_id, content=content)
    if _mark_video_sent_once(dedupe_key):
        log.info("reply_video_bytes: duplicate suppressed task_id=%s", task_id)
        return
    bio = BytesIO(content)
    bio.name = "result.mp4"
    if VIDEO_RESULT_SEND_AS_DOCUMENT:
        await update.effective_message.reply_document(document=InputFile(bio), caption=caption)
        return
    await update.effective_message.reply_video(video=InputFile(bio), caption=caption, supports_streaming=True)

def _ratio_for_aspect(aspect: str) -> str:
    """
    Р”Р»СҸ Runway API version 2024-11-06 ratio РҙРҫР»Р¶РөРҪ РұСӢСӮСҢ СҖР°Р·СҖРөСҲРөРҪРёРөРј,
    Р° РҪРө СҒСӮСҖРҫРәРҫР№ 9:16 / 16:9.
    """
    mapping = {
        "9:16": "768:1280",
        "16:9": "1280:768",
        "1:1": "960:960",
        "4:5": "768:960",
        "3:4": "768:1024",
        "4:3": "1024:768",
    }
    return mapping.get((aspect or "").strip(), "768:1280")

def _duration_for_engine(engine: str, duration_s: int) -> int:
    try:
        d = int(duration_s or 5)
    except Exception:
        d = 5
    engine = (engine or "").lower()
    if engine == "runway":
        return max(2, min(10, d))
    if engine == "kling":
        return 10 if d >= 7 else 5
    if engine == "sora":
        # Sora/Comet СҒСӮР°РұРёР»СҢРҪРөРө РҝСҖРёРҪРёРјР°РөСӮ seconds = 4/8/12.
        # 10 СҒРөРәСғРҪРҙ РёР· UI РҪРҫСҖРјР°Р»РёР·СғРөРј РІ РұР»РёР¶Р°Р№СҲРёР№ РҝРҫРҙРҙРөСҖР¶РёРІР°РөРјСӢР№ РІР°СҖРёР°РҪСӮ вҖ” 8,
        # РҙР»РёРҪРҪСӢРө Р·Р°РҝСҖРҫСҒСӢ вҖ” 12.
        if d <= 5:
            return 4
        if d <= 10:
            return 8
        return 12
    if engine == "luma":
        return 9 if d >= 7 else 5
    return max(5, min(15, d))

def _guess_aspect_from_image(img_bytes: bytes, fallback: str = "9:16") -> str:
    if Image is None:
        return fallback
    try:
        im = Image.open(BytesIO(img_bytes))
        w, h = im.size
        if h > w * 1.2:
            return "9:16"
        if w > h * 1.2:
            return "16:9"
        return "1:1"
    except Exception:
        return fallback

def _image_refs_for_i2v(update: Update, img_bytes: bytes) -> tuple[str, str]:
    """
    Р’РҫР·РІСҖР°СүР°РөРј СҒРҪР°СҮР°Р»Р° data_url, РҝРҫСӮРҫРј Telegram URL.
    Р”Р»СҸ Comet / Runway / Kling РұРөР·РҫРҝР°СҒРҪРөРө РҝРөСҖРІСӢРј РҝСҖРҫРұРҫРІР°СӮСҢ base64 data-url,
    РҝРҫСӮРҫРјСғ СҮСӮРҫ РІРҪРөСҲРҪРёРө API СҮР°СҒСӮРҫ РҪРө РјРҫРіСғСӮ РәРҫСҖСҖРөРәСӮРҪРҫ Р·Р°РұСҖР°СӮСҢ Telegram file_path.
    """
    data_url = (
        f"data:{sniff_image_mime(img_bytes)};base64,"
        f"{base64.b64encode(img_bytes).decode('ascii')}"
    )

    tg_url = ""
    try:
        tg_url = _get_cached_photo_url(update.effective_user.id)
    except Exception:
        tg_url = ""

    return data_url, tg_url

def _sora_size_for_aspect(aspect: str) -> tuple[str, int, int]:
    # Sora Videos API РҝСҖРёРҪРёРјР°РөСӮ РҪРө 9:16/16:9, Р° size.
    # Р”Р»СҸ СҒСӮР°РҪРҙР°СҖСӮРҪРҫРіРҫ sora-2 СҒСӮР°РұРёР»СҢРҪСӢРө СҖР°Р·РјРөСҖСӢ: 720x1280 РёР»Рё 1280x720.
    a = (aspect or "").strip()
    if a == "16:9":
        return "1280x720", 1280, 720
    return "720x1280", 720, 1280

def _prepare_sora_reference_image(img_bytes: bytes, aspect: str) -> tuple[bytes, str, str, str]:
    """
    Р“РҫСӮРҫРІРёСӮ РёР·РҫРұСҖР°Р¶РөРҪРёРө РҙР»СҸ Sora imageвҶ’video.
    Р’Р°Р¶РҪРҫ: РҫС„РёСҶРёР°Р»СҢРҪСӢР№ Videos API СӮСҖРөРұСғРөСӮ, СҮСӮРҫРұСӢ reference image СҒРҫРІРҝР°РҙР°Р»
    СҒ СҶРөР»РөРІСӢРј СҖР°Р·РјРөСҖРҫРј video size. РҹРҫСҚСӮРҫРјСғ РҙРөР»Р°РөРј center-crop + resize.
    Р’РҫР·РІСҖР°СүР°РөСӮ: (bytes, mime, data_url, size).
    """
    size, tw, th = _sora_size_for_aspect(aspect)

    if Image is None:
        mime = sniff_image_mime(img_bytes)
        data_url = f"data:{mime};base64,{base64.b64encode(img_bytes).decode('ascii')}"
        return img_bytes, mime, data_url, size

    try:
        im = Image.open(BytesIO(img_bytes))
        try:
            im = ImageOps.exif_transpose(im)
        except Exception:
            pass
        im = im.convert("RGB")
        w, h = im.size
        target_ratio = tw / th
        cur_ratio = w / max(1, h)

        if cur_ratio > target_ratio:
            # РЎР»РёСҲРәРҫРј СҲРёСҖРҫРәРҫРө вҖ” СҖРөР¶РөРј РәСҖР°СҸ.
            new_w = int(h * target_ratio)
            left = max(0, (w - new_w) // 2)
            im = im.crop((left, 0, left + new_w, h))
        elif cur_ratio < target_ratio:
            # РЎР»РёСҲРәРҫРј РІСӢСҒРҫРәРҫРө вҖ” СҖРөР¶РөРј РІРөСҖС…/РҪРёР·.
            new_h = int(w / target_ratio)
            top = max(0, (h - new_h) // 2)
            im = im.crop((0, top, w, top + new_h))

        resample = getattr(Image, "Resampling", Image).LANCZOS
        im = im.resize((tw, th), resample)
        out = BytesIO()
        im.save(out, format="JPEG", quality=92, optimize=True)
        prepared = out.getvalue()
        mime = "image/jpeg"
        data_url = f"data:{mime};base64,{base64.b64encode(prepared).decode('ascii')}"
        return prepared, mime, data_url, size
    except Exception as e:
        log.warning("Sora image prepare failed, using original bytes: %s", e)
        mime = sniff_image_mime(img_bytes)
        data_url = f"data:{mime};base64,{base64.b64encode(img_bytes).decode('ascii')}"
        return img_bytes, mime, data_url, size

async def _start_photo_revival(update: Update, context: ContextTypes.DEFAULT_TYPE, engine: str, img_bytes: bytes, prompt: str = ""):
    engine = (engine or "runway").lower().strip()
    if engine == "luma" and LUMA_TEMP_DISABLED:
        await update.effective_message.reply_text("вҡ пёҸ Luma РІСҖРөРјРөРҪРҪРҫ РҫСӮРәР»СҺСҮРөРҪР° Рё СҒРәСҖСӢСӮР° РёР· РјРөРҪСҺ. РҳСҒРҝРҫР»СҢР·СғР№СӮРө Runway, Kling РёР»Рё Sora 2 РұРөР· Р»СҺРҙРөР№.")
        return
    prompt = (prompt or "subtle lifelike animation, natural micro-movements, smooth cinematic camera motion").strip()
    dur, asp = parse_video_opts(prompt)
    if not re.search(r"(?:9:16|16:9|1:1|4:5|3:4|4:3)", prompt or "", re.I):
        asp = _guess_aspect_from_image(img_bytes, asp)
    dur = _duration_for_engine(engine, dur)

    pay_engine = "runway" if engine in ("runway", "kling", "sora") else "luma"
    est = _video_provider_cost_usd(engine, dur) if engine in ("runway", "kling", "sora") else 0.40

    async def _go():
        await update.effective_message.reply_text(
            f"вң… Р—Р°РҝСғСҒРәР°СҺ РҫР¶РёРІР»РөРҪРёРө С„РҫСӮРҫ: {engine.upper()} вҖў {dur} СҒРөРә вҖў {asp}."
        )
        if engine == "runway":
            return bool(await _run_runway_animate_photo(update, context, img_bytes, prompt=prompt, duration_s=dur, aspect=asp))
        if engine == "luma":
            return bool(await _run_luma_animate_photo(update, context, img_bytes, prompt=prompt, duration_s=dur, aspect=asp))
        if engine in ("sora", "kling"):
            return bool(await _run_comet_i2v(update, context, engine, img_bytes, prompt=prompt, duration_s=dur, aspect=asp))
        await update.effective_message.reply_text("вқҢ РқРөРёР·РІРөСҒСӮРҪСӢР№ РҙРІРёР¶РҫРә РҫР¶РёРІР»РөРҪРёСҸ С„РҫСӮРҫ.")
        return False

    await _try_pay_then_do(
        update, context, update.effective_user.id, pay_engine, est, _go,
        remember_kind=f"revive_photo_{engine}",
        remember_payload={"engine": engine, "duration": dur, "aspect": asp, "prompt": prompt},
    )

async def _poll_video_task_generic(
    update: Update,
    client: httpx.AsyncClient,
    headers: dict,
    base_url: str,
    status_paths: list[str],
    task_id: str,
    caption: str,
    max_wait_s: int = 1200,
    task_not_exist_soft_fail_s: int = 0,
    silent_soft_fail: bool = False,
) -> bool:
    """
    РЈРҪРёРІРөСҖСҒР°Р»СҢРҪСӢР№ polling РҙР»СҸ async-video Р·Р°РҙР°СҮ.

    Р’Р°Р¶РҪРҫ РҙР»СҸ Comet/Runway: РҫСӮРІРөСӮ task_not_exist РјРҫР¶РөСӮ РҝСҖРёС…РҫРҙРёСӮСҢ РҪРө РәР°Рә С„РёРҪР°Р»СҢРҪР°СҸ РҫСҲРёРұРәР°,
    Р° РәР°Рә СҒСӮР°РҙРёСҸ РҝРөСҖРІРёСҮРҪРҫР№ РёРҪРёСҶРёР°Р»РёР·Р°СҶРёРё Р·Р°РҙР°СҮРё. РҹРҫСҚСӮРҫРјСғ РјСӢ РҪРө СҒСҮРёСӮР°РөРј РөРіРҫ РјРіРҪРҫРІРөРҪРҪСӢРј
    РҝСҖРҫРІР°Р»РҫРј. РқРҫ РөСҒР»Рё РҫРҪ РҙРөСҖР¶РёСӮСҒСҸ РҙРҫР»СҢСҲРө task_not_exist_soft_fail_s, РІРҫР·РІСҖР°СүР°РөРј False,
    СҮСӮРҫРұСӢ РІРөСҖС…РҪРёР№ СғСҖРҫРІРөРҪСҢ РјРҫРі РҝРөСҖРөРәР»СҺСҮРёСӮСҢСҒСҸ РҪР° РҙСҖСғРіРҫР№ РҙРІРёР¶РҫРә/РјРҫРҙРөР»СҢ.
    """
    started = time.time()
    task_not_exist_seen_at: float | None = None

    while True:
        last_body = ""
        soft_not_exist_seen_this_round = False

        for path in status_paths:
            url = f"{base_url}{path}".format(id=task_id)
            try:
                rs = await client.get(url, headers=headers, timeout=60.0)
                body_preview = _api_error_preview(rs)

                if rs.status_code >= 400:
                    last_body = f"{rs.status_code}: {body_preview}"

                    # CometAPI/Runway soft-state: task created, but status storage is not ready yet.
                    if "task_not_exist" in (body_preview or "").lower():
                        soft_not_exist_seen_this_round = True
                        if task_not_exist_seen_at is None:
                            task_not_exist_seen_at = time.time()
                        if task_not_exist_soft_fail_s and (time.time() - task_not_exist_seen_at) >= task_not_exist_soft_fail_s:
                            log.warning(
                                "%s: task_not_exist persisted %.1fs for task_id=%s; soft fallback",
                                caption, time.time() - task_not_exist_seen_at, task_id,
                            )
                            if not silent_soft_fail:
                                await update.effective_message.reply_text(
                                    f"вҡ пёҸ {caption}: Р·Р°РҙР°СҮР° СҒР»РёСҲРәРҫРј РҙРҫР»РіРҫ РҪРө РҝРҫСҸРІР»СҸРөСӮСҒСҸ РІ Comet/Runway. РҹРөСҖРөРәР»СҺСҮР°СҺСҒСҢ РҪР° СҖРөР·РөСҖРІРҪСӢР№ РҝСғСӮСҢ."
                                )
                            return False
                        continue

                    continue

                try:
                    js = rs.json() or {}
                except Exception:
                    js = {}

            except Exception as e:
                last_body = str(e)
                continue

            st = str(js.get("status") or js.get("state") or js.get("task_status") or "").lower()

            # Comet/Runway РёРҪРҫРіРҙР° РҫСӮРҙР°С‘СӮ task_not_exist РІРҪСғСӮСҖРё JSON РҝСҖРё 200 OK.
            if st == "task_not_exist" or "task_not_exist" in json.dumps(js, ensure_ascii=False).lower():
                soft_not_exist_seen_this_round = True
                last_body = json.dumps(js, ensure_ascii=False)[:700]
                if task_not_exist_seen_at is None:
                    task_not_exist_seen_at = time.time()
                if task_not_exist_soft_fail_s and (time.time() - task_not_exist_seen_at) >= task_not_exist_soft_fail_s:
                    log.warning(
                        "%s: task_not_exist JSON persisted %.1fs for task_id=%s; soft fallback",
                        caption, time.time() - task_not_exist_seen_at, task_id,
                    )
                    if not silent_soft_fail:
                        await update.effective_message.reply_text(
                            f"вҡ пёҸ {caption}: Р·Р°РҙР°СҮР° СҒР»РёСҲРәРҫРј РҙРҫР»РіРҫ РҪРө РҝРҫСҸРІР»СҸРөСӮСҒСҸ РІ Comet/Runway. РҹРөСҖРөРәР»СҺСҮР°СҺСҒСҢ РҪР° СҖРөР·РөСҖРІРҪСӢР№ РҝСғСӮСҢ."
                        )
                    return False
                continue

            url = _extract_first_url(js.get("output")) or _extract_first_url(js.get("assets")) or _extract_first_url(js)
            if st in ("completed", "succeeded", "success", "finished", "ready", "done", "succeed") or (url and not st):
                if not url:
                    # OpenAI/Sora Videos API СҮР°СҒСӮРҫ РІРҫР·РІСҖР°СүР°РөСӮ completed РұРөР· URL.
                    # РӨРёРҪР°Р»СҢРҪСӢР№ MP4 РҪР°РҙРҫ Р·Р°РұСҖР°СӮСҢ РҫСӮРҙРөР»СҢРҪСӢРј GET /v1/videos/{id}/content.
                    if "sora" in (caption or "").lower() or "/v1/videos" in " ".join(status_paths):
                        try:
                            content_url = f"{base_url.rstrip("/")}/v1/videos/{task_id}/content"
                            cr = await client.get(content_url, headers=headers, timeout=240.0, follow_redirects=True)
                            if cr.status_code < 400 and cr.content and "application/json" not in (cr.headers.get("content-type") or "").lower():
                                await _reply_video_bytes(update, cr.content, f"{caption} вң…", task_id=task_id)
                                return True
                            log.warning("%s content download failed: %s %s", caption, cr.status_code, _api_error_preview(cr))
                        except Exception as e:
                            log.warning("%s content download exception: %s", caption, e)
                    await update.effective_message.reply_text(f"вҡ пёҸ {caption}: Р·Р°РҙР°СҮР° РіРҫСӮРҫРІР°, РҪРҫ СҒСҒСӢР»РәР°/MP4 РҪР° РІРёРҙРөРҫ РҪРө РҪР°Р№РҙРөРҪСӢ.")
                    return True
                await _reply_video_from_url(update, client, url, f"{caption} вң…", task_id=task_id)
                if "runway" in (caption or "").lower():
                    _provider_mark_success("runway_i2v")
                return True
            if st in ("failed", "fail", "error", "canceled", "cancelled", "rejected"):
                if "sora" in (caption or "").lower() and is_sora_people_moderation_error(js):
                    await update.effective_message.reply_text(_sora_people_moderation_text())
                    return True
                if "runway" in (caption or "").lower() and RUNWAY_HIDE_TECH_ERRORS:
                    _provider_mark_failure("runway_i2v", json.dumps(js, ensure_ascii=False)[:700])
                    return False
                await update.effective_message.reply_text(f"вқҢ {caption}: РҫСҲРёРұРәР° СҖРөРҪРҙРөСҖР°.\n{json.dumps(js, ensure_ascii=False)[:900]}")
                return True

        if time.time() - started > max_wait_s:
            # Р”Р»СҸ Runway/Comet timeout РҙРҫР»Р¶РөРҪ РҙР°СӮСҢ СҲР°РҪСҒ РІРөСҖС…РҪРөРјСғ fallback-СғСҖРҫРІРҪСҺ.
            if task_not_exist_seen_at is not None and silent_soft_fail:
                log.warning("%s: timeout with task_not_exist for task_id=%s; soft fallback", caption, task_id)
                return False
            if "runway" in (caption or "").lower() and RUNWAY_HIDE_TECH_ERRORS:
                _provider_mark_failure("runway_i2v", last_body[:700])
                return False
            await update.effective_message.reply_text(f"вҢӣ {caption}: РІСҖРөРјСҸ РҫР¶РёРҙР°РҪРёСҸ РІСӢСҲР»Рҫ. РҹРҫСҒР»РөРҙРҪРёР№ РҫСӮРІРөСӮ: {last_body[:500]}")
            return False

        # Р•СҒР»Рё РІСҒРө РҝСғСӮРё РҙР°Р»Рё СӮРҫР»СҢРәРҫ РјСҸРіРәРёР№ task_not_exist вҖ” РҝСҖРҫСҒСӮРҫ Р¶РҙС‘Рј СҒР»РөРҙСғСҺСүРёР№ СҶРёРәР».
        await asyncio.sleep(VIDEO_POLL_DELAY_S)

async def _create_and_poll_i2v(
    update: Update,
    base_url: str,
    api_key: str,
    create_payloads: list[tuple[str, dict]],
    status_paths: list[str],
    caption: str,
    task_not_exist_soft_fail_s: int = 0,
    silent_soft_fail: bool = False,
    max_wait_s: int | None = None,
) -> bool:
    if not api_key:
        await update.effective_message.reply_text(f"вқҢ {caption}: API-РәР»СҺСҮ РҪРө Р·Р°РҙР°РҪ РІ ENV.")
        return True

    auth_headers = {
        "Authorization": f"Bearer {api_key}",
        "Accept": "application/json",
    }

    last_err = ""
    all_errors: list[str] = []

    async with httpx.AsyncClient(timeout=90.0) as client:
        for path, payload in create_payloads:
            try:
                headers = dict(auth_headers)

                # Runway СҮРөСҖРөР· Comet СӮСҖРөРұСғРөСӮ РІРөСҖСҒРёСҺ API.
                if str(path).startswith("/runwayml/"):
                    headers["X-Runway-Version"] = RUNWAY_API_VERSION or "2024-11-06"

                # РЎРҝРөСҶ-СҖРөР¶РёРј РҙР»СҸ Sora/OpenAI Videos API: input_reference РәР°Рә С„Р°Р№Р»
                # РҙРҫР»Р¶РөРҪ СғС…РҫРҙРёСӮСҢ multipart/form-data, Р° РҪРө JSON. Р’ СҚСӮРҫРј СҖРөР¶РёРјРө
                # Content-Type РҪРө СҒСӮР°РІРёРј РІСҖСғСҮРҪСғСҺ вҖ” httpx СҒР°Рј РҙРҫРұР°РІРёСӮ boundary.
                if isinstance(payload, dict) and payload.get("__multipart"):
                    mp = payload.get("__multipart") or {}
                    data = mp.get("data") or {}
                    files = mp.get("files") or {}
                    r = await client.post(f"{base_url}{path}", headers=headers, data=data, files=files)
                else:
                    headers["Content-Type"] = "application/json"
                    r = await client.post(f"{base_url}{path}", headers=headers, json=payload)

                if r.status_code >= 400:
                    mode = "multipart" if isinstance(payload, dict) and payload.get("__multipart") else "json"
                    last_err = f"POST {path} [{mode}] вҶ’ {r.status_code}: {_api_error_preview(r)}"
                    all_errors.append(last_err)
                    log.warning("%s create failed: %s", caption, last_err)
                    continue

                try:
                    js = r.json() or {}
                except Exception:
                    js = {}

                ready_url = (
                    _extract_first_url(js.get("output"))
                    or _extract_first_url(js.get("outputs"))
                    or _extract_first_url(js.get("assets"))
                    or _extract_first_url(js.get("data"))
                    or _extract_first_url(js.get("result"))
                    or _extract_first_url(js.get("response"))
                    or _extract_first_url(js.get("payload"))
                    or _extract_first_url(js)
                )

                if ready_url:
                    await _reply_video_from_url(update, client, ready_url, f"{caption} вң…")
                    return True

                task_id = str(
                    js.get("id")
                    or js.get("task_id")
                    or js.get("generation_id")
                    or js.get("video_id")
                    or js.get("taskId")
                    or js.get("taskID")
                    or js.get("request_id")
                    or js.get("uuid")
                    or ""
                ).strip()

                if not task_id and isinstance(js.get("data"), dict):
                    d = js.get("data") or {}
                    task_id = str(
                        d.get("id")
                        or d.get("task_id")
                        or d.get("generation_id")
                        or d.get("video_id")
                        or d.get("taskId")
                        or d.get("taskID")
                        or d.get("request_id")
                        or d.get("uuid")
                        or ""
                    ).strip()

                if not task_id and isinstance(js.get("result"), dict):
                    d = js.get("result") or {}
                    task_id = str(
                        d.get("id")
                        or d.get("task_id")
                        or d.get("generation_id")
                        or d.get("video_id")
                        or d.get("taskId")
                        or d.get("taskID")
                        or d.get("request_id")
                        or d.get("uuid")
                        or ""
                    ).strip()

                if not task_id:
                    last_err = f"POST {path}: РҪРөСӮ id Р·Р°РҙР°СҮРё РІ РҫСӮРІРөСӮРө {json.dumps(js, ensure_ascii=False)[:700]}"
                    all_errors.append(last_err)
                    continue

                await update.effective_message.reply_text(f"вҸі {caption}: Р·Р°РҙР°СҮР° РҝСҖРёРҪСҸСӮР°, РҫР¶РёРҙР°СҺ СҖРөР·СғР»СҢСӮР°СӮвҖҰ")
                log.info("%s accepted: path=%s task_id=%s response=%s", caption, path, task_id, json.dumps(js, ensure_ascii=False)[:1200])

                return await _poll_video_task_generic(
                    update,
                    client,
                    headers,
                    base_url,
                    status_paths,
                    task_id,
                    caption,
                    max_wait_s=int(max_wait_s or max(LUMA_MAX_WAIT_S, RUNWAY_MAX_WAIT_S)),
                    task_not_exist_soft_fail_s=task_not_exist_soft_fail_s,
                    silent_soft_fail=silent_soft_fail,
                )

            except Exception as e:
                last_err = f"POST {path}: {e}"
                all_errors.append(last_err)
                log.warning("%s create exception: %s", caption, e)
                continue

    if all_errors:
        details = "\n".join(all_errors[-5:])
    else:
        details = last_err

    if "runway" in (caption or "").lower():
        if _is_runway_unavailable_text(details):
            _provider_mark_failure("runway_i2v", details)
        if silent_soft_fail or RUNWAY_HIDE_TECH_ERRORS:
            log.warning("%s hidden create failure: %s", caption, details[:1500])
            return False

    if "sora" in (caption or "").lower() and is_sora_people_moderation_error(details):
        await update.effective_message.reply_text(_sora_people_moderation_text())
        return False
    if "invalid api channeltype" in (details or "").lower():
        await update.effective_message.reply_text(
            f"вҡ пёҸ {caption}: Сғ СӮРөРәСғСүРөРіРҫ РҝСҖРҫРІР°Р№РҙРөСҖР°/РәР°РҪР°Р»Р° Sora СҒРөР№СҮР°СҒ РҪРөРҙРҫСҒСӮСғРҝРҪР° (invalid api channelType)."
        )
        return False
    await update.effective_message.reply_text(f"вқҢ {caption}: РҪРө СғРҙР°Р»РҫСҒСҢ СҒРҫР·РҙР°СӮСҢ Р·Р°РҙР°СҮСғ.\n{details[:900]}")
    return False

async def _run_luma_animate_photo(update: Update, context: ContextTypes.DEFAULT_TYPE, img_bytes: bytes, prompt: str, duration_s: int, aspect: str):
    if not LUMA_API_KEY:
        await update.effective_message.reply_text("вқҢ Luma: LUMA_API_KEY РҪРө Р·Р°РҙР°РҪ РІ ENV.")
        return False
    data_url, tg_url = _image_refs_for_i2v(update, img_bytes)
    image_ref = data_url or tg_url
    duration_s = _duration_for_engine("luma", duration_s)
    async with httpx.AsyncClient(timeout=60.0) as client:
        base = await _pick_luma_base(client)
    payloads = [
        (LUMA_CREATE_PATH, {
            "model": LUMA_MODEL,
            "prompt": prompt,
            "duration": f"{duration_s}s",
            "aspect_ratio": aspect,
            "keyframes": {"frame0": {"type": "image", "url": image_ref}},
        }),
        (LUMA_CREATE_PATH, {
            "model": LUMA_MODEL,
            "prompt": prompt,
            "duration": f"{duration_s}s",
            "aspect_ratio": aspect,
            "image_ref": image_ref,
        }),
    ]
    return bool(await _create_and_poll_i2v(update, base, LUMA_API_KEY, payloads, [LUMA_STATUS_PATH], "Luma imageвҶ’video"))

async def _run_comet_i2v(update: Update, context: ContextTypes.DEFAULT_TYPE, engine: str, img_bytes: bytes, prompt: str, duration_s: int, aspect: str):
    engine = (engine or "").lower()

    # Р”Р»СҸ Kling/Runway/Sora РҪРө РҫСӮРҝСҖР°РІР»СҸРөРј РҫРіСҖРҫРјРҪСӢРө СҒРәСҖРёРҪСҲРҫСӮСӢ РәР°Рә РөСҒСӮСҢ: РҪРҫСҖРјР°Р»РёР·СғРөРј JPEG Рё
    # РјСҸРіРәРҫ СғРұРёСҖР°РөРј СҮС‘СҖРҪСӢРө РҝРҫР»СҸ, РөСҒР»Рё СҚСӮРҫ РұРөР·РҫРҝР°СҒРҪРҫ.
    prepared_note = ""
    if engine in ("kling", "runway", "sora"):
        img_bytes, prepared_note = _prepare_i2v_source_image(img_bytes, aspect)
        if prepared_note:
            log.info("i2v source prepared for %s: %s", engine, prepared_note)

    data_url, tg_url = _image_refs_for_i2v(update, img_bytes)
    raw_b64 = base64.b64encode(img_bytes).decode("ascii")

    if engine == "sora":
        d = _duration_for_engine("sora", duration_s)

        # РҹРҫ С„Р°РәСӮСғ РІР°СҲРёС… СӮРөСҒСӮРҫРІ Comet/OpenAI /v1/videos РқР• РҝСҖРёРҪРёРјР°РөСӮ:
        #   - top-level image_url / image_urls
        #   - input_image
        #   - duration
        #   - aspect_ratio
        # РҹРҫ С„Р°РәСӮСғ Comet/OpenAI proxy РІ РІР°СҲРёС… Р»РҫРіР°С… РҫР¶РёРҙР°РөСӮ input_reference РәР°Рә РЎРўР РһРҡРЈ,
        # Р° РҪРө РәР°Рә РҫРұСҠРөРәСӮ. РҹРҫСҚСӮРҫРјСғ РҝСҖРҫРұСғРөРј РІ РҝРөСҖРІСғСҺ РҫСҮРөСҖРөРҙСҢ string data-url,
        # Р·Р°СӮРөРј string РұРөР· seconds, Р·Р°СӮРөРј multipart-С„Р°Р№Р» РәР°Рә Р·Р°РҝР°СҒРҪРҫР№ РІР°СҖРёР°РҪСӮ.
        sora_bytes, sora_mime, sora_data_url, size = _prepare_sora_reference_image(img_bytes, aspect)

        payloads = []

        def _add_json(payload: dict):
            for bad in ("input_image", "duration", "aspect_ratio", "image_url", "image_urls"):
                payload.pop(bad, None)
            payloads.append((SORA_CREATE_PATH, payload))

        def _add_multipart(seconds_value: str | None = None):
            data = {
                "model": SORA_MODEL,
                "prompt": prompt,
                "size": size,
            }
            if seconds_value:
                data["seconds"] = seconds_value
            files = {
                "input_reference": ("reference.jpg", sora_bytes, sora_mime or "image/jpeg"),
            }
            payloads.append((SORA_CREATE_PATH, {"__multipart": {"data": data, "files": files}}))

        # 1) РһСҒРҪРҫРІРҪРҫР№ РІР°СҖРёР°РҪСӮ: input_reference РәР°Рә string (data URL) + seconds + size.
        _add_json({
            "model": SORA_MODEL,
            "prompt": prompt,
            "input_reference": sora_data_url,
            "seconds": str(d),
            "size": size,
        })

        # 2) РўРҫ Р¶Рө РұРөР· seconds вҖ” РөСҒР»Рё РәР°РҪР°Р» СҒР°Рј РҪРҫСҖРјР°Р»РёР·СғРөСӮ РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ.
        _add_json({
            "model": SORA_MODEL,
            "prompt": prompt,
            "input_reference": sora_data_url,
            "size": size,
        })

        # 3) Р—Р°РҝР°СҒРҪРҫР№ РІР°СҖРёР°РҪСӮ СҮРөСҖРөР· РҝСғРұР»РёСҮРҪСӢР№ Telegram URL, РөСҒР»Рё proxy РҪРө Р»СҺРұРёСӮ data URL.
        if tg_url and tg_url.startswith("https://"):
            _add_json({
                "model": SORA_MODEL,
                "prompt": prompt,
                "input_reference": tg_url,
                "seconds": str(d),
                "size": size,
            })
            _add_json({
                "model": SORA_MODEL,
                "prompt": prompt,
                "input_reference": tg_url,
                "size": size,
            })

        # 4) Multipart-РІР°СҖРёР°РҪСӮ РәР°Рә Р·Р°РҝР°СҒРҪРҫР№ fallback.
        _add_multipart(str(d))
        _add_multipart(None)

        sora_ok = await _create_and_poll_i2v(
            update,
            COMET_BASE_URL,
            SORA_API_KEY,
            payloads,
            [SORA_STATUS_PATH, "/v1/videos/{id}", "/v1/tasks/{id}"],
            "Sora 2 imageвҶ’video (РұРөР· Р»СҺРҙРөР№)",
        )
        if (not sora_ok) and SORA_AUTO_FALLBACK_KLING:
            await update.effective_message.reply_text(
                "вҶӘпёҸ Sora СҒРөР№СҮР°СҒ РҪРөРҙРҫСҒСӮСғРҝРҪР° РёР»Рё РҪРө СҒРҫР·РҙР°Р»Р° Р·Р°РҙР°СҮСғ. РҹРөСҖРөРәР»СҺСҮР°СҺСҒСҢ РҪР° Kling imageвҶ’video РәР°Рә СҖРөР·РөСҖРІ."
            )
            return bool(await _run_comet_i2v(update, context, "kling", img_bytes, prompt, duration_s, aspect))
        return bool(sora_ok)

    if engine == "kling":
        # Kling СҮРөСҖРөР· Comet Р¶РҙС‘СӮ СҮРёСҒСӮСӢР№ base64 РұРөР· data:image/...;base64,
        # Рё duration СҒСӮСҖРҫРәРҫР№.
        d = str(_duration_for_engine("kling", duration_s))
        safe_prompt = (prompt or "").strip()
        if I2V_KLING_SAFE_PROMPT_SUFFIX and I2V_KLING_SAFE_PROMPT_SUFFIX.lower() not in safe_prompt.lower():
            safe_prompt = (safe_prompt + "; " + I2V_KLING_SAFE_PROMPT_SUFFIX).strip("; ")

        payloads = [
            (
                KLING_CREATE_PATH,
                {
                    "model": KLING_MODEL,
                    "prompt": safe_prompt,
                    "image": raw_b64,
                    "duration": d,
                    "aspect_ratio": aspect,
                },
            ),
            (
                "/kling/v1/videos/image2video",
                {
                    "model": KLING_MODEL,
                    "prompt": safe_prompt,
                    "image": raw_b64,
                    "duration": d,
                    "aspect_ratio": aspect,
                },
            ),
        ]

        return bool(await _create_and_poll_i2v(
            update,
            COMET_BASE_URL,
            KLING_API_KEY,
            payloads,
            ["/kling/v1/videos/image2video/{id}", KLING_STATUS_PATH, "/kling/v1/videos/{id}", "/v1/tasks/{id}", "/v1/videos/{id}"],
            "Kling imageвҶ’video",
        ))

    await update.effective_message.reply_text("вқҢ РқРөРёР·РІРөСҒСӮРҪСӢР№ Comet imageвҶ’video РҙРІРёР¶РҫРә.")
    return False


async def _run_comet_text_video(update: Update, context: ContextTypes.DEFAULT_TYPE, engine: str, prompt: str, duration_s: int, aspect: str) -> bool:
    """Text-to-video through CometAPI: Sora 2, Kling, or Runway."""
    engine = (engine or "").lower().strip(); prompt = (prompt or "").strip()
    if not prompt:
        await update.effective_message.reply_text("вқҢ РҹСғСҒСӮРҫР№ Р·Р°РҝСҖРҫСҒ РҙР»СҸ РІРёРҙРөРҫ.")
        return False
    if engine == "sora" and _prompt_likely_has_people(prompt):
        await update.effective_message.reply_text("вҡ пёҸ Sora 2 РҙРҫСҒСӮСғРҝРҪР° СӮРҫР»СҢРәРҫ РҙР»СҸ СҒСҶРөРҪ РұРөР· Р»СҺРҙРөР№. РҳСҒРҝРҫР»СҢР·СғР№СӮРө Kling РёР»Рё Runway.")
        return False
    await context.bot.send_chat_action(update.effective_chat.id, ChatAction.RECORD_VIDEO)
    if engine == "runway":
        return bool(await _run_runway_video(update, context, prompt, duration_s, aspect))
    if engine == "sora":
        d = _duration_for_engine("sora", duration_s); size, _, _ = _sora_size_for_aspect(aspect)
        payloads = [
            (SORA_CREATE_PATH, {"model": SORA_MODEL, "prompt": prompt, "seconds": str(d), "size": size}),
            (SORA_CREATE_PATH, {"model": SORA_MODEL, "prompt": prompt, "seconds": d, "size": size}),
            (SORA_CREATE_PATH, {"model": SORA_MODEL, "prompt": prompt, "size": size}),
        ]
        return bool(await _create_and_poll_i2v(update, COMET_BASE_URL, SORA_API_KEY, payloads, [SORA_STATUS_PATH, "/v1/videos/{id}", "/v1/tasks/{id}"], "Sora 2 textвҶ’video В· РұРөР· Р»СҺРҙРөР№"))
    if engine == "kling":
        d = str(_duration_for_engine("kling", duration_s))
        payloads = [
            (KLING_TEXT_CREATE_PATH, {"model": KLING_MODEL, "prompt": prompt, "duration": d, "aspect_ratio": aspect}),
            ("/kling/v1/videos/text2video", {"model": KLING_MODEL, "prompt": prompt, "duration": d, "aspect_ratio": aspect}),
            (KLING_TEXT_CREATE_PATH, {"prompt": prompt, "duration": d, "aspect_ratio": aspect}),
        ]
        return bool(await _create_and_poll_i2v(update, COMET_BASE_URL, KLING_API_KEY, payloads, [KLING_TEXT_STATUS_PATH, "/kling/v1/videos/text2video/{id}", "/kling/v1/videos/{id}", "/v1/tasks/{id}", "/v1/videos/{id}"], "Kling textвҶ’video"))
    await update.effective_message.reply_text("вқҢ РқРөРёР·РІРөСҒСӮРҪСӢР№ textвҶ’video РҙРІРёР¶РҫРә. Р”РҫСҒСӮСғРҝРҪСӢ Sora 2 РұРөР· Р»СҺРҙРөР№, Kling Рё Runway.")
    return False

async def _run_runway_comet_animate_photo(update: Update, context: ContextTypes.DEFAULT_TYPE, img_bytes: bytes, prompt: str, duration_s: int, aspect: str) -> bool:
    if not (RUNWAY_IMAGE2VIDEO_ENABLED and RUNWAY_USE_COMET and COMET_API_KEY):
        return False
    if not _provider_is_available("runway_i2v"):
        log.warning("Runway/Comet skipped: cooldown %ss", _provider_cooldown_left("runway_i2v"))
        return False

    img_bytes, prep_note = _prepare_i2v_source_image(img_bytes, aspect)
    if prep_note:
        log.info("Runway i2v source prepared: %s", prep_note)
    data_url, tg_url = _image_refs_for_i2v(update, img_bytes)

    duration = _duration_for_engine("runway", duration_s)
    ratio = _ratio_for_aspect(aspect)

    payloads = []
    # Р”Р»СҸ imageвҶ’video РҪР° Comet РҪРө РёСҒРҝРҫР»СҢР·СғРөРј gen4.5 РҝРөСҖРІСӢРј: Сғ РІР°СҒ РҫРҪ СҮР°СҒСӮРҫ РҫСӮРІРөСҮР°РөСӮ no available channel.
    # gen4_turbo/gen3a_turbo вҖ” РҪРҫСҖРјР°Р»СҢРҪСӢРө РәР°РҪРҙРёРҙР°СӮСӢ РҙР»СҸ РҫР¶РёРІР»РөРҪРёСҸ С„РҫСӮРҫ.
    for model in _runway_i2v_model_candidates():
        # РһСҒРҪРҫРІРҪРҫР№ С„РҫСҖРјР°СӮ Runway API 2024-11-06.
        payloads.append((RUNWAY_COMET_CREATE_PATH, {
            "model": model,
            "promptImage": data_url,
            "promptText": prompt,
            "duration": duration,
            "ratio": ratio,
            "watermark": False,
        }))
        # РӨРҫСҖРјР°СӮ promptImage РәР°Рә РјР°СҒСҒРёРІ СҒ position=first вҖ” СҚРәРІРёРІР°Р»РөРҪСӮРөРҪ string Рё СҒСӮР°РұРёР»СҢРҪРөРө РҪР° РҪРөРәРҫСӮРҫСҖСӢС… РҝСҖРҫРәСҒРё.
        payloads.append((RUNWAY_COMET_CREATE_PATH, {
            "model": model,
            "promptImage": [{"uri": data_url, "position": "first"}],
            "promptText": prompt,
            "duration": duration,
            "ratio": ratio,
            "watermark": False,
        }))
        # snake_case fallback РҙР»СҸ СҒРҫРІРјРөСҒСӮРёРјРҫСҒСӮРё СҒ СҖР°Р·РҪСӢРјРё РҝСҖРҫРәСҒРё.
        payloads.append((RUNWAY_COMET_CREATE_PATH, {
            "model": model,
            "prompt_image": data_url,
            "prompt_text": prompt,
            "duration": duration,
            "ratio": ratio,
            "watermark": False,
        }))

    # Р—Р°РҝР°СҒРҪРҫР№ РІР°СҖРёР°РҪСӮ СҮРөСҖРөР· Telegram URL, РөСҒР»Рё Comet РҪРө РҝСҖРёРјРөСӮ data-uri.
    if tg_url and tg_url.startswith("https://"):
        for model in _runway_i2v_model_candidates():
            payloads.append((RUNWAY_COMET_CREATE_PATH, {
                "model": model,
                "promptImage": tg_url,
                "promptText": prompt,
                "duration": duration,
                "ratio": ratio,
                "watermark": False,
            }))

    return await _create_and_poll_i2v(
        update,
        COMET_BASE_URL,
        COMET_API_KEY,
        payloads,
        [RUNWAY_COMET_STATUS_PATH, "/runwayml/v1/tasks/{id}", "/v1/tasks/{id}"],
        "Runway/Comet imageвҶ’video",
        task_not_exist_soft_fail_s=RUNWAY_TASK_NOT_EXIST_FALLBACK_S,
        silent_soft_fail=True,
    )

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ Runway: Р°РҪРёРјР°СҶРёСҸ Р·Р°РіСҖСғР¶РөРҪРҪРҫРіРҫ С„РҫСӮРҫ (imageвҶ’video) в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
async def _run_runway_direct_animate_photo(update: Update, context: ContextTypes.DEFAULT_TYPE, img_bytes: bytes, prompt: str, duration_s: int, aspect: str) -> bool:
    """Official Runway imageвҶ’video using ephemeral upload + task polling."""
    if not (RUNWAY_DIRECT_ENABLED and RUNWAY_API_KEY):
        return False

    img_bytes, prep_note = _prepare_i2v_source_image(img_bytes, aspect)
    if prep_note:
        log.info("Runway direct i2v source prepared: %s", prep_note)
    mime_type = sniff_image_mime(img_bytes)
    if mime_type not in {"image/jpeg", "image/png", "image/webp"}:
        mime_type = "image/jpeg"
    extension = {"image/jpeg": "jpg", "image/png": "png", "image/webp": "webp"}.get(mime_type, "jpg")
    filename = f"runway_input.{extension}"
    ratio = _runway_direct_ratio(aspect)
    duration = max(2, min(10, int(duration_s or 5)))
    last_err = ""

    for model in _runway_direct_i2v_model_candidates():
        try:
            async with _runway_official_client() as rw:
                task_id = await rw.create_image_to_video(
                    image_bytes=img_bytes,
                    filename=filename,
                    mime_type=mime_type,
                    prompt_text=prompt,
                    model=model,
                    ratio=ratio,
                    duration=duration,
                    endpoint=RUNWAY_I2V_PATH,
                    upload_endpoint=RUNWAY_UPLOAD_PATH,
                )
                await update.effective_message.reply_text(
                    f"вҸі Runway {model}: РёР·РҫРұСҖР°Р¶РөРҪРёРө Р·Р°РіСҖСғР¶РөРҪРҫ, Р·Р°РҙР°СҮР° РҝСҖРёРҪСҸСӮР°. РһР¶РёРҙР°СҺ СҖРөР·СғР»СҢСӮР°СӮвҖҰ"
                )
                result = await rw.wait_for_task(task_id, timeout_s=RUNWAY_MAX_WAIT_S)

            async with httpx.AsyncClient(timeout=240.0, follow_redirects=True) as dl_client:
                await _reply_video_from_url(
                    update, dl_client, result.first_output,
                    "вңЁ РһР¶РёРІРёР» С„РҫСӮРҫ вң… В· Powered by Runway",
                    task_id=task_id,
                )
            _provider_mark_success("runway_direct")
            return True
        except RunwayAPIError as e:
            last_err = str(e)
            log.warning("Runway direct i2v model=%s failed: %s", model, e)
            # Authentication, billing, moderation and invalid-input errors are not fixed by switching models.
            code = (e.failure_code or "").upper()
            if e.status_code in {400, 401, 402, 403} or code.startswith("SAFETY"):
                _provider_last_error["runway_direct"] = last_err[:700]
                context.user_data["_runway_direct_hard_stop"] = True
                await update.effective_message.reply_text(_runway_user_error_text(e))
                return False
            continue
        except Exception as e:
            last_err = str(e)
            log.warning("Runway direct i2v exception model=%s: %s", model, e)
            continue

    if last_err:
        _provider_last_error["runway_direct"] = last_err[:700]
        _provider_mark_failure("runway_direct", last_err)
    return False


async def _run_runway_animate_photo(update: Update, context: ContextTypes.DEFAULT_TYPE, img_bytes: bytes, prompt: str, duration_s: int, aspect: str):
    await context.bot.send_chat_action(update.effective_chat.id, ChatAction.RECORD_VIDEO)
    prompt = (prompt or "animate the input photo with subtle camera motion, lifelike micro-movements; keep the original person, do not transform identity, do not add phone frames or UI").strip()
    seconds = _duration_for_engine("runway", duration_s)
    ratio = _ratio_for_aspect(aspect)
    bad_src_note = _looks_like_screenshot_or_bad_i2v_source(img_bytes)
    if bad_src_note and I2V_WARN_BAD_SOURCE:
        with contextlib.suppress(Exception):
            if I2V_BAD_SOURCE_POLICY == "ask_clean":
                await update.effective_message.reply_text("в„№пёҸ " + bad_src_note + "\n\nРҹСҖРёСҲР»РёСӮРө СҮРёСҒСӮРҫРө С„РҫСӮРҫ, СҮСӮРҫРұСӢ РҝРҫР»СғСҮРёСӮСҢ СҒСӮР°РұРёР»СҢРҪСӢР№ СҖРөР·СғР»СҢСӮР°СӮ.")
                return False
            await update.effective_message.reply_text("в„№пёҸ " + bad_src_note + "\nРҹСҖРҫРҙРҫР»Р¶Р°СҺ РҫРұСҖР°РұРҫСӮРәСғ, РҪРҫ РәР°СҮРөСҒСӮРІРҫ РјРҫР¶РөСӮ РұСӢСӮСҢ С…СғР¶Рө.")

    # v88: РҫС„РёСҶРёР°Р»СҢРҪСӢР№ Runway Developer API вҖ” РҫСҒРҪРҫРІРҪРҫР№ РјР°СҖСҲСҖСғСӮ.
    context.user_data.pop("_runway_direct_hard_stop", None)
    if RUNWAY_DIRECT_FIRST and RUNWAY_DIRECT_ENABLED and RUNWAY_API_KEY:
        try:
            if await _run_runway_direct_animate_photo(update, context, img_bytes, prompt, seconds, aspect):
                _provider_mark_success("runway_direct")
                return True
            if context.user_data.pop("_runway_direct_hard_stop", False):
                return False
            _provider_mark_failure("runway_direct", _provider_last_error.get("runway_direct", "direct i2v failed"))
        except Exception as e:
            _provider_mark_failure("runway_direct", str(e))
            log.warning("Runway direct i2v failed; trying Comet: %s", e)

    # Comet is only a fallback; its Runway distributor channel can disappear independently.
    try:
        if await _run_runway_comet_animate_photo(update, context, img_bytes, prompt, seconds, aspect):
            return True
    except Exception as e:
        log.warning("Runway Comet route failed: %s", e)

    if (not RUNWAY_DIRECT_FIRST) and RUNWAY_DIRECT_ENABLED and RUNWAY_API_KEY:
        try:
            if await _run_runway_direct_animate_photo(update, context, img_bytes, prompt, seconds, aspect):
                _provider_mark_success("runway_direct")
                return True
            if context.user_data.pop("_runway_direct_hard_stop", False):
                return False
        except Exception as e:
            _provider_mark_failure("runway_direct", str(e))
            log.warning("Runway direct fallback failed: %s", e)

    if RUNWAY_AUTO_FALLBACK_KLING and COMET_API_KEY:
        await update.effective_message.reply_text(RUNWAY_PUBLIC_FALLBACK_TEXT)
        return bool(await _run_comet_i2v(update, context, "kling", img_bytes, prompt, seconds, aspect))

    if RUNWAY_API_KEY:
        await update.effective_message.reply_text("вҡ пёҸ РһС„РёСҶРёР°Р»СҢРҪСӢР№ Runway РҪРө РҝСҖРёРҪСҸР» Р·Р°РҙР°СҮСғ. РҹСҖРҫРІРөСҖСҢСӮРө RUNWAYML_API_SECRET (Environment/Secret File) Рё API-РәСҖРөРҙРёСӮСӢ: /diag_runway auth")
    else:
        await update.effective_message.reply_text("вҡ пёҸ Р”Р»СҸ РҝРҫСҒСӮРҫСҸРҪРҪРҫРіРҫ РҙРҫСҒСӮСғРҝР° Рә Runway РҙРҫРұР°РІСҢСӮРө RUNWAYML_API_SECRET РІ Render Secret File runway.env РёР»Рё Environment.")
    return False

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РҹРҫРәСғРҝРәРё/РёРҪРІРҫР№СҒСӢ в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def _plan_rub(tier: str, term: str) -> int:
    tier = (tier or "pro").lower()
    term = (term or "month").lower()
    return int(PLAN_PRICE_TABLE.get(tier, PLAN_PRICE_TABLE["pro"]).get(term, PLAN_PRICE_TABLE["pro"]["month"]))

def _plan_payload_and_amount(tier: str, months: int) -> tuple[str, int, str]:
    term = {1: "month", 3: "quarter", 12: "year"}.get(months, "month")
    amount = _plan_rub(tier, term)
    title = f"РҹРҫРҙРҝРёСҒРәР° {tier.upper()} ({term})"
    payload = f"sub:{tier}:{months}"
    return payload, amount, title

async def _send_invoice_rub(title: str, desc: str, amount_rub: int, payload: str, update: Update) -> bool:
    try:
        # РұРөСҖС‘Рј СӮРҫРәРөРҪ Рё РІР°Р»СҺСӮСғ РёР· РҙРІСғС… РёСҒСӮРҫСҮРҪРёРәРҫРІ (СҒСӮР°СҖСӢР№ PROVIDER_TOKEN РҳРӣРҳ РҪРҫРІСӢР№ YOOKASSA_PROVIDER_TOKEN)
        token = (PROVIDER_TOKEN or YOOKASSA_PROVIDER_TOKEN)
        curr  = (CURRENCY if (CURRENCY and CURRENCY != "RUB") else YOOKASSA_CURRENCY) or "RUB"

        if not token:
            await update.effective_message.reply_text("вҡ пёҸ Р®Kassa РҪРө РҪР°СҒСӮСҖРҫРөРҪР° (РҪРөСӮ СӮРҫРәРөРҪР°).")
            return False

        prices = [LabeledPrice(label=_ascii_label(title), amount=int(amount_rub) * 100)]

        await update.effective_message.reply_invoice(
            title=title,
            description=desc[:255],
            payload=payload,
            provider_token=token,
            currency=curr,
            prices=prices,
            need_email=False,
            need_name=False,
            need_phone_number=False,
            need_shipping_address=False,
            is_flexible=False
        )
        return True

    except Exception as e:
        log.exception("send_invoice error: %s", e)
        try:
            await update.effective_message.reply_text("РқРө СғРҙР°Р»РҫСҒСҢ РІСӢСҒСӮР°РІРёСӮСҢ СҒСҮС‘СӮ.")
        except Exception:
            pass
        return False

async def on_precheckout(update: Update, context: ContextTypes.DEFAULT_TYPE):
    try:
        q = update.pre_checkout_query
        await q.answer(ok=True)
    except Exception as e:
        log.exception("precheckout error: %s", e)

async def on_successful_payment(update: Update, context: ContextTypes.DEFAULT_TYPE):
    try:
        sp = update.message.successful_payment
        payload = sp.invoice_payload or ""
        total_minor = sp.total_amount or 0
        rub = total_minor / 100.0
        uid = update.effective_user.id

        if payload.startswith("sub:"):
            _, tier, months = payload.split(":", 2)
            months = int(months)
            until = activate_subscription_with_tier(uid, tier, months)
            await update.effective_message.reply_text(f"вң… РҹРҫРҙРҝРёСҒРәР° {tier.upper()} Р°РәСӮРёРІРёСҖРҫРІР°РҪР° РҙРҫ {until.strftime('%Y-%m-%d')}.\nрҹӘҷ РҡСҖРөРҙРёСӮСӢ РҪР°СҮРёСҒР»РөРҪСӢ: {SUBSCRIPTION_CREDITS.get((tier or "").lower(), 0) * int(months)} РәСҖ.")
            return

        if payload.startswith("topup:"):
            try:
                _, credits_s, rub_s = payload.split(":", 2)
                resolved = _credit_pack_resolve(int(credits_s), int(rub_s))
                if resolved:
                    credits, expected_rub = resolved
                    _wallet_total_add(uid, _credits_to_usd(credits))
                    await update.effective_message.reply_text(
                        f"вң… РһРҝР»Р°СӮР° РҝСҖРҫСҲР»Р° СғСҒРҝРөСҲРҪРҫ. РқР°СҮРёСҒР»РөРҪРҫ: {credits} РәСҖРөРҙРёСӮРҫРІ Р·Р° {expected_rub} вӮҪ."
                    )
                    return
            except Exception:
                log.exception("Failed to parse topup payload: %s", payload)

        # РӣСҺРұРҫРө РёРҪРҫРө payload вҖ” РҝРҫРҝРҫР»РҪРөРҪРёРө РөРҙРёРҪРҫРіРҫ РәРҫСҲРөР»СҢРәР°
        usd = _credits_to_usd(rub)
        _wallet_total_add(uid, usd)
        await update.effective_message.reply_text(f"рҹ’і РҹРҫРҝРҫР»РҪРөРҪРёРө: {rub:.0f} вӮҪ. РқР°СҮРёСҒР»РөРҪРҫ: {_credits_fmt_from_usd(usd)}.")
    except Exception as e:
        log.exception("successful_payment handler error: %s", e)


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ CryptoBot в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
CRYPTO_PAY_API_TOKEN = os.environ.get("CRYPTO_PAY_API_TOKEN", "").strip()
CRYPTO_BASE = "https://pay.crypt.bot/api"
TON_USD_RATE = float(os.environ.get("TON_USD_RATE", "5.0") or "5.0")  # Р·Р°РҝР°СҒРҪРҫР№ РәСғСҖСҒ

async def _crypto_create_invoice(usd_amount: float, asset: str = "USDT", description: str = "") -> tuple[str|None, str|None, float, str]:
    if not CRYPTO_PAY_API_TOKEN:
        return None, None, 0.0, asset
    try:
        payload = {"asset": asset, "amount": round(float(usd_amount), 2), "description": description or "Top-up"}
        headers = {"Crypto-Pay-API-Token": CRYPTO_PAY_API_TOKEN}
        async with httpx.AsyncClient(timeout=30.0) as client:
            r = await client.post(f"{CRYPTO_BASE}/createInvoice", headers=headers, json=payload)
            j = r.json()
            ok = j.get("ok") is True
            if not ok:
                return None, None, 0.0, asset
            res = j.get("result", {})
            return str(res.get("invoice_id")), res.get("pay_url"), float(res.get("amount", usd_amount)), res.get("asset") or asset
    except Exception as e:
        log.exception("crypto create error: %s", e)
        return None, None, 0.0, asset

async def _crypto_get_invoice(invoice_id: str) -> dict | None:
    if not CRYPTO_PAY_API_TOKEN:
        return None
    try:
        headers = {"Crypto-Pay-API-Token": CRYPTO_PAY_API_TOKEN}
        async with httpx.AsyncClient(timeout=20.0) as client:
            r = await client.get(f"{CRYPTO_BASE}/getInvoices?invoice_ids={invoice_id}", headers=headers)
            j = r.json()
            if not j.get("ok"):
                return None
            items = (j.get("result", {}) or {}).get("items", [])
            return items[0] if items else None
    except Exception as e:
        log.exception("crypto get error: %s", e)
        return None

async def _poll_crypto_invoice(context: ContextTypes.DEFAULT_TYPE, chat_id: int, message_id: int, user_id: int, invoice_id: str, usd_amount: float):
    try:
        for _ in range(120):  # ~12 РјРёРҪСғСӮ РҝСҖРё 6СҒ Р·Р°РҙРөСҖР¶РәРө
            inv = await _crypto_get_invoice(invoice_id)
            st = (inv or {}).get("status", "").lower() if inv else ""
            if st == "paid":
                _wallet_total_add(user_id, float(usd_amount))
                with contextlib.suppress(Exception):
                    await context.bot.edit_message_text(chat_id=chat_id, message_id=message_id,
                        text=f"вң… CryptoBot: РҝР»Р°СӮС‘Р¶ РҝРҫРҙСӮРІРөСҖР¶РҙС‘РҪ. РқР°СҮРёСҒР»РөРҪРҫ: {_credits_fmt_from_usd(float(usd_amount))}.")
                return
            if st in ("expired", "cancelled", "canceled", "failed"):
                with contextlib.suppress(Exception):
                    await context.bot.edit_message_text(chat_id=chat_id, message_id=message_id,
                        text=f"вқҢ CryptoBot: РҝР»Р°СӮС‘Р¶ РҪРө Р·Р°РІРөСҖСҲС‘РҪ (СҒСӮР°СӮСғСҒ: {st}).")
                return
            await asyncio.sleep(6.0)
        with contextlib.suppress(Exception):
            await context.bot.edit_message_text(chat_id=chat_id, message_id=message_id,
                text="вҢӣ CryptoBot: РІСҖРөРјСҸ РҫР¶РёРҙР°РҪРёСҸ РІСӢСҲР»Рҫ. РқР°Р¶РјРёСӮРө В«рҹ”Һ РҹСҖРҫРІРөСҖРёСӮСҢВ» РҝРҫР·Р¶Рө.")
    except Exception as e:
        log.exception("crypto poll error: %s", e)

async def _poll_crypto_sub_invoice(
    context: ContextTypes.DEFAULT_TYPE,
    chat_id: int,
    message_id: int,
    user_id: int,
    invoice_id: str,
    tier: str,
    months: int
):
    try:
        for _ in range(120):  # ~12 РјРёРҪСғСӮ РҝСҖРё Р·Р°РҙРөСҖР¶РәРө 6СҒ
            inv = await _crypto_get_invoice(invoice_id)
            st = (inv or {}).get("status", "").lower() if inv else ""
            if st == "paid":
                until = activate_subscription_with_tier(user_id, tier, months)
                with contextlib.suppress(Exception):
                    await context.bot.edit_message_text(
                        chat_id=chat_id, message_id=message_id,
                        text=f"вң… CryptoBot: РҝР»Р°СӮС‘Р¶ РҝРҫРҙСӮРІРөСҖР¶РҙС‘РҪ.\n"
                             f"РҹРҫРҙРҝРёСҒРәР° {tier.upper()} Р°РәСӮРёРІРҪР° РҙРҫ {until.strftime('%Y-%m-%d')}.\nрҹӘҷ РҡСҖРөРҙРёСӮСӢ РҪР°СҮРёСҒР»РөРҪСӢ: {SUBSCRIPTION_CREDITS.get((tier or "").lower(), 0) * int(months)} РәСҖ."
                    )
                return
            if st in ("expired", "cancelled", "canceled", "failed"):
                with contextlib.suppress(Exception):
                    await context.bot.edit_message_text(
                        chat_id=chat_id, message_id=message_id,
                        text=f"вқҢ CryptoBot: РҫРҝР»Р°СӮР° РҪРө Р·Р°РІРөСҖСҲРөРҪР° (СҒСӮР°СӮСғСҒ: {st})."
                    )
                return
            await asyncio.sleep(6.0)

        # РўР°Р№РјР°СғСӮ
        with contextlib.suppress(Exception):
            await context.bot.edit_message_text(
                chat_id=chat_id, message_id=message_id,
                text="вҢӣ CryptoBot: РІСҖРөРјСҸ РҫР¶РёРҙР°РҪРёСҸ РІСӢСҲР»Рҫ. РқР°Р¶РјРёСӮРө В«рҹ”Һ РҹСҖРҫРІРөСҖРёСӮСҢВ» РёР»Рё РҫРҝР»Р°СӮРёСӮРө Р·Р°РҪРҫРІРҫ."
            )
    except Exception as e:
        log.exception("crypto poll (subscription) error: %s", e)


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РҹСҖРөРҙР»РҫР¶РөРҪРёРө РҝРҫРҝРҫР»РҪРөРҪРёСҸ в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
async def _send_topup_menu(update: Update, context: ContextTypes.DEFAULT_TYPE):
    # v83: РҝРҫР»СҢР·РҫРІР°СӮРөР»СҺ РҝРҫРәР°Р·СӢРІР°РөРј РәСҖРөРҙРёСӮСӢ, РІРҪСғСӮСҖРё legacy-РұР°Р»Р°РҪСҒ С…СҖР°РҪРёСӮСҒСҸ РІ СӮРөС…РҪРёСҮРөСҒРәРҫРј СҚРәРІРёРІР°Р»РөРҪСӮРө.
    small_cr = int(os.environ.get("CREDIT_PACK_SMALL_CREDITS", "1000") or 1000)
    mid_cr = int(os.environ.get("CREDIT_PACK_MID_CREDITS", "3000") or 3000)
    big_cr = int(os.environ.get("CREDIT_PACK_BIG_CREDITS", "7000") or 7000)
    small_rub = int(os.environ.get("CREDIT_PACK_SMALL_RUB", "990") or 990)
    mid_rub = int(os.environ.get("CREDIT_PACK_MID_RUB", "2790") or 2790)
    big_rub = int(os.environ.get("CREDIT_PACK_BIG_RUB", "6290") or 6290)
    kb = InlineKeyboardMarkup([
        [InlineKeyboardButton(f"{small_cr} РәСҖ. вҖў {small_rub} вӮҪ", callback_data=f"topup:rub:{small_rub}"),
         InlineKeyboardButton(f"{mid_cr} РәСҖ. вҖў {mid_rub} вӮҪ", callback_data=f"topup:rub:{mid_rub}")],
        [InlineKeyboardButton(f"{big_cr} РәСҖ. вҖў {big_rub} вӮҪ", callback_data=f"topup:rub:{big_rub}")],
        [InlineKeyboardButton(f"Crypto ~{small_cr} РәСҖ.", callback_data=f"topup:crypto:{_credits_to_usd(small_cr):.2f}"),
         InlineKeyboardButton(f"Crypto ~{mid_cr} РәСҖ.", callback_data=f"topup:crypto:{_credits_to_usd(mid_cr):.2f}")],
    ])
    await update.effective_message.reply_text(
        "рҹӘҷ РҡСҖРөРҙРёСӮСӢ РёСҒРҝРҫР»СҢР·СғСҺСӮСҒСҸ РҙР»СҸ СӮСҸР¶С‘Р»СӢС… С„СғРҪРәСҶРёР№: РІРёРҙРөРҫ, РјСғР·СӢРәР°, AI-С„РҫСӮРҫ, FaceSwap, РіРҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ Рё РҝСҖРөРјРёСғРј-СҖРөРҪРҙРөСҖСӢ.\n"
        "1 РәСҖРөРҙРёСӮ = 1 вӮҪ. Р’СӢРұРөСҖРёСӮРө РҝР°РәРөСӮ:",
        reply_markup=kb,
    )


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РҹСҖРҫРјРҫ-РәРІРҫСӮСӢ РҝРҫ С„СғРҪРәСҶРёСҸРј в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def _promo_feature_key(engine: str, remember_kind: str = "") -> str:
    rk = (remember_kind or "").strip().lower()
    eng = (engine or "").strip().lower()
    if "vocal" in rk and ("clip" in rk or "lipsync" in rk):
        return "vocal_lipsync_clip"
    if "photo_music_clip" in rk or "photo_clip" in rk:
        return "photo_music_clip"
    if "talking_avatar" in rk or "avatar" in rk:
        return "talking_avatar"
    if rk.startswith("text_video") or rk.startswith("video_"):
        return rk
    if "revive_photo" in rk:
        return rk
    if "suno" in rk:
        return "suno_music"
    if "business_logo" in rk or "logo" in rk:
        return "business_logo"
    if "faceswap" in rk or "face_swap" in rk:
        return "faceswap"
    if "ai_selfie" in rk:
        return "ai_selfie"
    if "removebg" in rk or "remove_background" in rk:
        return "remove_background"
    if "replacebg" in rk or "replace_background" in rk:
        return "replace_background"
    if "outpaint" in rk:
        return "outpaint"
    if "image_retouch" in rk or "retouch" in rk:
        return "image_retouch"
    if "img_generate" in rk or "image_generate" in rk or eng == "img":
        return "image_generation"
    return rk or eng or "function"

def _promo_quota_kv_key(user_id: int, feature: str, ymd: str | None = None) -> str:
    safe = re.sub(r"[^a-z0-9_\-]+", "_", (feature or "function").lower())[:80]
    return f"promo5:{user_id}:{ymd or _today_ymd()}:{safe}"

def _try_consume_promo_daily5_quota(user_id: int, username: str | None, engine: str, remember_kind: str = "") -> tuple[bool, int, int, str]:
    if not is_promo_daily5_user(user_id, username):
        return False, 0, PROMO_DAILY5_PER_FUNCTION_LIMIT, ""
    feature = _promo_feature_key(engine, remember_kind)
    limit = max(0, int(PROMO_DAILY5_PER_FUNCTION_LIMIT))
    if limit <= 0:
        return False, 0, limit, feature
    key = _promo_quota_kv_key(user_id, feature)
    used = int(kv_get(key, "0") or "0")
    if used >= limit:
        return False, 0, limit, feature
    kv_set(key, str(used + 1))
    return True, max(0, limit - used - 1), limit, feature


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РҹРҫРҝСӢСӮРәР° РҫРҝР»Р°СӮРёСӮСҢ вҶ’ РІСӢРҝРҫР»РҪРёСӮСҢ в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
async def _try_pay_then_do(
    update: Update,
    context: ContextTypes.DEFAULT_TYPE,
    user_id: int,
    engine: str,
    est_cost_usd: float,
    coro_func,
    remember_kind: str = "",
    remember_payload: dict | None = None
):
    """Reserve capacity, run provider, and charge only after the action explicitly returns True."""
    username = (update.effective_user.username or "")

    promo_ok, promo_left, promo_limit, promo_feature = _try_consume_promo_daily5_quota(user_id, username, engine, remember_kind)
    if promo_ok:
        try:
            with contextlib.suppress(Exception):
                await update.effective_message.reply_text(f"рҹҺҒ РҹСҖРҫРјРҫ-РҙРҫСҒСӮСғРҝ: С„СғРҪРәСҶРёСҸ В«{promo_feature}В». РһСҒСӮР°Р»РҫСҒСҢ СҒРөРіРҫРҙРҪСҸ: {promo_left}/{promo_limit}.")
            result = await coro_func()
            if result is False:
                raise RuntimeError("provider returned unsuccessful result")
        except Exception as e:
            log.exception("promo daily5 action failed: %s", e)
            await update.effective_message.reply_text("вқҢ Р—Р°РҙР°СҮР° РҪРө РІСӢРҝРҫР»РҪРөРҪР°. РҹСҖРҫРјРҫ-РәСҖРөРҙРёСӮСӢ РҪРө СҒРҝРёСҒСӢРІР°СҺСӮСҒСҸ.")
        return
    if is_promo_daily5_user(user_id, username) and promo_feature:
        await update.effective_message.reply_text(f"РҹСҖРҫРјРҫ-Р»РёРјРёСӮ РҪР° С„СғРҪРәСҶРёСҺ В«{promo_feature}В» СҒРөРіРҫРҙРҪСҸ РёСҒСҮРөСҖРҝР°РҪ: {promo_limit}/{promo_limit}. GPT-СҮР°СӮ РҫСҒСӮР°С‘СӮСҒСҸ РұРөР·Р»РёРјРёСӮРҪСӢРј.")
        return

    free_kind = _free_quota_category(engine, remember_kind)
    if free_kind and get_subscription_tier(user_id) == "free" and not is_unlimited(user_id, username):
        q_ok, q_left, q_limit = _try_consume_free_daily_quota(user_id, username, free_kind)
        if q_ok:
            try:
                with contextlib.suppress(Exception): await update.effective_message.reply_text(f"рҹҺҒ Р‘РөСҒРҝР»Р°СӮРҪРҫРө РҙРөР№СҒСӮРІРёРө: {_free_quota_label(free_kind)}. РһСҒСӮР°Р»РҫСҒСҢ СҒРөРіРҫРҙРҪСҸ: {q_left}/{q_limit}.")
                result = await coro_func()
                if result is False: raise RuntimeError("provider returned unsuccessful result")
            except Exception as e:
                log.exception("free action failed: %s", e)
                await update.effective_message.reply_text("вқҢ Р—Р°РҙР°СҮР° РҪРө РІСӢРҝРҫР»РҪРөРҪР°. Р”РөРҪРөР¶РҪРҫРіРҫ СҒРҝРёСҒР°РҪРёСҸ РҪРө РұСӢР»Рҫ.")
            return
        await _send_free_quota_exhausted(update, context, free_kind); return

    if is_unlimited(user_id, username):
        try:
            result = await coro_func()
            if result is False: raise RuntimeError("provider returned unsuccessful result")
        except Exception as e:
            log.exception("unlimited action failed: %s", e)
            await update.effective_message.reply_text("вқҢ Р—Р°РҙР°СҮР° РҪРө РІСӢРҝРҫР»РҪРөРҪР°. РҹРҫРҝСҖРҫРұСғР№СӮРө РҝРҫР·Р¶Рө.")
        return

    provider_cost = max(0.0, float(est_cost_usd or 0.0))
    retail_usd = _retail_usd(provider_cost)
    price_credits = _retail_credits(provider_cost)
    tx_id = None
    try:
        tx_id, available_before = _credit_reserve(user_id, engine, remember_kind or engine, provider_cost, retail_usd, remember_payload)
    except Exception as e:
        log.exception("credit reserve failed: %s", e)
        await update.effective_message.reply_text("вқҢ РқРө СғРҙР°Р»РҫСҒСҢ РҝСҖРҫРІРөСҖРёСӮСҢ РәСҖРөРҙРёСӮРҪСӢР№ РұР°Р»Р°РҪСҒ. РҹРҫРҝСҖРҫРұСғР№СӮРө РөСүС‘ СҖР°Р·.")
        return

    if not tx_id:
        available_cr = int(round(_usd_to_credits(available_before)))
        missing_cr = max(0, price_credits - available_cr)
        await update.effective_message.reply_text(
            f"РқРөРҙРҫСҒСӮР°СӮРҫСҮРҪРҫ РәСҖРөРҙРёСӮРҫРІ. РЎСӮРҫРёРјРҫСҒСӮСҢ: {price_credits} РәСҖ. Р”РҫСҒСӮСғРҝРҪРҫ: {available_cr} РәСҖ. РқРө С…РІР°СӮР°РөСӮ: {missing_cr} РәСҖ.",
            reply_markup=InlineKeyboardMarkup([[InlineKeyboardButton("вӯҗ РўР°СҖРёС„СӢ", web_app=WebAppInfo(url=TARIFF_URL))],[InlineKeyboardButton("вһ• РҹРҫРҝРҫР»РҪРёСӮСҢ РұР°Р»Р°РҪСҒ", callback_data="topup")]])
        )
        return

    after_cr = max(0, int(round(_usd_to_credits(available_before - retail_usd))))
    with contextlib.suppress(Exception):
        await update.effective_message.reply_text(f"рҹӘҷ РЎСӮРҫРёРјРҫСҒСӮСҢ: {price_credits} РәСҖ. РЎРҝРёСҒР°РҪРёРө РҝСҖРҫРёР·РҫР№РҙС‘СӮ СӮРҫР»СҢРәРҫ РҝРҫСҒР»Рө СғСҒРҝРөСҲРҪРҫРіРҫ СҖРөР·СғР»СҢСӮР°СӮР°. РҹРҫСҒР»Рө РІСӢРҝРҫР»РҪРөРҪРёСҸ РҫСҒСӮР°РҪРөСӮСҒСҸ: {after_cr} РәСҖ.")

    try:
        result = await coro_func()
        if result is not True:
            _credit_release(tx_id, "released")
            await update.effective_message.reply_text("вҶ©пёҸ Р“РөРҪРөСҖР°СҶРёСҸ РҪРө Р·Р°РІРөСҖСҲРёР»Р°СҒСҢ вҖ” РәСҖРөРҙРёСӮСӢ РҪРө СҒРҝРёСҒР°РҪСӢ.")
            return
        if not _credit_commit(tx_id):
            _credit_release(tx_id, "released")
            await update.effective_message.reply_text("вҡ пёҸ Р РөР·СғР»СҢСӮР°СӮ РҝРҫР»СғСҮРөРҪ, РҪРҫ СҒРҝРёСҒР°РҪРёРө РҪРө Р·Р°С„РёРәСҒРёСҖРҫРІР°РҪРҫ. РһРұСҖР°СӮРёСӮРөСҒСҢ РІ РҝРҫРҙРҙРөСҖР¶РәСғ, РҝРҫРІСӮРҫСҖРҪРҫ Р·Р°РҝСғСҒРәР°СӮСҢ РҫРҝР»Р°СӮСғ РҪРө РҪСғР¶РҪРҫ.")
            return
        with contextlib.suppress(Exception):
            await update.effective_message.reply_text(f"вң… РЎРҝРёСҒР°РҪРҫ: {price_credits} РәСҖ.")
    except Exception as e:
        _credit_release(tx_id, "released")
        log.exception("paid action failed: %s", e)
        await update.effective_message.reply_text("вқҢ Р—Р°РҙР°СҮР° РҪРө РІСӢРҝРҫР»РҪРөРҪР°. РҡСҖРөРҙРёСӮСӢ РҪРө СҒРҝРёСҒР°РҪСӢ.")


async def cmd_diag_access(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user = update.effective_user
    user_id = user.id if user else 0
    username = user.username if user else ""
    await update.effective_message.reply_text(
        "рҹ”җ Access diagnostic\n"
        f"user_id: {user_id}\n"
        f"username: @{username or '-'}\n"
        f"unlimited: {is_unlimited(user_id, username)}\n"
        f"promo_unlim_gpt: {is_promo_unlim_gpt(user_id, username)}\n"
        f"promo_daily5: {is_promo_daily5_user(user_id, username)}\n"
        f"promo_limit_per_function: {PROMO_DAILY5_PER_FUNCTION_LIMIT}"
    )

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ /plans в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
async def cmd_plans(update: Update, context: ContextTypes.DEFAULT_TYPE):
    lines = ["вӯҗ РўР°СҖРёС„СӢ Рё РәСҖРөРҙРёСӮСӢ:", "РҹРҫРҙРҝРёСҒРәР° РҫСӮРәСҖСӢРІР°РөСӮ РҙРҫСҒСӮСғРҝ, РәСҖРөРҙРёСӮСӢ СҖР°СҒС…РҫРҙСғСҺСӮСҒСҸ РҪР° СӮСҸР¶С‘Р»СӢРө РіРөРҪРөСҖР°СҶРёРё."]
    for tier, terms in PLAN_PRICE_TABLE.items():
        lines.append(f"вҖ” {tier.upper()}: "
                     f"{terms['month']}вӮҪ/РјРөСҒ вҖў {terms['quarter']}вӮҪ/РәРІР°СҖСӮР°Р» вҖў {terms['year']}вӮҪ/РіРҫРҙ вҖў {SUBSCRIPTION_CREDITS.get(tier,0)} РәСҖ./РјРөСҒ")
    kb = InlineKeyboardMarkup([
        [InlineKeyboardButton("START 1 РјРөСҒ",    callback_data="buy:start:1"),
         InlineKeyboardButton("PRO 1 РјРөСҒ",      callback_data="buy:pro:1")],
        [InlineKeyboardButton("ULTIMATE 1 РјРөСҒ", callback_data="buy:ultimate:1")],
        [InlineKeyboardButton("РңРёРҪРё-РІРёСӮСҖРёРҪР°",    web_app=WebAppInfo(url=TARIFF_URL))]
    ])
    await update.effective_message.reply_text("\n".join(lines), reply_markup=kb)


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РһРұС‘СҖСӮРәР° РҙР»СҸ РҝРөСҖРөРҙР°СҮРё РҝСҖРҫРёР·РІРҫР»СҢРҪРҫРіРҫ СӮРөРәСҒСӮР° (РҪР°РҝСҖ. РёР· STT) в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
async def on_text_with_text(
    update: Update,
    context: ContextTypes.DEFAULT_TYPE,
    text: str,
):
    """
    РһРұС‘СҖСӮРәР° РҙР»СҸ РҝРөСҖРөРҙР°СҮРё СӮРөРәСҒСӮР° (РҪР°РҝСҖРёРјРөСҖ, РҝРҫСҒР»Рө STT) РІ on_text,
    РұРөР· РҝРҫРҝСӢСӮРҫРә РёР·РјРөРҪРёСӮСҢ update.message (read-only!).
    """
    text = (text or "").strip()
    if not text:
        await update.effective_message.reply_text("РқРө СғРҙР°Р»РҫСҒСҢ СҖР°СҒРҝРҫР·РҪР°СӮСҢ СӮРөРәСҒСӮ.")
        return

    await on_text(update, context, manual_text=text)


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РўРөРәСҒСӮРҫРІСӢР№ РІС…РҫРҙ в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
async def on_text(
    update: Update,
    context: ContextTypes.DEFAULT_TYPE,
    manual_text: str | None = None,
):
    # Р•СҒР»Рё СӮРөРәСҒСӮ РҝРөСҖРөРҙР°РҪ РёР·РІРҪРө вҶ’ РёСҒРҝРҫР»СҢР·СғРөРј РөРіРҫ
    # РёРҪР°СҮРө вҖ” РҫРұСӢСҮРҪСӢР№ СӮРөРәСҒСӮ СҒРҫРҫРұСүРөРҪРёСҸ
    if manual_text is not None:
        text = manual_text.strip()
    else:
        text = (update.message.text or "").strip()

    # Rename virtual chat before routing the message to GPT.
    rename_cid = context.user_data.pop("awaiting_chat_rename", None)
    if rename_cid:
        if _chat_rename(update.effective_user.id, update.effective_chat.id, int(rename_cid), text):
            await update.effective_message.reply_text("вң… РқР°Р·РІР°РҪРёРө СҮР°СӮР° РҫРұРҪРҫРІР»РөРҪРҫ.", reply_markup=_chat_list_kb(update.effective_user.id, update.effective_chat.id))
        else:
            await update.effective_message.reply_text("РқРө СғРҙР°Р»РҫСҒСҢ РҝРөСҖРөРёРјРөРҪРҫРІР°СӮСҢ СҮР°СӮ.")
        return

    if context.user_data.pop("awaiting_midjourney_prompt", None):
        await _start_midjourney_image(update, context, text)
        return

    # Active presentation/catalog project has priority over generic GPT routing.
    if await _presentation_studio_get().handle_text(update, context, text):
        return

    # Р’РҫРҝСҖРҫСҒСӢ Рҫ FaceSwap РҙРҫР»Р¶РҪСӢ РҫСӮРІРөСҮР°СӮСҢ РҫРҝРёСҒР°РҪРёРөРј С„СғРҪРәСҶРёРё, Р° РҪРө СҒСҖР°Р·Сғ Р·Р°РҝСғСҒРәР°СӮСҢ СҖРөР¶РёРј.
    if re.search(r"(РјРҫР¶(РөСҲСҢ|РөСӮРө|РҪРҫ)|СғРјРө(РөСҲСҢ|РөСӮРө)|СҒРҝРҫСҒРҫРұРөРҪ|РҝРҫРҙРҙРөСҖР¶РёРІР°РөСҲСҢ|РҙРөР»Р°РөСҲСҢ|РјРҫР¶РөСӮ\s+Р»Рё)", text or "", re.I) and re.search(r"(Р»РёСҶ|Р»РёСҶР°|Р»РёСҶРҫ|face|faceswap)", text or "", re.I):
        cap_early = capability_answer(text)
        if cap_early:
            await update.effective_message.reply_text(cap_early, reply_markup=main_kb)
            with contextlib.suppress(Exception):
                _chat_memory_add(update.effective_user.id, update.effective_chat.id, "user", text)
                _chat_memory_add(update.effective_user.id, update.effective_chat.id, "assistant", cap_early)
            return

    # Р—Р°РјРөРҪР° Р»РёСҶР°: РҫСӮРҙРөР»СҢРҪСӢР№ РҙРІСғС…СҲР°РіРҫРІСӢР№ СҖРөР¶РёРј.
    if _is_face_swap_request(text):
        _clear_medicine_wait(context)
        with contextlib.suppress(Exception):
            _mode_track_set(update.effective_user.id, "")
        await _start_faceswap_flow(update, context, None, use_cached=False)
        return

    # РЈРҙР°Р»РөРҪРёРө/Р·Р°РјРөРҪР° С„РҫРҪР°: РөСҒР»Рё РҝРҫР»СҢР·РҫРІР°СӮРөР»СҢ СғР¶Рө Р·Р°РіСҖСғР·РёР» С„РҫСӮРҫ РёР»Рё РұРҫСӮ Р¶РҙС‘СӮ СғСӮРҫСҮРҪРөРҪРёРө.
    if _is_replacebg_wait_text(context):
        img = _get_cached_photo(update.effective_user.id)
        if not img:
            _clear_replacebg_wait(context)
            await update.effective_message.reply_text("РЎРҪР°СҮР°Р»Р° РҝСҖРёСҲР»РёСӮРө С„РҫСӮРҫ, Р·Р°СӮРөРј РІСӢРұРөСҖРёСӮРө Р·Р°РјРөРҪСғ С„РҫРҪР°.", reply_markup=main_kb)
            return
        kind, prompt = _bg_kind_from_text(text)
        _clear_medicine_wait(context)
        _clear_replacebg_wait(context)
        with contextlib.suppress(Exception):
            _mode_track_set(update.effective_user.id, "")
        await _pedit_replacebg(update, context, img, kind=kind, prompt=prompt)
        return

    if _is_remove_bg_request(text):
        _clear_medicine_wait(context)
        with contextlib.suppress(Exception):
            _mode_track_set(update.effective_user.id, "")
        img = _get_cached_photo(update.effective_user.id)
        if img:
            await _pedit_removebg(update, context, img)
            return
        _set_waiting_removebg(context)
        await update.effective_message.reply_text("РҹСҖРёСҲР»РёСӮРө С„РҫСӮРҫ вҖ” СғРҙР°Р»СҺ С„РҫРҪ Рё РІРөСҖРҪСғ PNG СҒ РҝСҖРҫР·СҖР°СҮРҪРҫР№ РҝРҫРҙР»РҫР¶РәРҫР№.", reply_markup=main_kb)
        return

    if _is_replace_bg_request(text):
        _clear_medicine_wait(context)
        with contextlib.suppress(Exception):
            _mode_track_set(update.effective_user.id, "")
        img = _get_cached_photo(update.effective_user.id)
        kind, prompt = _bg_kind_from_text(text)
        if img:
            await _pedit_replacebg(update, context, img, kind=kind, prompt=prompt)
            return
        _set_waiting_replacebg(context, prompt=text)
        await update.effective_message.reply_text("РҹСҖРёСҲР»РёСӮРө С„РҫСӮРҫ вҖ” РІСӢСҖРөР¶Сғ РҫРұСҠРөРәСӮ Рё Р·Р°РјРөРҪСҺ СӮРҫР»СҢРәРҫ С„РҫРҪ. Р”Р»СҸ РҝСҖРөСҒРөСӮРҫРІ Рё СӮРөРәСҒСӮРҫРІРҫРіРҫ РҫРҝРёСҒР°РҪРёСҸ РҝРҫСҒСӮР°СҖР°СҺСҒСҢ СҒРҙРөР»Р°СӮСҢ СҖРөР·СғР»СҢСӮР°СӮ РәР°Рә РҪР°СҒСӮРҫСҸСүРөРө СҒРөР»С„Рё.", reply_markup=main_kb)
        return

    # Р РөСӮСғСҲСҢ СҒРҫРұСҒСӮРІРөРҪРҪРҫРіРҫ РёР·РҫРұСҖР°Р¶РөРҪРёСҸ: РөСҒР»Рё С„РҫСӮРҫ СғР¶Рө Р·Р°РіСҖСғР¶РөРҪРҫ Рё РұРҫСӮ Р¶РҙС‘СӮ СғСӮРҫСҮРҪРөРҪРёРө.
    if _is_retouch_wait_text(context):
        img = _get_cached_photo(update.effective_user.id)
        if not img:
            _clear_image_retouch_wait(context)
            await update.effective_message.reply_text(_retouch_user_hint_text(), reply_markup=main_kb)
            return
        instruction = text or context.user_data.get("retouch_prompt") or "СғРұСҖР°СӮСҢ Р»РёСҲРҪСҺСҺ РҪР°РҙРҝРёСҒСҢ/РІРҫРҙСҸРҪРҫР№ Р·РҪР°Рә"
        _clear_medicine_wait(context)
        _clear_image_retouch_wait(context)
        with contextlib.suppress(Exception):
            _mode_track_set(update.effective_user.id, "")
        await _start_image_retouch(update, context, img, instruction)
        return

    # AI-СҒРөР»С„Рё: С„РҫСӮРҫ СғР¶Рө Р·Р°РіСҖСғР¶РөРҪРҫ, Р¶РҙС‘Рј СҒСҶРөРҪСғ/Р·РҪР°РјРөРҪРёСӮРҫСҒСӮСҢ/РҝРөСҖСҒРҫРҪР°Р¶Р°.
    if context.user_data.get("awaiting_ai_selfie_prompt"):
        img = _get_cached_photo(update.effective_user.id)
        if not img:
            _clear_ai_selfie_wait(context)
            await update.effective_message.reply_text("РЎРҪР°СҮР°Р»Р° Р·Р°РіСҖСғР·РёСӮРө СҒРІРҫС‘ СҒРөР»С„Рё, Р·Р°СӮРөРј РҪР°Р¶РјРёСӮРө рҹӨі AI-СҒРөР»С„Рё СҒРҫ Р·РІРөР·РҙРҫР№.", reply_markup=main_kb)
            return
        preset = (context.user_data.pop("ai_selfie_preset_prompt", "") or "").strip()
        _clear_ai_selfie_wait(context)
        await _start_ai_selfie(update, context, img, text, preset)
        return

    # Р“РҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ: РҝРҫСҖСӮСҖРөСӮ СғР¶Рө Р·Р°РіСҖСғР¶РөРҪ, РҪРҫ РҝРөСҖРөРҙ СӮРөРәСҒСӮРҫРј СӮСҖРөРұСғРөСӮСҒСҸ РІСӢРұСҖР°СӮСҢ РіРҫР»РҫСҒ.
    if context.user_data.get("awaiting_avatar_voice_choice"):
        tl = (text or "").strip().lower()
        voice_aliases = {
            "nova": "nova", "РҪРҫРІР°": "nova",
            "onyx": "onyx", "РҫРҪРёРәСҒ": "onyx",
            "alloy": "alloy", "Р°Р»Р»РҫР№": "alloy",
            "shimmer": "shimmer", "СҲРёРјРјРөСҖ": "shimmer",
            "fable": "fable", "С„РөР№РұР»": "fable",
        }
        if tl in voice_aliases:
            chosen = voice_aliases[tl]
            context.user_data["avatar_tts_voice"] = chosen
            context.user_data.pop("awaiting_avatar_voice_choice", None)
            pending_script = (context.user_data.get("avatar_pending_script") or "").strip()
            img = _get_cached_photo(update.effective_user.id)
            if pending_script and img:
                context.user_data.pop("avatar_pending_script", None)
                _clear_avatar_wait(context)
                await update.effective_message.reply_text(f"вң… Р”Р»СҸ Р°РІР°СӮР°СҖР° РІСӢРұСҖР°РҪ РіРҫР»РҫСҒ: {_avatar_tts_voice_label(chosen)}. РўРөРәСҒСӮ СғР¶Рө РҝРҫР»СғСҮРөРҪ вҖ” Р·Р°РҝСғСҒРәР°СҺ РіРҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ.")
                await _start_talking_avatar(update, context, img, script_text=pending_script)
                return
            _set_avatar_wait(context)
            await update.effective_message.reply_text(f"вң… Р”Р»СҸ Р°РІР°СӮР°СҖР° РІСӢРұСҖР°РҪ РіРҫР»РҫСҒ: {_avatar_tts_voice_label(chosen)}. РўРөРҝРөСҖСҢ РҝСҖРёСҲР»РёСӮРө СӮРөРәСҒСӮ, РәРҫСӮРҫСҖСӢР№ РҙРҫР»Р¶РөРҪ РҝСҖРҫРёР·РҪРөСҒСӮРё Р°РІР°СӮР°СҖ.")
            return
        await update.effective_message.reply_text(_avatar_voice_choice_text(), reply_markup=_avatar_voice_choice_kb("act"))
        return

    # Р“РҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ: С„РҫСӮРҫ СғР¶Рө Р·Р°РіСҖСғР¶РөРҪРҫ, Р¶РҙС‘Рј СӮРөРәСҒСӮ/РіРҫР»РҫСҒ РҙР»СҸ СҖРөСҮРё.
    if context.user_data.get("awaiting_avatar_script"):
        img = _get_cached_photo(update.effective_user.id)
        if not img:
            _clear_avatar_wait(context)
            await update.effective_message.reply_text("РЎРҪР°СҮР°Р»Р° Р·Р°РіСҖСғР·РёСӮРө РҝРҫСҖСӮСҖРөСӮ СҮРөР»РҫРІРөРәР°, Р·Р°СӮРөРј РҪР°Р¶РјРёСӮРө рҹ—Ј Р“РҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ.", reply_markup=main_kb)
            return
        _clear_avatar_wait(context)
        await _start_talking_avatar(update, context, img, script_text=text)
        return

    # Р’РҫРәР°Р»СҢРҪСӢР№ РәР»РёРҝ: РҝРҫСҖСӮСҖРөСӮ СғР¶Рө Р·Р°РіСҖСғР¶РөРҪ, Р¶РҙС‘Рј РҫРҝРёСҒР°РҪРёРө РҝРөСҒРҪРё/РәР»РёРҝР°.
    if context.user_data.get("awaiting_vocal_clip_prompt"):
        img = _get_cached_photo(update.effective_user.id)
        if not img:
            _clear_vocal_clip_wait(context)
            await update.effective_message.reply_text("РЎРҪР°СҮР°Р»Р° Р·Р°РіСҖСғР·РёСӮРө РҝРҫСҖСӮСҖРөСӮ РҫРҙРҪРҫРіРҫ СҮРөР»РҫРІРөРәР°, Р·Р°СӮРөРј РҪР°Р¶РјРёСӮРө рҹҺӨ РҡР»РёРҝ СҒ РІРҫРәР°Р»РҫРј.", reply_markup=main_kb)
            return
        _clear_vocal_clip_wait(context)
        with contextlib.suppress(Exception):
            _mode_track_set(update.effective_user.id, "")
        await _start_vocal_clip(update, context, img, text)
        return

    # Р’РёРҙРөРҫ РҝРҫ СӮРөРәСҒСӮСғ/РіРҫР»РҫСҒСғ: Р¶РҙС‘Рј prompt РҝРҫСҒР»Рө РІСӢРұРҫСҖР° РҙРІРёР¶РәР°.
    if context.user_data.get("awaiting_text_video_prompt"):
        _clear_text_video_wait(context)
        await _start_text_video(update, context, text)
        return

    # РӨРҫСӮРҫвҶ’РІРёРҙРөРҫРәР»РёРҝ: С„РҫСӮРҫ СғР¶Рө Р·Р°РіСҖСғР¶РөРҪРҫ, Р¶РҙС‘Рј РҫРҝРёСҒР°РҪРёРө РәР»РёРҝР°.
    if context.user_data.get("awaiting_photo_clip_prompt"):
        img = _get_cached_photo(update.effective_user.id)
        if not img:
            _clear_photo_clip_wait(context)
            await update.effective_message.reply_text("РЎРҪР°СҮР°Р»Р° Р·Р°РіСҖСғР·РёСӮРө С„РҫСӮРҫ СҮРөР»РҫРІРөРәР°, Р·Р°СӮРөРј РҪР°Р¶РјРёСӮРө рҹҺө РӨРҫСӮРҫ вҶ’ РІРёРҙРөРҫРәР»РёРҝ.", reply_markup=main_kb)
            return
        _clear_photo_clip_wait(context)
        await _start_photo_music_clip(update, context, img, text)
        return

    # РўРөРәСҒСӮРҫРІСӢР№/РіРҫР»РҫСҒРҫРІРҫР№ Р·Р°РҝСҖРҫСҒ РҪР° СҖРөСӮСғСҲСҢ РҙРҫ Р·Р°РіСҖСғР·РәРё С„РҫСӮРҫ.
    # РӯСӮРҫ СҒРұСҖР°СҒСӢРІР°РөСӮ РјРөРҙРёСҶРёРҪСғ Рё РҝРөСҖРөРІРҫРҙРёСӮ СҒР»РөРҙСғСҺСүРөРө РёР·РҫРұСҖР°Р¶РөРҪРёРө РІ image-edit.
    if _is_image_retouch_request(text):
        _clear_medicine_wait(context)
        with contextlib.suppress(Exception):
            _mode_track_set(update.effective_user.id, "")
        img = _get_cached_photo(update.effective_user.id)
        if img and _has_own_image_confirmation(text):
            await _start_image_retouch(update, context, img, text)
            return
        _set_waiting_image_retouch(update, context, text)
        await update.effective_message.reply_text(_retouch_user_hint_text(), reply_markup=main_kb)
        return

    if _is_ai_selfie_intent(text):
        img = _get_cached_photo(update.effective_user.id)
        prompt = _clean_ai_selfie_prompt(text)
        if img and prompt:
            await _start_ai_selfie(update, context, img, prompt)
            return
        _set_ai_selfie_wait(context)
        await update.effective_message.reply_text(
            "Р”Р°, СҒРҙРөР»Р°СҺ AI-СҒРөР»С„Рё: Р·Р°РіСҖСғР·РёСӮРө СҒРІРҫС‘ С„РҫСӮРҫ, Р·Р°СӮРөРј РҪР°РҝРёСҲРёСӮРө, СҒ РәРөРј/РіРҙРө СҒРҙРөР»Р°СӮСҢ СҒСҶРөРҪСғ. РқР°РҝСҖРёРјРөСҖ: В«СҒРөР»С„Рё СҒ РёР·РІРөСҒСӮРҪСӢРј Р°РәСӮС‘СҖРҫРј РҪР° РәСҖР°СҒРҪРҫР№ РҙРҫСҖРҫР¶РәРө, iPhone selfie, 4:5В».",
            reply_markup=main_kb,
        )
        return

    if _is_avatar_intent(text):
        img = _get_cached_photo(update.effective_user.id)
        script = _clean_avatar_script(text)
        if img:
            if script and len(script) > 8:
                context.user_data["avatar_pending_script"] = script
            if context.user_data.get("avatar_tts_voice"):
                pending_script = (context.user_data.get("avatar_pending_script") or "").strip()
                if pending_script:
                    context.user_data.pop("avatar_pending_script", None)
                    await update.effective_message.reply_text(
                        f"вң… РҹРҫСҖСӮСҖРөСӮ РҪР°Р№РҙРөРҪ. Р“РҫР»РҫСҒ СғР¶Рө РІСӢРұСҖР°РҪ: {_avatar_tts_voice_label(_avatar_tts_voice_get(context))}. РўРөРәСҒСӮ СғР¶Рө РҝРҫР»СғСҮРөРҪ вҖ” Р·Р°РҝСғСҒРәР°СҺ РіРҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ."
                    )
                    await _start_talking_avatar(update, context, img, script_text=pending_script)
                else:
                    _set_avatar_wait(context)
                    await update.effective_message.reply_text(
                        f"вң… РҹРҫСҖСӮСҖРөСӮ РҪР°Р№РҙРөРҪ. Р“РҫР»РҫСҒ СғР¶Рө РІСӢРұСҖР°РҪ: {_avatar_tts_voice_label(_avatar_tts_voice_get(context))}. РўРөРҝРөСҖСҢ РҝСҖРёСҲР»РёСӮРө СӮРөРәСҒСӮ, voice РёР»Рё Р°СғРҙРёРҫС„Р°Р№Р» 2вҖ“60 СҒРөРәСғРҪРҙ."
                    )
            else:
                _set_avatar_voice_choice_wait(context)
                await update.effective_message.reply_text(_avatar_voice_choice_text(), reply_markup=_avatar_voice_choice_kb("act"))
        else:
            context.user_data["awaiting_avatar_photo"] = True
            if script and len(script) > 8:
                context.user_data["avatar_pending_script"] = script
            await update.effective_message.reply_text(
                "Р”Р°, СҒРҙРөР»Р°СҺ РіРҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ. РЎРҪР°СҮР°Р»Р° Р·Р°РіСҖСғР·РёСӮРө РҝРҫСҖСӮСҖРөСӮ СҮРөР»РҫРІРөРәР°. РҹРҫСҒР»Рө Р·Р°РіСҖСғР·РәРё СҸ РҫРұСҸР·Р°СӮРөР»СҢРҪРҫ РҝСҖРөРҙР»РҫР¶Сғ РІСӢРұСҖР°СӮСҢ РіРҫР»РҫСҒ, Р° Р·Р°СӮРөРј РёСҒРҝРҫР»СҢР·СғСҺ РІР°СҲ СӮРөРәСҒСӮ, voice РёР»Рё Р°СғРҙРёРҫС„Р°Р№Р» 2вҖ“60 СҒРөРәСғРҪРҙ.",
                reply_markup=main_kb,
            )
        return

    if _is_photo_clip_intent(text):
        img = _get_cached_photo(update.effective_user.id)
        prompt = _clean_photo_clip_prompt(text)
        if img and prompt:
            await _start_photo_music_clip(update, context, img, prompt)
            return
        _set_photo_clip_wait(context)
        await update.effective_message.reply_text(
            "Р”Р°, СҒРҙРөР»Р°СҺ РІРёРҙРөРҫРәР»РёРҝ РёР· С„РҫСӮРҫ. Р—Р°РіСҖСғР·РёСӮРө С„РҫСӮРҫ СҮРөР»РҫРІРөРәР°, Р·Р°СӮРөРј РҫРҝРёСҲРёСӮРө СҒСӮРёР»СҢ РәР»РёРҝР°/РјСғР·СӢРәРё, РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ Рё С„РҫСҖРјР°СӮ.",
            reply_markup=main_kb,
        )
        return

    # Р’РҫРҝСҖРҫСҒ/РәРҫРјР°РҪРҙР° РҝСҖРҫ РҫР¶РёРІР»РөРҪРёРө С„РҫСӮРҫ РҙРҫР»Р¶РҪСӢ СҒРұСҖР°СҒСӢРІР°СӮСҢ РјРөРҙРёСҶРёРҪСҒРәСғСҺ РІРөСӮРәСғ.
    # РҳРҪР°СҮРө РҝРҫСҒР»Рө РәРҪРҫРҝРәРё В«РңРөРҙРёСҶРёРҪР°В» СҒР»РөРҙСғСҺСүРөРө РҫРұСӢСҮРҪРҫРө С„РҫСӮРҫ РҫСҲРёРұРҫСҮРҪРҫ СғС…РҫРҙРёСӮ РІ РјРөРҙ. Р°РҪР°Р»РёР·.
    if _is_photo_revival_question(text) or _is_photo_revival_intent(text):
        _set_waiting_photo_revival(update, context)
        await update.effective_message.reply_text(_photo_revival_capability_text(), reply_markup=main_kb)
        return

    # Legacy presentation flags from older deployments are migrated into the v86 studio.
    if context.user_data.pop("awaiting_work_presentation_brief", None):
        await _presentation_studio_get().start(update, context, "presentation")
        if await _presentation_studio_get().handle_text(update, context, text):
            return

    if context.user_data.pop("awaiting_work_catalog_brief", None):
        await _presentation_studio_get().start(update, context, "catalog")
        if await _presentation_studio_get().handle_text(update, context, text):
            return

    if context.user_data.get("awaiting_work_logo_brief"):
        context.user_data.pop("awaiting_work_logo_brief", None)
        _mode_track_set(update.effective_user.id, "work_logo")
        await _generate_business_logo(update, context, text)
        return

    if context.user_data.get("awaiting_suno_brief"):
        context.user_data.pop("awaiting_suno_brief", None)
        _mode_track_set(update.effective_user.id, "suno_music")
        suno_brief = _prepare_suno_brief_from_context(context, text)
        await _run_suno_music(update, context, suno_brief)
        return

    # Р•СҒР»Рё РҝРҫР»СҢР·РҫРІР°СӮРөР»СҢ РІСӢРұСҖР°Р» Reels/С„РёР»СҢРј РІ РјРөРҪСҺ Рё СӮРөРҝРөСҖСҢ РҝРёСҲРөСӮ РІРІРҫРҙРҪСӢРө вҖ”
    # СҒРҪР°СҮР°Р»Р° РҙР°С‘Рј СҒСӮСҖСғРәСӮСғСҖРёСҖРҫРІР°РҪРҪСӢР№ СҒСҶРөРҪР°СҖРҪСӢР№ РҫСӮРІРөСӮ, Р° РҪРө РҫРұСүРёР№ СҮР°СӮ.
    if context.user_data.get("awaiting_reels_material"):
        context.user_data.pop("awaiting_reels_material", None)
        _mode_track_set(update.effective_user.id, "fun_reels")
        prompt = (
            "РўСӢ РҝСҖРҫРҙСҺСҒРөСҖ РәРҫСҖРҫСӮРәРёС… Reels/Shorts. РқР° СҖСғСҒСҒРәРҫРј РҝРҫРҙРіРҫСӮРҫРІСҢ: 1) С…СғРә, 2) СҒСҶРөРҪР°СҖРёР№ РҝРҫ СҒРөРәСғРҪРҙР°Рј, "
            "3) СӮРөРәСҒСӮ РҪР° СҚРәСҖР°РҪРө, 4) voice-over, 5) CTA, 6) РҝСҖРҫРјРҝСӮСӢ РҙР»СҸ Sora/Kling, "
            "7) РјРҫРҪСӮР°Р¶РҪСӢРө РҝРҫРҙСҒРәР°Р·РәРё. Р—Р°РҝСҖРҫСҒ РҝРҫР»СҢР·РҫРІР°СӮРөР»СҸ:\n" + text
        )
        reply = await ask_openai_text(prompt)
        await update.effective_message.reply_text(reply[:3900], reply_markup=main_kb)
        if len(reply) > 3900:
            await update.effective_message.reply_text(reply[3900:7800])
        await maybe_tts_reply(update, context, reply[:TTS_MAX_CHARS])
        return

    if context.user_data.get("awaiting_film_material"):
        context.user_data.pop("awaiting_film_material", None)
        _mode_track_set(update.effective_user.id, "fun_film")
        prompt = (
            "РўСӢ СҖРөР¶РёСҒСҒС‘СҖ Рё РҝСҖРҫРјРҝСӮ-РёРҪР¶РөРҪРөСҖ AI-РІРёРҙРөРҫ. РқР° СҖСғСҒСҒРәРҫРј РҝРҫРҙРіРҫСӮРҫРІСҢ РјРёРҪРё-С„РёР»СҢРј: "
            "1) Р»РҫРіР»Р°Р№РҪ, 2) СҒСӮСҖСғРәСӮСғСҖР° СҒСҶРөРҪ, 3) СҖР°СҒРәР°РҙСҖРҫРІРәР°, 4) РҝСҖРҫРјРҝСӮСӢ РҙР»СҸ РәРҫСҖРҫСӮРәРёС… РәР»РёРҝРҫРІ 5-10 СҒРөРә "
            "СҮРөСҖРөР· Sora 2 РұРөР· Р»СҺРҙРөР№, Kling РёР»Рё Runway, 5) РјРҫРҪСӮР°Р¶/Р·РІСғРә/СӮРёСӮСҖСӢ, 6) РҝР»Р°РҪ СҒРұРҫСҖРәРё. Р—Р°РҝСҖРҫСҒ РҝРҫР»СҢР·РҫРІР°СӮРөР»СҸ:\n" + text
        )
        reply = await ask_openai_text(prompt)
        await update.effective_message.reply_text(reply[:3900], reply_markup=main_kb)
        if len(reply) > 3900:
            await update.effective_message.reply_text(reply[3900:7800])
        await maybe_tts_reply(update, context, reply[:TTS_MAX_CHARS])
        return

    # Р’РҫРҝСҖРҫСҒСӢ Рҫ РІРҫР·РјРҫР¶РҪРҫСҒСӮСҸС….
    # Р•СҒР»Рё РІРҫРҝСҖРҫСҒ РҪРө РјРөРҙРёСҶРёРҪСҒРәРёР№, СҒРұСҖР°СҒСӢРІР°РөРј Р·Р°РІРёСҒСҲРёР№ РјРөРҙ. СҖРөР¶РёРј,
    # СҮСӮРҫРұСӢ СҒР»РөРҙСғСҺСүРёРө С„РҫСӮРҫ/РҙРҫРәСғРјРөРҪСӮСӢ РҪРө СғС…РҫРҙРёР»Рё РҫСҲРёРұРҫСҮРҪРҫ РІ РјРөРҙРёСҶРёРҪСғ.
    cap = capability_answer(text)
    if cap:
        if not _is_medical_capability_question(text):
            _clear_medicine_wait(context)
            with contextlib.suppress(Exception):
                _mode_track_set(update.effective_user.id, "")
        await update.effective_message.reply_text(cap, reply_markup=main_kb)
        return

    # РқР°РјС‘Рә РҪР° РјСғР·СӢРәСғ / РҝРөСҒРҪСҺ СҮРөСҖРөР· Suno
    if re.search(r"(?:СҒРҫР·РҙР°[Р№Рё]|СҒРҙРөР»Р°[Р№Рё]|СҒРіРөРҪРөСҖРёСҖСғ[Р№Рё]|РҪР°РҝРёСҲРё|Р·Р°РҝСғСҒСӮРё).{0,80}(РјСғР·СӢРә|РҝРөСҒРҪ|СӮСҖРөРә|РҙР¶РёРҪРіР»|РјРёРҪСғСҒРҫРІРә|suno)", text, re.I):
        _clear_medicine_wait(context)
        with contextlib.suppress(Exception):
            _mode_track_set(update.effective_user.id, "suno_music")
        await _run_suno_music(update, context, text)
        return

    # РқР°РјС‘Рә РҪР° РіРөРҪРөСҖР°СҶРёСҺ РІРёРҙРөРҫСҖРҫР»РёРәР°
    mtype, rest = detect_media_intent(text)
    if mtype == "video":
        _clear_medicine_wait(context)
        with contextlib.suppress(Exception):
            _mode_track_set(update.effective_user.id, "")
        duration, aspect = parse_video_opts(text)
        prompt = rest or re.sub(
            r"\b(\d+\s*(?:СҒРөРә|СҒ)\b|(?:9:16|16:9|1:1|4:5|3:4|4:3))",
            "",
            text,
            flags=re.I,
        ).strip(" ,.")

        if not prompt:
            await update.effective_message.reply_text(
                "РһРҝРёСҲРёСӮРө, СҮСӮРҫ РёРјРөРҪРҪРҫ СҒРҪСҸСӮСҢ, РҪР°РҝСҖ.: В«СҖРөСӮСҖРҫ-Р°РІСӮРҫ РҪР° РұРөСҖРөРіСғ, Р·Р°РәР°СӮВ»."
            )
            return

        aid = _new_aid()
        _pending_actions[aid] = {
            "prompt": prompt,
            "duration": duration,
            "aspect": aspect,
        }

        people_present = _prompt_likely_has_people(prompt)
        buttons = []
        if not people_present:
            buttons.append([InlineKeyboardButton(f"рҹҺһ Sora 2 В· РұРөР· Р»СҺРҙРөР№ В· {_video_price_credits('sora', duration)} РәСҖ.", callback_data=f"choose:sora:{aid}")])
        buttons.append([InlineKeyboardButton(f"рҹҺ¬ Kling В· {_video_price_credits('kling', duration)} РәСҖ.", callback_data=f"choose:kling:{aid}")])
        if TEXT_VIDEO_ALLOW_RUNWAY:
            buttons.append([InlineKeyboardButton(f"рҹҺҘ Runway В· {_video_price_credits('runway', duration)} РәСҖ.", callback_data=f"choose:runway:{aid}")])
        kb = InlineKeyboardMarkup(buttons)
        sora_note = "Р’ Р·Р°РҝСҖРҫСҒРө РҫРұРҪР°СҖСғР¶РөРҪ СҮРөР»РҫРІРөРә вҖ” Sora 2 СҒРәСҖСӢСӮР°." if people_present else "Sora 2 РҙРҫСҒСӮСғРҝРҪР° СӮРҫР»СҢРәРҫ РҙР»СҸ СҒСҶРөРҪ РұРөР· Р»СҺРҙРөР№."
        await update.effective_message.reply_text(
            f"Р§СӮРҫ РёСҒРҝРҫР»СҢР·РҫРІР°СӮСҢ?\nР”Р»РёСӮРөР»СҢРҪРҫСҒСӮСҢ: {duration} c вҖў РҗСҒРҝРөРәСӮ: {aspect}\nР—Р°РҝСҖРҫСҒ: В«{prompt}В»\n\n{sora_note}\nРЎСӮРҫРёРјРҫСҒСӮСҢ СғР¶Рө РІРәР»СҺСҮР°РөСӮ РјР°СҖР¶Сғ РұРҫСӮР° Рё РұСғРҙРөСӮ СҒРҝРёСҒР°РҪР° СӮРҫР»СҢРәРҫ РҝРҫСҒР»Рө СғСҒРҝРөСҲРҪРҫРіРҫ СҖРөР·СғР»СҢСӮР°СӮР°.",
            reply_markup=kb,
        )
        return



    # РқР°РјС‘Рә РҪР° РәР°СҖСӮРёРҪРәСғ
    if mtype == "image":
        _clear_medicine_wait(context)
        with contextlib.suppress(Exception):
            _mode_track_set(update.effective_user.id, "")
        prompt = rest or re.sub(
            r"^(img|image|picture)\s*[:\-]\s*",
            "",
            text,
            flags=re.I,
        ).strip()

        if not prompt:
            await update.effective_message.reply_text(
                "РӨРҫСҖРјР°СӮ: /img <РҫРҝРёСҒР°РҪРёРө РёР·РҫРұСҖР°Р¶РөРҪРёСҸ>"
            )
            return

        await _ask_image_engine_choice(update, context, prompt)
        return

    # Live-Р·Р°РҝСҖРҫСҒСӢ: РәСғСҖСҒСӢ, РҪРҫРІРҫСҒСӮРё, Р·Р°РәРҫРҪСӢ, РҝРҫРіРҫРҙР°, СҖРөР»РёР·СӢ Рё Р»СҺРұСӢРө Р°РәСӮСғР°Р»СҢРҪСӢРө РҙР°РҪРҪСӢРө.
    # Р“РҫР»РҫСҒРҫРІСӢРө Р·Р°РҝСҖРҫСҒСӢ РҝРҫРҝР°РҙР°СҺСӮ СҒСҺРҙР° СҮРөСҖРөР· on_text_with_text РҝРҫСҒР»Рө STT.
    if not _MEDICAL_TERMS_RE.search(text):
        _clear_medicine_wait(context)
        with contextlib.suppress(Exception):
            _mode_track_set(update.effective_user.id, "")
    if await maybe_handle_live_query(update, context, text):
        return

    # РһРұСӢСҮРҪСӢР№ СӮРөРәСҒСӮ вҶ’ GPT
    ok, _, _ = check_text_and_inc(
        update.effective_user.id,
        update.effective_user.username or "",
    )

    if not ok:
        await update.effective_message.reply_text(
            "РӣРёРјРёСӮ СӮРөРәСҒСӮРҫРІСӢС… Р·Р°РҝСҖРҫСҒРҫРІ РҪР° СҒРөРіРҫРҙРҪСҸ РёСҒСҮРөСҖРҝР°РҪ. "
            "РһС„РҫСҖРјРёСӮРө вӯҗ РҝРҫРҙРҝРёСҒРәСғ РёР»Рё РҝРҫРҝСҖРҫРұСғР№СӮРө Р·Р°РІСӮСҖР°."
        )
        return

    user_id = update.effective_user.id

    # Р РөР¶РёРјСӢ
    try:
        mode = _mode_get(user_id)
        track = _mode_track_get(user_id)
    except NameError:
        mode, track = "none", ""

    if mode and mode != "none":
        text_for_llm = f"[Р РөР¶РёРј: {mode}; РҹРҫРҙСҖРөР¶РёРј: {track or '-'}]\n{text}"
    else:
        text_for_llm = text

    if _MEDICAL_TERMS_RE.search(text):
        # РҜРІРҪСӢР№ РјРөРҙРёСҶРёРҪСҒРәРёР№ РІРҫРҝСҖРҫСҒ/СӮРөРәСҒСӮ вҖ” СҖР°Р·РұРёСҖР°РөРј РәР°Рә РјРөРҙ. РјР°СӮРөСҖРёР°Р».
        await _medical_analyze_text(update, context, text)
        _clear_medicine_wait(context)
        with contextlib.suppress(Exception):
            _mode_track_set(user_id, "")
        return

    # Р•СҒР»Рё РҝРҫР»СҢР·РҫРІР°СӮРөР»СҢ СҖР°РҪСҢСҲРө РҪР°Р¶Р°Р» РјРөРҙ. РҝРҫРҙРјРөРҪСҺ, РҪРҫ СӮРөРҝРөСҖСҢ СҒРҝСҖР°СҲРёРІР°РөСӮ РҝСҖРҫ PDF, РәРҪРёРіРё, С„РҫСӮРҫ,
    # РІРёРҙРөРҫ, РҪРҫРІРҫСҒСӮРё, РәСғСҖСҒ BTC Рё СӮ.Рҝ., РҪРө РҙРөСҖР¶РёРј РөРіРҫ РІ РјРөРҙРёСҶРёРҪСҒРәРҫР№ РІРөСӮРәРө.
    if (track or "").startswith("med_") and not context.user_data.get("medicine_waiting_for_material"):
        with contextlib.suppress(Exception):
            _mode_track_set(user_id, "")

    if mode == "РЈСҮС‘РұР°" and track:
        await study_process_text(update, context, text)
        return

    chat_id = update.effective_chat.id if update.effective_chat else 0
    # Р•СҒР»Рё СҚСӮРҫ РәРҫСҖРҫСӮРәРёР№ РҫСӮРІРөСӮ РҪР° РҝСҖРөРҙСӢРҙСғСүРёР№ СғСӮРҫСҮРҪСҸСҺСүРёР№ РІРҫРҝСҖРҫСҒ, СҖР°СҒСҲРёСҖСҸРөРј Р·Р°РҝСҖРҫСҒ РәРҫРҪСӮРөРәСҒСӮРҫРј.
    llm_input = _chat_memory_followup_query(user_id, chat_id, text_for_llm)
    reply = await ask_openai_text(llm_input, user_id=user_id, chat_id=chat_id)
    await update.effective_message.reply_text(reply)
    _chat_memory_add(user_id, chat_id, "user", text)
    _chat_memory_add(user_id, chat_id, "assistant", reply)
    await maybe_tts_reply(update, context, reply[:TTS_MAX_CHARS])

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РӨРҫСӮРҫ / Р”РҫРәСғРјРөРҪСӮСӢ / Р“РҫР»РҫСҒ в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
async def on_photo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    try:
        ph = update.message.photo[-1]
        f = await ph.get_file()
        data = await f.download_as_bytearray()
        img = bytes(data)
        _cache_photo(update.effective_user.id, img, getattr(f, "file_path", "") or "")

        user_id = update.effective_user.id
        caption = (update.message.caption or "").strip()

        # Presentation Studio: logo/product photo bulk upload.
        if await _presentation_studio_get().handle_photo(update, context, img, mime="image/jpeg", caption=caption):
            return

        # 0) Р—Р°РјРөРҪР° Р»РёСҶР°: РҙРІСғС…СҲР°РіРҫРІСӢР№ СҖРөР¶РёРј РҙРҫР»Р¶РөРҪ СҒСҖР°РұР°СӮСӢРІР°СӮСҢ СҖР°РҪСҢСҲРө РҫСҒСӮР°Р»СҢРҪСӢС… С„РҫСӮРҫ-РІРөСӮРҫРә.
        if context.user_data.get("faceswap_flow") == "await_target":
            await _maybe_choose_target_face(update, context, user_id, img)
            return
        if context.user_data.get("faceswap_flow") == "await_source":
            await _maybe_choose_source_face(update, context, user_id, img)
            return
        if caption and _is_face_swap_request(caption):
            _clear_medicine_wait(context)
            with contextlib.suppress(Exception):
                _mode_track_set(user_id, "")
            await _start_faceswap_flow(update, context, img)
            return

        # 0.5) Р“РҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ / С„РҫСӮРҫвҶ’РІРёРҙРөРҫРәР»РёРҝ РёР· РҝРҫРҙРҝРёСҒРё Рә С„РҫСӮРҫ.
        if caption and _is_avatar_intent(caption):
            script = _clean_avatar_script(caption)
            _clear_medicine_wait(context)
            if script and len(script) > 8:
                context.user_data["avatar_pending_script"] = script
            _set_avatar_voice_choice_wait(context)
            await update.effective_message.reply_text("РҹРҫСҖСӮСҖРөСӮ РҝРҫР»СғСҮРөРҪ. РўРөРҝРөСҖСҢ РІСӢРұРөСҖРёСӮРө РіРҫР»РҫСҒ РҙР»СҸ СӮРөРәСҒСӮРҫРІРҫР№ РҫР·РІСғСҮРәРё РёР»Рё РҝСҖРёСҲР»РёСӮРө СҒРІРҫР№ voice/audio.", reply_markup=_avatar_voice_choice_kb("act"))
            return

        if caption and _is_photo_clip_intent(caption):
            prompt = _clean_photo_clip_prompt(caption)
            _clear_medicine_wait(context)
            if prompt:
                await _start_photo_music_clip(update, context, img, prompt)
            else:
                _set_photo_clip_wait(context)
                await update.effective_message.reply_text("РӨРҫСӮРҫ РҝРҫР»СғСҮРөРҪРҫ. РўРөРҝРөСҖСҢ РҫРҝРёСҲРёСӮРө СҒСӮРёР»СҢ РІРёРҙРөРҫРәР»РёРҝР°, РјСғР·СӢРәСғ, РҙРІРёР¶РөРҪРёРө, РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ Рё С„РҫСҖРјР°СӮ.")
            return

        if caption and _is_ai_selfie_intent(caption):
            prompt = _clean_ai_selfie_prompt(caption)
            _clear_medicine_wait(context)
            if prompt:
                await _start_ai_selfie(update, context, img, prompt)
            else:
                _set_ai_selfie_wait(context)
                await update.effective_message.reply_text("РЎРөР»С„Рё РҝРҫР»СғСҮРөРҪРҫ. РўРөРҝРөСҖСҢ РҪР°РҝРёСҲРёСӮРө, СҒ РәРөРј/РіРҙРө СҒРҙРөР»Р°СӮСҢ AI-С„РҫСӮРҫ: Р·РҪР°РјРөРҪРёСӮРҫСҒСӮСҢ, РҝРөСҖСҒРҫРҪР°Р¶, РҝСҖРөРјСҢРөСҖР°, СҖРөРәР»Р°РјР°, travel/luxury.")
            return

        if context.user_data.get("awaiting_ai_selfie_photo"):
            context.user_data.pop("awaiting_ai_selfie_photo", None)
            preset = (context.user_data.get("ai_selfie_preset_prompt", "") or "").strip()
            _set_ai_selfie_wait(context)
            if preset:
                await update.effective_message.reply_text("рҹӨі РЎРөР»С„Рё РҝРҫР»СғСҮРөРҪРҫ. РҹСҖРөСҒРөСӮ РІСӢРұСҖР°РҪ. РўРөРҝРөСҖСҢ РҪР°РҝРёСҲРёСӮРө РёРјСҸ Р·РҪР°РјРөРҪРёСӮРҫСҒСӮРё/РҝРөСҖСҒРҫРҪР°Р¶Р° Рё РҙРөСӮР°Р»Рё СҒСҶРөРҪСӢ.")
            else:
                await update.effective_message.reply_text("рҹӨі РЎРөР»С„Рё РҝРҫР»СғСҮРөРҪРҫ. РўРөРҝРөСҖСҢ РҪР°РҝРёСҲРёСӮРө, СҒ РәРөРј/РіРҙРө СҒРҙРөР»Р°СӮСҢ AI-С„РҫСӮРҫ: Р·РҪР°РјРөРҪРёСӮРҫСҒСӮСҢ, РҝРөСҖСҒРҫРҪР°Р¶, РҝСҖРөРјСҢРөСҖР°, СҖРөРәР»Р°РјР°, travel/luxury.")
            return

        if context.user_data.get("awaiting_avatar_photo"):
            context.user_data.pop("awaiting_avatar_photo", None)
            # Р•СҒР»Рё РҝРҫР»СҢР·РҫРІР°СӮРөР»СҢ Р·Р°СҖР°РҪРөРө РІСӢРұСҖР°Р» РіРҫР»РҫСҒ вҖ” СҒСҖР°Р·Сғ Р¶РҙС‘Рј СӮРөРәСҒСӮ. РҳРҪР°СҮРө СҒРҪР°СҮР°Р»Р° РҝСҖРҫСҒРёРј РІСӢРұСҖР°СӮСҢ РіРҫР»РҫСҒ.
            if context.user_data.get("avatar_tts_voice"):
                _set_avatar_wait(context)
                await update.effective_message.reply_text(f"РҹРҫСҖСӮСҖРөСӮ РҝРҫР»СғСҮРөРҪ. Р“РҫР»РҫСҒ РІСӢРұСҖР°РҪ: {_avatar_tts_voice_label(_avatar_tts_voice_get(context))}. РўРөРҝРөСҖСҢ РҝСҖРёСҲР»РёСӮРө СӮРөРәСҒСӮ РёР»Рё voice/audio РҙР»СҸ РіРҫРІРҫСҖСҸСүРөРіРҫ Р°РІР°СӮР°СҖР°.")
            else:
                _set_avatar_voice_choice_wait(context)
                await update.effective_message.reply_text("РҹРҫСҖСӮСҖРөСӮ РҝРҫР»СғСҮРөРҪ. РўРөРҝРөСҖСҢ РІСӢРұРөСҖРёСӮРө РіРҫР»РҫСҒ РҙР»СҸ СӮРөРәСҒСӮРҫРІРҫР№ РҫР·РІСғСҮРәРё РёР»Рё РҝСҖРёСҲР»РёСӮРө СҒРІРҫР№ voice/audio.", reply_markup=_avatar_voice_choice_kb("act"))
            return

        if context.user_data.get("awaiting_photo_clip_photo"):
            context.user_data.pop("awaiting_photo_clip_photo", None)
            preset_prompt = (context.user_data.pop("photo_clip_preset_prompt", "") or "").strip()
            if preset_prompt:
                await update.effective_message.reply_text("РӨРҫСӮРҫ РҝРҫР»СғСҮРөРҪРҫ. Р—Р°РҝСғСҒРәР°СҺ РІСӢРұСҖР°РҪРҪСӢР№ РҝСҖРөСҒРөСӮ РІРёРҙРөРҫРәР»РёРҝР°.")
                await _start_photo_music_clip(update, context, img, preset_prompt)
            else:
                _set_photo_clip_wait(context)
                await update.effective_message.reply_text("РӨРҫСӮРҫ РҝРҫР»СғСҮРөРҪРҫ. РўРөРҝРөСҖСҢ РҫРҝРёСҲРёСӮРө СҒСӮРёР»СҢ РІРёРҙРөРҫРәР»РёРҝР°, РјСғР·СӢРәСғ, РҙРІРёР¶РөРҪРёРө, РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ Рё С„РҫСҖРјР°СӮ.")
            return

        # 0) РӨРҫСӮРҫ РҝСҖРёСҲР»Рҫ РҝРҫСҒР»Рө РІС…РҫРҙР° РёР· РјРөРҪСҺ В«Р Р°Р·РІР»РөСҮРөРҪРёСҸ вҶ’ Р—Р°РјРөРҪРёСӮСҢ С„РҫРҪВ».
        # РқРө Р·Р°РҝСғСҒРәР°РөРј Р·Р°РјРөРҪСғ СҒСҖР°Р·Сғ, Р° РҝРҫРәР°Р·СӢРІР°РөРј РІР°СҖРёР°РҪСӮСӢ С„РҫРҪР°.
        if context.user_data.get("photo_flow") == "replacebg_menu":
            context.user_data.pop("photo_flow", None)
            await update.effective_message.reply_text(
                "рҹ–ј РӨРҫСӮРҫ РҝРҫР»СғСҮРөРҪРҫ. Р’СӢРұРөСҖРёСӮРө РҪРҫРІСӢР№ С„РҫРҪ РёР»Рё РҪР°РҝРёСҲРёСӮРө СҒРІРҫР№ РІР°СҖРёР°РҪСӮ СӮРөРәСҒСӮРҫРј:",
                reply_markup=background_presets_kb(),
            )
            return

        # 0) РЈРҙР°Р»РөРҪРёРө/Р·Р°РјРөРҪР° С„РҫРҪР° РҙРҫР»Р¶РҪСӢ РҝРөСҖРөРұРёРІР°СӮСҢ РјРөРҙРёСҶРёРҪСҒРәРёР№ РәРҫРҪСӮРөРәСҒСӮ.
        if caption and _is_remove_bg_request(caption):
            _clear_medicine_wait(context)
            _clear_removebg_wait(context)
            with contextlib.suppress(Exception):
                _mode_track_set(user_id, "")
            await _pedit_removebg(update, context, img)
            return

        if _is_waiting_removebg(context):
            _clear_medicine_wait(context)
            _clear_removebg_wait(context)
            with contextlib.suppress(Exception):
                _mode_track_set(user_id, "")
            await _pedit_removebg(update, context, img)
            return

        if caption and _is_replace_bg_request(caption):
            kind, prompt = _bg_kind_from_text(caption)
            _clear_medicine_wait(context)
            _clear_replacebg_wait(context)
            with contextlib.suppress(Exception):
                _mode_track_set(user_id, "")
            await _pedit_replacebg(update, context, img, kind=kind, prompt=prompt)
            return

        if _is_waiting_replacebg(context):
            prompt = caption or context.user_data.get("replacebg_prompt") or "СҖР°Р·РјСӢСӮСӢР№ С„РҫРҪ"
            kind, prompt = _bg_kind_from_text(prompt)
            _clear_medicine_wait(context)
            _clear_replacebg_wait(context)
            with contextlib.suppress(Exception):
                _mode_track_set(user_id, "")
            await _pedit_replacebg(update, context, img, kind=kind, prompt=prompt)
            return

        # 0) Р РөСӮСғСҲСҢ СҒРҫРұСҒСӮРІРөРҪРҪРҫРіРҫ РёР·РҫРұСҖР°Р¶РөРҪРёСҸ / СғРҙР°Р»РөРҪРёРө Р»РёСҲРҪРөР№ РҪР°РҙРҝРёСҒРё / watermark.
        # РӯСӮРҫ РҙРҫР»Р¶РҪРҫ РҝРөСҖРөРұРёРІР°СӮСҢ РјРөРҙРёСҶРёРҪСҒРәРёР№ РәРҫРҪСӮРөРәСҒСӮ, РөСҒР»Рё РҝРҫР»СҢР·РҫРІР°СӮРөР»СҢ СҸРІРҪРҫ РҝСҖРҫСҒРёСӮ СҖРөСӮСғСҲСҢ.
        if caption and _is_image_retouch_request(caption):
            _clear_medicine_wait(context)
            _clear_image_retouch_wait(context)
            with contextlib.suppress(Exception):
                _mode_track_set(user_id, "")
            await _start_image_retouch(update, context, img, caption)
            return

        if _is_waiting_image_retouch(context):
            instruction = context.user_data.get("retouch_prompt") or caption or "СғРұСҖР°СӮСҢ Р»РёСҲРҪСҺСҺ РҪР°РҙРҝРёСҒСҢ/РІРҫРҙСҸРҪРҫР№ Р·РҪР°Рә"
            _clear_medicine_wait(context)
            _clear_image_retouch_wait(context)
            with contextlib.suppress(Exception):
                _mode_track_set(user_id, "")
            await _start_image_retouch(update, context, img, instruction)
            return

        # 1) РЎР°РјСӢР№ РІСӢСҒРҫРәРёР№ РҝСҖРёРҫСҖРёСӮРөСӮ: СҸРІРҪР°СҸ РәРҫРјР°РҪРҙР° РҫР¶РёРІРёСӮСҢ/Р°РҪРёРјРёСҖРҫРІР°СӮСҢ С„РҫСӮРҫ РІ РҝРҫРҙРҝРёСҒРё.
        # РӯСӮРҫ РҙРҫР»Р¶РҪРҫ РҝРөСҖРөРұРёРІР°СӮСҢ РҙР°Р¶Рө СҖР°РҪРөРө РҫСӮРәСҖСӢСӮСӢР№ СҖР°Р·РҙРөР» В«РңРөРҙРёСҶРёРҪР°В».
        if caption and _is_photo_revival_intent(caption):
            _set_waiting_photo_revival(update, context)
            engine = _revival_engine_from_text(caption, default="runway")
            prompt = _clean_revival_prompt(caption)
            _clear_photo_revival_wait(context)
            await _start_photo_revival(update, context, engine=engine, img_bytes=img, prompt=prompt)
            return

        # 2) Р•СҒР»Рё РҝРөСҖРөРҙ С„РҫСӮРҫ РҝРҫР»СҢР·РҫРІР°СӮРөР»СҢ РіРҫР»РҫСҒРҫРј/СӮРөРәСҒСӮРҫРј СҒРҝСҖРҫСҒРёР» РҝСҖРҫ РҫР¶РёРІР»РөРҪРёРө С„РҫСӮРҫ вҖ”
        # РҝРҫРәР°Р·СӢРІР°РөРј С„РҫСӮРҫ-РјР°СҒСӮРөСҖСҒРәСғСҺ, Р° РҪРө РјРөРҙРёСҶРёРҪСҒРәРёР№ Р°РҪР°Р»РёР·.
        if _is_waiting_photo_revival(context):
            _clear_photo_revival_wait(context)
            await update.effective_message.reply_text(
                "РӨРҫСӮРҫ РҝРҫР»СғСҮРөРҪРҫ. Р’СӢРұРөСҖРёСӮРө РҙРҫСҒСӮСғРҝРҪСӢР№ РҙРІРёР¶РҫРә РҙР»СҸ РҫР¶РёРІР»РөРҪРёСҸ:",
                reply_markup=photo_quick_actions_kb(),
            )
            return

        # 3) РңРөРҙРёСҶРёРҪСҒРәР°СҸ РІРөСӮРәР° вҖ” СӮРҫР»СҢРәРҫ СҸРІРҪСӢР№ РјРөРҙ. РҝРҫРҙСҖРөР¶РёРј/РҫР¶РёРҙР°РҪРёРө РёР»Рё РјРөРҙ. СҒР»РҫРІР° РІ РҝРҫРҙРҝРёСҒРё.
        if _should_route_medical(context, user_id, caption, "photo"):
            await _medical_analyze_image(update, context, img, goal=caption or None)
            _clear_medicine_wait(context)
            with contextlib.suppress(Exception):
                _mode_track_set(user_id, "")
            return

        if caption:
            tl = caption.lower()
            # РҫР¶РёРІРёСӮСҢ С„РҫСӮРҫ вҶ’ РІСӢРұСҖР°РҪРҪСӢР№ РҙРІРёР¶РҫРә РёР· РҝРҫРҙРҝРёСҒРё РёР»Рё Runway РҝРҫ СғРјРҫР»СҮР°РҪРёСҺ
            if any(k in tl for k in ("РҫР¶РёРІРё", "РҫР¶РёРІРёСӮСҢ", "Р°РҪРёРјРёСҖСғ", "Р°РҪРёРјРёСҖРҫРІР°СӮСҢ", "СҒРҙРөР»Р°Р№ РІРёРҙРөРҫ", "revive", "animate", "image to video", "i2v")):
                engine = _revival_engine_from_text(caption, default="runway")
                prompt = _clean_revival_prompt(caption)
                await _start_photo_revival(update, context, engine=engine, img_bytes=img, prompt=prompt)
                return

            # СҖРөСӮСғСҲСҢ / СғРұСҖР°СӮСҢ РІРҫРҙСҸРҪРҫР№ Р·РҪР°Рә / Р»РёСҲРҪСҺСҺ РҪР°РҙРҝРёСҒСҢ
            if _is_image_retouch_request(caption):
                await _start_image_retouch(update, context, img, caption); return

            # СғРҙР°Р»РёСӮСҢ С„РҫРҪ
            if any(k in tl for k in ("СғРҙР°Р»Рё С„РҫРҪ", "removebg", "СғРұСҖР°СӮСҢ С„РҫРҪ")):
                await _pedit_removebg(update, context, img); return

            # Р·Р°РјРөРҪРёСӮСҢ С„РҫРҪ
            if any(k in tl for k in ("Р·Р°РјРөРҪРё С„РҫРҪ", "replacebg", "СҖР°Р·РјСӢСӮСӢР№ С„РҫРҪ", "blur")):
                kind, prompt = _bg_kind_from_text(caption)
                await _pedit_replacebg(update, context, img, kind=kind, prompt=prompt); return

            # outpaint
            if "outpaint" in tl or "СҖР°СҒСҲРёСҖ" in tl:
                await _pedit_outpaint(update, context, img); return

            # СҖР°СҒРәР°РҙСҖРҫРІРәР°
            if "СҖР°СҒРәР°РҙСҖРҫРІ" in tl or "storyboard" in tl:
                await _pedit_storyboard(update, context, img); return

            # РәР°СҖСӮРёРҪРәР° РҝРҫ РҫРҝРёСҒР°РҪРёСҺ (Luma / С„РҫР»РұСҚРә OpenAI)
            if any(k in tl for k in ("РәР°СҖСӮРёРҪ", "РёР·РҫРұСҖР°Р¶РөРҪ", "image", "img")) and any(k in tl for k in ("СҒРіРөРҪРөСҖРёСҖСғ", "СҒРҫР·РҙР°", "СҒРҙРөР»Р°Р№")):
                await _start_luma_img(update, context, caption); return

        # РөСҒР»Рё СҸРІРҪРҫР№ РәРҫРјР°РҪРҙСӢ РІ РҝРҫРҙРҝРёСҒРё РҪРөСӮ вҖ” РҝРҫРәР°Р·СӢРІР°РөРј РұСӢСҒСӮСҖСӢРө РәРҪРҫРҝРәРё
        await update.effective_message.reply_text("РӨРҫСӮРҫ РҝРҫР»СғСҮРөРҪРҫ. Р§СӮРҫ СҒРҙРөР»Р°СӮСҢ?",
                                                  reply_markup=photo_quick_actions_kb())
    except Exception as e:
        log.exception("on_photo error: %s", e)
        with contextlib.suppress(Exception):
            await update.effective_message.reply_text("вқҢ РӨРҫСӮРҫ РҪРө СҖР°СҒРҝРҫР·РҪР°РҪРҫ, РҝРҫРҝСҖРҫРұСғР№СӮРө РөСүС‘ СҖР°Р·.")

async def on_doc(update: Update, context: ContextTypes.DEFAULT_TYPE):
    try:
        if not update.message or not update.message.document:
            return
        doc = update.message.document
        mt = (doc.mime_type or "").lower()
        tg_file = await doc.get_file()
        data = await tg_file.download_as_bytearray()
        raw = bytes(data)

        caption = (update.message.caption or "").strip()

        # Presentation Studio accepts image documents and ZIP archives with many photos.
        if await _presentation_studio_get().handle_document(
            update, context, raw, doc.file_name or "file", mt, caption=caption
        ):
            return

        if context.user_data.get("awaiting_avatar_script") and (mt.startswith("audio/") or (doc.file_name or "").lower().endswith((".mp3", ".wav", ".m4a", ".aac", ".ogg"))):
            img = _get_cached_photo(update.effective_user.id)
            if not img:
                _clear_avatar_wait(context)
                await update.effective_message.reply_text("РЎРҪР°СҮР°Р»Р° Р·Р°РіСҖСғР·РёСӮРө РҝРҫСҖСӮСҖРөСӮ СҮРөР»РҫРІРөРәР°, Р·Р°СӮРөРј РҪР°Р¶РјРёСӮРө рҹ—Ј Р“РҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ.")
                return
            _clear_avatar_wait(context)
            await _start_talking_avatar(
                update, context, img,
                script_text=caption if caption and not _is_avatar_intent(caption) else _clean_avatar_script(caption),
                audio_bytes=raw,
                audio_filename=doc.file_name or "audio",
                audio_file_url=getattr(tg_file, "file_path", "") or "",
                audio_mime=mt,
            )
            return

        if mt.startswith("image/"):
            _cache_photo(update.effective_user.id, raw, getattr(tg_file, "file_path", "") or "")

            # v70: РөСҒР»Рё РҝРҫР»СҢР·РҫРІР°СӮРөР»СҢ СғР¶Рө РІСӢРұСҖР°Р» СҖРөР¶РёРј В«РҡР»РёРҝ СҒ РІРҫРәР°Р»РҫРјВ»,
            # СҒР»РөРҙСғСҺСүР°СҸ С„РҫСӮРҫРіСҖР°С„РёСҸ РҙРҫР»Р¶РҪР° РҝСҖРҫРҙРҫР»Р¶Р°СӮСҢ СҚСӮРҫСӮ СҒСҶРөРҪР°СҖРёР№, Р° РҪРө РҫСӮРәСҖСӢРІР°СӮСҢ
            # РҫРұСүРөРө РјРөРҪСҺ В«РӨРҫСӮРҫ РҝРҫР»СғСҮРөРҪРҫ. Р§СӮРҫ СҒРҙРөР»Р°СӮСҢ?В». РӯСӮРҫ СҒСӮСҖР°С…СғРөСӮ СҒР»СғСҮР°Рё,
            # РәРҫРіРҙР° Telegram/РәР»РёРөРҪСӮ РҝРҫСӮРөСҖСҸР» transient flag, РҪРҫ mode_track СҒРҫС…СҖР°РҪРёР»СҒСҸ.
            try:
                _track_now = _mode_track_get(update.effective_user.id)
            except Exception:
                _track_now = ""
            if context.user_data.get("awaiting_vocal_clip_photo") or _track_now == "vocalclip":
                context.user_data.pop("awaiting_vocal_clip_photo", None)
                _set_mode_clean(update.effective_user.id, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "vocalclip")
                _set_vocal_clip_wait(context)
                await update.effective_message.reply_text(
                    "рҹҺӨ РҹРҫСҖСӮСҖРөСӮ РҝРҫР»СғСҮРөРҪ РҙР»СҸ РәР»РёРҝР° СҒ РІРҫРәР°Р»РҫРј. РўРөРҝРөСҖСҢ РҫРҝРёСҲРёСӮРө РҝРөСҒРҪСҺ/РәР»РёРҝ: СҒСӮРёР»СҢ, СҸР·СӢРә, РҪР°СҒСӮСҖРҫРөРҪРёРө, РҝСҖРёРҝРөРІ, РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ.\n\n"
                    "Р’Р°Р¶РҪРҫ: СҖРөР¶РёРј СҖР°СҒСҒСҮРёСӮР°РҪ РҪР° РҫРҙРҪРҫРіРҫ СҮРөР»РҫРІРөРәР° РІ РәР°РҙСҖРө."
                )
                return

            if context.user_data.get("faceswap_flow") == "await_target":
                await _maybe_choose_target_face(update, context, update.effective_user.id, raw)
                return
            if context.user_data.get("faceswap_flow") == "await_source":
                await _maybe_choose_source_face(update, context, update.effective_user.id, raw)
                return
            if caption and _is_face_swap_request(caption):
                _clear_medicine_wait(context)
                with contextlib.suppress(Exception):
                    _mode_track_set(update.effective_user.id, "")
                await _start_faceswap_flow(update, context, raw)
                return

            if caption and _is_avatar_intent(caption):
                script = _clean_avatar_script(caption)
                _clear_medicine_wait(context)
                if script and len(script) > 8:
                    context.user_data["avatar_pending_script"] = script
                _set_avatar_voice_choice_wait(context)
                await update.effective_message.reply_text("РҹРҫСҖСӮСҖРөСӮ РҝРҫР»СғСҮРөРҪ. РўРөРҝРөСҖСҢ РІСӢРұРөСҖРёСӮРө РіРҫР»РҫСҒ РҙР»СҸ СӮРөРәСҒСӮРҫРІРҫР№ РҫР·РІСғСҮРәРё РёР»Рё РҝСҖРёСҲР»РёСӮРө СҒРІРҫР№ voice/audio.", reply_markup=_avatar_voice_choice_kb("act"))
                return

            if caption and _is_photo_clip_intent(caption):
                prompt = _clean_photo_clip_prompt(caption)
                _clear_medicine_wait(context)
                if prompt:
                    await _start_photo_music_clip(update, context, raw, prompt)
                else:
                    _set_photo_clip_wait(context)
                    await update.effective_message.reply_text("РӨРҫСӮРҫ РҝРҫР»СғСҮРөРҪРҫ. РўРөРҝРөСҖСҢ РҫРҝРёСҲРёСӮРө СҒСӮРёР»СҢ РІРёРҙРөРҫРәР»РёРҝР°, РјСғР·СӢРәСғ, РҙРІРёР¶РөРҪРёРө, РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ Рё С„РҫСҖРјР°СӮ.")
                return

            if caption and _is_ai_selfie_intent(caption):
                prompt = _clean_ai_selfie_prompt(caption)
                _clear_medicine_wait(context)
                if prompt:
                    await _start_ai_selfie(update, context, raw, prompt)
                else:
                    _set_ai_selfie_wait(context)
                    await update.effective_message.reply_text("РЎРөР»С„Рё РҝРҫР»СғСҮРөРҪРҫ. РўРөРҝРөСҖСҢ РҪР°РҝРёСҲРёСӮРө, СҒ РәРөРј/РіРҙРө СҒРҙРөР»Р°СӮСҢ AI-С„РҫСӮРҫ: Р·РҪР°РјРөРҪРёСӮРҫСҒСӮСҢ, РҝРөСҖСҒРҫРҪР°Р¶, РҝСҖРөРјСҢРөСҖР°, СҖРөРәР»Р°РјР°, travel/luxury.")
                return

            if context.user_data.get("awaiting_ai_selfie_photo"):
                context.user_data.pop("awaiting_ai_selfie_photo", None)
                preset = (context.user_data.get("ai_selfie_preset_prompt", "") or "").strip()
                _set_ai_selfie_wait(context)
                if preset:
                    await update.effective_message.reply_text("рҹӨі РЎРөР»С„Рё РҝРҫР»СғСҮРөРҪРҫ. РҹСҖРөСҒРөСӮ РІСӢРұСҖР°РҪ. РўРөРҝРөСҖСҢ РҪР°РҝРёСҲРёСӮРө РёРјСҸ Р·РҪР°РјРөРҪРёСӮРҫСҒСӮРё/РҝРөСҖСҒРҫРҪР°Р¶Р° Рё РҙРөСӮР°Р»Рё СҒСҶРөРҪСӢ.")
                else:
                    await update.effective_message.reply_text("рҹӨі РЎРөР»С„Рё РҝРҫР»СғСҮРөРҪРҫ. РўРөРҝРөСҖСҢ РҪР°РҝРёСҲРёСӮРө, СҒ РәРөРј/РіРҙРө СҒРҙРөР»Р°СӮСҢ AI-С„РҫСӮРҫ: Р·РҪР°РјРөРҪРёСӮРҫСҒСӮСҢ, РҝРөСҖСҒРҫРҪР°Р¶, РҝСҖРөРјСҢРөСҖР°, СҖРөРәР»Р°РјР°, travel/luxury.")
                return

            if context.user_data.get("awaiting_vocal_clip_photo"):
                context.user_data.pop("awaiting_vocal_clip_photo", None)
                _set_vocal_clip_wait(context)
                await update.effective_message.reply_text(
                    "рҹҺӨ РҹРҫСҖСӮСҖРөСӮ РҝРҫР»СғСҮРөРҪ. РўРөРҝРөСҖСҢ РҫРҝРёСҲРёСӮРө РҝРөСҒРҪСҺ/РәР»РёРҝ: СҒСӮРёР»СҢ, СҸР·СӢРә, РҪР°СҒСӮСҖРҫРөРҪРёРө, РҝСҖРёРҝРөРІ, РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ.\n\nР’Р°Р¶РҪРҫ: СҖРөР¶РёРј СҖР°СҒСҒСҮРёСӮР°РҪ РҪР° РҫРҙРҪРҫРіРҫ СҮРөР»РҫРІРөРәР° РІ РәР°РҙСҖРө."
                )
                return

            if context.user_data.get("awaiting_avatar_photo"):
                context.user_data.pop("awaiting_avatar_photo", None)
                _set_avatar_voice_choice_wait(context)
                await update.effective_message.reply_text("РҹРҫСҖСӮСҖРөСӮ РҝРҫР»СғСҮРөРҪ. РўРөРҝРөСҖСҢ РІСӢРұРөСҖРёСӮРө РіРҫР»РҫСҒ РҙР»СҸ СӮРөРәСҒСӮРҫРІРҫР№ РҫР·РІСғСҮРәРё РёР»Рё РҝСҖРёСҲР»РёСӮРө СҒРІРҫР№ voice/audio.", reply_markup=_avatar_voice_choice_kb("act"))
                return

            if context.user_data.get("awaiting_photo_clip_photo"):
                context.user_data.pop("awaiting_photo_clip_photo", None)
                preset_prompt = (context.user_data.pop("photo_clip_preset_prompt", "") or "").strip()
                if preset_prompt:
                    await update.effective_message.reply_text("РӨРҫСӮРҫ РҝРҫР»СғСҮРөРҪРҫ. Р—Р°РҝСғСҒРәР°СҺ РІСӢРұСҖР°РҪРҪСӢР№ РҝСҖРөСҒРөСӮ РІРёРҙРөРҫРәР»РёРҝР°.")
                    await _start_photo_music_clip(update, context, raw, preset_prompt)
                else:
                    _set_photo_clip_wait(context)
                    await update.effective_message.reply_text("РӨРҫСӮРҫ РҝРҫР»СғСҮРөРҪРҫ. РўРөРҝРөСҖСҢ РҫРҝРёСҲРёСӮРө СҒСӮРёР»СҢ РІРёРҙРөРҫРәР»РёРҝР°, РјСғР·СӢРәСғ, РҙРІРёР¶РөРҪРёРө, РҙР»РёСӮРөР»СҢРҪРҫСҒСӮСҢ Рё С„РҫСҖРјР°СӮ.")
                return

            if context.user_data.get("photo_flow") == "replacebg_menu":
                context.user_data.pop("photo_flow", None)
                await update.effective_message.reply_text(
                    "рҹ–ј РҳР·РҫРұСҖР°Р¶РөРҪРёРө РҝРҫР»СғСҮРөРҪРҫ. Р’СӢРұРөСҖРёСӮРө РҪРҫРІСӢР№ С„РҫРҪ. Р—Р°РјРөРҪР° РІСӢРҝРҫР»РҪСҸРөСӮСҒСҸ РІ 2 СҚСӮР°РҝР°: РІСӢСҖРөР·Р°СҺ СҮРөР»РҫРІРөРәР°, Р·Р°СӮРөРј РҝРҫРҙСҒСӮР°РІР»СҸСҺ РҪРҫРІСӢР№ С„РҫРҪ РұРөР· РҝРөСҖРөСҖРёСҒРҫРІРәРё Р»РёСҶР° Рё РҫРҙРөР¶РҙСӢ.",
                    reply_markup=background_presets_kb(),
                )
                return

            if caption and _is_remove_bg_request(caption):
                _clear_medicine_wait(context)
                _clear_removebg_wait(context)
                with contextlib.suppress(Exception):
                    _mode_track_set(update.effective_user.id, "")
                await _pedit_removebg(update, context, raw)
                return

            if _is_waiting_removebg(context):
                _clear_medicine_wait(context)
                _clear_removebg_wait(context)
                with contextlib.suppress(Exception):
                    _mode_track_set(update.effective_user.id, "")
                await _pedit_removebg(update, context, raw)
                return

            if caption and _is_replace_bg_request(caption):
                kind, prompt = _bg_kind_from_text(caption)
                _clear_medicine_wait(context)
                _clear_replacebg_wait(context)
                with contextlib.suppress(Exception):
                    _mode_track_set(update.effective_user.id, "")
                await _pedit_replacebg(update, context, raw, kind=kind, prompt=prompt)
                return

            if _is_waiting_replacebg(context):
                prompt = caption or context.user_data.get("replacebg_prompt") or "СҖР°Р·РјСӢСӮСӢР№ С„РҫРҪ"
                kind, prompt = _bg_kind_from_text(prompt)
                _clear_medicine_wait(context)
                _clear_replacebg_wait(context)
                with contextlib.suppress(Exception):
                    _mode_track_set(update.effective_user.id, "")
                await _pedit_replacebg(update, context, raw, kind=kind, prompt=prompt)
                return

            if caption and _is_image_retouch_request(caption):
                _clear_medicine_wait(context)
                _clear_image_retouch_wait(context)
                with contextlib.suppress(Exception):
                    _mode_track_set(update.effective_user.id, "")
                await _start_image_retouch(update, context, raw, caption)
                return

            if _is_waiting_image_retouch(context):
                instruction = context.user_data.get("retouch_prompt") or caption or "СғРұСҖР°СӮСҢ Р»РёСҲРҪСҺСҺ РҪР°РҙРҝРёСҒСҢ/РІРҫРҙСҸРҪРҫР№ Р·РҪР°Рә"
                _clear_medicine_wait(context)
                _clear_image_retouch_wait(context)
                with contextlib.suppress(Exception):
                    _mode_track_set(update.effective_user.id, "")
                await _start_image_retouch(update, context, raw, instruction)
                return

            if caption and _is_photo_revival_intent(caption):
                _set_waiting_photo_revival(update, context)
                engine = _revival_engine_from_text(caption, default="runway")
                prompt = _clean_revival_prompt(caption)
                _clear_photo_revival_wait(context)
                await _start_photo_revival(update, context, engine=engine, img_bytes=raw, prompt=prompt)
                return

            if _is_waiting_photo_revival(context):
                _clear_photo_revival_wait(context)
                await update.effective_message.reply_text("РҳР·РҫРұСҖР°Р¶РөРҪРёРө РҝРҫР»СғСҮРөРҪРҫ РәР°Рә РҙРҫРәСғРјРөРҪСӮ. Р’СӢРұРөСҖРёСӮРө РҙРҫСҒСӮСғРҝРҪСӢР№ РҙРІРёР¶РҫРә РҙР»СҸ РҫР¶РёРІР»РөРҪРёСҸ:", reply_markup=photo_quick_actions_kb())
                return

            if _should_route_medical(context, update.effective_user.id, caption, doc.file_name or "image"):
                await _medical_analyze_image(update, context, raw, goal=caption or None)
                _clear_medicine_wait(context)
                with contextlib.suppress(Exception):
                    _mode_track_set(update.effective_user.id, "")
                return
            await update.effective_message.reply_text("РҳР·РҫРұСҖР°Р¶РөРҪРёРө РҝРҫР»СғСҮРөРҪРҫ РәР°Рә РҙРҫРәСғРјРөРҪСӮ. Р§СӮРҫ СҒРҙРөР»Р°СӮСҢ?", reply_markup=photo_quick_actions_kb())
            return

        text, kind = extract_text_from_document(raw, doc.file_name or "file")
        if not (text or "").strip():
            await update.effective_message.reply_text(f"РқРө СғРҙР°Р»РҫСҒСҢ РёР·РІР»РөСҮСҢ СӮРөРәСҒСӮ РёР· {kind}.")
            return

        # A text document can serve as the brief or revision for an active presentation project.
        if await _presentation_studio_get().handle_text(update, context, text):
            return

        goal = (update.message.caption or "").strip() or None
        if _should_route_medical(context, update.effective_user.id, caption, doc.file_name or "file"):
            await _medical_analyze_text(update, context, text, goal=goal)
            _clear_medicine_wait(context)
            with contextlib.suppress(Exception):
                _mode_track_set(update.effective_user.id, "")
            return

        await update.effective_message.reply_text(f"рҹ“„ РҳР·РІР»РөРәР°СҺ СӮРөРәСҒСӮ ({kind}), РіРҫСӮРҫРІР»СҺ РәРҫРҪСҒРҝРөРәСӮвҖҰ")
        summary = await summarize_long_text(text, query=goal)
        summary = summary or "Р“РҫСӮРҫРІРҫ."
        await update.effective_message.reply_text(summary)
        await maybe_tts_reply(update, context, summary[:TTS_MAX_CHARS])
    except Exception as e:
        log.exception("on_doc error: %s", e)
        with contextlib.suppress(Exception):
            await update.effective_message.reply_text("вқҢ РӨРҫСӮРҫ/РҙРҫРәСғРјРөРҪСӮ РҪРө СҖР°СҒРҝРҫР·РҪР°РҪ, РҝРҫРҝСҖРҫРұСғР№СӮРө РөСүС‘ СҖР°Р·.")

async def on_voice(update: Update, context: ContextTypes.DEFAULT_TYPE):
    try:
        if not update.message or not update.message.voice:
            return
        vf = await update.message.voice.get_file()
        bio = BytesIO(await vf.download_as_bytearray())
        bio.seek(0)
        setattr(bio, "name", f"voice.ogg")
        await context.bot.send_chat_action(update.effective_chat.id, ChatAction.TYPING)
        text = await transcribe_audio(bio, "voice.ogg")
        if not text:
            await update.effective_message.reply_text("РқРө СғРҙР°Р»РҫСҒСҢ СҖР°СҒРҝРҫР·РҪР°СӮСҢ СҖРөСҮСҢ.")
            return
        await on_text(update, context, manual_text=text)
    except Exception as e:
        log.exception("on_voice error: %s", e)
        with contextlib.suppress(Exception):
            await update.effective_message.reply_text("РһСҲРёРұРәР° РҝСҖРё РҫРұСҖР°РұРҫСӮРәРө voice.")

async def on_audio(update: Update, context: ContextTypes.DEFAULT_TYPE):
    try:
        if not update.message or not update.message.audio:
            return
        af = await update.message.audio.get_file()
        filename = update.message.audio.file_name or "audio.mp3"
        bio = BytesIO(await af.download_as_bytearray())
        bio.seek(0)
        setattr(bio, "name", filename)
        await context.bot.send_chat_action(update.effective_chat.id, ChatAction.TYPING)
        text = await transcribe_audio(bio, filename)
        if not text:
            await update.effective_message.reply_text("РқРө СғРҙР°Р»РҫСҒСҢ СҖР°СҒРҝРҫР·РҪР°СӮСҢ СҖРөСҮСҢ РёР· Р°СғРҙРёРҫ.")
            return
        await on_text(update, context, manual_text=text)
    except Exception as e:
        log.exception("on_audio error: %s", e)
        with contextlib.suppress(Exception):
            await update.effective_message.reply_text("РһСҲРёРұРәР° РҝСҖРё РҫРұСҖР°РұРҫСӮРәРө Р°СғРҙРёРҫ.")


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РһРұСҖР°РұРҫСӮСҮРёРә РҫСҲРёРұРҫРә PTB в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
async def on_error(update: object, context_: ContextTypes.DEFAULT_TYPE):
    log.exception("Unhandled error: %s", context_.error)
    try:
        if isinstance(update, Update) and update.effective_message:
            await update.effective_message.reply_text("РЈРҝСҒ, РҝСҖРҫРёР·РҫСҲР»Р° РҫСҲРёРұРәР°. РҜ СғР¶Рө СҖР°Р·РұРёСҖР°СҺСҒСҢ.")
    except Exception:
        pass


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ Р РҫСғСӮРөСҖСӢ РҙР»СҸ СӮРөРәСҒСӮРҫРІСӢС… РәРҪРҫРҝРҫРә/СҖРөР¶РёРјРҫРІ в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
async def on_btn_engines(update: Update, context: ContextTypes.DEFAULT_TYPE):
    return await cmd_engines(update, context)

async def on_btn_balance(update: Update, context: ContextTypes.DEFAULT_TYPE):
    return await cmd_balance(update, context)

async def on_btn_plans(update: Update, context: ContextTypes.DEFAULT_TYPE):
    user_id = update.effective_user.id
    text = _plans_overview_text(user_id)
    await update.effective_message.reply_text(text, reply_markup=plans_root_kb())

async def on_mode_school_text(update: Update, context: ContextTypes.DEFAULT_TYPE):
    txt = (
        "рҹҺ“ *РЈСҮС‘РұР°*\n"
        "РҹРҫРјРҫРіСғ: РәРҫРҪСҒРҝРөРәСӮСӢ РёР· PDF/EPUB/DOCX/TXT, СҖР°Р·РұРҫСҖ Р·Р°РҙР°СҮ РҝРҫСҲР°РіРҫРІРҫ, СҚСҒСҒРө/СҖРөС„РөСҖР°СӮСӢ, РјРёРҪРё-РәРІРёР·СӢ.\n\n"
        "_Р‘СӢСҒСӮСҖСӢРө РҙРөР№СҒСӮРІРёСҸ:_\n"
        "вҖў Р Р°Р·РҫРұСҖР°СӮСҢ PDF вҶ’ РәРҫРҪСҒРҝРөРәСӮ\n"
        "вҖў РЎРҫРәСҖР°СӮРёСӮСҢ РІ СҲРҝР°СҖРіР°Р»РәСғ\n"
        "вҖў РһРұСҠСҸСҒРҪРёСӮСҢ СӮРөРјСғ СҒ РҝСҖРёРјРөСҖР°РјРё\n"
        "вҖў РҹР»Р°РҪ РҫСӮРІРөСӮР° / РҝСҖРөР·РөРҪСӮР°СҶРёРё"
    )
    await update.effective_message.reply_text(txt, parse_mode="Markdown")

async def on_mode_work_text(update: Update, context: ContextTypes.DEFAULT_TYPE):
    txt = (
        "рҹ’ј *Р Р°РұРҫСӮР°/Р‘РёР·РҪРөСҒ*\n"
        "РҹРёСҒСҢРјР°, РҡРҹ, РҙРҫРіРҫРІРҫСҖРҪСӢРө СҮРөСҖРҪРҫРІРёРәРё, Р°РҪР°Р»РёСӮРёРәР°, РҝР»Р°РҪСӢ, РұСҖРёС„СӢ, РҝСҖРөР·РөРҪСӮР°СҶРёРё, PDF-РәР°СӮР°Р»РҫРіРё, Р»РҫРіРҫСӮРёРҝСӢ Рё СҖРөСӮСғСҲСҢ РұРёР·РҪРөСҒ-С„РҫСӮРҫ.\n\n"
        "_Р‘СӢСҒСӮСҖСӢРө РҙРөР№СҒСӮРІРёСҸ:_\n"
        "вҖў рҹ“„ РҹРёСҒСҢРјРҫ / РҙРҫРәСғРјРөРҪСӮ\n"
        "вҖў рҹ“Ҡ РЎРҫР·РҙР°СӮСҢ РҝСҖРөР·РөРҪСӮР°СҶРёСҺ\n"
        "вҖў рҹ“• РЎРҫР·РҙР°СӮСҢ PDF-РәР°СӮР°Р»РҫРі\n"
        "вҖў рҹҺЁ РЎРҫР·РҙР°СӮСҢ Р»РҫРіРҫСӮРёРҝ\n"
        "вҖў рҹ§Ҫ РЈРҙР°Р»РёСӮСҢ РІРҫРҙСҸРҪРҫР№ Р·РҪР°Рә / РҪР°РҙРҝРёСҒСҢ СҒ С„РҫСӮРҫ"
    )
    await update.effective_message.reply_text(txt, parse_mode="Markdown", reply_markup=_mode_kb("work"))

async def on_mode_fun_text(update: Update, context: ContextTypes.DEFAULT_TYPE):
    txt = (
        "рҹ”Ҙ *Р Р°Р·РІР»РөСҮРөРҪРёСҸ*\n"
        "Р—РҙРөСҒСҢ РұСӢСҒСӮСҖСӢРө СӮРІРҫСҖСҮРөСҒРәРёРө СҒСҶРөРҪР°СҖРёРё: РҫР¶РёРІРёСӮСҢ С„РҫСӮРҫРіСҖР°С„РёСҺ, СҒРҙРөР»Р°СӮСҢ РіРҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ, СҒРҫР·РҙР°СӮСҢ С„РҫСӮРҫвҶ’РІРёРҙРөРҫРәР»РёРҝ СҒ РјСғР·СӢРәРҫР№, РәР»РёРҝ СҒ РІРҫРәР°Р»РҫРј РҙР»СҸ 1 СҮРөР»РҫРІРөРәР°, РІРёРҙРөРҫ РҝРҫ СӮРөРәСҒСӮСғ/РіРҫР»РҫСҒСғ, "
        "Р·Р°РјРөРҪРёСӮСҢ Р»РёСҶРҫ, СғРҙР°Р»РёСӮСҢ РёР»Рё Р·Р°РјРөРҪРёСӮСҢ С„РҫРҪ, СҒРҙРөР»Р°СӮСҢ Reels/Shorts, СҒРҫР·РҙР°СӮСҢ РјРёРҪРё-С„РёР»СҢРј, РҝСҖРёРҙСғРјР°СӮСҢ РёРҙРөРё, СҒСҶРөРҪР°СҖРёР№, РёРіСҖСғ РёР»Рё РәРІРёР·.\n\n"
        "Р’СӢРұРөСҖРё РҙРөР№СҒСӮРІРёРө РҪРёР¶Рө РёР»Рё РҪР°РҝРёСҲРё СҒРІРҫРұРҫРҙРҪСӢР№ Р·Р°РҝСҖРҫСҒ."
    )
    await update.effective_message.reply_text(txt, parse_mode="Markdown", reply_markup=_fun_quick_kb())

# в”Җв”Җв”Җв”Җв”Җ РҡР»Р°РІРёР°СӮСғСҖР° В«Р Р°Р·РІР»РөСҮРөРҪРёСҸВ» СҒ РҪРҫРІСӢРјРё РәРҪРҫРҝРәР°РјРё в”Җв”Җв”Җв”Җв”Җ
def _fun_quick_kb() -> InlineKeyboardMarkup:
    rows = [
        [InlineKeyboardButton("рҹӘ„ РһР¶РёРІРёСӮСҢ С„РҫСӮРҫ", callback_data="fun:revive")],
        [InlineKeyboardButton("рҹ—Ј Р“РҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ", callback_data="fun:avatar")],
        [InlineKeyboardButton("рҹҺө РӨРҫСӮРҫ вҶ’ РІРёРҙРөРҫРәР»РёРҝ", callback_data="fun:photoclip")],
        [InlineKeyboardButton("рҹҺӨ РҡР»РёРҝ СҒ РІРҫРәР°Р»РҫРј (1 СҮРөР»РҫРІРөРә)", callback_data="fun:vocalclip")],
        [InlineKeyboardButton("рҹҺ¬ Р’РёРҙРөРҫ РҝРҫ СӮРөРәСҒСӮСғ/РіРҫР»РҫСҒСғ", callback_data="fun:textvideo")],
        [InlineKeyboardButton("рҹӨі AI-СҒРөР»С„Рё СҒРҫ Р·РІРөР·РҙРҫР№", callback_data="fun:aiselfie")],
        [InlineKeyboardButton("рҹҺӯ Р—Р°РјРөРҪР° Р»РёСҶР° РҪР° С„РҫСӮРҫ", callback_data="fun:faceswap")],
        [InlineKeyboardButton("рҹ§ј РЈРҙР°Р»РёСӮСҢ С„РҫРҪ РҪР° С„РҫСӮРҫ", callback_data="fun:removebg")],
        [InlineKeyboardButton("рҹ–ј Р—Р°РјРөРҪРёСӮСҢ С„РҫРҪ РҪР° С„РҫСӮРҫ", callback_data="fun:replacebg")],
        [InlineKeyboardButton("рҹ“ұ Reels / Shorts", callback_data="fun:reels")],
        [InlineKeyboardButton("рҹҺһ РЎРҫР·РҙР°СӮСҢ РјРёРҪРё-С„РёР»СҢРј", callback_data="fun:film")],
        [InlineKeyboardButton("рҹҺ¬ РЎСҶРөРҪР°СҖРёР№ / РәР°РҙСҖСӢ", callback_data="fun:storyboard")],
        [InlineKeyboardButton("рҹҺө РңСғР·СӢРәР° / РҝРөСҒРҪСҸ", callback_data="fun:music")],
        [InlineKeyboardButton("рҹҺ® РҳРіСҖСӢ / РәРІРёР·", callback_data="fun:quiz")],
        [InlineKeyboardButton("рҹҺӯ РҳРҙРөРё РҙР»СҸ РҙРҫСҒСғРіР°", callback_data="fun:ideas")],
        [InlineKeyboardButton("рҹ“қ РЎРІРҫРұРҫРҙРҪСӢР№ Р·Р°РҝСҖРҫСҒ", callback_data="fun:free")],
        [InlineKeyboardButton("в¬…пёҸ РқР°Р·Р°Рҙ", callback_data="fun:back")],
    ]
    return InlineKeyboardMarkup(rows)

# в”Җв”Җв”Җв”Җв”Җ РһРұСҖР°РұРҫСӮСҮРёРә РұСӢСҒСӮСҖСӢС… РҙРөР№СҒСӮРІРёР№ В«Р Р°Р·РІР»РөСҮРөРҪРёСҸВ» (fallback-friendly) в”Җв”Җв”Җв”Җв”Җ
async def on_cb_fun(update: Update, context: ContextTypes.DEFAULT_TYPE):
    q = update.callback_query
    data = (q.data or "").strip()
    action = data.split(":", 1)[1] if ":" in data else ""

    async def _try_call(*fn_names, **kwargs):
        fn = _pick_first_defined(*fn_names)
        if callable(fn):
            return await fn(update, context, **kwargs)
        return None

    if action == "avatar":
        await q.answer("Р“РҫРІРҫСҖСҸСүРёР№ Р°РІР°СӮР°СҖ")
        _clear_transient_flows(context)
        _set_mode_clean(q.from_user.id, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "avatar")
        await q.message.reply_text(_avatar_menu_text(), parse_mode="Markdown", reply_markup=_avatar_action_kb("fun"))
        return

    if action.startswith("av_voice_"):
        voice = action.rsplit("_", 1)[-1].strip()
        context.user_data["avatar_tts_voice"] = voice
        context.user_data.pop("awaiting_avatar_voice_choice", None)
        if _get_cached_photo(q.from_user.id):
            _set_avatar_wait(context)
            await q.answer(f"Р“РҫР»РҫСҒ: {voice}")
            await q.message.reply_text(f"вң… Р”Р»СҸ Р°РІР°СӮР°СҖР° РІСӢРұСҖР°РҪ РіРҫР»РҫСҒ: {_avatar_tts_voice_label(voice)}. РЁР°Рі 3/3: РҝСҖРёСҲР»РёСӮРө СӮРөРәСҒСӮ, РәРҫСӮРҫСҖСӢР№ РҙРҫР»Р¶РөРҪ РҝСҖРҫРёР·РҪРөСҒСӮРё Р°РІР°СӮР°СҖ.")
        else:
            context.user_data["awaiting_avatar_photo"] = True
            await q.answer(f"Р“РҫР»РҫСҒ: {voice}")
            await q.message.reply_text(f"вң… Р“РҫР»РҫСҒ РІСӢРұСҖР°РҪ: {_avatar_tts_voice_label(voice)}. РўРөРҝРөСҖСҢ РҝСҖРёСҲР»РёСӮРө РҝРҫСҖСӮСҖРөСӮ СҮРөР»РҫРІРөРәР°.")
        return

    if action == "avatar_upload":
        await q.answer("Р—Р°РіСҖСғР·РёСӮРө РҝРҫСҖСӮСҖРөСӮ")
        await _handle_avatar_upload_choice(update, context, q, prefix="fun")
        return

    if action in {"avatar_last", "avatar_text"}:
        await _handle_avatar_script_choice(update, context, q, prefix="fun", voice_mode=False)
        return

    if action == "avatar_voice":
        await _handle_avatar_script_choice(update, context, q, prefix="fun", voice_mode=True)
        return

    if action == "vocalclip":
        await q.answer("РҡР»РёРҝ СҒ РІРҫРәР°Р»РҫРј")
        _clear_transient_flows(context)
        _set_mode_clean(q.from_user.id, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "vocalclip")
        await q.message.reply_text(_vocal_clip_menu_text(), parse_mode="Markdown", reply_markup=_vocal_clip_action_kb("fun"))
        return

    if action == "vocalclip_upload":
        await q.answer("Р—Р°РіСҖСғР·РёСӮРө РҝРҫСҖСӮСҖРөСӮ")
        await _handle_vocalclip_upload_choice(update, context, q, prefix="fun")
        return

    if action in {"vocalclip_last", "vocalclip_prompt"}:
        await _handle_vocalclip_prompt_choice(update, context, q, prefix="fun")
        return

    if action == "textvideo":
        await q.answer("Р’РёРҙРөРҫ РҝРҫ СӮРөРәСҒСӮСғ/РіРҫР»РҫСҒСғ")
        _clear_transient_flows(context)
        _set_mode_clean(q.from_user.id, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "textvideo")
        await q.message.reply_text(_textvideo_menu_text(), parse_mode="Markdown", reply_markup=_textvideo_action_kb("fun"))
        return

    if action == "tv_engine_sora":
        await _handle_textvideo_engine_choice(update, context, q, "sora", prefix="fun")
        return

    if action == "tv_engine_kling":
        await _handle_textvideo_engine_choice(update, context, q, "kling", prefix="fun")
        return

    if action == "tv_engine_runway":
        await _handle_textvideo_engine_choice(update, context, q, "runway", prefix="fun")
        return

    if action == "tv_prompt":
        await _handle_textvideo_prompt_choice(update, context, q, prefix="fun")
        return

    if action == "photoclip":
        await q.answer("РӨРҫСӮРҫ вҶ’ РІРёРҙРөРҫРәР»РёРҝ")
        _clear_transient_flows(context)
        _set_mode_clean(q.from_user.id, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "photoclip")
        await q.message.reply_text(_photoclip_menu_text(), parse_mode="Markdown", reply_markup=_photoclip_action_kb("fun"))
        return

    if action == "photoclip_upload":
        await q.answer("Р—Р°РіСҖСғР·РёСӮРө С„РҫСӮРҫ")
        await _handle_photoclip_upload_choice(update, context, q, prefix="fun")
        return

    if action in {"photoclip_last", "photoclip_custom"}:
        await _handle_photoclip_prompt_choice(update, context, q, prefix="fun")
        return

    if action.startswith("pc_preset_"):
        kind = action.rsplit("_", 1)[-1]
        await _handle_photoclip_preset_choice(update, context, q, kind, prefix="fun")
        return

    if action == "faceswap":
        await q.answer("Р—Р°РјРөРҪР° Р»РёСҶР°")
        _clear_transient_flows(context)
        _set_mode_clean(q.from_user.id, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "faceswap")
        await _start_faceswap_flow(update, context, None, use_cached=False)
        return

    if action == "removebg":
        await q.answer("РЈРҙР°Р»РөРҪРёРө С„РҫРҪР°")
        _clear_transient_flows(context)
        _set_mode_clean(q.from_user.id, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "removebg")
        img = _get_cached_photo(q.from_user.id)
        if img:
            await q.message.reply_text("рҹ§ј РҳСҒРҝРҫР»СҢР·СғСҺ РҝРҫСҒР»РөРҙРҪРөРө Р·Р°РіСҖСғР¶РөРҪРҪРҫРө С„РҫСӮРҫ Рё СғРҙР°Р»СҸСҺ С„РҫРҪ.")
            await _pedit_removebg(update, context, img)
        else:
            _set_waiting_removebg(context)
            await q.message.reply_text("рҹ§ј РҹСҖРёСҲР»РёСӮРө С„РҫСӮРҫ вҖ” СғРҙР°Р»СҺ С„РҫРҪ Рё РІРөСҖРҪСғ PNG СҒ РҝСҖРҫР·СҖР°СҮРҪРҫР№ РҝРҫРҙР»РҫР¶РәРҫР№.", reply_markup=_fun_quick_kb())
        return

    if action == "replacebg":
        await q.answer("Р—Р°РјРөРҪР° С„РҫРҪР°")
        _clear_transient_flows(context)
        _set_mode_clean(q.from_user.id, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "replacebg")
        img = _get_cached_photo(q.from_user.id)
        if img:
            await q.message.reply_text("рҹ–ј Р’СӢРұРөСҖРёСӮРө РҪРҫРІСӢР№ С„РҫРҪ РҙР»СҸ РҝРҫСҒР»РөРҙРҪРөРіРҫ Р·Р°РіСҖСғР¶РөРҪРҪРҫРіРҫ С„РҫСӮРҫ:", reply_markup=background_presets_kb())
        else:
            context.user_data["photo_flow"] = "replacebg_menu"
            await q.message.reply_text("рҹ–ј РҹСҖРёСҲР»РёСӮРө С„РҫСӮРҫ. РҹРҫСҒР»Рө Р·Р°РіСҖСғР·РәРё СҸ РҝРҫРәР°Р¶Сғ РІР°СҖРёР°РҪСӮСӢ С„РҫРҪР°: РҝР»СҸР¶, РіРҫСҖСӢ, РҝСҖРёСҖРҫРҙР°, РіРҫСҖРҫРҙ РёР»Рё СҒРІРҫР№ СӮРөРәСҒСӮ.", reply_markup=_fun_quick_kb())
        return

    if action == "revive":
        if await _try_call("revive_old_photo_flow", "do_revive_photo"):
            return
        _set_waiting_photo_revival(update, context)
        await q.answer("РһР¶РёРІР»РөРҪРёРө С„РҫСӮРҫ")
        await q.message.reply_text(_fun_revive_help_text(), parse_mode="Markdown", reply_markup=_fun_quick_kb())
        return

    if action in {"smartreels", "reels"}:
        if await _try_call("smart_reels_from_video", "video_sense_reels"):
            return
        _clear_transient_flows(context)
        _set_mode_clean(q.from_user.id, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "fun_reels")
        context.user_data["awaiting_reels_material"] = True
        await q.answer("Reels / Shorts")
        await q.message.reply_text(_fun_reels_help_text(), parse_mode="Markdown", reply_markup=_fun_quick_kb())
        return

    if action == "film":
        _clear_transient_flows(context)
        _set_mode_clean(q.from_user.id, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "fun_film")
        context.user_data["awaiting_film_material"] = True
        await q.answer("РЎРҫР·РҙР°СӮСҢ С„РёР»СҢРј")
        await q.message.reply_text(_fun_film_help_text(), parse_mode="Markdown", reply_markup=_fun_quick_kb())
        return

    if action == "clip":
        if await _try_call("start_runway_flow", "luma_make_clip", "runway_make_clip"):
            return
        _clear_transient_flows(context)
        _set_mode_clean(q.from_user.id, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "fun_reels")
        context.user_data["awaiting_reels_material"] = True
        await q.answer()
        await q.message.reply_text(_fun_reels_help_text(), parse_mode="Markdown", reply_markup=_fun_quick_kb())
        return

    if action == "img":
        if await _try_call("cmd_img", "midjourney_flow", "images_make"):
            return
        await q.answer()
        await q.message.reply_text("Р’РІРөРҙРё /img Рё СӮРөРјСғ РәР°СҖСӮРёРҪРәРё, РёР»Рё РҝСҖРёСҲР»Рё СҖРөС„СӢ.", reply_markup=_fun_quick_kb())
        return

    if action == "storyboard":
        if await _try_call("start_storyboard", "storyboard_make"):
            return
        await q.answer()
        await q.message.reply_text("РқР°РҝРёСҲРё СӮРөРјСғ СҲРҫСҖСӮР° вҖ” РҪР°РәРёРҙР°СҺ СҒСӮСҖСғРәСӮСғСҖСғ Рё СҖР°СҒРәР°РҙСҖРҫРІРәСғ.", reply_markup=_fun_quick_kb())
        return

    if action == "music":
        _clear_transient_flows(context)
        _set_mode_clean(q.from_user.id, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "suno_music")
        context.user_data["awaiting_suno_brief"] = True
        await q.answer("РңСғР·СӢРәР° / Suno")
        await _show_suno_help_from_callback(q, context, reply_markup=_suno_menu_kb(), submenu=True)
        return

    if action in {"ideas", "quiz", "speech", "free", "back"}:
        await q.answer()
        await q.message.reply_text(
            "Р“РҫСӮРҫРІ! РқР°РҝРёСҲРё Р·Р°РҙР°СҮСғ РёР»Рё РІСӢРұРөСҖРё РәРҪРҫРҝРәСғ РІСӢСҲРө.",
            reply_markup=_fun_quick_kb()
        )
        return

    await q.answer()

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ Р РҫСғСӮРөСҖСӢ-РәРҪРҫРҝРәРё СҖРөР¶РёРјРҫРІ (РөРҙРёРҪР°СҸ СӮРҫСҮРәР° РІС…РҫРҙР°) в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
async def on_btn_study(update: Update, context: ContextTypes.DEFAULT_TYPE):
    _clear_transient_flows(context)
    _set_mode_clean(update.effective_user.id, "РЈСҮС‘РұР°", "")
    fn = globals().get("_send_mode_menu")
    if callable(fn):
        return await fn(update, context, "study")
    return await on_mode_school_text(update, context)

async def on_btn_work(update: Update, context: ContextTypes.DEFAULT_TYPE):
    _clear_transient_flows(context)
    _set_mode_clean(update.effective_user.id, "Р Р°РұРҫСӮР°/Р‘РёР·РҪРөСҒ", "")
    fn = globals().get("_send_mode_menu")
    if callable(fn):
        return await fn(update, context, "work")
    return await on_mode_work_text(update, context)

async def on_btn_fun(update: Update, context: ContextTypes.DEFAULT_TYPE):
    _clear_transient_flows(context)
    _set_mode_clean(update.effective_user.id, "Р Р°Р·РІР»РөСҮРөРҪРёСҸ", "")
    fn = globals().get("_send_mode_menu")
    if callable(fn):
        return await fn(update, context, "fun")
    return await on_mode_fun_text(update, context)

async def on_btn_medicine(update: Update, context: ContextTypes.DEFAULT_TYPE):
    _set_medical_waiting(update, context, "")
    await update.effective_message.reply_text(_medical_menu_text(), reply_markup=medicine_kb())

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ РҹРҫР·РёСӮРёРІРҪСӢР№ Р°РІСӮРҫ-РҫСӮРІРөСӮ РҝСҖРҫ РІРҫР·РјРҫР¶РҪРҫСҒСӮРё (СӮРөРәСҒСӮ/РіРҫР»РҫСҒ) в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
_CAPS_PATTERN = re.compile(
    r"(СғРјРөРөСҲСҢ|РјРҫР¶РөСҲСҢ|РҙРөР»Р°РөСҲСҢ|Р°РҪР°Р»РёР·РёСҖСғРөСҲСҢ|СҖР°РұРҫСӮР°РөСҲСҢ|РҝРҫРҙРҙРөСҖР¶РёРІР°РөСҲСҢ|СғРјРөРөСӮ\s+Р»Рё|РјРҫР¶РөСӮ\s+Р»Рё|РјРҫР¶РҪРҫ\s+Р»Рё)"
    r".{0,160}"
    r"(pdf|epub|fb2|docx|txt|РәРҪРёРі|РәРҪРёРіР°|РёР·РҫРұСҖР°Р¶РөРҪ|С„РҫСӮРҫ|С„РҫСӮРҫРіСҖР°С„|РәР°СҖСӮРёРҪ|РҫР¶РёРІ|Р°РҪРёРјРёСҖ|"
    r"image|jpeg|png|video|РІРёРҙРөРҫ|mp4|mov|Р°СғРҙРёРҫ|audio|mp3|wav|"
    r"РјРөРҙРёСҶРёРҪ|РјРөРҙРәР°СҖСӮ|РІСӢРҝРёСҒРә|Р°РҪР°РјРҪРөР·|Р°РҪР°Р»РёР·|СҒРҪРёРјРҫРә|РјСҖСӮ|РәСӮ|Р·Р°РәР»СҺСҮРөРҪРё|РІСҖР°СҮРөРұРҪ|РҙРёР°РіРҪРҫР·|СғР·Рё|СҖРөРҪСӮРіРөРҪ|Р»РёСҶРҫ|Р»РёСҶР°|Р»РёСҶ|faceswap|face\s*swap|РјСғР·СӢРә|РҝРөСҒРҪ|СӮСҖРөРә|suno)",
    re.I | re.S,
)

async def on_capabilities_qa(update: Update, context: ContextTypes.DEFAULT_TYPE):
    incoming_text = (getattr(update.effective_message, "text", "") or "").strip()
    if _is_photo_revival_question(incoming_text) or _is_photo_revival_intent(incoming_text):
        _set_waiting_photo_revival(update, context)
    cap = capability_answer(incoming_text)
    if cap:
        await update.effective_message.reply_text(cap, reply_markup=main_kb)
        return

    msg = (
        "Р”Р°, СғРјРөСҺ СҖР°РұРҫСӮР°СӮСҢ СҒ С„Р°Р№Р»Р°РјРё, РјРөРҙРёР° Рё РјРөРҙРёСҶРёРҪСҒРәРёРјРё РјР°СӮРөСҖРёР°Р»Р°РјРё:\n"
        "вҖў рҹ“„ Р”РҫРәСғРјРөРҪСӮСӢ: PDF/EPUB/FB2/DOCX/TXT вҖ” РәРҫРҪСҒРҝРөРәСӮ, СҖРөР·СҺРјРө, РёР·РІР»РөСҮРөРҪРёРө СӮР°РұР»РёСҶ, РҝСҖРҫРІРөСҖРәР° С„Р°РәСӮРҫРІ.\n"
        "вҖў рҹ–ј РҳР·РҫРұСҖР°Р¶РөРҪРёСҸ: Р°РҪР°Р»РёР·/РҫРҝРёСҒР°РҪРёРө, СғРҙР°Р»РөРҪРёРө Рё Р·Р°РјРөРҪР° С„РҫРҪР°, Р·Р°РјРөРҪР° Р»РёСҶР°, СҖРөСӮСғСҲСҢ, outpaint.\n"
        "вҖў вңЁ РһР¶РёРІР»РөРҪРёРө С„РҫСӮРҫ: Р·Р°РіСҖСғР·Рё С„РҫСӮРҫ вҖ” РјРҫР¶РҪРҫ РІСӢРұСҖР°СӮСҢ Runway, Kling РёР»Рё Sora 2 СӮРҫР»СҢРәРҫ РҙР»СҸ РәР°РҙСҖРҫРІ РұРөР· Р»СҺРҙРөР№.\n"
        "вҖў рҹҺһ Р’РёРҙРөРҫ: СҖР°Р·РұРҫСҖ СҒРјСӢСҒР»Р°, СӮР°Р№РјРәРҫРҙСӢ, *Reels РёР· РҙР»РёРҪРҪРҫРіРҫ РІРёРҙРөРҫ*, РёРҙРөРё/СҒРәСҖРёРҝСӮ, СҒСғРұСӮРёСӮСҖСӢ.\n"
        "вҖў рҹҺ§ РҗСғРҙРёРҫ/РәРҪРёРіРё: СӮСҖР°РҪСҒРәСҖРёРҝСҶРёСҸ, СӮРөР·РёСҒСӢ, РҝР»Р°РҪ.\n"
        "вҖў рҹ©ә РңРөРҙРёСҶРёРҪР°: РІСӢРҝРёСҒРәРё, Р°РҪР°РјРҪРөР·, Р·Р°РәР»СҺСҮРөРҪРёСҸ, Р°РҪР°Р»РёР·СӢ, СҒРҪРёРјРәРё, РңР Рў/РҡРў вҖ” СҒРҝСҖР°РІРҫСҮРҪСӢР№ СҖР°Р·РұРҫСҖ Рё РІРҫРҝСҖРҫСҒСӢ РІСҖР°СҮСғ.\n\n"
        "_РҹРҫРҙСҒРәР°Р·РәРё:_ РҝСҖРҫСҒСӮРҫ Р·Р°РіСҖСғР·РёСӮРө С„Р°Р№Р» РёР»Рё РҝСҖРёСҲР»РёСӮРө СҒСҒСӢР»РәСғ + РәРҫСҖРҫСӮРәРҫРө РўР—. "
        "Р”Р»СҸ С„РҫСӮРҫ вҖ” РјРҫР¶РҪРҫ РҪР°Р¶Р°СӮСҢ В«вңЁ РһР¶РёРІРёСӮСҢВ», РҙР»СҸ РІРёРҙРөРҫ вҖ” В«рҹҺ¬ Reels РёР· РҙР»РёРҪРҪРҫРіРҫ РІРёРҙРөРҫВ»."
    )
    await update.effective_message.reply_text(msg, parse_mode="Markdown", reply_markup=_fun_quick_kb())

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ Р’СҒРҝРҫРјРҫРіР°СӮРөР»СҢРҪРҫРө: РІР·СҸСӮСҢ РҝРөСҖРІСғСҺ РҫРұСҠСҸРІР»РөРҪРҪСғСҺ С„СғРҪРәСҶРёСҺ РҝРҫ РёРјРөРҪРё в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def _pick_first_defined(*names):
    for n in names:
        fn = globals().get(n)
        if callable(fn):
            return fn
    return None

# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ Telegram profile setup в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
async def _post_init_bot_profile(app):
    """РһРұРҪРҫРІР»СҸРөСӮ РҪР°Р·РІР°РҪРёРө/РҫРҝРёСҒР°РҪРёРө РұРҫСӮР° РІ Telegram РҝРөСҖРөРҙ /start, РөСҒР»Рё РІРәР»СҺСҮРөРҪРҫ AUTO_SET_BOT_PROFILE."""
    if not AUTO_SET_BOT_PROFILE:
        return
    try:
        if BOT_PUBLIC_NAME:
            await app.bot.set_my_name(name=BOT_PUBLIC_NAME)
        if BOT_SHORT_DESCRIPTION:
            await app.bot.set_my_short_description(short_description=BOT_SHORT_DESCRIPTION)
        if BOT_DESCRIPTION:
            await app.bot.set_my_description(description=BOT_DESCRIPTION)
        log.info("Telegram bot profile updated: %s", BOT_PUBLIC_NAME)
    except Exception as e:
        log.warning("Telegram bot profile update skipped: %s", e)


# в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ Р РөРіРёСҒСӮСҖР°СҶРёСҸ С…РөРҪРҙР»РөСҖРҫРІ Рё Р·Р°РҝСғСҒРә в”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җв”Җ
def build_application() -> "Application":
    if not BOT_TOKEN:
        raise RuntimeError("РқРө Р·Р°РҙР°РҪ BOT_TOKEN РІ РҝРөСҖРөРјРөРҪРҪСӢС… РҫРәСҖСғР¶РөРҪРёСҸ.")

    builder = ApplicationBuilder().token(BOT_TOKEN)
    if AUTO_SET_BOT_PROFILE:
        builder = builder.post_init(_post_init_bot_profile)
    app = builder.build()

    # РҡРҫРјР°РҪРҙСӢ
    app.add_handler(CommandHandler("start",        cmd_start))
    app.add_handler(CommandHandler("help",         cmd_help))
    app.add_handler(CommandHandler("examples",     cmd_examples))
    app.add_handler(CommandHandler("version",      cmd_version))
    app.add_handler(CommandHandler("engines",      cmd_engines))
    app.add_handler(CommandHandler("plans",        cmd_plans))
    app.add_handler(CommandHandler("balance",      cmd_balance))
    app.add_handler(CommandHandler("prices",       cmd_prices))
    app.add_handler(CommandHandler("set_welcome",  cmd_set_welcome))
    app.add_handler(CommandHandler("show_welcome", cmd_show_welcome))
    app.add_handler(CommandHandler("diag_limits",  cmd_diag_limits))
    app.add_handler(CommandHandler("diag_access",  cmd_diag_access))
    app.add_handler(CommandHandler("diag_stt",     cmd_diag_stt))
    app.add_handler(CommandHandler("diag_images",  cmd_diag_images))
    app.add_handler(CommandHandler("diag_video",   cmd_diag_video))
    app.add_handler(CommandHandler("diag_runway",  cmd_diag_runway))
    app.add_handler(CommandHandler("diag_yookassa", cmd_diag_yookassa))
    app.add_handler(CommandHandler("provider_status", cmd_provider_status))
    app.add_handler(CommandHandler("diag_bg",      cmd_diag_bg))
    app.add_handler(CommandHandler("diag_face",    cmd_diag_face))
    app.add_handler(CommandHandler("diag_suno",    cmd_diag_suno))
    app.add_handler(CommandHandler("img",          cmd_img))
    app.add_handler(CommandHandler("mj",           cmd_midjourney))
    app.add_handler(CommandHandler("midjourney",   cmd_midjourney))
    app.add_handler(CommandHandler("presentation", cmd_presentation))
    app.add_handler(CommandHandler("catalog",      cmd_catalog))
    app.add_handler(CommandHandler("diag_presentation", cmd_diag_presentation))
    app.add_handler(CommandHandler("chats",        cmd_chats))
    app.add_handler(CommandHandler("newchat",      cmd_newchat))
    app.add_handler(CommandHandler("music",       cmd_music))
    app.add_handler(CommandHandler("voice_on",     cmd_voice_on))
    app.add_handler(CommandHandler("voice_off",    cmd_voice_off))
    app.add_handler(CommandHandler("medicine",     cmd_mode_medicine))

    # РҹР»Р°СӮРөР¶Рё
    app.add_handler(PreCheckoutQueryHandler(on_precheckout))
    app.add_handler(MessageHandler(filters.SUCCESSFUL_PAYMENT, on_successful_payment))

    # >>> PATCH START вҖ” Handlers wiring (WebApp + callbacks + media + text) >>>

    # Р”Р°РҪРҪСӢРө РёР· РјРёРҪРё-РҝСҖРёР»РҫР¶РөРҪРёСҸ (WebApp)
    with contextlib.suppress(Exception):
        app.add_handler(MessageHandler(filters.StatusUpdate.WEB_APP_DATA, on_webapp_data))
    with contextlib.suppress(Exception):
        if hasattr(filters, "WEB_APP_DATA"):
            app.add_handler(MessageHandler(filters.WEB_APP_DATA, on_webapp_data))

    # === РҹРҗРўР§ 4: РҹРҫСҖСҸРҙРҫРә callback-С…РөРҪРҙР»РөСҖРҫРІ (СғР·РәРёРө вҶ’ РҫРұСүРёРө) ===
    # 1) РҹРҫРҙРҝРёСҒРәР°/РҫРҝР»Р°СӮСӢ
    app.add_handler(CallbackQueryHandler(on_cb_plans, pattern=r"^(?:plan:|pay:)$|^(?:plan:|pay:).+"))

    # 2) РқРҫРІСӢРө СҖРөР¶РёРјСӢ/РҝРҫРҙРјРөРҪСҺ: mode:* Рё act:* (РЈСҮС‘РұР°/Р Р°РұРҫСӮР°/Р Р°Р·РІР»РөСҮРөРҪРёСҸ/РңРөРҙРёСҶРёРҪР°)
    app.add_handler(CallbackQueryHandler(on_mode_cb, pattern=r"^(?:mode:|act:)"), group=0)

    # 2b) РЎСӮР°СҖСӢРө school:/work: callbacks, РөСҒР»Рё СӮР°РәРёРө РәРҪРҫРҝРәРё РөСүС‘ РіРҙРө-СӮРҫ РёСҒРҝРҫР»СҢР·СғСҺСӮСҒСҸ
    app.add_handler(CallbackQueryHandler(on_cb_mode, pattern=r"^(?:school:|work:)"), group=0)

    # 3) Р‘СӢСҒСӮСҖСӢРө СҖР°Р·РІР»РөСҮРөРҪРёСҸ (Р»СҺРұСӢРө fun:...)
    app.add_handler(CallbackQueryHandler(on_cb_fun,   pattern=r"^fun:[a-z_]+$"))

    # 3b) РҹРҫРҙРјРөРҪСҺ Suno: СҒРІРҫРұРҫРҙРҪСӢР№ Р·Р°РҝСҖРҫСҒ Рё РҝСҖРөСҒРөСӮСӢ
    app.add_handler(CallbackQueryHandler(on_cb_suno,  pattern=r"^suno:"), group=0)

    # 4) РһСҒСӮР°Р»СҢРҪРҫР№ catch-all (pedit/topup/engine/buy Рё СӮ.Рҝ.)
    # Р Р°Р·РјРөСүР°РөРј РІ РҝСҖРёРҫСҖРёСӮРөСӮРҪРҫР№ РіСҖСғРҝРҝРө, СҮСӮРҫРұСӢ РәРҫР»РұСҚРәРё РҫРұСҖР°РұР°СӮСӢРІР°Р»РёСҒСҢ СҒСҖР°Р·Сғ
    app.add_handler(CallbackQueryHandler(on_cb), group=0)

    # Р“РҫР»РҫСҒ/Р°СғРҙРёРҫ вҖ” РҫСӮРҪРҫСҒРёРј Рә РјРөРҙРёР°РіСҖСғРҝРҝРө (РёРҙС‘СӮ СҖР°РҪСҢСҲРө РҫРұСүРөРіРҫ СӮРөРәСҒСӮРҫРІРҫРіРҫ С…РөРҪРҙР»РөСҖР°)
    voice_fn = _pick_first_defined("handle_voice", "on_voice", "voice_handler")
    if voice_fn:
        app.add_handler(MessageHandler(filters.VOICE | filters.AUDIO, voice_fn), group=1)

    # РўРөРәСҒСӮРҫРІСӢРө РәРҪРҫРҝРәРё/СҸСҖР»СӢРәРё (РҫСҒСӮР°Р»СҢРҪСӢРө) вҖ” Р§РҳРЎРўРһ РұРөР· РҙСғРұР»РөР№
    import re

    # РЎСӮСҖРҫРіРёРө РҝР°СӮСӮРөСҖРҪСӢ: РҫРҙРҪРҫ РҪР°Р·РІР°РҪРёРө = РҫРҙРёРҪ С…РөРҪРҙР»РөСҖ (СҚРјРҫРҙР·Рё РҙРҫРҝСғСҒРәР°РөРј, Р»РёСҲРҪРёРө РҝСҖРҫРұРөР»СӢ вҖ” СӮРҫР¶Рө)
    BTN_ENGINES = re.compile(r"^\s*(?:рҹ§ \s*)?Р”РІРёР¶РәРё\s*$")
    BTN_BALANCE = re.compile(r"^\s*(?:рҹ’і|рҹ§ҫ)?\s*Р‘Р°Р»Р°РҪСҒ\s*$")
    BTN_PLANS   = re.compile(r"^\s*(?:вӯҗ\s*)?РҹРҫРҙРҝРёСҒРәР°(?:\s*[В·вҖў]\s*РҹРҫРјРҫСүСҢ)?\s*$")
    BTN_STUDY   = re.compile(r"^\s*(?:рҹҺ“\s*)?РЈСҮ[РөС‘]РұР°\s*$")
    BTN_WORK    = re.compile(r"^\s*(?:рҹ’ј\s*)?(?:Р Р°РұРҫСӮР°(?:\s*/\s*Р‘РёР·РҪРөСҒ)?|Р‘РёР·РҪРөСҒ)\s*$")
    BTN_FUN     = re.compile(r"^\s*(?:рҹ”Ҙ\s*)?Р Р°Р·РІР»РөСҮРөРҪРёСҸ\s*$")
    BTN_MED     = re.compile(r"^\s*(?:рҹ©ә|вҡ•пёҸ)?\s*РңРөРҙРёСҶРёРҪР°\s*$")
    BTN_CHATS   = re.compile(r"^\s*(?:рҹ’¬\s*)?РңРҫРё СҮР°СӮСӢ\s*$", re.I)
    BTN_NEWCHAT = re.compile(r"^\s*(?:вһ•\s*)?РқРҫРІСӢР№ СҮР°СӮ\s*$", re.I)

    # РҡРҪРҫРҝРәРё РІ РҝСҖРёРҫСҖРёСӮРөСӮРҪРҫР№ РіСҖСғРҝРҝРө (0), СҮСӮРҫРұСӢ РҫРҪРё СҒСҖР°РұР°СӮСӢРІР°Р»Рё СҖР°РҪСҢСҲРө Р»СҺРұСӢС… РҫРұСүРёС… РҫРұСҖР°РұРҫСӮСҮРёРәРҫРІ
    app.add_handler(MessageHandler(filters.Regex(BTN_ENGINES), on_btn_engines), group=0)
    app.add_handler(MessageHandler(filters.Regex(BTN_BALANCE), on_btn_balance), group=0)
    app.add_handler(MessageHandler(filters.Regex(BTN_PLANS),   on_btn_plans),   group=0)
    app.add_handler(MessageHandler(filters.Regex(BTN_STUDY),   on_btn_study),   group=0)
    app.add_handler(MessageHandler(filters.Regex(BTN_WORK),    on_btn_work),    group=0)
    app.add_handler(MessageHandler(filters.Regex(BTN_FUN),     on_btn_fun),     group=0)
    app.add_handler(MessageHandler(filters.Regex(BTN_MED),     on_btn_medicine), group=0)
    app.add_handler(MessageHandler(filters.Regex(BTN_CHATS),   cmd_chats), group=0)
    app.add_handler(MessageHandler(filters.Regex(BTN_NEWCHAT), cmd_newchat), group=0)

    # Р–С‘СҒСӮРәРёР№ РҝРөСҖРөС…РІР°СӮ В«РјРҫР¶РөСҲСҢ РҫР¶РёРІРёСӮСҢ С„РҫСӮРҫ?В» вҖ” РҙРҫ Р»СҺРұРҫРіРҫ GPT-РҫСӮРІРөСӮР°.
    # Р’РҗР–РқРһ: Р·РҙРөСҒСҢ РҪР°РјРөСҖРөРҪРҪРҫ РҪРөСӮ inline-С„Р»Р°РіРҫРІ РІРёРҙР° (?is), СҮСӮРҫРұСӢ Render/Python 3.12 РҪРө РҝР°РҙР°Р».
    photo_revive_capability_re = re.compile(
        r"(РјРҫР¶(РөСҲСҢ|РөСӮРө|РҪРҫ)|СғРјРө(РөСҲСҢ|РөСӮРө)|РјРҫР¶РөСӮ\s+Р»Рё|СҒРҝРҫСҒРҫРұРөРҪ|РҝРҫРҙРҙРөСҖР¶РёРІР°РөСҲСҢ|РҝРҫР»СғСҮРёСӮСҒСҸ|РҙРөР»Р°РөСҲСҢ)"
        r".{0,160}(РҫР¶РёРІ|Р°РҪРёРјРёСҖ|revive|animate)"
        r".{0,160}(С„РҫСӮРҫ|С„РҫСӮРҫРіСҖР°С„|РәР°СҖСӮРёРҪРә|РёР·РҫРұСҖР°Р¶РөРҪ|photo|image|picture)"
        r"|"
        r"(РҫР¶РёРІ|Р°РҪРёРјРёСҖ|revive|animate)"
        r".{0,160}(С„РҫСӮРҫ|С„РҫСӮРҫРіСҖР°С„|РәР°СҖСӮРёРҪРә|РёР·РҫРұСҖР°Р¶РөРҪ|photo|image|picture)"
        r".{0,80}\?",
        re.I | re.S,
    )
    app.add_handler(
        MessageHandler(filters.Regex(photo_revive_capability_re), on_photo_revival_capability),
        group=0,
    )
    # вһ• РҹРҫР·РёСӮРёРІРҪСӢР№ Р°РІСӮРҫ-РҫСӮРІРөСӮ РҪР° В«Р° СғРјРөРөСҲСҢ Р»РёвҖҰВ» вҖ” РҙРҫ РҫРұСүРөРіРҫ СӮРөРәСҒСӮР° (РҫСӮРҙРөР»СҢРҪР°СҸ РіСҖСғРҝРҝР°, РҪРёР¶Рө РәРҪРҫРҝРҫРә)
    app.add_handler(MessageHandler(filters.Regex(_CAPS_PATTERN), on_capabilities_qa), group=1)

    # РңРөРҙРёР° (С„РҫСӮРҫ/РҙРҫРәРё/РІРёРҙРөРҫ/РіРёС„) вҖ” СӮРҫР¶Рө РҝРөСҖРөРҙ РҫРұСүРёРј СӮРөРәСҒСӮРҫРј
    photo_fn = _pick_first_defined("handle_photo", "on_photo", "photo_handler", "handle_image_message")
    if photo_fn:
        app.add_handler(MessageHandler(filters.PHOTO, photo_fn), group=1)

    doc_fn = _pick_first_defined("handle_doc", "on_doc", "on_document", "handle_document", "doc_handler")
    if doc_fn:
        app.add_handler(MessageHandler(filters.Document.ALL, doc_fn), group=1)

    video_fn = _pick_first_defined("handle_video", "on_video", "video_handler")
    if video_fn:
        app.add_handler(MessageHandler(filters.VIDEO, video_fn), group=1)

    gif_fn = _pick_first_defined("handle_gif", "on_gif", "animation_handler")
    if gif_fn:
        app.add_handler(MessageHandler(filters.ANIMATION, gif_fn), group=1)

    # >>> PATCH END <<<

    # РһРұСүРёР№ СӮРөРәСҒСӮ вҖ” РЎРҗРңР«Рҷ РҝРҫСҒР»РөРҙРҪРёР№ (РҪРёР¶Рө РІСҒРөС… СҮР°СҒСӮРҪСӢС… РәРөР№СҒРҫРІ)
    text_fn = _pick_first_defined("handle_text", "on_text", "text_handler", "default_text_handler")
    if text_fn:
        btn_filters = (filters.Regex(BTN_ENGINES) | filters.Regex(BTN_BALANCE) |
                       filters.Regex(BTN_PLANS)   | filters.Regex(BTN_STUDY)   |
                       filters.Regex(BTN_WORK)    | filters.Regex(BTN_FUN) |
                       filters.Regex(BTN_MED)     | filters.Regex(BTN_CHATS) |
                       filters.Regex(BTN_NEWCHAT))
        app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND & ~btn_filters, text_fn), group=2)

    # РһСҲРёРұРәРё
    err_fn = _pick_first_defined("on_error", "handle_error")
    if err_fn:
        app.add_error_handler(err_fn)

    return app


# === main() СҒ РұРөР·РҫРҝР°СҒРҪРҫР№ РёРҪРёСҶРёР°Р»РёР·Р°СҶРёРөР№ Р‘Р” (РұРөР· РёР·РјРөРҪРөРҪРёР№ РҝРҫ СҒСғСӮРё) ===
def main():
    log.info("Starting bot patch version: %s", PATCH_VERSION)
    with contextlib.suppress(Exception):
        db_init()
    with contextlib.suppress(Exception):
        db_init_usage()
    with contextlib.suppress(Exception):
        _chat_memory_init()
    with contextlib.suppress(Exception):
        _db_init_prefs()

    app = build_application()

    if USE_WEBHOOK:
        _install_webapp_checkout_bridge(app)
        log.info("рҹҡҖ WEBHOOK mode. Public URL: %s  Path: %s  Port: %s", PUBLIC_URL, WEBHOOK_PATH, PORT)
        app.run_webhook(
            listen="0.0.0.0",
            port=PORT,
            url_path=WEBHOOK_PATH.lstrip("/"),
            webhook_url=f"{PUBLIC_URL.rstrip('/')}{WEBHOOK_PATH}",
            secret_token=(WEBHOOK_SECRET or None),
            allowed_updates=Update.ALL_TYPES,
        )
    else:
        log.info("рҹҡҖ POLLING mode.")
        with contextlib.suppress(Exception):
            asyncio.get_event_loop().run_until_complete(
                app.bot.delete_webhook(drop_pending_updates=True)
            )
        app.run_polling(
            close_loop=False,
            allowed_updates=Update.ALL_TYPES,
            drop_pending_updates=False,
        )


if __name__ == "__main__":
    main()
# === END PATCH ===
